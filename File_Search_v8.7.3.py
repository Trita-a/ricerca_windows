import os
import sys  # Aggiungo l'import di sys
import shutil
import ttkbootstrap as ttk
import tkinter as tk
from ttkbootstrap.constants import *
from tkinter import filedialog, messagebox, BooleanVar, StringVar, IntVar
import threading
import queue
from datetime import datetime
import getpass
import zipfile
from ttkbootstrap.dialogs import Querybox
import traceback
import time
import concurrent.futures
import mimetypes
import signal
import subprocess
import zipfile
import io
import csv

# Dizionario per tracciare il supporto alle librerie
file_format_support = {
    "docx": False,
    "pdf": False,
    "pptx": False,
    "excel": False,
    "odt": False,
    "rtf": False,
    "xls": False
}

# Elenco di librerie da installare se mancanti
missing_libraries = []

# Importazione di librerie opzionali con gestione degli errori
try:
    import docx
    file_format_support["docx"] = True
    print("Supporto Word (.docx) attivato")
except ImportError:
    missing_libraries.append("python-docx")
    print("Supporto Word (.docx) non disponibile")

try:
    import PyPDF2
    file_format_support["pdf"] = True
    print("Supporto PDF attivato")
except ImportError:
    missing_libraries.append("PyPDF2")
    print("Supporto PDF non disponibile")

try:
    import pptx
    file_format_support["pptx"] = True
    print("Supporto PowerPoint attivato")
except ImportError:
    missing_libraries.append("python-pptx")
    print("Supporto PowerPoint non disponibile")

try:
    import openpyxl
    file_format_support["excel"] = True
    print("Supporto Excel attivato")
except ImportError:
    missing_libraries.append("openpyxl")
    print("Supporto Excel non disponibile")

try:
    from striprtf.striprtf import rtf_to_text
    file_format_support["rtf"] = True
    print("Supporto RTF attivato")
except ImportError:
    missing_libraries.append("striprtf")
    print("Supporto RTF non disponibile")

try:
    import odfdo
    file_format_support["odt"] = True
    print("Supporto ODT attivato (libreria odfdo)")
except ImportError:
    try:
        from odf import opendocument, text, teletype
        file_format_support["odt"] = True
        print("Supporto ODT attivato (libreria odf)")
    except ImportError:
        missing_libraries.append("odfdo")
        print("Supporto ODT non disponibile")
try:
    import win32com.client
    file_format_support["doc"] = True
    print("Supporto Word (.doc) attivato tramite pywin32")
except ImportError:
    missing_libraries.append("pywin32")
    print("Supporto Word (.doc) non disponibile")

try:
    import win32com.client
    file_format_support["xls"] = True
    print("Supporto Excel (.xls) attivato tramite pywin32")
except ImportError:
    missing_libraries.append("pywin32")
    print("Supporto Excel (.xls) non disponibile")
try:
    import xlrd
    file_format_support["xls_native"] = True
    print("Supporto Excel (.xls) attivato tramite xlrd")
except ImportError:
    missing_libraries.append("xlrd")
    print("Supporto Excel (.xls) non disponibile via xlrd")

class FileSearchApp:
    def __init__(self, root):
        self.root = root
        self.root.title("File Search Tool V8.7.2. Nucleo Perugia")
        self.root.geometry("1300x850")
        
        # Inizializza le variabili per data/ora e utente
        self.datetime_var = StringVar()
        self.user_var = StringVar(value=getpass.getuser())
        
        # Variabili
        self.search_content = BooleanVar(value=True)
        self.search_path = StringVar()
        self.keywords = StringVar()
        self.search_results = []
        self.search_files = BooleanVar(value=True)
        self.search_folders = BooleanVar(value=True)
        self.ignore_hidden = BooleanVar(value=True)
        self.is_searching = False
        self.progress_queue = queue.Queue()
        self.search_executor = None
        self.exclude_system_files = BooleanVar(value=True)  # Per default escludiamo i file di sistema
        self.whole_word_search = BooleanVar(value=False)

        # Variabili per la ricerca a blocchi
        self.max_files_per_block = IntVar(value=1000)  # Numero massimo di file per blocco
        self.max_parallel_blocks = IntVar(value=4)  # Numero massimo di blocchi da elaborare in parallelo
        self.prioritize_user_folders = BooleanVar(value=True)  # Dare priorità alle cartelle utente
        self.block_size_auto_adjust = BooleanVar(value=True)  # Adatta automaticamente la dimensione dei blocchi
        
        # Variables for disk information
        self.dir_size_var = StringVar(value="")
        self.total_disk_var = StringVar(value="")
        self.used_disk_var = StringVar(value="")
        self.free_disk_var = StringVar(value="")

        # Variabili aggiuntive per i miglioramenti
        self.search_start_time = None
        self.stop_search = False
        self.max_depth = 0  # 0 = illimitato
        self.last_search_params = {}
        self.search_history = []
        self.advanced_filters = {
            "size_min": 0,
            "size_max": 0,
            "date_min": None,
            "date_max": None,
            "extensions": []
        }
        # Opzioni per la gestione dei permessi
        self.skip_permission_errors = BooleanVar(value=True)  # Salta silenziosamente errori di permesso
        self.excluded_paths = []  # Lista di percorsi da escludere dalla ricerca

        # Directory problematiche da escludere automaticamente
        self.problematic_dirs = [
            "Client Active Directory Rights Management Services",
            "Windows Resource Protection",
            "Windows Defender",
            "Microsoft Office",
            "System Volume Information",
            "$Recycle.Bin",
            "$WINDOWS.~BT",
            "$Windows.~WS"
        ]

        # Verifica dei privilegi di amministratore
        self.is_admin = False
        try:
            import ctypes
            self.is_admin = ctypes.windll.shell32.IsUserAnAdmin() != 0
        except:
            self.is_admin = False

        # Variabili per limiti di tempo e altre ottimizzazioni
        self.timeout_enabled = BooleanVar(value=False)
        self.timeout_seconds = IntVar(value=3600)  # Default 60 secondi
        self.max_files_to_check = IntVar(value=100000)  # Limite numero di file da controllare
        self.max_results = IntVar(value=50000)  # Limite numero di risultati
        self.chunk_size = 8192  # Grandezza del chunk per la lettura dei file
        self.max_file_size_mb = IntVar(value=100)  # Dimensione massima file da analizzare (in MB)
        self.worker_threads = IntVar(value=min(8, os.cpu_count() or 4))  # Numero di worker threads per elaborazione parallela
        self.use_indexing = BooleanVar(value=True)  # Usa indicizzazione per velocizzare ricerche future
        self.search_index = {}  # Dizionario per indicizzare file e contenuti
        
        # Lista di estensioni di file di sistema da escludere dalla ricerca nei contenuti
        self.system_file_extensions = [
            # File eseguibili e librerie
            '.exe', '.dll', '.sys', '.drv', '.ocx', '.vxd', '.com', '.bat', '.cmd', '.scr', '.app', '.dylib', '.exp', '.bpl',
            # File di configurazione
            '.ini', '.inf', '.reg', '.config', '.msi', '.msp', '.cab', '.plist', '.yaml', '.toml', '.rc', '.pid', '.lock',
            # File di cache e temporanei
            '.tmp', '.temp', '.dmp', '.cache', '.bak', '.swp', '.log', '.thumbs', '.crdownload', '.part', '.ds_store', '.desktop',
            # File di sistema Windows
            '.lnk', '.ico', '.cur', '.ani', '.cpl', '.evt', '.cat', '.mui', '.msstyles', '.gadget', '.theme', '.manifest',
            # File di sistema Linux/macOS
            '.so', '.o', '.a', '.la', '.ko', '.mod', '.rc.local', '.sh', '.service', '.desktop',
            # File di log e diagnostica
            '.log', '.dmp', '.etl', '.evtx', '.trace', '.stackdump',
            # Altri file binari comuni
            '.bin', '.dat', '.db', '.sqlite', '.sqlite3', '.qvm', '.pak', '.idx', '.wim', '.vmdk', '.qcow2',
            # File specifici di sistema
            '.efi', '.pf', '.etl', '.mui', '.res', '.tlb', '.pdb', '.bak', '.old', '.sav'
        ]

        # Info utente e datetime
        self.current_user = getpass.getuser()
        self.datetime_var = StringVar()
        self.update_datetime()
        self.create_widgets()

        # Verifica le librerie e mostra notifica se necessario
        self.check_and_notify_missing_libraries()
        self.debug_mode = True
        
        # Registra handler per CTRL+C
        self.register_interrupt_handler()
    
        # Imposta il tema iniziale
        self.update_theme_colors("dark")  # O "dark" se il tema predefinito è scuro
        
        # Aggiungi variabile per il percorso del log dei file saltati
        self.skipped_files_log_path = os.path.join(os.path.expanduser("~"), "skipped_files_log.txt")
    def manage_memory(self):
        """ Gestisce l'utilizzo della memoria durante la ricerca per evitare problemi"""
        import psutil
        import gc
        
        process = psutil.Process(os.getpid())
        memory_usage = process.memory_info().rss / 1024 / 1024  # MB
        
        # Se l'uso della memoria è troppo alto, esegui la garbage collection forzata
        if memory_usage > 500:  # 500 MB
            self.log_debug(f"Utilizzo memoria elevato: {memory_usage:.2f} MB. Esecuzione garbage collection.")
            gc.collect()
            
            # Dopo la GC, ricalcola l'utilizzo della memoria
            memory_usage = process.memory_info().rss / 1024 / 1024
            self.log_debug(f"Utilizzo memoria dopo GC: {memory_usage:.2f} MB")    

    def run_with_timeout(self, func, args=(), kwargs={}, timeout_sec=10):
        """ Esegue una funzione con un timeout, evitando blocchi indefiniti Returns: (result, completed)
            - result: risultato della funzione o None
            - completed: True se completata, False se interrotta per timeout
        """
        result = [None]
        completed = [False]
        
        def target():
            try:
                result[0] = func(*args, **kwargs)
            except Exception as e:
                self.log_debug(f"Eccezione nella funzione con timeout: {str(e)}")
            finally:
                completed[0] = True
        
        thread = threading.Thread(target=target)
        thread.daemon = True
        thread.start()
        thread.join(timeout_sec)
        
        return result[0], completed[0]
    
    def _get_all_descendants(self, widget):
        """Ottiene ricorsivamente tutti i widget discendenti"""
        descendants = []
        for child in widget.winfo_children():
            descendants.append(child)
            descendants.extend(self._get_all_descendants(child))
        return descendants
    
    def disable_all_controls(self):
        """Disabilita tutti i controlli UI durante la ricerca"""
        try:
            # Input e campi di testo
            self.path_entry["state"] = "disabled"
            self.keyword_entry["state"] = "disabled"
            
            # Pulsante Sfoglia
            if hasattr(self, 'browse_btn'):
                self.browse_btn["state"] = "disabled"
            
            # Pulsante Pulisci campi
            if hasattr(self, 'clear_btn'):
                self.clear_btn["state"] = "disabled"
            
            # Combobox e spinbox
            if hasattr(self, 'theme_combobox'):
                self.theme_combobox["state"] = "disabled"
            
            if hasattr(self, 'depth_spinbox'):
                self.depth_spinbox["state"] = "disabled"
            
            # Pulsanti principali
            if hasattr(self, 'search_button'):
                self.search_button["state"] = "disabled"
            
            # Gestisci i pulsanti dei filtri
            for btn_name in ['advanced_filters_btn', 'exclusions_btn', 'block_options_btn']:
                if hasattr(self, btn_name):
                    getattr(self, btn_name)["state"] = "disabled"
            
            # Gestisci specificamente il pulsante admin
            if hasattr(self, 'admin_button') and not self.is_admin:
                self.admin_button["state"] = "disabled"
            
            # Checkbox - ora usa hasattr per sicurezza
            if hasattr(self, 'perf_frame'):
                try:
                    for widget in self._get_all_descendants(self.perf_frame):
                        if isinstance(widget, ttk.Spinbox) or isinstance(widget, ttk.Checkbutton):
                            widget["state"] = "disabled"
                except Exception as e:
                    self.log_debug(f"Errore nella disabilitazione dei controlli in perf_frame: {str(e)}")
            
            # Pulsanti di azione sui risultati
            if hasattr(self, 'copy_button'):
                self.copy_button["state"] = "disabled"
            if hasattr(self, 'compress_button'):
                self.compress_button["state"] = "disabled"
                
        except Exception as e:
            # Registra l'errore senza interrompere l'esecuzione
            self.log_debug(f"Errore nella disabilitazione dei controlli: {str(e)}")

    def enable_all_controls(self):
        """Riabilita tutti i controlli UI dopo la ricerca"""
        # Input e campi di testo
        self.path_entry["state"] = "normal"
        self.keyword_entry["state"] = "normal"
        
        # Pulsante Sfoglia
        if hasattr(self, 'browse_btn'):
            self.browse_btn["state"] = "normal"
        
        # Pulsante Pulisci campi
        if hasattr(self, 'clear_btn'):
            self.clear_btn["state"] = "normal"
        
        # Combobox e spinbox
        if hasattr(self, 'theme_combobox'):
            self.theme_combobox["state"] = "readonly"
        self.depth_spinbox["state"] = "normal"
        
        # Pulsanti principali
        self.search_button["state"] = "normal"
        
        # Gestisci i pulsanti dei filtri
        for btn_name in ['advanced_filters_btn', 'exclusions_btn']:
            if hasattr(self, btn_name):
                getattr(self, btn_name)["state"] = "normal"

        # Gestisci i pulsanti  a blocchi
        if hasattr(self, 'block_options_btn'):
            self.block_options_btn["state"] = "normal"

        # Gestisci specificamente il pulsante admin
        if hasattr(self, 'admin_button') and not self.is_admin:
            self.admin_button["state"] = "normal"
        
        # Checkbox
        for widget in self.root.winfo_children():
            self._enable_checkbuttons_recursive(widget)
        
        # Pulsanti di azione sui risultati
        if hasattr(self, 'copy_button'):
            self.copy_button["state"] = "normal"
        if hasattr(self, 'compress_button'):
            self.compress_button["state"] = "normal"
        
        # Spinbox nelle opzioni di performance - usa il nuovo metodo helper
        for widget in self._get_all_descendants(self.perf_frame):
            if isinstance(widget, ttk.Spinbox) or isinstance(widget, ttk.Checkbutton):
                widget["state"] = "normal"
                
    def _disable_checkbuttons_recursive(self, widget):
        """Disabilita ricorsivamente tutte le checkbox nei widget"""
        if isinstance(widget, ttk.Checkbutton):
            widget["state"] = "disabled"
        
        # Processa i widget figli ricorsivamente
        if hasattr(widget, 'winfo_children'):
            for child in widget.winfo_children():
                self._disable_checkbuttons_recursive(child)
                
    def _enable_checkbuttons_recursive(self, widget):
        """Riabilita ricorsivamente tutte le checkbox nei widget"""
        if isinstance(widget, ttk.Checkbutton):
            widget["state"] = "normal"
        
        # Processa i widget figli ricorsivamente
        if hasattr(widget, 'winfo_children'):
            for child in widget.winfo_children():
                self._enable_checkbuttons_recursive(child)
    def show_content_search_warning(self):
        """Mostra un avviso se la ricerca nei contenuti è attivata"""
        if self.search_content.get():  # Se la checkbox è attivata
            return messagebox.askyesno(
                "Attenzione - Ricerca nei contenuti", 
                "Hai attivato la ricerca nei contenuti dei file.\n\n"
                "Questa operazione può richiedere molto più tempo, soprattutto con grandi quantità di file.\n\n"
                "Vuoi procedere con la ricerca nei contenuti?",
                icon="warning"
            )
        return True  # Se la ricerca nei contenuti non è attivata, procedi senza avviso    
    
    def register_interrupt_handler(self):
        """Registra il gestore degli interrupt (CTRL+C)"""
        def handle_interrupt(sig, frame):
            if self.is_searching:
                self.stop_search_process()
            else:
                self.root.quit()
        
        signal.signal(signal.SIGINT, handle_interrupt)

    def log_debug(self, message):
        """Funzione per logging, stampa solo quando debug_mode è True"""
        if self.debug_mode:
            print(f"[DEBUG] {message}")

    def check_and_notify_missing_libraries(self):
        """Verifica e notifica l'utente di eventuali librerie mancanti"""
        missing = []
        
        if not file_format_support["docx"]:
            missing.append("python-docx (per file Word)")
        if not file_format_support["pdf"]: 
            missing.append("PyPDF2 (per file PDF)")
        if not file_format_support["pptx"]:
            missing.append("python-pptx (per file PowerPoint)")
        if not file_format_support["excel"]:
            missing.append("openpyxl (per file Excel)")
        if not file_format_support["rtf"]:
            missing.append("striprtf (per file RTF)")
        if not file_format_support["odt"]:
            missing.append("odfdo (per file OpenDocument)")
        
        if missing:
            message = "Alcune funzionalità di ricerca nei contenuti sono disabilitate.\n\n"
            message += "Per abilitare il supporto completo ai vari formati di file, installa le seguenti librerie:\n\n"
            
            for lib in missing:
                message += f"- {lib}\n"
            
            message += "\nPuoi installarle con il comando:\n"
            message += "pip install " + " ".join([lib.split(" ")[0] for lib in missing])
            
            # Mostra la notifica dopo un breve ritardo per permettere all'UI di caricarsi
            self.root.after(1000, lambda: messagebox.showinfo("Librerie opzionali mancanti", message))

    def update_datetime(self):
        current_time = datetime.now().strftime('%d-%m-%Y %H:%M:%S')
        self.datetime_var.set(f"Data: {current_time} | Utente: {self.current_user}")
        self.root.after(1000, self.update_datetime)

    def browse_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.search_path.set(directory)
            self.update_disk_info(directory)

    def optimize_system_search(self, path):
        """Ottimizza la ricerca per percorsi di sistema come C:/ impostando parametri appropriati"""
        is_system_path = path.lower() in ["c:/", "c:\\", "d:/", "d:\\"]
        
        # Verifica se è un percorso di rete o server (inizia con \\ o //)
        is_network_path = path.startswith('\\\\') or path.startswith('//')
        
        # Verifica se è un percorso Unix root o home
        is_unix_root = path in ['/', '/home', '/usr', '/var', '/etc']
        
        if is_system_path or is_network_path or is_unix_root:
            # Salva i parametri attuali
            original_params = {
                "max_files": self.max_files_to_check.get(),
                "worker_threads": self.worker_threads.get(),
                "timeout": self.timeout_seconds.get() if self.timeout_enabled.get() else None
            }
            
            # Applica parametri ottimizzati
            self.max_files_to_check.set(1000000)  # Un milione di file
            self.worker_threads.set(min(12, os.cpu_count() or 4))  # Aumenta thread per server
            
            if self.timeout_enabled.get():
                self.timeout_seconds.set(max(3600, self.timeout_seconds.get()))  # Minimo 1 ora
            
            # Notifica l'utente
            messagebox.showinfo(
                "Ricerca su sistema o server",
                "Stai avviando una ricerca su un percorso di sistema o server.\n\n"
                "Per ottimizzare la ricerca, sono stati temporaneamente modificati alcuni parametri "
                "per consentire una scansione più approfondita.\n\n"
                "La ricerca potrebbe richiedere tempo e risorse significative."
            )
            
            return original_params
        return None
    def show_optimization_tips(self, path):
        """Mostra suggerimenti di ottimizzazione quando si cerca in un'unità di sistema"""
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]:
            optimization_done = False
            
            # Verifica se l'utente ha già implementato le ottimizzazioni
            if hasattr(self, 'excluded_paths'):
                windows_excluded = any("windows" in p.lower() for p in self.excluded_paths)
                program_files_excluded = any("program files (x86)" in p.lower() for p in self.excluded_paths)
                optimization_done = windows_excluded and program_files_excluded
            
            # Verifica la profondità di ricerca
            depth_optimized = self.max_depth >= 1 and self.max_depth <= 10
            
            # Se le ottimizzazioni non sono già state applicate, mostra il messaggio
            if not (optimization_done and depth_optimized):
                response = messagebox.askyesno(
                    "Ottimizza la ricerca",
                    "Stai per avviare una ricerca sull'intero disco C:\\\n\n"
                    "Per migliorare notevolmente le prestazioni, è consigliato:\n\n"
                    "1. Escludere le cartelle di sistema (Windows, Program Files)\n"
                    "2. Escludere le cartelle di altri utenti\n"
                    "3. Limitare la profondità di ricerca a 5-10 livelli\n\n"
                    "Vuoi applicare automaticamente queste ottimizzazioni?",
                    icon="question"
                )
                
                if response:
                    # Applica le ottimizzazioni
                    if not hasattr(self, 'excluded_paths'):
                        self.excluded_paths = []
                    
                    # Aggiungi esclusioni di sistema
                    system_paths = ["C:/Windows", "C:/Program Files", "C:/Program Files (x86)"]
                    for sys_path in system_paths:
                        if sys_path not in self.excluded_paths:
                            self.excluded_paths.append(sys_path)
                    
                    # Escludi altri utenti
                    current_user = getpass.getuser()
                    users_dir = "C:/Users"
                    if os.path.exists(users_dir):
                        try:
                            for user in os.listdir(users_dir):
                                user_path = os.path.join(users_dir, user)
                                if user.lower() != current_user.lower() and user.lower() not in ["public", "default", "all users"]:
                                    if user_path not in self.excluded_paths:
                                        self.excluded_paths.append(user_path)
                        except Exception as e:
                            self.log_debug(f"Errore nell'esclusione degli altri utenti: {str(e)}")
                    
                    # Imposta la profondità di ricerca a 7 (valore ragionevole)
                    self.depth_spinbox.set("7")
                    self.max_depth = 7
                    
                    # Notifica l'utente
                    messagebox.showinfo(
                        "Ottimizzazioni applicate",
                        f"Le ottimizzazioni sono state applicate:\n\n"
                        f"• Escluse {len(system_paths)} cartelle di sistema\n"
                        f"• Escluse le cartelle di altri utenti\n"
                        f"• Profondità di ricerca impostata a 7 livelli\n\n"
                        f"La ricerca sarà ora più veloce e stabile."
                    )
                    
                    return True
                else:
                    # NUOVA PARTE: Se l'utente risponde "No", imposta esplicitamente la profondità a 0
                    self.depth_spinbox.set("0")
                    self.max_depth = 0
                    
                    # Notifica l'utente che sta per iniziare una ricerca illimitata
                    self.log_debug("Utente ha rifiutato le ottimizzazioni, impostata profondità illimitata (0)")
                    
                    # Opzionalmente, puoi mostrare un messaggio di avviso
                    messagebox.showinfo(
                        "Ricerca illimitata",
                        "Stai per avviare una ricerca illimitata sull'intero disco.\n\n"
                        "La profondità di ricerca è stata impostata a illimitata (0).\n"
                        "Questa operazione potrebbe richiedere molto tempo e risorse di sistema."
                    )
                
                return False
            
            return False  # Per percorsi non di sistema, non fare nulla
    def debug_exclusions(self):
        """Visualizza lo stato corrente delle esclusioni per debug"""
        if hasattr(self, 'excluded_paths'):
            paths = '\n'.join(self.excluded_paths) if self.excluded_paths else "Nessun percorso escluso"
            messagebox.showinfo("Debug esclusioni", 
                            f"Stato esclusioni:\n{paths}\n\n"
                            f"Totale percorsi: {len(self.excluded_paths)}")
        else:
            messagebox.showinfo("Debug esclusioni", "Lista esclusioni non inizializzata")

    def manage_exclusions(self):
        """Apre una finestra di dialogo per gestire i percorsi esclusi dalla ricerca"""
        if not hasattr(self, 'excluded_paths'):
            self.excluded_paths = []
        dialog = ttk.Toplevel(self.root)
        dialog.title("Gestione percorsi esclusi")
        dialog.geometry("500x400")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Istruzioni
        ttk.Label(main_frame, text="Aggiungi cartelle da escludere dalla ricerca (es. C:/Windows)", 
                wraplength=480).pack(anchor=W, pady=(0, 10))
        
        # Frame per aggiungere nuovi percorsi
        add_frame = ttk.Frame(main_frame)
        add_frame.pack(fill=X, pady=5)
        
        path_var = StringVar()
        entry = ttk.Entry(add_frame, textvariable=path_var, width=40)
        entry.pack(side=LEFT, padx=(0, 5), fill=X, expand=YES)
        
        def browse_exclude_dir():
            directory = filedialog.askdirectory()
            if directory:
                path_var.set(directory)
        
        def add_exclusion():
            path = path_var.get().strip()
            if path and path not in self.excluded_paths:
                self.excluded_paths.append(path)
                update_list()
                path_var.set("")
        
        ttk.Button(add_frame, text="Sfoglia", command=browse_exclude_dir).pack(side=LEFT, padx=2)
        ttk.Button(add_frame, text="Aggiungi", command=add_exclusion).pack(side=LEFT, padx=2)
        
        # Lista dei percorsi esclusi
        ttk.Label(main_frame, text="Percorsi attualmente esclusi:").pack(anchor=W, pady=(10, 5))
        
        list_frame = ttk.Frame(main_frame)
        list_frame.pack(fill=BOTH, expand=YES, pady=5)
        
        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side=RIGHT, fill=Y)
        
        excluded_list = ttk.Treeview(list_frame, columns=("path",), show="headings", 
                                yscrollcommand=scrollbar.set, selectmode="extended")
        excluded_list.heading("path", text="Percorso")
        excluded_list.column("path", width=450)
        excluded_list.pack(fill=BOTH, expand=YES)
        
        scrollbar.config(command=excluded_list.yview)
        
        def update_list():
            # Pulisci la lista
            for item in excluded_list.get_children():
                excluded_list.delete(item)
            
            # Aggiungi i percorsi esclusi
            for path in self.excluded_paths:
                excluded_list.insert("", "end", values=(path,))
        
        def remove_selected():
            selected = excluded_list.selection()
            if selected:
                for item in selected:
                    values = excluded_list.item(item)["values"]
                    if values and values[0] in self.excluded_paths:
                        self.excluded_paths.remove(values[0])
                update_list()
        
        # Pulsanti per gestire la lista
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=10)
        
        ttk.Button(btn_frame, text="Rimuovi selezionati", command=remove_selected).pack(side=LEFT)
        
        # Preset di esclusioni comuni
        def add_common_exclusions():
            common_paths = [
                "C:/Windows", 
                "C:/Program Files", 
                "C:/Program Files (x86)",
                "C:/ProgramData",
                "C:/Users/All Users",
                "C:/Program Files (x86)/Client Active Directory Rights Management Services"
            ]
            added = 0
            for path in common_paths:
                if path not in self.excluded_paths:
                    self.excluded_paths.append(path)
                    added += 1
            
            if added > 0:
                update_list()
                messagebox.showinfo("Aggiunto", f"Aggiunti {added} percorsi comuni alle esclusioni")
        
        ttk.Button(btn_frame, text="Aggiungi esclusioni comuni", 
                command=add_common_exclusions).pack(side=LEFT, padx=5)
        
        # Aggiungi automaticamente altre cartelle utente
        def exclude_other_users():
            current_user = getpass.getuser()
            users_dir = "C:/Users"
            added = 0
            
            if os.path.exists(users_dir):
                try:
                    for user in os.listdir(users_dir):
                        user_path = os.path.join(users_dir, user)
                        # Escludi tutte le cartelle utente tranne quella dell'utente corrente
                        if user.lower() != current_user.lower() and user.lower() not in ["public", "default", "all users"]:
                            if user_path not in self.excluded_paths:
                                self.excluded_paths.append(user_path)
                                added += 1
                except Exception as e:
                    self.log_debug(f"Errore nell'esclusione degli altri utenti: {str(e)}")
            
            if added > 0:
                update_list()
                messagebox.showinfo("Aggiunto", f"Escluse {added} cartelle di altri utenti")
        
        ttk.Button(btn_frame, text="Escludi altri utenti", 
                command=exclude_other_users).pack(side=LEFT, padx=5)
        
        # Pulsanti OK/Annulla
        dialog_btn_frame = ttk.Frame(main_frame)
        dialog_btn_frame.pack(fill=X, pady=(10, 0))
        
        ttk.Button(dialog_btn_frame, text="OK", command=dialog.destroy).pack(side=RIGHT)
        
        # Inizializza la lista
        update_list()

        # Aggiorna la finestra per ottenere le dimensioni corrette
        dialog.update_idletasks()
        
        # Ottieni le dimensioni della finestra
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        
        # Calcola la posizione x,y per centrare la finestra
        screen_width = dialog.winfo_screenwidth()
        screen_height = dialog.winfo_screenheight()
        x = (screen_width // 2) - (width // 1)
        y = (screen_height // 2) - (height // 1)
        
        # Imposta la geometria con la posizione calcolata
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima per la finestra
        dialog.minsize(800, 600)

    def restart_as_admin(self):
        """Riavvia l'applicazione con privilegi di amministratore"""
        try:
            import ctypes, sys, os
            if sys.platform == 'win32':
                script = sys.argv[0]
                params = ' '.join(sys.argv[1:])
                self.log_debug("Tentativo di riavvio come amministratore")
                
                # Se il percorso di ricerca è impostato, passa quel percorso come parametro
                if self.search_path.get():
                    if not params:
                        params = f'"{self.search_path.get()}"'
                    
                # Avvia l'applicazione con privilegi elevati
                ctypes.windll.shell32.ShellExecuteW(None, "runas", sys.executable, f'"{script}" {params}', None, 1)
                
                # Chiude l'istanza corrente
                self.root.quit()
            else:
                messagebox.showinfo("Informazione", "Questa funzione è disponibile solo su Windows")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile riavviare come amministratore: {str(e)}")
            self.log_debug(f"Errore nel riavvio come admin: {str(e)}")

    def optimize_disk_search_order(self, path, directories):
        """Ottimizza l'ordine di esplorazione delle directory quando si cerca in un disco di sistema.
        Dà priorità alle cartelle più importanti come Users, Documents, ecc.
        """
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"]:
            prioritized = []
            normal = []
            
            # Cartelle ad alta priorità (esplora per prime)
            high_priority_folders = ["users", "documents", "documents and settings", "desktop", "downloads"]
            
            # Cartelle a bassa priorità (esplora per ultime)
            low_priority_folders = ["windows", "program files", "program files (x86)", "programdata", "$recycle.bin", "system volume information"]
            
            for dir_path in directories:
                dir_name = os.path.basename(dir_path).lower()
                
                if any(priority in dir_name for priority in high_priority_folders):
                    prioritized.insert(0, dir_path)  # Aggiungi all'inizio per priorità più alta
                elif any(low in dir_name for low in low_priority_folders):
                    normal.append(dir_path)  # Aggiungi alla fine
                else:
                    prioritized.append(dir_path)  # Priorità media
                    
            return prioritized + normal
        
        # Se non è un percorso di sistema, restituisci l'elenco originale
        return directories

    def start_search(self):
        # Pulisci risultati precedenti
        for item in self.results_list.get_children():
            self.results_list.delete(item)
    
        if not self.search_path.get() or not self.keywords.get():
            messagebox.showerror("Errore", "Inserisci directory e parole chiave")
            return
        
        # Mostra avviso se la ricerca nei contenuti è attivata
        if not self.show_content_search_warning():
            return  # Interrompi se l'utente annulla
        
        # Aggiorna le informazioni del disco qui, quando l'utente clicca su cerca
        self.update_disk_info()

        self.show_optimization_tips(self.search_path.get())
        
        # Reset per la nuova ricerca
        self.stop_search = False
        self.stop_button["state"] = "normal"
        original_params = self.optimize_system_search(self.search_path.get())
        
        # Aggiorna l'orario di avvio e resetta l'orario di fine
        current_time = datetime.now().strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")

        # Salva i parametri di ricerca
        self.last_search_params = {
            "path": self.search_path.get(),
            "keywords": self.keywords.get(),
            "search_files": self.search_files.get(),
            "search_folders": self.search_folders.get(),
            "search_content": self.search_content.get()
        }
        
        # Aggiungi alla cronologia di ricerca se non già presente
        if self.keywords.get() and self.keywords.get() not in [h["keywords"] for h in self.search_history]:
            self.search_history.append(self.last_search_params.copy())
            if len(self.search_history) > 10:  # Mantieni solo le ultime 10 ricerche
                self.search_history.pop(0)
        
        self.is_searching = True
        self.disable_all_controls()  # Disabilita tutti i controlli
        self.stop_button["state"] = "normal"  # Solo questo è abilitato durante la ricerca
        self.progress_bar["value"] = 0
        self.status_label["text"] = "Ricerca in corso..."
        
        # Aggiorna il valore della profondità massima
        try:
            self.max_depth = int(self.depth_spinbox.get())
        except ValueError:
            self.max_depth = 0  # Valore predefinito se non valido
        
        # Ottieni le parole chiave di ricerca
        search_terms = [term.strip() for term in self.keywords.get().split(',') if term.strip()]
            
        # Avvia la ricerca in un thread separato
        self.search_results = []  # Resetta i risultati
        search_thread = threading.Thread(target=self._search_thread, 
                                        args=(self.search_path.get(), search_terms, self.search_content.get()))
        search_thread.daemon = True
        search_thread.start()
        
        # Memorizza l'orario di inizio come oggetto datetime
        self.search_start_time = datetime.now()
        current_time = self.search_start_time.strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")
        self.total_time_label.config(text="--:--")  # Reset del tempo totale

        # Avvia l'aggiornamento della progress bar
        self.update_progress()
        # Salva original_params come attributo per ripristinarli dopo la ricerca
        if original_params:
            self.original_system_search_params = original_params

        current_theme = "dark"  # Sostituisci con la logica per determinare il tema corrente
        self.update_theme_colors(current_theme)

        # Aggiungi questo nuovo metodo per calcolare il tempo totale
    def update_total_time(self):
        """Calcola e visualizza il tempo totale trascorso tra inizio e fine ricerca"""
        if self.search_start_time:
            end_time = datetime.now()
            time_diff = end_time - self.search_start_time
            
            # Converti in minuti e secondi
            total_seconds = int(time_diff.total_seconds())
            minutes = total_seconds // 60
            seconds = total_seconds % 60
            
            # Formatta la stringa del tempo totale
            if minutes > 0:
                total_time_str = f"{minutes}min {seconds}sec"
            else:
                total_time_str = f"{seconds}sec"
            
            self.total_time_label.config(text=total_time_str) 
    # Esempio di utilizzo per cambiare il tema
    def change_theme(self, theme):
        self.update_theme_colors(theme)

    def _search_thread(self, path, keywords, search_content):
        try:
            # Inizializza i contatori di file e directory esaminati e il tempo di inizio
            files_checked = 0
            dirs_checked = 0
            start_time = time.time()
            timeout = self.timeout_seconds.get() if self.timeout_enabled.get() else None
            
            # Determina se si tratta di una ricerca completa del sistema (C:/ o simile)
            is_system_search = path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]

            # Per ricerche di sistema, adatta automaticamente i parametri
            if is_system_search:
                # Informa l'utente che la ricerca potrebbe richiedere molto tempo
                self.progress_queue.put(("status", "Ricerca completa del sistema in corso - potrebbe richiedere molto tempo"))
                
                # Temporaneamente aumenta i limiti per la ricerca di sistema
                original_max_files = self.max_files_to_check.get()
                original_timeout = timeout
                
                # Imposta limiti più elevati per la ricerca completa
                self.max_files_to_check.set(5000000)  # 5 milioni di file
                
                # Disabilita temporaneamente il timeout o aumentalo significativamente
                if timeout and timeout < 3600:
                    timeout = 3600 * 8  # 8 ore
                
                # Avviso all'utente
                self.progress_queue.put(("status", "Ricerca completa avviata - parametri adattati per ricerca approfondita"))
            
            # Variabili per gestire l'aggiornamento ogni secondo
            last_update_time = time.time()
            
            # Crea un executor per processare i file in parallelo in modo sicuro
            try:
                max_workers = max(1, min(32, self.worker_threads.get()))  # Garantisce un valore valido
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
            except Exception as e:
                self.log_debug(f"Errore nella creazione del ThreadPoolExecutor: {str(e)}")
                # Fallback a un valore sicuro
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
            
            futures = []
            
            # Insieme per tenere traccia delle cartelle visitate (evita loop infiniti con symlink)
            visited_dirs = set()
            
            # NUOVO: coda di priorità per i blocchi di ricerca
            block_queue = queue.PriorityQueue()

            # Funzione per elaborare un file
            def process_file(file_path, keywords):
                if self.stop_search:
                    return None
                    
                # Salta direttamente i file problematici
                if self.should_skip_file(file_path):
                    return None
                
                try:
                    # Implementazione timeout cross-platform tramite threading
                    result = [None]
                    exception = [None]
                    processing_completed = [False]
                    
                    # INSERISCI QUI LA FUNZIONE process_with_timeout()
                    def process_with_timeout():
                        try:
                            # Verifica filtri di dimensione
                            file_size = os.path.getsize(file_path)
                            if (self.advanced_filters["size_min"] > 0 and file_size < self.advanced_filters["size_min"]) or \
                            (self.advanced_filters["size_max"] > 0 and file_size > self.advanced_filters["size_max"]):
                                result[0] = None
                                return
                            
                            # Verifica filtri di data
                            if self.advanced_filters["date_min"] or self.advanced_filters["date_max"]:
                                mod_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                                
                                if self.advanced_filters["date_min"]:
                                    min_date = datetime.strptime(self.advanced_filters["date_min"], "%d-%m-%Y")
                                    if mod_time < min_date:
                                        result[0] = None
                                        return
                                        
                                if self.advanced_filters["date_max"]:
                                    max_date = datetime.strptime(self.advanced_filters["date_max"], "%d-%m-%Y")
                                    if mod_time > max_date:
                                        result[0] = None
                                        return
                            
                            # Verifica filtri estensione
                            if self.advanced_filters["extensions"] and not any(file_path.lower().endswith(ext.lower()) 
                                                                        for ext in self.advanced_filters["extensions"]):
                                result[0] = None
                                return
                                
                            if os.path.splitext(file_path)[1].lower() in ['.doc', '.xls']:
                                self.progress_queue.put(("progress", self.progress_bar["value"]))  # Forza aggiornamento
                                
                            # Verifica corrispondenza nel nome file
                            filename = os.path.basename(file_path)

                            matched = False
                            for keyword in keywords:
                                # Verifica se la ricerca di parole intere è attivata
                                if self.whole_word_search.get() and ' ' not in keyword:
                                    try:
                                        import re
                                        pattern = r'\b' + re.escape(keyword.lower()) + r'\b'
                                        if re.search(pattern, filename.lower()):
                                            matched = True
                                            break
                                    except:
                                        # Fallback in caso di errore con regex
                                        if keyword.lower() in filename.lower():
                                            matched = True
                                            break
                                # Se è un termine con spazi (contesto specifico)
                                elif ' ' in keyword:
                                    if keyword.lower() in filename.lower():
                                        matched = True
                                        break
                                # Ricerca normale
                                else:
                                    if keyword.lower() in filename.lower():
                                        matched = True
                                        break

                            if matched:
                                result[0] = self.create_file_info(file_path)
                                return
                            
                            # Verifica contenuto se richiesto
                            if search_content and self.should_search_content(file_path):
                                max_size_bytes = self.max_file_size_mb.get() * 1024 * 1024
                                
                                # Salta file troppo grandi
                                if file_size > max_size_bytes:
                                    self.log_debug(f"File {file_path} troppo grande per l'analisi del contenuto")
                                    result[0] = None
                                    return
                                    
                                content = self.get_file_content(file_path)
                                if content:
                                    matched = False
                                    for keyword in keywords:
                                        # Verifica se la ricerca di parole intere è attivata
                                        if self.whole_word_search.get() and ' ' not in keyword:
                                            try:
                                                import re
                                                pattern = r'\b' + re.escape(keyword.lower()) + r'\b'
                                                if re.search(pattern, content.lower()):
                                                    matched = True
                                                    break
                                            except:
                                                # Fallback in caso di errore con regex
                                                if keyword.lower() in content.lower():
                                                    matched = True
                                                    break
                                        # Se è un termine con spazi (contesto specifico)
                                        elif ' ' in keyword:
                                            if keyword.lower() in content.lower():
                                                matched = True
                                                break
                                        # Ricerca normale
                                        else:
                                            if keyword.lower() in content.lower():
                                                matched = True
                                                break
                                                
                                    if matched:
                                        result[0] = self.create_file_info(file_path)
                                        return
                            else:
                                # Log per i file di sistema esclusi
                                if search_content and os.path.splitext(file_path)[1].lower() in self.system_file_extensions:
                                    self.log_debug(f"File di sistema escluso dall'analisi del contenuto: {file_path}")
                                    
                            result[0] = None
                            
                        except Exception as e:
                            exception[0] = e
                            self.log_debug(f"Errore durante l'elaborazione del file {file_path}: {str(e)}")
                            result[0] = None
                        finally:
                            processing_completed[0] = True
                    
                    # Avvia il thread con un nome identificativo
                    import threading
                    worker_thread = threading.Thread(
                        target=process_with_timeout, 
                        name=f"Process-{os.path.basename(file_path)}"
                    )
                    worker_thread.daemon = True
                    worker_thread.start()
                    
                    # Attendi il completamento del thread con timeout ridotto
                    if os.path.splitext(file_path)[1].lower() in ['.doc', '.xls']:
                        # Timeout ridotto per file problematici
                        timeout_seconds = 5
                        self.log_debug(f"Utilizzo timeout ridotto (5s) per il file: {file_path}")
                    else:
                        timeout_seconds = 20
                    worker_thread.join(timeout_seconds)
                    
                    # Verifica se il thread è ancora in esecuzione (timeout raggiunto)
                    if not processing_completed[0]:
                        self.log_debug(f"Timeout nella elaborazione del file {file_path}")
                        # Non attendere più il thread, continua semplicemente
                        return None
                        
                    # Verifica se si è verificata un'eccezione
                    if exception[0]:
                        self.log_debug(f"Eccezione nell'elaborazione del file {file_path}: {str(exception[0])}")
                        return None
                        
                    # Restituisci il risultato dal thread
                    return result[0]
                    
                except Exception as e:
                    self.log_debug(f"Errore generale nella elaborazione del file {file_path}: {str(e)}")
                    return None
            
            # NUOVO: calcola la priorità del blocco (numerica, più bassa = più alta priorità)
            def calculate_block_priority(directory):
                # Se l'opzione è disabilitata, usa priorità standard per tutti
                if not self.prioritize_user_folders.get():
                    return 1
                    
                # Altrimenti, applica la prioritizzazione configurata
                # Dai priorità alle cartelle degli utenti
                if "users" in directory.lower() or "documenti" in directory.lower() or "desktop" in directory.lower():
                    return 0  # Alta priorità
                # Dai priorità medio-alta alle cartelle di dati
                elif "data" in directory.lower() or "database" in directory.lower() or "downloads" in directory.lower():
                    return 1  # Priorità media-alta
                # Dai priorità bassa a cartelle di sistema
                elif "windows" in directory.lower() or "program files" in directory.lower():
                    return 3  # Priorità bassa
                # Priorità standard per altre cartelle
                return 2
            
            # NUOVO: Inizializza la coda con blocchi iniziali
            def initialize_block_queue(root_path):
                try:
                    items = os.listdir(root_path)
                    
                    # Prima aggiungi le directory come blocchi separati
                    for item in items:
                        item_path = os.path.join(root_path, item)
                        try:
                            if os.path.isdir(item_path):
                                # Salta directory nascoste se richiesto
                                if self.ignore_hidden.get() and (item.startswith('.') or 
                                    (os.name == 'nt' and os.path.exists(item_path) and 
                                    os.stat(item_path).st_file_attributes & 2)):
                                    continue
                                
                                # Salta directory escluse
                                if hasattr(self, 'excluded_paths') and any(
                                    item_path.lower().startswith(excluded.lower()) 
                                    for excluded in self.excluded_paths):
                                    continue
                                    
                                # Calcola priorità e aggiungi alla coda
                                priority = calculate_block_priority(item_path)
                                block_queue.put((priority, item_path))
                                visited_dirs.add(os.path.realpath(item_path))
                        except Exception as e:
                            self.log_debug(f"Errore nell'aggiunta del blocco {item_path}: {str(e)}")
                    
                    # Poi processa i file nella directory principale
                    for item in items:
                        item_path = os.path.join(root_path, item)
                        try:
                            if not os.path.isdir(item_path):
                                if self.search_files.get():
                                    # Verifica file nascosti
                                    if self.ignore_hidden.get() and (item.startswith('.') or 
                                        (os.name == 'nt' and os.path.exists(item_path) and 
                                        os.stat(item_path).st_file_attributes & 2)):
                                        continue
                                    
                                    # Processa direttamente i file nella directory principale
                                    nonlocal files_checked
                                    files_checked += 1
                                    if files_checked > self.max_files_to_check.get():
                                        self.stop_search = True
                                        return
                                    
                                    future = self.search_executor.submit(process_file, item_path, keywords)
                                    futures.append(future)
                        except Exception as e:
                            self.log_debug(f"Errore nell'elaborazione del file {item_path}: {str(e)}")
                            
                except PermissionError:
                    self.log_debug(f"Permesso negato per la directory {root_path}")
                except Exception as e:
                    self.log_debug(f"Errore nell'inizializzazione dei blocchi da {root_path}: {str(e)}")
            
            # MODIFICATO: Ricerca breadth-first basata su blocchi
            def process_blocks():
                nonlocal files_checked, dirs_checked, last_update_time
    
                # Determina il numero massimo di file per blocco in base alle impostazioni
                max_files_in_block = self.max_files_per_block.get()
                
                # Adatta automaticamente la dimensione del blocco se richiesto
                if self.block_size_auto_adjust.get():
                    # Ottimizzazione per ricerche di sistema o cartelle molto grandi
                    if is_system_search:
                        max_files_in_block = max(2000, max_files_in_block)
                    elif files_checked > 100000:
                        max_files_in_block = max(5000, max_files_in_block)
                
                # Ottimizza il numero di blocchi paralleli in base al carico di sistema
                active_blocks = 0
                max_parallel = self.max_parallel_blocks.get()
                
                # Utilizziamo un set per tracciare i blocchi già processati più efficiente
                processed_blocks = set()
                
                while not block_queue.empty() and not self.stop_search:
                    # Verifica timeout
                    if timeout and time.time() - start_time > timeout:
                        self.progress_queue.put(("timeout", "Timeout raggiunto"))
                        return
                    
                    try:
                        # Prendi il blocco con priorità più alta
                        _, current_block = block_queue.get(block=False)
                        
                        # Salta blocchi già processati
                        if current_block in processed_blocks:
                            continue
                            
                        processed_blocks.add(current_block)
                        
                        # Aggiorna lo stato
                        current_time = time.time()
                        if current_time - last_update_time >= 0.5:  # Aggiorna ogni mezzo secondo
                            elapsed_time = current_time - start_time
                            self.progress_queue.put(("status", 
                                f"Analisi blocco: {current_block} (Cartelle: {dirs_checked}, File: {files_checked}, Tempo: {int(elapsed_time)}s)"))
                            self.progress_queue.put(("progress", 
                                min(90, int((files_checked / max(1, self.max_files_to_check.get())) * 100))))
                            last_update_time = current_time
                        
                        # Aggiungi blocchi di livello inferiore dalla directory corrente
                        try:
                            # Blocco attualmente in elaborazione
                            dirs_checked += 1
                            
                            # Lista elementi nella directory
                            items = os.listdir(current_block)
                            
                            # Prima processa le sottocartelle (aggiungi nuovi blocchi)
                            subfolders = []
                            for item in items:
                                if self.stop_search:
                                    return
                                    
                                item_path = os.path.join(current_block, item)
                                
                                # Salta file/cartelle nascoste
                                try:
                                    if self.ignore_hidden.get() and (item.startswith('.') or 
                                        (os.name == 'nt' and os.path.exists(item_path) and 
                                        os.stat(item_path).st_file_attributes & 2)):
                                        continue
                                except Exception as e:
                                    self.log_debug(f"Errore nel controllo hidden per {item_path}: {str(e)}")
                                    continue
                                    
                                # Gestione sottodirectory
                                try:
                                    if os.path.isdir(item_path):
                                        # Verifica se la cartella è già stata visitata
                                        try:
                                            real_path = os.path.realpath(item_path)
                                            if real_path in visited_dirs:
                                                continue
                                            visited_dirs.add(real_path)
                                        except:
                                            if item_path in visited_dirs:
                                                continue
                                            visited_dirs.add(item_path)
                                        
                                        # Verifica se il percorso deve essere escluso
                                        if hasattr(self, 'excluded_paths') and any(
                                            item_path.lower().startswith(excluded.lower()) 
                                            for excluded in self.excluded_paths):
                                            continue
                                        
                                        # Aggiungi alla lista delle sottocartelle
                                        subfolders.append(item_path)
                                        
                                        # Verifica corrispondenza nome cartella
                                        if self.search_folders.get():
                                            if any(keyword.lower() in item.lower() for keyword in keywords):
                                                folder_info = self.create_folder_info(item_path)
                                                self.search_results.append(folder_info)
                                except Exception as e:
                                    self.log_debug(f"Errore nell'analisi della directory {item_path}: {str(e)}")
                            
                            # Aggiunta sottocartelle alla coda con priorità calcolata
                            for subfolder in subfolders:
                                # Verifica limite di profondità
                                folder_depth = subfolder.count(os.path.sep) - path.count(os.path.sep)
                                if self.max_depth > 0 and folder_depth > self.max_depth:
                                    continue
                                    
                                priority = calculate_block_priority(subfolder)
                                block_queue.put((priority, subfolder))
                            
                            # Processa i file nella directory corrente
                            for item in items:
                                if self.stop_search:
                                    return
                                    
                                item_path = os.path.join(current_block, item)
                                
                                # Salta i file nascosti e le directory (già processate)
                                try:
                                    # Controllo file nascosto
                                    if self.ignore_hidden.get() and (item.startswith('.') or 
                                                (os.name == 'nt' and os.path.exists(item_path) and 
                                                    os.stat(item_path).st_file_attributes & 2)):
                                        continue
                                        
                                    # Salta le directory (già processate)
                                    if os.path.isdir(item_path):
                                        continue
                                except Exception as e:
                                    self.log_debug(f"Errore nel controllo del tipo per {item_path}: {str(e)}")
                                    continue
                                    
                                # Processa i file
                                if self.search_files.get():
                                    try:
                                        # Verifica limite file
                                        files_checked += 1
                                        if files_checked % 1000 == 0:
                                            self.manage_memory()
                                        if files_checked > self.max_files_to_check.get():
                                            self.stop_search = True
                                            self.progress_queue.put(("status", 
                                                f"Limite di {self.max_files_to_check.get()} file controllati raggiunto. "
                                                f"Aumenta il limite nelle opzioni per cercare più file."))
                                            return
                                        
                                        # Gestione sicura dell'executor
                                        try:
                                            if self.search_executor and not self.search_executor._shutdown:
                                                future = self.search_executor.submit(process_file, item_path, keywords)
                                                futures.append(future)
                                            else:
                                                result = process_file(item_path, keywords)
                                                if result:
                                                    self.search_results.append(result)
                                        except Exception as e:
                                            self.log_debug(f"Errore nell'elaborazione parallela del file {item_path}: {str(e)}")
                                            try:
                                                result = process_file(item_path, keywords)
                                                if result:
                                                    self.search_results.append(result)
                                            except:
                                                pass
                                            
                                    except Exception as e:
                                        self.log_debug(f"Errore nell'aggiunta del file {item_path} alla coda: {str(e)}")
                                        continue
                            
                        except PermissionError:
                            if self.skip_permission_errors.get():
                                self.log_debug(f"Saltata directory con permesso negato: {current_block}")
                            else:
                                dir_name = os.path.basename(current_block)
                                parent_dir = os.path.dirname(current_block)
                                
                                is_user_folder = (parent_dir.lower() in ["c:/users", "c:\\users"] and 
                                                dir_name.lower() != getpass.getuser().lower())
                                
                                if is_user_folder:
                                    self.log_debug(f"Cartella di un altro utente inaccessibile: {current_block}")
                                    self.progress_queue.put(("status", f"Saltata cartella utente protetta: {current_block}"))
                                else:
                                    self.log_debug(f"Permesso negato per la directory {current_block}")
                                    self.progress_queue.put(("status", f"Permesso negato: {current_block}"))
                        except Exception as e:
                            self.log_debug(f"Errore durante l'analisi della directory {current_block}: {str(e)}")
                            
                    except queue.Empty:
                        break

            # Aggiorna lo stato iniziale
            self.progress_queue.put(("status", f"Inizio ricerca a blocchi in: {path} (Profondità: {'illimitata' if self.max_depth == 0 else self.max_depth})"))
            
            # Assicurati che la lista di esclusioni sia inizializzata
            if not hasattr(self, 'excluded_paths'):
                self.excluded_paths = []
                
            # Avvia la ricerca a blocchi
            initialize_block_queue(path)
            process_blocks()
            
            # Ripristina i parametri originali se erano stati modificati
            if is_system_search and self.max_depth == 0 and 'original_max_files' in locals():
                self.max_files_to_check.set(original_max_files)
            
            # Aggiorna lo stato finale di analisi
            self.progress_queue.put(("status", f"Elaborazione risultati... (analizzati {files_checked} file in {dirs_checked} cartelle)"))
            
            # MODIFICA: Raccolta risultati dalle future con aggiornamento temporizzato
            completed = 0
            total_futures = len(futures)
            last_update_time = time.time()
            
            # Gestione sicura delle future
            if self.search_executor and not self.search_executor._shutdown and futures:
                try:
                    for future in concurrent.futures.as_completed(futures):
                        if self.stop_search:
                            break
                            
                        try:
                            result = future.result()
                            if result:
                                self.search_results.append(result)
                            
                            completed += 1
                            
                            # Aggiorna il progresso e il tempo più frequentemente
                            current_time = time.time()
                            if completed % 20 == 0 or current_time - last_update_time > 2:  # Aggiorna ogni 20 file o ogni 2 secondi
                                progress = 90 + min(10, int((completed / max(1, len(futures))) * 10))
                                self.progress_queue.put(("progress", progress))
                                
                                elapsed_time = current_time - start_time
                                self.progress_queue.put(("status", f"Elaborati {completed}/{len(futures)} file (tempo: {int(elapsed_time)}s)"))
                                
                                # Aggiorna anche il tempo totale durante la ricerca
                                if hasattr(self, 'search_start_time'):
                                    now = datetime.now()
                                    time_diff = now - self.search_start_time
                                    # Converti in minuti e secondi
                                    total_seconds = int(time_diff.total_seconds())
                                    minutes = total_seconds // 60
                                    seconds = total_seconds % 60
                                    
                                    if minutes > 0:
                                        total_time_str = f"{minutes}min {seconds}sec"
                                    else:
                                        total_time_str = f"{seconds}sec"
                                        
                                    self.progress_queue.put(("update_total_time", total_time_str))
                                
                                last_update_time = current_time
                        except Exception as e:
                            self.log_debug(f"Errore nell'elaborazione di un risultato: {str(e)}")
                except Exception as e:
                    self.log_debug(f"Errore nella raccolta dei risultati: {str(e)}")
            
            # Completa la ricerca in modo sicuro
            try:
                if self.search_executor and not self.search_executor._shutdown:
                    self.search_executor.shutdown(wait=False)
                self.search_executor = None
            except Exception as e:
                self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                self.search_executor = None
            
            # Ripristina i parametri originali se erano stati modificati per la ricerca di sistema
            if hasattr(self, 'original_system_search_params'):
                try:
                    self.max_files_to_check.set(self.original_system_search_params["max_files"])
                    self.worker_threads.set(self.original_system_search_params["worker_threads"])
                    delattr(self, 'original_system_search_params')
                except:
                    pass

            # Riporta il risultato finale
            elapsed_time = time.time() - start_time
            self.progress_queue.put(("status", 
                f"Ricerca completata! Analizzati {files_checked} file in {dirs_checked} cartelle in {int(elapsed_time)} secondi."))
            
            # Ordina i risultati per tipo e nome
            self.search_results.sort(key=lambda x: (x[0], x[1]))
            
            self.log_debug(f"Ricerca completata. Trovati {len(self.search_results)} risultati")
            self.progress_queue.put(("complete", "Ricerca completata"))
            
        except Exception as e:
            error_msg = f"Si è verificato un errore durante la ricerca: {str(e)}\n{traceback.format_exc()}"
            self.log_debug(error_msg)
            self.progress_queue.put(("error", error_msg))

    def search_current_user_only(self):
        """Imposta la ricerca solo nella cartella dell'utente corrente"""
        user_folder = os.path.join("C:/Users", getpass.getuser())
        if os.path.exists(user_folder):
            self.search_path.set(user_folder)
            
            # Mostra la conferma solo se l'utente non annulla il warning sulla ricerca nei contenuti
            if self.show_content_search_warning():
                messagebox.showinfo(
                    "Ricerca configurata", 
                    f"La ricerca è stata configurata per cercare solo nella tua cartella utente:\n{user_folder}\n\n"
                    "Questa modalità evita problemi di permesso con altre cartelle protette."
                )
                # Avvia direttamente la ricerca
                self.start_search()
        else:
            messagebox.showerror("Errore", "Impossibile trovare la tua cartella utente")

    def create_file_info(self, file_path):
        """Crea le informazioni del file per la visualizzazione"""
        try:
            file_size = os.path.getsize(file_path)
            modified_time = datetime.fromtimestamp(os.path.getmtime(file_path))
            created_time = datetime.fromtimestamp(os.path.getctime(file_path))
            file_name = os.path.basename(file_path)
            file_extension = os.path.splitext(file_name)[1].lower()
            
            # Determina il tipo di file
            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type:
                file_type = mime_type.split('/')[0].capitalize()
                if file_type == "Application":
                    if "pdf" in mime_type:
                        file_type = "PDF"
                    elif "word" in mime_type or file_extension == ".docx" or file_extension == ".doc":
                        file_type = "Word"
                    elif "excel" in mime_type or file_extension in [".xlsx", ".xls"]:
                        file_type = "Excel"
                    elif "powerpoint" in mime_type or file_extension in [".pptx", ".ppt"]:
                        file_type = "PowerPoint"
                    else:
                        file_type = "Documento"
            else:
                # Fallback basato sull'estensione
                if file_extension in ['.txt', '.md', '.rtf']:
                    file_type = "Testo"
                elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                    file_type = "Immagine"
                elif file_extension in ['.mp3', '.wav', '.ogg', '.flac']:
                    file_type = "Audio"
                elif file_extension in ['.mp4', '.avi', '.mkv', '.mov']:
                    file_type = "Video"
                elif file_extension in ['.exe', '.dll', '.bat']:
                    file_type = "Eseguibile"
                else:
                    file_type = "File"
            
            # Formatta dimensione file
            if file_size < 1024:
                size_str = f"{file_size} B"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.1f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.1f} MB"
            
            return (
                file_type,
                file_name,
                size_str,
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                file_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni del file {file_path}: {str(e)}")
            return (
                "File",
                os.path.basename(file_path),
                "N/A",
                "N/A",
                "N/A",
                file_path
            )

    def create_folder_info(self, folder_path):
        """Crea le informazioni della cartella per la visualizzazione"""
        try:
            modified_time = datetime.fromtimestamp(os.path.getmtime(folder_path))
            created_time = datetime.fromtimestamp(os.path.getctime(folder_path))
            folder_name = os.path.basename(folder_path)
            
            return (
                "Directory",
                folder_name,
                "",  # Le cartelle non hanno dimensione
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                folder_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni della cartella {folder_path}: {str(e)}")
            return (
                "Directory",
                os.path.basename(folder_path),
                "",
                "N/A",
                "N/A",
                folder_path
            )
            
    def should_search_content(self, file_path):
        """Determina se il contenuto del file dovrebbe essere analizzato"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.doc':
            self.log_debug(f"File .doc temporaneamente escluso dall'analisi: {file_path}")
            return False
        
        # Verifica se è un file di sistema - in questo caso salta la ricerca nei contenuti
        # solo se l'opzione di esclusione è attiva
        if self.exclude_system_files.get() and ext in self.system_file_extensions:
            self.log_debug(f"File di sistema escluso dall'analisi del contenuto: {file_path}")
            return False
        
        # Se l'esclusione dei file di sistema è disattivata, prova a leggere tutti i file
        # che non sono nelle estensioni supportate
        if not self.exclude_system_files.get():
            # Accetta tutti i file come leggibili se non è attiva l'esclusione
            return True
        
        # Se l'esclusione è attiva, usa solo le estensioni supportate
        return (
            (ext == '.txt') or
            (ext == '.md') or
            (ext == '.csv') or
            (ext == '.html') or
            (ext == '.htm') or
            (ext == '.xml') or
            (ext == '.json') or
            (ext == '.log') or
            (ext in ['.doc', '.docx'] and file_format_support["docx"]) or
            (ext == '.pdf' and file_format_support["pdf"]) or
            (ext in ['.ppt', '.pptx'] and file_format_support["pptx"]) or
            (ext in ['.xls', '.xlsx'] and file_format_support["excel"]) or
            (ext == '.rtf' and file_format_support["rtf"]) or
            (ext == '.odt' and file_format_support["odt"])
        )

    def should_skip_file(self, file_path):
        """Verifica se un file deve essere saltato durante l'analisi del contenuto"""
        ext = os.path.splitext(file_path)[1].lower()
        skip_type = "File di sistema" if ext in self.system_file_extensions else "File"
        skip_filename = os.path.basename(file_path)
        
        # Salta i file di Rights Management Services
        if "Rights Management Services" in file_path or "IRMProtectors" in file_path:
            self.log_debug(f"Saltato file protetto: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Rights Management Services")
            return True
                
        # Salta file con estensioni problematiche
        problematic_extensions = [".msoprotector.doc", ".msoprotector.ppt", ".msoprotector.xls"]
        if any(ext in file_path for ext in problematic_extensions):
            self.log_debug(f"Saltato file con formato problematico: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Formato problematico")
            return True

        # Salta file di sistema
        if self.exclude_system_files.get() and ext in self.system_file_extensions:
            self.log_debug(f"File di sistema escluso: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "File di sistema")
            return True
                
        return False
    def log_skipped_file(self, filepath, skiptype, filename, skipreason):
        """Registra i file saltati in un file di log"""
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_entry = f"{timestamp} - {skiptype} - {filename} - {filepath} - {skipreason}\n"
            
            with open(self.skipped_files_log_path, 'a', encoding='utf-8') as log_file:
                log_file.write(log_entry)
        except Exception as e:
            self.log_debug(f"Errore durante la scrittura del log dei file saltati: {str(e)}")

    def get_file_content(self, file_path):
        """Ottiene il contenuto del file in base all'estensione"""
        try:

            if self.should_skip_file(file_path):
                return ""
            elif ext == '.doc' and file_format_support["doc"]:
                # Salta i file .doc completamente per evitare blocchi
                self.log_debug(f"File .doc saltato per evitare blocchi: {file_path}")
                return ""
            import os
            ext = os.path.splitext(file_path)[1].lower()
            
            # File di testo semplice
            if ext in ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log']:
                try:
                    # Per file molto grandi, limita la lettura
                    file_size = os.path.getsize(file_path)
                    if file_size > 10 * 1024 * 1024:  # Se è più grande di 10MB
                        try:
                            with open(file_path, 'rb') as f:
                                # Leggi solo i primi e gli ultimi 500KB
                                header = f.read(512 * 1024).decode('utf-8', errors='ignore')
                                
                                # Posiziona il cursore per leggere la fine del file
                                f.seek(-min(512 * 1024, file_size), 2)
                                footer = f.read().decode('utf-8', errors='ignore')
                                
                                return header + "\n[...contenuto troncato...]\n" + footer
                        except Exception as e:
                            self.log_debug(f"Errore nella lettura parziale del file grande {file_path}: {str(e)}")
                    
                    # Per file di dimensioni normali, leggi tutto normalmente
                    with open(file_path, 'r', encoding='utf-8') as f:
                        return f.read()
                except UnicodeDecodeError:
                    # Fallback su altre codifiche
                    try:
                        with open(file_path, 'r', encoding='latin-1') as f:
                            return f.read()
                    except:
                        self.log_debug(f"Impossibile leggere il file {file_path} con le codifiche standard")
                        return ""
            # Word (.doc)
            elif ext == '.doc' and file_format_support["doc"]:
                try:
                    # Importa i moduli necessari
                    import pythoncom
                    import threading
                    import queue
                    
                    # Crea una coda per il risultato
                    result_queue = queue.Queue()
                    
                    # Funzione che verrà eseguita in un thread separato
                    def read_doc_in_thread():
                        try:
                            pythoncom.CoInitialize()  # Inizializza COM nel nuovo thread
                            word = win32com.client.Dispatch("Word.Application")
                            word.Visible = False
                            
                            # Disabilita tutti gli avvisi e dialoghi
                            word.DisplayAlerts = False
                            word.Options.ConfirmConversions = False
                            word.Options.CheckGrammarAsYouType = False
                            word.Options.CheckSpellingAsYouType = False
                            
                            try:
                                # Verifico che il file esista e ottengo il percorso assoluto e normalizzato
                                abs_path = os.path.abspath(file_path)
                                if not os.path.exists(abs_path):
                                    result_queue.put(f"ERRORE: Il file {abs_path} non esiste")
                                    return
                                
                                # Per sicurezza, converti tutte le barre in backslash in stile Windows
                                win_path = abs_path.replace('/', '\\')
                                
                                # Apri il documento con parametri minimi
                                try:
                                    doc = word.Documents.Open(win_path, ReadOnly=True)
                                    # Metodo 1: prova ad accedere direttamente al testo del documento
                                    text = doc.Content.Text if hasattr(doc, "Content") else ""
                                    result_queue.put(text)
                                    
                                    # Chiudi il documento
                                    doc.Close(SaveChanges=False)
                                except Exception as e:
                                    result_queue.put(f"ERRORE: {str(e)}")
                            except Exception as e:
                                result_queue.put(f"ERRORE: {str(e)}")
                            finally:
                                try:
                                    word.Quit()
                                except:
                                    pass
                                
                            pythoncom.CoUninitialize()
                        except Exception as e:
                            result_queue.put(f"ERRORE: {str(e)}")
                    
                    # Avvia il thread per la lettura con timeout più lungo
                    doc_thread = threading.Thread(target=read_doc_in_thread)
                    doc_thread.daemon = True
                    doc_thread.start()
                    
                    # Attendi il risultato con timeout (15 secondi)
                    doc_thread.join(15)
                    
                    # Se il thread è ancora vivo dopo il timeout, è bloccato
                    if doc_thread.is_alive():
                        self.log_debug(f"Timeout nella lettura del documento Word {file_path}")
                        return ""
                    
                    # Altrimenti prendi il risultato dalla coda
                    if not result_queue.empty():
                        result = result_queue.get()
                        if isinstance(result, str) and result.startswith("ERRORE:"):
                            self.log_debug(f"Errore nella lettura del documento Word {file_path}: {result}")
                            return ""
                        return result
                    return ""
                        
                except Exception as e:
                    self.log_debug(f"Errore generale nella lettura del documento Word {file_path}: {str(e)}")
                    return ""

            # Excel (.xls)
            elif ext == '.xls' and file_format_support["xls"]:
                try:
                    # Usa xlrd invece di win32com per i file .xls, è più stabile
                    import xlrd
                    
                    try:
                        wb = xlrd.open_workbook(file_path)
                        texts = []
                        for sheet_index in range(min(wb.nsheets, 5)):  # Limita a 5 fogli
                            sheet = wb.sheet_by_index(sheet_index)
                            for row_idx in range(min(sheet.nrows, 100)):  # Limita a 100 righe
                                row_texts = []
                                for col_idx in range(sheet.ncols):
                                    try:
                                        cell_value = sheet.cell(row_idx, col_idx).value
                                        if cell_value:
                                            row_texts.append(str(cell_value))
                                    except:
                                        continue
                                if row_texts:
                                    texts.append(" ".join(row_texts))
                        return "\n".join(texts)
                    except Exception as e:
                        self.log_debug(f"Errore nella lettura del foglio Excel {file_path}: {str(e)}")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore generale nella lettura del foglio Excel {file_path}: {str(e)}")
                    return ""
            # Word document
            elif ext in ['.docx', '.doc'] and file_format_support["docx"]:
                try:
                    doc = docx.Document(file_path)
                    return "\n".join([para.text for para in doc.paragraphs])
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del documento Word {file_path}: {str(e)}")
                    return ""
            
            # PDF
            elif ext == '.pdf' and file_format_support["pdf"]:
                try:
                    content = []
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        for page_num in range(min(10, len(reader.pages))):  # Limite di 10 pagine per prestazioni
                            page = reader.pages[page_num]
                            content.append(page.extract_text())
                    return "\n".join(content)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del PDF {file_path}: {str(e)}")
                    return ""
            
            # PowerPoint
            elif ext in ['.pptx', '.ppt'] and file_format_support["pptx"]:
                try:
                    prs = pptx.Presentation(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return "\n".join(text)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura della presentazione {file_path}: {str(e)}")
                    return ""
            
            # Excel
            elif ext in ['.xlsx', '.xls'] and file_format_support["excel"]:
                try:
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    texts = []
                    for sheet_name in wb.sheetnames[:5]:  # Limita a 5 fogli per prestazioni
                        sheet = wb[sheet_name]
                        for row in sheet.iter_rows(max_row=100):  # Limita a 100 righe per prestazioni
                            row_texts = [str(cell.value) for cell in row if cell.value is not None]
                            texts.append(" ".join(row_texts))
                    return "\n".join(texts)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del foglio Excel {file_path}: {str(e)}")
                    return ""
            
            # RTF
            elif ext == '.rtf' and file_format_support["rtf"]:
                try:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        rtf_text = f.read()
                    return rtf_to_text(rtf_text)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file RTF {file_path}: {str(e)}")
                    return ""
                    
            # ODT (OpenDocument Text)
            elif ext == '.odt' and file_format_support["odt"]:
                try:
                    # Prova prima con odfdo
                    if 'odfdo' in sys.modules:
                        doc = odfdo.Document(file_path)
                        return doc.get_formatted_text()
                    # Fallback su odf
                    else:
                        textdoc = opendocument.load(file_path)
                        allparas = textdoc.getElementsByType(text.P)
                        return "\n".join([teletype.extractText(para) for para in allparas])
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file ODT {file_path}: {str(e)}")
                    return ""
            
            if not self.exclude_system_files.get():
                try:
                    # Prima prova a leggere come testo UTF-8
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            return f.read()
                    except UnicodeDecodeError:
                        # Poi prova con latin-1 che può leggere qualsiasi byte
                        try:
                            with open(file_path, 'r', encoding='latin-1') as f:
                                return f.read()
                        except Exception as e:
                            # Se fallisce anche questo, prova a leggere i primi KB come binario
                            try:
                                with open(file_path, 'rb') as f:
                                    # Leggi solo i primi 10KB per evitare file troppo grandi
                                    binary_data = f.read(10240)
                                    # Converti in stringa ignorando caratteri non validi
                                    return str(binary_data)
                            except:
                                self.log_debug(f"Impossibile leggere il file {file_path} in nessun formato")
                except Exception as e:
                    self.log_debug(f"Errore nella lettura generica del file {file_path}: {str(e)}")
            
            return ""  # Tipo di file non supportato
        except Exception as e:
            self.log_debug(f"Errore generale nella lettura del file {file_path}: {str(e)}")
            return ""   

    def update_progress(self):
        if self.is_searching:
            try:
                # Processa tutti i messaggi nella coda
                messages_processed = 0
                max_messages_per_cycle = 50  # Limita il numero di messaggi processati per ciclo
                
                while messages_processed < max_messages_per_cycle:
                    try:
                        progress_type, value = self.progress_queue.get_nowait()
                        
                        # Processa il messaggio (codice esistente)
                        if progress_type == "update_total_time":
                            self.total_time_label.config(text=value)
                        elif progress_type == "progress":
                            self.progress_bar["value"] = value
                        elif progress_type == "status":
                            # Se il messaggio contiene informazioni sui file, separiamo le informazioni
                            if "analizzati" in value.lower() or "cartelle:" in value.lower() or "file:" in value.lower():
                                # Estrai il percorso se presente
                                if "analisi:" in value.lower():
                                    parts = value.split("(", 1)
                                    if len(parts) > 1:
                                        path_part = parts[0].strip()
                                        counts_part = "(" + parts[1]
                                        self.status_label["text"] = path_part
                                        self.analyzed_files_label["text"] = counts_part
                                    else:
                                        self.analyzed_files_label["text"] = value
                                else:
                                    self.analyzed_files_label["text"] = value
                            else:
                                # È solo un messaggio di stato semplice
                                self.status_label["text"] = value
                        elif progress_type == "complete":
                            self.is_searching = False
                            self.enable_all_controls()
                            self.stop_button["state"] = "disabled"
                            
                            # Aggiorna la lista dei risultati
                            self.update_results_list()
                            
                            # Aggiorna l'orario di fine e il tempo totale
                            current_time = datetime.now().strftime('%H:%M')
                            self.end_time_label.config(text=current_time)
                            self.update_total_time()  # Calcola e mostra il tempo totale
                            
                            # Calcola la dimensione del percorso alla fine della ricerca
                            self.dir_size_var.set("Calcolo in corso...")
                            path = self.search_path.get()
                            threading.Thread(target=self._calculate_dir_size_thread, args=(path,), daemon=True).start()

                            if len(self.search_results) == 0:
                                self.status_label["text"] = "Nessun file trovato per la ricerca effettuata"
                                self.root.after(100, lambda: messagebox.showinfo("Ricerca completata", "Nessun file trovato per la ricerca effettuata"))
                            else:
                                self.status_label["text"] = f"Ricerca completata! Trovati {len(self.search_results)} risultati."
                            
                            self.progress_bar["value"] = 100
                            return
                        elif progress_type == "error":
                            self.is_searching = False
                            self.enable_all_controls()
                            self.stop_button["state"] = "disabled"
                            current_time = datetime.now().strftime('%H:%M')
                            self.end_time_label.config(text=current_time)
                            self.update_total_time()  # Calcola e mostra il tempo totale
                            messagebox.showerror("Errore", value)
                            return
                        elif progress_type == "timeout":
                            self.is_searching = False
                            self.enable_all_controls()
                            self.stop_button["state"] = "disabled"
                            self.update_results_list()
                            current_time = datetime.now().strftime('%H:%M')
                            self.end_time_label.config(text=current_time)
                            self.update_total_time()  # Calcola e mostra il tempo totale
                            messagebox.showinfo("Timeout", "La ricerca è stata interrotta per timeout. Verranno mostrati i risultati parziali trovati.")
                            return
                        elif progress_type == "admin_prompt":
                            # Chiedi all'utente se desidera riavviare l'app come amministratore
                            response = messagebox.askyesno(
                                "Accesso limitato", 
                                "Alcune cartelle richiedono privilegi di amministratore per essere lette.\n\n" + 
                                "Vuoi riavviare l'applicazione come amministratore per ottenere accesso completo?",
                                icon="question"
                            )
                            if response:
                                self.stop_search_process()
                                self.root.after(1000, self.restart_as_admin)
                            return  # Non interrompere la ricerca se l'utente rifiuta
                        
                        messages_processed += 1
                        
                        # Forza l'aggiornamento dell'interfaccia ogni 10 messaggi
                        if messages_processed % 10 == 0:
                            self.root.update_idletasks()
                        
                    except queue.Empty:
                        break
                        
                # Aggiorna l'UI forzatamente dopo aver processato i messaggi
                self.root.update_idletasks()
                
                # Richiama se stesso più frequentemente per essere più reattivo
                self.root.after(100, self.update_progress)  # Ridotto da 200ms a 100ms
            except Exception as e:
                self.log_debug(f"Errore nell'aggiornamento del progresso: {str(e)}")
                self.root.after(500, self.update_progress)
            
    def stop_search_process(self):
        """Ferma il processo di ricerca in corso"""
        self.stop_search = True
        self.status_label["text"] = "Interrompendo la ricerca..."
        self.analyzed_files_label["text"] = "Ricerca interrotta dall'utente"
        current_time = datetime.now().strftime('%H:%M')
        self.end_time_label.config(text=current_time)
        self.update_total_time()
        
        # Chiusura più decisa dell'executor
        if hasattr(self, 'search_executor') and self.search_executor:
            try:
                self.search_executor.shutdown(wait=False, cancel_futures=True)
                self.search_executor = None
            except Exception as e:
                self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                self.search_executor = None
        
        # Ritardo per evitare problemi con l'interfaccia
        self.root.after(100, self.enable_all_controls)
        self.stop_button["state"] = "disabled"
    
        # Aggiornamento forzato dell'interfaccia
        self.root.update_idletasks()
        
        # Aggiorna l'interfaccia utente
        self.root.update_idletasks()

    def update_results_list(self):
        """Aggiorna la lista dei risultati con i risultati trovati"""
        # Pulisci la lista attuale
        for item in self.results_list.get_children():
            self.results_list.delete(item)
            
        # Aggiungi i risultati alla lista
        for result in self.search_results:
            item_type, name, size, modified, created, path = result
            
            # Applica stile in base al tipo di elemento
            if item_type == "Directory":
                tags = ("directory",)
            else:
                tags = ("file",)
                
            self.results_list.insert("", "end", values=result, tags=tags)
            
        # Aggiorna lo stato
        self.status_label["text"] = f"Trovati {len(self.search_results)} risultati"
        
    def update_theme_colors(self, theme="light"):
        """Aggiorna i colori del tema per evidenziare cartelle e file"""
        style = ttk.Style()
        
        # Configura i colori per il tema chiaro
        if theme == "light":
            style.configure("Treeview", background="#ffffff", foreground="#000000", fieldbackground="#ffffff")
            self.results_list.tag_configure("directory", background="#e6f2ff", foreground="#000000")
            self.results_list.tag_configure("file", background="#ffffff", foreground="#000000")
        # Configura i colori per il tema scuro
        elif theme == "dark":
            style.configure("Treeview", background="#333333", foreground="#ffffff", fieldbackground="#333333")
            self.results_list.tag_configure("directory", background="#4d4d4d", foreground="#ffffff")
            self.results_list.tag_configure("file", background="#333333", foreground="#ffffff")

    def show_block_options(self):
        """Mostra la finestra di dialogo per le opzioni avanzate di ricerca a blocchi"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Opzioni Avanzate Ricerca a Blocchi")
        dialog.geometry("560x600")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Descrizione
        ttk.Label(main_frame, text="Configura le opzioni per ottimizzare la ricerca a blocchi", 
             font=("", 10, "bold")).pack(anchor=W, pady=(0, 10))
        
        description = ttk.Label(main_frame, text="La ricerca a blocchi divide la scansione in blocchi più piccoli "
                                        "per migliorare le prestazioni e la reattività dell'interfaccia.",
                         wraplength=520, justify=LEFT)
        description.pack(fill=X, pady=(0, 10))
        
        # Frame per opzioni dei blocchi
        block_frame = ttk.LabelFrame(main_frame, text="Dimensione e Parallelismo", padding=10)
        block_frame.pack(fill=X, pady=5)
        
        # Grid per organizzare le opzioni
        grid = ttk.Frame(block_frame)
        grid.pack(fill=X)
        
        # Opzione: Max file per blocco
        ttk.Label(grid, text="Max file per blocco:").grid(row=0, column=0, sticky=W, padx=5, pady=5)
        files_per_block = ttk.Spinbox(grid, from_=100, to=10000, increment=100, width=7,
                                 textvariable=self.max_files_per_block)
        files_per_block.grid(row=0, column=1, padx=5, pady=5)
        
        # Opzione: Max blocchi paralleli
        ttk.Label(grid, text="Blocchi paralleli:").grid(row=0, column=2, sticky=W, padx=(20, 5), pady=5)
        parallel_blocks = ttk.Spinbox(grid, from_=1, to=16, increment=1, width=5,
                                 textvariable=self.max_parallel_blocks)
        parallel_blocks.grid(row=0, column=3, padx=5, pady=5)
        
        # Checkbox per opzioni aggiuntive
        auto_adjust = ttk.Checkbutton(grid, text="Adatta automaticamente la dimensione dei blocchi",
                                variable=self.block_size_auto_adjust)
        auto_adjust.grid(row=1, column=0, columnspan=4, sticky=W, padx=5, pady=5)
        self.create_tooltip(auto_adjust, "Regola automaticamente la dimensione dei blocchi in base al tipo di ricerca e al sistema")
        
        # Frame per opzioni di prioritizzazione
        priority_frame = ttk.LabelFrame(main_frame, text="Priorità di Ricerca", padding=10)
        priority_frame.pack(fill=X, pady=10)
        
        # Checkbox per dare priorità alle cartelle utente
        user_priority = ttk.Checkbutton(priority_frame, text="Dare priorità alle cartelle utente (Users, Documents, Desktop, ecc)",
                                variable=self.prioritize_user_folders)
        user_priority.pack(anchor=W, padx=5, pady=5)
        self.create_tooltip(user_priority, "Elabora prima le cartelle più importanti per trovare risultati utili più velocemente")
        
        # Frame per spiegazione avanzata
        info_frame = ttk.LabelFrame(main_frame, text="Come funziona", padding=10)
        info_frame.pack(fill=X, pady=10)
        
        info_text = "La ricerca a blocchi divide le cartelle in unità di lavoro più piccole (blocchi) " \
                "che vengono elaborate in base alla priorità. Questo permette di:\n\n" \
                "• Ottenere risultati utili più rapidamente\n" \
                "• Migliorare la reattività dell'interfaccia durante la ricerca\n" \
                "• Distribuire meglio il carico di lavoro sui thread\n" \
                "• Gestire meglio la memoria per ricerche su grandi volumi di dati"
        
        info_label = ttk.Label(info_frame, text=info_text, wraplength=520, justify=LEFT)
        info_label.pack(fill=X, padx=5, pady=5)
        
        # Pulsanti per chiudere/applicare
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        # Pulsante per ripristinare valori predefiniti
        ttk.Button(btn_frame, text="Valori predefiniti", 
              command=lambda: [self.max_files_per_block.set(1000), 
                              self.max_parallel_blocks.set(4),
                              self.prioritize_user_folders.set(True),
                              self.block_size_auto_adjust.set(True)]).pack(side=LEFT, padx=5)
        
        # Pulsante di chiusura
        ttk.Button(btn_frame, text="Chiudi", command=dialog.destroy).pack(side=RIGHT, padx=5)
        
        # Centra la finestra
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")

    def copy_selected(self):
        """Copia i file selezionati in una directory di destinazione"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da copiare")
            return
            
        # Chiedi la directory di destinazione
        dest_dir = filedialog.askdirectory(title="Seleziona la cartella di destinazione")
        if not dest_dir:
            return
            
        # Prepara variabili per tracciare l'avanzamento
        total = len(selected_items)
        copied = 0
        failed = 0
        skipped = 0
        
        try:
            for item in selected_items:
                values = self.results_list.item(item)['values']
                item_type, _, _, _, _, source_path = values
                
                # Ottieni il nome dell'elemento senza il percorso completo
                basename = os.path.basename(source_path)
                dest_path = os.path.join(dest_dir, basename)
                
                try:
                    if item_type == "Directory":
                        # Se la cartella esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("Cartella esistente", 
                                                f"La cartella {basename} esiste già. Vuoi sovrascriverla?"):
                                # Elimina la cartella esistente
                                shutil.rmtree(dest_path)
                            else:
                                skipped += 1
                                continue
                                
                        # Copia ricorsiva della cartella
                        shutil.copytree(source_path, dest_path)
                        copied += 1
                    else:
                        # Se il file esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("File esistente", 
                                                f"Il file {basename} esiste già. Vuoi sovrascriverlo?"):
                                # Continua con la sovrascrittura
                                pass
                            else:
                                skipped += 1
                                continue
                                
                        # Copia il file
                        shutil.copy2(source_path, dest_path)
                        copied += 1
                        
                except Exception as e:
                    failed += 1
                    self.log_debug(f"Errore durante la copia di {source_path}: {str(e)}")
                    
                # Aggiorna la barra di avanzamento
                progress = ((copied + failed + skipped) / total) * 100
                self.progress_bar["value"] = progress
                self.status_label["text"] = f"Copiati {copied}/{total} elementi..."
                self.root.update()
                
            # Messaggio di completamento
            if failed > 0:
                messagebox.showwarning("Copia completata con errori", 
                                    f"Copiati: {copied}\nSaltati: {skipped}\nFalliti: {failed}")
            else:
                messagebox.showinfo("Copia completata", 
                                 f"Copiati con successo: {copied}\nSaltati: {skipped}")
                
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'operazione: {str(e)}")
            
        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    def get_zip_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome del file ZIP"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Nome file ZIP")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome del file ZIP (senza estensione):", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = StringVar(value="archivio")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            zip_name = name_var.get().strip()
            if zip_name:
                result["name"] = zip_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per il file ZIP", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea ZIP", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]

    def compress_selected(self):
        """Versione avanzata della compressione con opzioni aggiuntive e log completo dei file"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da comprimere")
            return

        # Ottieni il nome del file ZIP
        zip_name = self.get_zip_name()

        if not zip_name:
            return

        # Ottieni il nome della cartella principale
        main_folder_name = self.get_main_folder_name()

        if main_folder_name is None:
            return  # L'utente ha annullato

        if not main_folder_name:
            main_folder_name = "files"  # Default se l'utente inserisce valore vuoto

        # Opzioni di compressione più dettagliate
        compression_dialog = tk.Toplevel(self.root)
        compression_dialog.title("Opzioni Compressione")
        compression_dialog.transient(self.root)
        compression_dialog.grab_set()

        frame = ttk.Frame(compression_dialog, padding=20)
        frame.pack(fill=tk.BOTH, expand=True)

        ttk.Label(frame, text="Scegli il tipo di compressione:").pack(anchor=tk.W, pady=(0,10))

        comp_var = tk.StringVar(value="nessuna")

        ttk.Radiobutton(frame, text="Nessuna (solo archiviazione)", 
            variable=comp_var, value="nessuna").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Minima (massima velocità) Liv.1", 
                    variable=comp_var, value="minima").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Standard (buon equilibrio) Liv.6", 
                    variable=comp_var, value="standard").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Massima (compressione ottimale, più lenta) Liv.9", 
                    variable=comp_var, value="massima").pack(anchor=tk.W)

        ttk.Label(frame, text="\nOpzioni per file di grandi dimensioni:").pack(anchor=tk.W, pady=(10,5))

        use_chunks = tk.BooleanVar(value=False)
        ttk.Checkbutton(frame, text="Elabora in blocchi (per file molto grandi)", 
                    variable=use_chunks).pack(anchor=tk.W)

        btn_frame = ttk.Frame(frame)
        btn_frame.pack(fill=tk.X, pady=(15,0))

        result = {"option": None, "chunks": False}

        def on_cancel():
            result["option"] = None
            compression_dialog.destroy()

        def on_ok():
            result["option"] = comp_var.get()
            result["chunks"] = use_chunks.get()
            compression_dialog.destroy()

        ttk.Button(btn_frame, text="Annulla", command=on_cancel).pack(side=tk.LEFT)
        ttk.Button(btn_frame, text="OK", command=on_ok).pack(side=tk.RIGHT)

        # Centra la finestra
        compression_dialog.update_idletasks()
        width = compression_dialog.winfo_reqwidth() + 50
        height = compression_dialog.winfo_reqheight() + 20
        x = (compression_dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (compression_dialog.winfo_screenheight() // 2) - (height // 2)
        compression_dialog.geometry(f"{width}x{height}+{x}+{y}")

        compression_dialog.wait_window()

        if result["option"] is None:
            return  # Utente ha annullato

        # Determina il metodo di compressione
        compression_level = 0  # Valore predefinito

        if result["option"] == "nessuna":
            compression_method = zipfile.ZIP_STORED
            compression_text = "nessuna"
            # compression_level non usato per ZIP_STORED
        elif result["option"] == "minima":
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "minima"
            compression_level = 1  # Compressione veloce, riduzione minima
        elif result["option"] == "massima":
            # Prova LZMA se disponibile, altrimenti BZIP2, altrimenti fallback su DEFLATED con livello massimo
            if hasattr(zipfile, 'ZIP_LZMA'):
                compression_method = zipfile.ZIP_LZMA
                compression_text = "massima (LZMA)"
                compression_level = 9
            elif hasattr(zipfile, 'ZIP_BZIP2'):
                compression_method = zipfile.ZIP_BZIP2
                compression_text = "massima (BZIP2)"
                compression_level = 9
            else:
                compression_method = zipfile.ZIP_DEFLATED
                compression_text = "massima (DEFLATED)"
                compression_level = 9
        else:  # standard
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "standard"
            compression_level = 6  # Compressione bilanciata
            
        # Log della scelta di compressione
        self.log_debug(f"Utilizzo compressione {compression_text} (metodo: {compression_method}, livello: {compression_level})")

        zip_path = filedialog.asksaveasfilename(
            defaultextension=".zip",
            initialfile=f"{zip_name}.zip",
            filetypes=[("ZIP files", "*.zip")],
            title="Salva file ZIP"
        )

        if not zip_path:
            return

        # Raccogli tutti i file delle cartelle selezionate
        files_in_folders = set()
        folder_paths = []
        single_files = []

        # Prima fase: raccogli informazioni su cartelle e file
        for item in selected_items:
            values = self.results_list.item(item)['values']
            item_type, _, _, _, _, source_path = values

            if item_type == "Directory":
                folder_paths.append(source_path)
                # Raccogli tutti i file nelle cartelle selezionate
                for root, _, files in os.walk(source_path):
                    for file in files:
                        full_path = os.path.join(root, file)
                        files_in_folders.add(os.path.abspath(full_path))
            else:
                single_files.append(source_path)

        # Filtra i file singoli che sono già presenti nelle cartelle
        filtered_single_files = [f for f in single_files if os.path.abspath(f) not in files_in_folders]

        total_items = len(folder_paths) + len(filtered_single_files)
        processed = 0

        # Lista per tenere traccia dei file (tutti, non solo omonimi)
        all_files_log = []
        # Lista per tenere traccia dei file omonimi (come prima)
        omonimi_log = []

        try:
            # Primo passaggio: analizza tutti i file e identifica gli omonimi
            self.log_debug(f"Analisi dei file da comprimere nella cartella {main_folder_name}...")

            # Dizionario {nome_file: [lista_percorsi]}
            file_names_map = {}

            # Raccogli i nomi dei file dalle cartelle
            for folder_path in folder_paths:
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        file_name = os.path.basename(file_path)

                        # Aggiungi al log di tutti i file
                        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                        modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                        
                        all_files_log.append({
                            "nome_file": file_name,
                            "percorso_originale": file_path,
                            "dimensione": self._format_size(file_size),
                            "ultima_modifica": modified_time,
                            "tipo": self._get_file_type(file_path)
                        })

                        if file_name not in file_names_map:
                            file_names_map[file_name] = []
                        file_names_map[file_name].append(file_path)

            # Raccogli i nomi dei file singoli
            for file_path in filtered_single_files:
                if os.path.exists(file_path):
                    file_name = os.path.basename(file_path)
                    
                    # Aggiungi al log di tutti i file
                    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                    modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                    
                    all_files_log.append({
                        "nome_file": file_name,
                        "percorso_originale": file_path,
                        "dimensione": self._format_size(file_size),
                        "ultima_modifica": modified_time,
                        "tipo": self._get_file_type(file_path)
                    })
                    
                    if file_name not in file_names_map:
                        file_names_map[file_name] = []
                    file_names_map[file_name].append(file_path)

            # Identifica i file omonimi (con lo stesso nome ma percorsi diversi)
            omonimi_files = {name: paths for name, paths in file_names_map.items() if len(paths) > 1}
            self.log_debug(f"Trovati {len(omonimi_files)} file con nomi duplicati")
            
            # Per l'elaborazione a blocchi
            if result["chunks"] and total_items > 50:
                chunk_size = 10  # Elabora 10 file alla volta
                self.log_debug(f"Elaborazione a blocchi attivata: {chunk_size} file per blocco")

            # Secondo passaggio: crea l'archivio ZIP
            with zipfile.ZipFile(zip_path, 'w', compression_method, compresslevel=compression_level) as zipf:
                # Tiene traccia dei percorsi già aggiunti al file ZIP
                added_zip_paths = set()

                # Comprimi le cartelle
                for folder_path in folder_paths:
                    base_folder = os.path.basename(folder_path)
                    for root, _, files in os.walk(folder_path):
                        for file in files:
                            file_path = os.path.join(root, file)
                            file_name = os.path.basename(file_path)

                            try:
                                # Salta file non leggibili
                                if not os.access(file_path, os.R_OK):
                                    self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                    continue

                                # Calcola il percorso relativo standard
                                rel_path = os.path.relpath(file_path, os.path.dirname(folder_path))

                                # Verifica se è un file omonimo
                                if file_name in omonimi_files:
                                    # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                    if omonimi_files[file_name][0] == file_path:
                                        # Usa il percorso all'interno della cartella principale
                                        zip_path_to_use = os.path.join(main_folder_name, rel_path)
                                    else:
                                        # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                        parent_folder = os.path.basename(os.path.dirname(file_path))
                                        unique_name = f"{parent_folder}_{file_name}"

                                        # Se anche questo nome è duplicato, aggiungi un contatore
                                        counter = 1
                                        while os.path.join("omonimi", unique_name) in added_zip_paths:
                                            unique_name = f"{parent_folder}_{counter}_{file_name}"
                                            counter += 1

                                        zip_path_to_use = os.path.join("omonimi", unique_name)

                                        # Registra questo file nel log degli omonimi
                                        omonimi_log.append({
                                            "nome_file": file_name,
                                            "percorso_originale": file_path,
                                            "primo_percorso": omonimi_files[file_name][0],
                                            "posizione_zip": zip_path_to_use
                                        })
                                else:
                                    # Non è un omonimo, inseriscilo nella cartella principale
                                    zip_path_to_use = os.path.join(main_folder_name, rel_path)

                                # Verifica se questo percorso ZIP è già stato usato
                                if zip_path_to_use in added_zip_paths:
                                    self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                    # Crea un nome alternativo
                                    alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)

                                    # Se anche questo nome è già usato, aggiungi un contatore
                                    counter = 1
                                    while zip_path_to_use in added_zip_paths:
                                        alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                        zip_path_to_use = os.path.join("omonimi", alt_name)
                                        counter += 1

                                # Aggiungi il file al ZIP e registra il percorso
                                zipf.write(file_path, zip_path_to_use)
                                added_zip_paths.add(zip_path_to_use)

                            except Exception as e:
                                self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Comprimi i file singoli
                for file_path in filtered_single_files:
                    if os.path.exists(file_path):
                        file_name = os.path.basename(file_path)

                        try:
                            # Salta file non leggibili
                            if not os.access(file_path, os.R_OK):
                                self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                continue

                            # Verifica se è un file omonimo
                            if file_name in omonimi_files:
                                # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                if omonimi_files[file_name][0] == file_path:
                                    # Usa il nome del file all'interno della cartella principale
                                    zip_path_to_use = os.path.join(main_folder_name, file_name)
                                else:
                                    # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                    parent_folder = os.path.basename(os.path.dirname(file_path))
                                    unique_name = f"{parent_folder}_{file_name}"

                                    # Se anche questo nome è duplicato, aggiungi un contatore
                                    counter = 1
                                    while os.path.join("omonimi", unique_name) in added_zip_paths:
                                        unique_name = f"{parent_folder}_{counter}_{file_name}"
                                        counter += 1

                                    zip_path_to_use = os.path.join("omonimi", unique_name)

                                    # Registra questo file nel log degli omonimi
                                    omonimi_log.append({
                                        "nome_file": file_name,
                                        "percorso_originale": file_path,
                                        "primo_percorso": omonimi_files[file_name][0],
                                        "posizione_zip": zip_path_to_use
                                    })
                            else:
                                # Non è un omonimo, inseriscilo nella cartella principale
                                zip_path_to_use = os.path.join(main_folder_name, file_name)

                            # Verifica se questo percorso ZIP è già stato usato
                            if zip_path_to_use in added_zip_paths:
                                self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                # Crea un nome alternativo
                                alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                zip_path_to_use = os.path.join("omonimi", alt_name)

                                # Se anche questo nome è già usato, aggiungi un contatore
                                counter = 1
                                while zip_path_to_use in added_zip_paths:
                                    alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)
                                    counter += 1

                            # Aggiungi il file al ZIP e registra il percorso
                            zipf.write(file_path, zip_path_to_use)
                            added_zip_paths.add(zip_path_to_use)

                        except Exception as e:
                            self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Crea log in formato CSV con collegamenti ipertestuali
                current_time_utc = datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')
                current_time_local = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                current_user = getpass.getuser()

                # Identifica quali file sono omonimi per poterli escludere dal log generale
                omonimo_paths = set()
                for entry in omonimi_log:
                    omonimo_paths.add(entry['percorso_originale'])

                # Lista di file normali (esclude gli omonimi)
                normal_files_log = [f for f in all_files_log if f['percorso_originale'] not in omonimo_paths]

                # ---------- CREAZIONE DEL CSV DEI FILE NORMALI ----------
                # Creazione intestazione per il CSV
                csv_header = f"LOG DEI FILE TROVATI - {current_time_local}\n"
                csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                csv_header += f"Totale file (esclusi omonimi): {len(normal_files_log)}\n\n"

                # Prepara l'output CSV
                csv_content = io.StringIO()
                csv_writer = csv.writer(csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)

                # Intestazione colonne
                csv_writer.writerow(["Nome file", "Tipo", "Dimensione", "Ultima modifica", "Percorso completo", "Link"])

                # Aggiungi i file normali (non omonimi)
                for file_entry in sorted(normal_files_log, key=lambda x: x["nome_file"]):
                    # Ottieni il percorso della directory che contiene il file
                    file_path = file_entry['percorso_originale']
                    directory_path = os.path.dirname(file_path).replace('\\', '/')
                    # Formatta il percorso della directory per il collegamento ipertestuale
                    folder_uri = f"file:///{directory_path}"
                    
                    # Aggiungi riga al CSV con la formula italiana per i collegamenti alla CARTELLA
                    csv_writer.writerow([
                        file_entry['nome_file'],
                        file_entry['tipo'],
                        file_entry['dimensione'],
                        file_entry['ultima_modifica'],
                        file_entry['percorso_originale'],
                        f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")'
    ])

                # Ottieni il contenuto del CSV come stringa
                csv_data = csv_header + csv_content.getvalue()

                # Aggiungi il file CSV dei file normali all'archivio
                zipf.writestr("Log_file_Trovati.csv", csv_data)

                # ---------- CREAZIONE DEL CSV DEGLI OMONIMI (SOLO SE CE NE SONO) ----------
                if omonimi_log:
                    # Creazione intestazione per il CSV degli omonimi
                    omonimi_csv_header = f"LOG DEI FILE OMONIMI - {current_time_local}\n"
                    omonimi_csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                    omonimi_csv_header += f"Totale file omonimi trovati: {len(omonimi_log)}\n\n"
                    
                    # Prepara l'output CSV per gli omonimi
                    omonimi_csv_content = io.StringIO()
                    omonimi_csv_writer = csv.writer(omonimi_csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
                    
                    # Intestazione colonne specifiche per gli omonimi
                    omonimi_csv_writer.writerow([
                        "Nome file", "Tipo", "Dimensione", "Ultima modifica", 
                        "Percorso originale", "Link", "Posizione nello ZIP", "In conflitto con"
                    ])
                    
                    # Estrai dai log completi solo i file omonimi
                    omonimi_files_details = [f for f in all_files_log if f['percorso_originale'] in omonimo_paths]
                    
                    # Crea un dizionario per unire le informazioni di omonimi_log e omonimi_files_details
                    omonimi_details_dict = {}
                    for entry in omonimi_log:
                        omonimi_details_dict[entry['percorso_originale']] = entry
                    
                    # Aggiungi i file omonimi al CSV
                    for file_entry in sorted(omonimi_files_details, key=lambda x: x["nome_file"]):
                        # Ottieni il percorso della directory che contiene il file
                        file_path = file_entry['percorso_originale']
                        directory_path = os.path.dirname(file_path).replace('\\', '/')
                        # Formatta il percorso della directory per il collegamento ipertestuale
                        folder_uri = f"file:///{directory_path}"
                        
                        # Trova le informazioni aggiuntive sugli omonimi
                        omonimo_info = omonimi_details_dict.get(file_entry['percorso_originale'], {})
                        posizione_zip = omonimo_info.get('posizione_zip', 'N/A')
                        in_conflitto = omonimo_info.get('primo_percorso', 'N/A')
                        
                        # Aggiungi riga al CSV degli omonimi con la formula italiana per i collegamenti alla CARTELLA
                        omonimi_csv_writer.writerow([
                            file_entry['nome_file'],
                            file_entry['tipo'],
                            file_entry['dimensione'],
                            file_entry['ultima_modifica'],
                            file_entry['percorso_originale'],
                            f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")',
                            posizione_zip,
                            in_conflitto
    ])
                    
                    # Ottieni il contenuto del CSV come stringa
                    omonimi_csv_data = omonimi_csv_header + omonimi_csv_content.getvalue()
                    
                    # Aggiungi il file CSV degli omonimi all'archivio
                    zipf.writestr("omonimi_log.csv", omonimi_csv_data)

            # Prepara il messaggio di completamento
            skipped_files = len(single_files) - len(filtered_single_files)
            message = f"Compressione completata!\nFile salvato in: {zip_path}\n"
            message += f"File organizzati nella cartella '{main_folder_name}'\n"
            message += f"Tipo di compressione utilizzata: {compression_text}\n"
            message += f"\nCreato file di log completo ('Log_file_Trovati.txt')"

            if result["chunks"]:
                message += "\nElaborazione a blocchi: Attiva"

            if omonimi_log:
                message += f"\nTrovati {len(omonimi_log)} file omonimi."
                message += f"\nCreato log dettagliato degli omonimi ('omonimi_log.txt')"

            if skipped_files > 0:
                message += f"\n{skipped_files} file saltati perché già presenti nelle cartelle"

            messagebox.showinfo("Completato", message)

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            self.log_debug(f"Errore dettagliato durante la compressione:\n{error_details}")

            messagebox.showerror("Errore", 
                f"Errore durante la compressione: {str(e)}\n\n"
                f"Tipo di errore: {type(e).__name__}")

        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."
            
    def get_directory_size(self, path):
        """Calculate the total size of a directory"""
        if not os.path.exists(path):
            return 0
            
        try:
            total_size = 0
            if os.path.isfile(path):
                return os.path.getsize(path)
                
            # Per evitare il blocco dell'interfaccia, imposta un timeout massimo
            start_time = time.time()
            max_time = 30  # massimo 30 secondi
            files_count = 0
            error_count = 0
                
            # For directories, walk through all files and subdirectories
            for dirpath, dirnames, filenames in os.walk(path):
                for f in filenames:
                    # Verifica se il timeout è scaduto
                    if time.time() - start_time > max_time:
                        self.log_debug(f"Timeout nel calcolo della dimensione per {path}")
                        return total_size
                        
                    try:
                        fp = os.path.join(dirpath, f)
                        
                        # Verifica esplicita che il file esiste ancora prima di tentare di leggerne la dimensione
                        if os.path.exists(fp) and not os.path.islink(fp):
                            total_size += os.path.getsize(fp)
                            files_count += 1
                    except FileNotFoundError:
                        # Ignora silenziosamente i file che non esistono più
                        pass
                    except (OSError, PermissionError) as e:
                        error_count += 1
                        # Limita il numero di errori da registrare per evitare spam nel log
                        if error_count < 100:  
                            self.log_debug(f"Error getting size of {fp}: {str(e)}")
                
                # Update the status periodically to show progress
                if files_count % 1000 == 0:  # Ogni 1000 file
                    self.root.after(0, lambda size=total_size: 
                        self.status_label.config(text=f"Calcolando dimensione: {self._format_size(size)}..."))
                    
            return total_size
        except Exception as e:
            self.log_debug(f"Error calculating directory size for {path}: {str(e)}")
            return 0
          
    def get_disk_space(self, path):
        """Get disk space information for the partition containing the path"""
        if not os.path.exists(path):
            return (0, 0, 0)
        
        try:
            if os.name == 'nt':  # Windows
                total, used, free = shutil.disk_usage(os.path.splitdrive(path)[0] + '\\')
            else:  # Linux/Mac
                total, used, free = shutil.disk_usage(os.path.abspath(os.path.join(path, os.pardir)))
            return (total, used, free)
        except Exception as e:
            self.log_debug(f"Error getting disk space for {path}: {str(e)}")
            return (0, 0, 0)
        
    def update_disk_info(self, path=None, calculate_dir_size=True):
        """Update the disk information for the current path"""
        if path is None:
            path = self.search_path.get()
            
        if not path or not os.path.exists(path):
            self.dir_size_var.set("N/A")
            self.total_disk_var.set("N/A")
            self.used_disk_var.set("N/A")
            self.free_disk_var.set("N/A")
            return
            
        # Imposta lo stato a "Calcolando"
        self.status_label["text"] = "Calcolando spazio disco..."
        self.root.update_idletasks()
        
        try:
            # Get disk space info
            total, used, free = self.get_disk_space(path)
            
            self.total_disk_var.set(self._format_size(total))
            self.used_disk_var.set(self._format_size(used))
            self.free_disk_var.set(self._format_size(free))
        except Exception as e:
            self.log_debug(f"Errore nel calcolo dello spazio disco: {str(e)}")
            self.total_disk_var.set("Errore")
            self.used_disk_var.set("Errore")
            self.free_disk_var.set("Errore")
        
        # Non calcolare la dimensione della cartella all'avvio della ricerca
        # perché potrebbe bloccare l'interfaccia troppo a lungo
        if calculate_dir_size:
            self.dir_size_var.set("Calcolo in corso...")
            # Get directory size in a separate thread
            threading.Thread(target=self._calculate_dir_size_thread, args=(path,), daemon=True).start()
        else:
            self.dir_size_var.set("Calcolo disattivato")
        
    def _calculate_dir_size_thread(self, path):
        """Thread function to calculate directory size"""
        dir_size = self.get_directory_size(path)
        
        # Update the UI from the main thread
        self.root.after(0, lambda: self.dir_size_var.set(self._format_size(dir_size)))
        self.root.after(0, lambda: self.status_label.config(text="In attesa..."))

    # Funzione helper per formattare la dimensione del file
    def _format_size(self, size_bytes):
        """Formatta la dimensione del file in modo leggibile"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.2f} KB"
        elif size_bytes < 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024):.2f} MB"
        else:
            return f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"

    # Funzione helper per determinare il tipo di file
    def _get_file_type(self, file_path):
        """Determina il tipo di file in base all'estensione"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.txt', '.md', '.rtf']:
            return "Documento di testo"
        elif ext in ['.doc', '.docx', '.odt']:
            return "Documento Word"
        elif ext in ['.xls', '.xlsx', '.ods']:
            return "Foglio di calcolo"
        elif ext in ['.pdf']:
            return "PDF"
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            return "Immagine"
        elif ext in ['.mp3', '.wav', '.ogg', '.flac']:
            return "Audio"
        elif ext in ['.mp4', '.avi', '.mkv', '.mov']:
            return "Video"
        elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
            return "Archivio"
        elif ext in ['.exe', '.dll', '.bat', '.cmd']:
            return "Eseguibile"
        else:
            return "File"
    
    def get_main_folder_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome della cartella principale"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Cartella Principale")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=tk.BOTH, expand=tk.YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome della cartella principale nell'archivio:", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = tk.StringVar(value="files")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=tk.X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=tk.X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            folder_name = name_var.get().strip()
            if folder_name:
                result["name"] = folder_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per la cartella", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea Cartella", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=tk.RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]
    
    def open_file_location(self, event=None):
        """Apre il percorso del file selezionato nel file explorer"""
        selected_items = self.results_list.selection()
        if not selected_items:
            return
            
        selected_item = selected_items[0]  # Prendi il primo elemento selezionato
        file_path = self.results_list.item(selected_item, "values")[5]  # Ottieni il percorso del file
        
        try:
            if os.path.exists(file_path):
                # Ottieni la directory contenente il file
                directory = os.path.dirname(file_path)
                
                if os.name == 'nt':  # Windows
                    # Converti eventuali forward slash in backslash per Windows
                    file_path = os.path.normpath(file_path)
                    # Usa il metodo più sicuro con subprocess invece di os.system
                    subprocess.run(['explorer', '/select,', file_path], shell=True)
                else:
                    # Per sistemi Linux/Unix
                    if shutil.which('xdg-open'):  # Verifica che xdg-open sia disponibile
                        subprocess.run(['xdg-open', directory])
                    elif shutil.which('open'):  # Per macOS
                        subprocess.run(['open', directory])
                    else:
                        self.log_debug("Nessun comando disponibile per aprire directory")
                        messagebox.showinfo("Informazione", f"Percorso del file: {directory}")
                
                self.log_debug(f"Apertura percorso: {file_path}")
            else:
                messagebox.showinfo("Informazione", f"Il percorso non esiste: {file_path}")
                self.log_debug(f"Percorso non esistente: {file_path}")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile aprire il percorso: {str(e)}")
            self.log_debug(f"Errore nell'apertura del percorso: {str(e)}")

    def show_advanced_filters_dialog(self):
        """Mostra la finestra di dialogo per i filtri di ricerca avanzati"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Filtri avanzati")
        dialog.geometry("530x440")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Filtri dimensione
        size_frame = ttk.LabelFrame(dialog, text="Dimensione file")
        size_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(size_frame, text="Min (KB):").grid(row=0, column=0, padx=5, pady=5)
        min_size = ttk.Entry(size_frame, width=10)
        min_size.grid(row=0, column=1, padx=5, pady=5)
        min_size.insert(0, str(self.advanced_filters["size_min"] // 1024))
        
        ttk.Label(size_frame, text="Max (KB):").grid(row=0, column=2, padx=5, pady=5)
        max_size = ttk.Entry(size_frame, width=10)
        max_size.grid(row=0, column=3, padx=5, pady=5)
        max_size.insert(0, str(self.advanced_filters["size_max"] // 1024 if self.advanced_filters["size_max"] else 0))
        
        # Filtri data - FORMAT DD-MM-YYYY
        date_frame = ttk.LabelFrame(dialog, text="Data modifica (DD-MM-YYYY)")
        date_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(date_frame, text="Da:").grid(row=0, column=0, padx=5, pady=5)
        min_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        min_date.grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(date_frame, text="A:").grid(row=0, column=2, padx=5, pady=5)
        max_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        max_date.grid(row=0, column=3, padx=5, pady=5)
        
        # Cancella sempre le date precedenti
        min_date.entry.delete(0, "end")
        max_date.entry.delete(0, "end")
        
        # Filtri estensione
        ext_frame = ttk.LabelFrame(dialog, text="Estensioni file (separate da virgola)")
        ext_frame.pack(fill=X, padx=10, pady=5)
        example_label = ttk.Label(ext_frame, text="Inserisci una o più estensioni da ricercare. Esempio: .pdf, .dot", 
                    font=("", 8), foreground="gray")
        example_label.pack(anchor="w", padx=5)
        
        extensions = ttk.Entry(ext_frame)
        extensions.pack(fill=X, padx=5, pady=5)
        if self.advanced_filters["extensions"]:
            extensions.insert(0, ", ".join(self.advanced_filters["extensions"]))
        
        # Aggiungiamo un frame di debug per vedere i filtri correnti
        debug_frame = ttk.LabelFrame(dialog, text="Debug - Stato filtri correnti")
        debug_frame.pack(fill=X, padx=10, pady=5)
        
        debug_text = ttk.Text(debug_frame, height=3, width=50)
        debug_text.pack(fill=X, padx=5, pady=5)
        debug_text.insert("1.0", f"Data min: {self.advanced_filters['date_min']}\n")
        debug_text.insert("2.0", f"Data max: {self.advanced_filters['date_max']}\n")
        debug_text.insert("3.0", f"Extensions: {self.advanced_filters['extensions']}")
        debug_text.config(state="disabled")
        
        # Pulsante Salva
        def save_filters():
            try:
                # Analizza i filtri di dimensione
                min_kb = int(min_size.get() or 0)
                max_kb = int(max_size.get() or 0)
                self.advanced_filters["size_min"] = min_kb * 1024
                self.advanced_filters["size_max"] = max_kb * 1024
                
                # Ottieni le date nel formato DD-MM-YYYY
                min_date_value = min_date.entry.get().strip()
                max_date_value = max_date.entry.get().strip()
                
                print(f"DEBUG - Date inserite: min={min_date_value}, max={max_date_value}")
                
                # Validazione date
                if min_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(min_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data minima non valido. Usa DD-MM-YYYY")
                        return
                
                if max_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(max_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data massima non valido. Usa DD-MM-YYYY")
                        return
                
                # Verifica che la data minima non sia successiva alla data massima
                if min_date_value and max_date_value:
                    min_date_obj = datetime.strptime(min_date_value, "%d-%m-%Y")
                    max_date_obj = datetime.strptime(max_date_value, "%d-%m-%Y")
                    
                    if min_date_obj > max_date_obj:
                        messagebox.showerror("Errore", 
                                            "La data di inizio non può essere successiva alla data di fine")
                        return
                
                # Salva le date validate
                self.advanced_filters["date_min"] = min_date_value
                self.advanced_filters["date_max"] = max_date_value
                
                # Analizza le estensioni
                exts = [e.strip() for e in extensions.get().split(",") if e.strip()]
                self.advanced_filters["extensions"] = [f".{e.lstrip('.')}" for e in exts]
                
                print(f"Filtri salvati: {self.advanced_filters}")  # Debug info
                dialog.destroy()
                
            except ValueError:
                messagebox.showerror("Errore", "Inserisci valori numerici validi per le dimensioni")
        
        # Creato un frame per i pulsanti nella stessa riga
        buttons_frame = ttk.Frame(dialog)
        buttons_frame.pack(fill=X, pady=10, padx=10)
        
        # Pulsante Annulla a sinistra
        ttk.Button(buttons_frame, text="Annulla", command=dialog.destroy).pack(side=RIGHT, padx=5)
        
        # Pulsante Salva a destra
        ttk.Button(buttons_frame, text="Salva", command=save_filters).pack(side=RIGHT, padx=5)

        dialog.update_idletasks()  # Aggiorna la finestra per ottenere le dimensioni corrette
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima per la finestra
        dialog.minsize(500, 400)

    def create_widgets(self):
        # Frame principale che conterrà tutto
        main_container = ttk.Frame(self.root)
        main_container.pack(fill=BOTH, expand=YES)
        
        # Intestazione (titolo e informazioni)
        header_frame = ttk.Frame(main_container, padding="10")
        header_frame.pack(fill=X)

        # Layout a tre colonne in una singola riga
        # 1. Tema a sinistra
        theme_frame = ttk.Frame(header_frame)
        theme_frame.pack(side=LEFT, fill=Y)

        ttk.Label(theme_frame, text="Tema:").pack(side=LEFT)
        themes = ttk.Style().theme_names()
        self.theme_combobox = ttk.Combobox(theme_frame, values=themes, width=15)
        self.theme_combobox.pack(side=LEFT, padx=5)
        self.theme_combobox.current(themes.index("darkly"))
        self.theme_combobox.bind("<<ComboboxSelected>>", lambda e: [ttk.Style().theme_use(self.theme_combobox.get()),self.update_theme_colors()])

        # 2. Titolo al centro
        title_frame = ttk.Frame(header_frame)
        title_frame.pack(side=LEFT, expand=True)

        title_label = ttk.Label(title_frame, text="File Search Tool.. Nucleo Perugia", 
                            font=("Helvetica", 14, "bold"))
        title_label.pack(anchor=CENTER)

        # 3. Data/ora e utente a destra
        datetime_frame = ttk.Frame(header_frame)
        datetime_frame.pack(side=RIGHT, fill=Y)

        datetime_label = ttk.Label(datetime_frame, textvariable=self.datetime_var, font=("Helvetica", 9))
        datetime_label.pack(side=RIGHT)

        # ==========================================================
        # SEZIONE DEI CONTROLLI DI RICERCA (organizzati per righe)
        # ==========================================================
        controls_frame = ttk.LabelFrame(main_container, text="Parametri di ricerca", padding=10)
        controls_frame.pack(fill=X, padx=10, pady=5)
        
        # ------------------------------------------------------
        # RIGA 1: Directory di ricerca
        # ------------------------------------------------------
        path_frame = ttk.Frame(controls_frame)
        path_frame.pack(fill=X, pady=5)
        
        path_label = ttk.Label(path_frame, text="Directory:", width=12, anchor=W)
        path_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(path_label, "Seleziona la directory per effettuare la ricerca dei file")

        self.path_entry = ttk.Entry(path_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        self.browse_btn = ttk.Button(path_frame, text="Sfoglia", command=self.browse_directory, width=10)
        self.browse_btn.pack(side=LEFT)
        
        # ------------------------------------------------------
        # RIGA 2: Parole chiave
        # ------------------------------------------------------
        keyword_frame = ttk.Frame(controls_frame)
        keyword_frame.pack(fill=X, pady=5)
        
        keyword_label = ttk.Label(keyword_frame, text="Parole chiave:", width=12, anchor=W)
        keyword_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(keyword_label, "Per la ricerca di più parole usa la virgola. Esempio: documento, fattura, contratto\n"
                                          "Attiva 'Parola intera' per cercare 'log' senza trovare 'login' o 'logo'")
        
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        info_label = ttk.Label(keyword_frame, text="Usa virgola per più termini", font=("", 8), foreground="gray")
        info_label.pack(side=LEFT, padx=5)
        
        # ------------------------------------------------------
        # RIGA 3: Opzioni base di ricerca (checkbox)
        # ------------------------------------------------------
        options_frame = ttk.Frame(controls_frame)
        options_frame.pack(fill=X, pady=5)
        
        options_label = ttk.Label(options_frame, text="Ricerca in:", width=12, anchor=W)
        options_label.pack(side=LEFT, padx=(0, 5))
        
        # Checkbox principali di ricerca allineate
        file_check = ttk.Checkbutton(options_frame, text="File", variable=self.search_files)
        file_check.pack(side=LEFT, padx=10)
        
        folder_check = ttk.Checkbutton(options_frame, text="Cartelle", variable=self.search_folders)
        folder_check.pack(side=LEFT, padx=10)
        
        content_check = ttk.Checkbutton(options_frame, text="Contenuti", variable=self.search_content)
        content_check.pack(side=LEFT, padx=10)
        self.create_tooltip(content_check, "Cerca le parole chiave anche all'interno dei file di testo")
        
        word_check = ttk.Checkbutton(options_frame, text="Parola intera", variable=self.whole_word_search)
        word_check.pack(side=LEFT, padx=10)
        self.create_tooltip(word_check, "Cerca solo parole intere (es. 'log' non troverà 'login' o 'logging')")
        
        sys_check = ttk.Checkbutton(options_frame, text="Escludi file di sistema", variable=self.exclude_system_files)
        sys_check.pack(side=LEFT, padx=10)
        self.create_tooltip(sys_check, "Esclude file di sistema come .exe, .dll, .sys, ecc. dalla ricerca")
        
        # ------------------------------------------------------
        # RIGA 4: Opzioni avanzate
        # ------------------------------------------------------
        advanced_frame = ttk.Frame(controls_frame)
        advanced_frame.pack(fill=X, pady=5)
        
        advanced_label = ttk.Label(advanced_frame, text="Opzioni:", width=12, anchor=W)
        advanced_label.pack(side=LEFT, padx=(0, 5))
        
        # Profondità di ricerca
        ttk.Label(advanced_frame, text="Profondità:").pack(side=LEFT)
        self.depth_spinbox = ttk.Spinbox(advanced_frame, from_=0, to=10, width=3)
        self.depth_spinbox.pack(side=LEFT, padx=5)
        self.depth_spinbox.set("0")
        ttk.Label(advanced_frame, text="(0 = illimitata)", font=("", 8), foreground="gray").pack(side=LEFT, padx=(0, 15))
        
        # Pulsanti di configurazione avanzata
        self.advanced_filters_btn = ttk.Button(advanced_frame, text="Filtri avanzati", 
                command=self.show_advanced_filters_dialog, width=15)
        self.advanced_filters_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(self.advanced_filters_btn, "Configura filtri per dimensione, data e tipo di file")

        self.exclusions_btn = ttk.Button(advanced_frame, text="Gestisci esclusi", 
                command=self.manage_exclusions, width=15)
        self.exclusions_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(self.exclusions_btn, "Gestisci le cartelle da escludere dalla ricerca")
        
        self.block_options_btn = ttk.Button(advanced_frame, text="Opzioni a Blocchi", 
                                    command=self.show_block_options, width=17)
        self.block_options_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(self.block_options_btn, "Configura le opzioni per la ricerca a blocchi (performance avanzata)")
        
        # ------------------------------------------------------
        # RIGA 5: Opzioni di performance
        # ------------------------------------------------------
        self.perf_frame = ttk.LabelFrame(controls_frame, text="Performance", padding=5)
        self.perf_frame.pack(fill=X, pady=5)

        # Prima riga di opzioni di performance
        perf_row1 = ttk.Frame(self.perf_frame)
        perf_row1.pack(fill=X, pady=2)

        # Timeout
        timeout_check = ttk.Checkbutton(perf_row1, text="Timeout ricerca", variable=self.timeout_enabled)
        timeout_check.pack(side=LEFT, padx=5)
        self.create_tooltip(timeout_check, "Interrompe automaticamente la ricerca dopo il tempo specificato")

        ttk.Label(perf_row1, text="Secondi:").pack(side=LEFT)
        timeout_spinbox = ttk.Spinbox(perf_row1, from_=10, to=600, width=4, textvariable=self.timeout_seconds)
        timeout_spinbox.pack(side=LEFT, padx=(0, 15))

        # Thread
        ttk.Label(perf_row1, text="Thread:").pack(side=LEFT)
        threads_spinbox = ttk.Spinbox(perf_row1, from_=1, to=16, width=2, textvariable=self.worker_threads)
        threads_spinbox.pack(side=LEFT, padx=(0, 15))
        self.create_tooltip(threads_spinbox, "Numero di thread paralleli per la ricerca (aumentando migliora la velocità ma usa più risorse)")

        # Dimensione massima file
        ttk.Label(perf_row1, text="Max file MB:").pack(side=LEFT)
        max_size_spinbox = ttk.Spinbox(perf_row1, from_=1, to=1000, width=4, textvariable=self.max_file_size_mb)
        max_size_spinbox.pack(side=LEFT, padx=(0, 15))
        self.create_tooltip(max_size_spinbox, "Dimensione massima in MB dei file di cui analizzare il contenuto")

        # Max risultati
        ttk.Label(perf_row1, text="Max risultati:").pack(side=LEFT)
        max_results_spinbox = ttk.Spinbox(perf_row1, from_=500, to=50000, width=6, textvariable=self.max_results)
        max_results_spinbox.pack(side=LEFT)
        self.create_tooltip(max_results_spinbox, "Numero massimo di risultati da mostrare")

        # Checkbox indicizzazione
        index_check = ttk.Checkbutton(perf_row1, text="Indicizzazione", variable=self.use_indexing)
        index_check.pack(side=LEFT, padx=15)
        self.create_tooltip(index_check, "Utilizza l'indicizzazione per velocizzare ricerche future")

        # Checkbox errori di permesso
        perm_check = ttk.Checkbutton(perf_row1, text="Ignora errori di permesso", variable=self.skip_permission_errors)
        perm_check.pack(side=LEFT)
        self.create_tooltip(perm_check, "Continua la ricerca anche quando alcune cartelle non possono essere lette")
        
        # ------------------------------------------------------
        # RIGA 6: Pulsanti di azione
        # ------------------------------------------------------
        button_frame = ttk.Frame(controls_frame)
        button_frame.pack(fill=X, pady=(10, 5))
        
        # Frame per centrare i pulsanti
        center_frame = ttk.Frame(button_frame)
        center_frame.pack(side=TOP)
        
        # Pulsante di ricerca (principale)
        self.search_button = ttk.Button(center_frame, text="CERCA", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=15)
        self.search_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.search_button, "Avvia la ricerca con i criteri specificati")
        
        # Pulsante per interrompere la ricerca
        self.stop_button = ttk.Button(center_frame, text="Interrompi ricerca",
                                    command=self.stop_search_process,
                                    style="danger.TButton", width=20,
                                    state="disabled")
        self.stop_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.stop_button, "Ferma immediatamente la ricerca in corso")
        
        # Pulsante per pulire i campi di ricerca
        self.clear_btn = ttk.Button(center_frame, text="Pulisci campi", 
                command=lambda: [self.search_path.set(""), self.keywords.set("")],
                style="secondary.Outline.TButton", width=15)
        self.clear_btn.pack(side=LEFT, padx=10)
        self.create_tooltip(self.clear_btn, "Cancella i campi di ricerca")
        
        # Pulsante admin solo su Windows
        if os.name == 'nt':
            self.admin_button = ttk.Button(center_frame, text="Avvia come Admin", 
                                command=self.restart_as_admin,
                                style="info.Outline.TButton", width=20)
            self.admin_button.pack(side=LEFT, padx=10)
            
            # Disabilita il pulsante se l'app è già avviata come amministratore
            if self.is_admin:
                self.admin_button.config(state="disabled")
                self.create_tooltip(self.admin_button, "L'applicazione è già in esecuzione come amministratore")
            else:
                self.create_tooltip(self.admin_button, "Riavvia l'applicazione con privilegi di amministratore per accedere a tutte le cartelle")
        
        # ==========================================================
        # SEZIONE INFORMATIVA
        # ==========================================================
        info_container = ttk.Frame(main_container)
        info_container.pack(fill=X, padx=10, pady=5)

        # Ridisegno completo della sezione stato ricerca per renderla fissa e più larga
        status_frame = ttk.LabelFrame(info_container, text="Stato ricerca", padding=5)
        status_frame.pack(fill=X, pady=5)  # Riempie tutto lo spazio orizzontale disponibile

        # Progress bar
        self.progress_bar = ttk.Progressbar(status_frame, mode='determinate')
        self.progress_bar.pack(fill=X, pady=5)

        # Layout a due colonne per le informazioni di stato
        status_grid = ttk.Frame(status_frame)
        status_grid.pack(fill=X, expand=YES)
        status_grid.columnconfigure(1, weight=1)  # La seconda colonna si espande

        # Riga 1: Status
        ttk.Label(status_grid, text="Analisi:", font=("", 9, "bold")).grid(row=0, column=0, sticky=W, padx=5, pady=2)
        self.status_label = ttk.Label(status_grid, text="In attesa...", wraplength=800)
        self.status_label.grid(row=0, column=1, sticky=W+E, padx=5, pady=2)

        # Riga 2: File analizzati
        ttk.Label(status_grid, text="File analizzati:", font=("", 9, "bold")).grid(row=1, column=0, sticky=W, padx=5, pady=2)
        self.analyzed_files_label = ttk.Label(status_grid, text="Nessuna ricerca avviata", wraplength=800)
        self.analyzed_files_label.grid(row=1, column=1, sticky=W+E, padx=5, pady=2)

        # Contenitore per le informazioni temporali e di spazio disco (sulla stessa riga)
        stats_container = ttk.Frame(info_container)
        stats_container.pack(fill=X, expand=YES, pady=5)

        # Frame per i tempi (metà sinistra)
        time_frame = ttk.LabelFrame(stats_container, text="Informazioni temporali", padding=5)
        time_frame.pack(side=LEFT, fill=X, expand=YES, padx=(0, 5))

        # Grid per i tempi con layout migliorato
        time_grid = ttk.Frame(time_frame)
        time_grid.pack(fill=X, pady=2)
        time_grid.columnconfigure(1, weight=1)
        time_grid.columnconfigure(3, weight=1)
        time_grid.columnconfigure(5, weight=1)

        ttk.Label(time_grid, text="Avvio:").grid(row=0, column=0, sticky=W, padx=5)
        self.start_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.start_time_label.grid(row=0, column=1, sticky=W, padx=5)

        ttk.Label(time_grid, text="Fine:").grid(row=0, column=2, sticky=W, padx=5)
        self.end_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.end_time_label.grid(row=0, column=3, sticky=W, padx=5)

        ttk.Label(time_grid, text="Durata:").grid(row=0, column=4, sticky=W, padx=5)
        self.total_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.total_time_label.grid(row=0, column=5, sticky=W, padx=5)

        # Frame per info disco (metà destra)
        disk_info_frame = ttk.LabelFrame(stats_container, text="Spazio disco", padding=5)
        disk_info_frame.pack(side=LEFT, fill=X, expand=YES)

        # Grid per info disco
        disk_grid = ttk.Frame(disk_info_frame)
        disk_grid.pack(fill=X, pady=2)

        ttk.Label(disk_grid, text="Directory:").grid(row=0, column=0, sticky=W, padx=5)
        ttk.Label(disk_grid, textvariable=self.dir_size_var, font=("", 9, "bold")).grid(row=0, column=1, sticky=W, padx=5)

        ttk.Label(disk_grid, text="Disco:").grid(row=0, column=2, sticky=W, padx=5)
        ttk.Label(disk_grid, textvariable=self.total_disk_var, font=("", 9, "bold")).grid(row=0, column=3, sticky=W, padx=5)

        ttk.Label(disk_grid, text="Usato:").grid(row=1, column=0, sticky=W, padx=5)
        ttk.Label(disk_grid, textvariable=self.used_disk_var, font=("", 9, "bold")).grid(row=1, column=1, sticky=W, padx=5)

        ttk.Label(disk_grid, text="Libero:").grid(row=1, column=2, sticky=W, padx=5)
        ttk.Label(disk_grid, textvariable=self.free_disk_var, font=("", 9, "bold")).grid(row=1, column=3, sticky=W, padx=5)
        
        # ==========================================================
        # SEZIONE RISULTATI (mantenuta come richiesto)
        # ==========================================================
        results_container = ttk.LabelFrame(main_container, text="Risultati di ricerca", padding=10)
        results_container.pack(fill=BOTH, expand=YES, padx=10, pady=(5, 10))
        
        # Frame per i pulsanti di azione sui risultati
        actions_frame = ttk.Frame(results_container)
        actions_frame.pack(fill=X, pady=(0, 5))
        
        # Pulsanti per la selezione
        selection_frame = ttk.Frame(actions_frame)
        selection_frame.pack(side=LEFT)
        
        select_all_btn = ttk.Button(selection_frame, text="Seleziona tutto", command=self.select_all)
        select_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(select_all_btn, "Seleziona tutti i risultati nella lista")
        
        deselect_all_btn = ttk.Button(selection_frame, text="Deseleziona tutto", command=self.deselect_all)
        deselect_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(deselect_all_btn, "Deseleziona tutti i risultati")
        
        invert_sel_btn = ttk.Button(selection_frame, text="Inverti selezione", command=self.invert_selection)
        invert_sel_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(invert_sel_btn, "Inverte la selezione corrente")
        
        # Pulsanti per le azioni
        action_frame = ttk.Frame(actions_frame)
        action_frame.pack(side=RIGHT)

        self.copy_button = ttk.Button(action_frame, text="Copia selezionati",
                                    command=self.copy_selected,
                                    style="TButton")
        self.copy_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.copy_button, "Copia i file selezionati nella directory specificata")
        
        self.compress_button = ttk.Button(action_frame, text="Comprimi selezionati",
                                        command=self.compress_selected,
                                        style="TButton")
        self.compress_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.compress_button, "Comprimi i file selezionati in un archivio ZIP")

        # TreeView con scrollbar
        treeview_container = ttk.Frame(results_container)
        treeview_container.pack(fill=BOTH, expand=True)

        # Creazione della TreeView
        self.results_list = ttk.Treeview(treeview_container, selectmode="extended",
                                    columns=("type", "author", "size", "modified", "created", "path"),
                                    show="headings")

        # Creazione delle scrollbar
        vsb = ttk.Scrollbar(treeview_container, orient="vertical", command=self.results_list.yview)
        hsb = ttk.Scrollbar(treeview_container, orient="horizontal", command=self.results_list.xview)

        # Configurazione delle scrollbar nella TreeView
        self.results_list.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

        # Posizionamento con grid
        self.results_list.grid(column=0, row=0, sticky='nsew')
        vsb.grid(column=1, row=0, sticky='ns')
        hsb.grid(column=0, row=1, sticky='ew')

        # Configura grid layout per l'espansione
        treeview_container.grid_columnconfigure(0, weight=1)
        treeview_container.grid_rowconfigure(0, weight=1)

        # Impostazione delle intestazioni delle colonne
        self.results_list.heading("type", text="Tipo")
        self.results_list.heading("author", text="Nome")  # Rinominato da "Autore" a "Nome"
        self.results_list.heading("size", text="Dimensione")
        self.results_list.heading("modified", text="Modificato")
        self.results_list.heading("created", text="Creato")
        self.results_list.heading("path", text="Percorso")

        # Impostazione delle larghezze delle colonne
        self.results_list.column("type", width=100, minwidth=50, stretch=NO)
        self.results_list.column("author", width=150, minwidth=80, stretch=NO)
        self.results_list.column("size", width=100, minwidth=80, stretch=NO) 
        self.results_list.column("modified", width=150, minwidth=120, stretch=NO)
        self.results_list.column("created", width=150, minwidth=120, stretch=NO)
        self.results_list.column("path", width=800, minwidth=200, stretch=YES)  # stretch=YES per utilizzare lo spazio disponibile
        
        # Aggiungi binding per l'evento di doppio clic
        self.results_list.bind("<Double-1>", self.open_file_location)
        self.create_tooltip(self.results_list, "Doppio click per aprire la posizione del file nel File Explorer")

        # Applica stili alle righe
        self.update_theme_colors()

        # Auto-focus sull'entry del percorso all'avvio
        self.path_entry.focus_set()

    def create_tooltip(self, widget, text, delay=500, fade=True):
        """Crea tooltip con ritardo, effetti di dissolvenza e larghezza automatica"""
        
        tooltip_timer = None
        fade_timer = None
        
        def show_tooltip():
            x = widget.winfo_rootx() + 25
            y = widget.winfo_rooty() + 25
            
            tooltip = ttk.Toplevel(widget)
            tooltip.wm_overrideredirect(True)
            tooltip.attributes("-alpha", 0.0)  # Inizia invisibile
            tooltip.attributes("-topmost", True)  # Mantiene sopra altre finestre
            
            # Crea un frame con bordo e padding
            frame = ttk.Frame(tooltip, borderwidth=1, relief="solid", padding=10)
            frame.pack(fill="both", expand=True)
            
            # Calcolo della larghezza del testo
            font = ("Segoe UI", 9)  # Puoi modificare questo per usare un altro font
            
            # Calcola la lunghezza approssimativa del testo
            temp_label = tk.Label(font=font)
            
            # Determina se il testo necessita di essere diviso in più righe
            lines = text.split("\n")
            max_line_width = 0
            
            for line in lines:
                temp_label.configure(text=line)
                line_width = temp_label.winfo_reqwidth()
                max_line_width = max(max_line_width, line_width)
            
            # Limita la larghezza massima a 500 pixel
            max_width = min(max_line_width, 500)
            
            wraplength = max_width
                
            # Crea l'etichetta con il testo
            label = ttk.Label(frame, text=text, justify="left", 
                            wraplength=wraplength, font=font)
            label.pack(fill="both", expand=True)
            
            # Distruggi il label temporaneo
            temp_label.destroy()
            
            # Aggiorna immediatamente per calcolare le dimensioni
            tooltip.update_idletasks()
            
            # Posiziona il tooltip in modo che non vada fuori dallo schermo
            tooltip_width = tooltip.winfo_reqwidth()
            tooltip_height = tooltip.winfo_reqheight()
            
            screen_width = tooltip.winfo_screenwidth()
            screen_height = tooltip.winfo_screenheight()
            
            # Aggiusta la posizione se il tooltip esce dallo schermo
            if x + tooltip_width > screen_width:
                x = screen_width - tooltip_width - 10
            
            if y + tooltip_height > screen_height:
                y = screen_height - tooltip_height - 10
                
            # Imposta la posizione finale
            tooltip.wm_geometry(f"+{x}+{y}")
            
            widget._tooltip = tooltip
            
            if fade:
                # Effetto dissolvenza in entrata
                def fade_in(alpha=0.0):
                    if not hasattr(widget, "_tooltip"):
                        return
                    
                    tooltip.attributes("-alpha", alpha)
                    if alpha < 1.0:
                        nonlocal fade_timer
                        fade_timer = widget.after(20, lambda: fade_in(alpha + 0.1))
                
                fade_in()
        
        def enter(event):
            nonlocal tooltip_timer
            # Avvia il timer per mostrare il tooltip dopo un certo ritardo
            tooltip_timer = widget.after(delay, show_tooltip)
        
        def leave(event):
            nonlocal tooltip_timer, fade_timer
            
            # Cancella il timer se esiste
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
                tooltip_timer = None
            
            # Cancella il timer di dissolvenza se esiste
            if fade_timer:
                widget.after_cancel(fade_timer)
                fade_timer = None
            
            # Rimuovi il tooltip se esiste
            if hasattr(widget, "_tooltip"):
                if fade:
                    # Effetto dissolvenza in uscita
                    def fade_out(alpha=1.0):
                        if alpha <= 0 or not hasattr(widget, "_tooltip"):
                            if hasattr(widget, "_tooltip"):
                                widget._tooltip.destroy()
                                del widget._tooltip
                        else:
                            widget._tooltip.attributes("-alpha", alpha)
                            nonlocal fade_timer
                            fade_timer = widget.after(20, lambda: fade_out(alpha - 0.1))
                    
                    fade_out()
                else:
                    widget._tooltip.destroy()
                    del widget._tooltip
        
        # Gestisce anche il caso in cui il widget venga distrutto
        def on_destroy(event):
            nonlocal tooltip_timer, fade_timer
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
            if fade_timer and hasattr(widget, "_tooltip"):
                widget.after_cancel(fade_timer)
                widget._tooltip.destroy()
        
        widget.bind("<Enter>", enter)
        widget.bind("<Leave>", leave)
        widget.bind("<Destroy>", on_destroy)

    def select_all(self):
        self.results_list.selection_set(self.results_list.get_children())
        
    def deselect_all(self):
        self.results_list.selection_remove(self.results_list.get_children())
        
    def invert_selection(self):
        all_items = self.results_list.get_children()
        selected_items = self.results_list.selection()
        self.results_list.selection_remove(selected_items)
        to_select = set(all_items) - set(selected_items)
        for item in to_select:
            self.results_list.selection_add(item)

# Funzione principale per eseguire l'applicazione
def main():
    import sys
    root = ttk.Window(themename="darkly")
    app = FileSearchApp(root)
    
    # Verifica se ci sono argomenti da linea di comando
    if len(sys.argv) > 1:
        # Se c'è un percorso fornito come argomento, lo imposta come percorso di ricerca
        app.search_path.set(sys.argv[1])
    
    root.mainloop()

if __name__ == "__main__":
    main()
