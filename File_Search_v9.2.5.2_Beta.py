# Importazioni essenziali per l'avvio
import io
import os
import tkinter as tk
import ttkbootstrap as ttk
from ttkbootstrap.constants import *
from tkinter import filedialog, messagebox, BooleanVar, StringVar, IntVar
import threading
import queue
from datetime import datetime
import getpass
import csv
import shutil
import zipfile
import traceback
import time
import concurrent.futures
import mimetypes
import signal
import re
import subprocess
import functools
import psutil
import gc
import os
import time
# Elenco di librerie da installare se mancanti
missing_libraries = []

# Definizione del decoratore con debug migliorato
def error_handler(func):
    @functools.wraps(func)
    def wrapper(self, *args, **kwargs):
        try:
            return func(self, *args, **kwargs)
        except Exception as e:
            error_type = type(e).__name__
            error_traceback = traceback.format_exc()
            
            # Registra l'errore con il tipo specifico
            self.log_error(
                f"Errore nell'esecuzione di {func.__name__}: [{error_type}]", 
                exception=e, 
                location=func.__name__, 
                traceback=error_traceback
            )
            
            # Aggiorna immediatamente la visualizzazione del debug log se è aperta
            if hasattr(self, 'debug_window') and self.debug_window and hasattr(self.debug_window, 'winfo_exists') and self.debug_window.winfo_exists():
                try:
                    # Assumendo che update_log_display sia un metodo accessibile
                    if hasattr(self.debug_window, 'update_log_display'):
                        self.debug_window.update_log_display()
                except:
                    pass
            
            # Mostra una finestra di errore all'utente
            import tkinter.messagebox as messagebox
            messagebox.showerror(
                "Errore applicazione", 
                f"Si è verificato un errore durante {func.__name__}.\n"
                f"Tipo errore: {error_type}\n"
                f"Dettagli: {str(e)}\n\n"
                f"L'errore è stato registrato nel log di debug."
            )
            return None
    return wrapper

# Dizionario per tracciare il supporto alle librerie - sarà popolato in seguito
file_format_support = {
    "docx": False, "pdf": False, "pptx": False, "excel": False,
    "odt": False, "rtf": False, "xls": False, "doc": False,
    "ods": False, "odp": False, "epub": False, "mobi": False, 
    "tex": True, "rst": True, "sqlite": True, "mdb": False,
    "odb": True, "tsv": True, "dbf": False, "dif": True,
    "executable": True, "code_files": True
}

class FileSearchApp:
    @error_handler
    def __init__(self, root):
        self.root = root
        self.root.title("File Search Tool V9.2.5 Beta Forensics G.di F.")
        
        # Imposta subito il debug mode per poter loggare
        self.debug_mode = True
        self.debug_log = []
        # Aggiungi questa riga per inizializzare current_user
        self.current_user = getpass.getuser()
        
        # Inizializza tutte le variabili in un passaggio
        self._init_essential_variables()
        self._init_remaining_variables()
        
        # Crea l'intera interfaccia in una volta sola
        self.create_widgets()
        
        # Applica il tema una sola volta
        self.update_theme_colors("dark")
        
        # Esegui attività di background dopo un breve ritardo
        self.root.after(500, self._delayed_startup_tasks)

    @error_handler
    def create_base_interface(self):
        """Crea solo l'interfaccia essenziale per un avvio veloce"""
        # Forza l'aggiornamento dell'interfaccia
        self.root.update_idletasks()

    @error_handler
    def complete_initialization(self):
        """Completa l'inizializzazione dell'applicazione"""
        try:
            # Rimuovi i componenti temporanei
            if hasattr(self, 'status_label'):
                self.status_label.destroy()
            if hasattr(self, 'progress_bar'):
                self.progress_bar.stop()
                self.progress_bar.destroy()
            
            # Inizializza le variabili rimanenti
            self._init_remaining_variables()
            
            # Crea l'interfaccia completa
            self.create_widgets()
            
            # Esegui le altre operazioni iniziali
            self._delayed_startup_tasks()
            
            # Avvia il monitoraggio della memoria
            self.root.after(30000, self.monitor_memory_usage)  # Primo controllo dopo 30 secondi
            
        except Exception as e:
            # In caso di errore, mostra un messaggio e prova a ripristinare l'applicazione
            messagebox.showerror("Errore di inizializzazione", 
                            f"Si è verificato un errore durante l'avvio: {str(e)}")
            # Tenta comunque di creare l'interfaccia
            try:
                self.create_widgets()
            except:
                pass

    @error_handler  
    def _init_essential_variables(self):
        """Inizializza solo le variabili essenziali per l'avvio"""
        # Inizializza datetime_var subito all'inizio per evitare errori di sequenza
        self.datetime_var = StringVar()
        
        # Variabili principali per la ricerca
        self.search_content = BooleanVar(value=True)
        self.search_path = StringVar()
        self.keywords = StringVar()
        self.search_results = []
        self.search_files = BooleanVar(value=True)
        self.search_folders = BooleanVar(value=True)
        self.is_searching = False
        self.progress_queue = queue.Queue()
        self.search_depth = StringVar(value="base") 

        # Variabili per data/ora e utente (datetime_var già inizializzato)
        self.user_var = StringVar(value=getpass.getuser())
        
        # Variabili essenziali per l'interfaccia
        self.ignore_hidden = BooleanVar(value=True)
        self.search_executor = None
        self.exclude_system_files = BooleanVar(value=True)
        self.whole_word_search = BooleanVar(value=False)
        self.dir_size_calculation = StringVar(value="disabilitato")
        
        # Variabili per la visualizzazione
        self.directory_calculation_enabled = False
        self.dir_size_var = StringVar(value="")
        self.total_disk_var = StringVar(value="")
        self.used_disk_var = StringVar(value="")
        self.free_disk_var = StringVar(value="")
        
        # Inizializza tutte le variabili utilizzate nel caricamento delle impostazioni
        self.timeout_enabled = tk.BooleanVar(value=False)
        self.timeout_seconds = tk.IntVar(value=3600)
        self.max_files_to_check = tk.IntVar(value=100000)
        self.max_results = tk.IntVar(value=50000)
        self.worker_threads = tk.IntVar(value=4)
        self.max_file_size_mb = tk.IntVar(value=100)
        self.use_indexing = tk.BooleanVar(value=True)
        self.skip_permission_errors = tk.BooleanVar(value=True)
        self.auto_memory_management = True
        self.memory_usage_percent = 75

        # Add a queue for debug logs
        self.debug_logs_queue = queue.Queue(maxsize=5000)  # Limit to 5000 entries to avoid memory issues
        
        # Aggiungi questa riga: lista permanente per i log completi dalla creazione dell'app
        self.complete_debug_log_history = []

    @error_handler
    def _init_remaining_variables(self):
        """Inizializza le variabili non essenziali per l'avvio"""
        # Variabili per la ricerca a blocchi
        self.max_files_per_block = IntVar(value=1000)
        self.max_parallel_blocks = IntVar(value=4)
        self.prioritize_user_folders = BooleanVar(value=True)
        self.block_size_auto_adjust = BooleanVar(value=True)
        
        # Configurazioni per la gestione della RAM
        self.auto_memory_management = True  # Attiva per default
        self.memory_usage_percent = 75      # Default 75% della RAM disponibile
        self.total_ram = psutil.virtual_memory().total

        # Variabili aggiuntive
        self.search_start_time = None
        self.stop_search = False
        self.max_depth = 0  # 0 = illimitato
        self.last_search_params = {}
        self.search_history = []
        
        # Filtri avanzati
        self.advanced_filters = {
            "size_min": 0,
            "size_max": 0,
            "date_min": None,
            "date_max": None,
            "extensions": []
        }
        
        # Opzioni per la gestione dei permessi
        self.skip_permission_errors = BooleanVar(value=True)
        self.excluded_paths = []
        self.load_settings_from_file()
        
        # Directory problematiche da escludere automaticamente
        self.problematic_dirs = [
            "Client Active Directory Rights Management Services",
            "Windows Resource Protection",
            "Windows Defender",
            "Microsoft Office",
            "System Volume Information",
            "$Recycle.Bin",
            "$WINDOWS.~BT",
            "$Windows.~WS"
        ]

        # Verifica dei privilegi di amministratore
        self.is_admin = False
        try:
            import ctypes
            self.is_admin = ctypes.windll.shell32.IsUserAnAdmin() != 0
        except:
            self.is_admin = False

        # Variabili per limiti di tempo e altre ottimizzazioni
        self.timeout_enabled = BooleanVar(value=False)
        self.timeout_seconds = IntVar(value=3600)
        self.max_files_to_check = IntVar(value=100000)
        self.max_results = IntVar(value=50000)
        self.chunk_size = 8192
        self.max_file_size_mb = IntVar(value=100)
        self.worker_threads = IntVar(value=min(8, os.cpu_count() or 4))
        self.use_indexing = BooleanVar(value=True)
        self.search_index = {}
        
        # Lista di estensioni di file di sistema da escludere dalla ricerca nei contenuti
        self.system_file_extensions = [
            # File eseguibili e librerie
            '.exe', '.dll', '.sys', '.drv', '.ocx', '.vxd', '.com', '.bat', '.cmd', '.scr', '.app', '.dylib', '.exp', '.bpl',
            # ... resto delle estensioni ... 
        ]

        # Percorso del log dei file saltati
        self.skipped_files_log_path = os.path.join(os.path.expanduser("~"), "skipped_files_log.txt")
        
        # Aggiorna l'orario
        self.update_datetime()
        pass

    @error_handler
    def _delayed_startup_tasks(self):
        """Esegue attività non essenziali all'avvio"""
        # Esegui queste operazioni in background
        threading.Thread(target=self._background_tasks, daemon=True).start()

    @error_handler
    def _background_tasks(self):
        """Esegue operazioni in background"""
        try:
            # Verifica librerie disponibili
            global file_format_support, missing_libraries
            
            # Controlla ogni libreria
            for module_name, format_key, import_name in [
                ("docx", "docx", "python-docx"),
                ("PyPDF2", "pdf", "PyPDF2"),
                # ... e così via per le altre librerie
            ]:
                try:
                    __import__(module_name)
                    file_format_support[format_key] = True
                except ImportError:
                    missing_libraries.append(import_name)
            
            # Mostra notifica dopo un ritardo
            if missing_libraries:
                self.root.after(2000, self.check_and_notify_missing_libraries)
            
            # Aggiorna informazioni disco senza calcolare dimensione directory
            if hasattr(self, 'search_path'):
                path = self.search_path.get()
                if path and os.path.exists(path):
                    self.update_disk_info(path, calculate_dir_size=False)
            
        except Exception as e:
            self.log_debug(f"Errore nelle operazioni in background: {str(e)}")

    @error_handler
    def _create_minimal_interface(self):
        """Crea solo i componenti essenziali dell'interfaccia"""
        # Frame principale che conterrà tutto
        self.main_container = ttk.Frame(self.root)
        self.main_container.pack(fill=BOTH, expand=YES)
                
        # Controlli minimi per la ricerca
        self.controls_frame = ttk.LabelFrame(self.main_container, text="Parametri di ricerca", padding=10)
        self.controls_frame.pack(fill=X, padx=10, pady=5)
        
        # Solo i controlli essenziali di ricerca
        self._create_essential_search_controls()
        
        # Barra di stato per informazioni
        status_frame = ttk.Frame(self.main_container, padding=5)
        status_frame.pack(fill=X, padx=10)
        
        self.status_label = ttk.Label(status_frame, text="Inizializzazione in corso...", font=("", 9))
        self.status_label.pack(side=LEFT)
        
        self.progress_bar = ttk.Progressbar(status_frame, mode='indeterminate')
        self.progress_bar.pack(side=LEFT, padx=10, fill=X, expand=YES)
        self.progress_bar.start(10)  # Avvia l'animazione di caricamento
        
        # Placeholder per i risultati
        self.results_container = ttk.LabelFrame(self.main_container, text="Risultati di ricerca", padding=10)
        self.results_container.pack(fill=BOTH, expand=YES, padx=10, pady=(5, 10))
        
        # Messaggio di caricamento
        self.loading_label = ttk.Label(self.results_container, text="Caricamento interfaccia in corso...", font=("", 12))
        self.loading_label.pack(expand=YES, pady=50)

    @error_handler
    def _create_essential_search_controls(self):
        """Crea solo i controlli essenziali per la ricerca"""
        # RIGA 1: Directory di ricerca
        path_frame = ttk.Frame(self.controls_frame)
        path_frame.pack(fill=X, pady=5)
        
        path_label = ttk.Label(path_frame, text="Directory:", width=12, anchor=W)
        path_label.pack(side=LEFT, padx=(0, 5))
        
        self.path_entry = ttk.Entry(path_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        self.browse_btn = ttk.Button(path_frame, text="Sfoglia", command=self.browse_directory, width=10)
        self.browse_btn.pack(side=LEFT)
        
        # RIGA 2: Parole chiave
        keyword_frame = ttk.Frame(self.controls_frame)
        keyword_frame.pack(fill=X, pady=5)
        
        keyword_label = ttk.Label(keyword_frame, text="Parole chiave:", width=12, anchor=W)
        keyword_label.pack(side=LEFT, padx=(0, 5))
        
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        # RIGA 3: Controllo della profondità di ricerca
        depth_frame = ttk.Frame(self.controls_frame)
        depth_frame.pack(fill=X, pady=5)
        
        depth_label = ttk.Label(depth_frame, text="Profondità max:", width=12, anchor=W)
        depth_label.pack(side=LEFT, padx=(0, 5))
        
        self.depth_spinbox = ttk.Spinbox(depth_frame, from_=0, to=20, width=5)
        self.depth_spinbox.pack(side=LEFT, padx=5)
        self.depth_spinbox.set("5")  # Valore predefinito
        
        depth_info_label = ttk.Label(depth_frame, text="(0 = illimitato, consigliato 5-10)")
        depth_info_label.pack(side=LEFT, padx=5)
        
        # Aggiunge un callback per aggiornare max_depth quando cambia il valore dello spinbox
        self.depth_spinbox.bind("<FocusOut>", lambda e: setattr(self, "max_depth", 
                            int(self.depth_spinbox.get()) if self.depth_spinbox.get().isdigit() else 5))
        
        # Solo il pulsante di ricerca
        button_frame = ttk.Frame(self.controls_frame)
        button_frame.pack(fill=X, pady=(10, 5))
        
        center_frame = ttk.Frame(button_frame)
        center_frame.pack(side=TOP)
        
        self.search_button = ttk.Button(center_frame, text="CERCA", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=15)
        self.search_button.pack(side=LEFT, padx=10)
        
        # Disabilita il pulsante fino al caricamento completo
        self.search_button["state"] = "disabled"

    @error_handler
    def _create_remaining_widgets(self):
        """Crea il resto dell'interfaccia in background"""
        try:
            # Aggiorna il messaggio di stato
            self.status_label["text"] = "Caricamento componenti avanzati..."
            
            # Rimuovi gli oggetti temporanei
            if hasattr(self, 'loading_label'):
                self.loading_label.destroy()
            
            # Ferma l'animazione della barra di progresso
            self.progress_bar.stop()
            self.progress_bar["mode"] = "determinate"
            
            # Crea tutti i componenti dell'interfaccia completa
            self.create_widgets()
            
            # Aggiorna lo stato
            self.status_label["text"] = "Applicazione pronta"
            
            # Riabilita il pulsante di ricerca
            self.search_button["state"] = "normal"
            
            # Imposta il focus sul campo di ricerca
            self.path_entry.focus_set()
            
            # Imposta il tema iniziale
            self.update_theme_colors("dark")
            
            self.log_debug("Interfaccia utente completamente caricata")
            
        except Exception as e:
            self.log_debug(f"Errore nel caricamento dell'interfaccia: {str(e)}")
            self.status_label["text"] = "Errore nel caricamento dell'interfaccia"

    @error_handler
    def _check_available_libraries(self):
        """Verifica la disponibilità delle librerie in background senza bloccare l'avvio"""
        # Esegui il controllo in un thread separato
        threading.Thread(target=self._async_check_libraries, daemon=True).start()

    @error_handler
    def _async_check_libraries(self):
        """Controlla le librerie in background senza bloccare l'interfaccia"""
        global file_format_support, missing_libraries
        
        # Funzione per verificare un singolo modulo
        def check_module(module_name, format_key, import_name=None):
            try:
                # Usa importlib per evitare bloccaggi
                import importlib
                importlib.import_module(module_name)
                file_format_support[format_key] = True
                self.log_debug(f"Supporto {format_key} attivato")
            except ImportError:
                missing_libraries.append(import_name or module_name)
                self.log_debug(f"Supporto {format_key} non disponibile")
        
        # Controlla ogni libreria
        check_module("docx", "docx", "python-docx")
        check_module("PyPDF2", "pdf", "PyPDF2")
        check_module("pptx", "pptx", "python-pptx")
        check_module("openpyxl", "excel", "openpyxl")
        check_module("striprtf.striprtf", "rtf", "striprtf")
        check_module("odf", "odt", "odfpy")
        check_module("ebooklib", "epub", "ebooklib")
        check_module("mobi", "mobi", "mobi")
        check_module("dbfread", "dbf", "dbfread")
        check_module("pyodbc", "mdb", "pyodbc")
        check_module("bs4", "epub_html", "beautifulsoup4")
        check_module("pefile", "pe_files", "pefile")
        
        # Controlla win32com (per i formati legacy di Office)
        if os.name == 'nt':  # Solo su Windows
            try:
                import win32com.client
                file_format_support["doc"] = True
                file_format_support["xls"] = True
                self.log_debug("Supporto formati legacy Office attivato (win32com)")
            except ImportError:
                missing_libraries.append("pywin32")
                self.log_debug("Supporto formati legacy Office non disponibile (manca pywin32)")
        
        # Controlla xlrd come alternativa per i file XLS
        try:
            import xlrd
            file_format_support["xls_native"] = True
            if not file_format_support.get("xls", False):
                file_format_support["xls"] = True
            self.log_debug("Supporto XLS attivato (xlrd)")
        except ImportError:
            if not file_format_support.get("xls", False):
                missing_libraries.append("xlrd")
                self.log_debug("Supporto XLS via xlrd non disponibile")
        
        # Controlla supporto ODT con librerie alternative
        try:
            from odf import opendocument
            file_format_support["odt"] = True
            self.log_debug("Supporto ODT attivato (libreria odf)")
        except ImportError:
            missing_libraries.append("odf")
            self.log_debug("Supporto ODT non disponibile")
        
        # Mostra notifica dopo un ritardo
        if missing_libraries:
            self.root.after(2000, self.check_and_notify_missing_libraries)

    @error_handler
    def process_file(self, file_path, keywords, search_content=True):
        """Processa un singolo file per verificare corrispondenze"""
        if self.stop_search:
            return None
        try:
            # Verifica nome file
            file_name = os.path.basename(file_path)
            matched = False
            
            # Verifica corrispondenze nel nome
            for keyword in keywords:
                if self.whole_word_search.get():
                    if self.is_whole_word_match(keyword, file_name):
                        matched = True
                        break
                elif keyword.lower() in file_name.lower():
                    matched = True
                    break
            
            content = ""
            # Verifica contenuto se richiesto e se non c'è già una corrispondenza nel nome
            if not matched and search_content and self.should_search_content(file_path):
                content = self.get_file_content(file_path)
                if content:
                    for keyword in keywords:
                        if self.whole_word_search.get():
                            if self.is_whole_word_match(keyword, content):
                                matched = True
                                break
                        elif keyword.lower() in content.lower():
                            matched = True
                            break
            
            if matched:
                # Verifica se il match è in un allegato di un file EMAIL (EML o MSG)
                _, ext = os.path.splitext(file_path)
                if ext.lower() in ['.eml', '.msg'] and search_content and "--- ALLEGATO" in content:
                    # Cerca il match dopo un'intestazione di allegato
                    attachment_sections = content.split("--- ALLEGATO")
                    for section in attachment_sections[1:]:  # Salta il primo che è l'intestazione email
                        for keyword in keywords:
                            if (self.whole_word_search.get() and self.is_whole_word_match(keyword, section)) or \
                            (not self.whole_word_search.get() and keyword.lower() in section.lower()):
                                # Match trovato in un allegato
                                self.log_debug(f"Match trovato in allegato di {file_path}")
                                return self.create_file_info(file_path, from_attachment=True)
                
                # Match normale (non in allegato o non in file EMAIL)
                return self.create_file_info(file_path)
                
        except Exception as e:
            self.log_debug(f"Errore nel processare il file {file_path}: {str(e)}")
            if self.debug_mode:
                import traceback
                self.log_debug(traceback.format_exc())
                
        return None

    @error_handler
    def process_file_with_timeout(self, file_path, keywords, search_content=True):
        """Process a file with timeout to prevent hanging"""
        # CORREZIONE: Rilevamento tipo file per ottimizzare timeout
        is_large_file = False
        is_binary_file = False
        is_network_file = False
        
        try:
            # Identifica se il file è su un percorso di rete
            if file_path.startswith('\\\\') or file_path.startswith('//'):
                is_network_file = True
                
            # Verifica la dimensione e imposta un flag per i file grandi
            if os.path.exists(file_path):
                file_size = os.path.getsize(file_path)
                # File sopra 5MB sono considerati grandi
                if file_size > 5 * 1024 * 1024:
                    is_large_file = True
                    
                # Rileva se è probabilmente un file binario
                ext = os.path.splitext(file_path)[1].lower()
                binary_extensions = ['.exe', '.dll', '.bin', '.obj', '.o', '.so', '.lib', '.sys', '.ocx']
                if ext in binary_extensions:
                    is_binary_file = True
        except:
            pass
        
        # CORREZIONE: Adatta il timeout in base al tipo di file
        timeout = 5.0  # Default 5 secondi
        if is_network_file:
            timeout = 10.0  # 10 secondi per file di rete
        if is_large_file:
            timeout += 5.0  # +5 secondi per file grandi
        if is_binary_file:
            timeout = min(timeout, 3.0)  # Limita a 3 secondi per file binari
        
        # CORREZIONE: Usa sistema a coda per i risultati invece di variabili condivise
        result_queue = queue.Queue()
        exception_queue = queue.Queue()
        completion_flag = threading.Event()
        
        # Create a thread to process the file
        def process_thread():
            try:
                result = self.process_file(file_path, keywords, search_content)
                result_queue.put(result)
            except Exception as e:
                exception_queue.put(e)
            finally:
                completion_flag.set()
        
        process_thread = threading.Thread(target=process_thread)
        process_thread.daemon = True
        process_thread.start()
        
        # Wait for the specified timeout
        completion = completion_flag.wait(timeout)
        
        if completion:
            # Processing completed normally
            if not exception_queue.empty():
                exception = exception_queue.get()
                self.log_debug(f"Eccezione nell'elaborazione del file {file_path}: {str(exception)}")
                return None
            return result_queue.get() if not result_queue.empty() else None
        else:
            # Processing timed out
            self.log_debug(f"Processing timed out for file: {file_path}")
            return None

    @error_handler
    def process_with_timeout(self, file_path, keywords, result, exception, processing_completed, search_content=True):
        try:
            # Verifica filtri di dimensione
            file_size = os.path.getsize(file_path)
            if (self.advanced_filters["size_min"] > 0 and file_size < self.advanced_filters["size_min"]) or \
            (self.advanced_filters["size_max"] > 0 and file_size > self.advanced_filters["size_max"]):
                result[0] = None
                return
            
            # Verifica filtri di data
            if self.advanced_filters["date_min"] or self.advanced_filters["date_max"]:
                mod_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                
                if self.advanced_filters["date_min"]:
                    min_date = datetime.strptime(self.advanced_filters["date_min"], "%d-%m-%Y")
                    if mod_time < min_date:
                        result[0] = None
                        return
                        
                if self.advanced_filters["date_max"]:
                    max_date = datetime.strptime(self.advanced_filters["date_max"], "%d-%m-%Y")
                    if mod_time > max_date:
                        result[0] = None
                        return
            
            # Verifica filtri estensione
            if self.advanced_filters["extensions"] and not any(file_path.lower().endswith(ext.lower()) 
                                                        for ext in self.advanced_filters["extensions"]):
                result[0] = None
                return
                
            if os.path.splitext(file_path)[1].lower() in ['.doc', '.xls']:
                self.progress_queue.put(("progress", self.progress_bar["value"]))  # Forza aggiornamento
                
            # Verifica corrispondenza nel nome file
            filename = os.path.basename(file_path)

            matched = False
            for keyword in keywords:
                # Verifica se la ricerca di parole intere è attivata
                if self.whole_word_search.get():
                    # Gestisce anche frasi con spazi usando la nuova funzione helper
                    if self.is_whole_word_match(keyword, filename):
                        matched = True
                        break
                # Se è un termine con spazi (contesto specifico)
                elif ' ' in keyword:
                    if keyword.lower() in filename.lower():
                        matched = True
                        break
                # Ricerca normale
                else:
                    if keyword.lower() in filename.lower():
                        matched = True
                        break
            
            if matched:
                if self.debug_mode and self.whole_word_search.get():
                    self.log_debug(f"Trovata corrispondenza per parola intera: '{keyword}' nel nome del file {os.path.basename(file_path)}")
                result[0] = self.create_file_info(file_path)
                return
            
            # Verifica contenuto se richiesto
            if search_content and self.should_search_content(file_path):
                max_size_bytes = self.max_file_size_mb.get() * 1024 * 1024
                
                # Salta file troppo grandi
                if file_size > max_size_bytes:
                    self.log_debug(f"File {file_path} troppo grande per l'analisi del contenuto")
                    result[0] = None
                    return
                    
                content = self.get_file_content(file_path)
                if content:
                    matched = False
                    for keyword in keywords:
                        # Verifica se la ricerca di parole intere è attivata
                        if self.whole_word_search.get():
                            # Gestisce anche frasi con spazi usando la nuova funzione helper
                            if self.is_whole_word_match(keyword, content):
                                matched = True
                                break
                        # Se è un termine con spazi (contesto specifico)
                        elif ' ' in keyword:
                            if keyword.lower() in content.lower():
                                matched = True
                                break
                        # Ricerca normale
                        else:
                            if keyword.lower() in content.lower():
                                matched = True
                                break
                                
                    if matched:
                        if self.debug_mode and self.whole_word_search.get():
                            self.log_debug(f"Trovata corrispondenza per parola intera: '{keyword}' nel contenuto del file {os.path.basename(file_path)}")
                        result[0] = self.create_file_info(file_path)
                        return
            else:
                # Log per i file di sistema esclusi
                if search_content and os.path.splitext(file_path)[1].lower() in self.system_file_extensions:
                    self.log_debug(f"File di sistema escluso dall'analisi del contenuto: {file_path}")
                    
            result[0] = None
            
        except Exception as e:
            exception[0] = e
            self.log_debug(f"Errore durante l'elaborazione del file {file_path}: {str(e)}")
            result[0] = None
        finally:
            processing_completed[0] = True

    @error_handler
    def manage_memory(self):
        """Gestisce la memoria dell'applicazione in base alle impostazioni configurate dall'utente.
        Rilascia memoria quando necessario per mantenere le prestazioni ottimali."""
        try:
            # Verifica se la gestione automatica è disattivata dall'utente
            if not getattr(self, 'auto_memory_management', True):
                # Modalità manuale: rispetta solo il limite di memoria configurato
                try:
                    # Ottieni informazioni sulla memoria
                    total_ram = psutil.virtual_memory().total
                    process = psutil.Process(os.getpid())
                    current_memory = process.memory_info().rss
                    
                    # Calcola il limite basato sulla percentuale configurata
                    percent = getattr(self, 'memory_usage_percent', 75)
                    max_memory = total_ram * (percent / 100)
                    
                    # Verifica se superiamo la soglia
                    if current_memory > max_memory:
                        self.log_debug(f"Limite memoria ({percent}%) superato: {current_memory/(1024**2):.2f} MB / " +
                                    f"{max_memory/(1024**2):.2f} MB - Eseguendo pulizia manuale")
                        
                        # Pulizia in due fasi
                        gc.collect(0)  # Pulizia leggera iniziale
                        time.sleep(0.05)
                        gc.collect(2)  # Pulizia completa
                        
                        # Supporto per Python 3.9+
                        if hasattr(gc, 'collect_step'):
                            gc.collect_step()
                        
                        # Riduzione della cache interna se necessario
                        memory_after_gc = process.memory_info().rss
                        if memory_after_gc > max_memory * 0.95 and hasattr(self, 'search_results') and len(self.search_results) > 5000:
                            self.log_debug(f"Memoria ancora alta dopo GC, riduzione cache risultati da {len(self.search_results)} a 5000 elementi")
                            self.search_results = self.search_results[:5000]  # Mantieni i risultati più recenti
                            gc.collect()  # Richiama GC dopo la riduzione cache
                except Exception as e:
                    self.log_debug(f"Errore nella gestione manuale della memoria: {str(e)}")
            else:
                # Modalità automatica: gestione proattiva della memoria
                try:
                    # Ottimizza la frequenza del garbage collector in base al carico
                    is_searching = getattr(self, 'is_search_running', False)
                    if is_searching:
                        # Durante la ricerca, raccolta meno frequente per non interrompere le operazioni
                        gc.set_threshold(700, 10, 10)
                    else:
                        # Quando inattivo, raccolta più aggressiva
                        gc.set_threshold(100, 5, 5)
                    
                    # Ottieni informazioni sulla memoria
                    process = psutil.Process(os.getpid())
                    memory_usage = process.memory_info().rss / (1024 * 1024)  # MB
                    
                    # Determina la soglia dinamica di memoria
                    system_memory = psutil.virtual_memory().total / (1024 * 1024)  # Total system RAM in MB
                    threshold = max(200, min(1000, system_memory * 0.3))  # 30% della RAM di sistema con limiti min/max
                    
                    # Log utilizzo memoria (solo ogni 10 chiamate per non sovraccaricare i log)
                    self._memory_log_counter = getattr(self, '_memory_log_counter', 0) + 1
                    if self._memory_log_counter % 10 == 0:
                        self.log_debug(f"Gestione memoria automatica: Utilizzo corrente {memory_usage:.2f} MB, soglia {threshold:.2f} MB")
                    
                    # Verifica se superiamo la soglia
                    if memory_usage > threshold:
                        self.log_debug(f"Utilizzo memoria elevato: {memory_usage:.2f} MB. Esecuzione garbage collection.")
                        
                        # Pulizia progressiva
                        gc.collect(0)  # Garbage collection leggera
                        memory_after_gc0 = process.memory_info().rss / (1024 * 1024)
                        
                        if memory_after_gc0 > threshold * 0.9:
                            self.log_debug("Prima pulizia insufficiente, esecuzione pulizia completa.")
                            gc.collect(2)  # Garbage collection completa
                            
                            # Riduzione della cache se la memoria è ancora alta
                            memory_final = process.memory_info().rss / (1024 * 1024)
                            if memory_final > threshold and hasattr(self, 'search_results') and len(self.search_results) > 10000:
                                self.log_debug(f"Riduzione cache risultati da {len(self.search_results)} a 10000.")
                                self.search_results = self.search_results[-10000:]
                                gc.collect()  # Altra pulizia dopo la riduzione della cache
                except Exception as e:
                    self.log_debug(f"Errore nella gestione automatica della memoria: {str(e)}")
        except ImportError:
            # psutil non disponibile
            self.log_debug("psutil non disponibile, usando solo garbage collection standard.")
            gc.collect()
        except Exception as e:
            self.log_debug(f"Errore generale nella gestione della memoria: {str(e)}")
            # Fallback alla garbage collection standard
            try:
                gc.collect()
            except Exception as e2:
                self.log_debug(f"Anche il fallback di garbage collection è fallito: {str(e2)}")
        
        # Pianifica il prossimo controllo della memoria
        if hasattr(self, 'root') and self.root:
            self.root.after(30000, self.manage_memory)  # Controlla ogni 30 secondi
    
    @error_handler
    def monitor_memory_usage(self):
        """Monitoraggio periodico dell'utilizzo della memoria. Rileva situazioni critiche e interviene preventivamente."""
        try:
            # Esegui solo ogni 30 secondi per ridurre l'overhead
            if not hasattr(self, '_last_memory_check'):
                self._last_memory_check = time.time()
                return
                
            current_time = time.time()
            if current_time - self._last_memory_check < 30:
                return
                
            self._last_memory_check = current_time
            
            # Ottieni l'utilizzo di memoria del processo corrente
            process = psutil.Process()
            memory_info = process.memory_info()
            memory_used_mb = memory_info.rss / (1024**2)
            
            # Ottieni la memoria totale del sistema
            total_ram = psutil.virtual_memory().total / (1024**2)
            
            # Calcola la percentuale utilizzata rispetto al limite configurato
            memory_percent = getattr(self, 'memory_usage_percent', 75)
            memory_limit = total_ram * (memory_percent / 100)
            usage_ratio = memory_used_mb / memory_limit
            
            # Se superiamo l'80% del limite configurato, avvisa e forza una pulizia
            if usage_ratio > 0.8:
                self.log_debug(f"Memoria in uso al {usage_ratio*100:.1f}% del limite configurato ({memory_used_mb:.1f}MB/{memory_limit:.1f}MB)")
                
                # Azioni preventive progressive
                if usage_ratio > 0.9:
                    self.log_debug("ATTENZIONE: Utilizzo memoria CRITICO - Esecuzione pulizia di emergenza")
                    # Pulizia aggressiva
                    gc.collect(2)
                    time.sleep(0.1)
                    gc.collect(2)
                    
                    # Riduzione cache se necessario
                    if hasattr(self, 'search_results') and len(self.search_results) > 5000:
                        original_len = len(self.search_results)
                        self.search_results = self.search_results[:5000]
                        self.log_debug(f"Riduzione di emergenza della cache risultati: {original_len} → 5000 elementi")
                else:
                    # Pulizia standard
                    gc.collect()
                    
        except Exception as e:
            self.log_debug(f"Errore nel monitoraggio della memoria: {str(e)}")
            
        finally:
            # Riprogram,a il prossimo controllo
            if hasattr(self, 'root') and self.root:
                self.root.after(30000, self.monitor_memory_usage)

    @error_handler
    def run_with_timeout(self, func, args=(), kwargs={}, timeout_sec=10):
        """ Esegue una funzione con un timeout, evitando blocchi indefiniti Returns: (result, completed)
            - result: risultato della funzione o None
            - completed: True se completata, False se interrotta per timeout
        """
        result = [None]
        completed = [False]
        
        def target():
            try:
                result[0] = func(*args, **kwargs)
            except Exception as e:
                self.log_debug(f"Eccezione nella funzione con timeout: {str(e)}")
            finally:
                completed[0] = True
        
        thread = threading.Thread(target=target)
        thread.daemon = True
        thread.start()
        thread.join(timeout_sec)
        
        return result[0], completed[0]
    
    @error_handler
    def _get_all_descendants(self, widget):
        """Ottiene ricorsivamente tutti i widget discendenti"""
        descendants = []
        for child in widget.winfo_children():
            descendants.append(child)
            descendants.extend(self._get_all_descendants(child))
        return descendants
    
    @error_handler
    def disable_all_controls(self):
        """Disabilita tutti i controlli UI durante la ricerca"""
        try:
            # Input e campi di testo
            self.path_entry["state"] = "disabled"
            self.keyword_entry["state"] = "disabled"
            
            # Pulsante Sfoglia
            if hasattr(self, 'browse_btn'):
                self.browse_btn["state"] = "disabled"
            
            # Pulsante Pulisci campi
            if hasattr(self, 'clear_btn'):
                self.clear_btn["state"] = "disabled"
            
            # Combobox e spinbox
            if hasattr(self, 'theme_combobox'):
                self.theme_combobox["state"] = "disabled"
            
            if hasattr(self, 'depth_spinbox'):
                self.depth_spinbox["state"] = "disabled"
            
            # Pulsanti principali
            if hasattr(self, 'search_button'):
                self.search_button["state"] = "disabled"
            
            # Gestisci i pulsanti dei filtri
            for btn_name in ['advanced_filters_btn', 'exclusions_btn', 'block_options_btn']:
                if hasattr(self, btn_name):
                    getattr(self, btn_name)["state"] = "disabled"
            
            # Gestisci specificamente il pulsante admin
            if hasattr(self, 'admin_button') and not self.is_admin:
                self.admin_button["state"] = "disabled"

            # Disabilita il pulsante delle opzioni di performance
            if hasattr(self, 'perf_options_btn'):
                self.perf_options_btn["state"] = "disabled"
            
            # Pulsanti di azione sui risultati
            if hasattr(self, 'copy_button'):
                self.copy_button["state"] = "disabled"
            if hasattr(self, 'compress_button'):
                self.compress_button["state"] = "disabled"
                
        except Exception as e:
            # Registra l'errore senza interrompere l'esecuzione
            self.log_debug(f"Errore nella disabilitazione dei controlli: {str(e)}")

    @error_handler
    def enable_all_controls(self):
        """Riabilita tutti i controlli UI dopo la ricerca"""
        # Input e campi di testo
        self.path_entry["state"] = "normal"
        self.keyword_entry["state"] = "normal"
        
        # Pulsante Sfoglia
        if hasattr(self, 'browse_btn'):
            self.browse_btn["state"] = "normal"
        
        # Pulsante Pulisci campi
        if hasattr(self, 'clear_btn'):
            self.clear_btn["state"] = "normal"
        
        # Combobox e spinbox
        if hasattr(self, 'theme_combobox'):
            self.theme_combobox["state"] = "readonly"
        if hasattr(self, 'depth_spinbox'):
            self.depth_spinbox["state"] = "normal"
        
        # Pulsanti principali
        self.search_button["state"] = "normal"
        
        # Gestisci i pulsanti dei filtri
        for btn_name in ['advanced_filters_btn', 'exclusions_btn']:
            if hasattr(self, btn_name):
                getattr(self, btn_name)["state"] = "normal"

        # Gestisci i pulsanti  a blocchi
        if hasattr(self, 'block_options_btn'):
            self.block_options_btn["state"] = "normal"

        # Gestisci specificamente il pulsante admin
        if hasattr(self, 'admin_button') and not self.is_admin:
            self.admin_button["state"] = "normal"
        
        # Checkbox
        for widget in self.root.winfo_children():
            self._enable_checkbuttons_recursive(widget)
        
        # Pulsanti di azione sui risultati
        if hasattr(self, 'copy_button'):
            self.copy_button["state"] = "normal"
        if hasattr(self, 'compress_button'):
            self.compress_button["state"] = "normal"
        
        # Abilita il pulsante delle opzioni di performance
        if hasattr(self, 'perf_options_btn'):
            self.perf_options_btn["state"] = "normal"

    @error_handler     
    def _disable_checkbuttons_recursive(self, widget):
        """Disabilita ricorsivamente tutte le checkbox nei widget"""
        if isinstance(widget, ttk.Checkbutton):
            widget["state"] = "disabled"
        
        # Processa i widget figli ricorsivamente
        if hasattr(widget, 'winfo_children'):
            for child in widget.winfo_children():
                self._disable_checkbuttons_recursive(child)
    
    @error_handler
    def _enable_checkbuttons_recursive(self, widget):
        """Riabilita ricorsivamente tutte le checkbox nei widget"""
        if isinstance(widget, ttk.Checkbutton):
            widget["state"] = "normal"
        
        # Processa i widget figli ricorsivamente
        if hasattr(widget, 'winfo_children'):
            for child in widget.winfo_children():
                self._enable_checkbuttons_recursive(child)

    @error_handler     
    def show_content_search_warning(self):
        """Mostra un avviso se la ricerca nei contenuti è attivata"""
        if self.search_content.get():  # Se la checkbox è attivata
            message = "Hai attivato la ricerca nei contenuti dei file."
            
            # MODIFICA QUI: Aggiungi un log di debug per verificare il valore effettivo
            self.log_debug(f"Livello di ricerca attuale: {self.search_depth.get()}")
            
            # Avviso specifico per la ricerca profonda
            if hasattr(self, 'search_depth') and self.search_depth.get() == "profonda":
                message += "\n\nATTENZIONE: Hai selezionato la ricerca PROFONDA che include tutti i tipi di file!\n" + \
                        "Questo potrebbe rallentare notevolmente la ricerca, soprattutto con grandi quantità di file."
            else:
                message += "\n\nQuesta operazione può richiedere molto più tempo, soprattutto con grandi quantità di file."
                
            message += "\n\nVuoi procedere con la ricerca nei contenuti?"
            
            return messagebox.askyesno(
                "Attenzione - Ricerca nei contenuti", 
                message,
                icon="warning"
            )
        return True   # Se la ricerca nei contenuti non è attivata, procedi senza avviso   
    
    @error_handler
    def log_debug(self, message):
        """Sistema di logging unificato che supporta sia la visualizzazione live che il logging storico
        e filtra i messaggi non necessari e le ripetizioni sulla soglia di memoria.
        """
        # Inizializza il set di messaggi già loggati se non esiste
        if not hasattr(self, 'already_logged_messages'):
            self.already_logged_messages = set()
        
        # Filtra messaggi specifici relativi alle estensioni
        filter_prefixes = [
            "Estensioni caricate per modalità base:",
            "Estensioni caricate per modalità avanzata:",
            "Estensioni caricate per modalità profonda:"
        ]
        
        # Se il messaggio inizia con uno dei prefissi da filtrare, saltalo
        if any(message.startswith(prefix) for prefix in filter_prefixes):
            return
        
        # Verifica se è un messaggio sulla soglia di memoria
        if message.startswith("Soglia memoria calcolata:"):
            # Controlla se abbiamo già registrato un messaggio simile
            if "soglia_memoria" in self.already_logged_messages:
                return  # Salta questo messaggio perché è una ripetizione
            else:
                # Registra che abbiamo loggato un messaggio sulla soglia di memoria
                self.already_logged_messages.add("soglia_memoria")
        
        # Inizializza gli attributi necessari se mancanti
        if not hasattr(self, 'complete_debug_log_history'):
            self.complete_debug_log_history = []
        
        if not hasattr(self, 'debug_logs_queue'):
            self.debug_logs_queue = queue.Queue(maxsize=5000)
        
        if not hasattr(self, 'debug_log'):
            self.debug_log = []
        
        # Creazione timestamp
        timestamp_full = datetime.now().strftime('%Y-%m-%d %H:%M:%S.%f')[:-3]  # Con millisecondi
        timestamp_short = timestamp_full.split(' ')[1]  # Solo l'ora per i log concisi
        
        # Rilevamento automatico dei messaggi di errore
        if "error" in message.lower() or "exception" in message.lower() or "failed" in message.lower():
            # Formattazione speciale per gli errori
            log_message_full = f"[WARNING] {timestamp_full} - {message}"
            log_message_short = f"[{timestamp_short}] [WARNING] {message}"
        else:
            # Formattazione standard per i messaggi di debug
            log_message_full = f"{timestamp_full} - {message}"
            log_message_short = f"[{timestamp_short}] {message}"
        
        try:
            # Aggiungi al debug_log (prima implementazione)
            self.debug_log.append(log_message_full)
            
            # Limita la dimensione del log per evitare problemi di memoria
            if len(self.debug_log) > 5000:
                self.debug_log = self.debug_log[-5000:]
            
            # Aggiungi alla cronologia completa (seconda implementazione)
            self.complete_debug_log_history.append(log_message_short)
            
            # Limita anche questa cronologia se diventa troppo grande
            if len(self.complete_debug_log_history) > 10000:
                self.complete_debug_log_history = self.complete_debug_log_history[-10000:]
            
            # Aggiungi alla coda per la visualizzazione live con thread-safety
            with threading.Lock():
                if self.debug_logs_queue.full():
                    try:
                        self.debug_logs_queue.get_nowait()
                    except queue.Empty:
                        pass
                self.debug_logs_queue.put_nowait(log_message_short)
        except Exception as e:
            # In caso di errori nel logging, almeno prova a stampare sulla console
            print(f"Error in logging system: {str(e)}")
        
        # Se siamo in modalità debug, stampa anche sulla console
        if getattr(self, 'debug_mode', True):  # Default a True se debug_mode non è definito
            print(f"[DEBUG] {message}")

    @error_handler
    def register_interrupt_handler(self):
        """Registra il gestore degli interrupt (CTRL+C)"""
        def handle_interrupt(sig, frame):
            if self.is_searching:
                self.stop_search_process()
            else:
                self.root.quit()
        
        signal.signal(signal.SIGINT, handle_interrupt)

    @error_handler
    def log_error(self, message, exception=None, location=None, traceback=None):
        """Registra un errore nel log di debug con dettagli aggiuntivi"""
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        
        # Crea un messaggio di errore dettagliato
        error_message = f"[ERROR] {timestamp} - {message}"
        
        print(f"LOG_ERROR: {error_message[:100]}...")
    
        # Registra nel log di debug
        if hasattr(self, 'debug_log'):
            self.debug_log.append(error_message)
            print(f"Debug log size: {len(self.debug_log)}")
            
        # Aggiungi informazioni sulla posizione dell'errore
        if location:
            error_message += f" | Location: {location}"
        
        # Aggiungi dettagli sull'eccezione
        if exception:
            exc_type = type(exception).__name__
            exc_details = str(exception)
            
            error_message += f" | Exception: [{exc_type}]: {exc_details}"
        
        # Gestisci il traceback (priorità al traceback esplicito se fornito)
        traceback_info = ""
        if traceback:
            # Usa il traceback fornito esplicitamente
            error_message += f"\n--- Traceback ---\n{traceback}\n-----------------"
        elif exception:
            # Ottieni traceback se disponibile e non già fornito
            try:
                import traceback as tb_module
                tb_info = tb_module.format_exc().split('\n')
                # Prendi solo le righe più rilevanti del traceback
                if len(tb_info) > 3:
                    traceback_info = " | " + " > ".join(tb_info[-4:-1])
                    error_message += traceback_info
            except:
                pass
        
        # Registra nel log di debug
        if hasattr(self, 'debug_log'):
            self.debug_log.append(error_message)
            
            # Limita la dimensione del log per evitare problemi di memoria
            if len(self.debug_log) > 1000:
                self.debug_log = self.debug_log[-1000:]
        
        # Aggiungi anche al log del sistema se possibile
        try:
            import logging
            logging.error(error_message)
        except:
            pass
        
        # Salva periodicamente il log su file (max una volta ogni 5 minuti)
        try:
            if not hasattr(self, 'last_error_log_save') or \
            (datetime.now() - self.last_error_log_save).total_seconds() > 300:
                self.last_error_log_save = datetime.now()
                
                # Salva nel file di log
                log_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logs")
                os.makedirs(log_dir, exist_ok=True)
                
                error_log_file = os.path.join(log_dir, f"error_log_{datetime.now().strftime('%Y%m%d')}.txt")
                with open(error_log_file, "a", encoding="utf-8") as f:
                    f.write(error_message + "\n")
        except:
            # Non interrompere il flusso dell'applicazione se il salvataggio fallisce
            pass

    @error_handler
    def log_file_processing(self, file_path, content_length=0, found=False, error=None):
        """Log dettagliato dell'elaborazione dei file per debug"""
        ext = os.path.splitext(file_path)[1].lower()
        message = f"Elaborazione {ext}: {os.path.basename(file_path)}"
        
        if error:
            message += f" | ERRORE: {str(error)}"
        elif content_length > 0:
            message += f" | Contenuto estratto: {content_length} caratteri"
            if found:
                message += " | TROVATO"
        else:
            message += " | Nessun contenuto estratto"
            
        self.log_debug(message)

    @error_handler
    def check_and_notify_missing_libraries(self):
        """Verifica e notifica l'utente di eventuali librerie mancanti"""
        missing = []
        
        if not file_format_support["docx"]:
            missing.append("python-docx (per file Word)")
        if not file_format_support["pdf"]: 
            missing.append("PyPDF2 (per file PDF)")
        if not file_format_support["pptx"]:
            missing.append("python-pptx (per file PowerPoint)")
        if not file_format_support["excel"]:
            missing.append("openpyxl (per file Excel)")
        if not file_format_support["rtf"]:
            missing.append("striprtf (per file RTF)")
        if missing:
            message = "Alcune funzionalità di ricerca nei contenuti sono disabilitate.\n\n"
            message += "Per abilitare il supporto completo ai vari formati di file, installa le seguenti librerie:\n\n"
            
            for lib in missing:
                message += f"- {lib}\n"
            
            message += "\nPuoi installarle con il comando:\n"
            message += "pip install " + " ".join([lib.split(" ")[0] for lib in missing])
            
            # Mostra la notifica dopo un breve ritardo per permettere all'UI di caricarsi
            self.root.after(1000, lambda: messagebox.showinfo("Librerie opzionali mancanti", message))

    @error_handler
    def update_datetime(self):
        current_time = datetime.now().strftime('%d/%m/%Y %H:%M:%S')
        self.datetime_var.set(f"Data: {current_time} | Utente: {self.user_var.get()}")
        self.root.after(1000, self.update_datetime)

    @error_handler
    def browse_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.search_path.set(directory)
            self.update_disk_info(directory)

    @error_handler
    def optimize_system_search(self, path):
        """Ottimizza la ricerca per percorsi di sistema come C:/ impostando parametri appropriati"""
        # Riconosci più tipi di percorsi di sistema
        is_system_path = path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"]
        
        # Verifica più accurata per percorsi di rete (Windows e Linux/Mac)
        is_network_path = path.startswith('\\\\') or path.startswith('//') or (':/' in path and not path[0].isalpha())
        
        # Verifica percorsi Unix
        is_unix_root = path in ['/', '/home', '/usr', '/var', '/etc'] or path.startswith('/mnt/') or path.startswith('/media/')
        
        if is_system_path or is_network_path or is_unix_root:
            # Salva i parametri attuali
            original_params = {
                "max_files": self.max_files_to_check.get(),
                "worker_threads": self.worker_threads.get(),
                "timeout": self.timeout_seconds.get() if self.timeout_enabled.get() else None,
                "max_file_size": self.max_file_size_mb.get(),
                "block_size": self.max_files_per_block.get()
            }
            
            # Applica parametri ottimizzati
            self.max_files_to_check.set(1000000)  # Un milione di file
            
            # Ottimizza il numero di thread in base al tipo di percorso
            if is_network_path:
                # Per percorsi di rete, usa meno thread per evitare congestione
                self.worker_threads.set(max(2, min(6, os.cpu_count() or 4)))
                # Aumenta la dimensione dei blocchi per ridurre overhead di rete
                self.max_files_per_block.set(2000)
                # Riduci la dimensione massima del file per percorsi di rete
                self.max_file_size_mb.set(min(50, self.max_file_size_mb.get()))
                self.log_debug("Applicata ottimizzazione per percorso di rete")
            else:
                # Per percorsi locali di sistema, usa più thread
                self.worker_threads.set(min(12, os.cpu_count() or 4))
                self.log_debug("Applicata ottimizzazione per percorso di sistema locale")
            
            if self.timeout_enabled.get():
                self.timeout_seconds.set(max(3600, self.timeout_seconds.get()))  # Minimo 1 ora
            
            # Notifica l'utente con informazioni più dettagliate
            messagebox.showinfo(
                "Ricerca su sistema o server",
                "Stai avviando una ricerca su un percorso di sistema o server.\n\n"
                f"Tipo percorso rilevato: {'Server/Rete' if is_network_path else 'Sistema'}\n"
                f"Thread di ricerca: {self.worker_threads.get()}\n"
                f"Limite file: {self.max_files_to_check.get():,}\n\n"
                "La ricerca potrebbe richiedere tempo e risorse significative."
            )
            
            return original_params
        return None
    
    @error_handler
    def is_network_path(self, path):
        """Verifica se il percorso è un percorso di rete"""
        # Percorsi UNC Windows (\\server\share)
        if path.startswith('\\\\') or path.startswith('//'):
            return True
            
        # Percorsi di rete mappati su Windows (Z:\ dove Z è una lettera mappata)
        if os.name == 'nt' and len(path) >= 2 and path[1] == ':':
            import subprocess
            try:
                # Usa net use per verificare se è un drive mappato
                drive_letter = path[0].upper() + ':'
                result = subprocess.run(['net', 'use', drive_letter], 
                                    capture_output=True, text=True, timeout=2)
                return "Remote name" in result.stdout or "Nome remoto" in result.stdout
            except:
                pass
                
        # Percorsi di rete su Unix/Linux (/mnt/*, /media/*)
        if os.name != 'nt':
            network_prefixes = ['/mnt/smb', '/mnt/cifs', '/media/smb', '/net/', '/nfs/']
            for prefix in network_prefixes:
                if path.startswith(prefix):
                    return True
                    
        return False

    @error_handler
    def show_optimization_tips(self, path):
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]:
            # Determina la lettera del disco corrente
            if os.name == 'nt':  # Windows
                if path.lower() in ["c:/", "c:\\"]:
                    drive_letter = "C:"
                elif path.lower() in ["d:/", "d:\\"]:
                    drive_letter = "D:"
                elif path.lower() in ["e:/", "e:\\"]:
                    drive_letter = "E:"
                else:
                    drive_letter = path[:2] if len(path) >= 2 else "Disco"
            else:
                drive_letter = "Disco"  # Per sistemi non Windows
            
            optimization_done = False
            
            # Verifica se l'utente ha già implementato le ottimizzazioni
            if hasattr(self, 'excluded_paths'):
                windows_excluded = any("windows" in p.lower() for p in self.excluded_paths)
                program_files_excluded = any("program files (x86)" in p.lower() for p in self.excluded_paths)
                optimization_done = windows_excluded and program_files_excluded
            
            # Verifica la profondità di ricerca
            depth_optimized = hasattr(self, 'max_depth') and self.max_depth >= 1 and self.max_depth <= 10
            
            # Se le ottimizzazioni non sono già state applicate, mostra il messaggio
            if not (optimization_done and depth_optimized):
                response = messagebox.askyesno(
                    "Ottimizza la ricerca",
                    f"Stai per avviare una ricerca sull'intero disco {drive_letter}\\\n\n"
                    "Per migliorare notevolmente le prestazioni, è consigliato:\n\n"
                    "1. Escludere le cartelle di sistema (Windows, Program Files)\n"
                    "2. Escludere le cartelle di altri utenti\n"
                    "3. Limitare la profondità di ricerca a 5-10 livelli\n\n"
                    "Vuoi applicare automaticamente queste ottimizzazioni?",
                    icon="question"
                )
                
                if response:
                    # Applica le ottimizzazioni
                    if not hasattr(self, 'excluded_paths'):
                        self.excluded_paths = []
                    
                    # Aggiungi esclusioni di sistema
                    system_paths = ["C:/Windows", "C:/Program Files", "C:/Program Files (x86)"]
                    for sys_path in system_paths:
                        if sys_path not in self.excluded_paths:
                            self.excluded_paths.append(sys_path)
                    
                    # Escludi altri utenti
                    current_user = getpass.getuser()
                    users_dir = "C:/Users"
                    if os.path.exists(users_dir):
                        try:
                            for user in os.listdir(users_dir):
                                user_path = os.path.join(users_dir, user)
                                if user.lower() != current_user.lower() and user.lower() not in ["public", "default", "all users"]:
                                    if user_path not in self.excluded_paths:
                                        self.excluded_paths.append(user_path)
                        except Exception as e:
                            self.log_debug(f"Errore nell'esclusione degli altri utenti: {str(e)}")
                    
                    # Imposta la profondità di ricerca a 7 (valore ragionevole)
                    # Usa una verifica prima di accedere a depth_spinbox
                    if hasattr(self, 'depth_spinbox'):
                        self.depth_spinbox.set("7")
                    self.max_depth = 7
                    
                    # Notifica l'utente
                    messagebox.showinfo(
                        "Ottimizzazioni applicate",
                        f"Le ottimizzazioni sono state applicate:\n\n"
                        f"• Escluse {len(system_paths)} cartelle di sistema\n"
                        f"• Escluse le cartelle di altri utenti\n"
                        f"• Profondità di ricerca impostata a 7 livelli\n\n"
                        f"La ricerca sarà ora più veloce e stabile."
                    )
                    
                    return True
                else:
                    # NUOVA PARTE: Se l'utente risponde "No", imposta esplicitamente la profondità a 0
                    # Usa una verifica prima di accedere a depth_spinbox
                    if hasattr(self, 'depth_spinbox'):
                        self.depth_spinbox.set("0")
                    self.max_depth = 0
                    
                    # Notifica l'utente che sta per iniziare una ricerca illimitata
                    self.log_debug("Utente ha rifiutato le ottimizzazioni, impostata profondità illimitata (0)")
                    
                    # Opzionalmente, puoi mostrare un messaggio di avviso
                    messagebox.showinfo(
                        "Ricerca illimitata",
                        "Stai per avviare una ricerca illimitata sull'intero disco.\n\n"
                        "La profondità di ricerca è stata impostata a illimitata (0).\n"
                        "Questa operazione potrebbe richiedere molto tempo e risorse di sistema."
                    )
                
                return False
            
            return False  # Per percorsi non di sistema, non fare nulla
        
    @error_handler
    def debug_exclusions(self):
        """Visualizza lo stato corrente delle esclusioni per debug"""
        if hasattr(self, 'excluded_paths'):
            paths = '\n'.join(self.excluded_paths) if self.excluded_paths else "Nessun percorso escluso"
            messagebox.showinfo("Debug esclusioni", 
                            f"Stato esclusioni:\n{paths}\n\n"
                            f"Totale percorsi: {len(self.excluded_paths)}")
        else:
            messagebox.showinfo("Debug esclusioni", "Lista esclusioni non inizializzata")

    @error_handler
    def restart_as_admin(self):
        """Riavvia l'applicazione con privilegi di amministratore"""
        try:
            import ctypes, sys, os
            if sys.platform == 'win32':
                script = sys.argv[0]
                params = ' '.join(sys.argv[1:])
                self.log_debug("Tentativo di riavvio come amministratore")
                
                # Se il percorso di ricerca è impostato, passa quel percorso come parametro
                if self.search_path.get():
                    if not params:
                        params = f'"{self.search_path.get()}"'
                    
                # Avvia l'applicazione con privilegi elevati
                ctypes.windll.shell32.ShellExecuteW(None, "runas", sys.executable, f'"{script}" {params}', None, 1)
                
                # Chiude l'istanza corrente
                self.root.quit()
            else:
                messagebox.showinfo("Informazione", "Questa funzione è disponibile solo su Windows")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile riavviare come amministratore: {str(e)}")
            self.log_debug(f"Errore nel riavvio come admin: {str(e)}")

    @error_handler
    def optimize_disk_search_order(self, path, directories):
        """Ottimizza l'ordine di esplorazione delle directory quando si cerca in un disco di sistema.
        Dà priorità alle cartelle più importanti come Users, Documents, ecc.
        """
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"]:
            prioritized = []
            normal = []
            
            # Cartelle ad alta priorità (esplora per prime)
            high_priority_folders = ["users", "documents", "documents and settings", "desktop", "downloads"]
            
            # Cartelle a bassa priorità (esplora per ultime)
            low_priority_folders = ["windows", "program files", "program files (x86)", "programdata", "$recycle.bin", "system volume information"]
            
            for dir_path in directories:
                dir_name = os.path.basename(dir_path).lower()
                
                if any(priority in dir_name for priority in high_priority_folders):
                    prioritized.insert(0, dir_path)  # Aggiungi all'inizio per priorità più alta
                elif any(low in dir_name for low in low_priority_folders):
                    normal.append(dir_path)  # Aggiungi alla fine
                else:
                    prioritized.append(dir_path)  # Priorità media
                    
            return prioritized + normal
        
        # Se non è un percorso di sistema, restituisci l'elenco originale
        return directories

    @error_handler
    def start_search(self):
        # Assicurati che qualsiasi ricerca precedente sia completamente terminata
        self.reset_search_state()
        
        # Reset the total files size label at the start of a new search
        if hasattr(self, 'total_files_size_label'):
            self.total_files_size_label.config(text="Dimensione totale: 0 B (0 file)")
        
        # Pulisci risultati precedenti
        for item in self.results_list.get_children():
            self.results_list.delete(item)
        
        # Ottieni i valori direttamente dai widget
        search_path = self.search_path.get().strip()
        keywords = self.keyword_entry.get().strip()
        
        # Stampa di debug (puoi rimuoverla in produzione)
        self.log_debug(f"Avvio ricerca - Percorso: '{search_path}', Parole chiave: '{keywords}'")
        
        # 1. Verifica il percorso
        if not search_path:
            messagebox.showerror("Errore", "Inserisci una directory di ricerca valida")
            return
        
        if not os.path.exists(search_path):
            messagebox.showerror("Errore", f"Il percorso specificato non esiste:\n{search_path}")
            return
        
        # 2. Verifica le parole chiave
        placeholder = "Scrivi la parola da ricercare..."
        keywords_raw = self.keyword_entry.get()  # Ottieni il valore grezzo senza strip()

        # Verifica se è vuoto, solo spazi, o è il placeholder
        if (not keywords_raw or 
            keywords_raw.strip() == "" or 
            keywords_raw == placeholder or 
            self.keyword_entry.cget("foreground") == "gray"):
            
            messagebox.showerror("Errore", "Inserisci le parole chiave da cercare")
            return

        # 3. Verifica colore (controllo extra per il placeholder)
        try:
            if self.keyword_entry.cget("foreground") == "gray":
                messagebox.showerror("Errore", "Inserisci le parole chiave da cercare")
                return
        except:
            pass
            
        # Mostra avviso se la ricerca nei contenuti è attivata
        if not self.show_content_search_warning():
            return  # Interrompi se l'utente annulla
        
        # Aggiorna le informazioni del disco qui, quando l'utente clicca su cerca
        calculation_mode = self.dir_size_calculation.get()
        
        # CORREZIONE: Determina se calcolare la dimensione in modo esplicito
        should_calculate = False
        if calculation_mode == "stimato" or calculation_mode == "accurato":
            should_calculate = True
            self.log_debug(f"Calcolo dimensione directory abilitato: {calculation_mode}")
        else:
            self.log_debug("Calcolo dimensione directory disabilitato")
        
        # Passa il flag corretto a update_disk_info
        self.update_disk_info(path=search_path, calculate_dir_size=should_calculate)

        self.show_optimization_tips(self.search_path.get())
        
        # Reset per la nuova ricerca
        self.stop_search = False
        
        # IMPORTANTE: Imposta is_searching PRIMA di disabilitare i controlli
        self.is_searching = True
        
        # CRITICO: Crea un nuovo executor per questa ricerca
        max_workers = max(1, min(32, self.worker_threads.get()))
        self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
        
        original_params = self.optimize_system_search(self.search_path.get())
        self.start_search_watchdog()
        
        # Aggiorna l'orario di avvio e resetta l'orario di fine
        current_time = datetime.now().strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")

        # Salva i parametri di ricerca
        self.last_search_params = {
            "path": self.search_path.get(),
            "keywords": self.keywords.get(),
            "search_files": self.search_files.get(),
            "search_folders": self.search_folders.get(),
            "search_content": self.search_content.get()
        }
        
        # Aggiungi alla cronologia di ricerca se non già presente
        if self.keywords.get() and self.keywords.get() not in [h["keywords"] for h in self.search_history]:
            self.search_history.append(self.last_search_params.copy())
            if len(self.search_history) > 10:  # Mantieni solo le ultime 10 ricerche
                self.search_history.pop(0)
        
        # Disabilita i controlli
        self.disable_all_controls()
        
        # IMPORTANTE: Il pulsante di interruzione deve essere abilitato DOPO aver disabilitato tutti i controlli
        if hasattr(self, 'stop_button'):
            self.stop_button["state"] = "normal"
            # Forza aggiornamento immediato del bottone
            self.stop_button.update()
        
        self.progress_bar["value"] = 0
        self.status_label["text"] = "Ricerca in corso..."
        
        # Aggiorna il valore della profondità massima - FIX
        try:
            if hasattr(self, 'depth_spinbox') and self.depth_spinbox.winfo_exists():
                try:
                    self.max_depth = int(self.depth_spinbox.get())
                except (ValueError, tk.TclError) as e:
                    self.log_debug(f"Errore nell'accesso al depth_spinbox: {str(e)}")
                    self.max_depth = getattr(self, 'max_depth', 0)  # Usa il valore esistente o default
            else:
                # Se il widget non esiste, usa il valore memorizzato come attributo 
                # o il valore predefinito se non esiste
                self.max_depth = getattr(self, 'max_depth', 0)
                self.log_debug(f"Widget depth_spinbox non disponibile, usando valore: {self.max_depth}")
        except Exception as e:
            self.log_debug(f"Errore generale nell'aggiornamento della profondità: {str(e)}")
            self.max_depth = 0  # Valore predefinito sicuro
        
        # Ottieni le parole chiave di ricerca
        search_terms = [term.strip() for term in self.keywords.get().split(',') if term.strip()]
        
        # DEBUG: Verifica che la ricerca sia correttamente impostata
        self.log_debug(f"STATO RICERCA: is_searching={self.is_searching}, stop_search={self.stop_search}")
        self.log_debug(f"Stato pulsante interruzione: {self.stop_button['state']}")
            
        # Avvia la ricerca in un thread separato
        self.search_results = []  # Resetta i risultati
        search_thread = threading.Thread(target=self._search_thread, 
                                        args=(self.search_path.get(), search_terms, self.search_content.get()))
        search_thread.daemon = True
        search_thread.start()
        
        # Memorizza l'orario di inizio come oggetto datetime
        self.search_start_time = datetime.now()
        current_time = self.search_start_time.strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")
        self.total_time_label.config(text="--:--")  # Reset del tempo totale

        # Avvia l'aggiornamento della progress bar
        self.update_progress()
        
        # Salva original_params come attributo per ripristinarli dopo la ricerca
        if original_params:
            self.original_system_search_params = original_params

        # IMPORTANTE: Doppio controllo del pulsante di interruzione dopo aver iniziato la ricerca
        if hasattr(self, 'stop_button'):
            # Usa un approccio a ritardo doppio per garantire che sia abilitato
            self.root.after(100, lambda: self.stop_button.configure(state="normal"))
            self.root.after(500, lambda: self.stop_button.configure(state="normal"))
        
        self.log_debug("Ricerca avviata correttamente con pulsante interruzione abilitato")

    @error_handler
    def start_search_watchdog(self):
        """Avvia un timer di controllo per rilevare se la ricerca si è bloccata"""
        self.last_progress_time = time.time()
        self.watchdog_active = True
        self.check_search_progress()

    @error_handler
    def check_search_progress(self):
        """Controlla se la ricerca ha fatto progressi recentemente"""
        if not self.is_searching or not hasattr(self, 'watchdog_active') or not self.watchdog_active:
            return
            
        current_time = time.time()
        elapsed_since_progress = current_time - self.last_progress_time
        
        # Se nessun progresso per 3 minuti, considera la ricerca bloccata
        if elapsed_since_progress > 180:  # 3 minuti
            self.log_debug("La ricerca sembra bloccata - tentativo di recupero")
            
            # Prova a recuperare forzando la chiusura dell'executor e riavviandolo
            if hasattr(self, 'search_executor') and self.search_executor:
                try:
                    self.search_executor.shutdown(wait=False)
                    self.search_executor = concurrent.futures.ThreadPoolExecutor(
                        max_workers=max(1, self.worker_threads.get())
                    )
                    self.last_progress_time = time.time()  # Reset timer
                    self.status_label["text"] = "Recupero dalla ricerca bloccata..."
                except Exception as e:
                    self.log_debug(f"Errore durante il recupero della ricerca: {str(e)}")
            
        # Controlla di nuovo tra 30 secondi
        self.root.after(30000, self.check_search_progress)

    @error_handler # Aggiungi questo nuovo metodo per calcolare il tempo totale
    def update_total_time(self):
        """Calcola e visualizza il tempo totale trascorso tra inizio e fine ricerca"""
        if self.search_start_time:
            end_time = datetime.now()
            time_diff = end_time - self.search_start_time
            
            # Converti in minuti e secondi
            total_seconds = int(time_diff.total_seconds())
            minutes = total_seconds // 60
            seconds = total_seconds % 60
            
            # Formatta la stringa del tempo totale
            if minutes > 0:
                total_time_str = f"{minutes}min {seconds}sec"
            else:
                total_time_str = f"{seconds}sec"
            
            self.total_time_label.config(text=total_time_str) 
    # Esempio di utilizzo per cambiare il tema
    def change_theme(self, theme):
        self.update_theme_colors(theme)

    @error_handler
    def calculate_block_priority(self, directory):
        """Calcola la priorità del blocco (numerica, più bassa = più alta priorità)"""
        # Se l'opzione è disabilitata, usa priorità standard per tutti
        if not self.prioritize_user_folders.get():
            return 1
            
        # Altrimenti, applica la prioritizzazione configurata
        # Dai priorità alle cartelle degli utenti
        if "users" in directory.lower() or "documenti" in directory.lower() or "desktop" in directory.lower():
            return 0  # Alta priorità
        # Dai priorità medio-alta alle cartelle di dati
        elif "data" in directory.lower() or "database" in directory.lower() or "downloads" in directory.lower():
            return 1  # Priorità media-alta
        # Dai priorità bassa a cartelle di sistema
        elif "windows" in directory.lower() or "program files" in directory.lower():
            return 3  # Priorità bassa
        # Priorità standard per altre cartelle
        return 2
    
    @error_handler
    def initialize_block_queue(self, root_path, block_queue, visited_dirs, files_checked, keywords, search_content, futures):
        """Inizializza la coda di blocchi con il percorso principale"""
        try:
            items = os.listdir(root_path)
            
            # Ottieni la profondità massima dalle impostazioni
            max_depth = self.depth_var.get() if hasattr(self, 'depth_var') else self.max_depth
            
            # Prima aggiungi le directory come blocchi separati
            for item in items:
                item_path = os.path.join(root_path, item)
                try:
                    if os.path.isdir(item_path):
                        # Salta directory nascoste se richiesto
                        if self.ignore_hidden.get() and (item.startswith('.') or 
                            (os.name == 'nt' and os.path.exists(item_path) and 
                            os.stat(item_path).st_file_attributes & 2)):
                            continue
                        
                        # Salta directory escluse
                        if hasattr(self, 'excluded_paths') and any(
                            item_path.lower().startswith(excluded.lower()) 
                            for excluded in self.excluded_paths):
                            continue
                        
                        # CORREZIONE: Se max_depth è 0 o non abbiamo raggiunto il limite, aggiungi la sottocartella
                        # Il valore 0 rappresenta "illimitato"
                        if max_depth == 0 or 1 <= max_depth:
                            # Calcola priorità e aggiungi alla coda con la profondità corretta
                            priority = self.calculate_block_priority(item_path)
                            block_queue.put((priority, item_path, 1))  # Profondità 1 per le sottocartelle dirette
                            visited_dirs.add(os.path.realpath(item_path))
                except Exception as e:
                    self.log_debug(f"Errore nell'aggiunta del blocco {item_path}: {str(e)}")
            
            # Poi processa i file nella directory principale
            for item in items:
                if self.stop_search:
                    return
                    
                item_path = os.path.join(root_path, item)
                try:
                    if not os.path.isdir(item_path):
                        if self.search_files.get():
                            # Verifica file nascosti
                            if self.ignore_hidden.get() and (item.startswith('.') or 
                                (os.name == 'nt' and os.path.exists(item_path) and 
                                os.stat(item_path).st_file_attributes & 2)):
                                continue
                            
                            # Processa direttamente i file nella directory principale
                            nonlocal_files_checked = files_checked[0]
                            nonlocal_files_checked += 1
                            files_checked[0] = nonlocal_files_checked
                            if nonlocal_files_checked > self.max_files_to_check.get():
                                self.stop_search = True
                                return
                            
                            future = self.search_executor.submit(self.process_file, item_path, keywords, search_content)
                            futures.append(future)
                except Exception as e:
                    self.log_debug(f"Errore nell'elaborazione del file {item_path}: {str(e)}")
                    
        except PermissionError:
            self.log_debug(f"Permesso negato per la directory {root_path}")
        except Exception as e:
            self.log_debug(f"Errore nell'inizializzazione dei blocchi da {root_path}: {str(e)}")

    @error_handler
    def process_file_batch(self, file_batch, files_checked, last_update_time, 
                      calculation_enabled, keywords, search_content, futures):
        """Processa un batch di file in modo ottimizzato"""
        for file_path in file_batch:
            try:
                # Verifica limite file
                files_checked[0] += 1
                
                # Verifica ottimizzata del limite di file
                if files_checked[0] > self.max_files_to_check.get():
                    self.stop_search = True
                    self.progress_queue.put(("status", 
                        f"Limite di {self.max_files_to_check.get():,} file controllati raggiunto. "
                        f"Aumenta il limite nelle opzioni per cercare più file."))
                    return
                    
                # Calcola dimensione per statistiche se abilitato
                if calculation_enabled:
                    try:
                        if os.path.isfile(file_path) and not os.path.islink(file_path):
                            file_size = os.path.getsize(file_path)
                            self.current_search_size += file_size
                            
                            # Aggiorna la dimensione mostrata periodicamente
                            current_time = time.time()
                            if files_checked[0] % 1000 == 0 or (current_time - last_update_time[0]) > 5:
                                self.progress_queue.put(("update_dir_size", self.current_search_size))
                                last_update_time[0] = current_time
                    except:
                        pass  # Ignora errori durante il calcolo della dimensione
                
                # Verifica se il processo può ancora eseguire submit (non è stato chiuso)
                try:
                    if self.search_executor and not self.search_executor._shutdown:
                        # Usa la versione con timeout per evitare blocchi
                        future = self.search_executor.submit(
                            self.process_file_with_timeout, 
                            file_path, keywords, search_content
                        )
                        futures.append(future)
                    else:
                        # Fallback diretto se l'executor è chiuso
                        result = self.process_file(file_path, keywords, search_content)
                        if result:
                            self.search_results.append(result)
                except Exception as e:
                    self.log_debug(f"Errore nell'elaborazione parallela del file {file_path}: {str(e)}")
                    # Fallback se l'executor fallisce
                    try:
                        result = self.process_file(file_path, keywords, search_content)
                        if result:
                            self.search_results.append(result)
                    except:
                        pass
                    
            except Exception as e:
                self.log_debug(f"Errore nell'aggiunta del file {file_path} alla coda: {str(e)}")

    @error_handler
    def process_blocks(self, block_queue, visited_dirs, start_time, timeout, is_system_search, 
                files_checked, dirs_checked, last_update_time, path, keywords, search_content, futures):
        """Elabora i blocchi dalla coda in base alla priorità"""
        # Determina il numero massimo di file per blocco in base alle impostazioni
        max_files_in_block = self.max_files_per_block.get()
        
        # Ottieni il valore della profondità massima
        max_depth = self.depth_var.get() if hasattr(self, 'depth_var') else self.max_depth
        
        # CORREZIONE: Modifica la logica di tracking - true se c'è un limite, false se è illimitato (0)
        depth_tracking_needed = max_depth > 0
        
        # CORREZIONE: Ottimizza il calcolo della dimensione
        calculation_enabled = self.dir_size_calculation.get() != "disabilitato"
        
        # Adatta automaticamente la dimensione del blocco
        if self.block_size_auto_adjust.get():
            # Verifica se è un percorso di rete
            is_network_path = path.startswith('\\\\') or path.startswith('//')
            
            # Ottimizzazione per percorsi di rete
            if is_network_path:
                # Blocchi più grandi per la rete per ridurre l'overhead
                max_files_in_block = max(2000, max_files_in_block)
                self.log_debug(f"Dimensione blocco adattata per rete: {max_files_in_block}")
            # Ottimizzazione per ricerche di sistema
            elif is_system_search:
                max_files_in_block = max(1500, max_files_in_block)
                self.log_debug(f"Dimensione blocco adattata per sistema: {max_files_in_block}")
            # Ottimizzazione per grandi quantità di file
            elif files_checked[0] > 50000:
                max_files_in_block = max(3000, max_files_in_block)
                self.log_debug(f"Dimensione blocco adattata per ricerca grande: {max_files_in_block}")
        
        # Ottimizza il numero di blocchi paralleli in base al carico
        max_parallel = self.max_parallel_blocks.get()
        
        # CORREZIONE: Adatta il parallelismo in base al percorso
        if path.startswith('\\\\') or path.startswith('//'):
            # Limita il parallelismo su percorsi di rete
            max_parallel = min(max_parallel, 3)
            self.log_debug(f"Parallelismo limitato per percorso di rete: {max_parallel}")
        
        # Utilizziamo un set per tracciare i blocchi già processati
        processed_blocks = set()
        
        # Aggiungi contatori per controllo delle interruzioni e gestione memoria
        interrupt_check_counter = 0
        memory_check_counter = 0
        last_watchdog_update = time.time()
        
        # NUOVO: Tracciamento delle prestazioni
        last_performance_check = time.time()
        files_at_last_check = files_checked[0]
        
        while not block_queue.empty() and not self.stop_search:
            # Verifica timeout
            current_time = time.time()
            if timeout and current_time - start_time > timeout:
                self.progress_queue.put(("timeout", "Timeout raggiunto"))
                return
            
            # CORREZIONE: Aggiorna il watchdog più frequentemente
            if hasattr(self, 'last_progress_time') and current_time - last_watchdog_update > 10:
                self.last_progress_time = current_time
                last_watchdog_update = current_time
                
                # NUOVO: Tracciamento delle prestazioni
                if current_time - last_performance_check >= 30:  # Ogni 30 secondi
                    # Calcola velocità di elaborazione
                    elapsed = current_time - last_performance_check
                    files_processed = files_checked[0] - files_at_last_check
                    
                    if elapsed > 0:
                        speed = files_processed / elapsed
                        self.log_debug(f"Velocità di elaborazione: {speed:.1f} file/sec")
                        
                        # Adatta dinamicamente la dimensione del blocco
                        if self.block_size_auto_adjust.get():
                            if speed < 10:  # Molto lento
                                new_size = max(max_files_in_block // 2, 100)
                                if new_size != max_files_in_block:
                                    max_files_in_block = new_size
                                    self.log_debug(f"Riduzione dimensione blocco a {max_files_in_block} per bassa velocità")
                            elif speed > 1000:  # Molto veloce
                                new_size = min(max_files_in_block * 2, 5000)
                                if new_size != max_files_in_block:
                                    max_files_in_block = new_size
                                    self.log_debug(f"Aumento dimensione blocco a {max_files_in_block} per alta velocità")
                    
                    # Aggiorna i contatori per il prossimo controllo
                    last_performance_check = current_time
                    files_at_last_check = files_checked[0]
            
            try:
                # Prendi il blocco con priorità più alta
                # CORREZIONE: Estrai correttamente la profondità se presente nella coda
                try:
                    # Controllo più robusto per la consistenza del formato della coda
                    queue_item = block_queue.queue[0]
                    if len(queue_item) >= 3:  # Nuovo formato con profondità
                        priority, current_block, current_depth = block_queue.get(block=False)
                    else:  # Vecchio formato senza profondità
                        priority, current_block = block_queue.get(block=False)
                        # Se non stiamo tracciando la profondità (max_depth = 0), imposta a 0
                        # altrimenti calcola la profondità basata sul percorso
                        current_depth = 0 if not depth_tracking_needed else current_block.count(os.path.sep) - path.count(os.path.sep)
                except IndexError:
                    # In caso di coda vuota, esci dal ciclo
                    break
                
                # CORREZIONE: Verifica la profondità massima, ma salta solo se non è illimitata (max_depth > 0)
                # Se max_depth = 0, non saltiamo mai perché è una ricerca illimitata
                if depth_tracking_needed and current_depth >= max_depth:
                    # Salta questa directory se abbiamo superato la profondità massima
                    continue
                
                # Salta blocchi già processati
                if current_block in processed_blocks:
                    continue
                    
                processed_blocks.add(current_block)
                
                # CORREZIONE: Aggiorna lo stato più frequentemente
                current_time = time.time()
                if current_time - last_update_time[0] >= 0.3:  # Ridotto da 0.5 a 0.3 secondi
                    elapsed_time = current_time - start_time
                    self.progress_queue.put(("status", 
                        f"Analisi blocco: {current_block} (Cartelle: {dirs_checked[0]}, File: {files_checked[0]}, Tempo: {int(elapsed_time)}s)"))
                    self.progress_queue.put(("progress", 
                        min(90, int((files_checked[0] / max(1, self.max_files_to_check.get())) * 100))))
                    last_update_time[0] = current_time
                
                # CORREZIONE: Implementa la verifica dei percorsi problematici prima di elaborare
                # Verifica se il percorso attuale è in una directory problematica o esclusa
                skip_block = False
                
                # Verifica più efficiente dei percorsi esclusi
                if hasattr(self, 'excluded_paths') and self.excluded_paths:
                    if any(current_block.lower().startswith(excluded.lower()) for excluded in self.excluded_paths):
                        self.log_debug(f"Salto blocco in percorso escluso: {current_block}")
                        skip_block = True
                
                # Verifica per directory problematiche note
                if not skip_block and hasattr(self, 'problematic_dirs'):
                    if any(problematic in current_block for problematic in self.problematic_dirs):
                        self.log_debug(f"Salto blocco in directory problematica: {current_block}")
                        skip_block = True
                
                if skip_block:
                    continue
                
                # Blocco attualmente in elaborazione
                dirs_checked[0] += 1
                
                # CORREZIONE: Usa try-except più granulare per gestire errori di accesso
                try:
                    items = os.listdir(current_block)
                except PermissionError:
                    if self.skip_permission_errors.get():
                        self.log_debug(f"Saltata directory con permesso negato: {current_block}")
                        continue
                    else:
                        # Gestione fallback per directory inaccessibili
                        dir_name = os.path.basename(current_block)
                        parent_dir = os.path.dirname(current_block)
                        is_user_folder = (parent_dir.lower() in ["c:/users", "c:\\users"] and 
                                        dir_name.lower() != getpass.getuser().lower())
                        if is_user_folder:
                            self.log_debug(f"Cartella di un altro utente inaccessibile: {current_block}")
                            self.progress_queue.put(("status", f"Saltata cartella utente protetta: {current_block}"))
                        else:
                            self.log_debug(f"Permesso negato per la directory {current_block}")
                            self.progress_queue.put(("status", f"Permesso negato: {current_block}"))
                        continue
                except Exception as e:
                    self.log_debug(f"Errore nell'accesso alla directory {current_block}: {str(e)}")
                    continue
                
                # Prima processa le sottocartelle (aggiungi nuovi blocchi)
                subfolders = []
                for item in items:
                    if self.stop_search:
                        return
                        
                    item_path = os.path.join(current_block, item)
                    
                    # Salta file/cartelle nascoste
                    try:
                        if self.ignore_hidden.get() and (item.startswith('.') or 
                            (os.name == 'nt' and os.path.exists(item_path) and 
                            os.stat(item_path).st_file_attributes & 2)):
                            continue
                    except Exception as e:
                        self.log_debug(f"Errore nel controllo hidden per {item_path}: {str(e)}")
                        continue
                        
                    # Gestione sottodirectory
                    try:
                        if os.path.isdir(item_path):
                            # Verifica se la cartella è già stata visitata
                            try:
                                real_path = os.path.realpath(item_path)
                                if real_path in visited_dirs:
                                    continue
                                visited_dirs.add(real_path)
                            except:
                                if item_path in visited_dirs:
                                    continue
                                visited_dirs.add(item_path)
                            
                            # CORREZIONE: Ottimizzazione verifica percorsi esclusi
                            # Verifica se il percorso deve essere escluso (più efficiente)
                            excluded = False
                            if hasattr(self, 'excluded_paths') and self.excluded_paths:
                                excluded = any(item_path.lower().startswith(excluded_path.lower()) 
                                            for excluded_path in self.excluded_paths)
                                
                            if excluded:
                                continue
                            
                            # Aggiungi alla lista delle sottocartelle
                            subfolders.append((item_path, current_depth + 1))  # CORREZIONE: Includi la profondità incrementata
                            
                            # Verifica corrispondenza nome cartella
                            if self.search_folders.get():
                                # CORREZIONE: Usa match parola intera dove appropriato
                                matched = False
                                for keyword in keywords:
                                    if self.whole_word_search.get():
                                        if self.is_whole_word_match(keyword, item):
                                            matched = True
                                            break
                                    elif keyword.lower() in item.lower():
                                        matched = True
                                        break
                                        
                                if matched:
                                    folder_info = self.create_folder_info(item_path)
                                    self.search_results.append(folder_info)
                                    
                    except Exception as e:
                        self.log_debug(f"Errore nell'analisi della directory {item_path}: {str(e)}")
                
                # CORREZIONE: Ottimizza l'ordine di elaborazione delle sottocartelle
                if self.prioritize_user_folders.get():
                    # Estrai solo i percorsi per la funzione optimize_disk_search_order
                    subfolder_paths = [item[0] for item in subfolders]
                    optimized_paths = self.optimize_disk_search_order(path, subfolder_paths)
                    
                    # Ricostruisci la lista con le profondità
                    optimized_subfolders = []
                    path_to_depth = {item[0]: item[1] for item in subfolders}
                    for optimized_path in optimized_paths:
                        optimized_subfolders.append((optimized_path, path_to_depth.get(optimized_path, current_depth + 1)))
                    subfolders = optimized_subfolders
                
                # CORREZIONE: Aggiunta sottocartelle alla coda con priorità calcolata
                # Usa sempre il formato con profondità per maggiore coerenza
                for subfolder_info in subfolders:
                    subfolder, folder_depth = subfolder_info
                    
                    priority = self.calculate_block_priority(subfolder)
                    
                    # CORREZIONE: Aggiungi sempre la profondità nella coda per coerenza
                    block_queue.put((priority, subfolder, folder_depth))
                
                # CORREZIONE: Ottimizza la gestione della memoria per blocchi grandi
                # Gestione della memoria più aggressiva
                memory_check_counter += 1
                if memory_check_counter >= 5:  # Controllo ogni 5 blocchi invece che per file
                    self.manage_memory()
                    memory_check_counter = 0
                
                # CORREZIONE: Processa i file nel blocco corrente con batch più piccoli
                file_batch = []
                for item in items:
                    if self.stop_search:
                        return
                        
                    item_path = os.path.join(current_block, item)
                    
                    # Salta i file nascosti e le directory (già processate)
                    try:
                        # Controllo file nascosto
                        if self.ignore_hidden.get() and (item.startswith('.') or 
                                    (os.name == 'nt' and os.path.exists(item_path) and 
                                        os.stat(item_path).st_file_attributes & 2)):
                            continue
                            
                        # Salta le directory (già processate)
                        if os.path.isdir(item_path):
                            continue
                    except Exception as e:
                        self.log_debug(f"Errore nel controllo del tipo per {item_path}: {str(e)}")
                        continue
                        
                    # Aggiungi file al batch se cerchiamo nei file
                    if self.search_files.get():
                        file_batch.append(item_path)
                        
                        # Quando il batch raggiunge dimensione massima, processalo
                        if len(file_batch) >= 100:  # Elabora a blocchi di 100 file
                            self.process_file_batch(file_batch, files_checked, last_update_time, 
                                                calculation_enabled, keywords, search_content, futures)
                            file_batch = []
                
                # Processa il batch finale di file
                if file_batch:
                    self.process_file_batch(file_batch, files_checked, last_update_time, 
                                        calculation_enabled, keywords, search_content, futures)
                    
            except queue.Empty:
                break

            # CORREZIONE: Controllo più frequente per interruzioni
            interrupt_check_counter += 1
            if interrupt_check_counter >= 3:  # Ridotto da 10 a 3 blocchi
                interrupt_check_counter = 0
                # Forza un aggiornamento dell'interfaccia più frequente
                try:
                    self.root.update_idletasks()
                except:
                    pass
    
    @error_handler
    # Funzione per calcolare il tempo rimanente stimato
    def calculate_remaining_time(self, files_processed, max_files, elapsed_time):
        """Calcola il tempo rimanente stimato in base al progresso attuale"""
        if files_processed == 0 or elapsed_time == 0:
            return "Calcolo..."
            
        try:
            # Calcola la velocità di elaborazione (files/secondo)
            files_per_second = files_processed / elapsed_time
            
            if files_per_second > 0:
                # Calcola il tempo rimanente stimato
                remaining_files = max_files - files_processed
                remaining_seconds = remaining_files / files_per_second
                
                # Formatta il risultato
                if remaining_seconds < 60:
                    return f"{int(remaining_seconds)} secondi"
                elif remaining_seconds < 3600:
                    minutes = int(remaining_seconds // 60)
                    seconds = int(remaining_seconds % 60)
                    return f"{minutes}m {seconds}s"
                else:
                    hours = int(remaining_seconds // 3600)
                    minutes = int((remaining_seconds % 3600) // 60)
                    return f"{hours}h {minutes}m"
            else:
                return "Calcolo..."
        except Exception as e:
            self.log_debug(f"Errore nel calcolo del tempo rimanente: {str(e)}")
            return "N/A"
        
    @error_handler
    def _search_thread(self, path, keywords, search_content):
        try:
            # Inizializza i contatori di file e directory esaminati e il tempo di inizio
            files_checked = [0]  # Uso una lista per poter modificare il valore nelle funzioni chiamate
            dirs_checked = [0]
            start_time = time.time()
            timeout = self.timeout_seconds.get() if self.timeout_enabled.get() else None
            
            # NUOVA RIGA: Verifica all'inizio se il calcolo della dimensione è abilitato
            calculation_enabled = self.dir_size_calculation.get() != "disabilitato"
            
            self.current_search_size = 0  # Dimensione totale dei file trovati
            last_size_update_time = time.time()  # Per aggiornamenti periodici

            # Determina se si tratta di una ricerca completa del sistema (C:/ o simile)
            is_system_search = path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]

            # Per ricerche di sistema, adatta automaticamente i parametri
            if is_system_search:
                # Informa l'utente che la ricerca potrebbe richiedere molto tempo
                self.progress_queue.put(("status", "Ricerca completa del sistema in corso - potrebbe richiedere molto tempo"))
                
                # Temporaneamente aumenta i limiti per la ricerca di sistema
                original_max_files = self.max_files_to_check.get()
                original_timeout = timeout
                
                # Imposta limiti più elevati per la ricerca completa
                self.max_files_to_check.set(5000000)  # 5 milioni di file
                
                # Disabilita temporaneamente il timeout o aumentalo significativamente
                if timeout and timeout < 3600:
                    timeout = 3600 * 8  # 8 ore
                
                # Avviso all'utente
                self.progress_queue.put(("status", "Ricerca completa avviata - parametri adattati per ricerca approfondita"))
            
            # Variabili per gestire l'aggiornamento ogni secondo
            last_update_time = [time.time()]
            
            # Crea un executor per processare i file in parallelo in modo sicuro
            try:
                max_workers = max(1, min(32, self.worker_threads.get()))  # Garantisce un valore valido
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
            except Exception as e:
                self.log_debug(f"Errore nella creazione del ThreadPoolExecutor: {str(e)}")
                # Fallback a un valore sicuro
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
            
            futures = []
            
            # Insieme per tenere traccia delle cartelle visitate (evita loop infiniti con symlink)
            visited_dirs = set()
            
            # Coda di priorità per i blocchi di ricerca
            block_queue = queue.PriorityQueue()

            # Aggiorna lo stato iniziale
            self.progress_queue.put(("status", f"Inizio ricerca a blocchi in: {path} (Profondità: {'illimitata' if self.max_depth == 0 else self.max_depth})"))
            
            # Assicurati che la lista di esclusioni sia inizializzata
            if not hasattr(self, 'excluded_paths'):
                self.excluded_paths = []
                
            # Avvia la ricerca a blocchi
            self.initialize_block_queue(path, block_queue, visited_dirs, files_checked, keywords, search_content, futures)
            self.process_blocks(block_queue, visited_dirs, start_time, timeout, is_system_search, 
                            files_checked, dirs_checked, last_update_time, path, keywords, search_content, futures)
            
            # Ripristina i parametri originali se erano stati modificati
            if is_system_search and self.max_depth == 0 and 'original_max_files' in locals():
                self.max_files_to_check.set(original_max_files)
            
            # Aggiorna lo stato finale di analisi
            self.progress_queue.put(("status", f"Elaborazione risultati... (analizzati {files_checked[0]} file in {dirs_checked[0]} cartelle)"))
            
            # Raccolta risultati dalle future con aggiornamento temporizzato
            completed = 0
            total_futures = len(futures)
            last_update_time[0] = time.time()
            
            # Gestione sicura delle future
            if self.search_executor and not self.search_executor._shutdown and futures:
                try:
                    for future in concurrent.futures.as_completed(futures):
                        if self.stop_search:
                            break
                            
                        try:
                            result = future.result()
                            if result:
                                self.search_results.append(result)
                            
                            completed += 1
                            
                            # Aggiorna il progresso e il tempo più frequentemente
                            current_time = time.time()
                            if completed % 20 == 0 or current_time - last_update_time[0] > 2:  # Aggiorna ogni 20 file o ogni 2 secondi
                                progress = 90 + min(10, int((completed / max(1, len(futures))) * 10))
                                self.progress_queue.put(("progress", progress))
                                
                                elapsed_time = current_time - start_time
                                self.progress_queue.put(("status", f"Elaborati {completed}/{len(futures)} file (tempo: {int(elapsed_time)}s)"))
                                
                                # Aggiorna anche il tempo totale durante la ricerca
                                if hasattr(self, 'search_start_time'):
                                    now = datetime.now()
                                    time_diff = now - self.search_start_time
                                    # Converti in minuti e secondi
                                    total_seconds = int(time_diff.total_seconds())
                                    minutes = total_seconds // 60
                                    seconds = total_seconds % 60
                                    
                                    if minutes > 0:
                                        total_time_str = f"{minutes}min {seconds}sec"
                                    else:
                                        total_time_str = f"{seconds}sec"
                                        
                                    self.progress_queue.put(("update_total_time", total_time_str))
                                
                                last_update_time[0] = current_time
                        except Exception as e:
                            self.log_debug(f"Errore nell'elaborazione di un risultato: {str(e)}")
                except Exception as e:
                    self.log_debug(f"Errore nella raccolta dei risultati: {str(e)}")
            
            # Completa la ricerca in modo sicuro
            try:
                if self.search_executor and not self.search_executor._shutdown:
                    self.search_executor.shutdown(wait=False)
                self.search_executor = None
            except Exception as e:
                self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                self.search_executor = None
            
            # Ripristina i parametri originali se erano stati modificati per la ricerca di sistema
            if hasattr(self, 'original_system_search_params'):
                try:
                    self.max_files_to_check.set(self.original_system_search_params["max_files"])
                    self.worker_threads.set(self.original_system_search_params["worker_threads"])
                    delattr(self, 'original_system_search_params')
                except:
                    pass

            # Riporta il risultato finale
            elapsed_time = time.time() - start_time
            self.progress_queue.put(("status", 
                f"Ricerca completata! Analizzati {files_checked[0]} file in {dirs_checked[0]} cartelle in {int(elapsed_time)} secondi."))
            
            # Ordina i risultati per tipo e nome
            self.search_results.sort(key=lambda x: (x[0], x[1]))
            
            self.log_debug(f"Ricerca completata. Trovati {len(self.search_results)} risultati")
            self.progress_queue.put(("complete", "Ricerca completata"))
            
        except Exception as e:
            error_msg = f"Si è verificato un errore durante la ricerca: {str(e)}\n{traceback.format_exc()}"
            self.log_debug(error_msg)
            self.progress_queue.put(("error", error_msg))

    @error_handler
    def is_whole_word_match(self, keyword, text):
        """Verifica se la keyword è presente nel testo come parola intera."""
        try:
            pattern = r'\b' + re.escape(keyword.lower()) + r'\b'
            return re.search(pattern, text.lower()) is not None
        except re.error as e:
            self.log_debug(f"Errore regex con il termine '{keyword}': {str(e)}")
            
            # Fallback manuale se la regex fallisce
            keyword_lower = keyword.lower()
            text_lower = text.lower()
            
            if keyword_lower not in text_lower:
                return False
                
            # Verifica manuale dei confini parola
            index = text_lower.find(keyword_lower)
            while index != -1:
                # Controlla il carattere prima della keyword
                has_char_before = index > 0 and text_lower[index-1].isalnum()
                
                # Controlla il carattere dopo la keyword
                end_pos = index + len(keyword_lower)
                has_char_after = end_pos < len(text_lower) and text_lower[end_pos].isalnum()
                
                # Se non ci sono caratteri alfanumerici prima e dopo, è una parola intera
                if not has_char_before and not has_char_after:
                    return True
                    
                # Cerca la prossima occorrenza
                index = text_lower.find(keyword_lower, index + 1)
            
            return False

    @error_handler  
    def search_current_user_only(self):
        """Imposta la ricerca solo nella cartella dell'utente corrente"""
        user_folder = os.path.join("C:/Users", getpass.getuser())
        if os.path.exists(user_folder):
            self.search_path.set(user_folder)
            
            # Mostra la conferma solo se l'utente non annulla il warning sulla ricerca nei contenuti
            if self.show_content_search_warning():
                messagebox.showinfo(
                    "Ricerca configurata", 
                    f"La ricerca è stata configurata per cercare solo nella tua cartella utente:\n{user_folder}\n\n"
                    "Questa modalità evita problemi di permesso con altre cartelle protette."
                )
                # Avvia direttamente la ricerca
                self.start_search()
        else:
            messagebox.showerror("Errore", "Impossibile trovare la tua cartella utente")

    @error_handler
    def create_file_info(self, file_path, from_attachment=False):
        """Crea le informazioni del file per la visualizzazione"""
        try:
            file_size = os.path.getsize(file_path)
            modified_time = datetime.fromtimestamp(os.path.getmtime(file_path))
            created_time = datetime.fromtimestamp(os.path.getctime(file_path))
            file_name = os.path.basename(file_path)
            file_extension = os.path.splitext(file_name)[1].lower()
            
            # Determina il tipo di file
            file_type = "File"  # Valore predefinito
            if os.path.isdir(file_path):
                file_type = "Directory"
            else:
                # Usa mimetypes per determinare il tipo
                import mimetypes
                mime_type, _ = mimetypes.guess_type(file_path)
                if mime_type:
                    file_type = mime_type.split('/')[0].capitalize()
                    if file_type == "Application":
                        if "pdf" in mime_type:
                            file_type = "PDF"
                        elif "word" in mime_type or file_extension == ".docx" or file_extension == ".doc":
                            file_type = "Word"
                        elif "excel" in mime_type or file_extension in [".xlsx", ".xls"]:
                            file_type = "Excel"
                        elif "powerpoint" in mime_type or file_extension in [".pptx", ".ppt"]:
                            file_type = "PowerPoint"
                        else:
                            file_type = "Documento"
                else:
                    # Fallback basato sull'estensione
                    if file_extension in ['.txt', '.md', '.rtf']:
                        file_type = "Testo"
                    elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                        file_type = "Immagine"
                    elif file_extension in ['.mp3', '.wav', '.ogg', '.flac']:
                        file_type = "Audio"
                    elif file_extension in ['.mp4', '.avi', '.mkv', '.mov']:
                        file_type = "Video"
                    elif file_extension in ['.exe', '.dll', '.bat']:
                        file_type = "Eseguibile"
                    elif file_extension == '.eml':
                        file_type = "Email"
            
            # Formatta dimensione file
            if file_size < 1024:
                size_str = f"{file_size} B"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.1f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.1f} MB"
            
            # Aggiungi il flag from_attachment alla tupla dei risultati
            return (
                file_type,
                file_name,
                size_str,
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                file_path,
                from_attachment  # Nuovo campo
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni del file {file_path}: {str(e)}")
            return (
                "File",
                os.path.basename(file_path),
                "N/A",
                "N/A",
                "N/A",
                file_path,
                False  # Default: non è un allegato
            )
    @error_handler
    def create_folder_info(self, folder_path):
        """Crea le informazioni della cartella per la visualizzazione"""
        try:
            modified_time = datetime.fromtimestamp(os.path.getmtime(folder_path))
            created_time = datetime.fromtimestamp(os.path.getctime(folder_path))
            folder_name = os.path.basename(folder_path)
            
            return (
                "Directory",
                folder_name,
                "",  # Le cartelle non hanno dimensione
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                folder_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni della cartella {folder_path}: {str(e)}")
            return (
                "Directory",
                os.path.basename(folder_path),
                "",
                "N/A",
                "N/A",
                folder_path
            )
    
    @error_handler
    def should_search_content(self, file_path):
        """Versione ottimizzata per determinare se analizzare il contenuto del file"""
        # Prima verifica le condizioni più veloci (principio fail-fast)
        if not self.search_content.get():
            return False
                    
        ext = os.path.splitext(file_path)[1].lower()
        
        # Attiva sempre la ricerca nei file Office, indipendentemente dal livello di ricerca
        if ext in ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']:
            self.log_debug(f"Analisi contenuto attivata per file Office: {os.path.basename(file_path)} ({ext})")
            return True
        
        # Seleziona il livello di ricerca attuale
        search_level = self.search_depth.get()
        
        # Ottieni le estensioni personalizzate dell'utente
        custom_extensions = self.get_extension_settings(search_level)
        
        # PRIORITÀ #1: Se l'estensione è stata aggiunta manualmente, cerca sempre il contenuto
        if ext in custom_extensions:
            self.log_debug(f"Ricerca contenuto in file con estensione personalizzata: {ext}")
            return True
        
        # PRIORITÀ #2: In modalità profonda senza estensioni personalizzate, cerca tutto
        if search_level == "profonda" and not custom_extensions:
            return True
        
        # Liste predefinite nel codice per ciascun livello
        base_extensions = ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.log', 
                        '.docx', '.doc', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls', '.rtf', '.odt', '.ods', '.odp',
                        '.csv','.eml', '.msg', '.emlx']

        advanced_extensions = base_extensions + ['.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', 
                                            '.vbs', '.js', '.config', '.ini', '.json', '.reg']
                                            
        deep_extensions = advanced_extensions + ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.mp3', '.mp4', 
                                            '.avi', '.mov', '.mkv', '.wav', '.flac', '.zip', '.rar', 
                                            '.7z', '.tar', '.gz', '.iso', '.psd', '.ai', '.svg']
        
        # Verifica il livello predefinito se non è nelle estensioni personalizzate
        if search_level == "base" and ext in base_extensions:
            return True
        elif search_level == "avanzata" and ext in advanced_extensions:
            return True
        elif search_level == "profonda" and ext in deep_extensions:
            return True
        
        return False

    @error_handler
    def should_skip_file(self, file_path):
        """Verifica se un file deve essere saltato durante l'analisi del contenuto"""
        ext = os.path.splitext(file_path)[1].lower()
        skip_type = "File di sistema" if ext in self.system_file_extensions else "File"
        skip_filename = os.path.basename(file_path)
        
        # IMPORTANTE: Verifica se il file è nelle estensioni personalizzate
        # Se è stato aggiunto manualmente, NON deve essere mai saltato
        search_level = self.search_depth.get()
        custom_extensions = self.get_extension_settings(search_level) 
        
        # Non saltare mai i file con estensioni aggiunte manualmente
        if ext in custom_extensions:
            # Aggiungi log per debug
            self.log_debug(f"File NON saltato perché in estensioni personalizzate: {file_path}")
            return False
        
        # Salta i file di Rights Management Services
        if "Rights Management Services" in file_path or "IRMProtectors" in file_path:
            self.log_debug(f"Saltato file protetto: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Rights Management Services")
            return True
                    
        # Salta file con estensioni problematiche
        problematic_extensions = [".msoprotector.doc", ".msoprotector.ppt", ".msoprotector.xls"]
        if any(ext in file_path for ext in problematic_extensions):
            self.log_debug(f"Saltato file con formato problematico: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Formato problematico")
            return True

        # Salta file di sistema solo se non sono nelle estensioni personalizzate
        if self.exclude_system_files.get() and ext in self.system_file_extensions:
            # Non saltare script files in ricerca avanzata/profonda
            script_files = ['.bat', '.cmd', '.ps1', '.vbs']
            if ext in script_files and search_level in ["avanzata", "profonda"]:
                return False
            
            self.log_debug(f"File di sistema escluso: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "File di sistema")
            return True
                    
        return False
    
    @error_handler
    def log_skipped_file(self, filepath, skiptype, filename, skipreason):
        """Registra i file saltati in un file di log"""
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_entry = f"{timestamp} - {skiptype} - {filename} - {filepath} - {skipreason}\n"
            
            with open(self.skipped_files_log_path, 'a', encoding='utf-8') as log_file:
                log_file.write(log_entry)
        except Exception as e:
            self.log_debug(f"Errore durante la scrittura del log dei file saltati: {str(e)}")

    @error_handler
    def export_skipped_files_log(self):
        """Esporta il log dei file saltati in un formato CSV"""
        try:
               
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da esportare.")
                return
                
            # Verifica che il file non sia vuoto
            if os.path.getsize(self.skipped_files_log_path) == 0:
                messagebox.showinfo("Informazione", "Il file di log esiste ma è vuoto. Non ci sono file esclusi da esportare.")
                return

            # Chiedi all'utente dove salvare il file esportato
            export_path = filedialog.asksaveasfilename(
                defaultextension=".csv",
                initialfile="file_esclusi_export.csv",
                filetypes=[("CSV files", "*.csv"), ("Text files", "*.txt"), ("All files", "*.*")],
                title="Salva il log dei file esclusi"
            )
            
            if not export_path:  # L'utente ha annullato
                return
                
            # Leggi il file di log originale
            with open(self.skipped_files_log_path, 'r', encoding='utf-8') as log_file:
                log_content = log_file.readlines()
                
            # Crea il file CSV
            with open(export_path, 'w', newline='', encoding='utf-8') as csv_file:
                csv_writer = csv.writer(csv_file, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
                
                # Intestazione
                csv_writer.writerow(["Data e Ora", "Tipo", "Nome File", "Percorso Completo", "Motivo Esclusione"])
                
                # Analizza e scrive ogni riga del log
                for log_line in log_content:
                    try:
                        # Formato tipico: 2023-01-01 12:34:56 - File di sistema - nomefile.exe - C:/percorso/nomefile.exe - File di sistema
                        parts = log_line.strip().split(" - ", 4)
                        if len(parts) >= 5:
                            csv_writer.writerow(parts)
                    except Exception as e:
                        self.log_debug(f"Errore nell'elaborazione della riga di log: {str(e)}")
                        
                # Aggiunge statistiche alla fine
                csv_writer.writerow([])
                csv_writer.writerow([f"Totale file esclusi: {len(log_content)}"])
                csv_writer.writerow([f"Esportazione eseguita il: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"])
                csv_writer.writerow([f"Utente: {self.current_user}"])
            
            # Aggiungi un link per aprire il file esportato
            open_export = messagebox.askyesno(
                "Esportazione completata", 
                f"Esportazione completata con successo!\n\nFile salvato in:\n{export_path}\n\nVuoi aprire il file?"
            )
            
            if open_export:
                try:
                    if os.name == 'nt':  # Windows
                        os.startfile(export_path)
                    else:  # macOS o Linux
                        subprocess.call(['xdg-open', export_path])
                except Exception as e:
                    self.log_debug(f"Errore nell'apertura del file esportato: {str(e)}")
                    messagebox.showinfo("Informazione", f"Il file è stato salvato in:\n{export_path}")
                    
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'esportazione: {str(e)}")
            self.log_debug(f"Errore nell'esportazione del log: {str(e)}")

    @error_handler
    @error_handler
    def view_skipped_files_log(self):
        """Visualizza il log dei file saltati in una finestra separata"""
        try:
            if not hasattr(self, 'search_start_time') or self.search_start_time is None:
                messagebox.showinfo("Informazione", "Non è stata ancora effettuata una ricerca per visualizzare i file esclusi.")
                return
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da visualizzare.")
                return
                
            # Verifica che il file non sia vuoto
            if os.path.getsize(self.skipped_files_log_path) == 0:
                messagebox.showinfo("Informazione", "Il file di log esiste ma è vuoto. Non ci sono file esclusi da visualizzare.")
                return
                
            # Crea una nuova finestra
            log_window = ttk.Toplevel(self.root)
            log_window.title("Log dei file esclusi")
            log_window.geometry("1000x800")
            log_window.transient(self.root)
            
            # Frame principale
            main_frame = ttk.Frame(log_window, padding=10)
            main_frame.pack(fill=BOTH, expand=YES)
            
            # Intestazione
            ttk.Label(main_frame, text="File esclusi durante la ricerca", 
                    font=("", 12, "bold")).pack(anchor=W, pady=(0, 10))
            
            # Area di testo per visualizzare i log
            text_frame = ttk.Frame(main_frame)
            text_frame.pack(fill=BOTH, expand=YES)
            
            scrollbar_y = ttk.Scrollbar(text_frame)
            scrollbar_y.pack(side=RIGHT, fill=Y)
            
            scrollbar_x = ttk.Scrollbar(text_frame, orient="horizontal")
            scrollbar_x.pack(side=BOTTOM, fill=X)
            
            log_text = tk.Text(text_frame, wrap="none", 
                            xscrollcommand=scrollbar_x.set,
                            yscrollcommand=scrollbar_y.set)
            log_text.pack(fill=BOTH, expand=YES)
            
            scrollbar_y.config(command=log_text.yview)
            scrollbar_x.config(command=log_text.xview)
            
            # Leggi e inserisci il contenuto del log
            with open(self.skipped_files_log_path, 'r', encoding='utf-8') as log_file:
                log_content = log_file.read()
                log_text.insert("1.0", log_content)
                
            # Rendi il testo di sola lettura
            log_text.config(state="disabled")
            
            # Pulsanti
            btn_frame = ttk.Frame(main_frame)
            btn_frame.pack(fill=X, pady=(10, 0))
            
            ttk.Button(btn_frame, text="Esporta CSV", 
                    command=self.export_skipped_files_log).pack(side=LEFT)
            
            ttk.Button(btn_frame, text="Svuota Log", 
                    command=self.clear_skipped_files_log).pack(side=LEFT, padx=5)
            
            ttk.Button(btn_frame, text="Chiudi", 
                    command=log_window.destroy).pack(side=RIGHT)
                    
            # Centra la finestra
            log_window.update_idletasks()
            width = log_window.winfo_width()
            height = log_window.winfo_height()
            x = (log_window.winfo_screenwidth() // 2) - (width // 2)
            y = (log_window.winfo_screenheight() // 2) - (height // 2)
            log_window.geometry(f"{width}x{height}+{x}+{y}")
            
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante la visualizzazione del log: {str(e)}")
            self.log_debug(f"Errore nella visualizzazione del log: {str(e)}")

    @error_handler
    def clear_skipped_files_log(self):
        """Svuota il log dei file esclusi"""
        try:
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da cancellare.")
                return
                
            # Svuota il file di log
            open(self.skipped_files_log_path, 'w', encoding='utf-8').close()
            
            # Reinizializza la lista dei file saltati
            self.skipped_files = []
            
            # Fornisci feedback all'utente
            messagebox.showinfo("Operazione completata", "Il log dei file esclusi è stato svuotato con successo.")
            
            # Se la finestra di log è attualmente aperta, aggiornala per mostrare il log vuoto
            for widget in self.root.winfo_children():
                if isinstance(widget, tk.Toplevel) and widget.winfo_exists():
                    # Cerca il widget Text nel toplevel
                    for child in widget.winfo_children():
                        if isinstance(child, ttk.Frame):  # main_frame
                            for frame in child.winfo_children():
                                if isinstance(frame, ttk.Frame):  # text_frame
                                    for text_widget in frame.winfo_children():
                                        if isinstance(text_widget, tk.Text):
                                            text_widget.config(state="normal")
                                            text_widget.delete("1.0", tk.END)
                                            text_widget.config(state="disabled")
                                            return
            
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante la cancellazione del log: {str(e)}")
            self.log_debug(f"Errore nella cancellazione del log: {str(e)}")

    @error_handler
    def get_file_content(self, file_path):
        """Estrae il contenuto testuale dai vari formati di file - versione completa"""
        try:
            # Inizializza le variabili all'inizio per evitare problemi di scope
            content = ""
            text_content = ""
            result = ""
            # Controlli preliminari
            if self.should_skip_file(file_path):
                return ""
                    
            ext = os.path.splitext(file_path)[1].lower()
            
            # Log per debug
            self.log_debug(f"Tentativo estrazione contenuto da: {os.path.basename(file_path)} ({ext})")
            
            # Controllo dimensione file
            try:
                file_size = os.path.getsize(file_path)
                if file_size > self.max_file_size_mb.get() * 1024 * 1024:
                    self.log_debug(f"File {file_path} troppo grande per l'analisi del contenuto")
                    return ""
            except Exception as e:
                self.log_debug(f"Errore nel controllo dimensione del file {file_path}: {str(e)}")
                return ""
            
            if self.stop_search:
                return ""
            # Word DOCX
            if ext == '.docx':
                try:
                    import docx
                    self.log_debug(f"Processando file DOCX: {file_path}")
                    doc = docx.Document(file_path)
                    content = []
                    for paragraph in doc.paragraphs:
                        if paragraph.text.strip():
                            content.append(paragraph.text)
                    
                    # Estrai anche il testo dalle tabelle
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                content.append(" | ".join(row_text))
                    
                    result = '\n'.join(content)
                    self.log_debug(f"Estratti {len(result)} caratteri da DOCX")
                    return result
                except ImportError:
                    self.log_debug("Libreria python-docx non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file DOCX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Word DOC (vecchio formato)
            elif ext == '.doc':
                try:
                    # Prova prima con pywin32 (solo Windows)
                    if os.name == 'nt':
                        self.log_debug("Tentativo di estrazione da DOC con win32com...")
                        try:
                            import win32com.client
                            import pythoncom
                            
                            # Inizializzazione necessaria per i thread
                            pythoncom.CoInitialize()
                            
                            try:
                                word = win32com.client.Dispatch("Word.Application")
                                word.Visible = False
                                word.DisplayAlerts = False
                                
                                # Apri il documento in modalità sola lettura
                                doc = word.Documents.Open(os.path.abspath(file_path), ReadOnly=True)
                                text = doc.Content.Text
                                doc.Close(SaveChanges=False)
                                word.Quit()
                                
                                pythoncom.CoUninitialize()
                                
                                if text:
                                    self.log_debug(f"Estratti {len(text)} caratteri da DOC")
                                    return text
                                else:
                                    self.log_debug("Nessun testo estratto dal file DOC")
                                    return ""
                            except Exception as e:
                                self.log_debug(f"Errore nell'apertura del DOC con win32com: {str(e)}")
                                # Cleanup in caso di errore
                                try:
                                    if 'doc' in locals() and doc:
                                        doc.Close(SaveChanges=False)
                                    if 'word' in locals() and word:
                                        word.Quit()
                                except:
                                    pass
                                
                                pythoncom.CoUninitialize()
                                return ""
                        except ImportError:
                            self.log_debug("win32com non disponibile per i file DOC")
                            return ""
                    else:
                        self.log_debug("Estrazione da DOC non supportata su questa piattaforma")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore generale nell'elaborazione del file DOC {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Excel XLSX
            elif ext == '.xlsx':
                try:
                    import openpyxl
                    self.log_debug(f"Processando file XLSX: {file_path}")
                    
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    texts = []
                    
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        sheet_texts = []
                        
                        # Prima verifica il range utilizzato
                        max_row = min(sheet.max_row, 500) if sheet.max_row else 0
                        max_col = min(sheet.max_column, 50) if sheet.max_column else 0
                        
                        if max_row > 0 and max_col > 0:
                            # Utilizza openpyxl 2.6+ con la sintassi più efficiente
                            try:
                                # Estrazione di righe specifiche
                                for row_idx in range(1, max_row + 1):
                                    row_values = []
                                    for col_idx in range(1, max_col + 1):
                                        cell = sheet.cell(row=row_idx, column=col_idx)
                                        if cell.value:
                                            row_values.append(str(cell.value))
                                    
                                    if row_values:
                                        sheet_texts.append(" ".join(row_values))
                            except Exception as e:
                                self.log_debug(f"Errore nell'iterazione del foglio {sheet_name}: {str(e)}")
                        
                        if sheet_texts:
                            texts.append(f"--- Foglio: {sheet_name} ---")
                            texts.append("\n".join(sheet_texts))
                    
                    result = "\n".join(texts)
                    self.log_debug(f"Estratti {len(result)} caratteri da XLSX")
                    return result
                    
                except ImportError:
                    self.log_debug("Libreria openpyxl non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file XLSX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Excel XLS (vecchio formato)
            elif ext == '.xls':
                try:
                    # Prima prova con xlrd
                    try:
                        import xlrd
                        self.log_debug(f"Processando file XLS con xlrd: {file_path}")
                        
                        book = xlrd.open_workbook(file_path, on_demand=True)
                        texts = []
                        
                        for sheet_idx in range(book.nsheets):
                            sheet = book.sheet_by_index(sheet_idx)
                            sheet_texts = []
                            
                            # Limite a 500 righe per prestazioni
                            for row_idx in range(min(sheet.nrows, 500)):
                                row_values = sheet.row_values(row_idx)
                                row_texts = [str(value) for value in row_values if value]
                                if row_texts:
                                    sheet_texts.append(" ".join(row_texts))
                            
                            if sheet_texts:
                                texts.append(f"--- Foglio: {sheet.name} ---")
                                texts.append("\n".join(sheet_texts))
                        
                        book.release_resources()
                        result = "\n".join(texts)
                        self.log_debug(f"Estratti {len(result)} caratteri da XLS")
                        return result
                        
                    except ImportError:
                        # Fallback a pywin32 se disponibile e su Windows
                        if os.name == 'nt':
                            try:
                                import win32com.client
                                import pythoncom
                                
                                # Inizializzazione necessaria per i thread
                                pythoncom.CoInitialize()
                                
                                self.log_debug("Processando file XLS con win32com")
                                
                                try:
                                    excel = win32com.client.Dispatch("Excel.Application")
                                    excel.Visible = False
                                    excel.DisplayAlerts = False
                                    
                                    workbook = excel.Workbooks.Open(os.path.abspath(file_path), ReadOnly=True)
                                    texts = []
                                    
                                    for i in range(1, workbook.Sheets.Count + 1):
                                        sheet = workbook.Sheets(i)
                                        texts.append(f"--- Foglio: {sheet.Name} ---")
                                        
                                        # Verifica se c'è un range utilizzato
                                        if sheet.UsedRange and sheet.UsedRange.Cells.Count > 0:
                                            used_range = sheet.UsedRange
                                            row_count = min(used_range.Rows.Count, 500)
                                            col_count = min(used_range.Columns.Count, 50)
                                            
                                            for row_idx in range(1, row_count + 1):
                                                row_texts = []
                                                for col_idx in range(1, col_count + 1):
                                                    cell_value = used_range.Cells(row_idx, col_idx).Value
                                                    if cell_value:
                                                        row_texts.append(str(cell_value))
                                                
                                                if row_texts:
                                                    texts.append(" ".join(row_texts))
                                    
                                    workbook.Close(False)
                                    excel.Quit()
                                    
                                    pythoncom.CoUninitialize()
                                    
                                    result = "\n".join(texts)
                                    self.log_debug(f"Estratti {len(result)} caratteri da XLS con win32com")
                                    return result
                                    
                                except Exception as e:
                                    self.log_debug(f"Errore nell'apertura dell'XLS con win32com: {str(e)}")
                                    
                                    # Cleanup in caso di errore
                                    try:
                                        if 'workbook' in locals() and workbook:
                                            workbook.Close(False)
                                        if 'excel' in locals() and excel:
                                            excel.Quit()
                                    except:
                                        pass
                                    
                                    pythoncom.CoUninitialize()
                                    return ""
                                    
                            except ImportError:
                                self.log_debug("Nessuna libreria disponibile per i file XLS")
                                return ""
                        else:
                            self.log_debug("Nessuna libreria disponibile per i file XLS su questa piattaforma")
                            return ""
                            
                except Exception as e:
                    self.log_debug(f"Errore generale nell'elaborazione del file XLS {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # PowerPoint PPTX
            elif ext == '.pptx':
                try:
                    import pptx
                    self.log_debug(f"Processando file PPTX: {file_path}")
                    
                    presentation = pptx.Presentation(file_path)
                    texts = []
                    
                    for i, slide in enumerate(presentation.slides):
                        slide_text = []
                        texts.append(f"--- Diapositiva {i+1} ---")
                        
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text.append(shape.text)
                        
                        if slide_text:
                            texts.append("\n".join(slide_text))
                    
                    result = "\n".join(texts)
                    self.log_debug(f"Estratti {len(result)} caratteri da PPTX")
                    return result
                    
                except ImportError:
                    self.log_debug("Libreria python-pptx non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file PPTX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # PowerPoint PPT (vecchio formato)
            elif ext == '.ppt':
                # Su Windows, prova con pywin32
                if os.name == 'nt':
                    try:
                        import win32com.client
                        import pythoncom
                        
                        # Inizializzazione necessaria per i thread
                        pythoncom.CoInitialize()
                        
                        self.log_debug(f"Processando file PPT: {file_path}")
                        
                        try:
                            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                            powerpoint.Visible = False
                            
                            presentation = powerpoint.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
                            texts = []
                            
                            for slide_idx in range(1, presentation.Slides.Count + 1):
                                slide = presentation.Slides.Item(slide_idx)
                                texts.append(f"--- Diapositiva {slide_idx} ---")
                                slide_text = []
                                
                                for shape_idx in range(1, slide.Shapes.Count + 1):
                                    shape = slide.Shapes.Item(shape_idx)
                                    if shape.HasTextFrame:
                                        if shape.TextFrame.HasText:
                                            slide_text.append(shape.TextFrame.TextRange.Text)
                                
                                if slide_text:
                                    texts.append("\n".join(slide_text))
                            
                            presentation.Close()
                            powerpoint.Quit()
                            
                            pythoncom.CoUninitialize()
                            
                            result = "\n".join(texts)
                            self.log_debug(f"Estratti {len(result)} caratteri da PPT")
                            return result
                            
                        except Exception as e:
                            self.log_debug(f"Errore nell'apertura del PPT con win32com: {str(e)}")
                            
                            # Cleanup in caso di errore
                            try:
                                if 'presentation' in locals() and presentation:
                                    presentation.Close()
                                if 'powerpoint' in locals() and powerpoint:
                                    powerpoint.Quit()
                            except:
                                pass
                            
                            pythoncom.CoUninitialize()
                            return ""
                            
                    except ImportError:
                        self.log_debug("win32com non disponibile per i file PPT")
                        return ""
                else:
                    self.log_debug("Estrazione di testo dai file PPT non supportata su questa piattaforma")
                    return ""
            
            if self.stop_search:
                return ""
            # Rich Text Format (RTF) - NUOVO
            elif ext == '.rtf':
                try:
                    from striprtf.striprtf import rtf_to_text
                    self.log_debug(f"Processando file RTF: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        rtf_content = f.read()
                        content = rtf_to_text(rtf_content)
                        self.log_debug(f"Estratti {len(content)} caratteri da RTF")
                        return content
                except ImportError:
                    self.log_debug("Libreria striprtf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file RTF {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""        
            # OpenDocument Text (ODT) - NUOVO
            elif ext == '.odt':
                try:
                    from odf import opendocument, text
                    self.log_debug(f"Processando file ODT: {file_path}")
                    
                    doc = opendocument.load(file_path)
                    paragraphs = []
                    
                    # Estrai tutto il testo dai paragrafi
                    for element in doc.getElementsByType(text.P):
                        paragraphs.append(element.firstChild.data if element.firstChild else "")
                        
                    content = "\n".join(paragraphs)
                    self.log_debug(f"Estratti {len(content)} caratteri da ODT")
                    return content
                except ImportError:
                    self.log_debug("Libreria odf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODT {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""       
            # OpenDocument Spreadsheet (ODS) - NUOVO
            elif ext == '.ods':
                try:
                    from odf import opendocument, table, text
                    self.log_debug(f"Processando file ODS: {file_path}")
                    
                    doc = opendocument.load(file_path)
                    sheets = []
                    
                    # Estrai tutte le tabelle e celle
                    for sheet in doc.getElementsByType(table.Table):
                        sheet_rows = []
                        sheet_name = sheet.getAttribute('table:name')
                        sheets.append(f"--- Foglio: {sheet_name} ---")
                        
                        for row in sheet.getElementsByType(table.TableRow):
                            row_cells = []
                            for cell in row.getElementsByType(table.TableCell):
                                cell_text = ""
                                for p in cell.getElementsByType(text.P):
                                    cell_text += p.firstChild.data if p.firstChild else ""
                                row_cells.append(cell_text)
                            if row_cells:
                                sheet_rows.append(" | ".join(row_cells))
                                
                        sheets.append("\n".join(sheet_rows))
                        
                    content = "\n".join(sheets)
                    self.log_debug(f"Estratti {len(content)} caratteri da ODS")
                    return content
                except ImportError:
                    self.log_debug("Libreria odf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODS {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # OpenDocument Presentation (ODP) - NUOVO
            elif ext == '.odp':
                try:
                    from odf import opendocument, draw, text
                    self.log_debug(f"Processando file ODP: {file_path}")
                    
                    doc = opendocument.load(file_path)
                    slides = []
                    
                    # Estrai tutte le diapositive
                    for page in doc.getElementsByType(draw.Page):
                        slide_text = []
                        slide_name = page.getAttribute('draw:name')
                        slides.append(f"--- Diapositiva: {slide_name} ---")
                        
                        # Estrai il testo dai frame e forme
                        for element in page.getElementsByType(draw.Frame):
                            for p in element.getElementsByType(text.P):
                                if p.firstChild:
                                    slide_text.append(p.firstChild.data)
                                    
                        slides.append("\n".join(slide_text))
                        
                    content = "\n".join(slides)
                    self.log_debug(f"Estratti {len(content)} caratteri da ODP")
                    return content
                except ImportError:
                    self.log_debug("Libreria odf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODP {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""        
            # EPUB - NUOVO
            elif ext == '.epub':
                try:
                    import ebooklib
                    from ebooklib import epub
                    from bs4 import BeautifulSoup
                    self.log_debug(f"Processando file EPUB: {file_path}")
                    
                    # Funzione per estrarre testo dall'HTML
                    def chapter_to_text(content):
                        soup = BeautifulSoup(content, 'html.parser')
                        return soup.get_text()
                    
                    book = epub.read_epub(file_path)
                    chapters = []
                    
                    # Estrai metadati
                    title = book.get_metadata('DC', 'title')
                    if title:
                        chapters.append(f"Titolo: {title[0][0]}")
                    
                    authors = book.get_metadata('DC', 'creator')
                    if authors:
                        chapters.append(f"Autore: {authors[0][0]}")
                        
                    # Estrai contenuto
                    for item in book.get_items():
                        if item.get_type() == ebooklib.ITEM_DOCUMENT:
                            chapters.append(chapter_to_text(item.get_content()))
                            
                    content = "\n".join(chapters)
                    self.log_debug(f"Estratti {len(content)} caratteri da EPUB")
                    return content
                except ImportError:
                    self.log_debug("Librerie ebooklib o bs4 non disponibili")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file EPUB {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # MOBI - NUOVO
            elif ext == '.mobi':
                try:
                    import mobi
                    import tempfile
                    import shutil
                    self.log_debug(f"Processando file MOBI: {file_path}")
                    
                    tempdir = tempfile.mkdtemp()
                    try:
                        # Estrai il contenuto del file MOBI
                        extractor = mobi.Mobi(file_path)
                        extractor.extract(tempdir)
                        
                        # Leggi il testo estratto
                        text_content = []
                        for filename in os.listdir(tempdir):
                            if filename.endswith('.txt'):
                                with open(os.path.join(tempdir, filename), 'r', encoding='utf-8', errors='replace') as f:
                                    text_content.append(f.read())
                        
                        content = "\n".join(text_content)
                        self.log_debug(f"Estratti {len(content)} caratteri da MOBI")
                        return content
                    finally:
                        # Pulisci i file temporanei
                        shutil.rmtree(tempdir)
                except ImportError:
                    self.log_debug("Libreria mobi non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file MOBI {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""       
            # LaTeX - NUOVO
            elif ext == '.tex':
                try:
                    self.log_debug(f"Processando file LaTeX: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        
                    # Rimuovi i comandi LaTeX più comuni
                    import re
                    # Rimuovi i comandi
                    content = re.sub(r'\\[a-zA-Z]+(\{[^}]*\}|\[[^\]]*\])*', ' ', content)
                    # Rimuovi gli ambienti
                    content = re.sub(r'\\begin\{[^}]*\}(.*?)\\end\{[^}]*\}', ' ', content, flags=re.DOTALL)
                    # Rimuovi i commenti
                    content = re.sub(r'%.*?(\n|$)', ' ', content)
                    # Rimuovi le graffe
                    content = re.sub(r'\{|\}', '', content)
                    # Sostituisci più spazi con uno solo
                    content = re.sub(r'\s+', ' ', content)
                    
                    self.log_debug(f"Estratti {len(content)} caratteri da LaTeX")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file LaTeX {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # reStructuredText - NUOVO
            elif ext == '.rst':
                try:
                    self.log_debug(f"Processando file reStructuredText: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        
                    # Rimuovi gli elementi di markup più comuni
                    import re
                    # Rimuovi i titoli
                    content = re.sub(r'(=+|-+|~+|\^+|"+)\n', ' ', content)
                    # Rimuovi i link
                    content = re.sub(r'`[^`]*`_', ' ', content)
                    # Rimuovi i riferimenti alle direttive
                    content = re.sub(r'\.\. [a-z]+::', ' ', content)
                    
                    self.log_debug(f"Estratti {len(content)} caratteri da RST")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file RST {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""        
            # SQLite database (.db, .sqlite, .sqlite3) - NUOVO
            elif ext in ['.db', '.sqlite', '.sqlite3']:
                try:
                    import sqlite3
                    self.log_debug(f"Processando file SQLite: {file_path}")
                    
                    # Verifica che sia un file SQLite valido
                    if not os.path.getsize(file_path) > 100:
                        return ""
                        
                    try:
                        # Connetti al database
                        conn = sqlite3.connect(file_path)
                        cursor = conn.cursor()
                        
                        # Ottieni la lista delle tabelle
                        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                        tables = cursor.fetchall()
                        
                        content_parts = []
                        content_parts.append(f"Database SQLite: {os.path.basename(file_path)}")
                        content_parts.append(f"Numero di tabelle: {len(tables)}")
                        
                        # Estrai struttura e campione di dati da ogni tabella
                        for table in tables:
                            table_name = table[0]
                            content_parts.append(f"\n--- Tabella: {table_name} ---")
                            
                            # Ottieni struttura della tabella
                            cursor.execute(f"PRAGMA table_info({table_name});")
                            columns = cursor.fetchall()
                            col_names = [col[1] for col in columns]
                            content_parts.append("Colonne: " + ", ".join(col_names))
                            
                            # Ottieni un campione di dati (massimo 10 righe)
                            try:
                                cursor.execute(f"SELECT * FROM {table_name} LIMIT 10;")
                                rows = cursor.fetchall()
                                if rows:
                                    content_parts.append(f"Campione dati ({len(rows)} righe):")
                                    for row in rows:
                                        content_parts.append(str(row))
                            except:
                                content_parts.append("Errore nell'estrazione del campione dati")
                        
                        conn.close()
                        content = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(content)} caratteri da SQLite")
                        return content
                        
                    except sqlite3.Error:
                        # Non è un database SQLite valido o è cifrato
                        return ""
                except ImportError:
                    self.log_debug("Libreria sqlite3 non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file SQLite {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # Microsoft Access (.mdb, .accdb) - NUOVO
            elif ext in ['.mdb', '.accdb']:
                try:
                    import pyodbc
                    self.log_debug(f"Processando file Access: {file_path}")
                    
                    # Verifica che siamo su Windows
                    if os.name != 'nt':
                        self.log_debug("L'accesso ai file MDB/ACCDB è supportato solo su Windows")
                        return ""
                    
                    # Connetti al database Access
                    driver = "Microsoft Access Driver (*.mdb, *.accdb)"
                    conn_str = f"Driver={{{driver}}};DBQ={file_path};"
                    
                    try:
                        conn = pyodbc.connect(conn_str)
                        cursor = conn.cursor()
                        
                        # Ottieni l'elenco delle tabelle
                        tables = []
                        for row in cursor.tables():
                            if row.table_type == 'TABLE':
                                tables.append(row.table_name)
                        
                        content_parts = []
                        content_parts.append(f"Database Access: {os.path.basename(file_path)}")
                        content_parts.append(f"Numero di tabelle: {len(tables)}")
                        
                        # Estrai struttura e campione di dati da ogni tabella
                        for table in tables:
                            content_parts.append(f"\n--- Tabella: {table} ---")
                            
                            # Ottieni struttura della tabella
                            columns = cursor.columns(table=table)
                            col_names = [col.column_name for col in columns]
                            content_parts.append("Colonne: " + ", ".join(col_names))
                            
                            # Ottieni un campione di dati (massimo 10 righe)
                            try:
                                cursor.execute(f"SELECT * FROM [{table}] LIMIT 10")
                                rows = cursor.fetchall()
                                if rows:
                                    content_parts.append(f"Campione dati ({len(rows)} righe):")
                                    for row in rows:
                                        content_parts.append(str(row))
                            except:
                                content_parts.append("Errore nell'estrazione del campione dati")
                        
                        conn.close()
                        content = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(content)} caratteri da Access")
                        return content
                        
                    except pyodbc.Error as e:
                        self.log_debug(f"Errore di accesso al database: {str(e)}")
                        return ""
                except ImportError:
                    self.log_debug("Libreria pyodbc non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Access {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # OpenDocument Database (.odb) - NUOVO
            elif ext == '.odb':
                try:
                    import zipfile
                    import xml.etree.ElementTree as ET
                    self.log_debug(f"Processando file ODB: {file_path}")
                    
                    # ODB è essenzialmente un file ZIP con file XML all'interno
                    if zipfile.is_zipfile(file_path):
                        with zipfile.ZipFile(file_path, 'r') as zip_ref:
                            content_parts = []
                            
                            # Estrai il file contenente lo schema
                            try:
                                with zip_ref.open('content.xml') as content_file:
                                    tree = ET.parse(content_file)
                                    root = tree.getroot()
                                    
                                    # Cerca namespace
                                    ns = {'db': 'urn:oasis:names:tc:opendocument:xmlns:database:1.0',
                                        'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'}
                                    
                                    # Estrai informazioni sulle tabelle
                                    tables = root.findall('.//db:table', ns)
                                    content_parts.append(f"Database ODB: {os.path.basename(file_path)}")
                                    content_parts.append(f"Numero di tabelle trovate: {len(tables)}")
                                    
                                    for table in tables:
                                        if 'db:name' in table.attrib:
                                            table_name = table.get('{urn:oasis:names:tc:opendocument:xmlns:database:1.0}name')
                                            content_parts.append(f"\n--- Tabella: {table_name} ---")
                                            
                                            # Estrai colonne
                                            columns = table.findall('.//db:column', ns)
                                            col_names = []
                                            for column in columns:
                                                if 'db:name' in column.attrib:
                                                    col_name = column.get('{urn:oasis:names:tc:opendocument:xmlns:database:1.0}name')
                                                    col_names.append(col_name)
                                            
                                            content_parts.append("Colonne: " + ", ".join(col_names))
                            except Exception as e:
                                content_parts.append(f"Errore nell'estrazione dello schema: {str(e)}")
                            
                            content = "\n".join(content_parts)
                            self.log_debug(f"Estratti {len(content)} caratteri da ODB")
                            return content
                    else:
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODB {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""       
            # Tab-Separated Values (.tsv) - NUOVO
            elif ext == '.tsv':
                try:
                    import csv
                    self.log_debug(f"Processando file TSV: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        reader = csv.reader(f, delimiter='\t')
                        rows = []
                        
                        # Limita a 1000 righe per file grandi
                        for i, row in enumerate(reader):
                            if i >= 1000:
                                rows.append("... (file troncato, troppe righe)")
                                break
                            rows.append("\t".join(row))
                            
                    content = "\n".join(rows)
                    self.log_debug(f"Estratti {len(content)} caratteri da TSV")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file TSV {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # dBase format (.dbf) - NUOVO
            elif ext == '.dbf':
                try:
                    import dbfread
                    self.log_debug(f"Processando file DBF: {file_path}")
                    
                    table = dbfread.DBF(file_path)
                    records = []
                    
                    # Ottieni nomi delle colonne
                    headers = table.field_names
                    records.append("Colonne: " + ", ".join(headers))
                    
                    # Ottieni un campione di dati
                    for i, record in enumerate(table):
                        if i >= 50:  # Limita a 50 record
                            records.append("... (file troncato, troppi record)")
                            break
                            
                        record_data = []
                        for field in headers:
                            record_data.append(f"{field}: {record[field]}")
                        records.append(" | ".join(record_data))
                        
                    content = "\n".join(records)
                    self.log_debug(f"Estratti {len(content)} caratteri da DBF")
                    return content
                except ImportError:
                    self.log_debug("Libreria dbfread non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file DBF {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # Data Interchange Format (.dif) - NUOVO
            elif ext == '.dif':
                try:
                    self.log_debug(f"Processando file DIF: {file_path}")
                    
                    # I file DIF hanno una struttura specifica
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        lines = f.readlines()
                        
                    if len(lines) < 3:
                        return ""
                        
                    content_lines = []
                    i = 0
                    
                    # Salta l'intestazione
                    while i < len(lines) and not lines[i].strip().startswith('DATA'):
                        i += 1
                        
                    # Estrai i dati
                    data_mode = False
                    current_row = []
                    
                    while i < len(lines):
                        line = lines[i].strip()
                        
                        if line.startswith('BOT'):  # Beginning of tuple
                            current_row = []
                        elif line.startswith('EOD'):  # End of data
                            break
                        elif not data_mode and line.startswith('1,0'):
                            data_mode = True
                        elif data_mode and line.startswith('V') or line.startswith('C'):
                            # La prossima riga contiene il valore
                            i += 1
                            if i < len(lines):
                                value = lines[i].strip().strip('"')
                                current_row.append(value)
                        elif line.startswith('EOT'):  # End of tuple
                            content_lines.append(",".join(current_row))
                            data_mode = False
                            
                        i += 1
                        
                    content = "\n".join(content_lines)
                    self.log_debug(f"Estratti {len(content)} caratteri da DIF")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file DIF {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Testo semplice
            elif ext in ['.txt', '.csv', '.log', '.ini', '.xml', '.json', '.md', '.html', '.htm',
                        '.py', '.js', '.java', '.cpp', '.c', '.cs', '.php', '.rb', '.go', '.swift', 
                        '.sql', '.sh', '.bat', '.ps1', '.vbs', '.pl', '.ts', '.kt', '.scala',
                        '.h', '.hpp', '.vb', '.lua', '.rs', '.groovy', '.yml', '.yaml', '.toml',
                        '.properties', '.conf', '.config', '.cfg', '.reg']:
                try:
                    # Apri con diverse codifiche per essere robusto
                    encodings = ['utf-8', 'latin-1', 'windows-1252']
                    content = ""
                    
                    for encoding in encodings:
                        try:
                            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                                content = f.read(1024*1024)  # Leggi al massimo 1 MB
                                break
                        except UnicodeDecodeError:
                            continue
                        except Exception as e:
                            self.log_debug(f"Errore con codifica {encoding}: {str(e)}")
                    
                    if content:
                        self.log_debug(f"Estratti {len(content)} caratteri da file di testo")
                        return content
                    else:
                        self.log_debug("Nessun contenuto estratto dal file di testo")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file di testo {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # PDF (aggiunto per completezza)
            elif ext == '.pdf':
                try:
                    import PyPDF2
                    self.log_debug(f"Processando file PDF: {file_path}")
                    
                    content = []
                    with open(file_path, 'rb') as f:
                        try:
                            reader = PyPDF2.PdfReader(f)
                            num_pages = min(len(reader.pages), 50)  # Limita a 50 pagine
                            
                            for page_num in range(num_pages):
                                try:
                                    page_text = reader.pages[page_num].extract_text()
                                    if page_text and page_text.strip():
                                        content.append(f"--- Pagina {page_num+1} ---")
                                        content.append(page_text)
                                except Exception as e:
                                    self.log_debug(f"Errore nell'estrazione testo pagina {page_num}: {str(e)}")
                        except Exception as e:
                            self.log_debug(f"Errore nell'apertura del PDF: {str(e)}")
                    
                    result = "\n".join(content)
                    self.log_debug(f"Estratti {len(result)} caratteri da PDF")
                    return result
                    
                except ImportError:
                    self.log_debug("Libreria PyPDF2 non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file PDF: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # ===== EMAIL E CALENDARIO =====
            # File Email (.eml)
            elif ext == '.eml':
                try:
                    self.log_debug(f"Processando file Email: {file_path}")
                    import email
                    import base64
                    import quopri
                    from email.header import decode_header
                    
                    # Leggi il file come binario per gestire correttamente tutti i tipi di codifica
                    with open(file_path, 'rb') as f:
                        msg_data = f.read()
                        
                    # Verifica che abbiamo letto dei dati
                    if not msg_data:
                        self.log_debug(f"File EML vuoto: {file_path}")
                        return ""
                        
                    self.log_debug(f"Letti {len(msg_data)} bytes dal file EML")
                    
                    # Analizza il messaggio email
                    msg = email.message_from_bytes(msg_data)
                    content_parts = []
                    
                    # Estrai intestazioni
                    self.log_debug("Estrazione intestazioni email")
                    for header in ['From', 'To', 'Subject', 'Date']:
                        if msg[header]:
                            # Decodifica l'intestazione se necessario
                            try:
                                header_parts = decode_header(msg[header])
                                decoded_header = ""
                                for part, encoding in header_parts:
                                    if isinstance(part, bytes):
                                        decoded_header += part.decode(encoding or 'utf-8', errors='replace')
                                    else:
                                        decoded_header += part
                                content_parts.append(f"{header}: {decoded_header}")
                            except:
                                content_parts.append(f"{header}: {msg[header]}")
                    
                    # Aggiungi un separatore dopo le intestazioni
                    if content_parts:
                        content_parts.append("-" * 40)
                    
                    # Estrai corpo del messaggio e allegati
                    attachment_count = 0
                    self.log_debug("Inizio analisi parti dell'email")
                    
                    # Elabora tutte le parti del messaggio
                    for part in msg.walk():
                        content_type = part.get_content_type()
                        self.log_debug(f"Elaborazione parte email: {content_type}")
                        
                        # Estrai testo semplice dal corpo dell'email
                        if content_type == "text/plain" and not part.get_filename():
                            try:
                                payload = part.get_payload(decode=True)
                                if payload:
                                    charset = part.get_content_charset() or 'utf-8'
                                    text = payload.decode(charset, errors='replace')
                                    content_parts.append(text)
                                    self.log_debug(f"Estratti {len(text)} caratteri di testo")
                            except Exception as e:
                                self.log_debug(f"Errore nell'estrazione del testo: {str(e)}")
                                try:
                                    # Fallback a utf-8
                                    content_parts.append(payload.decode('utf-8', errors='replace'))
                                except:
                                    pass
                                    
                        # Estrai HTML dal corpo dell'email
                        elif content_type == "text/html" and not part.get_filename():
                            try:
                                payload = part.get_payload(decode=True)
                                if payload:
                                    charset = part.get_content_charset() or 'utf-8'
                                    html_content = payload.decode(charset, errors='replace')
                                    
                                    # Opzionalmente, prova a estrarre il testo dall'HTML
                                    try:
                                        from bs4 import BeautifulSoup
                                        soup = BeautifulSoup(html_content, 'html.parser')
                                        text_content = soup.get_text(separator=' ', strip=True)
                                        content_parts.append(text_content)
                                        self.log_debug(f"Estratti {len(text_content)} caratteri da HTML")
                                    except ImportError:
                                        # Se BeautifulSoup non è disponibile, usa l'HTML grezzo
                                        content_parts.append(html_content)
                                        self.log_debug(f"BeautifulSoup non disponibile, usato HTML grezzo ({len(html_content)} caratteri)")
                            except Exception as e:
                                self.log_debug(f"Errore nell'estrazione HTML: {str(e)}")
                        
                        # Gestisci gli allegati
                        elif part.get('Content-Disposition') or part.get_filename():
                            try:
                                # Ottieni il nome dell'allegato
                                filename = part.get_filename()
                                attachment_count += 1
                                
                                # Normalizza il nome dell'allegato
                                if not filename:
                                    filename = f"allegato_{attachment_count}"
                                
                                # Decodifica il nome dell'allegato se necessario
                                if isinstance(filename, bytes):
                                    try:
                                        filename = filename.decode('utf-8', errors='replace')
                                    except:
                                        filename = f"allegato_{attachment_count}"
                                
                                # Decodifica header encodati se necessario
                                try:
                                    decoded_parts = decode_header(filename)
                                    decoded_filename = ""
                                    for part_data, charset in decoded_parts:
                                        if isinstance(part_data, bytes):
                                            decoded_filename += part_data.decode(charset or 'utf-8', errors='replace')
                                        else:
                                            decoded_filename += part_data
                                    filename = decoded_filename
                                except:
                                    pass
                                
                                self.log_debug(f"Allegato trovato: {filename}")
                                
                                # Aggiungi informazioni sull'allegato
                                content_parts.append(f"\n--- ALLEGATO {attachment_count}: {filename} ---\n")
                                
                                # METODO 1: Usa get_payload(decode=True) che dovrebbe gestire automaticamente la decodifica
                                payload = part.get_payload(decode=True)
                                
                                # METODO 2: se il primo metodo fallisce, prova con decodifica manuale
                                if not payload:
                                    encoding = part.get('Content-Transfer-Encoding', '').lower()
                                    raw_payload = part.get_payload()
                                    
                                    if encoding == 'base64' and isinstance(raw_payload, str):
                                        try:
                                            payload = base64.b64decode(raw_payload)
                                            self.log_debug("Allegato decodificato con base64")
                                        except:
                                            self.log_debug("Errore nella decodifica base64")
                                    elif encoding == 'quoted-printable' and isinstance(raw_payload, str):
                                        try:
                                            payload = quopri.decodestring(raw_payload.encode('utf-8'))
                                            self.log_debug("Allegato decodificato con quoted-printable")
                                        except:
                                            self.log_debug("Errore nella decodifica quoted-printable")
                                    elif isinstance(raw_payload, str):
                                        try:
                                            payload = raw_payload.encode('utf-8', errors='replace')
                                            self.log_debug("Allegato convertito da stringa a bytes")
                                        except:
                                            self.log_debug("Errore nella conversione della stringa a bytes")
                                
                                if payload and len(payload) > 0:
                                    self.log_debug(f"Analisi contenuto allegato: {filename} ({len(payload)} bytes)")
                                    
                                    # IMPORTANTE: CHIAMA IL METODO PER ELABORARE L'ALLEGATO
                                    attachment_content = self.process_email_attachment(
                                        payload, filename, part.get_content_type())
                                    
                                    if attachment_content:
                                        content_parts.append(attachment_content)
                                        self.log_debug(f"Contenuto allegato {filename} aggiunto ({len(attachment_content)} caratteri)")
                                    else:
                                        content_parts.append(f"[Allegato {filename}: nessun contenuto estraibile]")
                                        self.log_debug(f"Nessun contenuto estraibile dall'allegato {filename}")
                                else:
                                    content_parts.append(f"[Allegato {filename}: vuoto o non decodificabile]")
                                    self.log_debug(f"Allegato vuoto o non decodificabile: {filename}")
                                    
                            except Exception as e:
                                self.log_debug(f"Errore nell'elaborazione dell'allegato: {str(e)}")
                                content_parts.append(f"[Errore nell'elaborazione dell'allegato: {str(e)}]")
                    
                    self.log_debug(f"Totale allegati trovati: {attachment_count}")
                    
                    # Se non è stato trovato alcun contenuto, prova un metodo più semplice
                    if not content_parts:
                        self.log_debug("Nessun contenuto trovato, tentativo con metodo alternativo")
                        try:
                            if msg.is_multipart():
                                for part in msg.walk():
                                    content_type = part.get_content_type()
                                    if content_type == "text/plain":
                                        payload = part.get_payload(decode=True)
                                        if payload:
                                            content_parts.append(payload.decode('utf-8', errors='replace'))
                            else:
                                payload = msg.get_payload(decode=True)
                                if payload:
                                    content_parts.append(payload.decode('utf-8', errors='replace'))
                        except Exception as e:
                            self.log_debug(f"Errore nel metodo alternativo: {str(e)}")
                    
                    content = "\n".join(content_parts)
                    self.log_debug(f"Estratti in totale {len(content)} caratteri da EML (inclusi allegati)")
                    return content
                except ImportError as e:
                    self.log_debug(f"Modulo necessario non disponibile: {str(e)}")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file EML {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File vCard (.vcf)
            elif ext == '.vcf':
                try:
                    self.log_debug(f"Processando file vCard: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        # Estrai campi più importanti per la ricerca
                        processed_content = []
                        for line in content.splitlines():
                            if line.startswith(('FN:', 'N:', 'EMAIL:', 'TEL:', 'ADR:', 'ORG:', 'TITLE:', 'NOTE:')):
                                processed_content.append(line)
                        
                        result = "\n".join(processed_content)
                        self.log_debug(f"Estratti {len(result)} caratteri da vCard")
                        return result
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file vCard {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File iCalendar (.ics)
            elif ext == '.ics':
                try:
                    self.log_debug(f"Processando file iCalendar: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        # Estrai campi più importanti per la ricerca
                        processed_content = []
                        current_event = []
                        in_event = False
                        
                        for line in content.splitlines():
                            line = line.strip()
                            if line == "BEGIN:VEVENT":
                                in_event = True
                                current_event = []
                            elif line == "END:VEVENT":
                                in_event = False
                                processed_content.append("\n".join(current_event))
                                processed_content.append("---")
                            elif in_event and line.startswith(('SUMMARY:', 'DESCRIPTION:', 'LOCATION:', 
                                                            'ORGANIZER:', 'ATTENDEE:', 'DTSTART:', 
                                                            'DTEND:', 'CATEGORIES:')):
                                current_event.append(line)
                        
                        result = "\n".join(processed_content)
                        self.log_debug(f"Estratti {len(result)} caratteri da iCalendar")
                        return result
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file iCalendar {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # ===== PRESENTAZIONI =====
            # File PowerPoint Show (.pps)
            elif ext == '.pps':
                try:
                    # Utilizziamo la stessa implementazione di PPT dato che hanno lo stesso formato
                    if os.name == 'nt':
                        import win32com.client
                        import pythoncom
                        
                        # Inizializzazione necessaria per i thread
                        pythoncom.CoInitialize()
                        
                        self.log_debug(f"Processando file PowerPoint Show: {file_path}")
                        
                        try:
                            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                            powerpoint.Visible = False
                            
                            presentation = powerpoint.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
                            texts = []
                            
                            for slide_idx in range(1, presentation.Slides.Count + 1):
                                slide = presentation.Slides.Item(slide_idx)
                                texts.append(f"--- Diapositiva {slide_idx} ---")
                                slide_text = []
                                
                                for shape_idx in range(1, slide.Shapes.Count + 1):
                                    shape = slide.Shapes.Item(shape_idx)
                                    if shape.HasTextFrame:
                                        if shape.TextFrame.HasText:
                                            slide_text.append(shape.TextFrame.TextRange.Text)
                                
                                if slide_text:
                                    texts.append("\n".join(slide_text))
                            
                            presentation.Close()
                            powerpoint.Quit()
                            
                            pythoncom.CoUninitialize()
                            
                            result = "\n".join(texts)
                            self.log_debug(f"Estratti {len(result)} caratteri da PPS")
                            return result
                            
                        except Exception as e:
                            self.log_debug(f"Errore nell'apertura del PPS con win32com: {str(e)}")
                            
                            # Cleanup in caso di errore
                            try:
                                if 'presentation' in locals() and presentation:
                                    presentation.Close()
                                if 'powerpoint' in locals() and powerpoint:
                                    powerpoint.Quit()
                            except:
                                pass
                            
                            pythoncom.CoUninitialize()
                            return ""
                    else:
                        self.log_debug("Estrazione di testo dai file PPS non supportata su questa piattaforma")
                        return ""
                except ImportError:
                    self.log_debug("win32com non disponibile per i file PPS")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore generale nell'elaborazione del file PPS {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File Keynote (.key)
            elif ext == '.key':
                try:
                    self.log_debug(f"Processando file Keynote: {file_path}")
                    # Keynote è essenzialmente un pacchetto compresso con file XML all'interno
                    if not zipfile.is_zipfile(file_path):
                        self.log_debug(f"File Keynote non valido (non è un file zip): {file_path}")
                        return ""
                        
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        content_parts = []
                        
                        # Cerca file di testo nel pacchetto Keynote
                        for file_info in zip_ref.infolist():
                            if file_info.filename.endswith(('.xml', '.txt')):
                                try:
                                    with zip_ref.open(file_info) as content_file:
                                        content = content_file.read().decode('utf-8', errors='replace')
                                        # Estrai solo il testo dalle presentazioni, rimuovendo i tag
                                        import re
                                        text_content = re.sub(r'<[^>]+>', ' ', content)
                                        text_content = re.sub(r'\s+', ' ', text_content).strip()
                                        if text_content:
                                            content_parts.append(text_content)
                                except:
                                    continue
                                    
                        result = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(result)} caratteri da Keynote")
                        return result
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Keynote {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # ===== FILE DI CONFIGURAZIONE =====
            # File YAML (.yml, .yaml)
            elif ext in ['.yml', '.yaml']:
                try:
                    self.log_debug(f"Processando file YAML: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da YAML")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file YAML {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File TOML (.toml)
            elif ext == '.toml':
                try:
                    self.log_debug(f"Processando file TOML: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da TOML")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file TOML {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File Registry Windows (.reg)
            elif ext == '.reg':
                try:
                    self.log_debug(f"Processando file Registry: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da Registry")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Registry {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File plist (.plist)
            elif ext == '.plist':
                try:
                    self.log_debug(f"Processando file Property List: {file_path}")
                    # Prova a leggere come XML
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        # Rimuovi i tag XML per estrarre solo il testo
                        import re
                        text_content = re.sub(r'<[^>]+>', ' ', content)
                        text_content = re.sub(r'\s+', ' ', text_content).strip()
                        self.log_debug(f"Estratti {len(text_content)} caratteri da plist")
                        return text_content
                except UnicodeDecodeError:
                    try:
                        # Potrebbe essere un file plist binario
                        self.log_debug("Tentativo di lettura come plist binario")
                        with open(file_path, 'rb') as f:
                            content = f.read().decode('latin-1', errors='replace')
                            # Estrai stringhe leggibili
                            printable = ''.join(c for c in content if c.isprintable() and len(c.strip()) > 0)
                            self.log_debug(f"Estratti {len(printable)} caratteri da plist binario")
                            return printable
                    except Exception as inner_e:
                        self.log_debug(f"Errore nella lettura del plist binario: {str(inner_e)}")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file plist {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File Properties (.properties)
            elif ext == '.properties':
                try:
                    self.log_debug(f"Processando file Properties: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da Properties")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Properties {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File htaccess (.htaccess)
            elif file_path.endswith('.htaccess'):
                try:
                    self.log_debug(f"Processando file htaccess: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da htaccess")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file htaccess {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # ===== LINGUAGGI DI PROGRAMMAZIONE =====
            # I vari linguaggi di programmazione possono usare lo stesso parser di testo
            elif ext in ['.h', '.hpp', '.vb', '.lua', '.rs', '.groovy']:
                try:
                    self.log_debug(f"Processando file di codice {ext}: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da file {ext}")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file {ext} {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File MSG (Outlook)
            elif ext == '.msg':
                try:
                    self.log_debug(f"Processando file MSG Outlook: {file_path}")
                    try:
                        import extract_msg
                        msg = extract_msg.Message(file_path)
                        
                        content_parts = []
                        
                        # Estrai intestazioni principali
                        from_address = msg.sender if hasattr(msg, 'sender') else "N/A"
                        to_address = msg.to if hasattr(msg, 'to') else "N/A"
                        subject = msg.subject if hasattr(msg, 'subject') else "N/A"
                        msg_date = msg.date if hasattr(msg, 'date') else "N/A"
                        
                        # Aggiungi intestazioni al contenuto
                        content_parts.append(f"Da: {from_address}")
                        content_parts.append(f"A: {to_address}")
                        content_parts.append(f"Oggetto: {subject}")
                        content_parts.append(f"Data: {msg_date}")
                        content_parts.append("-" * 40)  # Separatore
                        
                        # Estrai corpo del messaggio
                        msg_body = msg.body if hasattr(msg, 'body') else ""
                        if msg_body:
                            content_parts.append(msg_body)
                        
                        # Contatori per monitoraggio allegati
                        attachment_count = 0
                        
                        # Processa allegati esattamente come fatto per gli EML
                        for attachment in msg.attachments:
                            if not hasattr(attachment, 'name') or not attachment.name:
                                continue
                            
                            attachment_count += 1
                            filename = attachment.name
                            content_parts.append(f"\n--- ALLEGATO {attachment_count}: {filename} ---\n")
                            
                            try:
                                # Ottieni i dati binari dell'allegato
                                attachment_data = None
                                if hasattr(attachment, 'data'):
                                    attachment_data = attachment.data
                                elif hasattr(attachment, 'getBytes'):
                                    attachment_data = attachment.getBytes()
                                
                                if not attachment_data:
                                    content_parts.append(f"[Allegato vuoto o non leggibile]")
                                    continue
                                
                                # Ottieni il tipo di contenuto (MIME type)
                                content_type = ""
                                if hasattr(attachment, 'mimetype') and attachment.mimetype:
                                    content_type = attachment.mimetype
                                
                                # IMPORTANTE: Usa esattamente la stessa funzione di EML
                                # per processare l'allegato ed estrarre il contenuto
                                attachment_content = self.process_email_attachment(
                                    attachment_data, filename, content_type)
                                
                                if attachment_content:
                                    # Aggiungi il contenuto estratto ai risultati
                                    content_parts.append(attachment_content)
                                    self.log_debug(f"Estratto contenuto da allegato '{filename}': {len(attachment_content)} caratteri")
                                else:
                                    content_parts.append(f"[Allegato {filename}: nessun contenuto estraibile]")
                                    
                            except Exception as e:
                                self.log_debug(f"Errore nell'elaborazione dell'allegato {filename}: {str(e)}")
                                content_parts.append(f"[Errore nell'elaborazione dell'allegato: {str(e)}]")
                        
                        content = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(content)} caratteri da MSG (inclusi {attachment_count} allegati)")
                        return content
                        
                    except ImportError:
                        self.log_debug("Libreria extract_msg non disponibile")
                        # Metodo fallback - estrai contenuto usando regexp base
                        with open(file_path, 'rb') as f:
                            binary_content = f.read()
                            text_content = ""
                            # Cerca stringhe ASCII leggibili nel file binario
                            import re
                            text_chunks = re.findall(b'[\x20-\x7E\r\n]{4,}', binary_content)
                            for chunk in text_chunks:
                                try:
                                    text_content += chunk.decode('utf-8', errors='replace') + "\n"
                                except:
                                    pass
                            
                            if text_content:
                                self.log_debug(f"Estratti {len(text_content)} caratteri da MSG (metodo fallback)")
                                return text_content
                            return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file MSG {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File PST/OST (Outlook database) - versione semplificata senza dipendenze esterne
            elif ext in ['.pst', '.ost']:
                try:
                    self.log_debug(f"Processando file {ext} Outlook: {file_path}")
                    
                    # Estrai metadati di base
                    file_size = os.path.getsize(file_path)
                    content_parts = []
                    
                    content_parts.append(f"File {ext.upper()} Outlook")
                    content_parts.append(f"Percorso: {file_path}")
                    content_parts.append(f"Dimensione: {self._format_size(file_size)}")
                    content_parts.append(f"Data modifica: {datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M')}")
                    
                    # Estrai stringhe di testo leggibili (approccio base)
                    try:
                        with open(file_path, 'rb') as f:
                            binary_content = f.read(1024*1024)  # Leggi solo il primo MB
                            
                            # Usa espressione regolare per trovare stringhe ASCII leggibili
                            import re
                            # Trova stringhe alfanumeriche con spazi/punteggiatura di almeno 5 caratteri
                            pattern = re.compile(b'[a-zA-Z0-9\\s\\.,@\\-_:;\'"/]{5,}')
                            matches = pattern.findall(binary_content)
                            
                            # Filtra e converte le stringhe trovate
                            strings = []
                            for match in matches:
                                try:
                                    text = match.decode('utf-8', errors='replace')
                                    # Filtra stringhe che sembrano email valide o messaggi
                                    if ('@' in text or 
                                        text.startswith("To:") or 
                                        text.startswith("From:") or 
                                        text.startswith("Subject:")):
                                        strings.append(text)
                                    # O stringhe abbastanza lunghe da essere significative
                                    elif len(text) > 15:
                                        strings.append(text)
                                except:
                                    pass
                            
                            # Aggiungi i risultati all'output
                            if strings:
                                content_parts.append("\n--- Contenuto estratto ---\n")
                                content_parts.extend(strings)
                    except Exception as e:
                        self.log_debug(f"Errore nell'estrazione del testo: {str(e)}")
                    
                    content = "\n".join(content_parts)
                    self.log_debug(f"Estratti {len(content)} caratteri da {ext}")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file {ext} {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File MBOX - usa il modulo mailbox standard
            elif ext == '.mbox':
                try:
                    self.log_debug(f"Processando file MBOX: {file_path}")
                    import mailbox
                    import email
                    
                    # Apri il file MBOX
                    mbox = mailbox.mbox(file_path)
                    content_parts = []
                    
                    # Limita il numero di messaggi da processare
                    max_messages = 50
                    processed = 0
                    
                    content_parts.append(f"File MBOX: {os.path.basename(file_path)}")
                    content_parts.append(f"Totale messaggi: {len(mbox)}")
                    content_parts.append("-" * 40)
                    
                    # Itera attraverso i messaggi
                    for key, msg in mbox.items():
                        if processed >= max_messages:
                            content_parts.append(f"[Limitato a {max_messages} messaggi]")
                            break
                        
                        try:
                            # Estrai le intestazioni principali
                            headers = []
                            for header in ['From', 'To', 'Subject', 'Date']:
                                if msg[header]:
                                    headers.append(f"{header}: {msg[header]}")
                            
                            content_parts.append(f"\n--- MESSAGGIO {processed+1} ---")
                            content_parts.extend(headers)
                            
                            # Estrai corpo del messaggio
                            if msg.is_multipart():
                                for part in msg.walk():
                                    content_type = part.get_content_type()
                                    if content_type == "text/plain" and not part.get_filename():
                                        try:
                                            payload = part.get_payload(decode=True)
                                            if payload:
                                                charset = part.get_content_charset() or 'utf-8'
                                                text = payload.decode(charset, errors='replace')
                                                content_parts.append(text)
                                        except:
                                            pass
                            else:
                                payload = msg.get_payload(decode=True)
                                if payload:
                                    charset = msg.get_content_charset() or 'utf-8'
                                    try:
                                        text = payload.decode(charset, errors='replace')
                                        content_parts.append(text)
                                    except:
                                        pass
                            
                            processed += 1
                        except Exception as e:
                            self.log_debug(f"Errore nel processare il messaggio MBOX: {str(e)}")
                    
                    mbox.close()
                    content = "\n".join(content_parts)
                    self.log_debug(f"Estratti {len(content)} caratteri da MBOX")
                    return content
                except ImportError:
                    self.log_debug("Libreria mailbox non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file MBOX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File EMLX (Apple Mail)
            elif ext == '.emlx':
                try:
                    self.log_debug(f"Processando file EMLX: {file_path}")
                    import email
                    
                    # I file EMLX hanno un formato particolare: prima riga è un numero, poi segue il messaggio email
                    with open(file_path, 'rb') as f:
                        content = f.read()
                        
                    # Separa il numero iniziale dal resto del contenuto email
                    try:
                        # La prima riga è un numero seguito da una nuova riga
                        parts = content.split(b'\n', 1)
                        if len(parts) > 1:
                            email_content = parts[1]
                            
                            # Parse dell'email con il modulo email standard
                            msg = email.message_from_bytes(email_content)
                            
                            # Da qui in poi possiamo usare lo stesso codice per gestire il messaggio
                            # come facciamo per i file EML standard
                            content_parts = []
                            
                            # Estrai intestazioni
                            for header in ['From', 'To', 'Subject', 'Date']:
                                if msg[header]:
                                    content_parts.append(f"{header}: {msg[header]}")
                            
                            # Estrai corpo del messaggio
                            if msg.is_multipart():
                                for part in msg.walk():
                                    content_type = part.get_content_type()
                                    if content_type == "text/plain" and not part.get_filename():
                                        payload = part.get_payload(decode=True)
                                        if payload:
                                            charset = part.get_content_charset() or 'utf-8'
                                            try:
                                                text = payload.decode(charset, errors='replace')
                                                content_parts.append(text)
                                            except:
                                                pass
                            else:
                                payload = msg.get_payload(decode=True)
                                if payload:
                                    try:
                                        text = payload.decode('utf-8', errors='replace')
                                        content_parts.append(text)
                                    except:
                                        pass
                            
                            result = "\n".join(content_parts)
                            self.log_debug(f"Estratti {len(result)} caratteri da EMLX")
                            return result
                    except Exception as e:
                        self.log_debug(f"Errore nell'analisi del file EMLX: {str(e)}")
                    
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file EMLX {file_path}: {str(e)}")
                    return ""
            
        except Exception as e:
            self.log_debug(f"Errore generale nella lettura del file {file_path}: {str(e)}")
            return ""
    
    @error_handler
    def process_email_attachment(self, attachment_data, attachment_name, content_type):
        """Process an email attachment and extract its content"""
        temp_dir = None
        temp_file_path = None
        
        try:
            import tempfile
            import os

            # Log dettagliato per il debug
            self.log_debug(f"Inizio elaborazione allegato: {attachment_name} ({len(attachment_data)} bytes)")
            
            # Sanitizza il nome del file per evitare caratteri problematici
            safe_name = ''.join(c for c in attachment_name if c.isalnum() or c in '._- ')
            if not safe_name:
                safe_name = f"attachment_{hash(attachment_data) % 10000}.bin"
            
            # Se il nome file non ha estensione ma abbiamo un content_type, aggiungila
            if '.' not in safe_name and content_type:
                if content_type == 'application/pdf':
                    safe_name += '.pdf'
                elif content_type == 'application/msword':
                    safe_name += '.doc'
                elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    safe_name += '.docx'
                elif content_type == 'application/vnd.ms-excel':
                    safe_name += '.xls'
                elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                    safe_name += '.xlsx'
                elif content_type == 'text/csv':
                    safe_name += '.csv'
                elif content_type == 'application/vnd.oasis.opendocument.text':
                    safe_name += '.odt'
                elif content_type == 'application/rtf' or content_type == 'text/rtf':
                    safe_name += '.rtf'
                elif content_type == 'application/xml' or content_type == 'text/xml':
                    safe_name += '.xml'
                elif content_type == 'text/plain':
                    safe_name += '.txt'
            
            # Crea il file temporaneo
            temp_dir = tempfile.mkdtemp(prefix="email_att_")
            temp_file_path = os.path.join(temp_dir, safe_name)
            
            # Salva l'allegato su file temporaneo
            with open(temp_file_path, 'wb') as f:
                f.write(attachment_data)
            
            self.log_debug(f"File temporaneo creato: {temp_file_path}")
            
            # Ottieni l'estensione del file
            _, ext = os.path.splitext(temp_file_path)
            ext = ext.lower()
            
            # IMPORTANTE: Verifica se l'estensione è supportata - log dettagliato
            supports_content_search = self.should_search_content(temp_file_path)
            self.log_debug(f"L'estensione {ext} supporta la ricerca nei contenuti? {supports_content_search}")
            
            # Force include common document types regardless of extension settings
            force_include = ext in ['.pdf', '.doc', '.docx', '.xls', '.xlsx', '.csv', '.odt', '.rtf', '.xml', '.txt']
            
            # Check if we should process this type of file
            if supports_content_search or force_include:
                # Skip nested email files to prevent recursion
                if ext == '.eml':
                    self.log_debug(f"Allegato email rilevato, evito ricorsione: {temp_file_path}")
                    return "[Contenuto allegato email non elaborato per evitare ricorsione]"
                
                # For supported formats, use the existing content extraction
                try:
                    self.log_debug(f"Inizio estrazione contenuto da {safe_name} ({ext})")
                    
                    # Gestione specifica per tipo di file
                    content = ""
                    
                    # File di testo semplice
                    if ext in ['.txt', '.csv', '.log', '.ini', '.xml', '.json', '.md']:
                        try:
                            with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                content = f.read()
                            self.log_debug(f"Letti {len(content)} caratteri da file di testo {ext}")
                        except Exception as text_err:
                            self.log_debug(f"Errore lettura testo: {str(text_err)}")
                    
                    # PDF
                    elif ext == '.pdf':
                        try:
                            import PyPDF2
                            with open(temp_file_path, 'rb') as f:
                                try:
                                    reader = PyPDF2.PdfReader(f)
                                    pdf_text = []
                                    for page_num in range(len(reader.pages)):
                                        page = reader.pages[page_num]
                                        pdf_text.append(page.extract_text())
                                    content = "\n".join(pdf_text)
                                    self.log_debug(f"Estratti {len(content)} caratteri da PDF")
                                except Exception as pdf_err:
                                    self.log_debug(f"Errore PDF: {str(pdf_err)}")
                        except ImportError:
                            self.log_debug("PyPDF2 non disponibile")
                    
                    # Word DOCX
                    elif ext == '.docx':
                        try:
                            import docx
                            doc = docx.Document(temp_file_path)
                            doc_text = []
                            for para in doc.paragraphs:
                                if para.text:
                                    doc_text.append(para.text)
                            content = "\n".join(doc_text)
                            self.log_debug(f"Estratti {len(content)} caratteri da DOCX")
                        except ImportError:
                            self.log_debug("python-docx non disponibile")
                        except Exception as docx_err:
                            self.log_debug(f"Errore DOCX: {str(docx_err)}")
                    
                    # Excel XLSX - NUOVO
                    elif ext == '.xlsx':
                        try:
                            import openpyxl
                            wb = openpyxl.load_workbook(temp_file_path, read_only=True, data_only=True)
                            xlsx_texts = []
                            
                            for sheet_name in wb.sheetnames:
                                sheet = wb[sheet_name]
                                xlsx_texts.append(f"--- Foglio: {sheet_name} ---")
                                
                                # Max 500 righe e 50 colonne per evitare problemi con file molto grandi
                                max_row = min(sheet.max_row, 500) if sheet.max_row else 0
                                max_col = min(sheet.max_column, 50) if sheet.max_column else 0
                                
                                for row_idx in range(1, max_row + 1):
                                    row_values = []
                                    for col_idx in range(1, max_col + 1):
                                        cell = sheet.cell(row=row_idx, column=col_idx)
                                        if cell.value:
                                            row_values.append(str(cell.value))
                                    if row_values:
                                        xlsx_texts.append(" ".join(row_values))
                                        
                            content = "\n".join(xlsx_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da XLSX")
                        except ImportError:
                            self.log_debug("Libreria openpyxl non disponibile")
                        except Exception as xlsx_err:
                            self.log_debug(f"Errore XLSX: {str(xlsx_err)}")
                    
                    # Excel XLS (vecchio formato) - NUOVO
                    elif ext == '.xls':
                        try:
                            import xlrd
                            book = xlrd.open_workbook(temp_file_path)
                            xls_texts = []
                            
                            for sheet_idx in range(book.nsheets):
                                sheet = book.sheet_by_index(sheet_idx)
                                xls_texts.append(f"--- Foglio: {sheet.name} ---")
                                
                                # Limita il numero di righe per prestazioni
                                for row_idx in range(min(sheet.nrows, 500)):
                                    row_values = sheet.row_values(row_idx)
                                    row_texts = [str(value) for value in row_values if value]
                                    if row_texts:
                                        xls_texts.append(" ".join(row_texts))
                                        
                            content = "\n".join(xls_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da XLS")
                        except ImportError:
                            self.log_debug("Libreria xlrd non disponibile")
                        except Exception as xls_err:
                            self.log_debug(f"Errore XLS: {str(xls_err)}")
                    
                    # CSV - NUOVO
                    elif ext == '.csv':
                        try:
                            import csv
                            csv_texts = []
                            
                            with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                reader = csv.reader(f)
                                # Limita a 1000 righe per file grandi
                                for i, row in enumerate(reader):
                                    if i >= 1000:
                                        csv_texts.append("... (file troncato, troppe righe)")
                                        break
                                    csv_texts.append("\t".join(row))
                                    
                            content = "\n".join(csv_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da CSV")
                        except Exception as csv_err:
                            self.log_debug(f"Errore CSV: {str(csv_err)}")
                    
                    # OpenDocument Text (ODT) - NUOVO
                    elif ext == '.odt':
                        try:
                            from odf import opendocument, text
                            
                            doc = opendocument.load(temp_file_path)
                            paragraphs = []
                            
                            # Estrai tutto il testo dai paragrafi
                            for element in doc.getElementsByType(text.P):
                                if element.firstChild:
                                    paragraphs.append(element.firstChild.data)
                                else:
                                    paragraphs.append("")
                                    
                            content = "\n".join(paragraphs)
                            self.log_debug(f"Estratti {len(content)} caratteri da ODT")
                        except ImportError:
                            self.log_debug("Libreria odf non disponibile")
                        except Exception as odt_err:
                            self.log_debug(f"Errore ODT: {str(odt_err)}")
                    
                    # Rich Text Format (RTF) - NUOVO
                    elif ext == '.rtf':
                        try:
                            from striprtf.striprtf import rtf_to_text
                            
                            with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                rtf_content = f.read()
                                content = rtf_to_text(rtf_content)
                                self.log_debug(f"Estratti {len(content)} caratteri da RTF")
                        except ImportError:
                            self.log_debug("Libreria striprtf non disponibile")
                        except Exception as rtf_err:
                            self.log_debug(f"Errore RTF: {str(rtf_err)}")
                    
                    # XML - NUOVO
                    elif ext == '.xml':
                        try:
                            import xml.etree.ElementTree as ET
                            
                            tree = ET.parse(temp_file_path)
                            root = tree.getroot()
                            
                            # Funzione ricorsiva per estrarre tutti i testi
                            def extract_text_from_element(element):
                                result = []
                                if element.text and element.text.strip():
                                    result.append(element.text.strip())
                                for child in element:
                                    result.extend(extract_text_from_element(child))
                                if element.tail and element.tail.strip():
                                    result.append(element.tail.strip())
                                return result
                            
                            xml_texts = extract_text_from_element(root)
                            content = "\n".join(xml_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da XML")
                        except Exception as xml_err:
                            self.log_debug(f"Errore XML: {str(xml_err)}")
                            # Fallback a lettura semplice del file
                            try:
                                with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                    content = f.read()
                                    self.log_debug(f"XML letto come testo semplice: {len(content)} caratteri")
                            except:
                                pass
                    
                    # Se è un altro tipo di file, usa get_file_content
                    else:
                        try:
                            content = self.get_file_content(temp_file_path)
                            self.log_debug(f"Contenuto estratto via get_file_content: {len(content)} caratteri")
                        except Exception as get_file_err:
                            self.log_debug(f"Errore get_file_content: {str(get_file_err)}")
                    
                    if content and len(content) > 0:
                        self.log_debug(f"Estratti con successo {len(content)} caratteri dall'allegato {safe_name}")
                        return content
                    else:
                        self.log_debug(f"Nessun contenuto estratto dall'allegato {safe_name}")
                        return f"[Allegato {safe_name}: nessun contenuto estratto]"
                except Exception as e:
                    self.log_debug(f"Errore nell'estrazione del contenuto dell'allegato {safe_name}: {str(e)}")
                    return f"[Errore nell'estrazione: {str(e)}]"
            else:
                self.log_debug(f"Allegato {safe_name} non elaborabile secondo le impostazioni attuali")
                return f"[Allegato {safe_name} ({ext}) non elaborabile]"
            
        except Exception as e:
            self.log_debug(f"Errore generale nell'elaborazione dell'allegato {attachment_name}: {str(e)}")
            return ""
        finally:
            # Clean up temp directory
            if temp_dir and os.path.exists(temp_dir):
                try:
                    import shutil
                    shutil.rmtree(temp_dir)
                    self.log_debug(f"Directory temporanea rimossa: {temp_dir}")
                except Exception as e:
                    self.log_debug(f"Errore nella pulizia dei file temporanei: {str(e)}")

    @error_handler
    def update_progress(self):
        if self.is_searching:
            try:
                # Process only a limited number of messages per cycle
                max_messages_per_cycle = 15  # Reduced from 20 to 15 for better responsiveness
                messages_processed = 0
                
                start_time = time.time()
                # Prevent processing for too long
                max_processing_time = 0.03  # Reduced from 0.05 to 0.03 seconds for more responsive UI
                
                while messages_processed < max_messages_per_cycle:
                    try:
                        # Processa tutti i messaggi nella coda
                        progress_type, value = self.progress_queue.get_nowait()
                        
                        # CORREZIONE: Ottimizza l'elaborazione dei messaggi di stato
                        if progress_type == "update_total_time":
                            if hasattr(self, 'total_time_label') and self.total_time_label.winfo_exists():
                                self.total_time_label.config(text=value)
                        elif progress_type == "progress":
                            if hasattr(self, 'progress_bar') and self.progress_bar.winfo_exists():
                                current_value = self.progress_bar["value"]
                                # Evita aggiornamenti inutili se il valore non è cambiato significativamente
                                if abs(current_value - value) >= 1:
                                    self.progress_bar["value"] = value
                        elif progress_type == "status":
                            # CORREZIONE: Migliora la gestione dei messaggi di stato
                            if ("analizzati" in value.lower() or "cartelle:" in value.lower() or 
                                "file:" in value.lower()) and hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                # Estrai il percorso se presente
                                if "analisi:" in value.lower():
                                    parts = value.split("(", 1)
                                    if len(parts) > 1:
                                        path_part = parts[0].strip()
                                        counts_part = "(" + parts[1]
                                        if hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                            self.status_label["text"] = path_part
                                        if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                            self.analyzed_files_label["text"] = counts_part
                                    else:
                                        if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                            self.analyzed_files_label["text"] = value
                                else:
                                    # È solo un messaggio di stato semplice
                                    if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                        self.analyzed_files_label["text"] = value
                            elif hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                # CORREZIONE: Evita aggiornamenti ripetitivi dello stesso messaggio
                                if not hasattr(self, 'last_status_message') or self.last_status_message != value:
                                    self.status_label["text"] = value
                                    self.last_status_message = value
                        elif progress_type == "update_dir_size":
                            # CORREZIONE: Aggiorna solo se abilitato e necessario
                            calculation_mode = self.dir_size_calculation.get()
                            if hasattr(self, 'dir_size_var') and calculation_mode != "disabilitato":
                                current_text = self.dir_size_var.get()
                                new_text = self._format_size(value)
                                # Aggiorna solo se il valore è cambiato significativamente
                                if current_text != new_text:
                                    self.dir_size_var.set(new_text)
                            elif hasattr(self, 'dir_size_var') and calculation_mode == "disabilitato":
                                self.dir_size_var.set("Calcolo disattivato")
                        elif progress_type == "complete":
                            self.is_searching = False
                            self.enable_all_controls()
                            if hasattr(self, 'stop_button') and self.stop_button.winfo_exists():
                                self.stop_button["state"] = "disabled"
                            
                            # CORREZIONE: Aggiungi un breve ritardo prima di aggiornare i risultati
                            # per garantire che l'interfaccia sia reattiva
                            self.root.after(100, self.update_results_list)
                            
                            # Aggiorna il tempo finale
                            if hasattr(self, 'search_start_time') and self.search_start_time:
                                # Imposta il timestamp di fine ricerca
                                self.search_end_time = datetime.now()
                                current_time = self.search_end_time.strftime('%H:%M')
                                
                                if hasattr(self, 'end_time_label') and self.end_time_label.winfo_exists():
                                    self.end_time_label.config(text=current_time)
                                    
                                # Aggiorna il tempo totale
                                self.update_total_time()
                                
                                # Aggiorna il debug log
                                self.log_debug(f"Ricerca completata. Trovati {len(self.search_results)} risultati")
                        
                        messages_processed += 1
                        
                        # Check if we've been processing for too long
                        if time.time() - start_time > max_processing_time:
                            break
                            
                    except queue.Empty:
                        break
                    except tk.TclError as e:
                        # Widget non più esistente o distrutto - non grave, ignora e continua
                        self.log_debug(f"Widget non più disponibile: {str(e)}")
                        continue
                    except Exception as e:
                        # Log altri errori ma non interrompere l'aggiornamento
                        self.log_debug(f"Errore durante l'elaborazione messaggio: {str(e)}")
                        continue
                
                # Force UI update after processing messages
                try:
                    self.root.update_idletasks()
                except:
                    pass
                
                # CORREZIONE: Adatta meglio la frequenza di aggiornamento
                if hasattr(self.progress_queue, 'qsize'):
                    queue_size = self.progress_queue.qsize()
                    if queue_size > 100:
                        # Many messages in queue, update more frequently
                        self.root.after(20, self.update_progress)  # 20ms (50Hz)
                    elif queue_size > 20:
                        # Moderate queue, moderate update rate
                        self.root.after(50, self.update_progress)  # 50ms (20Hz)
                    else:
                        # Few messages, standard update rate
                        self.root.after(100, self.update_progress)  # 100ms (10Hz)
                else:
                    # No qsize method, use standard rate
                    self.root.after(100, self.update_progress)
                    
            except tk.TclError as e:
                self.log_debug(f"TclError nell'aggiornamento del progresso: {str(e)}")
                self.root.after(500, self.update_progress)
            except Exception as e:
                self.log_debug(f"Errore nell'aggiornamento del progresso: {str(e)}")
                self.root.after(500, self.update_progress)
    
    @error_handler
    def reset_search_state(self):
        """Reimpostazione completa dello stato di ricerca per garantire coerenza"""
        self.log_debug("Reimpostazione dello stato di ricerca")
        
        # Variabili di controllo principali
        self.is_searching = False
        self.stop_search = False
        
        # Elimina flag temporanei di interruzione se presenti
        if hasattr(self, '_stopping_in_progress'):
            delattr(self, '_stopping_in_progress')
        
        # Chiudi executor se ancora esistente
        if hasattr(self, 'search_executor') and self.search_executor:
            try:
                self.search_executor.shutdown(wait=False)
            except:
                pass
            self.search_executor = None
        
        # Ferma il watchdog
        self.watchdog_active = False
        
        # Svuota le code
        try:
            while not self.progress_queue.empty():
                try:
                    self.progress_queue.get_nowait()
                    self.progress_queue.task_done()
                except:
                    break
        except:
            pass
        
        # Aggiorna UI per riflettere lo stato corretto
        if hasattr(self, 'stop_button'):
            self.stop_button["state"] = "disabled"
        
        self.log_debug("Stato di ricerca reimpostato correttamente")

    @error_handler
    def stop_search_process(self):
        """Ferma il processo di ricerca in corso in modo aggressivo"""
        # Verificare se una ricerca è effettivamente in corso
        if not self.is_searching:
            self.log_debug("Tentativo di interruzione ma nessuna ricerca in corso")
            return
        
        self.log_debug("INTERRUZIONE: Stop ricerca richiesto dall'utente")
        
        # 1. Imposta i flag di interruzione
        self.stop_search = True
        self.is_searching = False
        
        # 2. Aggiorna l'interfaccia
        self.status_label["text"] = "Interruzione ricerca in corso... attendere"
        self.analyzed_files_label["text"] = "Ricerca interrotta dall'utente"
        self.progress_bar["value"] = 100
        
        # 3. Cattura il tempo di fine
        self.search_end_time = datetime.now()
        current_time = self.search_end_time.strftime('%H:%M')
        self.end_time_label.config(text=current_time)
        self.update_total_time()  # Aggiorna il tempo totale
        
        # 4. CRITICO: Ferma tutti i thread in modo aggressivo
        if hasattr(self, 'search_executor') and self.search_executor:
            try:
                self.log_debug("INTERRUZIONE: Chiusura dell'executor...")
                
                # Ferma l'executor nel modo più aggressivo possibile
                try:
                    import sys
                    if sys.version_info >= (3, 9):
                        self.search_executor.shutdown(wait=False, cancel_futures=True)
                    else:
                        self.search_executor.shutdown(wait=False)
                        
                    # Crea un nuovo executor per le future ricerche
                    self.search_executor = None
                    
                except Exception as e:
                    self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                    self.search_executor = None  # Forza nullificazione
                
            except Exception as e:
                self.log_debug(f"Errore generale nell'interruzione: {str(e)}")
                self.search_executor = None
        
        # 5. NUOVO: Visualizza risultati parziali trovati
        self.root.after(500, self.update_results_list)
        
        # 6. Riabilita l'interfaccia utente dopo l'interruzione
        self.root.after(1000, self._complete_interrupt_process)
        
        # 7. Forza aggiornamento GUI
        try:
            self.root.update_idletasks()
        except:
            pass
        
        self.log_debug("INTERRUZIONE: Processo di interruzione ricerca completato")
    
    @error_handler
    def _complete_interrupt_process(self):
        """Completa il processo di interruzione ripristinando l'interfaccia"""
        try:
            # Riabilita tutti i controlli
            self.enable_all_controls()
            
            # Reimpostazione completa dello stato
            self.reset_search_state()
            
            self.status_label["text"] = "Ricerca interrotta dall'utente"
            self.log_debug("Interfaccia ripristinata dopo interruzione")
        except Exception as e:
            self.log_debug(f"Errore nel completamento dell'interruzione: {str(e)}")

    @error_handler
    def update_results_list(self):
        """Aggiorna la lista dei risultati con i risultati trovati"""
        # Pulisci la lista attuale
        for item in self.results_list.get_children():
            self.results_list.delete(item)
        
        # Aggiungi i risultati alla lista
        for result in self.search_results:
            # Gestisci il nuovo formato del risultato con 7 elementi (incluso from_attachment)
            if len(result) >= 7:
                item_type, author, size, modified, created, path, from_attachment = result
            else:
                item_type, author, size, modified, created, path = result
                from_attachment = False
            
            # Imposta l'icona della graffetta nella colonna dedicata
            attachment_icon = "📎" if from_attachment else ""
            
            # Applica stile in base al tipo di elemento
            if item_type == "Directory":
                tags = ("directory",)
            elif from_attachment:
                tags = ("attachment",)  # Tag speciale per gli allegati
            else:
                tags = ("file",)
            
            display_values = (
                item_type,        # Tipo
                attachment_icon,  # Icona allegato
                size,             # Dimensione
                modified,         # Data modifica
                created,          # Data creazione
                author,           # Nome/Autore
                path              # Percorso
            )
                
            # Inserisci l'elemento nella TreeView con i tag appropriati
            self.results_list.insert("", "end", values=display_values, tags=tags)
        
        # Aggiorna lo stato
        self.status_label["text"] = f"Trovati {len(self.search_results)} risultati"
        
        # Aggiorna la dimensione totale dei file trovati
        self.update_total_files_size()
        
        # Se ci sono risultati, seleziona il primo elemento
        if self.search_results:
            first_item = self.results_list.get_children()[0]
            self.results_list.selection_set(first_item)
            self.results_list.focus(first_item)
            
            # Opzionalmente, mostra il contenuto del primo risultato
            if hasattr(self, 'auto_preview') and self.auto_preview.get():
                self.show_file_contents()
    
    @error_handler
    def update_total_files_size(self):
        """Calcola e aggiorna la dimensione totale dei file trovati"""
        try:
            total_size = 0
            file_count = 0
            
            # Output debug info
            self.log_debug(f"Calculating total size for {len(self.search_results)} results")
            
            # Calcola la dimensione totale dai risultati
            for result in self.search_results:
                try:
                    # Handle both 6-element and 7-element results
                    if len(result) >= 7:  # If it includes from_attachment flag
                        item_type, _, size_str, _, _, _, _ = result
                    else:
                        item_type, _, size_str, _, _, _ = result
                    
                    # Conta solo i file, non le directory
                    if item_type != "Directory":
                        file_count += 1
                        
                        # Estrai il valore numerico dalla stringa della dimensione
                        if size_str and isinstance(size_str, str):
                            if 'KB' in size_str:
                                size_value = float(size_str.split()[0]) * 1024
                            elif 'MB' in size_str:
                                size_value = float(size_str.split()[0]) * 1024 * 1024
                            elif 'GB' in size_str:
                                size_value = float(size_str.split()[0]) * 1024 * 1024 * 1024
                            else:
                                # Assume B or other unit
                                size_value = float(size_str.split()[0])
                            
                            total_size += size_value
                except Exception as e:
                    self.log_debug(f"Error processing file size: {str(e)}")
                    continue
            
            # Formatta la dimensione totale
            formatted_size = self._format_size(total_size)
            
            # Aggiorna il label con la dimensione totale e il numero di file
            if hasattr(self, 'total_files_size_label') and self.total_files_size_label.winfo_exists():
                new_text = f"Dimensione totale: {formatted_size} ({file_count} file)"
                self.log_debug(f"Updating size label to: {new_text}")
                self.total_files_size_label.config(text=new_text)
                
                # Forza aggiornamento dell'interfaccia
                self.root.update_idletasks()
        except Exception as e:
            self.log_debug(f"Error in update_total_files_size: {str(e)}")

    @error_handler
    def update_selected_files_size(self, event=None):
        """Calcola e visualizza la dimensione totale dei file selezionati"""
        selected_items = self.results_list.selection()
        
        total_size = 0
        file_count = 0
        
        for item in selected_items:
            values = self.results_list.item(item)['values']
            # Verifica che sia un file (non una directory)
            if values and values[0] != "Directory":
                file_count += 1
                size_str = values[2]  # La dimensione è nella terza colonna (indice 2)
                
                # Estrai il valore numerico dalla stringa della dimensione
                if size_str and isinstance(size_str, str):
                    try:
                        if 'KB' in size_str:
                            size_value = float(size_str.split()[0]) * 1024
                        elif 'MB' in size_str:
                            size_value = float(size_str.split()[0]) * 1024 * 1024
                        elif 'GB' in size_str:
                            size_value = float(size_str.split()[0]) * 1024 * 1024 * 1024
                        else:
                            # Assume B o altra unità
                            size_value = float(size_str.split()[0])
                        
                        total_size += size_value
                    except:
                        pass  # Ignora errori di parsing
        
        # Formatta la dimensione e aggiorna l'etichetta
        if file_count > 0:
            formatted_size = self._format_size(total_size)
            self.selected_files_size_label.config(text=f"Selezionati: {formatted_size} ({file_count} file)")
        else:
            self.selected_files_size_label.config(text="Selezionati: 0 (0 file)")

    @error_handler
    def update_theme_colors(self, theme="light"):
        """Aggiorna i colori del tema per evidenziare cartelle e file"""
        style = ttk.Style()
    
        # Configura i colori per il tema chiaro
        if theme == "light":
            # Codice esistente per configurare il tema chiaro
            style.configure("Treeview", background="#ffffff", foreground="#000000", fieldbackground="#ffffff")
            self.results_list.tag_configure("directory", background="#e6f2ff", foreground="#000000")
            self.results_list.tag_configure("file", background="#ffffff", foreground="#000000")
            # AGGIUNGI QUI: Nuovo stile per gli allegati
            self.results_list.tag_configure("attachment", background="#f8f8e0", foreground="#000000")
        # Configura i colori per il tema scuro
        elif theme == "dark":
            # Codice esistente per configurare il tema scuro
            style.configure("Treeview", background="#333333", foreground="#ffffff", fieldbackground="#333333")
            self.results_list.tag_configure("directory", background="#4d4d4d", foreground="#ffffff")
            self.results_list.tag_configure("file", background="#333333", foreground="#ffffff")
            # AGGIUNGI QUI: Nuovo stile per gli allegati
            self.results_list.tag_configure("attachment", background="#464630", foreground="#ffffff")

    @error_handler
    def show_block_options(self):
        """Mostra la finestra di dialogo per le opzioni avanzate di ricerca a blocchi"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Opzioni Avanzate Ricerca a Blocchi")
        dialog.geometry("560x600")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Descrizione
        ttk.Label(main_frame, text="Configura le opzioni per ottimizzare la ricerca a blocchi", 
             font=("", 10, "bold")).pack(anchor=W, pady=(0, 10))
        
        description = ttk.Label(main_frame, text="La ricerca a blocchi divide la scansione in blocchi più piccoli "
                                        "per migliorare le prestazioni e la reattività dell'interfaccia.",
                         wraplength=520, justify=LEFT)
        description.pack(fill=X, pady=(0, 10))
        
        # Frame per opzioni dei blocchi
        block_frame = ttk.LabelFrame(main_frame, text="Dimensione e Parallelismo", padding=10)
        block_frame.pack(fill=X, pady=5)
        
        # Grid per organizzare le opzioni
        grid = ttk.Frame(block_frame)
        grid.pack(fill=X)
        
        # Opzione: Max file per blocco
        ttk.Label(grid, text="Max file per blocco:").grid(row=0, column=0, sticky=W, padx=5, pady=5)
        files_per_block = ttk.Spinbox(grid, from_=100, to=10000, increment=100, width=7,
                                 textvariable=self.max_files_per_block)
        files_per_block.grid(row=0, column=1, padx=5, pady=5)
        
        # Opzione: Max blocchi paralleli
        ttk.Label(grid, text="Blocchi paralleli:").grid(row=0, column=2, sticky=W, padx=(20, 5), pady=5)
        parallel_blocks = ttk.Spinbox(grid, from_=1, to=16, increment=1, width=5,
                                 textvariable=self.max_parallel_blocks)
        parallel_blocks.grid(row=0, column=3, padx=5, pady=5)
        
        # Checkbox per opzioni aggiuntive
        auto_adjust = ttk.Checkbutton(grid, text="Adatta automaticamente la dimensione dei blocchi",
                                variable=self.block_size_auto_adjust)
        auto_adjust.grid(row=1, column=0, columnspan=4, sticky=W, padx=5, pady=5)
        self.create_tooltip(auto_adjust, "Regola automaticamente la dimensione dei blocchi in base al tipo di ricerca e al sistema")
        
        # Frame per opzioni di prioritizzazione
        priority_frame = ttk.LabelFrame(main_frame, text="Priorità di Ricerca", padding=10)
        priority_frame.pack(fill=X, pady=10)
        
        # Checkbox per dare priorità alle cartelle utente
        user_priority = ttk.Checkbutton(priority_frame, text="Dare priorità alle cartelle utente (Users, Documents, Desktop, ecc)",
                                variable=self.prioritize_user_folders)
        user_priority.pack(anchor=W, padx=5, pady=5)
        self.create_tooltip(user_priority, "Elabora prima le cartelle più importanti per trovare risultati utili più velocemente")
        
        # Frame per spiegazione avanzata
        info_frame = ttk.LabelFrame(main_frame, text="Come funziona", padding=10)
        info_frame.pack(fill=X, pady=10)
        
        info_text = "La ricerca a blocchi divide le cartelle in unità di lavoro più piccole (blocchi) " \
                "che vengono elaborate in base alla priorità. Questo permette di:\n\n" \
                "• Ottenere risultati utili più rapidamente\n" \
                "• Migliorare la reattività dell'interfaccia durante la ricerca\n" \
                "• Distribuire meglio il carico di lavoro sui thread\n" \
                "• Gestire meglio la memoria per ricerche su grandi volumi di dati"
        
        info_label = ttk.Label(info_frame, text=info_text, wraplength=520, justify=LEFT)
        info_label.pack(fill=X, padx=5, pady=5)
        
        # Pulsanti per chiudere/applicare
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        # Pulsante per ripristinare valori predefiniti
        ttk.Button(btn_frame, text="Valori predefiniti", 
              command=lambda: [self.max_files_per_block.set(1000), 
                              self.max_parallel_blocks.set(4),
                              self.prioritize_user_folders.set(True),
                              self.block_size_auto_adjust.set(True)]).pack(side=LEFT, padx=5)
        
        # Pulsante di chiusura
        ttk.Button(btn_frame, text="Chiudi", command=dialog.destroy).pack(side=RIGHT, padx=5)
        
        # Centra la finestra
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")

    @error_handler
    def copy_selected(self):
        """Copia i file selezionati in una directory di destinazione"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da copiare")
            return
            
        # Chiedi la directory di destinazione
        dest_dir = filedialog.askdirectory(title="Seleziona la cartella di destinazione")
        if not dest_dir:
            return
            
        # Prepara variabili per tracciare l'avanzamento
        total = len(selected_items)
        copied = 0
        failed = 0
        skipped = 0
        
        try:
            for item in selected_items:
                values = self.results_list.item(item)['values']
                item_type, _, _, _, _, source_path = values
                
                # Ottieni il nome dell'elemento senza il percorso completo
                basename = os.path.basename(source_path)
                dest_path = os.path.join(dest_dir, basename)
                
                try:
                    if item_type == "Directory":
                        # Se la cartella esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("Cartella esistente", 
                                                f"La cartella {basename} esiste già. Vuoi sovrascriverla?"):
                                # Elimina la cartella esistente
                                shutil.rmtree(dest_path)
                            else:
                                skipped += 1
                                continue
                                
                        # Copia ricorsiva della cartella
                        shutil.copytree(source_path, dest_path)
                        copied += 1
                    else:
                        # Se il file esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("File esistente", 
                                                f"Il file {basename} esiste già. Vuoi sovrascriverlo?"):
                                # Continua con la sovrascrittura
                                pass
                            else:
                                skipped += 1
                                continue
                                
                        # Copia il file
                        shutil.copy2(source_path, dest_path)
                        copied += 1
                        
                except Exception as e:
                    failed += 1
                    self.log_debug(f"Errore durante la copia di {source_path}: {str(e)}")
                    
                # Aggiorna la barra di avanzamento
                progress = ((copied + failed + skipped) / total) * 100
                self.progress_bar["value"] = progress
                self.status_label["text"] = f"Copiati {copied}/{total} elementi..."
                self.root.update()
                
            # Messaggio di completamento
            if failed > 0:
                messagebox.showwarning("Copia completata con errori", 
                                    f"Copiati: {copied}\nSaltati: {skipped}\nFalliti: {failed}")
            else:
                messagebox.showinfo("Copia completata", 
                                 f"Copiati con successo: {copied}\nSaltati: {skipped}")
                
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'operazione: {str(e)}")
            
        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    @error_handler
    def get_zip_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome del file ZIP"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Nome file ZIP")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome del file ZIP (senza estensione):", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = StringVar(value="archivio")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            zip_name = name_var.get().strip()
            if zip_name:
                result["name"] = zip_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per il file ZIP", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea ZIP", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]

    @error_handler
    def compress_selected(self):
        """Versione avanzata della compressione che mantiene la struttura originale delle directory"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da comprimere")
            return

        # Ottieni il nome del file ZIP
        zip_name = self.get_zip_name()
        if not zip_name:
            return

        # Ottieni il nome della cartella principale
        main_folder_name = self.get_main_folder_name()
        if main_folder_name is None:
            return  # L'utente ha annullato
        if not main_folder_name:
            main_folder_name = "files"  # Default se l'utente inserisce valore vuoto

        # Opzioni di compressione più dettagliate
        compression_dialog = tk.Toplevel(self.root)
        compression_dialog.title("Opzioni Compressione")
        compression_dialog.transient(self.root)
        compression_dialog.grab_set()

        frame = ttk.Frame(compression_dialog, padding=20)
        frame.pack(fill=tk.BOTH, expand=True)

        ttk.Label(frame, text="Scegli il tipo di compressione:").pack(anchor=tk.W, pady=(0,10))

        comp_var = tk.StringVar(value="nessuna")
        preserve_var = tk.BooleanVar(value=True)  # Nuova opzione per preservare la struttura

        ttk.Radiobutton(frame, text="Nessuna (solo archiviazione)", 
            variable=comp_var, value="nessuna").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Minima (massima velocità) Liv.1", 
                    variable=comp_var, value="minima").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Standard (buon equilibrio) Liv.6", 
                    variable=comp_var, value="standard").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Massima (compressione ottimale, più lenta) Liv.9", 
                    variable=comp_var, value="massima").pack(anchor=tk.W)

        # Frame per le opzioni di hash
        hash_frame = ttk.LabelFrame(frame, text="Calcolo hash del file ZIP")
        hash_frame.pack(fill=tk.X, pady=(10,5))

        ttk.Label(hash_frame, text="Seleziona gli algoritmi di hash da calcolare sul file compresso:").pack(anchor=tk.W, pady=(5,5))

        # Variabili per gli algoritmi di hash
        hash_md5 = tk.BooleanVar(value=True)
        hash_sha1 = tk.BooleanVar(value=False)
        hash_sha256 = tk.BooleanVar(value=False)

        # Checkbox per gli algoritmi
        ttk.Checkbutton(hash_frame, text="MD5 (più veloce)", variable=hash_md5).pack(anchor=tk.W, padx=20)
        ttk.Checkbutton(hash_frame, text="SHA1", variable=hash_sha1).pack(anchor=tk.W, padx=20)
        ttk.Checkbutton(hash_frame, text="SHA256 (più sicuro)", variable=hash_sha256).pack(anchor=tk.W, padx=20)

        ttk.Label(frame, text="\nOpzioni struttura:", font=("", 9, "bold")).pack(anchor=tk.W, pady=(10,5))

        ttk.Checkbutton(frame, text="Preserva la struttura delle directory", 
                    variable=preserve_var).pack(anchor=tk.W)
        ttk.Label(frame, text="Se attivato, mantiene l'alberatura originale delle cartelle", 
                    font=("", 8), foreground="gray").pack(anchor=tk.W, padx=(20, 5))

        ttk.Label(frame, text="\nOpzioni per file di grandi dimensioni:").pack(anchor=tk.W, pady=(10,5))

        use_chunks = tk.BooleanVar(value=False)
        ttk.Checkbutton(frame, text="Elabora in blocchi (per file molto grandi)", 
                    variable=use_chunks).pack(anchor=tk.W)

        btn_frame = ttk.Frame(frame)
        btn_frame.pack(fill=tk.X, pady=(15,0))

        result = {"option": None, "chunks": False, "preserve": True}

        def on_cancel():
            result["option"] = None
            compression_dialog.destroy()

        def on_ok():
            result["option"] = comp_var.get()
            result["chunks"] = use_chunks.get()
            result["preserve"] = preserve_var.get()
            
            # Aggiungi gli algoritmi selezionati
            result["hash_algorithms"] = []
            if hash_md5.get():
                result["hash_algorithms"].append("md5")
            if hash_sha1.get():
                result["hash_algorithms"].append("sha1")
            if hash_sha256.get():
                result["hash_algorithms"].append("sha256")
                
            compression_dialog.destroy()

        ttk.Button(btn_frame, text="Annulla", command=on_cancel).pack(side=tk.LEFT)
        ttk.Button(btn_frame, text="OK", command=on_ok).pack(side=tk.RIGHT)

        # Centra la finestra
        compression_dialog.update_idletasks()
        width = compression_dialog.winfo_reqwidth() + 50
        height = compression_dialog.winfo_reqheight() + 20
        x = (compression_dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (compression_dialog.winfo_screenheight() // 2) - (height // 2)
        compression_dialog.geometry(f"{width}x{height}+{x}+{y}")

        compression_dialog.wait_window()

        if result["option"] is None:
            return  # Utente ha annullato

        # Determina il metodo di compressione
        compression_level = 0  # Valore predefinito

        if result["option"] == "nessuna":
            compression_method = zipfile.ZIP_STORED
            compression_text = "nessuna"
        elif result["option"] == "minima":
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "minima"
            compression_level = 1  # Compressione veloce, riduzione minima
        elif result["option"] == "massima":
            # Prova LZMA se disponibile, altrimenti BZIP2, altrimenti fallback su DEFLATED con livello massimo
            if hasattr(zipfile, 'ZIP_LZMA'):
                compression_method = zipfile.ZIP_LZMA
                compression_text = "massima (LZMA)"
                compression_level = 9
            elif hasattr(zipfile, 'ZIP_BZIP2'):
                compression_method = zipfile.ZIP_BZIP2
                compression_text = "massima (BZIP2)"
                compression_level = 9
            else:
                compression_method = zipfile.ZIP_DEFLATED
                compression_text = "massima (DEFLATED)"
                compression_level = 9
        else:  # standard
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "standard"
            compression_level = 6  # Compressione bilanciata
            
        # Log della scelta di compressione
        self.log_debug(f"Utilizzo compressione {compression_text} (metodo: {compression_method}, livello: {compression_level})")
        self.log_debug(f"Preserva struttura directory: {result['preserve']}")

        zip_path = filedialog.asksaveasfilename(
            defaultextension=".zip",
            initialfile=f"{zip_name}.zip",
            filetypes=[("ZIP files", "*.zip")],
            title="Salva file ZIP"
        )

        if not zip_path:
            return

        # Determinare il percorso base per i calcoli relativi
        base_path = self._find_common_base_path(selected_items)
        self.log_debug(f"Percorso base per la struttura: {base_path}")

        # Raccogli tutti i file delle cartelle selezionate
        files_in_folders = set()
        folder_paths = []
        single_files = []

        # Prima fase: raccogli informazioni su cartelle e file
        for item in selected_items:
            values = self.results_list.item(item)['values']
            item_type, _, _, _, _, _, source_path = values

            if item_type == "Directory":
                folder_paths.append(source_path)
                # Raccogli tutti i file nelle cartelle selezionate
                for root, _, files in os.walk(source_path):
                    for file in files:
                        full_path = os.path.join(root, file)
                        files_in_folders.add(os.path.abspath(full_path))
            else:
                single_files.append(source_path)

        # Filtra i file singoli che sono già presenti nelle cartelle
        filtered_single_files = [f for f in single_files if os.path.abspath(f) not in files_in_folders]

        total_items = len(folder_paths) + len(filtered_single_files)
        processed = 0

        # Lista per tenere traccia dei file (tutti, non solo omonimi)
        all_files_log = []
        # Lista per tenere traccia dei file omonimi (come prima)
        omonimi_log = []

        # Ottieni informazioni sulla ricerca corrente per i log
        search_directory = self.search_path.get() if hasattr(self, 'search_path') else "N/A"
        search_keywords = self.keywords.get() if hasattr(self, 'keywords') else "N/A"

        try:
            # Primo passaggio: analizza tutti i file e identifica gli omonimi
            self.log_debug(f"Analisi dei file da comprimere nella cartella {main_folder_name}...")

            # Dizionario {nome_file: [lista_percorsi]}
            file_names_map = {}

            # Raccogli i nomi dei file dalle cartelle
            for folder_path in folder_paths:
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        file_name = os.path.basename(file_path)

                        # Aggiungi al log di tutti i file
                        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                        modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                        
                        all_files_log.append({
                            "nome_file": file_name,
                            "percorso_originale": file_path,
                            "dimensione": self._format_size(file_size),
                            "ultima_modifica": modified_time,
                            "tipo": self._get_file_type(file_path)
                        })

                        if file_name not in file_names_map:
                            file_names_map[file_name] = []
                        file_names_map[file_name].append(file_path)

            # Raccogli i nomi dei file singoli
            for file_path in filtered_single_files:
                if os.path.exists(file_path):
                    file_name = os.path.basename(file_path)
                    
                    # Aggiungi al log di tutti i file
                    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                    modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                    
                    all_files_log.append({
                        "nome_file": file_name,
                        "percorso_originale": file_path,
                        "dimensione": self._format_size(file_size),
                        "ultima_modifica": modified_time,
                        "tipo": self._get_file_type(file_path)
                    })
                    
                    if file_name not in file_names_map:
                        file_names_map[file_name] = []
                    file_names_map[file_name].append(file_path)

            # Identifica i file omonimi (con lo stesso nome ma percorsi diversi)
            omonimi_files = {name: paths for name, paths in file_names_map.items() if len(paths) > 1}
            self.log_debug(f"Trovati {len(omonimi_files)} file con nomi duplicati")
            
            # Per l'elaborazione a blocchi
            if result["chunks"] and total_items > 50:
                chunk_size = 10  # Elabora 10 file alla volta
                self.log_debug(f"Elaborazione a blocchi attivata: {chunk_size} file per blocco")

            # Secondo passaggio: crea l'archivio ZIP
            with zipfile.ZipFile(zip_path, 'w', compression_method, compresslevel=compression_level) as zipf:
                # Tiene traccia dei percorsi già aggiunti al file ZIP
                added_zip_paths = set()

                # Comprimi le cartelle
                for folder_path in folder_paths:
                    # Utilizzo la struttura originale o piatta in base alla scelta dell'utente
                    for root, _, files in os.walk(folder_path):
                        for file in files:
                            file_path = os.path.join(root, file)
                            file_name = os.path.basename(file_path)

                            try:
                                # Salta file non leggibili
                                if not os.access(file_path, os.R_OK):
                                    self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                    continue

                                # Calcola percorso relativo in base all'opzione selezionata
                                if result["preserve"]:
                                    # Calcola il percorso relativo rispetto alla base
                                    rel_path = os.path.relpath(file_path, base_path)
                                    zip_path_to_use = os.path.join(main_folder_name, rel_path)
                                else:
                                    # Comportamento originale senza preservare la struttura
                                    rel_path = os.path.relpath(file_path, os.path.dirname(folder_path))
                                    
                                    # Verifica se è un file omonimo
                                    if file_name in omonimi_files:
                                        # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                        if omonimi_files[file_name][0] == file_path:
                                            # Usa il percorso all'interno della cartella principale
                                            zip_path_to_use = os.path.join(main_folder_name, rel_path)
                                        else:
                                            # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                            parent_folder = os.path.basename(os.path.dirname(file_path))
                                            unique_name = f"{parent_folder}_{file_name}"

                                            # Se anche questo nome è duplicato, aggiungi un contatore
                                            counter = 1
                                            while os.path.join("omonimi", unique_name) in added_zip_paths:
                                                unique_name = f"{parent_folder}_{counter}_{file_name}"
                                                counter += 1

                                            zip_path_to_use = os.path.join("omonimi", unique_name)

                                            # Registra questo file nel log degli omonimi
                                            omonimi_log.append({
                                                "nome_file": file_name,
                                                "percorso_originale": file_path,
                                                "primo_percorso": omonimi_files[file_name][0],
                                                "posizione_zip": zip_path_to_use
                                            })
                                    else:
                                        # Non è un omonimo, inseriscilo nella cartella principale
                                        zip_path_to_use = os.path.join(main_folder_name, rel_path)

                                # Verifica se questo percorso ZIP è già stato usato (potrebbe accadere anche con la struttura preservata)
                                if zip_path_to_use in added_zip_paths:
                                    self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                    # Crea un nome alternativo
                                    alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)

                                    # Se anche questo nome è già usato, aggiungi un contatore
                                    counter = 1
                                    while zip_path_to_use in added_zip_paths:
                                        alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                        zip_path_to_use = os.path.join("omonimi", alt_name)
                                        counter += 1

                                # Aggiungi il file al ZIP e registra il percorso
                                zipf.write(file_path, zip_path_to_use)
                                added_zip_paths.add(zip_path_to_use)

                            except Exception as e:
                                self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Comprimi i file singoli
                for file_path in filtered_single_files:
                    if os.path.exists(file_path):
                        file_name = os.path.basename(file_path)

                        try:
                            # Salta file non leggibili
                            if not os.access(file_path, os.R_OK):
                                self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                continue

                            # Calcola percorso relativo in base all'opzione selezionata
                            if result["preserve"]:
                                # Calcola il percorso relativo rispetto alla base
                                rel_path = os.path.relpath(file_path, base_path)
                                zip_path_to_use = os.path.join(main_folder_name, rel_path)
                            else:
                                # Comportamento originale senza preservare la struttura
                                # Verifica se è un file omonimo
                                if file_name in omonimi_files:
                                    # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                    if omonimi_files[file_name][0] == file_path:
                                        # Usa il nome del file all'interno della cartella principale
                                        zip_path_to_use = os.path.join(main_folder_name, file_name)
                                    else:
                                        # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                        parent_folder = os.path.basename(os.path.dirname(file_path))
                                        unique_name = f"{parent_folder}_{file_name}"

                                        # Se anche questo nome è duplicato, aggiungi un contatore
                                        counter = 1
                                        while os.path.join("omonimi", unique_name) in added_zip_paths:
                                            unique_name = f"{parent_folder}_{counter}_{file_name}"
                                            counter += 1

                                        zip_path_to_use = os.path.join("omonimi", unique_name)

                                        # Registra questo file nel log degli omonimi
                                        omonimi_log.append({
                                            "nome_file": file_name,
                                            "percorso_originale": file_path,
                                            "primo_percorso": omonimi_files[file_name][0],
                                            "posizione_zip": zip_path_to_use
                                        })
                                else:
                                    # Non è un omonimo, inseriscilo nella cartella principale
                                    zip_path_to_use = os.path.join(main_folder_name, file_name)

                            # Verifica se questo percorso ZIP è già stato usato
                            if zip_path_to_use in added_zip_paths:
                                self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                # Crea un nome alternativo
                                alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                zip_path_to_use = os.path.join("omonimi", alt_name)

                                # Se anche questo nome è già usato, aggiungi un contatore
                                counter = 1
                                while zip_path_to_use in added_zip_paths:
                                    alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)
                                    counter += 1

                            # Aggiungi il file al ZIP e registra il percorso
                            zipf.write(file_path, zip_path_to_use)
                            added_zip_paths.add(zip_path_to_use)

                        except Exception as e:
                            self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Crea log in formato CSV con collegamenti ipertestuali
                current_time_local = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                current_user = getpass.getuser()

                # Identifica quali file sono omonimi per poterli escludere dal log generale
                omonimo_paths = set()
                for entry in omonimi_log:
                    omonimo_paths.add(entry['percorso_originale'])

                # Lista di file normali (esclude gli omonimi)
                normal_files_log = [f for f in all_files_log if f['percorso_originale'] not in omonimo_paths]

                # ---------- CREAZIONE DEL CSV DEI FILE NORMALI ----------
                # Creazione intestazione per il CSV con informazioni sulla ricerca
                csv_header = f"LOG DEI FILE TROVATI - {current_time_local}\n"
                csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                csv_header += f"Directory di ricerca: {search_directory}\n"
                csv_header += f"Parole chiave ricercate: {search_keywords}\n"
                csv_header += f"Opzione struttura: {'Preservata' if result['preserve'] else 'Piatta'}\n"
                csv_header += f"Totale file: {len(all_files_log)}\n\n"

                # Prepara l'output CSV
                csv_content = io.StringIO()
                csv_writer = csv.writer(csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)

                # Intestazione colonne
                csv_writer.writerow(["Nome file", "Tipo", "Dimensione", "Ultima modifica", "Percorso completo", "Link"])

                # Aggiungi i file normali (non omonimi)
                for file_entry in sorted(all_files_log, key=lambda x: x["nome_file"]):
                    # Ottieni il percorso della directory che contiene il file
                    file_path = file_entry['percorso_originale']
                    directory_path = os.path.dirname(file_path).replace('\\', '/')
                    # Formatta il percorso della directory per il collegamento ipertestuale
                    folder_uri = f"file:///{directory_path}"
                    
                    # Aggiungi riga al CSV con la formula italiana per i collegamenti alla CARTELLA
                    csv_writer.writerow([
                        file_entry['nome_file'],
                        file_entry['tipo'],
                        file_entry['dimensione'],
                        file_entry['ultima_modifica'],
                        file_entry['percorso_originale'],
                        f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")'
                    ])

                # Ottieni il contenuto del CSV come stringa
                csv_data = csv_header + csv_content.getvalue()

                # Aggiungi il file CSV dei file normali all'archivio
                zipf.writestr("Log_file_Trovati.csv", csv_data)

                # ---------- CREAZIONE DEL CSV DEGLI OMONIMI (SOLO SE CE NE SONO) ----------
                if omonimi_log:
                    # Creazione intestazione per il CSV degli omonimi con info di ricerca
                    omonimi_csv_header = f"LOG DEI FILE OMONIMI - {current_time_local}\n"
                    omonimi_csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                    omonimi_csv_header += f"Directory di ricerca: {search_directory}\n"
                    omonimi_csv_header += f"Parole chiave ricercate: {search_keywords}\n"
                    omonimi_csv_header += f"Totale file omonimi trovati: {len(omonimi_log)}\n\n"
                    
                    # Prepara l'output CSV per gli omonimi
                    omonimi_csv_content = io.StringIO()
                    omonimi_csv_writer = csv.writer(omonimi_csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
                    
                    # Intestazione colonne specifiche per gli omonimi
                    omonimi_csv_writer.writerow([
                        "Nome file", "Tipo", "Dimensione", "Ultima modifica", 
                        "Percorso originale", "Link", "Posizione nello ZIP", "In conflitto con"
                    ])
                    
                    # Estrai dai log completi solo i file omonimi
                    omonimi_files_details = [f for f in all_files_log if f['percorso_originale'] in omonimo_paths]
                    
                    # Crea un dizionario per unire le informazioni di omonimi_log e omonimi_files_details
                    omonimi_details_dict = {}
                    for entry in omonimi_log:
                        omonimi_details_dict[entry['percorso_originale']] = entry
                    
                    # Aggiungi i file omonimi al CSV
                    for file_entry in sorted(omonimi_files_details, key=lambda x: x["nome_file"]):
                        # Ottieni il percorso della directory che contiene il file
                        file_path = file_entry['percorso_originale']
                        directory_path = os.path.dirname(file_path).replace('\\', '/')
                        # Formatta il percorso della directory per il collegamento ipertestuale
                        folder_uri = f"file:///{directory_path}"
                        
                        # Trova le informazioni aggiuntive sugli omonimi
                        omonimo_info = omonimi_details_dict.get(file_entry['percorso_originale'], {})
                        posizione_zip = omonimo_info.get('posizione_zip', 'N/A')
                        in_conflitto = omonimo_info.get('primo_percorso', 'N/A')
                        
                        # Aggiungi riga al CSV degli omonimi con la formula italiana per i collegamenti alla CARTELLA
                        omonimi_csv_writer.writerow([
                            file_entry['nome_file'],
                            file_entry['tipo'],
                            file_entry['dimensione'],
                            file_entry['ultima_modifica'],
                            file_entry['percorso_originale'],
                            f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")',
                            posizione_zip,
                            in_conflitto
                        ])
                    
                    # Ottieni il contenuto del CSV come stringa
                    omonimi_csv_data = omonimi_csv_header + omonimi_csv_content.getvalue()
                    
                    # Aggiungi il file CSV degli omonimi all'archivio
                    zipf.writestr("omonimi_log.csv", omonimi_csv_data)

            # Verifica se ci sono algoritmi di hash selezionati
            hash_algorithms = result.get("hash_algorithms", [])
            hash_results = {}
            hash_success = False

            # Calcola gli hash del file ZIP se richiesto
            if hash_algorithms:
                try:
                    self.status_label["text"] = f"Calcolo hash del file ZIP..."
                    self.root.update_idletasks()
                    
                    # Calcola gli hash sul file ZIP
                    hash_results = self.calculate_file_hash(zip_path, hash_algorithms)
                    hash_success = True
                    
                    # Crea un file di report per gli hash
                    hash_report_path = os.path.splitext(zip_path)[0] + "_hash.txt"
                    with open(hash_report_path, 'w', encoding='utf-8') as hash_file:
                        hash_file.write(f"REPORT HASH PER IL FILE: {os.path.basename(zip_path)}\n")
                        hash_file.write(f"Data e ora: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n")
                        hash_file.write(f"Utente: {self.current_user}\n\n")
                        
                        for algorithm in hash_algorithms:
                            hash_value = hash_results.get(algorithm, "Errore")
                            hash_file.write(f"{algorithm.upper()}: {hash_value}\n")
                
                except Exception as e:
                    self.log_debug(f"Errore nel calcolo hash del file ZIP: {str(e)}")

            # Prepara il messaggio di completamento
            skipped_files = len(single_files) - len(filtered_single_files)
            message = f"Compressione completata!\nFile salvato in: {zip_path}\n"
            message += f"File organizzati nella cartella '{main_folder_name}'\n"
            message += f"Tipo di compressione utilizzata: {compression_text}\n"
            
            # Aggiungi informazioni sugli hash calcolati
            if hash_algorithms:
                if hash_success:
                    message += f"\nHash calcolati sul file ZIP:\n"
                    for algorithm in hash_algorithms:
                        hash_value = hash_results.get(algorithm, "Errore")
                        message += f"{algorithm.upper()}: {hash_value}\n"
                    message += f"\nI dettagli sono stati salvati in: {os.path.basename(hash_report_path)}"
                else:
                    message += f"\nErrore nel calcolo degli hash richiesti."

            if result["preserve"]:
                message += f"Struttura delle directory originali: Preservata\n"
            else:
                message += f"Struttura delle directory originali: Non preservata\n"
                
            message += f"\nCreato file di log completo ('Log_file_Trovati.csv')"

            if result["chunks"]:
                message += "\nElaborazione a blocchi: Attiva"

            if omonimi_log:
                message += f"\nTrovati {len(omonimi_log)} file omonimi."
                message += f"\nCreato log dettagliato degli omonimi ('omonimi_log.csv')"

            if skipped_files > 0:
                message += f"\n{skipped_files} file saltati perché già presenti nelle cartelle"

            messagebox.showinfo("Completato", message)

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            self.log_debug(f"Errore dettagliato durante la compressione:\n{error_details}")

            messagebox.showerror("Errore", 
                f"Errore durante la compressione: {str(e)}\n\n"
                f"Tipo di errore: {type(e).__name__}")

        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    @error_handler
    def _find_common_base_path(self, selected_items):
        """Trova il percorso base comune a tutti i file selezionati"""
        if not selected_items:
            return ""
        
        paths = []
        
        # Raccogli tutti i percorsi
        for item in selected_items:
            values = self.results_list.item(item)['values']
            paths.append(values[5])  # Il percorso è nella sesta colonna
        
        # Trova il percorso comune più lungo
        def common_path(paths):
            if not paths:
                return ""
            
            # Normalizza i percorsi per gestire i diversi separatori di directory
            paths = [os.path.normpath(p) for p in paths]
            
            # Dividi ogni percorso in componenti
            components = [p.split(os.path.sep) for p in paths]
            
            # Se ci sono percorsi di unità diverse (es. C:\ e D:\) su Windows
            if os.name == 'nt' and len(set(c[0] for c in components if c)) > 1:
                return ""  # Non c'è un percorso comune tra unità diverse
            
            common = []
            for i in range(min(len(c) for c in components)):
                if len(set(c[i] for c in components)) == 1:
                    common.append(components[0][i])
                else:
                    break
            
            # Se non c'è un elemento comune, restituisci la radice/unità
            if not common and os.name == 'nt' and components and components[0]:
                return components[0][0] + os.path.sep  # es. "C:\"
            elif not common:
                return os.path.sep  # radice Unix "/"
                
            return os.path.sep.join(common)
        
        base_path = common_path(paths)
        
        # Se il percorso base finisce con un separatore, va bene
        # altrimenti dobbiamo aggiungere il separatore se base_path non è vuoto
        if base_path and not base_path.endswith(os.path.sep):
            base_path = os.path.dirname(base_path)
        
        # Se non abbiamo trovato un percorso comune, usa il percorso di ricerca corrente
        if not base_path and hasattr(self, 'search_path'):
            base_path = self.search_path.get()
        
        return base_path
    
    @error_handler
    def get_directory_size(self, path):
        """Calculate the total size of a directory"""
        if not os.path.exists(path):
            return 0
            
        try:
            total_size = 0
            if os.path.isfile(path):
                return os.path.getsize(path)
                
            # Per evitare il blocco dell'interfaccia, imposta un timeout massimo
            start_time = time.time()
            max_time = 30  # massimo 30 secondi
            files_count = 0
            error_count = 0
                
            # For directories, walk through all files and subdirectories
            for dirpath, dirnames, filenames in os.walk(path):
                for f in filenames:
                    # Verifica se il timeout è scaduto
                    if time.time() - start_time > max_time:
                        self.log_debug(f"Timeout nel calcolo della dimensione per {path}")
                        return total_size
                        
                    try:
                        fp = os.path.join(dirpath, f)
                        
                        # Verifica esplicita che il file esiste ancora prima di tentare di leggerne la dimensione
                        if os.path.exists(fp) and not os.path.islink(fp):
                            total_size += os.path.getsize(fp)
                            files_count += 1
                    except FileNotFoundError:
                        # Ignora silenziosamente i file che non esistono più
                        pass
                    except (OSError, PermissionError) as e:
                        error_count += 1
                        # Limita il numero di errori da registrare per evitare spam nel log
                        if error_count < 100:  
                            self.log_debug(f"Error getting size of {fp}: {str(e)}")
                
                # Update the status periodically to show progress
                if files_count % 1000 == 0:  # Ogni 1000 file
                    self.root.after(0, lambda size=total_size: 
                        self.status_label.config(text=f"Calcolando dimensione: {self._format_size(size)}..."))
                    
            return total_size
        except Exception as e:
            self.log_debug(f"Error calculating directory size for {path}: {str(e)}")
            return 0
    
    @error_handler
    def get_directory_size_system(self, path):
        """Utilizza comandi di sistema per ottenere dimensioni di directory molto grandi"""
        try:
            if os.name == 'nt':  # Windows
                # Usa PowerShell per calcolare la dimensione (molto più veloce per directory grandi)
                cmd = f'powershell -command "Get-ChildItem -Path \'{path}\' -Recurse -Force -ErrorAction SilentlyContinue | Measure-Object -Property Length -Sum | Select-Object -ExpandProperty Sum"'
                result = subprocess.check_output(cmd, shell=True, stderr=subprocess.STDOUT)
                size = int(result.strip())
                return size
            else:  # Linux/Unix
                # Usa il comando du che è ottimizzato per il calcolo delle dimensioni
                result = subprocess.check_output(['du', '-sb', path])
                size = int(result.split()[0])
                return size
        except Exception as e:
            self.log_debug(f"Errore nel calcolo della dimensione tramite comando di sistema: {str(e)}")
            # Fallback al metodo standard
            return self.get_directory_size(path)
    
    @error_handler
    def estimate_directory_size(self, path, sample_size=100):
        """Stima la dimensione di una directory campionando alcuni file"""
        import random
        
        if not os.path.exists(path) or os.path.isfile(path):
            return self.get_directory_size(path)  # Usa il metodo esatto per file o percorsi non validi
        
        try:
            # Ottieni un conteggio rapido dei file (questo è veloce)
            total_files = 0
            sampled_files = 0
            total_sampled_size = 0
            
            # Prima passata veloce per contare i file
            for root, _, files in os.walk(path, topdown=True):
                total_files += len(files)
                # Limita il tempo della prima passata
                if total_files > 10000:  # Se ci sono più di 10000 file, passiamo alla stima
                    break
            
            # Se pochi file, usare metodo preciso
            if total_files < 1000:
                return self.get_directory_size(path)
                
            # Seconda passata per il campionamento
            for root, _, files in os.walk(path, topdown=True):
                for file in files:
                    # Campiona casualmente 1 file ogni X
                    if random.randint(1, max(1, total_files // sample_size)) == 1:
                        try:
                            file_path = os.path.join(root, file)
                            if os.path.exists(file_path) and not os.path.islink(file_path):
                                total_sampled_size += os.path.getsize(file_path)
                                sampled_files += 1
                        except:
                            pass
                    
                    # Se abbiamo campionato abbastanza file, calcola la stima
                    if sampled_files >= sample_size:
                        break
                
                if sampled_files >= sample_size:
                    break
            
            # Calcola la stima finale
            if sampled_files > 0:
                avg_file_size = total_sampled_size / sampled_files
                estimated_size = avg_file_size * total_files
                self.log_debug(f"Dimensione stimata per {path}: {self._format_size(estimated_size)} (basata su {sampled_files} campioni)")
                return estimated_size
            else:
                return self.get_directory_size(path)  # Fallback al metodo standard
                
        except Exception as e:
            self.log_debug(f"Errore nella stima della dimensione: {str(e)}")
            return 0
    @error_handler
    def get_disk_space(self, path):
        """Get disk space information for the partition containing the path"""
        if not os.path.exists(path):
            return (0, 0, 0)
        
        try:
            if os.name == 'nt':  # Windows
                total, used, free = shutil.disk_usage(os.path.splitdrive(path)[0] + '\\')
            else:  # Linux/Mac
                total, used, free = shutil.disk_usage(os.path.abspath(os.path.join(path, os.pardir)))
            return (total, used, free)
        except Exception as e:
            self.log_debug(f"Error getting disk space for {path}: {str(e)}")
            return (0, 0, 0)

    @error_handler
    def update_disk_info(self, path=None, calculate_dir_size=True):
        """Aggiorna le informazioni sul disco e sulla directory"""
        if path is None:
            path = self.search_path.get()
        
        if not path or not os.path.exists(path):
            return

        # Aggiorna la variabile di controllo globale in base al parametro
        self.directory_calculation_enabled = calculate_dir_size
        
        # Solo per debug
        self.log_debug(f"Aggiornamento info disco: path={path}, calculation_enabled={self.directory_calculation_enabled}")
        
        # Se il calcolo è disabilitato, aggiorna solo le informazioni del disco
        if not self.directory_calculation_enabled:
            self.log_debug("Calcolo dimensione directory disabilitato nelle impostazioni")
            
            # Aggiorniamo solo le informazioni del disco senza calcolare la dimensione della directory
            try:
                total, used, free = self.get_disk_space(path)
                percent_used = (used / total) * 100 if total > 0 else 0
                
                total_formatted = self._format_size(total)
                free_formatted = self._format_size(free)
                
                # CORREZIONE: Verifica se il widget esiste prima di usarlo
                if hasattr(self, 'disk_info_label') and self.disk_info_label.winfo_exists():
                    self.disk_info_label.config(text=f"Disco: {total_formatted} totali, {free_formatted} liberi ({percent_used:.1f}% usato)")
                else:
                    self.log_debug("Widget disk_info_label non disponibile")
                
                # Importante: NON avviare il thread di calcolo della directory
                # CORREZIONE: Verifica se il widget esiste prima di usarlo
                if hasattr(self, 'dir_size_label') and self.dir_size_label.winfo_exists():
                    self.dir_size_label.config(text="Dimensione directory: calcolo disabilitato")
                else:
                    self.log_debug("Widget dir_size_label non disponibile")
            except Exception as e:
                self.log_debug(f"Errore nell'aggiornamento info disco: {str(e)}")
            
            return
        
        # Se arriviamo qui, il calcolo è abilitato e dobbiamo procedere
        self._async_update_disk_info(path, calculate_dir_size)
    
    @error_handler
    def _update_disk_info_thread(self, path, calculate_dir_size):
        """Thread per calcolare le informazioni del disco"""
        try:
            # Get disk space info
            total, used, free = self.get_disk_space(path)
            
            # Aggiorna l'UI dal thread principale
            self.root.after(0, lambda: self.total_disk_var.set(self._format_size(total)))
            self.root.after(0, lambda: self.used_disk_var.set(self._format_size(used)))
            self.root.after(0, lambda: self.free_disk_var.set(self._format_size(free)))
            self.root.after(0, lambda: self.status_label.config(text="In attesa..."))
            
            # Log per debug
            self.log_debug(f"Info disco aggiornate: Totale={self._format_size(total)}, Libero={self._format_size(free)}")
            
        except Exception as e:
            self.log_debug(f"Errore nel calcolo dello spazio disco: {str(e)}")
            self.root.after(0, lambda: self.total_disk_var.set("Errore"))
            self.root.after(0, lambda: self.used_disk_var.set("Errore"))
            self.root.after(0, lambda: self.free_disk_var.set("Errore"))
        
        # Controlla esplicitamente le impostazioni per il calcolo della directory
        calculation_mode = self.dir_size_calculation.get()
        
        # Gestisce il caso in cui non si deve calcolare la dimensione
        if not calculate_dir_size or calculation_mode == "disabilitato" or self.is_searching:
            # Se il calcolo è disabilitato o stiamo cercando, mostra "N/D" invece di calcolare
            self.log_debug("Calcolo dimensione directory saltato: " + 
                        ("parametro disabilitato" if not calculate_dir_size else 
                        "impostazione disabilitata" if calculation_mode == "disabilitato" else 
                        "ricerca in corso"))
            
            self.root.after(0, lambda: self.dir_size_var.set("N/D"))
            return
        
        # Se arriviamo qui, dobbiamo calcolare la dimensione
        try:
            self.log_debug(f"Avvio calcolo dimensione directory: {path} (modalità: {calculation_mode})")
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            self._calculate_dir_size_thread(path)
        except Exception as e:
            self.log_debug(f"Errore nell'avvio del calcolo dimensione directory: {str(e)}")
            self.root.after(0, lambda: self.dir_size_var.set("Errore"))

    @error_handler
    def _async_update_disk_info(self, path, calculate_dir_size):
        """Esegue il calcolo delle informazioni del disco in background"""
        try:
            # Get disk space info
            total, used, free = self.get_disk_space(path)
            
            # Aggiorna l'UI dal thread principale
            self.root.after(0, lambda: self.total_disk_var.set(self._format_size(total)))
            self.root.after(0, lambda: self.used_disk_var.set(self._format_size(used)))
            self.root.after(0, lambda: self.free_disk_var.set(self._format_size(free)))
        except Exception as e:
            self.log_debug(f"Errore nel calcolo dello spazio disco: {str(e)}")
            self.root.after(0, lambda: self.total_disk_var.set("Errore"))
            self.root.after(0, lambda: self.used_disk_var.set("Errore"))
            self.root.after(0, lambda: self.free_disk_var.set("Errore"))
        
        # Verifica modalità di calcolo dimensione
        calculation_mode = self.dir_size_calculation.get()
        
        # Rispetta sempre l'impostazione disabilitato, indipendentemente dal valore di calculate_dir_size
        if calculation_mode == "disabilitato":
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo disattivato"))
            self.log_debug("Calcolo dimensione directory disattivato nelle impostazioni")
            return
        
        # Se è incrementale e siamo in fase di ricerca, non fare niente qui
        if calculation_mode == "incrementale" and self.is_searching:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            return
                
        # Per le altre modalità o se non siamo in ricerca
        if calculate_dir_size:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            # Il calcolo avviene già in un thread separato
            self._calculate_dir_size_thread(path)
        else:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo disattivato"))
    
    @error_handler
    def _calculate_dir_size_thread(self, path):
        """Thread function to calculate directory size"""
        # Limita il carico di lavoro per evitare blocchi del sistema
        import threading
        try:
            # Imposta una priorità più bassa per questo thread
            if hasattr(threading.current_thread(), "setName"):
                threading.current_thread().name = "LowPriority_DirSize"
        except:
            pass
            
        calculation_mode = self.dir_size_calculation.get()
        dir_size = 0
        
        try:
            if calculation_mode == "preciso":
                dir_size = self.get_directory_size(path)
            elif calculation_mode == "stimato":
                dir_size = self.estimate_directory_size(path)
            elif calculation_mode == "sistema":
                dir_size = self.get_directory_size_system(path)
            else:  # incrementale o fallback
                dir_size = self.get_directory_size(path)
                
            # Update the UI from the main thread - CORREZIONE con check
            if self.root and self.root.winfo_exists():
                self.root.after(0, lambda: self.dir_size_var.set(self._format_size(dir_size)))
                self.root.after(0, lambda: self.status_label.config(text="In attesa...") if hasattr(self, 'status_label') and self.status_label.winfo_exists() else None)
        except Exception as e:
            self.log_debug(f"Errore nel calcolo della dimensione: {str(e)}")
            if self.root and self.root.winfo_exists():
                self.root.after(0, lambda: self.dir_size_var.set("Errore"))
                self.root.after(0, lambda: self.status_label.config(text="In attesa...") if hasattr(self, 'status_label') and self.status_label.winfo_exists() else None)

    @error_handler
    def refresh_directory_size(self):
        """Aggiorna manualmente il calcolo della dimensione della directory"""
        # Ottieni il percorso corrente
        path = self.search_path.get()
        
        # Verifica che il percorso esista
        if not path or not os.path.exists(path):
            messagebox.showinfo("Informazione", "Seleziona prima un percorso valido")
            return
            
        # Verifica la modalità di calcolo
        calculation_mode = self.dir_size_calculation.get()
        if calculation_mode == "disabilitato":
            response = messagebox.askyesno("Calcolo disabilitato", 
                                        "Il calcolo della dimensione è attualmente disabilitato.\n\n" +
                                        "Vuoi attivarlo e procedere con il calcolo?")
            if response:
                # Seleziona la modalità "preciso"
                self.dir_size_calculation.set("preciso")
            else:
                return
        
        # Aggiorna lo stato
        self.dir_size_var.set("Calcolo in corso...")
        self.status_label.config(text="Calcolo dimensione directory...")
        
        # Esegui il calcolo in un thread separato
        threading.Thread(target=self._calculate_dir_size_thread, args=(path,), daemon=True).start()

    @error_handler # Funzione helper per formattare la dimensione del file
    def _format_size(self, size_bytes):
        """Formatta la dimensione del file in modo leggibile"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.2f} KB"
        elif size_bytes < 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024):.2f} MB"
        else:
            return f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"
    
    @error_handler
    def calculate_file_hash(self, file_path, algorithms=None):
        """Calcola gli hash di un file usando gli algoritmi specificati :param file_path: Percorso del file
        :param algorithms: Lista degli algoritmi da usare ('md5', 'sha1', 'sha256'):return: Dizionario con gli hash calcolati"""
        import hashlib
        
        if algorithms is None:
            algorithms = ['md5']
        
        results = {}
        
        try:
            for algorithm in algorithms:
                if algorithm == 'md5':
                    hash_obj = hashlib.md5()
                elif algorithm == 'sha1':
                    hash_obj = hashlib.sha1()
                elif algorithm == 'sha256':
                    hash_obj = hashlib.sha256()
                else:
                    continue
                    
                # Leggi il file in blocchi per gestire file di grandi dimensioni
                with open(file_path, 'rb') as f:
                    for chunk in iter(lambda: f.read(4096), b''):
                        hash_obj.update(chunk)
                    
                results[algorithm] = hash_obj.hexdigest()
                
                # Log per debug
                self.log_debug(f"Calcolato {algorithm} per {os.path.basename(file_path)}: {results[algorithm]}")
                
        except Exception as e:
            self.log_debug(f"Errore nel calcolo hash {algorithm} per {file_path}: {str(e)}")
            # In caso di errore, restituisci None per quell'algoritmo
            for algorithm in algorithms:
                if algorithm not in results:
                    results[algorithm] = "Errore"
        
        return results

    @error_handler # Funzione helper per determinare il tipo di file
    def _get_file_type(self, file_path):
        """Determina il tipo di file in base all'estensione"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.txt', '.md', '.rtf']:
            return "Documento di testo"
        elif ext in ['.doc', '.docx', '.odt']:
            return "Documento Word"
        elif ext in ['.xls', '.xlsx', '.ods']:
            return "Foglio di calcolo"
        elif ext in ['.pdf']:
            return "PDF"
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            return "Immagine"
        elif ext in ['.mp3', '.wav', '.ogg', '.flac']:
            return "Audio"
        elif ext in ['.mp4', '.avi', '.mkv', '.mov']:
            return "Video"
        elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
            return "Archivio"
        elif ext in ['.exe', '.dll', '.bat', '.cmd']:
            return "Eseguibile"
        else:
            return "File"
    
    @error_handler
    def get_main_folder_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome della cartella principale"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Cartella Principale")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=tk.BOTH, expand=tk.YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome della cartella principale nell'archivio:", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = tk.StringVar(value="files")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=tk.X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=tk.X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            folder_name = name_var.get().strip()
            if folder_name:
                result["name"] = folder_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per la cartella", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea Cartella", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=tk.RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]
    
    @error_handler
    def open_file_location(self, event=None):
        """Apre il percorso del file selezionato nel file explorer"""
            
        selected_items = self.results_list.selection()
        if not selected_items:
            return
            
        selected_item = selected_items[0]  # Prendi il primo elemento selezionato
        file_path = self.results_list.item(selected_item, "values")[6]  # Modificato da 5 a 6
        
        try:
            if os.path.exists(file_path):
                # Ottieni la directory contenente il file
                directory = os.path.dirname(file_path)
                
                if os.name == 'nt':  # Windows
                    # Converti eventuali forward slash in backslash per Windows
                    file_path = os.path.normpath(file_path)
                    # Usa il metodo più sicuro con subprocess invece di os.system
                    subprocess.run(['explorer', '/select,', file_path], shell=True)
                else:
                    # Per sistemi Linux/Unix
                    if shutil.which('xdg-open'):  # Verifica che xdg-open sia disponibile
                        subprocess.run(['xdg-open', directory])
                    elif shutil.which('open'):  # Per macOS
                        subprocess.run(['open', directory])
                    else:
                        self.log_debug("Nessun comando disponibile per aprire directory")
                        messagebox.showinfo("Informazione", f"Percorso del file: {directory}")
                
                self.log_debug(f"Apertura percorso: {file_path}")
            else:
                messagebox.showinfo("Informazione", f"Il percorso non esiste: {file_path}")
                self.log_debug(f"Percorso non esistente: {file_path}")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile aprire il percorso: {str(e)}")
            self.log_debug(f"Errore nell'apertura del percorso: {str(e)}")

    @error_handler
    def show_advanced_filters_dialog(self):
        """Mostra la finestra di dialogo per i filtri di ricerca avanzati"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Filtri avanzati")
        dialog.geometry("530x250")  # Increased height to accommodate the new field
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Filtri dimensione
        size_frame = ttk.LabelFrame(dialog, text="Dimensione file")
        size_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(size_frame, text="Min (KB):").grid(row=0, column=0, padx=5, pady=5)
        min_size = ttk.Entry(size_frame, width=10)
        min_size.grid(row=0, column=1, padx=5, pady=5)
        min_size.insert(0, str(self.advanced_filters["size_min"] // 1024))
        self.create_tooltip(min_size, "Quando viene imposto un valore maggiore di zero, solo i file con dimensione superiore a questo valore verranno inclusi nei risultati")
        
        ttk.Label(size_frame, text="Max (KB):").grid(row=0, column=2, padx=5, pady=5)
        max_size = ttk.Entry(size_frame, width=10)
        max_size.grid(row=0, column=3, padx=5, pady=5)
        max_size.insert(0, str(self.advanced_filters["size_max"] // 1024 if self.advanced_filters["size_max"] else 0))
        self.create_tooltip(max_size, "Quando viene imposto un valore maggiore di zero, solo i file con dimensione inferiore a questo valore verranno inclusi nei risultati")

        # Filtri data - FORMAT DD-MM-YYYY
        date_frame = ttk.LabelFrame(dialog, text="Data modifica (DD-MM-YYYY)")
        date_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(date_frame, text="Da:").grid(row=0, column=0, padx=5, pady=5)
        min_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        min_date.grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(date_frame, text="A:").grid(row=0, column=2, padx=5, pady=5)
        max_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        max_date.grid(row=0, column=3, padx=5, pady=5)
        
        # Cancella sempre le date precedenti
        min_date.entry.delete(0, "end")
        max_date.entry.delete(0, "end")
              
        # Aggiungiamo un frame di debug per vedere i filtri correnti
        debug_frame = ttk.LabelFrame(dialog, text="Debug - Stato filtri correnti")
        debug_frame.pack(fill=X, padx=10, pady=5)
        
        debug_text = ttk.Text(debug_frame, height=3, width=50)
        debug_text.pack(fill=X, padx=5, pady=5)
        debug_text.insert("1.0", f"Data min: {self.advanced_filters['date_min']}\n")
        debug_text.insert("2.0", f"Data max: {self.advanced_filters['date_max']}\n")
        debug_text.insert("3.0", f"Extensions: {self.advanced_filters['extensions']}")
        debug_text.config(state="disabled")
        
        # Pulsante Salva
        def save_filters():
            try:
                # Analizza i filtri di dimensione
                min_kb = int(min_size.get() or 0)
                max_kb = int(max_size.get() or 0)
                self.advanced_filters["size_min"] = min_kb * 1024
                self.advanced_filters["size_max"] = max_kb * 1024
                
                # Ottieni le date nel formato DD-MM-YYYY
                min_date_value = min_date.entry.get().strip()
                max_date_value = max_date.entry.get().strip()
                
                print(f"DEBUG - Date inserite: min={min_date_value}, max={max_date_value}")
                
                # Validazione date
                if min_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(min_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data minima non valido. Usa DD-MM-YYYY")
                        return
                
                if max_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(max_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data massima non valido. Usa DD-MM-YYYY")
                        return
                
                # Verifica che la data minima non sia successiva alla data massima
                if min_date_value and max_date_value:
                    min_date_obj = datetime.strptime(min_date_value, "%d-%m-%Y")
                    max_date_obj = datetime.strptime(max_date_value, "%d-%m-%Y")
                    
                    if min_date_obj > max_date_obj:
                        messagebox.showerror("Errore", 
                                        "La data di inizio non può essere successiva alla data di fine")
                        return
                
                # Salva le date validate
                self.advanced_filters["date_min"] = min_date_value
                self.advanced_filters["date_max"] = max_date_value
                
                print(f"Filtri salvati: {self.advanced_filters}")  # Debug info
                dialog.destroy()
                
            except ValueError:
                messagebox.showerror("Errore", "Inserisci valori numerici validi per le dimensioni")
        
        # Creato un frame per i pulsanti nella stessa riga
        buttons_frame = ttk.Frame(dialog)
        buttons_frame.pack(fill=X, pady=10, padx=10)
        
        # Pulsante Annulla a sinistra
        ttk.Button(buttons_frame, text="Annulla", command=dialog.destroy).pack(side=RIGHT, padx=5)
        
        # Pulsante Salva a destra
        ttk.Button(buttons_frame, text="Salva", command=save_filters).pack(side=RIGHT, padx=5)

        dialog.update_idletasks()  # Aggiorna la finestra per ottenere le dimensioni corrette
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima per la finestra
        dialog.minsize(500, 400)

    @error_handler
    def configure_extensions(self, mode="base"):
        """Dialog to configure file extensions for different search modes"""
        dialog = ttk.Toplevel(self.root)
        dialog.title(f"Configura estensioni - Modalità {mode.capitalize()}")
        dialog.geometry("1000x450")
        dialog.transient(self.root)
        dialog.grab_set()
        
        main_frame = ttk.Frame(dialog, padding=15)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Heading text
        ttk.Label(main_frame, text=f"Seleziona le estensioni di file da includere nella ricerca per la modalità {mode}", 
                font=("", 10)).pack(anchor=W, pady=(0, 15))
        
        # Category tabs
        notebook = ttk.Notebook(main_frame)
        notebook.pack(fill=BOTH, expand=YES, pady=(0, 15))
        
        # Initialize dictionary to store checkbutton variables
        ext_vars = {}
        
        # Define extensions by category - EXPANDED with more extensions
        extension_categories = {
        "Documenti": [
            (".txt", "File di testo"),
            (".doc", "Word vecchio"),
            (".docx", "Word"),
            (".pdf", "PDF"),
            (".rtf", "Rich Text"),
            (".odt", "OpenDoc Text"),
            (".md", "Markdown"),
            (".xml", "XML"),
            (".html", "HTML"),
            (".htm", "HTM"),
            (".log", "Log file"),
            (".tex", "LaTeX"),
            (".rst", "reStructuredText"),
            (".epub", "E-book EPUB"),
            (".mobi", "E-book Mobi"),
            (".vcf", "vCard"),
            (".ics", "iCalendar"),
        ],
        "Email": [
            (".eml", "Email standard"),
            (".msg", "Email formato Outlook"),
            (".pst", "Archivio Outlook"),
            (".ost", "Archivio Outlook Offline"),
            (".mbox", "Mailbox Unix/Linux"),
            (".emlx", "Email formato Apple Mail"),
            (".mbx", "Mailbox formato Eudora/Thunderbird"),
            (".dbx", "Archivio Outlook Express"),
            (".wab", "Windows Address Book"),
            (".nws", "Email formato Outlook Express"),
            (".mht", "MIME HTML Archive"),
            (".mhtml", "MIME HTML Archive"),
            (".imapmbox", "IMAP Mailbox"),
            (".email", "Email generica")
        ],
        "Fogli calcolo": [
            (".xls", "Excel vecchio"),
            (".xlsx", "Excel"),
            (".ods", "OpenCalc"),
            (".csv", "CSV (valori separati da virgola)"),
            (".tsv", "TSV (valori separati da tab)"),
            (".dbf", "Database File"),
            (".dif", "Data Interchange Format")
        ],
        "Presentazioni": [
            (".ppt", "PowerPoint vecchio"),
            (".pptx", "PowerPoint"),
            (".odp", "OpenImpress"),
            (".key", "Keynote"),
            (".pps", "PowerPoint Show")
        ],
        "Database": [
            (".db", "Database generico"),
            (".sqlite", "SQLite"),
            (".sqlite3", "SQLite3"),
            (".mdb", "Access DB"),
            (".accdb", "Access DB nuovo"),
            (".odb", "OpenOffice DB"),
            (".sql", "SQL script")
        ],
        "Immagini": [
            (".jpg", "JPEG"),
            (".jpeg", "JPEG"),
            (".png", "PNG"),
            (".gif", "GIF"),
            (".bmp", "Bitmap"),
            (".tiff", "TIFF"),
            (".tif", "TIF"),
            (".svg", "SVG"),
            (".webp", "WebP"),
            (".ico", "Icon"),
            (".raw", "Raw"),
            (".psd", "Photoshop"),
            (".ai", "Illustrator"),
            (".odg", "OpenOffice Draw"),
            (".xcf", "GIMP"),
            (".heic", "HEIC")
        ],
        "Audio": [
            (".mp3", "MP3"),
            (".wav", "WAV"),
            (".ogg", "OGG"),
            (".flac", "FLAC"),
            (".aac", "AAC"),
            (".m4a", "M4A"),
            (".wma", "WMA"),
            (".mid", "MIDI"),
            (".midi", "MIDI"),
            (".aiff", "AIFF"),
            (".opus", "Opus")
        ],
        "Video": [
            (".mp4", "MP4"),
            (".avi", "AVI"),
            (".mkv", "MKV"),
            (".mov", "MOV"),
            (".wmv", "WMV"),
            (".flv", "FLV"),
            (".webm", "WebM"),
            (".m4v", "M4V"),
            (".mpg", "MPEG"),
            (".mpeg", "MPEG"),
            (".3gp", "3GP"),
            (".ogv", "OGV"),
            (".ts", "TS (Transport Stream)")
        ],
        "Archivi": [
            (".zip", "ZIP"),
            (".rar", "RAR"),
            (".7z", "7-Zip"),
            (".tar", "TAR"),
            (".gz", "GZip"),
            (".bz2", "BZip2"),
            (".iso", "ISO"),
            (".tgz", "Tar GZipped"),
            (".xz", "XZ"),
            (".cab", "Cabinet"),
            (".jar", "Java Archive")
        ],
        "Eseguibili": [
            (".exe", "Eseguibile"),
            (".dll", "Libreria"),
            (".bat", "Batch"),
            (".cmd", "Command"),
            (".ps1", "PowerShell"),
            (".vbs", "VBScript"),
            (".sh", "Shell script"),
            (".msi", "Installer Windows"),
            (".app", "Applicazione macOS"),
            (".deb", "Pacchetto Debian"),
            (".rpm", "Red Hat Package"),
            (".apk", "Android Package")
        ],
        "Configurazione": [
            (".ini", "INI"),
            (".config", "Config"),
            (".conf", "Config"),
            (".reg", "Registry"),
            (".cfg", "Config"),
            (".properties", "Properties"),
            (".yml", "YAML"),
            (".yaml", "YAML"),
            (".json", "JSON Config"),
            (".toml", "TOML"),
            (".env", "Environment"),
            (".htaccess", "Apache Config"),
            (".plist", "macOS Property List")
        ],
        "Programmazione": [
            (".c", "C"),
            (".cpp", "C++"),
            (".cs", "C#"),
            (".java", "Java"),
            (".py", "Python"),
            (".js", "JavaScript"),
            (".php", "PHP"),
            (".rb", "Ruby"),
            (".go", "Golang"),
            (".rs", "Rust"),
            (".swift", "Swift"),
            (".pl", "Perl"),
            (".lua", "Lua"),
            (".h", "Header C"),
            (".hpp", "Header C++"),
            (".vb", "Visual Basic"),
            (".ts", "TypeScript"),
            (".scala", "Scala"),
            (".groovy", "Groovy"),
            (".kt", "Kotlin")
        ]
    }
        
        # Create a list of all extensions for "profonda" mode
        all_extensions = []
        for category_extensions in extension_categories.values():
            for ext, _ in category_extensions:
                all_extensions.append(ext.lower())
        
        # Load current settings (this would load from your saved settings)
        current_settings = self.get_extension_settings(mode)
        
        # Extensions that should be included in each search level
        base_extensions = ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.log', 
                        '.docx', '.doc', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls', '.rtf', '.odt', '.ods', '.odp',
                        '.csv','.eml', '.msg', '.emlx']
                        
        advanced_extensions = base_extensions + ['.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', 
                                            '.vbs', '.js', '.config', '.ini', '.reg',
                                            '.pst', '.ost', '.mbox']
        
        # Per modalità profonda, usa tutte le estensioni definite
        if mode == "profonda":
            # Sovrascrive le impostazioni correnti per far sì che tutte le estensioni siano selezionate
            current_settings = all_extensions
        
        # Dizionario per tenere traccia dei pulsanti per categoria
        category_buttons = {}
        
        # Create tabs for each category
        for category, extensions in extension_categories.items():
            # Create a frame for this category
            category_frame = ttk.Frame(notebook, padding=10)
            notebook.add(category_frame, text=category)
            
            # Aggiungi pulsante "Seleziona tutto" in alto
            button_frame = ttk.Frame(category_frame)
            button_frame.pack(fill=X, pady=(0, 10))
            
            # Funzione per selezionare tutte le estensioni in questa categoria
            def select_all_category(cat=category):
                for ext, desc in extension_categories[cat]:
                    if ext.lower() in ext_vars:
                        ext_vars[ext.lower()].set(True)
            
            # Funzione per deselezionare tutte le estensioni in questa categoria
            def deselect_all_category(cat=category):
                for ext, desc in extension_categories[cat]:
                    if ext.lower() in ext_vars:
                        ext_vars[ext.lower()].set(False)
            
            # Aggiungi i pulsanti
            select_all_btn = ttk.Button(button_frame, text=f"Seleziona tutto", 
                                    command=select_all_category)
            select_all_btn.pack(side=LEFT, padx=(0, 5))
            
            deselect_all_btn = ttk.Button(button_frame, text=f"Deseleziona tutti", 
                                        command=deselect_all_category)
            deselect_all_btn.pack(side=LEFT)
            
            # Salva i pulsanti nel dizionario per riferimento futuro
            category_buttons[category] = (select_all_btn, deselect_all_btn)
            
            # Create a grid layout
            content_frame = ttk.Frame(category_frame)
            content_frame.pack(fill=BOTH, expand=YES)
            row, col = 0, 0
            
            # Add checkboxes for each extension in this category
            for ext, desc in extensions:
                # Check if this extension should be selected based on mode
                is_selected = ext.lower() in current_settings
                
                # Create a variable for this checkbox
                var = BooleanVar(value=is_selected)
                ext_vars[ext.lower()] = var
                
                # Create the checkbox
                cb = ttk.Checkbutton(content_frame, text=f"{ext} ({desc})", variable=var)
                cb.grid(row=row, column=col, sticky=W, padx=5, pady=3)
                
                # Update grid position
                col += 1
                if col > 2:  # 3 columns per row
                    col = 0
                    row += 1
        
        # Button frame
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        # Helper functions
        def select_all():
            for var in ext_vars.values():
                var.set(True)
                
        def deselect_all():
            for var in ext_vars.values():
                var.set(False)
                
        def restore_defaults():
            # Modified restore_defaults function
            if mode == "profonda":
                # For deep mode, select ALL extensions
                for var in ext_vars.values():
                    var.set(True)
            else:
                # Reset to default extensions for base/avanzata modes
                default_list = base_extensions if mode == "base" else advanced_extensions
                for ext, var in ext_vars.items():
                    var.set(ext in default_list)
        
        def save_settings():
            # Save the selected extensions
            selected_extensions = [ext for ext, var in ext_vars.items() if var.get()]
            self.save_extension_settings(mode, selected_extensions)
            dialog.destroy()
            
            # Update the search_depth combo if needed
            current_depth = self.search_depth.get()
            if current_depth == mode:
                # Refresh the UI to reflect the new settings
                self.log_debug(f"Aggiornate le estensioni per la modalità {mode} ({len(selected_extensions)} estensioni)")
        
        # Buttons
        ttk.Button(btn_frame, text="Seleziona tutti", command=select_all).pack(side=LEFT, padx=5)
        ttk.Button(btn_frame, text="Deseleziona tutti", command=deselect_all).pack(side=LEFT, padx=5)
        ttk.Button(btn_frame, text="Ripristina default", command=restore_defaults).pack(side=LEFT, padx=5)
        
        ttk.Button(btn_frame, text="Annulla", command=dialog.destroy).pack(side=RIGHT, padx=5)
        ttk.Button(btn_frame, text="Salva", command=save_settings).pack(side=RIGHT, padx=5)
        
        # Center the dialog on screen
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Initial focus
        notebook.select(0)  # Focus first tab
    
    @error_handler
    def get_default_extensions(self, mode="base"):
        """Get default extensions for the specified search mode"""
        if mode == "base":
            return [
                # Documenti di testo
                '.txt', '.md', '.html', '.htm', '.xml', '.json', '.log',
                # Documenti Office
                '.docx', '.doc', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls',
                # Formati OpenDocument
                '.rtf', '.odt', '.ods', '.odp', 
                # Fogli di calcolo
                '.csv',
                # Formati email di base
                '.eml', '.msg', '.emlx'
            ]
        elif mode == "avanzata":
            # Il resto del codice rimane uguale
            base_exts = self.get_default_extensions("base")
            advanced_only = [
                # Estensioni avanzate...
                '.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', '.vbs', 
                '.config', '.ini', '.reg',
                '.py', '.java', '.php', '.cs', '.cpp', '.c', '.h', '.rb', '.js',
                '.db', '.sqlite', '.sqlite3',
                '.pst', '.ost', '.mbox', '.mbx', '.dbx',
                '.env', '.yml', '.yaml', '.toml', '.json5',
                '.bak', '.old', '.tmp', '.temp',
                '.epub', '.tex', '.rst',
                '.css', '.less', '.scss', '.jsp', '.asp', '.aspx',
                '.dot', '.dotx', '.xlt', '.xltx', '.pot', '.potx', '.ppsx'
            ]
            return base_exts + [ext for ext in advanced_only if ext not in base_exts]
        else:  # profonda
            return []  # Il valore effettivo viene determinato in configure_extensions
    
    @error_handler
    def get_extension_settings(self, mode="base"):
        """Load saved extension settings for the specified search mode"""
        if not hasattr(self, 'extension_settings'):
            # Initialize with defaults
            self.extension_settings = {
                "base": self.get_default_extensions("base"),
                "avanzata": self.get_default_extensions("avanzata"),
                "profonda": self.get_default_extensions("profonda")
            }
        
        extensions = self.extension_settings.get(mode)
        
        # Se non abbiamo estensioni per questa modalità o la lista è vuota, 
        # usa i valori predefiniti
        if extensions is None or len(extensions) == 0:
            extensions = self.get_default_extensions(mode)
            # Salva queste estensioni predefinite per usi futuri
            if not hasattr(self, 'extension_settings'):
                self.extension_settings = {}
            self.extension_settings[mode] = extensions
        
        self.log_debug(f"Estensioni caricate per modalità {mode}: {', '.join(extensions)}")
        return extensions
        
    @error_handler
    def save_extension_settings(self, mode, extensions):
        """Save extension settings for the specified search mode"""
        if not hasattr(self, 'extension_settings'):
            self.extension_settings = {}
        
        # Ensure all extensions have a leading dot and are lowercase
        normalized_extensions = []
        for ext in extensions:
            if not ext.startswith('.'):
                ext = '.' + ext
            normalized_extensions.append(ext.lower())
        
        self.extension_settings[mode] = normalized_extensions
        
        # Log what was saved
        self.log_debug(f"Saved {len(normalized_extensions)} extensions for {mode} mode")
        self.log_debug(f"Extensions: {', '.join(normalized_extensions)}")
        
        # Forza un aggiornamento dell'interfaccia se necessario
        if hasattr(self, 'search_depth') and self.search_depth.get() == mode:
            self.log_debug(f"Aggiornata UI per modalità {mode}")
            
        # Qui potresti aggiungere codice per salvare le impostazioni su file
        self.save_settings_to_file()
    
    @error_handler
    def save_settings_to_file(self):
        """Salva le impostazioni dell'applicazione su un file JSON"""
        try:
            import json
            import os

            # Cartella per le impostazioni
            settings_dir = os.path.join(os.path.expanduser("~"), ".file_search_tool")
            if not os.path.exists(settings_dir):
                os.makedirs(settings_dir)
            
            # File per salvare le impostazioni
            settings_file = os.path.join(settings_dir, "application_settings.json")
            
            # Creazione del dizionario delle impostazioni
            settings = {
                # Impostazioni esistenti
                "extension_settings": getattr(self, 'extension_settings', {}),
                "dir_size_calculation": self.dir_size_calculation.get(),
                "timeout_enabled": self.timeout_enabled.get(),
                "timeout_seconds": self.timeout_seconds.get(),
                "max_files_to_check": self.max_files_to_check.get(),
                "max_results": self.max_results.get(),
                "worker_threads": self.worker_threads.get(),
                "max_file_size_mb": self.max_file_size_mb.get(),
                "use_indexing": self.use_indexing.get(),
                "skip_permission_errors": self.skip_permission_errors.get(),
                
                # Nuove impostazioni per la gestione della RAM
                "auto_memory_management": getattr(self, 'auto_memory_management', True),
                "memory_usage_percent": getattr(self, 'memory_usage_percent', 75)
            }
            
            # Salva le impostazioni su file
            with open(settings_file, 'w') as f:
                json.dump(settings, f, indent=4)
            
            self.log_debug(f"Impostazioni dell'applicazione salvate in {settings_file}")
        except Exception as e:
            self.log_error("Errore durante il salvataggio delle impostazioni", e)
    
    @error_handler
    def load_settings_from_file(self):
        """Carica le impostazioni dell'applicazione da un file JSON"""
        try:
            import json
            import os

            # File per le impostazioni
            settings_file = os.path.join(os.path.expanduser("~"), ".file_search_tool", "application_settings.json")

            # Controlla se il file esiste
            if os.path.exists(settings_file):
                with open(settings_file, 'r') as f:
                    settings = json.load(f)
                    self.log_debug(f"Impostazioni dell'applicazione caricate da {settings_file}")

                    # Carica le impostazioni esistenti
                    self.extension_settings = settings.get("extension_settings", {
                        "base": self.get_default_extensions("base"),
                        "avanzata": self.get_default_extensions("avanzata"),
                        "profonda": self.get_default_extensions("profonda")
                    })
                    self.dir_size_calculation.set(settings.get("dir_size_calculation", "disabilitato"))
                    self.timeout_enabled.set(settings.get("timeout_enabled", False))
                    self.timeout_seconds.set(settings.get("timeout_seconds", 3600))
                    self.max_files_to_check.set(settings.get("max_files_to_check", 100000))
                    self.max_results.set(settings.get("max_results", 50000))
                    self.worker_threads.set(settings.get("worker_threads", min(8, os.cpu_count() or 4)))
                    self.max_file_size_mb.set(settings.get("max_file_size_mb", 100))
                    self.use_indexing.set(settings.get("use_indexing", True))
                    self.skip_permission_errors.set(settings.get("skip_permission_errors", True))

                    # Carica le impostazioni per la gestione della memoria
                    self.auto_memory_management = settings.get("auto_memory_management", True)
                    self.memory_usage_percent = settings.get("memory_usage_percent", 75)
            else:
                # Inizializza con valori predefiniti se il file non esiste
                self.log_debug(f"File delle impostazioni non trovato: {settings_file}. Inizializzo con i valori predefiniti.")
                self.extension_settings = {
                    "base": self.get_default_extensions("base"),
                    "avanzata": self.get_default_extensions("avanzata"),
                    "profonda": self.get_default_extensions("profonda")
                }
                self.dir_size_calculation.set("disabilitato")
                self.timeout_enabled.set(False)
                self.timeout_seconds.set(3600)
                self.max_files_to_check.set(100000)
                self.max_results.set(50000)
                self.worker_threads.set(min(8, os.cpu_count() or 4))
                self.max_file_size_mb.set(100)
                self.use_indexing.set(True)
                self.skip_permission_errors.set(True)
                self.auto_memory_management = True
                self.memory_usage_percent = 75
        except Exception as e:
            self.log_error("Errore durante il caricamento delle impostazioni", e)
            # Fallback ai valori predefiniti in caso di errore
            self.extension_settings = {
                "base": self.get_default_extensions("base"),
                "avanzata": self.get_default_extensions("avanzata"),
                "profonda": self.get_default_extensions("profonda")
            }
            self.dir_size_calculation.set("disabilitato")
            self.timeout_enabled.set(False)
            self.timeout_seconds.set(3600)
            self.max_files_to_check.set(100000)
            self.max_results.set(50000)
            self.worker_threads.set(min(8, os.cpu_count() or 4))
            self.max_file_size_mb.set(100)
            self.use_indexing.set(True)
            self.skip_permission_errors.set(True)
            self.auto_memory_management = True
            self.memory_usage_percent = 75

    @error_handler
    def create_widgets(self):
        # Frame principale che conterrà tutto
        main_container = ttk.Frame(self.root)
        main_container.pack(fill=BOTH, expand=YES)
        
        # Intestazione (titolo e informazioni)
        header_frame = ttk.Frame(main_container, padding="10")
        header_frame.pack(fill=X)

        # Layout a tre colonne in una singola riga
        # 1. Titolo e icona a sinistra
        title_frame = ttk.Frame(header_frame)
        title_frame.pack(side=LEFT, fill=Y)

        try:
            # Usa PIL per il supporto di più formati e il ridimensionamento
            from PIL import Image, ImageTk
            
            # Sostituisci 'logo.png' con il percorso della tua immagine
            current_dir = os.path.dirname(os.path.abspath(__file__))
            image_path = os.path.join(current_dir, "logo.png")
            
            # Carica e ridimensiona l'immagine
            original_image = Image.open(image_path)
            resized_image = original_image.resize((84, 84))
            self.logo_image = ImageTk.PhotoImage(resized_image)
            
            # Crea un label per l'immagine
            logo_label = ttk.Label(title_frame, image=self.logo_image)
            logo_label.pack(side=LEFT, padx=(0, 10))
        except Exception as e:
            self.log_debug(f"Impossibile caricare l'immagine del logo: {str(e)}")
            logo_label = None

        # Container per titolo e nome impilati verticalmente
        title_text_container = ttk.Frame(title_frame)
        title_text_container.pack(side=LEFT)

        title_label = ttk.Label(title_text_container, text="File Search Tool.. Forensics G.di F.", 
                            font=("Helvetica", 14, "bold"))
        title_label.pack(anchor=W)

        # Aggiunta del testo "Antonino" sotto il titolo
        antonino_label = ttk.Label(title_text_container, text="APS.QS Antonino Tessio", 
                            font=("Helvetica", 12))
        antonino_label.pack(anchor=W)

        # 2. Tema al centro
        theme_frame = ttk.Frame(header_frame)
        theme_frame.pack(side=LEFT, expand=True, fill=Y)

        ttk.Label(theme_frame, text="Tema:").pack(side=LEFT)
        themes = ttk.Style().theme_names()
        self.theme_combobox = ttk.Combobox(theme_frame, values=themes, width=15)
        self.theme_combobox.pack(side=LEFT, padx=5)
        self.theme_combobox.current(themes.index("darkly"))
        self.theme_combobox.bind("<<ComboboxSelected>>", lambda e: [ttk.Style().theme_use(self.theme_combobox.get()),self.update_theme_colors()])

        # 3. Data/ora e utente a destra - rimane invariato
        datetime_frame = ttk.Frame(header_frame)
        datetime_frame.pack(side=RIGHT, fill=Y)

        datetime_label = ttk.Label(datetime_frame, textvariable=self.datetime_var, font=("Helvetica", 9))
        datetime_label.pack(side=RIGHT)

        # ==========================================================
        # SEZIONE DEI CONTROLLI DI RICERCA (organizzati per righe)
        # ==========================================================
        controls_frame = ttk.LabelFrame(main_container, text="Parametri di ricerca", padding=10)
        controls_frame.pack(fill=X, padx=10, pady=5)
        
        # ------------------------------------------------------
        # RIGA 1: Directory di ricerca
        # ------------------------------------------------------
        path_frame = ttk.Frame(controls_frame)
        path_frame.pack(fill=X, pady=5)
        
        path_label = ttk.Label(path_frame, text="Directory:", width=12, anchor=W)
        path_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(path_label, "Seleziona la directory per effettuare la ricerca dei file")
        
        self.path_entry = ttk.Entry(path_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        self.browse_btn = ttk.Button(path_frame, text="Sfoglia", command=self.browse_directory, width=10)
        self.browse_btn.pack(side=LEFT)
        
        # ------------------------------------------------------
        # RIGA 2: Parole chiave
        # ------------------------------------------------------
        keyword_frame = ttk.Frame(controls_frame)
        keyword_frame.pack(fill=X, pady=5)
        
        keyword_label = ttk.Label(keyword_frame, text="Parole chiave:", width=12, anchor=W)
        keyword_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(keyword_label, "Per la ricerca di più parole usa la virgola. Esempio: documento, fattura, contratto\n"
                                        "Attiva 'Parola intera' per cercare 'log' senza trovare 'login' o 'logo'")
        
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)

        # ------------------------------------------------------
        # RIGA 3: Opzioni base di ricerca (checkbox)
        # ------------------------------------------------------
        options_frame = ttk.Frame(controls_frame)
        options_frame.pack(fill=X, pady=5)

        # Prima parte: Livelli di ricerca
        search_options = ttk.Frame(options_frame)
        search_options.pack(side=LEFT, fill=Y)

        ttk.Label(search_options, text="Livelli di ricerca:").pack(side=LEFT, padx=(0, 5))
        search_depth_combo = ttk.Combobox(search_options, textvariable=self.search_depth, 
                                        values=["base", "avanzata", "profonda"], 
                                        width=10, state="readonly")
        search_depth_combo.pack(side=LEFT, padx=5)
        search_depth_combo.current(0)

        extensions_btn = ttk.Button(search_options, text="Configura estensioni", 
                            command=lambda: self.configure_extensions(self.search_depth.get()))
        extensions_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(extensions_btn, "Configura quali estensioni di file includere nella ricerca")

        # Pulsante impostazioni avanzate
        advanced_options_btn = ttk.Button(search_options, text="Impostazioni avanzate", 
                                        command=self.show_advanced_options)
        advanced_options_btn.pack(side=LEFT, padx=10)
        self.create_tooltip(advanced_options_btn, "Configura tutte le impostazioni avanzate (profondità, filtri, esclusioni, performance)")

        # Separatore per creare più spazio
        separator = ttk.Frame(options_frame, width=30)
        separator.pack(side=LEFT)

        # Pulsanti di azione (spostati qui dopo lo spazio)
        action_buttons = ttk.Frame(options_frame)
        action_buttons.pack(side=LEFT, fill=Y)

        # Pulsante di ricerca (principale)
        self.search_button = ttk.Button(action_buttons, text="CERCA", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=15)
        self.search_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.search_button, "Avvia la ricerca con i criteri specificati")

        # Pulsante per interrompere la ricerca
        self.stop_button = ttk.Button(action_buttons, text="Interrompi ricerca",
                                    command=self.stop_search_process,
                                    style="danger.TButton", width=20,
                                    state="disabled")
        self.stop_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.stop_button, "Ferma immediatamente la ricerca in corso")

        # Pulsante per pulire i campi di ricerca
        self.clear_btn = ttk.Button(action_buttons, text="Pulisci campi", 
                    command=lambda: [self.search_path.set(""), self.keywords.set("")],
                    style="secondary.Outline.TButton", width=15)
        self.clear_btn.pack(side=LEFT, padx=10)
        self.create_tooltip(self.clear_btn, "Cancella i campi di ricerca")

        # Pulsante admin solo su Windows
        if os.name == 'nt':
            self.admin_button = ttk.Button(action_buttons, text="Avvia come Admin", 
                                    command=self.restart_as_admin,
                                    style="info.Outline.TButton", width=20)
            self.admin_button.pack(side=LEFT, padx=10)
            
            # Disabilita il pulsante se l'app è già avviata come amministratore
            if self.is_admin:
                self.admin_button.config(state="disabled")
                self.create_tooltip(self.admin_button, "L'applicazione è già in esecuzione come amministratore")
            else:
                self.create_tooltip(self.admin_button, "Riavvia l'applicazione con privilegi di amministratore per accedere a tutte le cartelle")
        
        # ==========================================================
        # SEZIONE INFORMAZIONI TEMPORALI E SPAZIO DISCO (SOPRA LA LISTA RISULTATI)
        # ==========================================================
        info_bar = ttk.Frame(main_container)
        info_bar.pack(fill=X, padx=10, pady=5)
        
        # Frame per le informazioni temporali a sinistra
        time_frame = ttk.LabelFrame(info_bar, text="Informazioni temporali", padding=5)
        time_frame.pack(side=LEFT, fill=X, padx=(0, 5), expand=YES)
        
        time_grid = ttk.Frame(time_frame)
        time_grid.pack(fill=X, pady=2)
        
        ttk.Label(time_grid, text="Avvio:").grid(row=0, column=0, sticky=W, padx=5)
        self.start_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.start_time_label.grid(row=0, column=1, sticky=W, padx=5)
        
        ttk.Label(time_grid, text="Fine:").grid(row=0, column=2, sticky=W, padx=5)
        self.end_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.end_time_label.grid(row=0, column=3, sticky=W, padx=5)
        
        ttk.Label(time_grid, text="Durata:").grid(row=0, column=4, sticky=W, padx=5)
        self.total_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.total_time_label.grid(row=0, column=5, sticky=W, padx=5)
        
        # Frame per info disco a destra
        disk_info_frame = ttk.LabelFrame(info_bar, text="Spazio disco", padding=5)
        disk_info_frame.pack(side=LEFT, fill=X, expand=YES)

        # Uso layout semplice con pack per ordinare gli elementi in linea
        disk_grid = ttk.Frame(disk_info_frame)
        disk_grid.pack(fill=X, pady=2)

        # 1. Usato
        ttk.Label(disk_grid, text="Usato:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.used_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 2. Libero
        ttk.Label(disk_grid, text="Libero:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.free_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 3. Totale
        ttk.Label(disk_grid, text="Totale:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.total_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 4. Directory
        ttk.Label(disk_grid, text="Directory:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.dir_size_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 5. Pulsante aggiorna
        refresh_size_btn = ttk.Button(disk_grid, text="Aggiorna", command=self.refresh_directory_size, 
                                width=8, style="info.TButton")
        refresh_size_btn.pack(side=LEFT, padx=5)
        
        # ==========================================================
        # SEZIONE STATO RICERCA
        # ==========================================================
        # Contenitore per le informazioni di stato sopra i risultati
        status_frame = ttk.LabelFrame(main_container, text="Stato ricerca", padding=5)
        status_frame.pack(fill=X, padx=10, pady=5)

        # Progress bar
        self.progress_bar = ttk.Progressbar(status_frame, mode='determinate')
        self.progress_bar.pack(fill=X, pady=5)

        # Status grid with improved layout
        status_grid = ttk.Frame(status_frame)
        status_grid.pack(fill=X)

        # Row 1: Status with debug button
        status_row = ttk.Frame(status_grid)
        status_row.pack(fill=X, pady=2)

        ttk.Label(status_row, text="Analisi:", width=12, anchor=W, font=("", 9, "bold")).pack(side=LEFT, padx=5)
        self.status_label = ttk.Label(status_row, text="In attesa...", wraplength=1000)
        self.status_label.pack(side=LEFT, fill=X, expand=YES, padx=5)

        # Add Debug Button on the right side of the status row
        self.debug_button = ttk.Button(status_row, text="Debug Log", command=self.show_debug_log, 
                                    style="info.Outline.TButton", width=12)
        self.debug_button.pack(side=RIGHT, padx=5)
        self.create_tooltip(self.debug_button, "Mostra la finestra di debug con log dettagliati sulla ricerca in corso")

        # Row 2: File analyzed
        files_row = ttk.Frame(status_grid)
        files_row.pack(fill=X, pady=2)

        ttk.Label(files_row, text="File analizzati:", width=12, anchor=W, font=("", 9, "bold")).pack(side=LEFT, padx=5)
        self.analyzed_files_label = ttk.Label(files_row, text="Nessuna ricerca avviata", wraplength=2000)
        self.analyzed_files_label.pack(side=LEFT, fill=X, expand=YES, padx=5)
        # ==========================================================
        # SEZIONE RISULTATI 
        # ==========================================================
        results_container = ttk.LabelFrame(main_container, text="Risultati di ricerca", padding=10)
        results_container.pack(fill=BOTH, expand=YES, padx=10, pady=(5, 10))
        
        # Frame per i pulsanti di azione sui risultati
        actions_frame = ttk.Frame(results_container)
        actions_frame.pack(fill=X, pady=(0, 5))

        # Pulsanti per la selezione
        selection_frame = ttk.Frame(actions_frame)
        selection_frame.pack(side=LEFT)

        select_all_btn = ttk.Button(selection_frame, text="Seleziona tutto", command=self.select_all)
        select_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(select_all_btn, "Seleziona tutti i risultati nella lista")

        deselect_all_btn = ttk.Button(selection_frame, text="Deseleziona tutto", command=self.deselect_all)
        deselect_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(deselect_all_btn, "Deseleziona tutti i risultati")

        # Crea frame centrale che si espande per riempire lo spazio disponibile
        center_frame = ttk.Frame(actions_frame)
        center_frame.pack(side=LEFT, fill=X, expand=YES)

        # Creiamo un sotto-frame per contenere entrambi i label affiancati
        labels_frame = ttk.Frame(center_frame)
        labels_frame.pack(anchor=CENTER)

        # Label per la dimensione totale (esistente)
        self.total_files_size_label = ttk.Label(labels_frame, text="Dimensione totale: 0  (0 file)", font=("", 9, "bold"))
        self.total_files_size_label.pack(side=LEFT, padx=10)

        # NUOVO: Label per la dimensione dei file selezionati
        self.selected_files_size_label = ttk.Label(labels_frame, text="Selezionati: 0  (0 file)", font=("", 9, "bold"))
        self.selected_files_size_label.pack(side=LEFT, padx=10)

        # Pulsanti per le azioni
        action_frame = ttk.Frame(actions_frame)
        action_frame.pack(side=RIGHT)

        self.copy_button = ttk.Button(action_frame, text="Copia selezionati",
                                    command=self.copy_selected,
                                    style="TButton")
        self.copy_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.copy_button, "Copia i file selezionati nella directory specificata")

        self.compress_button = ttk.Button(action_frame, text="Comprimi selezionati",
                                        command=self.compress_selected,
                                        style="TButton")
        self.compress_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.compress_button, "Comprimi i file selezionati in un archivio ZIP")

        self.view_log_button = ttk.Button(action_frame, text="Visualizza file esclusi",
                                        command=self.view_skipped_files_log,
                                        style="secondary.TButton")
        self.view_log_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.view_log_button, "Visualizza il log dei file esclusi dalla ricerca")

        # TreeView con scrollbar
        treeview_container = ttk.Frame(results_container)
        treeview_container.pack(fill=BOTH, expand=True)

        # Creazione della TreeView
        self.results_list = ttk.Treeview(treeview_container, selectmode="extended",
                                columns=("type", "attachment", "size", "modified", "created", "author", "path"), 
                                show="headings")

        # Creazione delle scrollbar
        vsb = ttk.Scrollbar(treeview_container, orient="vertical", command=self.results_list.yview)
        hsb = ttk.Scrollbar(treeview_container, orient="horizontal", command=self.results_list.xview)

        # Configurazione delle scrollbar nella TreeView
        self.results_list.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

        # Posizionamento con grid
        self.results_list.grid(column=0, row=0, sticky='nsew')
        vsb.grid(column=1, row=0, sticky='ns')
        hsb.grid(column=0, row=1, sticky='ew')

        # Configura grid layout per l'espansione
        treeview_container.grid_columnconfigure(0, weight=1)
        treeview_container.grid_rowconfigure(0, weight=1)

        # Impostazione delle intestazioni delle colonne
        self.results_list.heading("type", text="Tipo")
        self.results_list.heading("attachment", text="Allegato")
        self.results_list.heading("size", text="Dimensione")
        self.results_list.heading("modified", text="Modificato")
        self.results_list.heading("created", text="Creato")
        self.results_list.heading("author", text="Nome")
        self.results_list.heading("path", text="Percorso")

        # Impostazione delle intestazioni delle colonne con funzione di ordinamento
        self.results_list.heading("type", text="Tipo", 
                            command=lambda: self.treeview_sort_column(self.results_list, "type", False))
        self.results_list.heading("size", text="Dimensione", 
                            command=lambda: self.treeview_sort_column(self.results_list, "size", False))
        self.results_list.heading("modified", text="Modificato", 
                            command=lambda: self.treeview_sort_column(self.results_list, "modified", False))
        self.results_list.heading("created", text="Creato", 
                            command=lambda: self.treeview_sort_column(self.results_list, "created", False))
        self.results_list.heading("author", text="Nome", 
                            command=lambda: self.treeview_sort_column(self.results_list, "author", False))
        self.results_list.heading("path", text="Percorso", 
                            command=lambda: self.treeview_sort_column(self.results_list, "path", False))

        # Impostazione delle larghezze fisse delle colonne
        self.results_list.column("type", width=120, minwidth=50, stretch=NO, anchor="center")
        self.results_list.column("attachment", width=70, minwidth=30, stretch=NO, anchor="center")
        self.results_list.column("size", width=100, minwidth=80, stretch=NO, anchor="center")
        self.results_list.column("modified", width=150, minwidth=120, stretch=NO, anchor="center")
        self.results_list.column("created", width=150, minwidth=120, stretch=NO, anchor="center")
        self.results_list.column("author", width=400, minwidth=80, stretch=NO, anchor="w")
        self.results_list.column("path", width=600, minwidth=200, stretch=YES, anchor="w")

        # Aggiungi binding per l'evento di doppio clic
        self.results_list.bind("<Double-1>", self.open_file_location)

        # Aggiungi binding per l'evento di selezione
        self.results_list.bind("<<TreeviewSelect>>", self.update_selected_files_size)

        # Applica stili alle righe
        self.update_theme_colors()

        # Auto-focus sull'entry del percorso all'avvio
        self.path_entry.focus_set()

    @error_handler
    def show_advanced_options(self):
        """Mostra una finestra di dialogo unificata per tutte le opzioni avanzate"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Impostazioni avanzate")
        dialog.geometry("800x650")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Crea un notebook con schede per le diverse categorie di opzioni
        notebook = ttk.Notebook(main_frame)
        notebook.pack(fill=BOTH, expand=YES, pady=10)
        
        # ================= Scheda 1: Opzioni di ricerca =================
        search_options_frame = ttk.Frame(notebook, padding=15)
        notebook.add(search_options_frame, text="Opzioni di ricerca")
        
        # Profondità di ricerca
        depth_frame = ttk.LabelFrame(search_options_frame, text="Profondità di ricerca", padding=10)
        depth_frame.pack(fill=X, pady=10)
        
        depth_desc = ttk.Label(depth_frame, 
                        text="Imposta il numero massimo di livelli di cartella da esplorare.\n"
                            "Il valore 0 indica profondità illimitata.",
                        wraplength=700)
        depth_desc.pack(anchor=W, pady=(0, 10))
        
        depth_control = ttk.Frame(depth_frame)
        depth_control.pack(fill=X)
        
        ttk.Label(depth_control, text="Profondità:").pack(side=LEFT)
        
        # Usa variabili per i controlli
        depth_var = IntVar(value=self.max_depth)
        depth_spinbox = ttk.Spinbox(depth_control, from_=0, to=20, width=3, textvariable=depth_var)
        depth_spinbox.pack(side=LEFT, padx=5)
        ttk.Label(depth_control, text="(0 = illimitata)", foreground="gray").pack(side=LEFT)
        
        # AGGIUNTO TOOLTIP
        self.create_tooltip(depth_spinbox, 
                    "Controlla quanto a fondo cercare nelle sottocartelle.\n"
                    "Un valore più alto aumenta il tempo di ricerca ma trova file in cartelle più profonde.\n"
                    "Il valore 0 cerca in tutte le sottocartelle senza limiti di profondità.")
        
        # Contenuti da cercare
        content_frame = ttk.LabelFrame(search_options_frame, text="Contenuti da cercare", padding=10)
        content_frame.pack(fill=X, pady=10)
        
        search_files_var = BooleanVar(value=self.search_files.get())
        search_files_cb = ttk.Checkbutton(content_frame, text="Cerca nei file", variable=search_files_var)
        search_files_cb.pack(anchor=W, pady=2)
        self.create_tooltip(search_files_cb, "Attiva per includere i file nei risultati di ricerca")
        
        search_folders_var = BooleanVar(value=self.search_folders.get())
        search_folders_cb = ttk.Checkbutton(content_frame, text="Cerca nelle cartelle", variable=search_folders_var)
        search_folders_cb.pack(anchor=W, pady=2)
        self.create_tooltip(search_folders_cb, "Attiva per includere le cartelle nei risultati di ricerca")
        
        search_content_var = BooleanVar(value=self.search_content.get())
        search_content_cb = ttk.Checkbutton(content_frame, text="Cerca nei contenuti dei file", variable=search_content_var)
        search_content_cb.pack(anchor=W, pady=2)
        self.create_tooltip(search_content_cb, 
                        "Attiva per cercare le parole chiave all'interno dei file.\n"
                        "Questo rallenta la ricerca ma permette di trovare i file\n"
                        "che contengono il testo cercato anche se non è nel nome.")
        
        whole_word_var = BooleanVar(value=self.whole_word_search.get())
        whole_word_cb = ttk.Checkbutton(content_frame, text="Cerca parole intere", variable=whole_word_var)
        whole_word_cb.pack(anchor=W, pady=2)
        self.create_tooltip(whole_word_cb, 
                        "Quando attivato, cerca solo corrispondenze esatte delle parole.\n"
                        "Esempio: cercando 'log' NON troverà 'login' o 'catalogo'.")
        
        # ================= Scheda 2: Filtri avanzati =================
        filters_frame = ttk.Frame(notebook, padding=15)
        notebook.add(filters_frame, text="Filtri avanzati")
        
        # Filtri dimensione
        size_frame = ttk.LabelFrame(filters_frame, text="Dimensione file", padding=10)
        size_frame.pack(fill=X, pady=10)
        
        ttk.Label(size_frame, text="Filtra i file in base alla dimensione:").pack(anchor=W, pady=(0, 10))
        
        size_grid = ttk.Frame(size_frame)
        size_grid.pack(fill=X)
        
        size_min_label = ttk.Label(size_grid, text="Dimensione minima (KB):")
        size_min_label.grid(row=0, column=0, padx=5, pady=5, sticky=W)
        min_size_var = StringVar(value=str(self.advanced_filters["size_min"] // 1024))
        min_size = ttk.Entry(size_grid, width=10, textvariable=min_size_var)
        min_size.grid(row=0, column=1, padx=5, pady=5)
        self.create_tooltip(min_size, 
                    "Filtra i file più piccoli della dimensione specificata in KB.\n"
                    "Esempio: impostando 100 verranno ignorati i file minori di 100KB.")
        
        size_max_label = ttk.Label(size_grid, text="Dimensione massima (KB):")
        size_max_label.grid(row=1, column=0, padx=5, pady=5, sticky=W)
        max_size_var = StringVar(value=str(self.advanced_filters["size_max"] // 1024 if self.advanced_filters["size_max"] else 0))
        max_size = ttk.Entry(size_grid, width=10, textvariable=max_size_var)
        max_size.grid(row=1, column=1, padx=5, pady=5)
        self.create_tooltip(max_size, 
                    "Filtra i file più grandi della dimensione specificata in KB.\n"
                    "Esempio: impostando 1000 verranno ignorati i file maggiori di 1000KB (1MB).\n"
                    "Impostare a 0 per nessun limite.")
        
        # Filtri data
        date_frame = ttk.LabelFrame(filters_frame, text="Data di modifica", padding=10)
        date_frame.pack(fill=X, pady=10)
        
        ttk.Label(date_frame, text="Filtra i file in base alla data di modifica:").pack(anchor=W, pady=(0, 10))
        
        date_grid = ttk.Frame(date_frame)
        date_grid.pack(fill=X)
        
        date_min_label = ttk.Label(date_grid, text="Data inizio (DD-MM-YYYY):")
        date_min_label.grid(row=0, column=0, padx=5, pady=5, sticky=W)
        min_date = ttk.DateEntry(date_grid, dateformat="%d-%m-%Y")
        min_date.grid(row=0, column=1, padx=5, pady=5)
        min_date.entry.delete(0, 'end')
        if self.advanced_filters["date_min"]:
            min_date.entry.insert(0, self.advanced_filters["date_min"])
        self.create_tooltip(min_date, 
                    "Includi solo file modificati dopo questa data.\n"
                    "Lascia vuoto per nessun limite di data minima.")
        
        date_max_label = ttk.Label(date_grid, text="Data fine (DD-MM-YYYY):")
        date_max_label.grid(row=1, column=0, padx=5, pady=5, sticky=W)
        max_date = ttk.DateEntry(date_grid, dateformat="%d-%m-%Y")
        max_date.grid(row=1, column=1, padx=5, pady=5)
        max_date.entry.delete(0, 'end')
        if self.advanced_filters["date_max"]:
            max_date.entry.insert(0, self.advanced_filters["date_max"])
        self.create_tooltip(max_date, 
                    "Includi solo file modificati prima di questa data.\n"
                    "Lascia vuoto per nessun limite di data massima.")
        
        # ================= Scheda 3: Gestione esclusioni =================
        exclusions_frame = ttk.Frame(notebook, padding=15)
        notebook.add(exclusions_frame, text="Esclusioni")

        ttk.Label(exclusions_frame, text="Aggiungi cartelle da escludere dalla ricerca:", wraplength=700).pack(anchor=W, pady=(0, 10))

        system_exclusions_frame = ttk.Frame(exclusions_frame)
        system_exclusions_frame.pack(fill=X, pady=5)

        # Definizione dei percorsi di sistema
        system_paths = [
            "C:/Windows", 
            "C:/Program Files", 
            "C:/Program Files (x86)",
            "C:/ProgramData",
            "C:/Users/All Users",
            "C:/Program Files (x86)/Client Active Directory Rights Management Services"
        ]

        system_exclusions_var = tk.BooleanVar(value=False)

        # Imposta il valore iniziale della checkbox in base alle esclusioni attuali
        if hasattr(self, 'excluded_paths') and all(path in self.excluded_paths for path in system_paths):
            system_exclusions_var.set(True)

        # Funzione per gestire il toggle delle esclusioni di sistema
        def toggle_system_exclusions():
            if not hasattr(self, 'excluded_paths'):
                self.excluded_paths = []
                
            if system_exclusions_var.get():
                # Aggiungi percorsi di sistema
                added = 0
                for path in system_paths:
                    if path not in self.excluded_paths:
                        self.excluded_paths.append(path)
                        added += 1
                
                # Aggiorna la lista visualizzata
                if added > 0:
                    for item in excluded_list.get_children():
                        excluded_list.delete(item)
                    for path in self.excluded_paths:
                        excluded_list.insert("", "end", values=(path,))
                    messagebox.showinfo("Esclusioni sistema", f"Aggiunti {added} percorsi di sistema alle esclusioni")
            else:
                # Rimuovi percorsi di sistema
                removed = 0
                for path in system_paths[:]:
                    if path in self.excluded_paths:
                        self.excluded_paths.remove(path)
                        removed += 1
                
                # Aggiorna la lista visualizzata
                if removed > 0:
                    for item in excluded_list.get_children():
                        excluded_list.delete(item)
                    for path in self.excluded_paths:
                        excluded_list.insert("", "end", values=(path,))
                    messagebox.showinfo("Esclusioni sistema", f"Rimossi {removed} percorsi di sistema dalle esclusioni")

        # Checkbox per esclusioni di sistema
        system_check = ttk.Checkbutton(system_exclusions_frame, text="Escludi cartelle di sistema (Windows, Program Files, ecc.)", 
            variable=system_exclusions_var, command=toggle_system_exclusions)
        system_check.pack(anchor=W)
        self.create_tooltip(system_check, 
                            "Quando selezionato, esclude automaticamente cartelle di sistema come:\n"
                            "- Windows\n- Program Files\n- ProgramData\n"
                            "Questo velocizzerà notevolmente le ricerche sui dischi di sistema.")

        # Lista dei percorsi esclusi
        list_frame = ttk.Frame(exclusions_frame)
        list_frame.pack(fill=BOTH, expand=YES, pady=5)  # Cambiato expand da NO a YES

        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side=RIGHT, fill=Y)

        # Aggiungi il parametro height per specificare il numero di righe visibili
        excluded_list = ttk.Treeview(list_frame, columns=("path",), show="headings", 
                                yscrollcommand=scrollbar.set, selectmode="extended", height=15)  # Aggiunto height=15
        excluded_list.heading("path", text="Percorso")
        excluded_list.column("path", width=450)
        excluded_list.pack(fill=BOTH, expand=YES)
        self.create_tooltip(excluded_list, 
                    "Elenco di cartelle escluse dalla ricerca.\n"
                    "Le sottocartelle di questi percorsi non verranno analizzate.")
        
        scrollbar.config(command=excluded_list.yview)
        
        # Aggiungi i percorsi esclusi alla lista
        if hasattr(self, 'excluded_paths'):
            for path in self.excluded_paths:
                excluded_list.insert("", "end", values=(path,))
        
        # Frame per aggiungere nuovi percorsi
        add_frame = ttk.Frame(exclusions_frame)
        add_frame.pack(fill=X, pady=10)
        
        path_var = StringVar()
        path_entry = ttk.Entry(add_frame, textvariable=path_var, width=50)
        path_entry.pack(side=LEFT, padx=(0, 5), fill=X, expand=YES)
        self.create_tooltip(path_entry, 
                    "Inserisci il percorso completo della cartella da escludere.\n"
                    "Esempio: C:/Windows o C:/Program Files")
        
        def browse_exclude():
            directory = filedialog.askdirectory()
            if directory:
                path_var.set(directory)
        
        browse_btn = ttk.Button(add_frame, text="Sfoglia", command=browse_exclude)
        browse_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(browse_btn, "Seleziona la cartella da escludere usando una finestra di dialogo")
        
        def add_exclusion():
            path = path_var.get().strip()
            if path:
                excluded_list.insert("", "end", values=(path,))
                path_var.set("")
        
        add_btn = ttk.Button(add_frame, text="Aggiungi", command=add_exclusion)
        add_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(add_btn, "Aggiungi il percorso specificato alla lista delle esclusioni")
        
        # Pulsante per rimuovere elementi selezionati
        def remove_selected():
            selected = excluded_list.selection()
            if selected:
                for item in selected:
                    excluded_list.delete(item)
        
        remove_btn = ttk.Button(exclusions_frame, text="Rimuovi selezionati", command=remove_selected)
        remove_btn.pack(anchor=W, pady=5)
        self.create_tooltip(remove_btn, "Rimuovi i percorsi selezionati dalla lista delle esclusioni")

        
        # ================= Scheda 4: Opzioni a blocchi =================
        blocks_frame = ttk.Frame(notebook, padding=15)
        notebook.add(blocks_frame, text="Opzioni a blocchi")
        
        desc_label = ttk.Label(blocks_frame, text="La ricerca a blocchi divide le cartelle in unità di lavoro più piccole per migliorare le prestazioni e la reattività.",
                        wraplength=700)
        desc_label.pack(fill=X, pady=(0, 15))
        
        # Dimensione blocchi
        size_frame = ttk.LabelFrame(blocks_frame, text="Dimensione e parallelismo", padding=10)
        size_frame.pack(fill=X, pady=10)
        
        size_grid = ttk.Frame(size_frame)
        size_grid.pack(fill=X)
        
        files_block_label = ttk.Label(size_grid, text="Max file per blocco:")
        files_block_label.grid(row=0, column=0, sticky=W, padx=5, pady=5)
        files_block_var = IntVar(value=self.max_files_per_block.get())
        files_block = ttk.Spinbox(size_grid, from_=100, to=10000, increment=100, width=7, textvariable=files_block_var)
        files_block.grid(row=0, column=1, padx=5, pady=5)
        self.create_tooltip(files_block, 
                    "Numero massimo di file da processare in un singolo blocco.\n"
                    "Valori più bassi aumentano la reattività dell'interfaccia ma\n"
                    "possono ridurre leggermente la velocità di ricerca complessiva.")
        
        parallel_label = ttk.Label(size_grid, text="Blocchi paralleli:")
        parallel_label.grid(row=0, column=2, sticky=W, padx=(20, 5), pady=5)
        parallel_var = IntVar(value=self.max_parallel_blocks.get())
        parallel = ttk.Spinbox(size_grid, from_=1, to=16, increment=1, width=5, textvariable=parallel_var)
        parallel.grid(row=0, column=3, padx=5, pady=5)
        self.create_tooltip(parallel, 
                    "Numero di blocchi da elaborare contemporaneamente.\n"
                    "Aumentare questo valore su sistemi con molte CPU può\n"
                    "velocizzare la ricerca, ma consuma più risorse.")
        
        # Opzioni aggiuntive
        options_frame = ttk.LabelFrame(blocks_frame, text="Ottimizzazioni", padding=10)
        options_frame.pack(fill=X, pady=10)
        
        auto_adjust_var = BooleanVar(value=self.block_size_auto_adjust.get())
        auto_adjust_cb = ttk.Checkbutton(options_frame, text="Adatta automaticamente la dimensione dei blocchi", 
                    variable=auto_adjust_var)
        auto_adjust_cb.pack(anchor=W, pady=2)
        self.create_tooltip(auto_adjust_cb, 
                    "Quando attivato, la dimensione dei blocchi viene adattata\n"
                    "automaticamente in base al tipo di ricerca e alla velocità del sistema.\n"
                    "Utile per bilanciare prestazioni e reattività.")
        
        prioritize_var = BooleanVar(value=self.prioritize_user_folders.get())
        prioritize_cb = ttk.Checkbutton(options_frame, text="Dare priorità alle cartelle utente", 
                    variable=prioritize_var)
        prioritize_cb.pack(anchor=W, pady=2)
        self.create_tooltip(prioritize_cb, 
                    "Analizza prima le cartelle più importanti come Documenti, Desktop,\n"
                    "per trovare rapidamente i file più rilevanti.\n"
                    "Utile quando si cercano file personali.")
        
        # ================= Scheda 5: Performance =================
        performance_frame = ttk.Frame(notebook, padding=15)
        notebook.add(performance_frame, text="Performance")

        # Timeout e limiti
        timeout_frame = ttk.LabelFrame(performance_frame, text="Timeout e limiti", padding=10)
        timeout_frame.pack(fill=X, pady=10)

        # Rimuoviamo il posizionamento separato del checkbox e lo inseriamo direttamente nella griglia
        timeout_grid = ttk.Frame(timeout_frame)
        timeout_grid.pack(fill=X, pady=5)

        # Riga 0: Checkbox e secondi nella stessa riga
        timeout_enabled_var = BooleanVar(value=self.timeout_enabled.get())
        timeout_check = ttk.Checkbutton(timeout_grid, text="Attiva timeout ricerca", variable=timeout_enabled_var)
        timeout_check.grid(row=0, column=0, sticky=W, padx=5, pady=2)
        self.create_tooltip(timeout_check, 
                    "Interrompe automaticamente la ricerca dopo il tempo specificato.\n"
                    "Utile per evitare ricerche troppo lunghe su grandi volumi di dati.")

        timeout_label = ttk.Label(timeout_grid, text="Secondi:")
        timeout_label.grid(row=0, column=1, sticky=W, padx=(55, 5), pady=2)
        timeout_seconds_var = IntVar(value=self.timeout_seconds.get())
        timeout_spin = ttk.Spinbox(timeout_grid, from_=10, to=3600, width=5, textvariable=timeout_seconds_var)
        timeout_spin.grid(row=0, column=2, padx=5, pady=2, sticky=W)
        self.create_tooltip(timeout_spin, 
                    "Durata massima della ricerca in secondi prima dell'interruzione automatica.\n"
                    "Es. 300 = 5 minuti, 3600 = 1 ora")
        
        max_files_label = ttk.Label(timeout_grid, text="Max file da controllare:")
        max_files_label.grid(row=1, column=0, sticky=W, padx=5, pady=5)
        max_files_var = IntVar(value=self.max_files_to_check.get())
        max_files = ttk.Spinbox(timeout_grid, from_=1000, to=10000000, width=8, textvariable=max_files_var)
        max_files.grid(row=1, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_files, 
                    "Numero massimo di file da controllare prima di terminare la ricerca.\n"
                    "Limita le ricerche molto estese per evitare tempi di attesa eccessivi.\n"
                    "Valori consigliati: 10000 per ricerche rapide, 100000 per ricerche approfondite.")
        
        max_results_label = ttk.Label(timeout_grid, text="Max risultati:")
        max_results_label.grid(row=1, column=2, sticky=W, padx=5, pady=5)
        max_results_var = IntVar(value=self.max_results.get())
        max_results = ttk.Spinbox(timeout_grid, from_=500, to=100000, width=8, textvariable=max_results_var)
        max_results.grid(row=1, column=3, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_results, 
                    "Numero massimo di risultati da mostrare.\n"
                    "Limita il numero di risultati per migliorare le prestazioni\n"
                    "dell'interfaccia in caso di ricerche con molte corrispondenze.")
        
        # Processamento
        process_frame = ttk.LabelFrame(performance_frame, text="Processamento", padding=10)
        process_frame.pack(fill=X, pady=10)
        
        process_grid = ttk.Frame(process_frame)
        process_grid.pack(fill=X)
        
        threads_label = ttk.Label(process_grid, text="Thread paralleli:")
        threads_label.grid(row=0, column=0, sticky=W, padx=5, pady=5)
        threads_var = IntVar(value=self.worker_threads.get())
        threads = ttk.Spinbox(process_grid, from_=1, to=16, width=3, textvariable=threads_var)
        threads.grid(row=0, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(threads, 
                    "Numero di thread paralleli per la ricerca.\n"
                    "Più thread accelerano la ricerca ma usano più CPU.\n"
                    "Consigliato: 4-8 thread su PC moderni.")
        
        max_size_label = ttk.Label(process_grid, text="Dimensione max file (MB):")
        max_size_label.grid(row=0, column=2, sticky=W, padx=5, pady=5)
        max_size_mb_var = IntVar(value=self.max_file_size_mb.get())
        max_size_mb = ttk.Spinbox(process_grid, from_=1, to=1000, width=5, textvariable=max_size_mb_var)
        max_size_mb.grid(row=0, column=3, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_size_mb, 
                    "Dimensione massima in MB dei file di cui analizzare il contenuto.\n"
                    "I file più grandi saranno considerati solo in base al nome.\n"
                    "Un valore più basso aumenta la velocità di ricerca.")
        
        # Calcolo dimensioni
        calc_frame = ttk.LabelFrame(performance_frame, text="Calcolo dimensioni", padding=10)
        calc_frame.pack(fill=X, pady=10)
        
        calc_label = ttk.Label(calc_frame, text="Modalità di calcolo:")
        calc_label.pack(side=LEFT, padx=5)
        dir_size_calc_var = StringVar(value=self.dir_size_calculation.get())
        calc_combo = ttk.Combobox(calc_frame, textvariable=dir_size_calc_var, 
                            values=["incrementale", "preciso", "stimato", "sistema", "disabilitato"], 
                            width=12, state="readonly")
        calc_combo.pack(side=LEFT, padx=5)
        self.create_tooltip(calc_combo, 
                    "Modalità di calcolo della dimensione delle directory:\n"
                    "• incrementale: aggiorna durante la ricerca\n"
                    "• preciso: calcolo completo ma più lento\n"
                    "• stimato: più veloce ma approssimato\n"
                    "• sistema: usa comandi di sistema esterni\n"
                    "• disabilitato: non calcolare le dimensioni")
        # Gestione memoria (Aggiunta)
        memory_frame = ttk.LabelFrame(performance_frame, text="Gestione memoria", padding=10)
        memory_frame.pack(fill=tk.X, pady=10)

        auto_memory_var = BooleanVar(value=self.auto_memory_management)
        memory_percent_var = IntVar(value=self.memory_usage_percent)

        auto_memory_check = ttk.Checkbutton(memory_frame, text="Gestione automatica della memoria", variable=auto_memory_var)
        auto_memory_check.pack(anchor=tk.W, pady=5)

        def toggle_memory_slider():
            if auto_memory_var.get():
                # When auto memory is enabled, reset slider to default 75%
                memory_slider.set(75)
                memory_slider.config(state="disabled")
            else:
                # When manual control is enabled, enable slider
                memory_slider.config(state="normal")

        auto_memory_check.config(command=toggle_memory_slider)

        memory_slider_label = ttk.Label(memory_frame, text="Soglia utilizzo RAM (%):")
        memory_slider_label.pack(anchor=tk.W, padx=5, pady=5)

        memory_slider = ttk.Scale(memory_frame, from_=10, to=95, orient=tk.HORIZONTAL,
                                variable=memory_percent_var, command=lambda v: memory_percent_var.set(int(float(v))))
        memory_slider.pack(fill=tk.X, padx=5, pady=5)

        def update_memory_details():
            percent = memory_percent_var.get()
            try:
                total_ram = psutil.virtual_memory().total / (1024 ** 3)
                ram_usage = total_ram * (percent / 100)
                memory_details_label.config(
                    text=f"RAM utilizzabile: {ram_usage:.2f} GB di {total_ram:.2f} GB totali"
                )
            except Exception as e:
                memory_details_label.config(
                    text="Errore nel calcolo della RAM. Assicurati che psutil sia installato."
                )
                print(f"Errore nel calcolo della RAM: {e}")

        memory_percent_var.trace("w", lambda *args: update_memory_details())
        memory_details_label = ttk.Label(memory_frame, text="")
        memory_details_label.pack(anchor=tk.W, padx=5, pady=5)
        update_memory_details()

        toggle_memory_slider()
        # Pulsanti finali per la finestra
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(15, 0))

        # Aggiungi questa funzione prima della definizione dei pulsanti finali
        def restore_defaults():
            # Valori predefiniti per le opzioni di ricerca
            depth_var.set(0)  # Profondità illimitata
            search_files_var.set(True)
            search_folders_var.set(True)
            search_content_var.set(True)  # Per default non cerchiamo nei contenuti
            whole_word_var.set(False)
            
            # Valori predefiniti per i filtri avanzati
            min_size_var.set("0")
            max_size_var.set("0")
            min_date.entry.delete(0, 'end')
            max_date.entry.delete(0, 'end')
            
            # Ripristino esclusioni
            system_exclusions_var.set(True)
            for item in excluded_list.get_children():
                excluded_list.delete(item)
            
            # Percorsi di sistema standard
            system_paths = [
                "C:/Windows", 
                "C:/Program Files", 
                "C:/Program Files (x86)",
                "C:/ProgramData",
                "C:/Users/All Users",
                "C:/Program Files (x86)/Client Active Directory Rights Management Services"
            ]
            for path in system_paths:
                excluded_list.insert("", "end", values=(path,))
            
            # Opzioni a blocchi predefinite
            files_block_var.set(500)  # Valore predefinito
            parallel_var.set(2)      # Valore predefinito
            auto_adjust_var.set(True)
            prioritize_var.set(True)
            
            # Opzioni di performance predefinite
            timeout_enabled_var.set(True)
            timeout_seconds_var.set(3600)
            max_files_var.set(100000)
            max_results_var.set(50000)
            threads_var.set(4)
            max_size_mb_var.set(50)
            dir_size_calc_var.set("disabilitato")
            
            # Opzioni memoria predefinite
            auto_memory_var.set(True)
            memory_percent_var.set(75)
            update_memory_details()
            toggle_memory_slider()
            
            # Mostra conferma all'utente
            messagebox.showinfo("Ripristino", "I valori predefiniti sono stati ripristinati.")

        # E poi, modifica il pulsante per usare la funzione appena creata:
        defaults_btn = ttk.Button(btn_frame, text="Ripristina valori predefiniti", command=restore_defaults)
        defaults_btn.pack(side=LEFT)
        self.create_tooltip(defaults_btn, "Ripristina tutte le impostazioni ai valori predefiniti")
        
        def save_options():
            # Log dell'inizio dell'operazione
            self.log_debug("===== INIZIO MODIFICA IMPOSTAZIONI AVANZATE =====")
            
            try:
                # Memorizza i valori precedenti per confronto
                old_depth = self.max_depth
                old_search_files = self.search_files.get()
                old_search_folders = self.search_folders.get()
                old_search_content = self.search_content.get()
                old_whole_word = self.whole_word_search.get()
                old_size_min = self.advanced_filters["size_min"]
                old_size_max = self.advanced_filters["size_max"]
                old_date_min = self.advanced_filters["date_min"]
                old_date_max = self.advanced_filters["date_max"]
                old_extensions = self.advanced_filters.get("extensions", [])
                
                # Memorizza i valori precedenti di gestione memoria
                old_auto_memory = getattr(self, 'auto_memory_management', True)
                old_memory_percent = getattr(self, 'memory_usage_percent', 75)
                
                # Inizializzazione sicura per excluded_paths
                if hasattr(self, 'excluded_paths') and self.excluded_paths is not None:
                    old_excluded_paths = self.excluded_paths.copy()
                else:
                    old_excluded_paths = []
                    self.excluded_paths = []  # Inizializza se non esiste
                    
                old_max_files_block = self.max_files_per_block.get()
                old_max_parallel = self.max_parallel_blocks.get()
                old_auto_adjust = self.block_size_auto_adjust.get()
                old_prioritize = self.prioritize_user_folders.get()
                old_timeout_enabled = self.timeout_enabled.get()
                old_timeout_seconds = self.timeout_seconds.get()
                old_max_files = self.max_files_to_check.get()
                old_max_results = self.max_results.get()
                old_worker_threads = self.worker_threads.get()
                old_max_file_size = self.max_file_size_mb.get()
                old_dir_size_calc = self.dir_size_calculation.get()
                
                # Salva le opzioni di ricerca
                self.max_depth = int(depth_var.get())
                self.search_files.set(search_files_var.get())
                self.search_folders.set(search_folders_var.get())
                self.search_content.set(search_content_var.get())
                self.whole_word_search.set(whole_word_var.get())
                
                # Salva i filtri avanzati
                min_kb = int(min_size_var.get() or 0)
                max_kb = int(max_size_var.get() or 0)
                self.advanced_filters["size_min"] = min_kb * 1024
                self.advanced_filters["size_max"] = max_kb * 1024 if max_kb > 0 else 0
                
                self.advanced_filters["date_min"] = min_date.entry.get().strip()
                self.advanced_filters["date_max"] = max_date.entry.get().strip()
                
                # Salva i percorsi esclusi
                new_excluded_paths = []
                for item in excluded_list.get_children():
                    values = excluded_list.item(item)["values"]
                    if values:
                        new_excluded_paths.append(values[0])
                self.excluded_paths = new_excluded_paths
                
                # Salva le opzioni a blocchi
                self.max_files_per_block.set(files_block_var.get())
                self.max_parallel_blocks.set(parallel_var.get())
                self.block_size_auto_adjust.set(auto_adjust_var.get())
                self.prioritize_user_folders.set(prioritize_var.get())
                
                # Salva le opzioni di performance
                self.timeout_enabled.set(timeout_enabled_var.get())
                self.timeout_seconds.set(timeout_seconds_var.get())
                self.max_files_to_check.set(max_files_var.get())
                self.max_results.set(max_results_var.get())
                self.worker_threads.set(threads_var.get())
                self.max_file_size_mb.set(max_size_mb_var.get())
                self.dir_size_calculation.set(dir_size_calc_var.get())
                
                # Salva le opzioni di gestione della memoria
                self.auto_memory_management = auto_memory_var.get()
                self.memory_usage_percent = memory_percent_var.get()
                
                # Log delle modifiche
                self.log_debug(f"Profondità ricerca: {old_depth} -> {self.max_depth}")
                self.log_debug(f"Cerca nei file: {old_search_files} -> {self.search_files.get()}")
                self.log_debug(f"Cerca nelle cartelle: {old_search_folders} -> {self.search_folders.get()}")
                self.log_debug(f"Cerca nei contenuti: {old_search_content} -> {self.search_content.get()}")
                self.log_debug(f"Parole intere: {old_whole_word} -> {self.whole_word_search.get()}")
                self.log_debug(f"Dimensione min (KB): {old_size_min//1024} -> {self.advanced_filters['size_min']//1024}")
                self.log_debug(f"Dimensione max (KB): {old_size_max//1024} -> {self.advanced_filters['size_max']//1024}")
                self.log_debug(f"Data min: '{old_date_min}' -> '{self.advanced_filters['date_min']}'")
                self.log_debug(f"Data max: '{old_date_max}' -> '{self.advanced_filters['date_max']}'")

                # Log delle modifiche alla gestione memoria
                self.log_debug(f"Gestione memoria automatica: {old_auto_memory} -> {self.auto_memory_management}")
                self.log_debug(f"Percentuale utilizzo memoria: {old_memory_percent}% -> {self.memory_usage_percent}%")
                
                # Log percorsi esclusi
                added_paths = [p for p in self.excluded_paths if p not in old_excluded_paths]
                removed_paths = [p for p in old_excluded_paths if p not in self.excluded_paths]
                if added_paths:
                    self.log_debug(f"Percorsi esclusi aggiunti: {', '.join(added_paths)}")
                if removed_paths:
                    self.log_debug(f"Percorsi esclusi rimossi: {', '.join(removed_paths)}")
                    
                self.log_debug(f"File per blocco: {old_max_files_block} -> {self.max_files_per_block.get()}")
                self.log_debug(f"Blocchi paralleli: {old_max_parallel} -> {self.max_parallel_blocks.get()}")
                self.log_debug(f"Auto-adattamento blocchi: {old_auto_adjust} -> {self.block_size_auto_adjust.get()}")
                self.log_debug(f"Priorità cartelle utente: {old_prioritize} -> {self.prioritize_user_folders.get()}")
                self.log_debug(f"Timeout attivo: {old_timeout_enabled} -> {self.timeout_enabled.get()}")
                self.log_debug(f"Secondi timeout: {old_timeout_seconds} -> {self.timeout_seconds.get()}")
                self.log_debug(f"Max file da controllare: {old_max_files} -> {self.max_files_to_check.get()}")
                self.log_debug(f"Max risultati: {old_max_results} -> {self.max_results.get()}")
                self.log_debug(f"Thread paralleli: {old_worker_threads} -> {self.worker_threads.get()}")
                self.log_debug(f"Dimensione max file (MB): {old_max_file_size} -> {self.max_file_size_mb.get()}")
                self.log_debug(f"Modalità calcolo dimensioni: '{old_dir_size_calc}' -> '{self.dir_size_calculation.get()}'")
                
                # CORREZIONE: Salva effettivamente le impostazioni su file
                if hasattr(self, 'save_settings_to_file'):
                    self.save_settings_to_file()
                    self.log_debug("Impostazioni avanzate salvate su file permanente")
                else:
                    self.log_debug("AVVISO: Metodo save_settings_to_file non disponibile, salvataggio permanente non effettuato")
                
                self.log_debug("===== FINE MODIFICA IMPOSTAZIONI AVANZATE - SALVATAGGIO COMPLETATO =====")
                messagebox.showinfo("Impostazioni", "Opzioni salvate con successo!")
                dialog.destroy()
                
            except ValueError as e:
                error_msg = f"Errore di valore non valido: {str(e)}"
                self.log_debug(f"ERRORE SALVATAGGIO: {error_msg}")
                messagebox.showerror("Errore", error_msg)
                return
            except Exception as e:
                error_msg = f"Errore durante il salvataggio: {str(e)}"
                self.log_debug(f"ERRORE SALVATAGGIO: {error_msg}")
                import traceback
                self.log_debug(traceback.format_exc())  # Registra il traceback completo dell'errore
                messagebox.showerror("Errore", error_msg)
                return
        
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=dialog.destroy)
        cancel_btn.pack(side=RIGHT, padx=5)
        self.create_tooltip(cancel_btn, "Chiudi la finestra senza salvare le modifiche")
        
        save_btn = ttk.Button(btn_frame, text="Salva", command=save_options)
        save_btn.pack(side=RIGHT, padx=5)
        self.create_tooltip(save_btn, "Salva tutte le impostazioni e chiudi la finestra")
        
        # Centra la finestra
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
    
    @error_handler
    def create_tooltip(self, widget, text, delay=500, fade=True):
        """Crea tooltip con ritardo, effetti di dissolvenza e larghezza automatica"""
        
        tooltip_timer = None
        fade_timer = None
        
        def show_tooltip():
            x = widget.winfo_rootx() + 25
            y = widget.winfo_rooty() + 25
            
            tooltip = ttk.Toplevel(widget)
            tooltip.wm_overrideredirect(True)
            tooltip.attributes("-alpha", 0.0)  # Inizia invisibile
            tooltip.attributes("-topmost", True)  # Mantiene sopra altre finestre
            
            # Crea un frame con bordo e padding
            frame = ttk.Frame(tooltip, borderwidth=1, relief="solid", padding=10)
            frame.pack(fill="both", expand=True)
            
            # Calcolo della larghezza del testo
            font = ("Segoe UI", 9)  # Puoi modificare questo per usare un altro font
            
            # Calcola la lunghezza approssimativa del testo
            temp_label = tk.Label(font=font)
            
            # Determina se il testo necessita di essere diviso in più righe
            lines = text.split("\n")
            max_line_width = 0
            
            for line in lines:
                temp_label.configure(text=line)
                line_width = temp_label.winfo_reqwidth()
                max_line_width = max(max_line_width, line_width)
            
            # Limita la larghezza massima a 500 pixel
            max_width = min(max_line_width, 500)
            
            wraplength = max_width
                
            # Crea l'etichetta con il testo
            label = ttk.Label(frame, text=text, justify="left", 
                            wraplength=wraplength, font=font)
            label.pack(fill="both", expand=True)
            
            # Distruggi il label temporaneo
            temp_label.destroy()
            
            # Aggiorna immediatamente per calcolare le dimensioni
            tooltip.update_idletasks()
            
            # Posiziona il tooltip in modo che non vada fuori dallo schermo
            tooltip_width = tooltip.winfo_reqwidth()
            tooltip_height = tooltip.winfo_reqheight()
            
            screen_width = tooltip.winfo_screenwidth()
            screen_height = tooltip.winfo_screenheight()
            
            # Aggiusta la posizione se il tooltip esce dallo schermo
            if x + tooltip_width > screen_width:
                x = screen_width - tooltip_width - 10
            
            if y + tooltip_height > screen_height:
                y = screen_height - tooltip_height - 10
                
            # Imposta la posizione finale
            tooltip.wm_geometry(f"+{x}+{y}")
            
            widget._tooltip = tooltip
            
            if fade:
                # Effetto dissolvenza in entrata
                def fade_in(alpha=0.0):
                    if not hasattr(widget, "_tooltip"):
                        return
                    
                    tooltip.attributes("-alpha", alpha)
                    if alpha < 1.0:
                        nonlocal fade_timer
                        fade_timer = widget.after(20, lambda: fade_in(alpha + 0.1))
                
                fade_in()
        
        def enter(event):
            nonlocal tooltip_timer
            # Avvia il timer per mostrare il tooltip dopo un certo ritardo
            tooltip_timer = widget.after(delay, show_tooltip)
        
        def leave(event):
            nonlocal tooltip_timer, fade_timer
            
            # Cancella il timer se esiste
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
                tooltip_timer = None
            
            # Cancella il timer di dissolvenza se esiste
            if fade_timer:
                widget.after_cancel(fade_timer)
                fade_timer = None
            
            # Rimuovi il tooltip se esiste
            if hasattr(widget, "_tooltip"):
                if fade:
                    # Effetto dissolvenza in uscita
                    def fade_out(alpha=1.0):
                        if alpha <= 0 or not hasattr(widget, "_tooltip"):
                            if hasattr(widget, "_tooltip"):
                                widget._tooltip.destroy()
                                del widget._tooltip
                        else:
                            widget._tooltip.attributes("-alpha", alpha)
                            nonlocal fade_timer
                            fade_timer = widget.after(20, lambda: fade_out(alpha - 0.1))
                    
                    fade_out()
                else:
                    widget._tooltip.destroy()
                    del widget._tooltip
        
        # Gestisce anche il caso in cui il widget venga distrutto
        def on_destroy(event):
            nonlocal tooltip_timer, fade_timer
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
            if fade_timer and hasattr(widget, "_tooltip"):
                widget.after_cancel(fade_timer)
                widget._tooltip.destroy()
        
        widget.bind("<Enter>", enter)
        widget.bind("<Leave>", leave)
        widget.bind("<Destroy>", on_destroy)

    def select_all(self):
        self.results_list.selection_set(self.results_list.get_children())
        
    def deselect_all(self):
        self.results_list.selection_remove(self.results_list.get_children())

    @error_handler   
    def invert_selection(self):
        all_items = self.results_list.get_children()
        selected_items = self.results_list.selection()
        self.results_list.selection_remove(selected_items)
        to_select = set(all_items) - set(selected_items)
        for item in to_select:
            self.results_list.selection_add(item)
    
    @error_handler
    def treeview_sort_column(self, tv, col, reverse):
        """Ordina il TreeView in base alla colonna cliccata"""
        # Ottieni la lista di item con i loro valori
        l = [(tv.set(k, col), k) for k in tv.get_children('')]
        
        # Determina il tipo di ordinamento in base alla colonna
        try:
            if col == "size":
                # Gestione speciale per le dimensioni (KB, MB, GB)
                def extract_size(size_str):
                    if not size_str:
                        return 0
                    try:
                        # Estrai il valore numerico e converti tutto in byte per confrontare correttamente
                        value = float(size_str.split()[0])
                        if "KB" in size_str:
                            return value * 1024
                        elif "MB" in size_str:
                            return value * 1024 * 1024
                        elif "GB" in size_str:
                            return value * 1024 * 1024 * 1024
                        else:
                            return value  # Assumiamo che sia in byte
                    except:
                        return 0
                        
                l.sort(key=lambda x: extract_size(x[0]), reverse=reverse)
            elif col in ("modified", "created"):
                # Gestione per le date (assumendo formato DD/MM/YYYY HH:MM)
                from datetime import datetime
                def parse_date(date_str):
                    try:
                        if date_str and date_str != "N/A":
                            return datetime.strptime(date_str, "%d/%m/%Y %H:%M")
                        return datetime(1900, 1, 1)  # Data di default per valori vuoti
                    except:
                        return datetime(1900, 1, 1)
                        
                l.sort(key=lambda x: parse_date(x[0]), reverse=reverse)
            else:
                # Ordinamento standard alfanumerico
                l.sort(reverse=reverse)
        except Exception as e:
            self.log_debug(f"Errore durante l'ordinamento: {str(e)}")
            # Fallback all'ordinamento alfanumerico semplice
            l.sort(reverse=reverse)
            
        # Riorganizza gli elementi nel TreeView
        for index, (val, k) in enumerate(l):
            tv.move(k, '', index)

        # Memorizza la colonna ordinata e la direzione per il prossimo click
        tv.heading(col, command=lambda _col=col: self.treeview_sort_column(tv, _col, not reverse))
        
        # Aggiorna l'indicatore visivo di ordinamento (aggiunge freccia all'intestazione)
        for c in tv["columns"]:
            if c == col:
                tv.heading(c, text=f"{tv.heading(c, 'text').split(' ')[0]} {'▼' if reverse else '▲'}")
            else:
                # Rimuovi eventuali indicatori di ordinamento da altre colonne
                tv.heading(c, text=tv.heading(c, 'text').split(' ')[0])

    @error_handler # Add this to the FileSearchApp class to track log messages
    def _init_essential_variables(self):
        """Inizializza solo le variabili essenziali per l'avvio"""
        # Inizializza datetime_var subito all'inizio per evitare errori di sequenza
        self.datetime_var = StringVar()
        self.max_depth = 5
        # Variabili principali per la ricerca
        self.search_content = BooleanVar(value=True)
        self.search_path = StringVar()
        self.keywords = StringVar()
        self.search_results = []
        self.search_files = BooleanVar(value=True)
        self.search_folders = BooleanVar(value=True)
        self.is_searching = False
        self.progress_queue = queue.Queue()
        self.search_depth = StringVar(value="base") 

        # Inizializza impostazioni per la gestione della RAM
        self.auto_memory_management = True
        self.memory_usage_percent = 75
        # Inizializza tutte le variabili utilizzate nelle impostazioni
        self.timeout_enabled = tk.BooleanVar(value=False)
        self.timeout_seconds = tk.IntVar(value=3600)
        self.max_files_to_check = tk.IntVar(value=100000)
        self.max_results = tk.IntVar(value=50000)
        self.worker_threads = tk.IntVar(value=4)
        self.max_file_size_mb = tk.IntVar(value=100)
        self.use_indexing = tk.BooleanVar(value=True)
        self.skip_permission_errors = tk.BooleanVar(value=True)

        # Variabili per data/ora e utente (datetime_var già inizializzato)
        self.user_var = StringVar(value=getpass.getuser())
        
        # Variabili essenziali per l'interfaccia
        self.ignore_hidden = BooleanVar(value=True)
        self.search_executor = None
        self.exclude_system_files = BooleanVar(value=True)
        self.whole_word_search = BooleanVar(value=False)
        self.dir_size_calculation = StringVar(value="disabilitato")
        
        # Variabili per la visualizzazione
        self.directory_calculation_enabled = False
        self.dir_size_var = StringVar(value="")
        self.total_disk_var = StringVar(value="")
        self.used_disk_var = StringVar(value="")
        self.free_disk_var = StringVar(value="")
        
        # Add a queue for debug logs
        self.debug_logs_queue = queue.Queue(maxsize=5000)  # Limit to 5000 entries to avoid memory issues
        
    @error_handler # Show debug log window with export functionality
    def show_debug_log(self):
        """Displays a window with debug logs and export functionality"""
        if not hasattr(self, 'debug_window') or not self.debug_window.winfo_exists():
            self.debug_window = tk.Toplevel(self.root)
            self.debug_window.title("Debug Log")
            self.debug_window.geometry("1000x700")
            
            frame = ttk.Frame(self.debug_window)
            frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
            
            # Header frame con conteggio messaggi
            header_frame = ttk.Frame(frame)
            header_frame.pack(fill=tk.X, pady=(0, 10))
            
            # Verifica che il debug log sia inizializzato
            if not hasattr(self, 'debug_log'):
                self.debug_log = []
            
            # Etichetta informativa con conteggio
            self.log_count_label = ttk.Label(
                header_frame, 
                text=f"Registro di debug dell'applicazione: {len(self.debug_log)} messaggi"
            )
            self.log_count_label.pack(side=tk.LEFT)
            
            # Opzione auto-scroll
            self.autoscroll_var = tk.BooleanVar(value=True)
            ttk.Checkbutton(
                header_frame, 
                text="Auto scorrimento", 
                variable=self.autoscroll_var
            ).pack(side=tk.RIGHT, padx=5)
            
            # Crea un text widget con scrollbar
            text_frame = ttk.Frame(frame)
            text_frame.pack(fill=tk.BOTH, expand=True)
            
            # Scrollbar verticale
            v_scrollbar = ttk.Scrollbar(text_frame)
            v_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
            
            # Scrollbar orizzontale
            h_scrollbar = ttk.Scrollbar(text_frame, orient=tk.HORIZONTAL)
            h_scrollbar.pack(side=tk.BOTTOM, fill=tk.X)
            
            # Text widget con sfondo scuro per migliore leggibilità
            self.debug_text = tk.Text(
                text_frame, 
                wrap=tk.NONE,  # Permette scroll orizzontale
                width=80, 
                height=20,
                bg="#1e1e1e",  # Sfondo scuro
                fg="#e0e0e0",  # Testo chiaro
                font=("Consolas", 10),  # Font a larghezza fissa
                xscrollcommand=h_scrollbar.set,
                yscrollcommand=v_scrollbar.set
            )
            
            self.debug_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
            
            # Configura le scrollbar
            v_scrollbar.config(command=self.debug_text.yview)
            h_scrollbar.config(command=self.debug_text.xview)
            
            # Aggiungi pulsanti di utilità
            btn_frame = ttk.Frame(self.debug_window)
            btn_frame.pack(fill=tk.X, padx=10, pady=5)
            
            ttk.Button(
                btn_frame, 
                text="Aggiorna", 
                command=self.update_log_display
            ).pack(side=tk.LEFT, padx=5)
            
            ttk.Button(
                btn_frame, 
                text="Pulisci Log", 
                command=self.clear_log
            ).pack(side=tk.LEFT, padx=5)
            
            ttk.Button(
                btn_frame, 
                text="Esporta in TXT", 
                command=self.export_log_to_txt
            ).pack(side=tk.LEFT, padx=5)
            
            # Posiziona la finestra al centro
            self.debug_window.update_idletasks()
            width = self.debug_window.winfo_width()
            height = self.debug_window.winfo_height()
            x = (self.debug_window.winfo_screenwidth() // 2) - (width // 2)
            y = (self.debug_window.winfo_screenheight() // 2) - (height // 2)
            self.debug_window.geometry(f"{width}x{height}+{x}+{y}")
            
            # Imposta una dimensione minima ragionevole
            self.debug_window.minsize(800, 500)
            
            # Mostra il log corrente
            self.update_log_display()

        else:
            # Se la finestra esiste già, portala in primo piano
            self.debug_window.lift()
            self.debug_window.focus_force()
            # Aggiorna il contenuto
            self.update_log_display()

    @error_handler
    def update_log_display(self):
        """Aggiorna la visualizzazione dei log nella finestra di debug"""
        if hasattr(self, 'debug_window') and hasattr(self, 'debug_text') and self.debug_window.winfo_exists():
            # Salva la posizione corrente dello scroll verticale
            current_pos = self.debug_text.yview()[0]
            
            # Cancella il contenuto attuale
            self.debug_text.config(state=tk.NORMAL)
            self.debug_text.delete(1.0, tk.END)
            
            # Verifica che il debug log sia inizializzato
            if not hasattr(self, 'debug_log'):
                self.debug_log = []
                
            # Aggiorna l'etichetta con il conteggio dei messaggi
            if hasattr(self, 'log_count_label'):
                self.log_count_label.config(text=f"Registro di debug dell'applicazione: {len(self.debug_log)} messaggi")
                
            if self.debug_log:
                # Limita la visualizzazione a 5000 messaggi per non rallentare l'interfaccia
                max_display = 5000
                if len(self.debug_log) > max_display:
                    self.debug_text.insert(tk.END, f"[Mostrando solo gli ultimi {max_display} di {len(self.debug_log)} messaggi...]\n\n")
                    log_entries = self.debug_log[-max_display:]
                else:
                    log_entries = self.debug_log
                    
                # Inserisci i log
                log_text = "\n".join(log_entries)
                self.debug_text.insert(tk.END, log_text)
                
                # Evidenzia gli errori con colore rosso
                self.highlight_errors()
                
                # Forza un aggiornamento dell'interfaccia
                self.debug_window.update_idletasks()
                
                # Prima imposta l'indice orizzontale all'inizio (più a sinistra)
                self.debug_text.xview_moveto(0.0)
                
                # Poi imposta anche il primo carattere visibile nella prima colonna
                self.debug_text.mark_set("insert", "1.0")
                
                # Aggiungi un ritardo per garantire che lo scroll sia applicato
                def force_left_scroll():
                    self.debug_text.xview_moveto(0.0)
                    # Modifica la posizione fisica dello scrollbar
                    if hasattr(self.debug_text, 'xscrollcommand'):
                        self.debug_text.xscrollcommand(0.0, 1.0)
                
                # Esegui con un piccolo ritardo per assicurarti che venga applicato dopo il rendering
                self.debug_window.after(50, force_left_scroll)
                
                # Decidi se mantenere la posizione verticale precedente o scorrere alla fine
                if hasattr(self, 'autoscroll_var') and self.autoscroll_var.get():
                    self.debug_text.see(tk.END)  # Scorri alla fine verticalmente
                    # Ma mantieni comunque lo scroll orizzontale a sinistra
                    self.debug_window.after(100, lambda: self.debug_text.xview_moveto(0.0))
                else:
                    # Torna alla posizione precedente dello scroll verticale
                    self.debug_text.yview_moveto(current_pos)
            else:
                self.debug_text.insert(tk.END, "Nessun messaggio di debug disponibile.")
                
            # Rendi il testo di nuovo sola lettura
            self.debug_text.config(state=tk.DISABLED)

    @error_handler
    def clear_log(self):
        """Pulisce la visualizzazione del log"""
        if hasattr(self, 'debug_text') and self.debug_text.winfo_exists():
            self.debug_text.config(state=tk.NORMAL)
            self.debug_text.delete(1.0, tk.END)
            self.debug_text.insert(tk.END, "Log cancellato. I nuovi messaggi appariranno qui.")
            self.debug_text.config(state=tk.DISABLED)

    @error_handler
    def export_log_to_txt(self):
        """Esporta il log completo in un file di testo"""
        # Chiedi all'utente dove salvare il file
        file_path = tk.filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt"), ("All files", "*.*")],
            title="Salva log come file di testo"
        )
        
        if not file_path:  # Utente ha annullato
            return
            
        try:
            # Verifica che il debug log sia inizializzato
            if not hasattr(self, 'debug_log'):
                self.debug_log = []
                
            log_content = "\n".join(self.debug_log)
            
            # Aggiungi intestazione con timestamp e info utente
            import getpass
            header = f"Debug Log - Esportato il {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n"
            header += f"Utente: {getpass.getuser()}\n"
            header += f"Applicazione: File Search Tool V9.2.5 Beta\n"
            header += f"Numero totale messaggi: {len(self.debug_log)}\n"
            header += "-" * 80 + "\n\n"
            
            # Scrivi su file
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(header + log_content)

            tk.messagebox.showinfo("Esportazione completata", 
                            f"Log salvato con successo in:\n{file_path}\n\n"
                            f"Il file contiene {len(self.debug_log)} messaggi di log.")
        except Exception as e:
            tk.messagebox.showerror("Errore esportazione", 
                            f"Si è verificato un errore durante l'esportazione:\n{str(e)}")

    @error_handler
    def highlight_errors(self):
        """Evidenzia le righe di errore nel log con colore rosso"""
        if hasattr(self, 'debug_text'):
            # Cerca tutte le righe con "[ERROR]"
            start = "1.0"
            while True:
                start = self.debug_text.search("[ERROR]", start, tk.END)
                if not start:
                    break
                    
                # Trova la fine della riga
                line_end = self.debug_text.search("\n", start, tk.END)
                if not line_end:
                    line_end = tk.END
                    
                # Crea un tag per questa riga
                self.debug_text.tag_add("error", start, line_end)
                
                # Imposta il prossimo inizio ricerca
                start = line_end
                
            # Configura il tag "error" per mostrare il testo in rosso
            self.debug_text.tag_config("error", foreground="#ff6b6b")

# Funzione principale per eseguire l'applicazione
def main():
    import sys
    
    # Crea la finestra principale con il tema desiderato
    root = ttk.Window(themename="darkly")
    root.withdraw()  # Nascondi completamente la finestra durante l'inizializzazione
    
    # Crea la schermata di splash
    splash = create_splash_screen(root)
    
    # Inizializza l'applicazione (ma l'interfaccia rimane nascosta)
    app = FileSearchApp(root)
    
    # Controlla se ci sono argomenti da linea di comando
    if len(sys.argv) > 1:
        app.search_path.set(sys.argv[1])
    
    # Funzione per completare l'avvio e mostrare la finestra principale
    def finish_startup():
        splash.destroy()
        # Configura le dimensioni della finestra prima di mostrarla
        root.geometry("1300x850")
        # Mostra la finestra completamente costruita
        root.deiconify()
        
    # Completa l'avvio dopo un breve ritardo
    root.after(1500, finish_startup)
    
    root.mainloop()

def create_splash_screen(parent):
    splash_win = tk.Toplevel(parent)
    splash_win.title("")
    splash_win.overrideredirect(True)
    splash_win.attributes("-topmost", True)
    
    # Dimensioni dello splash
    width, height = 500, 250
    screen_width = splash_win.winfo_screenwidth()
    screen_height = splash_win.winfo_screenheight()
    x = (screen_width // 2) - (width // 2)
    y = (screen_height // 2) - (height // 2)
    splash_win.geometry(f"{width}x{height}+{x}+{y}")
    
    # Contenuto dello splash
    frame = ttk.Frame(splash_win, padding=20)
    frame.pack(fill=tk.BOTH, expand=tk.YES)
    
    ttk.Label(frame, text="File Search Tool V9.2.5 Beta", 
            font=("Helvetica", 18, "bold")).pack(pady=(10, 5))
    ttk.Label(frame, text="Forensics G.di F.", 
            font=("Helvetica", 14)).pack(pady=(0, 20))
    ttk.Label(frame, text="Caricamento applicazione in corso...").pack()
    
    progress = ttk.Progressbar(frame, mode="indeterminate")
    progress.pack(fill=tk.X, pady=10)
    progress.start(10)
    
    return splash_win

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        import traceback
        with open("error_log.txt", "w", encoding="utf-8") as f:
            f.write(f"Si è verificato un errore: {str(e)}\n")
            f.write(traceback.format_exc())
        print(f"Si è verificato un errore: {str(e)}")
        
        # Se stai usando tkinter, puoi anche mostrare un messaggio all'utente
        try:
            import tkinter as tk
            from tkinter import messagebox
            root = tk.Tk()
            root.withdraw()
            messagebox.showerror("Errore durante l'avvio", 
                f"Si è verificato un errore durante l'avvio dell'applicazione.\n\n"
                f"Errore: {str(e)}\n\n"
                f"I dettagli sono stati salvati nel file error_log.txt")
        except:
            pass  # Se anche la visualizzazione del messaggio fallisce, continua
