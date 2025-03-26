import os
import sys  # Aggiungo l'import di sys
import shutil
import ttkbootstrap as ttk
from ttkbootstrap.constants import *
from tkinter import filedialog, messagebox, BooleanVar, StringVar, IntVar
import threading
import queue
from datetime import datetime
import getpass
import zipfile
from ttkbootstrap.dialogs import Querybox
import traceback
import time
import concurrent.futures
import mimetypes
import signal

# Dizionario per tracciare il supporto alle librerie
file_format_support = {
    "docx": False,
    "pdf": False,
    "pptx": False,
    "excel": False,
    "odt": False,
    "rtf": False
}

# Elenco di librerie da installare se mancanti
missing_libraries = []

# Importazione di librerie opzionali con gestione degli errori
try:
    import docx
    file_format_support["docx"] = True
    print("Supporto Word (.docx) attivato")
except ImportError:
    missing_libraries.append("python-docx")
    print("Supporto Word (.docx) non disponibile")

try:
    import PyPDF2
    file_format_support["pdf"] = True
    print("Supporto PDF attivato")
except ImportError:
    missing_libraries.append("PyPDF2")
    print("Supporto PDF non disponibile")

try:
    import pptx
    file_format_support["pptx"] = True
    print("Supporto PowerPoint attivato")
except ImportError:
    missing_libraries.append("python-pptx")
    print("Supporto PowerPoint non disponibile")

try:
    import openpyxl
    file_format_support["excel"] = True
    print("Supporto Excel attivato")
except ImportError:
    missing_libraries.append("openpyxl")
    print("Supporto Excel non disponibile")

try:
    from striprtf.striprtf import rtf_to_text
    file_format_support["rtf"] = True
    print("Supporto RTF attivato")
except ImportError:
    missing_libraries.append("striprtf")
    print("Supporto RTF non disponibile")

try:
    import odfdo
    file_format_support["odt"] = True
    print("Supporto ODT attivato (libreria odfdo)")
except ImportError:
    try:
        from odf import opendocument, text, teletype
        file_format_support["odt"] = True
        print("Supporto ODT attivato (libreria odf)")
    except ImportError:
        missing_libraries.append("odfdo")
        print("Supporto ODT non disponibile")
try:
    import win32com.client
    file_format_support["doc"] = True
    print("Supporto Word (.doc) attivato tramite pywin32")
except ImportError:
    missing_libraries.append("pywin32")
    print("Supporto Word (.doc) non disponibile")

try:
    import win32com.client
    file_format_support["xls"] = True
    print("Supporto Excel (.xls) attivato tramite pywin32")
except ImportError:
    missing_libraries.append("pywin32")
    print("Supporto Excel (.xls) non disponibile")
    
class FileSearchApp:
    def __init__(self, root):
        self.root = root
        self.root.title("File Search Tool.. Nucleo Perugia")
        self.root.geometry("1300x850")
        
        # Inizializza le variabili per data/ora e utente
        self.datetime_var = StringVar()
        self.user_var = StringVar(value=getpass.getuser())
        
        # Variabili
        self.search_content = BooleanVar(value=False)
        self.search_path = StringVar()
        self.keywords = StringVar()
        self.search_results = []
        self.search_files = BooleanVar(value=True)
        self.search_folders = BooleanVar(value=True)
        self.ignore_hidden = BooleanVar(value=True)
        self.is_searching = False
        self.progress_queue = queue.Queue()
        self.search_executor = None
        
        # Variabili aggiuntive per i miglioramenti
        self.search_start_time = None
        self.stop_search = False
        self.max_depth = 0  # 0 = illimitato
        self.last_search_params = {}
        self.search_history = []
        self.advanced_filters = {
            "size_min": 0,
            "size_max": 0,
            "date_min": None,
            "date_max": None,
            "extensions": []
        }
        
        # Variabili per limiti di tempo e altre ottimizzazioni
        self.timeout_enabled = BooleanVar(value=False)
        self.timeout_seconds = IntVar(value=3600)  # Default 60 secondi
        self.max_files_to_check = IntVar(value=100000)  # Limite numero di file da controllare
        self.max_results = IntVar(value=1000)  # Limite numero di risultati
        self.chunk_size = 8192  # Grandezza del chunk per la lettura dei file
        self.max_file_size_mb = IntVar(value=100)  # Dimensione massima file da analizzare (in MB)
        self.worker_threads = IntVar(value=min(8, os.cpu_count() or 4))  # Numero di worker threads per elaborazione parallela
        self.use_indexing = BooleanVar(value=True)  # Usa indicizzazione per velocizzare ricerche future
        self.search_index = {}  # Dizionario per indicizzare file e contenuti
        
        # Info utente e datetime
        self.current_user = getpass.getuser()
        self.datetime_var = StringVar()
        self.update_datetime()
        self.create_widgets()

        # Verifica le librerie e mostra notifica se necessario
        self.check_and_notify_missing_libraries()
        self.debug_mode = True
        
        # Registra handler per CTRL+C
        self.register_interrupt_handler()
    
        # Imposta il tema iniziale
        self.update_theme_colors("dark")  # O "dark" se il tema predefinito è scuro
        
    def register_interrupt_handler(self):
        """Registra il gestore degli interrupt (CTRL+C)"""
        def handle_interrupt(sig, frame):
            if self.is_searching:
                self.stop_search_process()
            else:
                self.root.quit()
        
        signal.signal(signal.SIGINT, handle_interrupt)

    def log_debug(self, message):
        """Funzione per logging, stampa solo quando debug_mode è True"""
        if self.debug_mode:
            print(f"[DEBUG] {message}")

    def check_and_notify_missing_libraries(self):
        """Verifica e notifica l'utente di eventuali librerie mancanti"""
        missing = []
        
        if not file_format_support["docx"]:
            missing.append("python-docx (per file Word)")
        if not file_format_support["pdf"]: 
            missing.append("PyPDF2 (per file PDF)")
        if not file_format_support["pptx"]:
            missing.append("python-pptx (per file PowerPoint)")
        if not file_format_support["excel"]:
            missing.append("openpyxl (per file Excel)")
        if not file_format_support["rtf"]:
            missing.append("striprtf (per file RTF)")
        if not file_format_support["odt"]:
            missing.append("odfdo (per file OpenDocument)")
        
        if missing:
            message = "Alcune funzionalità di ricerca nei contenuti sono disabilitate.\n\n"
            message += "Per abilitare il supporto completo ai vari formati di file, installa le seguenti librerie:\n\n"
            
            for lib in missing:
                message += f"- {lib}\n"
            
            message += "\nPuoi installarle con il comando:\n"
            message += "pip install " + " ".join([lib.split(" ")[0] for lib in missing])
            
            # Mostra la notifica dopo un breve ritardo per permettere all'UI di caricarsi
            self.root.after(1000, lambda: messagebox.showinfo("Librerie opzionali mancanti", message))

    def update_datetime(self):
        current_time = datetime.now().strftime('%d-%m-%Y %H:%M:%S')
        self.datetime_var.set(f"Data: {current_time} | Utente: {self.current_user}")
        self.root.after(1000, self.update_datetime)

    def browse_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.search_path.set(directory)
            
    def optimize_system_search(self, path):
        """Ottimizza la ricerca per percorsi di sistema come C:/ impostando parametri appropriati"""
        is_system_path = path.lower() in ["c:/", "c:\\", "d:/", "d:\\"]
        
        # Verifica se è un percorso di rete o server (inizia con \\ o //)
        is_network_path = path.startswith('\\\\') or path.startswith('//')
        
        # Verifica se è un percorso Unix root o home
        is_unix_root = path in ['/', '/home', '/usr', '/var', '/etc']
        
        if is_system_path or is_network_path or is_unix_root:
            # Salva i parametri attuali
            original_params = {
                "max_files": self.max_files_to_check.get(),
                "worker_threads": self.worker_threads.get(),
                "timeout": self.timeout_seconds.get() if self.timeout_enabled.get() else None
            }
            
            # Applica parametri ottimizzati
            self.max_files_to_check.set(1000000)  # Un milione di file
            self.worker_threads.set(min(12, os.cpu_count() or 4))  # Aumenta thread per server
            
            if self.timeout_enabled.get():
                self.timeout_seconds.set(max(3600, self.timeout_seconds.get()))  # Minimo 1 ora
            
            # Notifica l'utente
            messagebox.showinfo(
                "Ricerca su sistema o server",
                "Stai avviando una ricerca su un percorso di sistema o server.\n\n"
                "Per ottimizzare la ricerca, sono stati temporaneamente modificati alcuni parametri "
                "per consentire una scansione più approfondita.\n\n"
                "La ricerca potrebbe richiedere tempo e risorse significative."
            )
            
            return original_params
        return None
    
    def optimize_disk_search_order(self, path, directories):
        """Ottimizza l'ordine di esplorazione delle directory quando si cerca in un disco di sistema.
        Dà priorità alle cartelle più importanti come Users, Documents, ecc.
        """
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"]:
            prioritized = []
            normal = []
            
            # Cartelle ad alta priorità (esplora per prime)
            high_priority_folders = ["users", "documents", "documents and settings", "desktop", "downloads"]
            
            # Cartelle a bassa priorità (esplora per ultime)
            low_priority_folders = ["windows", "program files", "program files (x86)", "programdata", "$recycle.bin", "system volume information"]
            
            for dir_path in directories:
                dir_name = os.path.basename(dir_path).lower()
                
                if any(priority in dir_name for priority in high_priority_folders):
                    prioritized.insert(0, dir_path)  # Aggiungi all'inizio per priorità più alta
                elif any(low in dir_name for low in low_priority_folders):
                    normal.append(dir_path)  # Aggiungi alla fine
                else:
                    prioritized.append(dir_path)  # Priorità media
                    
            return prioritized + normal
        
        # Se non è un percorso di sistema, restituisci l'elenco originale
        return directories

    def start_search(self):
        # Pulisci risultati precedenti
        for item in self.results_list.get_children():
            self.results_list.delete(item)
        
        if not self.search_path.get() or not self.keywords.get():
            messagebox.showerror("Errore", "Inserisci directory e parole chiave")
            return
        
        # Reset per la nuova ricerca
        self.stop_search = False
        self.stop_button["state"] = "normal"
        original_params = self.optimize_system_search(self.search_path.get())

        # Aggiorna l'orario di avvio e resetta l'orario di fine
        current_time = datetime.now().strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")

        # Salva i parametri di ricerca
        self.last_search_params = {
            "path": self.search_path.get(),
            "keywords": self.keywords.get(),
            "search_files": self.search_files.get(),
            "search_folders": self.search_folders.get(),
            "search_content": self.search_content.get()
        }
        
        # Aggiungi alla cronologia di ricerca se non già presente
        if self.keywords.get() and self.keywords.get() not in [h["keywords"] for h in self.search_history]:
            self.search_history.append(self.last_search_params.copy())
            if len(self.search_history) > 10:  # Mantieni solo le ultime 10 ricerche
                self.search_history.pop(0)
        
        self.is_searching = True
        self.search_button["state"] = "disabled"
        self.progress_bar["value"] = 0
        self.status_label["text"] = "Ricerca in corso..."
        
        # Aggiorna il valore della profondità massima
        try:
            self.max_depth = int(self.depth_spinbox.get())
        except ValueError:
            self.max_depth = 0  # Valore predefinito se non valido
        
        # Ottieni le parole chiave di ricerca
        search_terms = [term.strip() for term in self.keywords.get().split(',') if term.strip()]
            
        # Avvia la ricerca in un thread separato
        self.search_results = []  # Resetta i risultati
        search_thread = threading.Thread(target=self._search_thread, 
                                        args=(self.search_path.get(), search_terms, self.search_content.get()))
        search_thread.daemon = True
        search_thread.start()
        
        # Memorizza l'orario di inizio come oggetto datetime
        self.search_start_time = datetime.now()
        current_time = self.search_start_time.strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")
        self.total_time_label.config(text="--:--")  # Reset del tempo totale

        # Avvia l'aggiornamento della progress bar
        self.update_progress()
        # Salva original_params come attributo per ripristinarli dopo la ricerca
        if original_params:
            self.original_system_search_params = original_params

        current_theme = "dark"  # Sostituisci con la logica per determinare il tema corrente
        self.update_theme_colors(current_theme)

        # Aggiungi questo nuovo metodo per calcolare il tempo totale
    def update_total_time(self):
        """Calcola e visualizza il tempo totale trascorso tra inizio e fine ricerca"""
        if self.search_start_time:
            end_time = datetime.now()
            time_diff = end_time - self.search_start_time
            
            # Converti in minuti e secondi
            total_seconds = int(time_diff.total_seconds())
            minutes = total_seconds // 60
            seconds = total_seconds % 60
            
            # Formatta la stringa del tempo totale
            if minutes > 0:
                total_time_str = f"{minutes}min {seconds}sec"
            else:
                total_time_str = f"{seconds}sec"
            
            self.total_time_label.config(text=total_time_str)   
    # Esempio di utilizzo per cambiare il tema
    def change_theme(self, theme):
        self.update_theme_colors(theme)

    def _search_thread(self, path, keywords, search_content):
        try:
            # Inizializza i contatori di file e directory esaminati e il tempo di inizio
            files_checked = 0
            dirs_checked = 0
            start_time = time.time()
            timeout = self.timeout_seconds.get() if self.timeout_enabled.get() else None
            
            # Determina se si tratta di una ricerca completa del sistema (C:/ o simile)
            is_system_search = path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]

            # Per ricerche di sistema, adatta automaticamente i parametri
            if is_system_search:
                # Informa l'utente che la ricerca potrebbe richiedere molto tempo
                self.progress_queue.put(("status", "Ricerca completa del sistema in corso - potrebbe richiedere molto tempo"))
                
                # Temporaneamente aumenta i limiti per la ricerca di sistema
                original_max_files = self.max_files_to_check.get()
                original_timeout = timeout
                
                # Imposta limiti più elevati per la ricerca completa
                self.max_files_to_check.set(5000000)  # 5 milioni di file
                
                # Disabilita temporaneamente il timeout o aumentalo significativamente
                if timeout and timeout < 3600:
                    timeout = 3600 * 8  # 8 ore
                
                # Avviso all'utente
                self.progress_queue.put(("status", "Ricerca completa avviata - parametri adattati per ricerca approfondita"))
            
            # Variabili per gestire l'aggiornamento ogni secondo
            last_update_time = time.time()
            
            # Crea un executor per processare i file in parallelo in modo sicuro
            try:
                max_workers = max(1, min(32, self.worker_threads.get()))  # Garantisce un valore valido
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
            except Exception as e:
                self.log_debug(f"Errore nella creazione del ThreadPoolExecutor: {str(e)}")
                # Fallback a un valore sicuro
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
            
            futures = []
            
            # Insieme per tenere traccia delle cartelle visitate (evita loop infiniti con symlink)
            visited_dirs = set()

            # Funzione per elaborare un file
            def process_file(file_path, keywords):
                if self.stop_search:
                    return None
                try:
                        # Aggiungi un timeout complessivo per l'elaborazione di ogni file
                        import signal
                        
                        # Funzione per gestire il timeout
                        def timeout_handler(signum, frame):
                            raise TimeoutError(f"Tempo scaduto durante l'elaborazione del file {file_path}")
                        
                        # Imposta un allarme di 20 secondi per ogni file
                        if os.name != 'nt':  # signal.SIGALRM non è disponibile su Windows
                            signal.signal(signal.SIGALRM, timeout_handler)
                            signal.alarm(20)
                        
                        # Resto del codice di elaborazione...
                        
                        # Disattiva l'allarme se tutto è andato bene
                        if os.name != 'nt':
                            signal.alarm(0)
                            
                except TimeoutError as e:
                    self.log_debug(str(e))
                    return None
                except Exception as e:
                    self.log_debug(f"Errore durante l'elaborazione del file {file_path}: {str(e)}")
                    return None
                try:
                    # Verifica filtri di dimensione
                    file_size = os.path.getsize(file_path)
                    if (self.advanced_filters["size_min"] > 0 and file_size < self.advanced_filters["size_min"]) or \
                    (self.advanced_filters["size_max"] > 0 and file_size > self.advanced_filters["size_max"]):
                        return None
                    
                    # Verifica filtri di data
                    if self.advanced_filters["date_min"] or self.advanced_filters["date_max"]:
                        mod_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                        mod_date_str = mod_time.strftime("%d-%m-%Y")
                        
                        if self.advanced_filters["date_min"]:
                            min_date = datetime.strptime(self.advanced_filters["date_min"], "%d-%m-%Y")
                            if mod_time < min_date:
                                return None
                                
                        if self.advanced_filters["date_max"]:
                            max_date = datetime.strptime(self.advanced_filters["date_max"], "%d-%m-%Y")
                            if mod_time > max_date:
                                return None
                    
                    # Verifica filtri estensione
                    if self.advanced_filters["extensions"] and not any(file_path.lower().endswith(ext.lower()) 
                                                                for ext in self.advanced_filters["extensions"]):
                        return None
                    if os.path.splitext(file_path)[1].lower() in ['.doc', '.xls']:
                        self.progress_queue.put(("progress", self.progress_bar["value"]))  # Forza aggiornamento
                        
                    # Verifica corrispondenza nel nome file
                    filename = os.path.basename(file_path)
                    if any(keyword.lower() in filename.lower() for keyword in keywords):
                        return self.create_file_info(file_path)
                    
                    # Verifica contenuto se richiesto
                    if search_content and self.should_search_content(file_path):
                        max_size_bytes = self.max_file_size_mb.get() * 1024 * 1024
                        
                        # Salta file troppo grandi
                        if file_size > max_size_bytes:
                            self.log_debug(f"File {file_path} troppo grande per l'analisi del contenuto")
                            return None
                            
                        content = self.get_file_content(file_path)
                        if content and any(keyword.lower() in content.lower() for keyword in keywords):
                            return self.create_file_info(file_path)
                    
                    return None
                    
                except Exception as e:
                    self.log_debug(f"Errore durante l'elaborazione del file {file_path}: {str(e)}")
                    return None
            
            # Funzione per visitare una directory
            def visit_directory(directory, current_depth=0):
                nonlocal dirs_checked, files_checked, visited_dirs, last_update_time
                
                if self.stop_search:
                    return
                    
                # MODIFICA: Aggiornamento dello stato ogni secondo
                current_time = time.time()
                if current_time - last_update_time >= 1.0:  # Aggiorna ogni secondo
                    elapsed_time = current_time - start_time
                    self.progress_queue.put(("status", 
                        f"Analisi: {directory} (Cartelle: {dirs_checked}, File: {files_checked}, Tempo: {int(elapsed_time)}s)"))
                    self.progress_queue.put(("progress", 
                        min(90, int((files_checked / max(1, self.max_files_to_check.get())) * 100))))
                    last_update_time = current_time
                
                # Verifica timeout
                if timeout and time.time() - start_time > timeout:
                    self.progress_queue.put(("timeout", "Timeout raggiunto"))
                    return
                
                # Verifica limite di profondità (0 = illimitato)
                if self.max_depth > 0 and current_depth > self.max_depth:
                    return
                    
                # Gestione dei percorsi per evitare loop infiniti
                try:
                    real_path = os.path.realpath(directory)
                    if real_path in visited_dirs:
                        return
                    visited_dirs.add(real_path)
                except Exception as e:
                    self.log_debug(f"Impossibile risolvere il percorso reale per {directory}: {str(e)}")
                    if directory in visited_dirs:
                        return
                    visited_dirs.add(directory)

                try:
                    # Incrementa il contatore delle cartelle analizzate
                    dirs_checked += 1
                    
                    # Gestione listdir in modo sicuro
                    try:
                        items = os.listdir(directory)
                    except PermissionError:
                        self.log_debug(f"Permesso negato per la directory {directory}")
                        return
                    except FileNotFoundError:
                        self.log_debug(f"Directory non trovata: {directory} (potrebbe essere stata rimossa durante la scansione)")
                        return
                    except OSError as e:
                        self.log_debug(f"Errore OS nell'accesso alla directory {directory}: {str(e)}")
                        return
                    except Exception as e:
                        self.log_debug(f"Errore generico nell'accesso alla directory {directory}: {str(e)}")
                        return
                    
                    # Prima processa tutte le directory, poi i file (priorità alla profondità)
                    directories = []
                    
                    # Loop per le directory
                    for item in items:
                        if self.stop_search:
                            return
                            
                        item_path = os.path.join(directory, item)
                        
                        # Salta file/cartelle nascosti se richiesto
                        try:
                            if self.ignore_hidden.get() and (item.startswith('.') or 
                                                    (os.name == 'nt' and os.path.exists(item_path) and 
                                                        os.stat(item_path).st_file_attributes & 2)):
                                continue
                        except Exception as e:
                            self.log_debug(f"Errore nel controllo hidden per {item_path}: {str(e)}")
                            continue
                        
                        # Gestione directory
                        try:
                            if os.path.isdir(item_path):
                                directories.append(item_path)
                                
                                # Verifica corrispondenza nome cartella
                                if self.search_folders.get():
                                    if any(keyword.lower() in item.lower() for keyword in keywords):
                                        folder_info = self.create_folder_info(item_path)
                                        self.search_results.append(folder_info)
                        except Exception as e:
                            self.log_debug(f"Errore nell'analisi della directory {item_path}: {str(e)}")
                            continue
                    
                    # Loop per i file
                    for item in items:
                        if self.stop_search:
                            return
                            
                        item_path = os.path.join(directory, item)
                        
                        # Salta i file nascosti e le directory (già processate)
                        try:
                            # Controllo file nascosto
                            if self.ignore_hidden.get() and (item.startswith('.') or 
                                                (os.name == 'nt' and os.path.exists(item_path) and 
                                                    os.stat(item_path).st_file_attributes & 2)):
                                continue
                                
                            # Salta le directory (già processate)
                            if os.path.isdir(item_path):
                                continue
                        except Exception as e:
                            self.log_debug(f"Errore nel controllo del tipo per {item_path}: {str(e)}")
                            continue
                            
                        # Processa i file
                        if self.search_files.get():
                            try:
                                # Verifica limite file
                                files_checked += 1
                                if files_checked > self.max_files_to_check.get():
                                    self.stop_search = True
                                    self.progress_queue.put(("status", 
                                        f"Limite di {self.max_files_to_check.get()} file controllati raggiunto. "
                                        f"Aumenta il limite nelle opzioni per cercare più file."))
                                    return
                                
                                # MODIFICA: Gestione più sicura dell'executor
                                try:
                                    # Verifica che l'executor sia disponibile
                                    if self.search_executor and not self.search_executor._shutdown:
                                        future = self.search_executor.submit(process_file, item_path, keywords)
                                        futures.append(future)
                                    else:
                                        # Fallback elaborazione diretta
                                        result = process_file(item_path, keywords)
                                        if result:
                                            self.search_results.append(result)
                                except Exception as e:
                                    self.log_debug(f"Errore nell'elaborazione parallela del file {item_path}: {str(e)}")
                                    # Fallback elaborazione diretta
                                    try:
                                        result = process_file(item_path, keywords)
                                        if result:
                                            self.search_results.append(result)
                                    except:
                                        pass
                                
                                # MODIFICATO: Rimosso l'aggiornamento basato sul conteggio
                                # Ora utilizziamo solo l'aggiornamento temporale
                                
                            except Exception as e:
                                self.log_debug(f"Errore nell'aggiunta del file {item_path} alla coda: {str(e)}")
                                continue
                    
                    # Ottimizza ordine directory e visita ricorsivamente
                    directories = self.optimize_disk_search_order(directory, directories)
                    for dir_path in directories:
                        if self.stop_search:
                            return
                            
                        # Visita le sottodirectory con gestione delle eccezioni
                        try:
                            visit_directory(dir_path, current_depth + 1)
                        except Exception as e:
                            self.log_debug(f"Errore nella visita ricorsiva di {dir_path}: {str(e)}")
                            continue
                            
                except PermissionError:
                    self.log_debug(f"Permesso negato per la directory {directory}")
                except Exception as e:
                    self.log_debug(f"Errore durante l'analisi della directory {directory}: {str(e)}")
            
            # Aggiorna lo stato iniziale
            self.progress_queue.put(("status", f"Inizio ricerca in: {path} (Profondità: {'illimitata' if self.max_depth == 0 else self.max_depth})"))
            
            # Avvia la ricerca
            visit_directory(path)
            
            # Ripristina i parametri originali se erano stati modificati
            if is_system_search and self.max_depth == 0 and 'original_max_files' in locals():
                self.max_files_to_check.set(original_max_files)
            
            # Aggiorna lo stato finale di analisi
            self.progress_queue.put(("status", f"Elaborazione risultati... (analizzati {files_checked} file in {dirs_checked} cartelle)"))
            
            # MODIFICA: Raccolta risultati dalle future con aggiornamento temporizzato
            completed = 0
            total_futures = len(futures)
            last_update_time = time.time()
            
            # Gestione sicura delle future
            if self.search_executor and not self.search_executor._shutdown and futures:
                try:
                    for future in concurrent.futures.as_completed(futures):
                        if self.stop_search:
                            break
                            
                        try:
                            result = future.result()
                            if result:
                                self.search_results.append(result)
                            
                            completed += 1
                            
                            # Aggiorna il progresso e il tempo più frequentemente
                            current_time = time.time()
                            if completed % 20 == 0 or current_time - last_update_time > 2:  # Aggiorna ogni 20 file o ogni 2 secondi
                                progress = 90 + min(10, int((completed / max(1, len(futures))) * 10))
                                self.progress_queue.put(("progress", progress))
                                
                                elapsed_time = current_time - start_time
                                self.progress_queue.put(("status", f"Elaborati {completed}/{len(futures)} file (tempo: {int(elapsed_time)}s)"))
                                
                                # Aggiorna anche il tempo totale durante la ricerca
                                if hasattr(self, 'search_start_time'):
                                    now = datetime.now()
                                    time_diff = now - self.search_start_time
                                    # Converti in minuti e secondi
                                    total_seconds = int(time_diff.total_seconds())
                                    minutes = total_seconds // 60
                                    seconds = total_seconds % 60
                                    
                                    if minutes > 0:
                                        total_time_str = f"{minutes}min {seconds}sec"
                                    else:
                                        total_time_str = f"{seconds}sec"
                                        
                                    self.progress_queue.put(("update_total_time", total_time_str))
                                
                                last_update_time = current_time
                        except Exception as e:
                            self.log_debug(f"Errore nell'elaborazione di un risultato: {str(e)}")
                except Exception as e:
                    self.log_debug(f"Errore nella raccolta dei risultati: {str(e)}")
            
            # Completa la ricerca in modo sicuro
            try:
                if self.search_executor and not self.search_executor._shutdown:
                    self.search_executor.shutdown(wait=False)
                self.search_executor = None
            except Exception as e:
                self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                self.search_executor = None
            
            # Ripristina i parametri originali se erano stati modificati per la ricerca di sistema
            if hasattr(self, 'original_system_search_params'):
                try:
                    self.max_files_to_check.set(self.original_system_search_params["max_files"])
                    self.worker_threads.set(self.original_system_search_params["worker_threads"])
                    delattr(self, 'original_system_search_params')
                except:
                    pass

            # Riporta il risultato finale
            elapsed_time = time.time() - start_time
            self.progress_queue.put(("status", 
                f"Ricerca completata! Analizzati {files_checked} file in {dirs_checked} cartelle in {int(elapsed_time)} secondi."))
            
            # Ordina i risultati per tipo e nome
            self.search_results.sort(key=lambda x: (x[0], x[1]))
            
            self.log_debug(f"Ricerca completata. Trovati {len(self.search_results)} risultati")
            self.progress_queue.put(("complete", "Ricerca completata"))
            
        except Exception as e:
            error_msg = f"Si è verificato un errore durante la ricerca: {str(e)}\n{traceback.format_exc()}"
            self.log_debug(error_msg)
            self.progress_queue.put(("error", error_msg))

    def create_file_info(self, file_path):
        """Crea le informazioni del file per la visualizzazione"""
        try:
            file_size = os.path.getsize(file_path)
            modified_time = datetime.fromtimestamp(os.path.getmtime(file_path))
            created_time = datetime.fromtimestamp(os.path.getctime(file_path))
            file_name = os.path.basename(file_path)
            file_extension = os.path.splitext(file_name)[1].lower()
            
            # Determina il tipo di file
            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type:
                file_type = mime_type.split('/')[0].capitalize()
                if file_type == "Application":
                    if "pdf" in mime_type:
                        file_type = "PDF"
                    elif "word" in mime_type or file_extension == ".docx" or file_extension == ".doc":
                        file_type = "Word"
                    elif "excel" in mime_type or file_extension in [".xlsx", ".xls"]:
                        file_type = "Excel"
                    elif "powerpoint" in mime_type or file_extension in [".pptx", ".ppt"]:
                        file_type = "PowerPoint"
                    else:
                        file_type = "Documento"
            else:
                # Fallback basato sull'estensione
                if file_extension in ['.txt', '.md', '.rtf']:
                    file_type = "Testo"
                elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                    file_type = "Immagine"
                elif file_extension in ['.mp3', '.wav', '.ogg', '.flac']:
                    file_type = "Audio"
                elif file_extension in ['.mp4', '.avi', '.mkv', '.mov']:
                    file_type = "Video"
                elif file_extension in ['.exe', '.dll', '.bat']:
                    file_type = "Eseguibile"
                else:
                    file_type = "File"
            
            # Formatta dimensione file
            if file_size < 1024:
                size_str = f"{file_size} B"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.1f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.1f} MB"
            
            return (
                file_type,
                file_name,
                size_str,
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                file_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni del file {file_path}: {str(e)}")
            return (
                "File",
                os.path.basename(file_path),
                "N/A",
                "N/A",
                "N/A",
                file_path
            )

    def create_folder_info(self, folder_path):
        """Crea le informazioni della cartella per la visualizzazione"""
        try:
            modified_time = datetime.fromtimestamp(os.path.getmtime(folder_path))
            created_time = datetime.fromtimestamp(os.path.getctime(folder_path))
            folder_name = os.path.basename(folder_path)
            
            return (
                "Directory",
                folder_name,
                "",  # Le cartelle non hanno dimensione
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                folder_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni della cartella {folder_path}: {str(e)}")
            return (
                "Directory",
                os.path.basename(folder_path),
                "",
                "N/A",
                "N/A",
                folder_path
            )
            
    def should_search_content(self, file_path):
        """Determina se il contenuto del file dovrebbe essere analizzato"""
        ext = os.path.splitext(file_path)[1].lower()
        
        # Verifica direttamente i tipi di file supportati per la ricerca nei contenuti
        return (
            (ext == '.txt') or
            (ext == '.md') or
            (ext == '.csv') or
            (ext == '.html') or
            (ext == '.htm') or
            (ext == '.xml') or
            (ext == '.json') or
            (ext == '.log') or
            (ext in ['.doc', '.docx'] and file_format_support["docx"]) or
            (ext == '.pdf' and file_format_support["pdf"]) or
            (ext in ['.ppt', '.pptx'] and file_format_support["pptx"]) or
            (ext in ['.xls', '.xlsx'] and file_format_support["excel"]) or
            (ext == '.rtf' and file_format_support["rtf"]) or
            (ext == '.odt' and file_format_support["odt"])
        )

    def get_file_content(self, file_path):
        """Ottiene il contenuto del file in base all'estensione"""
        try:
            # Importa i moduli necessari all'inizio della funzione
            import os  # Importa os qui per assicurarsi che sia disponibile in tutti i blocchi
            ext = os.path.splitext(file_path)[1].lower()
            
            # File di testo semplice
            if ext in ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log']:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        return f.read()
                except UnicodeDecodeError:
                    # Fallback su altre codifiche
                    try:
                        with open(file_path, 'r', encoding='latin-1') as f:
                            return f.read()
                    except:
                        self.log_debug(f"Impossibile leggere il file {file_path} con le codifiche standard")
                        return ""
            # Word (.doc)
            elif ext == '.doc' and file_format_support["doc"]:
                try:
                    # Importa i moduli necessari
                    import pythoncom
                    import threading
                    import queue
                    
                    # Crea una coda per il risultato
                    result_queue = queue.Queue()
                    
                    # Funzione che verrà eseguita in un thread separato
                    def read_doc_in_thread():
                        try:
                            pythoncom.CoInitialize()  # Inizializza COM nel nuovo thread
                            word = win32com.client.Dispatch("Word.Application")
                            word.Visible = False
                            
                            # Disabilita tutti gli avvisi e dialoghi
                            word.DisplayAlerts = False
                            word.Options.ConfirmConversions = False
                            word.Options.CheckGrammarAsYouType = False
                            word.Options.CheckSpellingAsYouType = False
                            
                            try:
                                # Verifico che il file esista e ottengo il percorso assoluto e normalizzato
                                abs_path = os.path.abspath(file_path)
                                if not os.path.exists(abs_path):
                                    result_queue.put(f"ERRORE: Il file {abs_path} non esiste")
                                    return
                                
                                # Per sicurezza, converti tutte le barre in backslash in stile Windows
                                win_path = abs_path.replace('/', '\\')
                                
                                # Apri il documento con parametri minimi
                                try:
                                    doc = word.Documents.Open(win_path, ReadOnly=True)
                                    # Metodo 1: prova ad accedere direttamente al testo del documento
                                    text = doc.Content.Text if hasattr(doc, "Content") else ""
                                    result_queue.put(text)
                                    
                                    # Chiudi il documento
                                    doc.Close(SaveChanges=False)
                                except Exception as e:
                                    result_queue.put(f"ERRORE: {str(e)}")
                            except Exception as e:
                                result_queue.put(f"ERRORE: {str(e)}")
                            finally:
                                try:
                                    word.Quit()
                                except:
                                    pass
                                
                            pythoncom.CoUninitialize()
                        except Exception as e:
                            result_queue.put(f"ERRORE: {str(e)}")
                    
                    # Avvia il thread per la lettura con timeout più lungo
                    doc_thread = threading.Thread(target=read_doc_in_thread)
                    doc_thread.daemon = True
                    doc_thread.start()
                    
                    # Attendi il risultato con timeout (15 secondi)
                    doc_thread.join(15)
                    
                    # Se il thread è ancora vivo dopo il timeout, è bloccato
                    if doc_thread.is_alive():
                        self.log_debug(f"Timeout nella lettura del documento Word {file_path}")
                        return ""
                    
                    # Altrimenti prendi il risultato dalla coda
                    if not result_queue.empty():
                        result = result_queue.get()
                        if isinstance(result, str) and result.startswith("ERRORE:"):
                            self.log_debug(f"Errore nella lettura del documento Word {file_path}: {result}")
                            return ""
                        return result
                    return ""
                        
                except Exception as e:
                    self.log_debug(f"Errore generale nella lettura del documento Word {file_path}: {str(e)}")
                    return ""

            # Excel (.xls)
            elif ext == '.xls' and file_format_support["xls"]:
                try:
                    # Usa xlrd invece di win32com per i file .xls, è più stabile
                    import xlrd
                    
                    try:
                        wb = xlrd.open_workbook(file_path)
                        texts = []
                        for sheet_index in range(min(wb.nsheets, 5)):  # Limita a 5 fogli
                            sheet = wb.sheet_by_index(sheet_index)
                            for row_idx in range(min(sheet.nrows, 100)):  # Limita a 100 righe
                                row_texts = []
                                for col_idx in range(sheet.ncols):
                                    try:
                                        cell_value = sheet.cell(row_idx, col_idx).value
                                        if cell_value:
                                            row_texts.append(str(cell_value))
                                    except:
                                        continue
                                if row_texts:
                                    texts.append(" ".join(row_texts))
                        return "\n".join(texts)
                    except Exception as e:
                        self.log_debug(f"Errore nella lettura del foglio Excel {file_path}: {str(e)}")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore generale nella lettura del foglio Excel {file_path}: {str(e)}")
                    return ""
            # Word document
            elif ext in ['.docx', '.doc'] and file_format_support["docx"]:
                try:
                    doc = docx.Document(file_path)
                    return "\n".join([para.text for para in doc.paragraphs])
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del documento Word {file_path}: {str(e)}")
                    return ""
            
            # PDF
            elif ext == '.pdf' and file_format_support["pdf"]:
                try:
                    content = []
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        for page_num in range(min(10, len(reader.pages))):  # Limite di 10 pagine per prestazioni
                            page = reader.pages[page_num]
                            content.append(page.extract_text())
                    return "\n".join(content)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del PDF {file_path}: {str(e)}")
                    return ""
            
            # PowerPoint
            elif ext in ['.pptx', '.ppt'] and file_format_support["pptx"]:
                try:
                    prs = pptx.Presentation(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return "\n".join(text)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura della presentazione {file_path}: {str(e)}")
                    return ""
            
            # Excel
            elif ext in ['.xlsx', '.xls'] and file_format_support["excel"]:
                try:
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    texts = []
                    for sheet_name in wb.sheetnames[:5]:  # Limita a 5 fogli per prestazioni
                        sheet = wb[sheet_name]
                        for row in sheet.iter_rows(max_row=100):  # Limita a 100 righe per prestazioni
                            row_texts = [str(cell.value) for cell in row if cell.value is not None]
                            texts.append(" ".join(row_texts))
                    return "\n".join(texts)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del foglio Excel {file_path}: {str(e)}")
                    return ""
            
            # RTF
            elif ext == '.rtf' and file_format_support["rtf"]:
                try:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        rtf_text = f.read()
                    return rtf_to_text(rtf_text)
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file RTF {file_path}: {str(e)}")
                    return ""
                    
            # ODT (OpenDocument Text)
            elif ext == '.odt' and file_format_support["odt"]:
                try:
                    # Prova prima con odfdo
                    if 'odfdo' in sys.modules:
                        doc = odfdo.Document(file_path)
                        return doc.get_formatted_text()
                    # Fallback su odf
                    else:
                        textdoc = opendocument.load(file_path)
                        allparas = textdoc.getElementsByType(text.P)
                        return "\n".join([teletype.extractText(para) for para in allparas])
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file ODT {file_path}: {str(e)}")
                    return ""
            return ""  # Tipo di file non supportato
        except Exception as e:
            self.log_debug(f"Errore generale nella lettura del file {file_path}: {str(e)}")
            return ""
        
    def update_progress(self):
        if self.is_searching:
            try:
                while True:
                    progress_type, value = self.progress_queue.get_nowait()
                    if progress_type == "update_total_time":
                        self.total_time_label.config(text=value)
                    if progress_type == "progress":
                        self.progress_bar["value"] = value
                    elif progress_type == "status":
                        # Se il messaggio contiene informazioni sui file, separiamo le informazioni
                        if "analizzati" in value.lower() or "cartelle:" in value.lower() or "file:" in value.lower():
                            # Estrai il percorso se presente
                            if "analisi:" in value.lower():
                                parts = value.split("(", 1)
                                if len(parts) > 1:
                                    path_part = parts[0].strip()
                                    counts_part = "(" + parts[1]
                                    self.status_label["text"] = path_part
                                    self.analyzed_files_label["text"] = counts_part
                                else:
                                    self.analyzed_files_label["text"] = value
                            else:
                                self.analyzed_files_label["text"] = value
                        else:
                            # È solo un messaggio di stato semplice
                            self.status_label["text"] = value
                    elif progress_type == "complete":
                        self.is_searching = False
                        self.search_button["state"] = "normal"
                        self.stop_button["state"] = "disabled"
                        
                        # Aggiorna la lista dei risultati
                        self.update_results_list()
                        
                        # Aggiorna l'orario di fine e il tempo totale
                        current_time = datetime.now().strftime('%H:%M')
                        self.end_time_label.config(text=current_time)
                        self.update_total_time()  # Calcola e mostra il tempo totale
                        
                        if len(self.search_results) == 0:
                            self.status_label["text"] = "Nessun file trovato per la ricerca effettuata"
                            self.root.after(100, lambda: messagebox.showinfo("Ricerca completata", "Nessun file trovato per la ricerca effettuata"))
                        else:
                            self.status_label["text"] = f"Ricerca completata! Trovati {len(self.search_results)} risultati."
                        
                        self.progress_bar["value"] = 100
                        return
                    elif progress_type == "error":
                        self.is_searching = False
                        self.search_button["state"] = "normal"
                        self.stop_button["state"] = "disabled"
                        current_time = datetime.now().strftime('%H:%M')
                        self.end_time_label.config(text=current_time)
                        self.update_total_time()  # Calcola e mostra il tempo totale
                        messagebox.showerror("Errore", value)
                        return
                    elif progress_type == "timeout":
                        self.is_searching = False
                        self.search_button["state"] = "normal"
                        self.stop_button["state"] = "disabled"
                        self.update_results_list()
                        current_time = datetime.now().strftime('%H:%M')
                        self.end_time_label.config(text=current_time)
                        self.update_total_time()  # Calcola e mostra il tempo totale
                        messagebox.showinfo("Timeout", "La ricerca è stata interrotta per timeout. Verranno mostrati i risultati parziali trovati.")
                        return
            except queue.Empty:
                pass
                    
            self.root.after_idle(self.update_progress)
            
    def stop_search_process(self):
        """Ferma il processo di ricerca in corso"""
        self.stop_search = True
        self.status_label["text"] = "Interrompendo la ricerca..."
        self.analyzed_files_label["text"] = "Ricerca interrotta dall'utente"
        current_time = datetime.now().strftime('%H:%M')
        self.end_time_label.config(text=current_time)
        self.update_total_time()  # Calcola e mostra il tempo totale
        
        # Chiudi l'executor se esiste
        if self.search_executor:
            self.search_executor.shutdown(wait=False, cancel_futures=True)
        
        # Aggiorna l'orario di fine quando la ricerca viene interrotta
        self.search_button["state"] = "normal"
        self.stop_button["state"] = "disabled"

    def update_results_list(self):
        """Aggiorna la lista dei risultati con i risultati trovati"""
        # Pulisci la lista attuale
        for item in self.results_list.get_children():
            self.results_list.delete(item)
            
        # Aggiungi i risultati alla lista
        for result in self.search_results:
            item_type, name, size, modified, created, path = result
            
            # Applica stile in base al tipo di elemento
            if item_type == "Directory":
                tags = ("directory",)
            else:
                tags = ("file",)
                
            self.results_list.insert("", "end", values=result, tags=tags)
            
        # Aggiorna lo stato
        self.status_label["text"] = f"Trovati {len(self.search_results)} risultati"
        
    def update_theme_colors(self, theme="light"):
        """Aggiorna i colori del tema per evidenziare cartelle e file"""
        style = ttk.Style()
        
        # Configura i colori per il tema chiaro
        if theme == "light":
            style.configure("Treeview", background="#ffffff", foreground="#000000", fieldbackground="#ffffff")
            self.results_list.tag_configure("directory", background="#e6f2ff", foreground="#000000")
            self.results_list.tag_configure("file", background="#ffffff", foreground="#000000")
        # Configura i colori per il tema scuro
        elif theme == "dark":
            style.configure("Treeview", background="#333333", foreground="#ffffff", fieldbackground="#333333")
            self.results_list.tag_configure("directory", background="#4d4d4d", foreground="#ffffff")
            self.results_list.tag_configure("file", background="#333333", foreground="#ffffff")


    def copy_selected(self):
        """Copia i file selezionati in una directory di destinazione"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da copiare")
            return
            
        # Chiedi la directory di destinazione
        dest_dir = filedialog.askdirectory(title="Seleziona la cartella di destinazione")
        if not dest_dir:
            return
            
        # Prepara variabili per tracciare l'avanzamento
        total = len(selected_items)
        copied = 0
        failed = 0
        skipped = 0
        
        try:
            for item in selected_items:
                values = self.results_list.item(item)['values']
                item_type, _, _, _, _, source_path = values
                
                # Ottieni il nome dell'elemento senza il percorso completo
                basename = os.path.basename(source_path)
                dest_path = os.path.join(dest_dir, basename)
                
                try:
                    if item_type == "Directory":
                        # Se la cartella esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("Cartella esistente", 
                                                f"La cartella {basename} esiste già. Vuoi sovrascriverla?"):
                                # Elimina la cartella esistente
                                shutil.rmtree(dest_path)
                            else:
                                skipped += 1
                                continue
                                
                        # Copia ricorsiva della cartella
                        shutil.copytree(source_path, dest_path)
                        copied += 1
                    else:
                        # Se il file esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("File esistente", 
                                                f"Il file {basename} esiste già. Vuoi sovrascriverlo?"):
                                # Continua con la sovrascrittura
                                pass
                            else:
                                skipped += 1
                                continue
                                
                        # Copia il file
                        shutil.copy2(source_path, dest_path)
                        copied += 1
                        
                except Exception as e:
                    failed += 1
                    self.log_debug(f"Errore durante la copia di {source_path}: {str(e)}")
                    
                # Aggiorna la barra di avanzamento
                progress = ((copied + failed + skipped) / total) * 100
                self.progress_bar["value"] = progress
                self.status_label["text"] = f"Copiati {copied}/{total} elementi..."
                self.root.update()
                
            # Messaggio di completamento
            if failed > 0:
                messagebox.showwarning("Copia completata con errori", 
                                    f"Copiati: {copied}\nSaltati: {skipped}\nFalliti: {failed}")
            else:
                messagebox.showinfo("Copia completata", 
                                 f"Copiati con successo: {copied}\nSaltati: {skipped}")
                
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'operazione: {str(e)}")
            
        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    def compress_selected(self):
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da comprimere")
            return
            
        zip_name = Querybox.get_string(
            prompt="Inserisci il nome del file ZIP (senza estensione):",
            title="Nome file ZIP",
            initialvalue="archivio"
        )
        
        if not zip_name:
            return
            
        zip_path = filedialog.asksaveasfilename(
            defaultextension=".zip",
            initialfile=f"{zip_name}.zip",
            filetypes=[("ZIP files", "*.zip")],
            title="Salva file ZIP"
        )
        
        if not zip_path:
            return
            
        # Raccogli tutti i file delle cartelle selezionate
        files_in_folders = set()
        folder_paths = []
        single_files = []
        
        # Prima fase: raccogli informazioni su cartelle e file
        for item in selected_items:
            values = self.results_list.item(item)['values']
            item_type, _, _, _, _, source_path = values
            
            if item_type == "Directory":
                folder_paths.append(source_path)
                # Raccogli tutti i file nelle cartelle selezionate
                for root, _, files in os.walk(source_path):
                    for file in files:
                        full_path = os.path.join(root, file)
                        files_in_folders.add(os.path.abspath(full_path))
            else:
                single_files.append(source_path)
        
        # Filtra i file singoli che sono già presenti nelle cartelle
        filtered_single_files = [f for f in single_files if os.path.abspath(f) not in files_in_folders]
        
        total_items = len(folder_paths) + len(filtered_single_files)
        processed = 0
        
        try:
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                # Prima comprimi le cartelle
                for folder_path in folder_paths:
                    base_folder = os.path.basename(folder_path)
                    for root, _, files in os.walk(folder_path):
                        for file in files:
                            file_path = os.path.join(root, file)
                            # Mantieni la struttura delle cartelle nel zip
                            rel_path = os.path.relpath(file_path, os.path.dirname(folder_path))
                            zipf.write(file_path, rel_path)
                    
                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()
                
                # Poi comprimi i file singoli (solo quelli non già presenti nelle cartelle)
                for file_path in filtered_single_files:
                    if os.path.exists(file_path):  # Verifica che il file esista ancora
                        zipf.write(file_path, os.path.basename(file_path))
                        
                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()
            
            # Prepara il messaggio di completamento
            skipped_files = len(single_files) - len(filtered_single_files)
            message = f"Compressione completata!\nFile salvato in: {zip_path}"
            if skipped_files > 0:
                message += f"\n{skipped_files} file saltati perché già presenti nelle cartelle"
                
            messagebox.showinfo("Completato", message)
            
        except Exception as e:
            messagebox.showerror("Errore", f"Errore durante la compressione: {str(e)}")
            
        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    def show_advanced_filters_dialog(self):
        """Mostra la finestra di dialogo per i filtri di ricerca avanzati"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Filtri avanzati")
        dialog.geometry("530x480")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Filtri dimensione
        size_frame = ttk.LabelFrame(dialog, text="Dimensione file")
        size_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(size_frame, text="Min (KB):").grid(row=0, column=0, padx=5, pady=5)
        min_size = ttk.Entry(size_frame, width=10)
        min_size.grid(row=0, column=1, padx=5, pady=5)
        min_size.insert(0, str(self.advanced_filters["size_min"] // 1024))
        
        ttk.Label(size_frame, text="Max (KB):").grid(row=0, column=2, padx=5, pady=5)
        max_size = ttk.Entry(size_frame, width=10)
        max_size.grid(row=0, column=3, padx=5, pady=5)
        max_size.insert(0, str(self.advanced_filters["size_max"] // 1024 if self.advanced_filters["size_max"] else 0))
        
        # Filtri data - FORMAT DD-MM-YYYY
        date_frame = ttk.LabelFrame(dialog, text="Data modifica (DD-MM-YYYY)")
        date_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(date_frame, text="Da:").grid(row=0, column=0, padx=5, pady=5)
        min_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        min_date.grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(date_frame, text="A:").grid(row=0, column=2, padx=5, pady=5)
        max_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        max_date.grid(row=0, column=3, padx=5, pady=5)
        
        # Prepopola le date se disponibili
        if self.advanced_filters["date_min"]:
            min_date.entry.delete(0, "end")
            min_date.entry.insert(0, self.advanced_filters["date_min"])
                
        if self.advanced_filters["date_max"]:
            max_date.entry.delete(0, "end")
            max_date.entry.insert(0, self.advanced_filters["date_max"])
        
        # Filtri estensione
        ext_frame = ttk.LabelFrame(dialog, text="Estensioni file (separate da virgola)")
        ext_frame.pack(fill=X, padx=10, pady=5)
        example_label = ttk.Label(ext_frame, text="Inserisci una o più estensioni da ricercare. Esempio: .pdf, .dot", 
                       font=("", 8), foreground="gray")
        example_label.pack(anchor="w", padx=5)
        
        extensions = ttk.Entry(ext_frame)
        extensions.pack(fill=X, padx=5, pady=5)
        if self.advanced_filters["extensions"]:
            extensions.insert(0, ", ".join(self.advanced_filters["extensions"]))
        
        # Aggiungiamo un frame di debug per vedere i filtri correnti
        debug_frame = ttk.LabelFrame(dialog, text="Debug - Stato filtri correnti")
        debug_frame.pack(fill=X, padx=10, pady=5)
        
        debug_text = ttk.Text(debug_frame, height=3, width=50)
        debug_text.pack(fill=X, padx=5, pady=5)
        debug_text.insert("1.0", f"Data min: {self.advanced_filters['date_min']}\n")
        debug_text.insert("2.0", f"Data max: {self.advanced_filters['date_max']}\n")
        debug_text.insert("3.0", f"Extensions: {self.advanced_filters['extensions']}")
        debug_text.config(state="disabled")
        
        # Pulsante Salva
        def save_filters():
            try:
                # Analizza i filtri di dimensione
                min_kb = int(min_size.get() or 0)
                max_kb = int(max_size.get() or 0)
                self.advanced_filters["size_min"] = min_kb * 1024
                self.advanced_filters["size_max"] = max_kb * 1024
                
                # Ottieni le date nel formato DD-MM-YYYY
                min_date_value = min_date.entry.get().strip()
                max_date_value = max_date.entry.get().strip()
                
                print(f"DEBUG - Date inserite: min={min_date_value}, max={max_date_value}")
                
                # Validazione date
                if min_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(min_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data minima non valido. Usa DD-MM-YYYY")
                        return
                
                if max_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(max_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data massima non valido. Usa DD-MM-YYYY")
                        return
                
                # Verifica che la data minima non sia successiva alla data massima
                if min_date_value and max_date_value:
                    min_date_obj = datetime.strptime(min_date_value, "%d-%m-%Y")
                    max_date_obj = datetime.strptime(max_date_value, "%d-%m-%Y")
                    
                    if min_date_obj > max_date_obj:
                        messagebox.showerror("Errore", 
                                            "La data di inizio non può essere successiva alla data di fine")
                        return
                
                # Salva le date validate
                self.advanced_filters["date_min"] = min_date_value
                self.advanced_filters["date_max"] = max_date_value
                
                # Analizza le estensioni
                exts = [e.strip() for e in extensions.get().split(",") if e.strip()]
                self.advanced_filters["extensions"] = [f".{e.lstrip('.')}" for e in exts]
                
                print(f"Filtri salvati: {self.advanced_filters}")  # Debug info
                dialog.destroy()
                
            except ValueError:
                messagebox.showerror("Errore", "Inserisci valori numerici validi per le dimensioni")
        
        ttk.Button(dialog, text="Salva", command=save_filters).pack(pady=10)
        ttk.Button(dialog, text="Annulla", command=dialog.destroy).pack(pady=5)

        dialog.update_idletasks()  # Aggiorna la finestra per ottenere le dimensioni corrette
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima per la finestra
        dialog.minsize(500, 400)

    def create_widgets(self):
        # Frame principale che conterrà tutto tranne la barra di stato
        main_container = ttk.Frame(self.root)
        main_container.pack(fill=BOTH, expand=YES)
        
        # Frame principale per il contenuto
        main_frame = ttk.Frame(main_container, padding="10")
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Header con titolo e info utente
        header_frame = ttk.Frame(main_frame)
        header_frame.pack(fill=X, pady=(0, 10))
        
        # Titolo a sinistra e info data/utente a destra
        title_label = ttk.Label(header_frame, text="File Search Tool.. Nucleo Perugia", 
                            font=("Helvetica", 14, "bold"))
        title_label.pack(side=LEFT)
        
        datetime_label = ttk.Label(header_frame, textvariable=self.datetime_var, font=("Helvetica", 9))
        datetime_label.pack(side=RIGHT)
        
        # Selezione tema - posizionata prima dei parametri di ricerca
        theme_frame = ttk.Frame(main_frame)
        theme_frame.pack(fill=X, pady=5)
        
        ttk.Label(theme_frame, text="Tema:").pack(side=LEFT, padx=5)
        themes = ttk.Style().theme_names()
        theme_combobox = ttk.Combobox(theme_frame, values=themes, width=15)
        theme_combobox.pack(side=LEFT, padx=5)
        theme_combobox.current(themes.index("darkly"))
        theme_combobox.bind("<<ComboboxSelected>>", lambda e: [ttk.Style().theme_use(theme_combobox.get()),self.update_theme_colors()])
        
        # Utilizziamo un pannello a due colonne per ottimizzare lo spazio
        main_panel = ttk.PanedWindow(main_frame, orient=HORIZONTAL)
        main_panel.pack(fill=BOTH, expand=YES, pady=5)
        
        # Pannello sinistro: versione semplice senza etichetta
        left_frame = ttk.Frame(main_panel)
        main_panel.add(left_frame, weight=50)
        
        # Pannello destro: contiene i risultati
        right_frame = ttk.LabelFrame(main_panel, text="Risultati")
        main_panel.add(right_frame, weight=60)
        
        # ---- Configurazione pannello sinistro ----
        
        # Directory frame - più compatto
        path_frame = ttk.LabelFrame(left_frame, text="Directory di ricerca", padding=5)
        path_frame.pack(fill=X, pady=5)
        
        # Label di aiuto
        example_label = ttk.Label(path_frame, text="Seleziona la directory per effettuare la ricerca dei file", 
                                font=("", 8), foreground="gray")
        example_label.pack(anchor="w", padx=5)
        
        entry_browse_frame = ttk.Frame(path_frame)
        entry_browse_frame.pack(fill=X, pady=5)
        
        self.path_entry = ttk.Entry(entry_browse_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        browse_btn = ttk.Button(entry_browse_frame, text="Sfoglia", command=self.browse_directory)
        browse_btn.pack(side=LEFT)
        
        # Keywords frame
        keyword_frame = ttk.LabelFrame(left_frame, text="Parole chiave", padding=5)
        keyword_frame.pack(fill=X, pady=5)
        
        # Label di aiuto per keywords
        example_label = ttk.Label(keyword_frame, text="Per la ricerca di più parole usa la virgola. Esempio: documento, fattura, contratto", 
                                font=("", 8), foreground="gray")
        example_label.pack(anchor="w", padx=5)
        
        # Crea l'entry delle parole chiave
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(fill=X, padx=5, pady=5)
        
        # Frame per i filtri di ricerca principali
        options_frame = ttk.LabelFrame(left_frame, text="Opzioni di ricerca", padding=5)
        options_frame.pack(fill=X, pady=5)
        
        example_label = ttk.Label(options_frame, text="Seleziona una o più opzioni per la ricerca", 
                                font=("", 8), foreground="gray")
        example_label.pack(anchor="w", padx=5)
        
         # Prima riga di opzioni
        options_row1 = ttk.Frame(options_frame)
        options_row1.pack(fill=X, pady=2)
        
        ttk.Checkbutton(options_row1, text="Cerca file", 
                        variable=self.search_files).pack(side=LEFT, padx=5)
        ttk.Checkbutton(options_row1, text="Cerca cartelle", 
                        variable=self.search_folders).pack(side=LEFT, padx=5)
        content_checkbox = ttk.Checkbutton(options_row1, text="Cerca nei contenuti", 
                                        variable=self.search_content, 
                                        command=self.warn_content_search)
        content_checkbox.pack(side=LEFT, padx=5)
        self.create_tooltip(content_checkbox, "Cerca le parole chiave anche all'interno dei file di testo")
        
        # Pulsante filtri avanzati spostato qui, dopo "Cerca nei contenuti"
        ttk.Button(options_row1, text="Filtri avanzati", 
                command=self.show_advanced_filters_dialog).pack(side=LEFT, padx=10)
        
        # Seconda riga per profondità
        depth_frame = ttk.Frame(options_frame)
        depth_frame.pack(fill=X, pady=5, padx=5, anchor=W)
        
        ttk.Label(depth_frame, text="Profondità max:").pack(side=LEFT)
        self.depth_spinbox = ttk.Spinbox(depth_frame, from_=0, to=10, width=3)
        self.depth_spinbox.pack(side=LEFT, padx=5)
        self.depth_spinbox.set("0")
        ttk.Label(depth_frame, text="(0 = illimitata)", 
                font=("", 8), foreground="gray").pack(side=LEFT)
               
        # Frame per le opzioni di performance
        perf_frame = ttk.LabelFrame(left_frame, text="Opzioni di Performance", padding=5)
        perf_frame.pack(fill=X, pady=5)
        
        # Grid per organizzare le opzioni di performance in modo più compatto
        perf_grid = ttk.Frame(perf_frame)
        perf_grid.pack(fill=X)
        
        # Prima riga: timeout e threads
        ttk.Checkbutton(perf_grid, text="Timeout ricerca", 
                    variable=self.timeout_enabled).grid(row=0, column=0, sticky=W)
        
        timeout_frame = ttk.Frame(perf_grid)
        timeout_frame.grid(row=0, column=1, sticky=W)
        ttk.Label(timeout_frame, text="Secondi:").pack(side=LEFT)
        timeout_spinbox = ttk.Spinbox(timeout_frame, from_=10, to=600, width=4, 
                                    textvariable=self.timeout_seconds)
        timeout_spinbox.pack(side=LEFT, padx=5)
        
        threads_frame = ttk.Frame(perf_grid)
        threads_frame.grid(row=0, column=2, sticky=W, padx=(10,0))
        ttk.Label(threads_frame, text="Thread:").pack(side=LEFT)
        threads_spinbox = ttk.Spinbox(threads_frame, from_=1, to=16, width=2, 
                                    textvariable=self.worker_threads)
        threads_spinbox.pack(side=LEFT, padx=5)
        
        # Seconda riga: dimensione file e max risultati
        filesize_frame = ttk.Frame(perf_grid)
        filesize_frame.grid(row=1, column=0, sticky=W, pady=5)
        ttk.Label(filesize_frame, text="Max file MB:").pack(side=LEFT)
        max_size_spinbox = ttk.Spinbox(filesize_frame, from_=1, to=1000, width=4,
                                    textvariable=self.max_file_size_mb)
        max_size_spinbox.pack(side=LEFT, padx=5)
        
        max_results_frame = ttk.Frame(perf_grid)
        max_results_frame.grid(row=1, column=1, columnspan=2, sticky=W, pady=5)
        ttk.Label(max_results_frame, text="Max risultati:").pack(side=LEFT)
        max_results_spinbox = ttk.Spinbox(max_results_frame, from_=100, to=100000, width=6,
                                        textvariable=self.max_results)
        max_results_spinbox.pack(side=LEFT, padx=5)
        
        # Checkbox per indicizzazione
        ttk.Checkbutton(perf_grid, text="Indicizzazione", 
                    variable=self.use_indexing).grid(row=2, column=0, columnspan=3, sticky=W)
        
        # Pulsanti di azione
        button_frame = ttk.Frame(left_frame)
        button_frame.pack(pady=10, fill=X)
        
        # Pulsante di ricerca
        self.search_button = ttk.Button(button_frame, text="Cerca", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=12)
        self.search_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.search_button, "Avvia la ricerca con i criteri specificati")
        
        # Pulsante per interrompere la ricerca
        self.stop_button = ttk.Button(button_frame, 
                                    text="Interrompi ricerca",
                                    command=self.stop_search_process,
                                    style="danger.TButton", width=15,
                                    state="disabled")
        self.stop_button.pack(side=LEFT, padx=5)
        
        # Pulsante per pulire i campi di ricerca
        ttk.Button(button_frame, 
                text="Pulisci campi", 
                command=lambda: [self.search_path.set(""), self.keywords.set("")],
                style="secondary.Outline.TButton", width=12).pack(side=LEFT, padx=5)
        
        # Stato ricerca
        status_container = ttk.Frame(left_frame)
        status_container.pack(fill=X, pady=5)
        
        # Frame per i tempi di ricerca
        time_frame = ttk.LabelFrame(status_container, text="Informazioni temporali", padding=5)
        time_frame.pack(fill=X, pady=5)
        
        times_grid = ttk.Frame(time_frame)
        times_grid.pack(fill=X)
        
        # Organizziamo il tempo usando un layout a griglia per allineamento migliore
        times_grid.columnconfigure(1, weight=1)
        times_grid.columnconfigure(3, weight=1)
        
        ttk.Label(times_grid, text="Avvio ricerca:", 
                font=("", 9)).grid(row=0, column=0, sticky=W, padx=(5,2))
        self.start_time_label = ttk.Label(times_grid, text="--:--", 
                                        font=("", 9, "bold"))
        self.start_time_label.grid(row=0, column=1, sticky=W)
        
        ttk.Label(times_grid, text="Fine ricerca:", 
                font=("", 9)).grid(row=0, column=2, sticky=W, padx=(15,2))
        self.end_time_label = ttk.Label(times_grid, text="--:--", 
                                    font=("", 9, "bold"))
        self.end_time_label.grid(row=0, column=3, sticky=W)
        
        ttk.Label(times_grid, text="Durata:", 
                font=("", 9)).grid(row=0, column=4, sticky=W, padx=(25,2))
        self.total_time_label = ttk.Label(times_grid, text="--:--", 
                                    font=("", 9, "bold"))
        self.total_time_label.grid(row=0, column=5, sticky=W)


        # Frame per lo stato dell'analisi
        analysis_frame = ttk.LabelFrame(status_container, text="Stato analisi", padding=5)
        analysis_frame.pack(fill=X, pady=5)
        
        # Creiamo un layout a griglia per una separazione più chiara
        analysis_frame.columnconfigure(0, weight=1)  # Colonna espandibile
        
        # Label per il percorso corrente dell'analisi
        ttk.Label(analysis_frame, text="Analisi:", font=("", 9, "bold")).grid(row=0, column=0, sticky="w", padx=5, pady=(5, 0))
        self.status_label = ttk.Label(analysis_frame, text="In attesa...", justify=LEFT, wraplength=1000)
        self.status_label.grid(row=1, column=0, sticky="w", padx=5, pady=2)
        
        # Separatore orizzontale tra le due sezioni
        ttk.Separator(analysis_frame, orient='horizontal').grid(row=2, column=0, sticky="ew", pady=5, padx=5)
        
        # Label per il contatore dei file analizzati
        ttk.Label(analysis_frame, text="File analizzati:", font=("", 9, "bold")).grid(row=3, column=0, sticky="w", padx=5, pady=(0, 0))
        self.analyzed_files_label = ttk.Label(analysis_frame, text="Nessuna ricerca avviata", justify=LEFT)
        self.analyzed_files_label.grid(row=4, column=0, sticky="w", padx=5, pady=(0, 5))
        
        # Progress bar
        progress_frame = ttk.Frame(left_frame)
        progress_frame.pack(fill=X, pady=5)
        
        self.progress_bar = ttk.Progressbar(progress_frame, mode='determinate')
        self.progress_bar.pack(fill=X)
        
        # ---- Configurazione pannello destro (risultati) ----
        
         # Frame per i pulsanti di azione
        action_buttons_frame = ttk.Frame(right_frame)
        action_buttons_frame.pack(fill=X, pady=(0, 5))
        
        # Pulsanti per la selezione in un sottogruppo
        selection_frame = ttk.Frame(action_buttons_frame)
        selection_frame.pack(side=LEFT, fill=Y)
        
        ttk.Button(selection_frame, text="Seleziona tutto", 
                command=self.select_all).pack(side=LEFT, padx=2)
        ttk.Button(selection_frame, text="Deseleziona tutto", 
                command=self.deselect_all).pack(side=LEFT, padx=2)
        ttk.Button(selection_frame, text="Inverti selezione", 
                command=self.invert_selection).pack(side=LEFT, padx=2)
        
        # Scrollbar per la lista risultati
        scrollbar = ttk.Scrollbar(right_frame)
        scrollbar.pack(side=RIGHT, fill=Y)
        
        # TreeView per i risultati
        self.results_list = ttk.Treeview(right_frame, selectmode="extended",
                                        columns=("type", "author", "size", "modified", "created", "path"),
                                        show="headings",
                                        yscrollcommand=scrollbar.set)
        
        # Imposta le colonne con larghezze ottimizzate
        self.results_list.column("type", width=80, anchor="center")
        self.results_list.column("author", width=200, anchor="w")
        self.results_list.column("size", width=80, anchor="center")
        self.results_list.column("modified", width=120, anchor="center")
        self.results_list.column("created", width=120, anchor="center")
        self.results_list.column("path", width=350, anchor="w")
        
        # Imposta le intestazioni
        self.results_list.heading("type", text="Tipo File")
        self.results_list.heading("author", text="Nome File")
        self.results_list.heading("size", text="Dimensione")
        self.results_list.heading("modified", text="Modificato")
        self.results_list.heading("created", text="Creato")
        self.results_list.heading("path", text="Percorso")
        
        self.results_list.pack(fill=BOTH, expand=YES)
        
        # Configura scrollbar
        scrollbar.config(command=self.results_list.yview)
        
        # Frame per i pulsanti di azione principali sotto la lista
        main_buttons_frame = ttk.Frame(right_frame)
        main_buttons_frame.pack(fill=X, pady=5)
        
        # Pulsanti copia e comprimi
        self.copy_button = ttk.Button(main_buttons_frame, text="Copia selezionati",
                                    command=self.copy_selected,
                                    style="TButton")
        self.copy_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.copy_button, "Copia i file selezionati nella directory specificata")
        
        self.compress_button = ttk.Button(main_buttons_frame, text="Comprimi selezionati",
                                        command=self.compress_selected,
                                        style="TButton")
        self.compress_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.compress_button, "Comprimi i file selezionati in un archivio ZIP")

        # Applica stili alle righe
        self.update_theme_colors()
    
    # Avviso sulla spunta della ricerca dei contenuti
    def warn_content_search(self):
        """Mostra un avviso se l'utente attiva la ricerca nei contenuti"""
        if self.search_content.get():  # Se la checkbox è stata appena attivata
            messagebox.showwarning(
                "Attenzione", 
                "Attivando questa opzione rallenterà la ricerca e l'attesa sarà maggiore.",
                parent=self.root
            )

    def create_tooltip(self, widget, text, delay=500, fade=True):
        """Crea tooltip con ritardo, effetti di dissolvenza e larghezza automatica"""
        
        tooltip_timer = None
        fade_timer = None
        
        def show_tooltip():
            x = widget.winfo_rootx() + 25
            y = widget.winfo_rooty() + 25
            
            tooltip = ttk.Toplevel(widget)
            tooltip.wm_overrideredirect(True)
            tooltip.attributes("-alpha", 0.0)  # Inizia invisibile
            tooltip.attributes("-topmost", True)  # Mantiene sopra altre finestre
            
            # Crea un frame con bordo e padding
            frame = ttk.Frame(tooltip, borderwidth=1, relief="solid", padding=10)
            frame.pack(fill="both", expand=True)
            
            # Calcolo della larghezza del testo
            font = ("Segoe UI", 9)  # Puoi modificare questo per usare un altro font
            
            # Calcola la lunghezza approssimativa del testo
            import tkinter as tk
            temp_label = tk.Label(font=font)
            
            # Determina se il testo necessita di essere diviso in più righe
            lines = text.split("\n")
            max_line_width = 0
            
            for line in lines:
                temp_label.configure(text=line)
                line_width = temp_label.winfo_reqwidth()
                max_line_width = max(max_line_width, line_width)
            
            # Limita la larghezza massima a 500 pixel
            max_width = min(max_line_width, 500)
            
            # Se il testo è più corto della larghezza massima, non imposta wraplength
            # altrimenti imposta wraplength per dividere il testo in righe
            if max_line_width > 500:
                wraplength = 500
            else:
                wraplength = max_line_width
                
            # Crea l'etichetta con il testo
            label = ttk.Label(frame, text=text, justify="left", 
                            wraplength=wraplength, font=font)
            label.pack(fill="both", expand=True)
            
            # Distruggi il label temporaneo
            temp_label.destroy()
            
            # Aggiorna immediatamente per calcolare le dimensioni
            tooltip.update_idletasks()
            
            # Posiziona il tooltip in modo che non vada fuori dallo schermo
            tooltip_width = tooltip.winfo_reqwidth()
            tooltip_height = tooltip.winfo_reqheight()
            
            screen_width = tooltip.winfo_screenwidth()
            screen_height = tooltip.winfo_screenheight()
            
            # Aggiusta la posizione se il tooltip esce dallo schermo
            if x + tooltip_width > screen_width:
                x = screen_width - tooltip_width - 10
            
            if y + tooltip_height > screen_height:
                y = screen_height - tooltip_height - 10
                
            # Imposta la posizione finale
            tooltip.wm_geometry(f"+{x}+{y}")
            
            widget._tooltip = tooltip
            
            if fade:
                # Effetto dissolvenza in entrata
                def fade_in(alpha=0.0):
                    if not hasattr(widget, "_tooltip"):
                        return
                    
                    tooltip.attributes("-alpha", alpha)
                    if alpha < 1.0:
                        nonlocal fade_timer
                        fade_timer = widget.after(20, lambda: fade_in(alpha + 0.1))
                
                fade_in()
        
        def enter(event):
            nonlocal tooltip_timer
            # Avvia il timer per mostrare il tooltip dopo un certo ritardo
            tooltip_timer = widget.after(delay, show_tooltip)
        
        def leave(event):
            nonlocal tooltip_timer, fade_timer
            
            # Cancella il timer se esiste
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
                tooltip_timer = None
            
            # Cancella il timer di dissolvenza se esiste
            if fade_timer:
                widget.after_cancel(fade_timer)
                fade_timer = None
            
            # Rimuovi il tooltip se esiste
            if hasattr(widget, "_tooltip"):
                if fade:
                    # Effetto dissolvenza in uscita
                    def fade_out(alpha=1.0):
                        if alpha <= 0 or not hasattr(widget, "_tooltip"):
                            if hasattr(widget, "_tooltip"):
                                widget._tooltip.destroy()
                                del widget._tooltip
                        else:
                            widget._tooltip.attributes("-alpha", alpha)
                            nonlocal fade_timer
                            fade_timer = widget.after(20, lambda: fade_out(alpha - 0.1))
                    
                    fade_out()
                else:
                    widget._tooltip.destroy()
                    del widget._tooltip
        
        # Gestisce anche il caso in cui il widget venga distrutto
        def on_destroy(event):
            nonlocal tooltip_timer, fade_timer
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
            if fade_timer and hasattr(widget, "_tooltip"):
                widget.after_cancel(fade_timer)
                widget._tooltip.destroy()
        
        widget.bind("<Enter>", enter)
        widget.bind("<Leave>", leave)
        widget.bind("<Destroy>", on_destroy)

    def select_all(self):
        self.results_list.selection_set(self.results_list.get_children())
        
    def deselect_all(self):
        self.results_list.selection_remove(self.results_list.get_children())
        
    def invert_selection(self):
        all_items = self.results_list.get_children()
        selected_items = self.results_list.selection()
        self.results_list.selection_remove(selected_items)
        to_select = set(all_items) - set(selected_items)
        for item in to_select:
            self.results_list.selection_add(item)

# Funzione principale per eseguire l'applicazione
def main():
    import sys
    root = ttk.Window(themename="darkly")
    app = FileSearchApp(root)
    
    # Verifica se ci sono argomenti da linea di comando
    if len(sys.argv) > 1:
        # Se c'è un percorso fornito come argomento, lo imposta come percorso di ricerca
        app.search_path.set(sys.argv[1])
    
    root.mainloop()

if __name__ == "__main__":
    main()
