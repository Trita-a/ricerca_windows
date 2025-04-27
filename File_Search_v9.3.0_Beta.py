import concurrent.futures
import csv
import functools
import gc
import getpass
import hashlib
import io
import json
import mimetypes
import os
import platform
import queue
import re
import shutil
import signal
import subprocess
import tarfile
import tempfile
import threading
import time
import tkinter as tk
import traceback
import uuid
import webbrowser
import zipfile
from datetime import datetime
from tkinter import filedialog, messagebox, BooleanVar, StringVar, IntVar

# Third-party imports
import psutil
import pythoncom
import rarfile
import ttkbootstrap as ttk
import win32com.client
from ttkbootstrap.constants import *

# Import necessari per Windows Search
try:
    import win32com.client
    import pythoncom
    WINDOWS_SEARCH_AVAILABLE = True
except ImportError:
    WINDOWS_SEARCH_AVAILABLE = False

# Definizione della costante per nascondere le finestre CMD in Windows
if platform.system() == "Windows":
    CREATE_NO_WINDOW = 0x08000000  # Per Python < 3.7

# Informazioni sulla versione dell'applicazione
APP_VERSION = "V9.2.8"
APP_STAGE = "Beta"
APP_NAME = "File Search Tool"
APP_FULL_NAME = f"{APP_NAME} {APP_VERSION} {APP_STAGE}"
APP_TITLE = f"{APP_NAME} {APP_VERSION} {APP_STAGE} Forensics"

# Elenco di librerie da installare se mancanti
missing_libraries = []

# Definizione del decoratore con debug migliorato
def error_handler(func):
    @functools.wraps(func)
    def wrapper(self, *args, **kwargs):
        try:
            return func(self, *args, **kwargs)
        except Exception as e:
            error_type = type(e).__name__
            error_traceback = traceback.format_exc()
            
            # Registra l'errore con il tipo specifico
            self.log_error(
                f"Errore nell'esecuzione di {func.__name__}: [{error_type}]", 
                exception=e, 
                location=func.__name__, 
                traceback=error_traceback
            )
            
            # Aggiorna immediatamente la visualizzazione del debug log se è aperta
            if hasattr(self, 'debug_window') and self.debug_window and hasattr(self.debug_window, 'winfo_exists') and self.debug_window.winfo_exists():
                try:
                    # Assumendo che update_log_display sia un metodo accessibile
                    if hasattr(self.debug_window, 'update_log_display'):
                        self.debug_window.update_log_display()
                except:
                    pass
            
            # Mostra una finestra di errore all'utente
            import tkinter.messagebox as messagebox
            messagebox.showerror(
                "Errore applicazione", 
                f"Si è verificato un errore durante {func.__name__}.\n"
                f"Tipo errore: {error_type}\n"
                f"Dettagli: {str(e)}\n\n"
                f"L'errore è stato registrato nel log di debug."
            )
            return None
    return wrapper

file_format_support = {
    "docx": False, "pdf": False, "pptx": False, "excel": False,
    "odt": False, "rtf": False, "xls": False, "doc": False,
    "ods": False, "odp": False, "epub": False, "mobi": False, 
    "tex": True, "rst": True, "sqlite": True, "mdb": True,
    "odb": True, "tsv": True, "dbf": False, "dif": True,
    "executable": True, "code_files": True, "accdb": True  
}

class PathUtils:
    """Classe di utilità per operazioni sui percorsi di file e cartelle.
    Contiene metodi relativi all'identificazione e gestione di percorsi di rete."""
    
    @staticmethod
    def is_network_path(path):
        """Determina se un percorso è un percorso di rete."""
        if not path:
            return False
        
        # Percorsi UNC Windows (\\server\share)
        if path.startswith('\\\\') or path.startswith('//'):
            return True
            
        # Percorsi di rete mappati su Windows (Z:\ dove Z è una lettera mappata)
        if os.name == 'nt' and len(path) >= 2 and path[1] == ':':
            try:
                # Usa net use per verificare se è un drive mappato
                drive_letter = path[0].upper() + ':'
                
                # Aggiungiamo il flag per nascondere la finestra CMD
                CREATE_NO_WINDOW = 0x08000000  # Per versioni di Python precedenti alla 3.7
                
                result = subprocess.run(
                    ['net', 'use', drive_letter], 
                    capture_output=True, 
                    text=True, 
                    timeout=2,
                    creationflags=CREATE_NO_WINDOW)
                
                return "Remote name" in result.stdout or "Nome remoto" in result.stdout
            except:
                pass
                
        # Riconoscimento di percorsi HTTP/FTP (addizionale rispetto al codice originale)
        if path.startswith(('http://', 'https://', 'ftp://')):
            return True
        
        return False
    
    def get_network_path_info(path):
        """Restituisce informazioni su un percorso di rete (estensione per funzionalità future).
        dict: Informazioni sul percorso di rete o None se non è un percorso di rete"""
        if not PathUtils.is_network_path(path):
            return None
            
        return {"is_network": True, "path": path}
    
class WindowsSearchHelper:
    """Classe per utilizzare Windows Search Service tramite COM per ricerche veloci
    con supporto migliorato per percorsi di rete e file di grandi dimensioni"""
    
    def __init__(self, logger=None):
        self.available = self._check_service_availability()
        self.logger = logger
        self.query_cache = {}  # Cache per i risultati delle query
        self.cache_timeout = 300  # Tempo di validità della cache in secondi
        self.connection_timeout = 10  # Timeout per le connessioni in secondi
        self.retry_attempts = 3  # Tentativi per le operazioni COM su rete
        self.is_network_optimized = True  # Abilita ottimizzazioni per rete
        
    def log(self, message, level="info"):
        if self.logger:
            if hasattr(self.logger, level):
                # Logger standard con metodi info, error, etc.
                getattr(self.logger, level)(message)
            elif hasattr(self.logger, 'log_debug') and level == "info":
                # FileSearchApp ha log_debug invece di info
                self.logger.log_debug(message)
            elif hasattr(self.logger, 'log_error') and level == "error":
                # FileSearchApp ha log_error invece di error
                self.logger.log_error(message)
            else:
                # Fallback: prova a usare print in caso di emergenza
                print(f"[{level.upper()}] {message}")
    
    def _check_service_availability(self):
        """Verifica se il servizio Windows Search è disponibile e attivo"""
        if not WINDOWS_SEARCH_AVAILABLE:
            return False
            
        try:
            # Inizializza COM per il thread corrente
            pythoncom.CoInitialize()
            # Prova a connettere con il servizio Windows Search
            connection = win32com.client.Dispatch("ADODB.Connection")
            connection.Open("Provider=Search.CollatorDSO;Extended Properties='Application=Windows';")
            connection.Close()
            return True
        except Exception:
            return False
        finally:
            pythoncom.CoUninitialize()
    
    def _is_network_path(self, path):
        """Determina se un percorso è su rete"""
        return PathUtils.is_network_path(path)
    
    def _get_connection(self, network_path=False):
        """Crea una connessione COM con gestione migliorata per rete"""
        retry_count = self.retry_attempts if network_path else 1
        
        for attempt in range(retry_count):
            try:
                connection = win32com.client.Dispatch("ADODB.Connection")
                connection.ConnectionTimeout = self.connection_timeout
                connection.Open("Provider=Search.CollatorDSO;Extended Properties='Application=Windows';")
                return connection
            except Exception as e:
                if attempt == retry_count - 1:
                    raise
                time.sleep(0.5)  # Attendi prima di riprovare
        
        return None  # Non dovrebbe mai arrivare qui
    
    def search_files(self, search_path, keywords, file_extensions=None, max_results=1000, use_cache=True):
        """Esegue la ricerca di file utilizzando Windows Search"""
        if not self.available:
            self.log("Servizio Windows Search non disponibile. Passando alla ricerca standard.", "warning")
            return []
        
        # Calcola un hash per questa query per la cache
        query_key = f"{search_path}:{','.join(sorted(keywords))}:{','.join(sorted(file_extensions or []))}"
        query_hash = hashlib.md5(query_key.encode()).hexdigest()
        
        # Controlla se abbiamo risultati in cache
        if use_cache and query_hash in self.query_cache:
            cache_entry = self.query_cache[query_hash]
            cache_time, cache_results = cache_entry
            
            # Verifica se la cache è ancora valida
            if (time.time() - cache_time) < self.cache_timeout:
                self.log(f"Risultati recuperati dalla cache per {search_path}", "info")
                return cache_results
        
        # Determina se stiamo cercando su un percorso di rete
        is_network = self._is_network_path(search_path) if self.is_network_optimized else False
        
        try:
            pythoncom.CoInitialize()
            
            # Prepara la stringa di connessione con gestione ottimizzata per rete
            connection = self._get_connection(network_path=is_network)
            
            # Prepara la query SQL
            recordset = win32com.client.Dispatch("ADODB.Recordset")
            
            # Costruisci la condizione per le parole chiave
            keyword_conditions = []
            for keyword in keywords:
                # Cerca sia nel nome file che nel contenuto
                keyword_conditions.append(f"CONTAINS(System.FileName, '\"{keyword}\"') OR CONTAINS(System.Search.Contents, '\"{keyword}\"')")
            
            keyword_query = " AND ".join(f"({condition})" for condition in keyword_conditions)
            
            # Aggiungi condizione per le estensioni file se specificate
            extension_condition = ""
            if file_extensions and len(file_extensions) > 0:
                ext_list = ", ".join(f"'{ext.replace('.', '')}'" for ext in file_extensions)
                extension_condition = f" AND System.FileExtension IN ({ext_list})"
            
            # Ottimizzazione della query per percorsi di rete
            if is_network:
                # Per i percorsi di rete, prima cerca solo nei nomi dei file per velocità
                scope_condition = f"SCOPE = '{search_path}'"
                query = f"SELECT System.ItemPathDisplay FROM SystemIndex WHERE {scope_condition} AND ({keyword_query}){extension_condition}"
                self.log(f"Query Windows Search ottimizzata per rete: {query}")
            else:
                # Query standard
                scope_condition = f"SCOPE = '{search_path}'"
                query = f"SELECT System.ItemPathDisplay FROM SystemIndex WHERE {scope_condition} AND ({keyword_query}){extension_condition}"
                self.log(f"Query Windows Search: {query}")
            
            # Imposta timeout per l'esecuzione della query (più lungo per rete)
            recordset.CursorLocation = 3  # adUseClient
            recordset.MaxRecords = max_results
            
            # Esegui la query con gestione migliorata degli errori
            try:
                recordset.Open(query, connection)
            except Exception as e:
                if is_network:
                    # Per errori su rete, prova una query più semplice
                    self.log(f"Errore nella query complessa su rete, tentativo con query semplificata: {str(e)}", "warning")
                    query = f"SELECT System.ItemPathDisplay FROM SystemIndex WHERE {scope_condition} AND CONTAINS(System.FileName, '*')"
                    try:
                        recordset.Open(query, connection)
                    except:
                        raise  # Se fallisce anche la query semplificata, propaga l'errore
                else:
                    raise
            
            # Raccogli i risultati
            results = []
            count = 0
            
            while not recordset.EOF and count < max_results:
                try:
                    file_path = recordset.Fields.Item("System.ItemPathDisplay").Value
                    if os.path.exists(file_path):  # Verifica che il file esista ancora
                        results.append(file_path)
                    count += 1
                    recordset.MoveNext()
                except Exception as e:
                    self.log(f"Errore durante l'accesso al record: {str(e)}", "warning")
                    recordset.MoveNext()
            
            # Chiudi le connessioni
            recordset.Close()
            connection.Close()
            
            self.log(f"Ricerca Windows Search completata. Trovati {len(results)} risultati.")
            
            # Salva i risultati in cache
            self.query_cache[query_hash] = (time.time(), results)
            
            # Pulisci la cache se è diventata troppo grande
            if len(self.query_cache) > 50:  # Limite di 50 query in cache
                oldest_keys = sorted(self.query_cache.keys(), 
                                   key=lambda k: self.query_cache[k][0])[:10]
                for key in oldest_keys:
                    del self.query_cache[key]
            
            return results
            
        except Exception as e:
            self.log(f"Errore durante la ricerca Windows Search: {str(e)}", "error")
            return []
        finally:
            pythoncom.CoUninitialize()
    
    def search_files_async(self, search_path, keywords, callback, file_extensions=None, max_results=1000):
        """Esegue la ricerca in modo asincrono per non bloccare l'interfaccia"""
        def search_thread():
            results = self.search_files(search_path, keywords, file_extensions, max_results)
            if callback:
                callback(results)
        
        thread = threading.Thread(target=search_thread)
        thread.daemon = True
        thread.start()
        return thread
    
    def index_status(self, path):
        """Verifica lo stato di indicizzazione di un percorso"""
        if not self.available:
            return False
        
        # Controlla se è un percorso di rete
        is_network = self._is_network_path(path) if self.is_network_optimized else False
        
        try:
            pythoncom.CoInitialize()
            
            # Ottieni lo stato di indicizzazione dal servizio
            connection = self._get_connection(network_path=is_network)
            
            recordset = win32com.client.Dispatch("ADODB.Recordset")
            query = f"SELECT System.Search.CatalogName FROM SystemIndex WHERE SCOPE = '{path}'"
            
            recordset.Open(query, connection)
            indexed = not recordset.EOF
            
            recordset.Close()
            connection.Close()
            
            return indexed
        except Exception as e:
            self.log(f"Errore durante la verifica dello stato di indicizzazione: {str(e)}", "warning")
            return False
        finally:
            pythoncom.CoUninitialize()
            
    def optimize_query_for_path(self, path, keywords, file_extensions=None):
        """Ottimizza la strategia di query in base al tipo di percorso"""
        is_network = self._is_network_path(path) if self.is_network_optimized else False
        
        # Costruisci la condizione per le parole chiave
        keyword_conditions = []
        for keyword in keywords:
            if is_network:
                # Su rete, inizialmente cerca solo nei nomi dei file per velocità
                keyword_conditions.append(f"CONTAINS(System.FileName, '\"{keyword}\"')")
            else:
                # Altrimenti cerca ovunque
                keyword_conditions.append(f"CONTAINS(System.FileName, '\"{keyword}\"') OR CONTAINS(System.Search.Contents, '\"{keyword}\"')")
        
        keyword_query = " AND ".join(f"({condition})" for condition in keyword_conditions)
        
        # Aggiungi condizione per le estensioni file se specificate
        extension_condition = ""
        if file_extensions and len(file_extensions) > 0:
            ext_list = ", ".join(f"'{ext.replace('.', '')}'" for ext in file_extensions)
            extension_condition = f" AND System.FileExtension IN ({ext_list})"
        
        # Costruisci la query completa
        scope_condition = f"SCOPE = '{path}'"
        query = f"SELECT System.ItemPathDisplay FROM SystemIndex WHERE {scope_condition} AND ({keyword_query}){extension_condition}"
        
        return query, is_network
    
    def clear_cache(self):
        """Pulisce la cache delle query"""
        self.query_cache.clear()
        self.log("Cache delle query Windows Search pulita")
    
    def set_network_optimization(self, enabled=True):
        """Abilita o disabilita le ottimizzazioni di rete"""
        self.is_network_optimized = enabled
        self.log(f"Ottimizzazione di rete {'abilitata' if enabled else 'disabilitata'}")
    
    def set_timeout(self, timeout_seconds):
        """Imposta il timeout per le connessioni"""
        self.connection_timeout = timeout_seconds
        self.log(f"Timeout di connessione impostato a {timeout_seconds} secondi")
    
    def set_retry_attempts(self, attempts):
        """Imposta il numero di tentativi per le operazioni COM su rete"""
        self.retry_attempts = attempts
        self.log(f"Tentativi per operazioni COM impostati a {attempts}")
    
    def set_cache_timeout(self, timeout_seconds):
        """Imposta il tempo di validità della cache"""
        self.cache_timeout = timeout_seconds
        self.log(f"Timeout della cache impostato a {timeout_seconds} secondi")

class NetworkSearchOptimizer:
    """Classe per ottimizzare la ricerca su percorsi di rete"""
    
    def __init__(self, logger=None):
        self.logger = logger
        self.network_cache = {}  # Cache dei risultati di rete
        self.network_connections = {}  # Stato delle connessioni di rete
        self.retry_count = 3  # Numero di tentativi per le operazioni di rete
        self.chunk_size = 8 * 1024 * 1024  # 8MB per trasferimento di rete
        self.timeout_multiplier = 2.5  # Moltiplicatore di timeout per percorsi di rete
    
    def log(self, message, level="info"):
        if self.logger:
            if level == "debug":
                self.logger.debug(message)
            elif level == "info":
                self.logger.info(message)
            elif level == "warning":
                self.logger.warning(message)
            elif level == "error":
                self.logger.error(message)
    
    def optimize_network_path(self, path):
        """Ottimizza un percorso di rete per le prestazioni"""
        if not self.is_network_path(path):
            return path
            
        # Normalizza percorso di rete
        normalized_path = self._normalize_network_path(path)
        
        # Verifica la connessione e pre-autentica se necessario
        self._ensure_network_connection(normalized_path)
        
        return normalized_path
    
    def _normalize_network_path(self, path):
        """Normalizza un percorso di rete per un accesso più coerente"""
        # Sostituisci più backslash con uno solo
        normalized = re.sub(r'\\{2,}', r'\\', path)
        
        # Assicurati che i percorsi UNC inizino con \\
        if normalized.startswith('\\') and not normalized.startswith('\\\\'):
            normalized = '\\' + normalized
        
        return normalized
    
    def _ensure_network_connection(self, path):
        """Assicura che la connessione di rete sia attiva e autenticata"""
        # Estrai server dal percorso di rete
        server_match = re.match(r'\\\\([^\\]+)', path)
        if not server_match:
            return False
            
        server = server_match.group(1)
        
        # Se abbiamo già verificato questa connessione di rete, ritorna il risultato memorizzato
        if server in self.network_connections:
            return self.network_connections[server]
            
        try:
            # Verifica che il server sia raggiungibile
            for i in range(self.retry_count):
                try:
                    subprocess.run(["ping", "-n", "1", "-w", "1000", server], 
                                  capture_output=True, check=True, timeout=2)
                    break
                except (subprocess.SubprocessError, subprocess.TimeoutExpired):
                    if i == self.retry_count - 1:
                        self.log(f"Impossibile connettersi al server di rete: {server}", "warning")
                        self.network_connections[server] = False
                        return False
            
            # Memorizza il risultato
            self.network_connections[server] = True
            return True
        except Exception as e:
            self.log(f"Errore durante la verifica della connessione di rete a {server}: {str(e)}", "error")
            self.network_connections[server] = False
            return False
    
    def is_network_path(self, path):
        """Determina se un percorso è un percorso di rete"""
        return PathUtils.is_network_path(path)
    
    def get_network_files(self, path, pattern="*", recursive=True, max_depth=3):
        """Ottiene un elenco di file da un percorso di rete con il pattern specificato"""
        cache_key = f"{path}:{pattern}:{recursive}:{max_depth}"
        
        # Controlla se i risultati sono nella cache
        if cache_key in self.network_cache:
            self.log(f"Risultati di rete recuperati dalla cache per {path}", "debug")
            return self.network_cache[cache_key]
        
        # Assicurati che il percorso di rete sia ottimizzato
        path = self.optimize_network_path(path)
        
        files = []
        try:
            # Usa robocopy per elencare i file di rete in modo efficiente
            # Robocopy ha una migliore gestione degli errori di rete rispetto a os.walk
            temp_output = tempfile.NamedTemporaryFile(delete=False, suffix='.txt')
            temp_output.close()
            
            depth_param = f"/LEV:{max_depth}" if recursive else "/LEV:1"
            cmd = [
                "robocopy", 
                path, 
                os.devnull, 
                pattern,
                "/L",      # Solo elenco, non copia
                "/NJH",    # No Job Header
                "/NJS",    # No Job Summary
                "/NC",     # No Class
                "/NS",     # No Size
                depth_param,
                f"/LOG:{temp_output.name}"
            ]
            
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
            
            # Leggi i risultati dal file temporaneo
            with open(temp_output.name, 'r', encoding='utf-8', errors='ignore') as f:
                for line in f:
                    line = line.strip()
                    if line and not line.startswith('New File') and not line.endswith('\\'):
                        # Estrai il percorso del file
                        file_path = os.path.join(path, line.strip())
                        files.append(file_path)
            
            # Pulisci il file temporaneo
            os.unlink(temp_output.name)
            
            # Memorizza i risultati nella cache
            self.network_cache[cache_key] = files
            
            return files
        except Exception as e:
            self.log(f"Errore durante l'enumerazione dei file di rete in {path}: {str(e)}", "error")
            return []
    
    def read_network_file_in_chunks(self, file_path, keywords, chunk_size=None):
        """Legge un file di rete in blocchi per ridurre l'utilizzo della memoria"""
        if chunk_size is None:
            chunk_size = self.chunk_size
            
        try:
            found_keywords = set()
            file_size = os.path.getsize(file_path)
            
            # Per file piccoli, leggi tutto in una volta
            if file_size < chunk_size:
                with open(file_path, 'rb') as f:
                    content = f.read()
                try:
                    text = content.decode('utf-8', errors='ignore')
                    for keyword in keywords:
                        if keyword.lower() in text.lower():
                            found_keywords.add(keyword)
                except:
                    pass
                return len(found_keywords) > 0, found_keywords
            
            # Per file grandi, leggi a blocchi
            with open(file_path, 'rb') as f:
                buffer = b""
                while True:
                    chunk = f.read(chunk_size)
                    if not chunk:
                        break
                    
                    # Aggiungi il chunk al buffer esistente
                    buffer += chunk
                    
                    # Decodifica e cerca le parole chiave
                    try:
                        text = buffer.decode('utf-8', errors='ignore')
                        for keyword in keywords:
                            if keyword.lower() in text.lower():
                                found_keywords.add(keyword)
                    except:
                        pass
                    
                    # Se abbiamo trovato tutte le parole chiave, possiamo fermarci
                    if len(found_keywords) == len(keywords):
                        return True, found_keywords
                    
                    # Mantieni l'ultima parte del buffer per gestire le parole chiave divise tra i chunk
                    buffer = buffer[-100:]  # Mantieni gli ultimi 100 byte
            
            return len(found_keywords) > 0, found_keywords
        except Exception as e:
            self.log(f"Errore durante la lettura del file di rete {file_path}: {str(e)}", "error")
            return False, set()
    
    def get_network_timeout(self, standard_timeout):
        """Calcola il timeout ottimale per le operazioni di rete"""
        return standard_timeout * self.timeout_multiplier
    
class LargeFileHandler:
    """Classe per gestire in modo efficiente i file di grandi dimensioni"""
    
    def __init__(self, logger=None):
        self.logger = logger
        self.large_file_threshold = 50 * 1024 * 1024  # 50 MB
        self.huge_file_threshold = 500 * 1024 * 1024  # 500 MB
        self.read_chunk_size = 4 * 1024 * 1024  # 4 MB
        self.max_preview_size = 10 * 1024  # 10 KB per l'anteprima
        self.supported_parsers = {}
        
        # Inizializza i parser specifici per tipo di file
        self._initialize_file_parsers()
    
    def log(self, message, level="info"):
        if self.logger:
            if level == "debug":
                self.logger.debug(message)
            elif level == "info":
                self.logger.info(message)
            elif level == "warning":
                self.logger.warning(message)
            elif level == "error":
                self.logger.error(message)
    
    def _initialize_file_parsers(self):
        """Inizializza i parser specializzati per i vari tipi di file"""
        # Aggiungiamo parser per tipi di file comuni che potrebbero essere grandi
        try:
            import xml.etree.ElementTree as ET
            self.supported_parsers['.xml'] = self._parse_xml
        except ImportError:
            self.log("Modulo xml.etree.ElementTree non disponibile", "warning")
        
        try:
            import csv
            self.supported_parsers['.csv'] = self._parse_csv
        except ImportError:
            self.log("Modulo csv non disponibile", "warning")
        
        try:
            import json
            self.supported_parsers['.json'] = self._parse_json
        except ImportError:
            self.log("Modulo json non disponibile", "warning")
    
    def get_file_size_category(self, file_path):
        """ Categorizza un file in base alla sua dimensione.Categoria del file ("normal", "medium", "large", "huge", o "gigantic")"""
        try:
            # Verifica che il file esista
            if not os.path.exists(file_path) or not os.path.isfile(file_path):
                return "normal"
                
            # Ottieni la dimensione del file una sola volta
            file_size = os.path.getsize(file_path)
            
            # Categorizza in base alle soglie
            if file_size >= self.gigantic_file_threshold:
                return "gigantic"
            elif file_size >= self.huge_file_threshold:
                return "huge"
            elif file_size >= self.large_file_threshold:
                return "large"
            elif file_size >= self.medium_file_threshold:
                return "medium"
            else:
                return "normal"
        except Exception as e:
            self.log(f"Errore nel determinare la categoria del file {file_path}: {str(e)}", "error")
            return "normal"  # In caso di errore, tratta come file normale

    def is_large_file(self, file_path):
        """Determina se un file è considerato 'grande'"""
        return self.get_file_size_category(file_path) in ["large", "huge", "gigantic"]
        
    def is_huge_file(self, file_path):
        """Determina se un file è considerato 'enorme'"""
        return self.get_file_size_category(file_path) in ["huge", "gigantic"]
    
    def search_in_large_file(self, file_path, keywords, is_whole_word=False):
        """Cerca keywords in un file di grandi dimensioni in modo ottimizzato"""
        extension = os.path.splitext(file_path)[1].lower()
        
        # Usa un parser specifico se disponibile per questo tipo di file
        if extension in self.supported_parsers:
            return self.supported_parsers[extension](file_path, keywords, is_whole_word)
        
        # Altrimenti usa la ricerca generica a blocchi
        return self._chunk_search(file_path, keywords, is_whole_word)
    
    def _chunk_search(self, file_path, keywords, is_whole_word=False):
        """Cerca keywords in un file leggendolo a blocchi"""
        keywords_lower = [k.lower() for k in keywords]
        found_keywords = set()
        overlap = max(len(max(keywords, key=len)) * 2, 200)  # Sovrapponi i blocchi per evitare di perdere keyword spezzate
        
        try:
            with open(file_path, 'rb') as f:
                file_size = os.path.getsize(file_path)
                
                # Per file molto grandi, aumenta la dimensione del chunk
                chunk_size = self.read_chunk_size
                if file_size > self.huge_file_threshold:
                    chunk_size = self.read_chunk_size * 2
                
                # Buffer per sovrapporre i blocchi
                last_data = b""
                position = 0
                
                while position < file_size:
                    # Leggi un nuovo blocco
                    f.seek(position)
                    new_data = f.read(chunk_size)
                    if not new_data:
                        break
                    
                    # Combina con i dati precedenti per gestire le keyword divise
                    data = last_data + new_data
                    
                    # Converti in testo e cerca
                    try:
                        text = data.decode('utf-8', errors='ignore')
                        text_lower = text.lower()
                        
                        for i, keyword in enumerate(keywords_lower):
                            if is_whole_word:
                                # Cerca parole intere usando espressioni regolari
                                import re
                                pattern = r'\b' + re.escape(keyword) + r'\b'
                                if re.search(pattern, text_lower, re.IGNORECASE):
                                    found_keywords.add(keywords[i])
                            else:
                                if keyword in text_lower:
                                    found_keywords.add(keywords[i])
                    except Exception as e:
                        self.log(f"Errore nella decodifica del testo in {file_path}: {str(e)}", "error")
                    
                    # Se abbiamo trovato tutte le keywords, fermiamoci
                    if len(found_keywords) == len(keywords):
                        return True, found_keywords
                    
                    # Salva gli ultimi overlap byte per la prossima iterazione
                    last_data = new_data[-overlap:] if len(new_data) > overlap else new_data
                    
                    # Avanza nella posizione
                    position += chunk_size - overlap
            
            return len(found_keywords) > 0, found_keywords
        
        except Exception as e:
            self.log(f"Errore durante la ricerca a blocchi in {file_path}: {str(e)}", "error")
            return False, set()
    
    def _parse_xml(self, file_path, keywords, is_whole_word=False):
        """Cerca in modo efficiente in file XML di grandi dimensioni"""
        import xml.sax
        import re
        
        class XMLHandler(xml.sax.ContentHandler):
            def __init__(self, keywords, is_whole_word):
                self.keywords = [k.lower() for k in keywords]
                self.original_keywords = keywords
                self.found_keywords = set()
                self.current_text = ""
                self.is_whole_word = is_whole_word
            
            def characters(self, content):
                self.current_text += content
            
            def endElement(self, name):
                text_lower = self.current_text.lower()
                
                for i, keyword in enumerate(self.keywords):
                    if self.is_whole_word:
                        pattern = r'\b' + re.escape(keyword) + r'\b'
                        if re.search(pattern, text_lower, re.IGNORECASE):
                            self.found_keywords.add(self.original_keywords[i])
                    else:
                        if keyword in text_lower:
                            self.found_keywords.add(self.original_keywords[i])
                
                self.current_text = ""
        
        try:
            handler = XMLHandler(keywords, is_whole_word)
            parser = xml.sax.make_parser()
            parser.setContentHandler(handler)
            parser.parse(file_path)
            return len(handler.found_keywords) > 0, handler.found_keywords
        except Exception as e:
            self.log(f"Errore durante il parsing XML di {file_path}: {str(e)}", "warning")
            # Fallback alla ricerca a blocchi
            return self._chunk_search(file_path, keywords, is_whole_word)
    
    def _parse_csv(self, file_path, keywords, is_whole_word=False):
        """Cerca in modo efficiente in file CSV di grandi dimensioni"""
        import csv
        import re
        
        found_keywords = set()
        keywords_lower = [k.lower() for k in keywords]
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore', newline='') as csvfile:
                reader = csv.reader(csvfile)
                
                # Leggi solo le prime 1000 righe per file molto grandi
                if self.is_huge_file(file_path):
                    rows_to_read = 1000
                else:
                    rows_to_read = float('inf')
                
                for i, row in enumerate(reader):
                    if i >= rows_to_read:
                        break
                    
                    for cell in row:
                        cell_lower = cell.lower()
                        for i, keyword in enumerate(keywords_lower):
                            if is_whole_word:
                                pattern = r'\b' + re.escape(keyword) + r'\b'
                                if re.search(pattern, cell_lower, re.IGNORECASE):
                                    found_keywords.add(keywords[i])
                            else:
                                if keyword in cell_lower:
                                    found_keywords.add(keywords[i])
                    
                    # Se abbiamo trovato tutte le keywords, fermiamoci
                    if len(found_keywords) == len(keywords):
                        return True, found_keywords
            
            return len(found_keywords) > 0, found_keywords
        except Exception as e:
            self.log(f"Errore durante il parsing CSV di {file_path}: {str(e)}", "warning")
            # Fallback alla ricerca a blocchi
            return self._chunk_search(file_path, keywords, is_whole_word)
    
    def _parse_json(self, file_path, keywords, is_whole_word=False):
        """Cerca in modo efficiente in file JSON di grandi dimensioni"""
        import json
        import re
        
        try:
            # Per file JSON enormi, usa una strategia di parsing a blocchi
            if self.is_huge_file(file_path):
                return self._chunk_search(file_path, keywords, is_whole_word)
            
            # Per file JSON di dimensioni gestibili, carica tutto
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)
            
            # Converti in testo per cercare le keywords
            text = json.dumps(data, ensure_ascii=False)
            text_lower = text.lower()
            
            found_keywords = set()
            for i, keyword in enumerate(keywords):
                keyword_lower = keyword.lower()
                if is_whole_word:
                    pattern = r'\b' + re.escape(keyword_lower) + r'\b'
                    if re.search(pattern, text_lower, re.IGNORECASE):
                        found_keywords.add(keyword)
                else:
                    if keyword_lower in text_lower:
                        found_keywords.add(keyword)
            
            return len(found_keywords) > 0, found_keywords
        except Exception as e:
            self.log(f"Errore durante il parsing JSON di {file_path}: {str(e)}", "warning")
            # Fallback alla ricerca a blocchi
            return self._chunk_search(file_path, keywords, is_whole_word)
    
    def get_file_preview(self, file_path):
        """Ottiene un'anteprima di un file di grandi dimensioni"""
        try:
            with open(file_path, 'rb') as f:
                data = f.read(self.max_preview_size)
                try:
                    preview = data.decode('utf-8', errors='ignore')
                    return preview + ("..." if os.path.getsize(file_path) > self.max_preview_size else "")
                except:
                    return f"[Anteprima non disponibile: contenuto binario]"
        except Exception as e:
            return f"[Errore nell'apertura del file: {str(e)}]"
        
class FileSearchApp:
    @error_handler
    def __init__(self, root):
            self.root = root
            self.root.title(APP_TITLE)
            
            # Reset monitoraggio memoria
            self.search_in_progress = False
            self.memory_monitor_id = None
            
            # Imposta subito il debug mode per poter loggare
            self.debug_mode = True
            self.debug_log = []
            # Aggiungi questa riga per inizializzare current_user
            self.current_user = getpass.getuser()
            
            # Imposta icona se disponibile
            try:
                icon_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "search_icon.ico")
                if os.path.exists(icon_path):
                    self.root.iconbitmap(icon_path)
            except:
                pass
        
            # Assicurati che all'avvio non ci siano timer attivi da esecuzioni precedenti
            for attr in dir(self):
                if attr.startswith('_after_id_') or attr.endswith('_timer_id'):
                    try:
                        after_id = getattr(self, attr)
                        if isinstance(after_id, int) and after_id > 0:
                            self.root.after_cancel(after_id)
                            self.debug_log.append(f"[DEBUG] Cancellato timer {attr}: {after_id}")
                    except Exception:
                        pass  # Ignora errori nel cleanup iniziale
            
            # Inizializza tutte le variabili in un passaggio
            self._init_essential_variables()
            self._init_remaining_variables()
            
            # Crea l'intera interfaccia in una volta sola
            self.create_widgets()
            
            # Esegui attività di background dopo un breve ritardo
            self.root.after(500, self._delayed_startup_tasks)
            # Controlla gli aggiornamenti all'avvio (dopo l'inizializzazione completa)
            self.root.after(5000, self.check_for_updates_on_startup)

    @error_handler
    def create_base_interface(self):
        """Crea solo l'interfaccia essenziale per un avvio veloce"""
        # Forza l'aggiornamento dell'interfaccia
        self.root.update_idletasks()

    @error_handler
    def complete_initialization(self):
        """Completa l'inizializzazione dell'applicazione"""
        try:
            # Rimuovi i componenti temporanei
            if hasattr(self, 'status_label'):
                self.status_label.destroy()
            if hasattr(self, 'progress_bar'):
                self.progress_bar.stop()
                self.progress_bar.destroy()
            
            # Inizializza le variabili rimanenti
            self._init_remaining_variables()
            
            # Crea l'interfaccia completa
            self.create_widgets()
            
            # Esegui le altre operazioni iniziali
            self._delayed_startup_tasks()
            
            # Avvia il monitoraggio della memoria
            self.root.after(30000, self.monitor_memory_usage)  # Primo controllo dopo 30 secondi
            
        except Exception as e:
            # In caso di errore, mostra un messaggio e prova a ripristinare l'applicazione
            messagebox.showerror("Errore di inizializzazione", 
                            f"Si è verificato un errore durante l'avvio: {str(e)}")
            # Tenta comunque di creare l'interfaccia
            try:
                self.create_widgets()
            except:
                pass

    @error_handler  
    def _init_essential_variables(self):
        """Inizializza solo le variabili essenziali per l'avvio"""
        # Inizializza datetime_var subito all'inizio per evitare errori di sequenza
        self.datetime_var = StringVar()
        self.max_depth = 5
        
        # Definisci max_workers prima di utilizzarlo
        self.max_workers = os.cpu_count() or 4  # Utilizza il numero di CPU disponibili o 4 come fallback

        import concurrent.futures
        self.executor = concurrent.futures.ThreadPoolExecutor(max_workers=self.max_workers)

        # Tracciamento file già processati (per evitare duplicati)
        self.processed_files = set()

        # Variabili principali per la ricerca
        self.search_content = BooleanVar(value=True)
        self.search_path = StringVar()
        self.keywords = StringVar()
        self.search_results = []
        self.search_files = BooleanVar(value=True)
        self.search_folders = BooleanVar(value=True)
        self.is_searching = False
        self.progress_queue = queue.Queue()
        self.search_depth = StringVar(value="base") 
        self.excluded_dirs = []
        
        # Variabili per sistema e configurazione
        self.app_name = APP_NAME
        self.app_version = APP_VERSION
        self.app_stage = APP_STAGE
        self.app_full_name = APP_FULL_NAME
        self.debug_mode = True
        self.is_search_running = False
        self.search_interrupted = False

        # Inizializza impostazioni per la gestione della RAM
        self.auto_memory_management = True
        self.memory_usage_percent = 75
        
        # Inizializza tutte le variabili utilizzate nelle impostazioni
        self.timeout_enabled = tk.BooleanVar(value=False)
        self.timeout_seconds = tk.IntVar(value=3600)
        self.max_files_to_check = tk.IntVar(value=100000)
        self.max_results = tk.IntVar(value=50000)
        self.worker_threads = tk.IntVar(value=4)
        self.max_file_size_mb = tk.IntVar(value=100)
        self.use_indexing = tk.BooleanVar(value=True)
        self.skip_permission_errors = tk.BooleanVar(value=True)
        
        # Variabili per data/ora e utente
        self.user_var = StringVar(value=getpass.getuser())
        
        # Variabili essenziali per l'interfaccia
        self.ignore_hidden = BooleanVar(value=True)
        self.search_executor = None
        self.exclude_system_files = BooleanVar(value=True)
        self.whole_word_search = BooleanVar(value=False)
        self.dir_size_calculation = StringVar(value="disabilitato")
        
        # Variabili per la visualizzazione
        self.directory_calculation_enabled = False
        self.dir_size_var = StringVar(value="")
        self.total_disk_var = StringVar(value="")
        self.used_disk_var = StringVar(value="")
        self.free_disk_var = StringVar(value="")
        
        # Variabili per il monitoraggio e tracking
        self.files_checked = 0
        self.files_matched = 0
        self.dirs_checked = 0
        self.skipped_files = 0
        self.search_start_time = None
        self.total_search_time = 0
        self.search_completed = False
        self.last_status_update = 0
        self.status_update_interval = 0.2  # secondi
        
        # Variabili per i file e le estensioni
        self.all_file_extensions = set()
        self.extension_categories = {}
        self.selected_extensions = {}
        self.skipped_extensions = set()
        self.excluded_extensions = set()
        self.excluded_folders = set()
        self.total_files_size = 0
        self.selected_files_size = 0
        self.path_disk_info = {}

        # Inizializza le classi di ottimizzazione
        self.network_optimizer = NetworkSearchOptimizer(logger=self)
        self.large_file_handler = LargeFileHandler(logger=None)
        self.windows_search_helper = WindowsSearchHelper(logger=self)

        # Configura le opzioni di rete
        self.network_retry_count = 3
        self.network_search_enabled = True
        self.network_parallel_searches = 4  # Numero di ricerche parallele su rete
        
        # Configura le opzioni per file di grandi dimensioni
        self.large_file_search_enabled = True
        self.medium_file_threshold = 10 * 1024 * 1024  # 10 MB
        self.large_file_threshold = 50 * 1024 * 1024   # 50 MB
        self.huge_file_threshold = 500 * 1024 * 1024   # 500 MB
        self.gigantic_file_threshold = 2 * 1024 * 1024 * 1024  # 2 GB

        # Configura anche la classe WindowsSearchHelper con le ottimizzazioni
        if hasattr(self, 'windows_search_helper'):
            self.windows_search_helper.set_network_optimization(self.network_search_enabled)
            self.windows_search_helper.set_retry_attempts(self.network_retry_count)
            self.windows_search_helper.set_timeout(10)
        
        # Variabili per logging e debug
        self.log_messages = []
        self.log_count = 0
        self.error_count = 0
        self.warning_count = 0
        self.info_count = 0
        self.log_filter_text = ""
        
        # Debug logs queue (rimuovo la duplicazione)
        self.debug_logs_queue = queue.Queue(maxsize=5000)  # Limit to 5000 entries to avoid memory issues
        
        # Aggiungi questa riga: lista permanente per i log completi dalla creazione dell'app
        self.complete_debug_log_history = []
            
        # Altre variabili che potrebbero essere necessarie
        self.search_whole_words = False
        self.search_case_sensitive = False
        self.search_timeout = 120  # secondi
        self.system_search_running = False
        self.search_was_timeout = False
        self.extension_filter_enabled = True
        self.skipped_files_log = []

    @error_handler
    def _init_remaining_variables(self):
        """Inizializza le variabili non essenziali per l'avvio"""
        # Variabili per la ricerca a blocchi
        self.max_files_per_block = IntVar(value=1000)
        self.max_parallel_blocks = IntVar(value=4)
        self.prioritize_user_folders = BooleanVar(value=True)
        self.block_size_auto_adjust = BooleanVar(value=True)
        
        # Configurazioni per la gestione della RAM
        self.auto_memory_management = True  # Attiva per default
        self.memory_usage_percent = 75      # Default 75% della RAM disponibile
        self.total_ram = psutil.virtual_memory().total

        # Variabili aggiuntive
        self.search_start_time = None
        self.stop_search = False
        self.max_depth = 0  # 0 = illimitato
        self.last_search_params = {}
        self.search_history = []
        
        # Filtri avanzati
        self.advanced_filters = {
            "size_min": 0,
            "size_max": 0,
            "date_min": None,
            "date_max": None,
            "extensions": []
        }
        
        # Opzioni per la gestione dei permessi
        self.skip_permission_errors = BooleanVar(value=True)
        self.excluded_paths = []
        self.load_settings_from_file()
        
        # Directory problematiche da escludere automaticamente
        self.problematic_dirs = [
            "Client Active Directory Rights Management Services",
            "Windows Resource Protection",
            "Windows Defender",
            "Microsoft Office",
            "System Volume Information",
            "$Recycle.Bin",
            "$WINDOWS.~BT",
            "$Windows.~WS"
        ]

        # Verifica dei privilegi di amministratore
        self.is_admin = False
        try:
            import ctypes
            self.is_admin = ctypes.windll.shell32.IsUserAnAdmin() != 0
        except:
            self.is_admin = False

        # Variabili per limiti di tempo e altre ottimizzazioni
        self.timeout_enabled = BooleanVar(value=False)
        self.timeout_seconds = IntVar(value=3600)
        self.max_files_to_check = IntVar(value=100000)
        self.max_results = IntVar(value=50000)
        self.chunk_size = 8192
        self.max_file_size_mb = IntVar(value=100)
        self.worker_threads = IntVar(value=min(8, os.cpu_count() or 4))
        self.use_indexing = BooleanVar(value=True)
        self.search_index = {}
        
        # Lista di estensioni di file di sistema da escludere dalla ricerca nei contenuti
        self.system_file_extensions = [
            # File eseguibili e librerie
            '.exe', '.dll', '.sys', '.drv', '.ocx', '.vxd', '.com', '.bat', '.cmd', '.scr', '.app', '.dylib', '.exp', '.bpl',
            # ... resto delle estensioni ... 
        ]

        # Percorso del log dei file saltati
        self.skipped_files_log_path = os.path.join(os.path.expanduser("~"), "skipped_files_log.txt")
        
        # Aggiorna l'orario
        self.update_datetime()
        pass

    @error_handler
    def _delayed_startup_tasks(self):
        """Esegue attività non essenziali all'avvio"""
        # Esegui queste operazioni in background
        threading.Thread(target=self._background_tasks, daemon=True).start()

    @error_handler
    def _background_tasks(self):
        """Esegue operazioni in background"""
        try:
            # Verifica librerie disponibili
            global file_format_support, missing_libraries
            
            # Controlla ogni libreria
            for module_name, format_key, import_name in [
                ("docx", "docx", "python-docx"),
                ("PyPDF2", "pdf", "PyPDF2"),
                # ... e così via per le altre librerie
            ]:
                try:
                    __import__(module_name)
                    file_format_support[format_key] = True
                except ImportError:
                    missing_libraries.append(import_name)
            
            # Mostra notifica dopo un ritardo
            if missing_libraries:
                self.root.after(2000, self.check_and_notify_missing_libraries)
            
            # Aggiorna informazioni disco senza calcolare dimensione directory
            if hasattr(self, 'search_path'):
                path = self.search_path.get()
                if path and os.path.exists(path):
                    self.update_disk_info(path, calculate_dir_size=False)
            
        except Exception as e:
            self.log_debug(f"Errore nelle operazioni in background: {str(e)}")

    @error_handler
    def _create_minimal_interface(self):
        """Crea solo i componenti essenziali dell'interfaccia"""
        # Frame principale che conterrà tutto
        self.main_container = ttk.Frame(self.root)
        self.main_container.pack(fill=BOTH, expand=YES)
                
        # Controlli minimi per la ricerca
        self.controls_frame = ttk.LabelFrame(self.main_container, text="Parametri di ricerca", padding=10)
        self.controls_frame.pack(fill=X, padx=10, pady=5)
        
        # Solo i controlli essenziali di ricerca
        self._create_essential_search_controls()
        
        # Barra di stato per informazioni
        status_frame = ttk.Frame(self.main_container, padding=5)
        status_frame.pack(fill=X, padx=10)
        
        self.status_label = ttk.Label(status_frame, text="Inizializzazione in corso...", font=("", 9))
        self.status_label.pack(side=LEFT)
        
        self.progress_bar = ttk.Progressbar(status_frame, mode='indeterminate')
        self.progress_bar.pack(side=LEFT, padx=10, fill=X, expand=YES)
        self.progress_bar.start(10)  # Avvia l'animazione di caricamento
        
        # Placeholder per i risultati
        self.results_container = ttk.LabelFrame(self.main_container, text="Risultati di ricerca", padding=10)
        self.results_container.pack(fill=BOTH, expand=YES, padx=10, pady=(5, 10))
        
        # Messaggio di caricamento
        self.loading_label = ttk.Label(self.results_container, text="Caricamento interfaccia in corso...", font=("", 12))
        self.loading_label.pack(expand=YES, pady=50)

    @error_handler
    def _create_essential_search_controls(self):
        """Crea solo i controlli essenziali per la ricerca"""
        # RIGA 1: Directory di ricerca
        path_frame = ttk.Frame(self.controls_frame)
        path_frame.pack(fill=X, pady=5)
        
        path_label = ttk.Label(path_frame, text="Directory:", width=12, anchor=W)
        path_label.pack(side=LEFT, padx=(0, 5))
        
        self.path_entry = ttk.Entry(path_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        self.browse_btn = ttk.Button(path_frame, text="Sfoglia", command=self.browse_directory, width=10)
        self.browse_btn.pack(side=LEFT)
        
        # RIGA 2: Parole chiave
        keyword_frame = ttk.Frame(self.controls_frame)
        keyword_frame.pack(fill=X, pady=5)
        
        keyword_label = ttk.Label(keyword_frame, text="Parole chiave:", width=12, anchor=W)
        keyword_label.pack(side=LEFT, padx=(0, 5))
        
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        # RIGA 3: Controllo della profondità di ricerca
        depth_frame = ttk.Frame(self.controls_frame)
        depth_frame.pack(fill=X, pady=5)
        
        depth_label = ttk.Label(depth_frame, text="Profondità max:", width=12, anchor=W)
        depth_label.pack(side=LEFT, padx=(0, 5))
        
        self.depth_spinbox = ttk.Spinbox(depth_frame, from_=0, to=20, width=5)
        self.depth_spinbox.pack(side=LEFT, padx=5)
        self.depth_spinbox.set("5")  # Valore predefinito
        
        depth_info_label = ttk.Label(depth_frame, text="(0 = illimitato, consigliato 5-10)")
        depth_info_label.pack(side=LEFT, padx=5)
        
        # Aggiunge un callback per aggiornare max_depth quando cambia il valore dello spinbox
        self.depth_spinbox.bind("<FocusOut>", lambda e: setattr(self, "max_depth", 
                            int(self.depth_spinbox.get()) if self.depth_spinbox.get().isdigit() else 5))
        
        # Solo il pulsante di ricerca
        button_frame = ttk.Frame(self.controls_frame)
        button_frame.pack(fill=X, pady=(10, 5))
        
        center_frame = ttk.Frame(button_frame)
        center_frame.pack(side=TOP)
        
        self.search_button = ttk.Button(center_frame, text="CERCA", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=15)
        self.search_button.pack(side=LEFT, padx=10)
        
        # Disabilita il pulsante fino al caricamento completo
        self.search_button["state"] = "disabled"

    @error_handler
    def _create_remaining_widgets(self):
        """Crea il resto dell'interfaccia in background"""
        try:
            # Aggiorna il messaggio di stato
            self.status_label["text"] = "Caricamento componenti avanzati..."
            
            # Rimuovi gli oggetti temporanei
            if hasattr(self, 'loading_label'):
                self.loading_label.destroy()
            
            # Ferma l'animazione della barra di progresso
            self.progress_bar.stop()
            self.progress_bar["mode"] = "determinate"
            
            # Crea tutti i componenti dell'interfaccia completa
            self.create_widgets()
            
            # Aggiorna lo stato
            self.status_label["text"] = "Applicazione pronta"
            
            # Riabilita il pulsante di ricerca
            self.search_button["state"] = "normal"
            
            # Imposta il focus sul campo di ricerca
            self.path_entry.focus_set()
            
            self.log_debug("Interfaccia utente completamente caricata")
            
        except Exception as e:
            self.log_debug(f"Errore nel caricamento dell'interfaccia: {str(e)}")
            self.status_label["text"] = "Errore nel caricamento dell'interfaccia"

    @error_handler
    def _check_available_libraries(self):
        """Verifica la disponibilità delle librerie in background senza bloccare l'avvio"""
        # Esegui il controllo in un thread separato
        threading.Thread(target=self._async_check_libraries, daemon=True).start()

        # Aggiungi verifica per le librerie di gestione degli archivi compressi
        self._async_check_libraries()
        
        # Verifica librerie per archivi compressi
        try:
            import zipfile
        except ImportError:
            missing_libraries.append("zipfile")
        
        try:
            import rarfile
        except ImportError:
            missing_libraries.append("rarfile")
        
        try:
            import shutil
            import os
            
            seven_zip_found = False
            seven_zip_paths = [
                r"C:\Program Files\7-Zip\7z.exe",
                r"C:\Program Files (x86)\7-Zip\7z.exe",
                "7z"  # Se è nel PATH di sistema
            ]
            
            for path in seven_zip_paths:
                if os.path.isfile(path) or shutil.which(path):
                    seven_zip_found = True
                    break
                    
            if not seven_zip_found:
                missing_libraries.append("7-Zip")
        except Exception:
            missing_libraries.append("7-Zip")
    
        try:
            import tarfile
        except ImportError:
            missing_libraries.append("tarfile")
            
    @error_handler
    def _async_check_libraries(self):
        """Controlla le librerie in background senza bloccare l'interfaccia"""
        global file_format_support, missing_libraries
        
        # Funzione per verificare un singolo modulo
        def check_module(module_name, format_key, import_name=None):
            try:
                # Usa importlib per evitare bloccaggi
                import importlib
                importlib.import_module(module_name)
                file_format_support[format_key] = True
                self.log_debug(f"Supporto {format_key} attivato")
            except ImportError:
                missing_libraries.append(import_name or module_name)
                self.log_debug(f"Supporto {format_key} non disponibile")
        
        # Controlla ogni libreria
        check_module("docx", "docx", "python-docx")
        check_module("PyPDF2", "pdf", "PyPDF2")
        check_module("pptx", "pptx", "python-pptx")
        check_module("openpyxl", "excel", "openpyxl")
        check_module("striprtf.striprtf", "rtf", "striprtf")
        check_module("odf", "odt", "odfpy")
        check_module("ebooklib", "epub", "ebooklib")
        check_module("mobi", "mobi", "mobi")
        check_module("dbfread", "dbf", "dbfread")
        check_module("bs4", "epub_html", "beautifulsoup4")
        check_module("pefile", "pe_files", "pefile")
        check_module('pyodbc', 'mdb', 'accdb', 'pyodbc')
        check_module("pywin32", "windows_search", "win32com.client")

        # Controlla win32com (per i formati legacy di Office)
        if os.name == 'nt':  # Solo su Windows
            try:
                import win32com.client
                file_format_support["doc"] = True
                file_format_support["xls"] = True
                self.log_debug("Supporto formati legacy Office attivato (win32com)")
            except ImportError:
                missing_libraries.append("pywin32")
                self.log_debug("Supporto formati legacy Office non disponibile (manca pywin32)")
        
        # Controlla xlrd come alternativa per i file XLS
        try:
            import xlrd
            file_format_support["xls_native"] = True
            if not file_format_support.get("xls", False):
                file_format_support["xls"] = True
            self.log_debug("Supporto XLS attivato (xlrd)")
        except ImportError:
            if not file_format_support.get("xls", False):
                missing_libraries.append("xlrd")
                self.log_debug("Supporto XLS via xlrd non disponibile")
        
        # Controlla supporto ODT con librerie alternative
        try:
            from odf import opendocument
            file_format_support["odt"] = True
            self.log_debug("Supporto ODT attivato (libreria odf)")
        except ImportError:
            missing_libraries.append("odf")
            self.log_debug("Supporto ODT non disponibile")
        
        # Mostra notifica dopo un ritardo
        if missing_libraries:
            self.root.after(2000, self.check_and_notify_missing_libraries)

    @error_handler
    def process_file(self, file_path, keywords, search_content=True):
        """Processa un singolo file per verificare corrispondenze"""
        if self.stop_search:
            return None
        
        # NUOVA LOGICA: Verifica se il file è già stato processato
        # Normalizza il percorso per essere sicuri
        normalized_path = os.path.normpath(os.path.abspath(file_path))
        if normalized_path in self.processed_files:
            self.log_debug(f"File già processato, saltato: {normalized_path}")
            return None
        
        try:
            # Verifica nome file
            file_name = os.path.basename(file_path)
            matched = False
            
            # Verifica corrispondenze nel nome
            for keyword in keywords:
                if self.whole_word_search.get():
                    if self.is_whole_word_match(keyword, file_name):
                        matched = True
                        break
                elif keyword.lower() in file_name.lower():
                    matched = True
                    break
            
            content = ""
            # Verifica contenuto se richiesto e se non c'è già una corrispondenza nel nome
            if not matched and search_content and self.should_search_content(file_path):
                # NUOVA LOGICA: Controlla se il file è marcato per analisi parziale
                is_partial_analysis = hasattr(self, '_partial_analysis_files') and file_path in self._partial_analysis_files
                
                if is_partial_analysis:
                    # Usa l'analisi parziale per file giganteschi
                    self.log_debug(f"Applicando analisi parziale per file gigantesco: {os.path.basename(file_path)}")
                    matched = self._partial_content_search(file_path, keywords)
                else:
                    # Continua con l'analisi normale
                    content = self.get_file_content(file_path)
                    
                    # NUOVA GESTIONE: verifica se content è un dizionario (archivi compressi)
                    if isinstance(content, dict):
                        for file_in_archive, file_content in content.items():
                            # Controlla match nel nome del file interno
                            for keyword in keywords:
                                if self.whole_word_search.get():
                                    if self.is_whole_word_match(keyword, file_in_archive):
                                        self.log_debug(f"Match trovato nel nome del file interno: {file_in_archive}")
                                        matched = True
                                        break
                                elif keyword.lower() in file_in_archive.lower():
                                    self.log_debug(f"Match trovato nel nome del file interno: {file_in_archive}")
                                    matched = True
                                    break
                            
                            # Se non ha già trovato match nel nome, controlla nel contenuto
                            if not matched and isinstance(file_content, str):
                                for keyword in keywords:
                                    if self.whole_word_search.get():
                                        if self.is_whole_word_match(keyword, file_content):
                                            self.log_debug(f"Match trovato nel contenuto del file interno: {file_in_archive}")
                                            matched = True
                                            break
                                    elif keyword.lower() in file_content.lower():
                                        self.log_debug(f"Match trovato nel contenuto del file interno: {file_in_archive}")
                                        matched = True
                                        break
                            
                            # Interrompe il ciclo se ha già trovato una corrispondenza
                            if matched:
                                break
                    # Contenuto standard (stringa)
                    elif isinstance(content, str):
                        for keyword in keywords:
                            if self.whole_word_search.get():
                                if self.is_whole_word_match(keyword, content):
                                    matched = True
                                    break
                            elif keyword.lower() in content.lower():
                                matched = True
                                break
            
            if matched:
                # NUOVA LOGICA: Aggiungi il file all'elenco dei processati
                self.processed_files.add(normalized_path)
                
                # Verifica se il match è in un allegato di un file EMAIL (EML o MSG)
                _, ext = os.path.splitext(file_path)
                if ext.lower() in ['.eml', '.msg'] and search_content and isinstance(content, str) and "--- ALLEGATO" in content:
                    # Cerca il match dopo un'intestazione di allegato
                    attachment_sections = content.split("--- ALLEGATO")
                    for section in attachment_sections[1:]:  # Salta il primo che è l'intestazione email
                        for keyword in keywords:
                            if (self.whole_word_search.get() and self.is_whole_word_match(keyword, section)) or \
                            (not self.whole_word_search.get() and keyword.lower() in section.lower()):
                                # Match trovato in un allegato
                                self.log_debug(f"Match trovato in allegato di {file_path}")
                                return self.create_file_info(file_path, from_attachment=True)
                
                # Match normale (non in allegato o non in file EMAIL)
                return self.create_file_info(file_path)
                    
            # NUOVA LOGICA: Se non è stato trovato match, aggiungi comunque il file all'elenco dei processati
            self.processed_files.add(normalized_path)
                    
        except Exception as e:
            self.log_error(f"Errore nel processare il file {file_path}", exception=e)
            if self.debug_mode:
                import traceback
                self.log_error(traceback.format_exc())
                
        return None
    
    @error_handler
    def process_file_with_timeout(self, file_path, keywords, search_content=True):
        """Process a file with timeout to prevent hanging"""
        # CORREZIONE: Rilevamento tipo file per ottimizzare timeout
        is_large_file = False
        is_binary_file = False
        is_network_file = False
        is_email_file = False  # Nuovo flag per file email
        
        try:
            # Identifica se il file è su un percorso di rete
            if file_path.startswith('\\\\') or file_path.startswith('//'):
                is_network_file = True
                
            # Verifica la dimensione e imposta un flag per i file grandi
            if os.path.exists(file_path):
                file_size = os.path.getsize(file_path)
                # File sopra 5MB sono considerati grandi
                if file_size > 5 * 1024 * 1024:
                    is_large_file = True
                    
                # Rileva se è probabilmente un file binario
                ext = os.path.splitext(file_path)[1].lower()
                binary_extensions = ['.exe', '.dll', '.bin', '.obj', '.o', '.so', '.lib', '.sys', '.ocx']
                if ext in binary_extensions:
                    is_binary_file = True
                    
                # Identifica file email che possono avere allegati
                if ext.lower() in ['.msg', '.eml']:
                    is_email_file = True
        except:
            pass
        
        # CORREZIONE: Adatta il timeout in base al tipo di file
        timeout = 5.0  # Default 5 secondi
        if is_network_file:
            timeout = 10.0  # 10 secondi per file di rete
        if is_large_file:
            timeout += 5.0  # +5 secondi per file grandi
        if is_binary_file:
            timeout = min(timeout, 3.0)  # Limita a 3 secondi per file binari
        if is_email_file:
            timeout += 15.0  # Incremento significativo per file email con allegati
        
        # Manteniamo liste invece di code per compatibilità
        result = [None]
        exception = [None]
        processing_completed = [False]
        
        # CORREZIONE: Utilizzo di una funzione wrapper per gestire correttamente gli argomenti
        def process_wrapper():
            try:
                # Chiama process_file con i parametri corretti
                res = self.process_file(file_path, keywords, search_content)
                # Salva il risultato
                result[0] = res
            except Exception as e:
                # Cattura eventuali eccezioni
                exception[0] = e
            finally:
                # Segna il completamento
                processing_completed[0] = True
        
        # Usa la funzione wrapper come target
        thread = threading.Thread(
            target=process_wrapper
        )
        thread.daemon = True
        thread.start()
        thread.join(timeout)
        
        # Se il thread è ancora in esecuzione ma abbiamo superato il timeout
        if not processing_completed[0] and thread.is_alive():
            self.log_debug(f"Processing timed out for file: {file_path}")
            
            # PARTE NUOVA: Gestione dei risultati tardivi per file email
            if is_email_file:
                def check_late_results():
                    if processing_completed[0] or not thread.is_alive():
                        # Il thread ha terminato, controlla se ci sono nuovi risultati
                        if result[0] is not None:
                            self.log_debug(f"Risultati tardivi trovati in: {file_path} - Aggiornamento UI")
                            # Aggiungi i risultati tardivi alla UI
                            self.add_search_result(result[0])  # Sostituisci con il tuo metodo
                            self.update_results_list()  # Aggiorna la UI
                    else:
                        # Il thread sta ancora lavorando, ricontrolla più tardi
                        self.root.after(1000, check_late_results)
                
                # Inizia a controllare i risultati tardivi
                self.root.after(2000, check_late_results)
        
        if exception[0]:
            raise exception[0]
        
        return result[0] if result[0] is not None else []

    @error_handler
    def manage_memory(self):
        """Gestisce la memoria dell'applicazione in base alle impostazioni configurate dall'utente.
        Rilascia memoria quando necessario per mantenere le prestazioni ottimali."""
        try:
            # Verifica se la gestione automatica è disattivata dall'utente
            if not getattr(self, 'auto_memory_management', True):
                # Modalità manuale: rispetta solo il limite di memoria configurato
                try:
                    # Ottieni informazioni sulla memoria
                    total_ram = psutil.virtual_memory().total
                    process = psutil.Process(os.getpid())
                    current_memory = process.memory_info().rss
                    
                    # Calcola il limite basato sulla percentuale configurata
                    percent = getattr(self, 'memory_usage_percent', 75)
                    max_memory = total_ram * (percent / 100)
                    
                    # Verifica se superiamo la soglia
                    if current_memory > max_memory:
                        self.log_debug(f"Limite memoria ({percent}%) superato: {current_memory/(1024**2):.2f} MB / " +
                                    f"{max_memory/(1024**2):.2f} MB - Eseguendo pulizia manuale")
                        
                        # Pulizia in due fasi
                        gc.collect(0)  # Pulizia leggera iniziale
                        time.sleep(0.05)
                        gc.collect(2)  # Pulizia completa
                        
                        # Supporto per Python 3.9+
                        if hasattr(gc, 'collect_step'):
                            gc.collect_step()
                        
                        # Riduzione della cache interna se necessario
                        memory_after_gc = process.memory_info().rss
                        if memory_after_gc > max_memory * 0.95 and hasattr(self, 'search_results') and len(self.search_results) > 5000:
                            self.log_debug(f"Memoria ancora alta dopo GC, riduzione cache risultati da {len(self.search_results)} a 5000 elementi")
                            self.search_results = self.search_results[:5000]  # Mantieni i risultati più recenti
                            gc.collect()  # Richiama GC dopo la riduzione cache
                except Exception as e:
                    self.log_debug(f"Errore nella gestione manuale della memoria: {str(e)}")
            else:
                # Modalità automatica: gestione proattiva della memoria
                try:
                    # Ottimizza la frequenza del garbage collector in base al carico
                    is_searching = getattr(self, 'is_search_running', False)
                    if is_searching:
                        # Durante la ricerca, raccolta meno frequente per non interrompere le operazioni
                        gc.set_threshold(700, 10, 10)
                    else:
                        # Quando inattivo, raccolta più aggressiva
                        gc.set_threshold(100, 5, 5)
                    
                    # Ottieni informazioni sulla memoria
                    process = psutil.Process(os.getpid())
                    memory_usage = process.memory_info().rss / (1024 * 1024)  # MB
                    
                    # Determina la soglia dinamica di memoria
                    system_memory = psutil.virtual_memory().total / (1024 * 1024)  # Total system RAM in MB
                    threshold = max(200, min(1000, system_memory * 0.3))  # 30% della RAM di sistema con limiti min/max
                    
                    # Log utilizzo memoria (solo ogni 10 chiamate per non sovraccaricare i log)
                    self._memory_log_counter = getattr(self, '_memory_log_counter', 0) + 1
                    if self._memory_log_counter % 20 == 0:
                        self.log_debug(f"Gestione memoria automatica: Utilizzo corrente {memory_usage:.2f} MB, soglia {threshold:.2f} MB")
                    
                    # Verifica se superiamo la soglia
                    if memory_usage > threshold:
                        self.log_debug(f"Utilizzo memoria elevato: {memory_usage:.2f} MB. Esecuzione garbage collection.")
                        
                        # Pulizia progressiva
                        gc.collect(0)  # Garbage collection leggera
                        memory_after_gc0 = process.memory_info().rss / (1024 * 1024)
                        
                        if memory_after_gc0 > threshold * 0.9:
                            self.log_debug("Prima pulizia insufficiente, esecuzione pulizia completa.")
                            gc.collect(2)  # Garbage collection completa
                            
                            # Riduzione della cache se la memoria è ancora alta
                            memory_final = process.memory_info().rss / (1024 * 1024)
                            if memory_final > threshold and hasattr(self, 'search_results') and len(self.search_results) > 10000:
                                self.log_debug(f"Riduzione cache risultati da {len(self.search_results)} a 10000.")
                                self.search_results = self.search_results[-10000:]
                                gc.collect()  # Altra pulizia dopo la riduzione della cache
                except Exception as e:
                    self.log_debug(f"Errore nella gestione automatica della memoria: {str(e)}")
        except ImportError:
            # psutil non disponibile
            self.log_debug("psutil non disponibile, usando solo garbage collection standard.")
            gc.collect()
        except Exception as e:
            self.log_debug(f"Errore generale nella gestione della memoria: {str(e)}")
            # Fallback alla garbage collection standard
            try:
                gc.collect()
            except Exception as e2:
                self.log_debug(f"Anche il fallback di garbage collection è fallito: {str(e2)}")
        
        # Pianifica il prossimo controllo della memoria
        if hasattr(self, 'root') and self.root:
            self.root.after(30000, self.manage_memory)  # Controlla ogni 30 secondi
    
    @error_handler
    def monitor_memory_usage(self):
        """Monitoraggio periodico dell'utilizzo della memoria. Rileva situazioni critiche e interviene preventivamente."""
        try:
            # Verifica se il monitoraggio dovrebbe essere attivo
            if not hasattr(self, 'search_in_progress') or not self.search_in_progress:
                self.log_debug("Monitoraggio memoria terminato: ricerca non più attiva")
                self.memory_monitor_id = None
                return
                
            # Controlla anche il flag is_searching come doppia verifica
            if not hasattr(self, 'is_searching') or not self.is_searching:
                self.log_debug("Monitoraggio memoria terminato: flag is_searching è False")
                self.search_in_progress = False
                self.memory_monitor_id = None
                return
                
            # Esegui solo ogni 30 secondi per ridurre l'overhead
            if not hasattr(self, '_last_memory_check'):
                self._last_memory_check = time.time()
                return
                
            current_time = time.time()
            if current_time - self._last_memory_check < 30:
                return
                
            self._last_memory_check = current_time
            
            # Ottieni l'utilizzo di memoria del processo corrente
            process = psutil.Process()
            memory_info = process.memory_info()
            memory_used_mb = memory_info.rss / (1024**2)
            
            # Ottieni la memoria totale del sistema
            total_ram = psutil.virtual_memory().total / (1024**2)
            
            # Calcola la percentuale utilizzata rispetto al limite configurato
            memory_percent = getattr(self, 'memory_usage_percent', 75)
            memory_limit = total_ram * (memory_percent / 100)
            usage_ratio = memory_used_mb / memory_limit
            
            # Log standard dell'utilizzo della memoria (sempre mostrato)
            self.log_debug(f"Gestione memoria automatica: Utilizzo corrente {memory_used_mb:.2f} MB, soglia {memory_limit:.2f} MB")
            
            # Se superiamo l'80% del limite configurato, avvisa e forza una pulizia
            if usage_ratio > 0.8:
                self.log_debug(f"Memoria in uso al {usage_ratio*100:.1f}% del limite configurato ({memory_used_mb:.1f}MB/{memory_limit:.1f}MB)")
                
                # Azioni preventive progressive
                if usage_ratio > 0.9:
                    self.log_debug("ATTENZIONE: Utilizzo memoria CRITICO - Esecuzione pulizia di emergenza")
                    # Pulizia aggressiva
                    gc.collect(2)
                    time.sleep(0.1)
                    gc.collect(2)
                    
                    # Riduzione cache se necessario
                    if hasattr(self, 'search_results') and len(self.search_results) > 5000:
                        original_len = len(self.search_results)
                        self.search_results = self.search_results[:5000]
                        self.log_debug(f"Riduzione di emergenza della cache risultati: {original_len} → 5000 elementi")
                else:
                    # Pulizia standard
                    gc.collect()
                    
        except Exception as e:
            self.log_debug(f"Errore nel monitoraggio della memoria: {str(e)}")
            
        finally:
            # Riprogramma il prossimo controllo solo se il monitoraggio è ancora attivo
            if (hasattr(self, 'search_in_progress') and self.search_in_progress and 
                hasattr(self, 'is_searching') and self.is_searching and 
                hasattr(self, 'root') and self.root):
                self.memory_monitor_id = self.root.after(30000, self.monitor_memory_usage)
            else:
                self.log_debug("Monitoraggio memoria non riprogrammato: ricerca non più attiva")
                # Assicuriamoci che i flag siano coerenti
                self.search_in_progress = False
                # Azzera l'ID del timer
                self.memory_monitor_id = None

    @error_handler
    def start_memory_monitoring(self):
        """Avvia il monitoraggio della memoria con sicurezza"""
        # Prima ferma qualsiasi monitoraggio esistente per evitare duplicati
        self.stop_memory_monitoring()
        
        # Inizializza lo stato e il timestamp
        self.search_in_progress = True
        self._last_memory_check = time.time()
        
        # Log diagnosticov
        self.log_debug("AVVIO: Monitoraggio memoria automatica iniziato")
        
        # Avvia il ciclo di monitoraggio
        self.memory_monitor_id = self.root.after(5000, self.monitor_memory_usage)
        
        # Verifica che l'ID sia salvato correttamente
        self.log_debug(f"ID Monitoraggio Memoria: {self.memory_monitor_id}")

    @error_handler
    def stop_memory_monitoring(self):
        """Ferma il monitoraggio della memoria senza interferire con la ricerca"""
        self.log_debug("RICHIESTA STOP: Fermando monitoraggio memoria")
        
        # Imposta SOLO il flag di monitoraggio memoria a False
        self.search_in_progress = False
        
        # Cancella timer in modo più aggressivo
        if hasattr(self, 'memory_monitor_id') and self.memory_monitor_id:
            try:
                # Tentativo principale
                self.root.after_cancel(self.memory_monitor_id)
                self.log_debug(f"Timer memoria cancellato: {self.memory_monitor_id}")
            except Exception as e:
                self.log_debug(f"Errore cancellazione timer: {str(e)}")
            
            # Sempre resetta l'ID
            self.memory_monitor_id = None
        else:
            self.log_debug("Nessun monitoraggio memoria attivo da fermare")
        
        # FORZA una pulizia di memoria
        gc.collect()
        self.log_debug("STOP COMPLETO: Monitoraggio memoria fermato")

    @error_handler
    def run_with_timeout(self, func, args=(), kwargs={}, timeout_sec=10):
        """ Esegue una funzione con un timeout, evitando blocchi indefiniti Returns: (result, completed)
            - result: risultato della funzione o None
            - completed: True se completata, False se interrotta per timeout
        """
        result = [None]
        completed = [False]
        
        def target():
            try:
                result[0] = func(*args, **kwargs)
            except Exception as e:
                self.log_debug(f"Eccezione nella funzione con timeout: {str(e)}")
            finally:
                completed[0] = True
        
        thread = threading.Thread(target=target)
        thread.daemon = True
        thread.start()
        thread.join(timeout_sec)
        
        return result[0], completed[0]
    
    @error_handler
    def _get_all_descendants(self, widget):
        """Ottiene ricorsivamente tutti i widget discendenti"""
        descendants = []
        for child in widget.winfo_children():
            descendants.append(child)
            descendants.extend(self._get_all_descendants(child))
        return descendants
    
    @error_handler
    def set_controls_state(self, enabled=True, control_type=None, widget=None, exclude_widgets=None):
        """Imposta lo stato abilitato/disabilitato per i controlli dell'interfaccia utente."""
        # Stabilisci lo stato da applicare
        state = "normal" if enabled else "disabled"
        
        # Se exclude_widgets non è specificato, inizializzalo come lista vuota
        if exclude_widgets is None:
            exclude_widgets = []
        
        # Usa il widget root se non è specificato un widget di partenza
        if widget is None:
            widget = self.root
        
        # Ottieni tutti i widget discendenti, incluso il widget di partenza
        descendants = self._get_all_descendants(widget)
        
        # Traccia il tipo di controllo gestito
        control_type_lower = control_type.lower() if control_type else None
        
        for w in descendants:
            # Salta questo widget se è nella lista dei widget da escludere
            if w in exclude_widgets:
                continue
                
            widget_type = w.winfo_class().lower()
            
            # Salta il widget se è richiesto un tipo specifico e questo non corrisponde
            if control_type_lower and not (
                (control_type_lower == 'checkbutton' and widget_type == 'ttk::checkbutton') or
                (control_type_lower == 'button' and widget_type in ('ttk::button', 'button')) or
                (control_type_lower == 'entry' and widget_type in ('ttk::entry', 'entry')) or
                (control_type_lower == 'combobox' and widget_type == 'ttk::combobox') or
                (control_type_lower == 'spinbox' and widget_type in ('ttk::spinbox', 'spinbox'))
            ):
                continue
            
            try:
                # Gestisci i widget ttk in modo specifico 
                if widget_type.startswith('ttk::'):
                    w.configure(state=state)
                # Gestisci i widget Text separatamente
                elif widget_type == 'text':
                    if state == "normal":
                        w.configure(state="normal")
                    else:
                        # Salva il contenuto corrente
                        content = w.get("1.0", "end-1c")
                        w.configure(state="normal")
                        w.delete("1.0", "end")
                        w.insert("1.0", content)
                        w.configure(state="disabled")
                # Gestisci i widget standard Tkinter
                elif hasattr(w, 'config') and not widget_type in ['frame', 'toplevel', 'canvas']:
                    w.config(state=state)
                
                # Gestione speciale per le Listbox
                if widget_type == 'listbox' or widget_type == 'ttk::treeview':
                    if not enabled:
                        w.bind('<Button-1>', lambda e: 'break')
                    else:
                        w.unbind('<Button-1>')
                        
                # Log per debug se necessario
                if self.debug_mode:
                    self.log_debug(f"{'Abilitato' if enabled else 'Disabilitato'} controllo: {widget_type}")
                    
            except Exception as e:
                self.log_debug(f"Errore nel configurare lo stato del widget {widget_type}: {str(e)}")
                
    @error_handler
    def disable_all_controls(self):
        """Disabilita tutti i controlli UI durante la ricerca eccetto il debug log button"""
        # Assumendo che self.debug_log_button sia il riferimento al pulsante debug log
        self.set_controls_state(enabled=False, exclude_widgets=[self.debug_button])       
        
    @error_handler
    def enable_all_controls(self):
        """Abilita tutti i controlli UI dopo la ricerca"""
        self.set_controls_state(enabled=True)

    @error_handler     
    def _disable_checkbuttons_recursive(self, widget):
        """Disabilita ricorsivamente tutte le checkbox nei widget"""
        self.set_controls_state(enabled=False, control_type="checkbutton", widget=widget)
    
    @error_handler
    def _enable_checkbuttons_recursive(self, widget):
        """Abilita tutti i checkbutton in modo ricorsivo"""
        self.set_controls_state(enabled=True, control_type="checkbutton", widget=widget)

    @error_handler     
    def show_content_search_warning(self):
        """Mostra un avviso se la ricerca nei contenuti è attivata"""
        if self.search_content.get():  # Se la checkbox è attivata
            message = "Hai attivato la ricerca nei contenuti dei file."
            
            # MODIFICA QUI: Aggiungi un log di debug per verificare il valore effettivo
            self.log_debug(f"Livello di ricerca attuale: {self.search_depth.get()}")
            
            # Avviso specifico per la ricerca profonda
            if hasattr(self, 'search_depth') and self.search_depth.get() == "profonda":
                message += "\n\nATTENZIONE: Hai selezionato la ricerca PROFONDA che include tutti i tipi di file!\n" + \
                        "Questo potrebbe rallentare notevolmente la ricerca, soprattutto con grandi quantità di file."
            else:
                message += "\n\nQuesta operazione può richiedere molto più tempo, soprattutto con grandi quantità di file."
                
            message += "\n\nVuoi procedere con la ricerca nei contenuti?"
            
            return messagebox.askyesno(
                "Attenzione - Ricerca nei contenuti", 
                message,
                icon="warning"
            )
        return True   # Se la ricerca nei contenuti non è attivata, procedi senza avviso   
    
    @error_handler
    def log_debug(self, message):
        """Sistema di logging unificato che supporta sia la visualizzazione live che il logging storico
        e filtra i messaggi non necessari e le ripetizioni sulla soglia di memoria.
        """
        # Inizializza il set di messaggi già loggati se non esiste
        if not hasattr(self, 'already_logged_messages'):
            self.already_logged_messages = set()
        
        # Inizializza il contatore dei messaggi se non esiste
        if not hasattr(self, 'last_displayed_log_index'):
            self.last_displayed_log_index = 0
            
        # Filtra messaggi specifici relativi alle estensioni
        filter_prefixes = [
            "Estensioni caricate per modalità base:",
            "Estensioni caricate per modalità avanzata:",
            "Estensioni caricate per modalità profonda:"
        ]
        
        # Se il messaggio inizia con uno dei prefissi da filtrare, saltalo
        if any(message.startswith(prefix) for prefix in filter_prefixes):
            return
        
        # Verifica se è un messaggio sulla soglia di memoria
        if message.startswith("Soglia memoria calcolata:"):
            # Controlla se abbiamo già registrato un messaggio simile
            if "soglia_memoria" in self.already_logged_messages:
                return  # Salta questo messaggio perché è una ripetizione
            else:
                # Registra che abbiamo loggato un messaggio sulla soglia di memoria
                self.already_logged_messages.add("soglia_memoria")
        
        # Inizializza gli attributi necessari se mancanti
        if not hasattr(self, 'complete_debug_log_history'):
            self.complete_debug_log_history = []
        
        if not hasattr(self, 'debug_logs_queue'):
            self.debug_logs_queue = queue.Queue(maxsize=5000)
        
        if not hasattr(self, 'debug_log'):
            self.debug_log = []
        
        # Creazione timestamp
        timestamp_full = datetime.now().strftime('%Y-%m-%d %H:%M:%S.%f')[:-3]  # Con millisecondi
        timestamp_short = timestamp_full.split(' ')[1]  # Solo l'ora per i log concisi
        
        # Rilevamento automatico dei messaggi di errore o warning
        if any(keyword in message.lower() for keyword in ["error", "exception", "failed", "errore", "eccezione", "fallito"]):
            # Formattazione speciale per gli errori
            log_message_full = f"[ERRORE] {timestamp_full} - {message}"
            log_message_short = f"[{timestamp_short}] [ERRORE] {message}"
        elif any(keyword in message.lower() for keyword in ["warning", "warn", "attenzione", "avviso"]):
            # Formattazione per i warning
            log_message_full = f"[AVVISO] {timestamp_full} - {message}"
            log_message_short = f"[{timestamp_short}] [AVVISO] {message}"
        else:
            # Formattazione standard per i messaggi di info
            log_message_full = f"[INFO] {timestamp_full} - {message}"
            log_message_short = f"[{timestamp_short}] [INFO] {message}"
        
        try:
            # Aggiungi al debug_log (prima implementazione)
            self.debug_log.append(log_message_full)
            
            # Limita la dimensione del log per evitare problemi di memoria
            if len(self.debug_log) > 5000:
                self.debug_log = self.debug_log[-5000:]
                # Se abbiamo dovuto ridurre il log, aggiorniamo l'indice dell'ultimo visualizzato
                self.last_displayed_log_index = max(0, self.last_displayed_log_index - (len(self.debug_log) - 5000))
            
            # Aggiungi alla cronologia completa (seconda implementazione)
            self.complete_debug_log_history.append(log_message_short)
            
            # Limita anche questa cronologia se diventa troppo grande
            if len(self.complete_debug_log_history) > 10000:
                self.complete_debug_log_history = self.complete_debug_log_history[-10000:]
            
            # Aggiungi alla coda per la visualizzazione live con thread-safety
            with threading.Lock():
                if self.debug_logs_queue.full():
                    try:
                        self.debug_logs_queue.get_nowait()
                    except queue.Empty:
                        pass
                self.debug_logs_queue.put_nowait(log_message_short)
                
            # Verifica se la finestra di debug è aperta
            if (hasattr(self, 'debug_window') and 
                hasattr(self, 'debug_text') and
                self.debug_window.winfo_exists()):
                
                # Usa try/except per verificare che siamo nel thread principale
                try:
                    # Verifica se siamo nel thread principale
                    self.root.winfo_exists()  # Lancia eccezione se non siamo nel thread principale

                    # Controlliamo se c'è già una chiamata di aggiornamento prevista
                    if not hasattr(self, 'update_scheduled') or not self.update_scheduled:
                        self.update_scheduled = True
                        # Pianifica l'aggiornamento con un ritardo per ridurre la frequenza degli aggiornamenti
                        self.root.after(100, self.add_new_logs_to_display)

                except:
                    # Se non siamo nel thread principale, usa after_idle che è sicuro per i thread
                    try:
                        if not hasattr(self, 'update_scheduled') or not self.update_scheduled:
                            self.update_scheduled = True
                            self.root.after_idle(self.add_new_logs_to_display)
                    except:
                        # In caso di errore, non aggiornare l'interfaccia
                        pass
                
        except Exception as e:
            # In caso di errori nel logging, almeno prova a stampare sulla console
            print(f"Error in logging system: {str(e)}")
        
        # Se siamo in modalità debug, stampa anche sulla console
        if getattr(self, 'debug_mode', True):  # Default a True se debug_mode non è definito
            print(f"[DEBUG] {message}")

    @error_handler
    def log_current_settings(self, context="ricerca"):
        """Funzione centralizzata per registrare le impostazioni correnti nel log."""
        # Determina l'intestazione in base al contesto
        if context == "ricerca":
            header = "===== INIZIO RICERCA CON LE SEGUENTI IMPOSTAZIONI ====="
            footer = "===== FINE LOGGING IMPOSTAZIONI - INIZIO RICERCA ====="
        
        # Registra l'intestazione
        self.log_debug(header)
        
        # Log delle estensioni
        mode = getattr(self, 'extension_mode', 'base')
        try:
            extensions = self.get_extension_settings(mode)
            self.log_debug(f"Saved {len(extensions)} extensions for {mode} mode")
            self.log_debug(f"Extensions: {', '.join(extensions)}")
            self.log_debug(f"Aggiornata UI per modalità {mode}")
        except Exception as e:
            self.log_debug(f"Errore nel recupero delle estensioni: {str(e)}")
        
        # Log del percorso delle impostazioni
        settings_path = os.path.join(os.path.expanduser('~'), '.file_search_tool', 'application_settings.json')
        self.log_debug(f"Impostazioni dell'applicazione salvate in {settings_path}")
        
        # Log delle impostazioni di ricerca
        self.log_debug(f"Profondità ricerca: {getattr(self, 'max_depth', 0)}")
        self.log_debug(f"Cerca nei file: {self.search_files.get()}")
        self.log_debug(f"Cerca nelle cartelle: {self.search_folders.get()}")
        self.log_debug(f"Cerca nei contenuti: {self.search_content.get()}")
        self.log_debug(f"Parole intere: {getattr(self, 'whole_word_search', tk.BooleanVar()).get()}")
        
        # Log dei filtri avanzati
        size_min_kb = getattr(self, 'advanced_filters', {}).get('size_min', 0) // 1024
        size_max_kb = getattr(self, 'advanced_filters', {}).get('size_max', 0) // 1024
        self.log_debug(f"Dimensione min (KB): {size_min_kb}")
        self.log_debug(f"Dimensione max (KB): {size_max_kb}")
        
        date_min = getattr(self, 'advanced_filters', {}).get('date_min', '')
        date_max = getattr(self, 'advanced_filters', {}).get('date_max', '')
        self.log_debug(f"Data min: '{date_min}'")
        self.log_debug(f"Data max: '{date_max}'")
        
        # Log delle impostazioni di gestione memoria
        auto_memory = getattr(self, 'auto_memory_management', False)
        memory_percent = getattr(self, 'memory_usage_percent', 38)
        self.log_debug(f"Gestione memoria automatica: {auto_memory}")
        self.log_debug(f"Percentuale utilizzo memoria: {memory_percent}%")
        
        # Log delle impostazioni di blocco
        files_per_block = getattr(self, 'max_files_per_block', tk.IntVar()).get()
        parallel_blocks = getattr(self, 'max_parallel_blocks', tk.IntVar()).get()
        auto_adjust = getattr(self, 'block_size_auto_adjust', tk.BooleanVar()).get()
        prioritize = getattr(self, 'prioritize_user_folders', tk.BooleanVar()).get()
        self.log_debug(f"File per blocco: {files_per_block}")
        self.log_debug(f"Blocchi paralleli: {parallel_blocks}")
        self.log_debug(f"Auto-adattamento blocchi: {auto_adjust}")
        self.log_debug(f"Priorità cartelle utente: {prioritize}")
        
        # Log delle impostazioni di timeout e limiti
        timeout_enabled = getattr(self, 'timeout_enabled', tk.BooleanVar()).get()
        timeout_seconds = getattr(self, 'timeout_seconds', tk.IntVar()).get()
        max_files = getattr(self, 'max_files_to_check', tk.IntVar()).get()
        max_results = getattr(self, 'max_results', tk.IntVar()).get()
        worker_threads = getattr(self, 'worker_threads', tk.IntVar()).get()
        max_file_size = getattr(self, 'max_file_size_mb', tk.IntVar()).get()
        dir_size_calc = getattr(self, 'dir_size_calculation', tk.StringVar()).get()
        
        self.log_debug(f"Timeout attivo: {timeout_enabled}")
        self.log_debug(f"Secondi timeout: {timeout_seconds}")
        self.log_debug(f"Max file da controllare: {max_files}")
        self.log_debug(f"Max risultati: {max_results}")
        self.log_debug(f"Thread paralleli: {worker_threads}")
        self.log_debug(f"Dimensione max file (MB): {max_file_size}")
        self.log_debug(f"Modalità calcolo dimensioni: '{dir_size_calc}'")
        
        # Log dei percorsi esclusi
        excluded_paths = getattr(self, 'excluded_paths', [])
        if excluded_paths:
            self.log_debug(f"Numero percorsi esclusi: {len(excluded_paths)}")
            self.log_debug("Percorsi esclusi:")
            for idx, path in enumerate(excluded_paths, 1):
                self.log_debug(f"  {idx}. {path}")
        else:
            self.log_debug("Nessun percorso escluso configurato")
        
        # Registra il footer
        self.log_debug(footer)
    
    @error_handler
    def add_new_logs_to_display(self):
        """Aggiunge solo i nuovi messaggi di log alla visualizzazione rispettando il filtro corrente"""
        self.update_scheduled = False
        
        if not hasattr(self, 'debug_window') or not hasattr(self, 'debug_text') or not self.debug_window.winfo_exists():
            return
            
        # Verifica che il debug log sia inizializzato
        if not hasattr(self, 'debug_log'):
            self.debug_log = []
            self.last_displayed_log_index = 0
            
        if not hasattr(self, 'last_displayed_log_index'):
            self.last_displayed_log_index = 0
        
        # Verifica se ci sono nuovi messaggi da visualizzare
        if self.last_displayed_log_index >= len(self.debug_log):
            return
        
        # Aggiorna la lista completa di tutti i messaggi per il filtraggio
        if not hasattr(self, 'all_log_messages'):
            self.all_log_messages = self.debug_log.copy()
        else:
            # Aggiungi solo i nuovi messaggi all'elenco completo
            self.all_log_messages.extend(self.debug_log[self.last_displayed_log_index:])
        
        # Verifica se è attivo un filtro
        filter_active = hasattr(self, 'current_filter') and hasattr(self, 'filter_var') and self.filter_var.get() != "Tutti"
        
        # Se c'è un filtro attivo, riapplica il filtro su tutti i messaggi
        if filter_active:
            selected_filter = self.filter_var.get()
            # Mappa la selezione al prefisso corrispondente
            filter_map = {
                "Errore": ["[ERRORE]", "[ERROR]"],
                "Avviso": ["[AVVISO]", "[WARNING]"],
                "Info": ["[INFO]"]
            }
            
            prefixes = filter_map.get(selected_filter, [])
            
            # Filtra i messaggi in base ai prefissi
            filtered_messages = []
            if prefixes:
                filtered_messages = [msg for msg in self.all_log_messages 
                                if any(prefix in msg for prefix in prefixes)]
            
            # Aggiorna l'etichetta con il conteggio dei messaggi filtrati
            if hasattr(self, 'log_count_label'):
                self.log_count_label.config(
                    text=f"Registro di debug dell'applicazione: {len(filtered_messages)} messaggi ({selected_filter})"
                )
            
            # Pulisci il testo esistente e inserisci i messaggi filtrati
            self.debug_text.config(state=tk.NORMAL)
            self.debug_text.delete("1.0", tk.END)
            
            # Limita la visualizzazione a 5000 messaggi
            max_display = 5000
            if len(filtered_messages) > max_display:
                self.debug_text.insert(tk.END, f"[Mostrando solo gli ultimi {max_display} di {len(filtered_messages)} messaggi filtrati...]\n\n")
                filtered_messages = filtered_messages[-max_display:]
            
            # Inserisci i messaggi filtrati
            for message in filtered_messages:
                self.debug_text.insert(tk.END, message + "\n")
        else:
            # Comportamento standard senza filtro
            # Aggiorna l'etichetta con il conteggio dei messaggi
            if hasattr(self, 'log_count_label'):
                self.log_count_label.config(text=f"Registro di debug dell'applicazione: {len(self.debug_log)} messaggi")
            
            # Inserisci solo i nuovi log
            self.debug_text.config(state=tk.NORMAL)
            
            # Limita la visualizzazione a 5000 messaggi
            max_display = 5000
            
            # Se il log è troppo grande, mostra un messaggio informativo
            if len(self.debug_log) > max_display and self.last_displayed_log_index == 0:
                self.debug_text.insert(tk.END, f"[Mostrando solo gli ultimi {max_display} di {len(self.debug_log)} messaggi...]\n\n")
                # Aggiorna l'indice iniziale per iniziare dai messaggi più recenti
                self.last_displayed_log_index = len(self.debug_log) - max_display
            
            # Inserisci i nuovi log
            for i in range(self.last_displayed_log_index, len(self.debug_log)):
                self.debug_text.insert(tk.END, self.debug_log[i] + "\n")
        
        # Aggiorna l'indice dell'ultimo messaggio visualizzato
        self.last_displayed_log_index = len(self.debug_log)
        
        # Evidenzia gli errori con colori appropriati
        self.highlight_errors()
        
        # Decidi se scorrere alla fine se l'autoscroll è attivato
        if hasattr(self, 'autoscroll_var') and self.autoscroll_var.get():
            self.debug_text.see(tk.END)
            # Mantieni comunque lo scroll orizzontale a sinistra
            self.debug_text.xview_moveto(0.0)
        
        # Rendi il testo di nuovo sola lettura
        self.debug_text.config(state=tk.DISABLED)

    @error_handler
    def register_interrupt_handler(self):
        """Registra il gestore degli interrupt (CTRL+C)"""
        def handle_interrupt(sig, frame):
            if self.is_searching:
                self.stop_search_process()
            else:
                self.root.quit()
        
        signal.signal(signal.SIGINT, handle_interrupt)

    @error_handler
    def log_error(self, message, exception=None, location=None, traceback=None):
        """Registra un errore nel log di debug con dettagli aggiuntivi"""
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        
        # Crea un messaggio di errore dettagliato con prefisso in italiano
        error_message = f"[ERRORE] {timestamp} - {message}"
        
        print(f"LOG_ERROR: {error_message[:100]}...")

        # Registra nel log di debug
        if hasattr(self, 'debug_log'):
            self.debug_log.append(error_message)
            print(f"Debug log size: {len(self.debug_log)}")
            
        # Aggiungi informazioni sulla posizione dell'errore
        if location:
            error_message += f" | Posizione: {location}"
        
        # Aggiungi dettagli sull'eccezione
        if exception:
            exc_type = type(exception).__name__
            exc_details = str(exception)
            
            error_message += f" | Eccezione: [{exc_type}]: {exc_details}"
        
        # Gestisci il traceback (priorità al traceback esplicito se fornito)
        traceback_info = ""
        if traceback:
            # Usa il traceback fornito esplicitamente
            error_message += f"\n--- Traceback ---\n{traceback}\n-----------------"
        elif exception:
            # Ottieni traceback se disponibile e non già fornito
            try:
                import traceback as tb_module
                tb_info = tb_module.format_exc().split('\n')
                # Prendi solo le righe più rilevanti del traceback
                if len(tb_info) > 3:
                    traceback_info = " | " + " > ".join(tb_info[-4:-1])
                    error_message += traceback_info
            except:
                pass
        
        # Registra nel log di debug
        if hasattr(self, 'debug_log'):
            self.debug_log.append(error_message)
            
            # Limita la dimensione del log per evitare problemi di memoria
            if len(self.debug_log) > 1000:
                self.debug_log = self.debug_log[-1000:]
        
        # Aggiungi anche al log del sistema se possibile
        try:
            import logging
            logging.error(error_message)
        except:
            pass
        
        # Salva periodicamente il log su file (max una volta ogni 5 minuti)
        try:
            if not hasattr(self, 'last_error_log_save') or \
            (datetime.now() - self.last_error_log_save).total_seconds() > 300:
                self.last_error_log_save = datetime.now()
                
                # Salva nel file di log
                log_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logs")
                os.makedirs(log_dir, exist_ok=True)
                
                error_log_file = os.path.join(log_dir, f"error_log_{datetime.now().strftime('%Y%m%d')}.txt")
                with open(error_log_file, "a", encoding="utf-8") as f:
                    f.write(error_message + "\n")
        except:
            # Non interrompere il flusso dell'applicazione se il salvataggio fallisce
            pass
    
    @error_handler
    def highlight_errors(self):
        """Evidenzia errori e avvisi nel testo con colori diversi"""
        if not hasattr(self, 'debug_text') or not self.debug_text:
            return
               
        # Prima rimuovi tutti i tag esistenti (più efficiente che rimuoverli riga per riga)
        self.debug_text.tag_remove("error", "1.0", "end")
        self.debug_text.tag_remove("warning", "1.0", "end")
        self.debug_text.tag_remove("info", "1.0", "end")
        
        # Cerca tutti i messaggi di errore (sia in inglese che in italiano)
        start_index = "1.0"
        while True:
            # Cerca i prefissi di errore
            error_pos = self.debug_text.search(r"\[(ERROR|ERRORE)\]", start_index, "end", regexp=True)
            if not error_pos:
                break
            
            # Trova la fine della riga
            line_end = self.debug_text.index(f"{error_pos} lineend")
            
            # Applica il tag di errore
            self.debug_text.tag_add("error", error_pos, line_end)
            
            # Passa alla posizione successiva
            start_index = line_end
        
        # Cerca tutti i messaggi di avviso (sia in inglese che in italiano)
        start_index = "1.0"
        while True:
            # Cerca i prefissi di avviso
            warning_pos = self.debug_text.search(r"\[(WARNING|AVVISO)\]", start_index, "end", regexp=True)
            if not warning_pos:
                break
            
            # Trova la fine della riga
            line_end = self.debug_text.index(f"{warning_pos} lineend")
            
            # Applica il tag di warning
            self.debug_text.tag_add("warning", warning_pos, line_end)
            
            # Passa alla posizione successiva
            start_index = line_end
        
        # Cerca tutti i messaggi informativi
        start_index = "1.0"
        while True:
            # Cerca i prefissi di info
            info_pos = self.debug_text.search(r"\[INFO\]", start_index, "end", regexp=True)
            if not info_pos:
                break
            
            # Trova la fine della riga
            line_end = self.debug_text.index(f"{info_pos} lineend")
            
            # Applica il tag di info
            self.debug_text.tag_add("info", info_pos, line_end)
            
            # Passa alla posizione successiva
            start_index = line_end

    @error_handler
    def check_and_notify_missing_libraries(self):
        """Verifica e notifica l'utente di eventuali librerie mancanti"""
        missing = []
        
        if not file_format_support["docx"]:
            missing.append("python-docx (per file Word)")
        if not file_format_support["pdf"]: 
            missing.append("PyPDF2 (per file PDF)")
        if not file_format_support["pptx"]:
            missing.append("python-pptx (per file PowerPoint)")
        if not file_format_support["excel"]:
            missing.append("openpyxl (per file Excel)")
        if not file_format_support["rtf"]:
            missing.append("striprtf (per file RTF)")
        if missing:
            message = "Alcune funzionalità di ricerca nei contenuti sono disabilitate.\n\n"
            message += "Per abilitare il supporto completo ai vari formati di file, installa le seguenti librerie:\n\n"
            
            for lib in missing:
                message += f"- {lib}\n"
            
            message += "\nPuoi installarle con il comando:\n"
            message += "pip install " + " ".join([lib.split(" ")[0] for lib in missing])
            
            # Mostra la notifica dopo un breve ritardo per permettere all'UI di caricarsi
            self.root.after(1000, lambda: messagebox.showinfo("Librerie opzionali mancanti", message))

    @error_handler
    def update_datetime(self):
        current_time = datetime.now().strftime('%d/%m/%Y %H:%M:%S')
        self.datetime_var.set(f"Data: {current_time} | Utente: {self.user_var.get()}")
        self.root.after(1000, self.update_datetime)

    @error_handler
    def browse_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.search_path.set(directory)
            self.update_disk_info(directory)

    @error_handler
    def optimize_system_search(self, path):
        """Ottimizza la ricerca per percorsi di sistema come C:/ impostando parametri appropriati"""
        # Riconosci più tipi di percorsi di sistema
        is_system_path = path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"]
        
        # Verifica più accurata per percorsi di rete (Windows e Linux/Mac)
        is_network_path = path.startswith('\\\\') or path.startswith('//') or (':/' in path and not path[0].isalpha())
        
        # Verifica percorsi Unix
        is_unix_root = path in ['/', '/home', '/usr', '/var', '/etc'] or path.startswith('/mnt/') or path.startswith('/media/')
        
        if is_system_path or is_network_path or is_unix_root:
            # Salva i parametri attuali
            original_params = {
                "max_files": self.max_files_to_check.get(),
                "worker_threads": self.worker_threads.get(),
                "timeout": self.timeout_seconds.get() if self.timeout_enabled.get() else None,
                "max_file_size": self.max_file_size_mb.get(),
                "block_size": self.max_files_per_block.get()
            }
            
            # Applica parametri ottimizzati
            self.max_files_to_check.set(1000000)  # Un milione di file
            
            # Ottimizza il numero di thread in base al tipo di percorso
            if is_network_path:
                # Per percorsi di rete, usa meno thread per evitare congestione
                self.worker_threads.set(max(2, min(6, os.cpu_count() or 4)))
                # Aumenta la dimensione dei blocchi per ridurre overhead di rete
                self.max_files_per_block.set(2000)
                # Riduci la dimensione massima del file per percorsi di rete
                self.max_file_size_mb.set(min(50, self.max_file_size_mb.get()))
                self.log_debug("Applicata ottimizzazione per percorso di rete")
            else:
                # Per percorsi locali di sistema, usa più thread
                self.worker_threads.set(min(12, os.cpu_count() or 4))
                self.log_debug("Applicata ottimizzazione per percorso di sistema locale")
            
            if self.timeout_enabled.get():
                self.timeout_seconds.set(max(3600, self.timeout_seconds.get()))  # Minimo 1 ora
            
            # Notifica l'utente con informazioni più dettagliate
            messagebox.showinfo(
                "Ricerca su sistema o server",
                "Stai avviando una ricerca su un percorso di sistema o server.\n\n"
                f"Tipo percorso rilevato: {'Server/Rete' if is_network_path else 'Sistema'}\n"
                f"Thread di ricerca: {self.worker_threads.get()}\n"
                f"Limite file: {self.max_files_to_check.get():,}\n\n"
                "La ricerca potrebbe richiedere tempo e risorse significative."
            )
            
            return original_params
        return None
    
    @error_handler
    def is_network_path(self, path):
        """Verifica se il percorso è un percorso di rete"""
        # Percorsi UNC Windows (\\server\share)
        return PathUtils.is_network_path(path)

    @error_handler
    def show_optimization_tips(self, path):
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]:
            # Determina la lettera del disco corrente
            if os.name == 'nt':  # Windows
                if path.lower() in ["c:/", "c:\\"]:
                    drive_letter = "C:"
                elif path.lower() in ["d:/", "d:\\"]:
                    drive_letter = "D:"
                elif path.lower() in ["e:/", "e:\\"]:
                    drive_letter = "E:"
                else:
                    drive_letter = path[:2] if len(path) >= 2 else "Disco"
            else:
                drive_letter = "Disco"  # Per sistemi non Windows
            
            optimization_done = False
            
            # Verifica se l'utente ha già implementato le ottimizzazioni
            if hasattr(self, 'excluded_paths'):
                windows_excluded = any("windows" in p.lower() for p in self.excluded_paths)
                program_files_excluded = any("program files (x86)" in p.lower() for p in self.excluded_paths)
                optimization_done = windows_excluded and program_files_excluded
            
            # Verifica la profondità di ricerca
            depth_optimized = hasattr(self, 'max_depth') and self.max_depth >= 1 and self.max_depth <= 10
            
            # Se le ottimizzazioni non sono già state applicate, mostra il messaggio
            if not (optimization_done and depth_optimized):
                response = messagebox.askyesno(
                    "Ottimizza la ricerca",
                    f"Stai per avviare una ricerca sull'intero disco {drive_letter}\\\n\n"
                    "Per migliorare notevolmente le prestazioni, è consigliato:\n\n"
                    "1. Escludere le cartelle di sistema (Windows, Program Files)\n"
                    "2. Escludere le cartelle di altri utenti\n"
                    "3. Limitare la profondità di ricerca a 5-10 livelli\n\n"
                    "Vuoi applicare automaticamente queste ottimizzazioni?",
                    icon="question"
                )
                
                if response:
                    # Applica le ottimizzazioni
                    if not hasattr(self, 'excluded_paths'):
                        self.excluded_paths = []
                    
                    # Aggiungi esclusioni di sistema
                    system_paths = ["C:/Windows", "C:/Program Files", "C:/Program Files (x86)"]
                    for sys_path in system_paths:
                        if sys_path not in self.excluded_paths:
                            self.excluded_paths.append(sys_path)
                    
                    # Escludi altri utenti
                    current_user = getpass.getuser()
                    users_dir = "C:/Users"
                    if os.path.exists(users_dir):
                        try:
                            for user in os.listdir(users_dir):
                                user_path = os.path.join(users_dir, user)
                                if user.lower() != current_user.lower() and user.lower() not in ["public", "default", "all users"]:
                                    if user_path not in self.excluded_paths:
                                        self.excluded_paths.append(user_path)
                        except Exception as e:
                            self.log_debug(f"Errore nell'esclusione degli altri utenti: {str(e)}")
                    
                    # Imposta la profondità di ricerca a 7 (valore ragionevole)
                    # Usa una verifica prima di accedere a depth_spinbox
                    if hasattr(self, 'depth_spinbox'):
                        self.depth_spinbox.set("7")
                    self.max_depth = 7
                    
                    # Notifica l'utente
                    messagebox.showinfo(
                        "Ottimizzazioni applicate",
                        f"Le ottimizzazioni sono state applicate:\n\n"
                        f"• Escluse {len(system_paths)} cartelle di sistema\n"
                        f"• Escluse le cartelle di altri utenti\n"
                        f"• Profondità di ricerca impostata a 7 livelli\n\n"
                        f"La ricerca sarà ora più veloce e stabile."
                    )
                    
                    return True
                else:
                    # NUOVA PARTE: Se l'utente risponde "No", imposta esplicitamente la profondità a 0
                    # Usa una verifica prima di accedere a depth_spinbox
                    if hasattr(self, 'depth_spinbox'):
                        self.depth_spinbox.set("0")
                    self.max_depth = 0
                    
                    # Notifica l'utente che sta per iniziare una ricerca illimitata
                    self.log_debug("Utente ha rifiutato le ottimizzazioni, impostata profondità illimitata (0)")
                    
                    # Opzionalmente, puoi mostrare un messaggio di avviso
                    messagebox.showinfo(
                        "Ricerca illimitata",
                        "Stai per avviare una ricerca illimitata sull'intero disco.\n\n"
                        "La profondità di ricerca è stata impostata a illimitata (0).\n"
                        "Questa operazione potrebbe richiedere molto tempo e risorse di sistema."
                    )
                
                return False
            
            return False  # Per percorsi non di sistema, non fare nulla
        
    @error_handler
    def debug_exclusions(self):
        """Visualizza lo stato corrente delle esclusioni per debug"""
        if hasattr(self, 'excluded_paths'):
            paths = '\n'.join(self.excluded_paths) if self.excluded_paths else "Nessun percorso escluso"
            messagebox.showinfo("Debug esclusioni", 
                            f"Stato esclusioni:\n{paths}\n\n"
                            f"Totale percorsi: {len(self.excluded_paths)}")
        else:
            messagebox.showinfo("Debug esclusioni", "Lista esclusioni non inizializzata")

    @error_handler
    def restart_as_admin(self):
        """Riavvia l'applicazione con privilegi di amministratore"""
        try:
            import ctypes, sys, os
            if sys.platform == 'win32':
                script = sys.argv[0]
                params = ' '.join(sys.argv[1:])
                self.log_debug("Tentativo di riavvio come amministratore")
                
                # Se il percorso di ricerca è impostato, passa quel percorso come parametro
                if self.search_path.get():
                    if not params:
                        params = f'"{self.search_path.get()}"'
                    
                # Avvia l'applicazione con privilegi elevati
                ctypes.windll.shell32.ShellExecuteW(None, "runas", sys.executable, f'"{script}" {params}', None, 1)
                
                # Chiude l'istanza corrente
                self.root.quit()
            else:
                messagebox.showinfo("Informazione", "Questa funzione è disponibile solo su Windows")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile riavviare come amministratore: {str(e)}")
            self.log_debug(f"Errore nel riavvio come admin: {str(e)}")

    @error_handler
    def optimize_disk_search_order(self, path, directories):
        """Ottimizza l'ordine di esplorazione delle directory quando si cerca in un disco di sistema.
        Dà priorità alle cartelle più importanti come Users, Documents, ecc.
        """
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"]:
            prioritized = []
            normal = []
            
            # Cartelle ad alta priorità (esplora per prime)
            high_priority_folders = ["users", "documents", "documents and settings", "desktop", "downloads"]
            
            # Cartelle a bassa priorità (esplora per ultime)
            low_priority_folders = ["windows", "program files", "program files (x86)", "programdata", "$recycle.bin", "system volume information"]
            
            for dir_path in directories:
                dir_name = os.path.basename(dir_path).lower()
                
                if any(priority in dir_name for priority in high_priority_folders):
                    prioritized.insert(0, dir_path)  # Aggiungi all'inizio per priorità più alta
                elif any(low in dir_name for low in low_priority_folders):
                    normal.append(dir_path)  # Aggiungi alla fine
                else:
                    prioritized.append(dir_path)  # Priorità media
                    
            return prioritized + normal
        
        # Se non è un percorso di sistema, restituisci l'elenco originale
        return directories

    @error_handler
    def start_search(self):
        # Assicurati che qualsiasi ricerca precedente sia completamente terminata
        self.reset_search_state()
        
        # Reset the total files size label at the start of a new search
        if hasattr(self, 'total_files_size_label'):
            self.total_files_size_label.config(text="Dimensione totale: 0 B (0 file)")
        
        # Pulisci risultati precedenti
        for item in self.results_list.get_children():
            self.results_list.delete(item)
        
        # Ottieni i valori direttamente dai widget
        search_path = self.search_path.get().strip()
        keywords = self.keyword_entry.get().strip()
        
        # Stampa di debug (puoi rimuoverla in produzione)
        self.log_debug(f"Avvio ricerca - Percorso: '{search_path}', Parole chiave: '{keywords}'")
        
        # 1. Verifica il percorso
        if not search_path:
            messagebox.showerror("Errore", "Inserisci una directory di ricerca valida")
            return
        
        if not os.path.exists(search_path):
            messagebox.showerror("Errore", f"Il percorso specificato non esiste:\n{search_path}")
            return
        
        # 2. Verifica le parole chiave
        placeholder = "Scrivi la parola da ricercare..."
        keywords_raw = self.keyword_entry.get()  # Ottieni il valore grezzo senza strip()

        # Verifica se è vuoto, solo spazi, o è il placeholder
        if (not keywords_raw or 
            keywords_raw.strip() == "" or 
            keywords_raw == placeholder or 
            self.keyword_entry.cget("foreground") == "gray"):
            
            messagebox.showerror("Errore", "Inserisci le parole chiave da cercare")
            return

        # 3. Verifica colore (controllo extra per il placeholder)
        try:
            if self.keyword_entry.cget("foreground") == "gray":
                messagebox.showerror("Errore", "Inserisci le parole chiave da cercare")
                return
        except:
            pass
        
        # Log delle impostazioni correnti
        self.log_current_settings(context="ricerca")

        # Mostra avviso se la ricerca nei contenuti è attivata
        if not self.show_content_search_warning():
            return  # Interrompi se l'utente annulla
        
        # Aggiorna le informazioni del disco qui, quando l'utente clicca su cerca
        calculation_mode = self.dir_size_calculation.get()
        
        # CORREZIONE: Determina se calcolare la dimensione in modo esplicito
        should_calculate = False
        if calculation_mode == "stimato" or calculation_mode == "accurato":
            should_calculate = True
            self.log_debug(f"Calcolo dimensione directory abilitato: {calculation_mode}")
        else:
            self.log_debug("Calcolo dimensione directory disabilitato")
        
        # Passa il flag corretto a update_disk_info
        self.update_disk_info(path=search_path, calculate_dir_size=should_calculate)

        self.show_optimization_tips(self.search_path.get())
        
        # Reset per la nuova ricerca
        self.stop_search = False
        
        # Imposta is_searching PRIMA di disabilitare i controlli
        self.is_searching = True

        # Avvia il monitoraggio della memoria
        self.start_memory_monitoring()
        
        # CRITICO: Crea un nuovo executor per questa ricerca
        max_workers = max(1, min(32, self.worker_threads.get()))
        self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
        
        original_params = self.optimize_system_search(self.search_path.get())
        self.start_search_watchdog()
        
        # Aggiorna l'orario di avvio e resetta l'orario di fine
        current_time = datetime.now().strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")

        # Salva i parametri di ricerca
        self.last_search_params = {
            "path": self.search_path.get(),
            "keywords": self.keywords.get(),
            "search_files": self.search_files.get(),
            "search_folders": self.search_folders.get(),
            "search_content": self.search_content.get()
        }
        
        # Aggiungi alla cronologia di ricerca se non già presente
        if self.keywords.get() and self.keywords.get() not in [h["keywords"] for h in self.search_history]:
            self.search_history.append(self.last_search_params.copy())
            if len(self.search_history) > 10:  # Mantieni solo le ultime 10 ricerche
                self.search_history.pop(0)
        
        # Disabilita i controlli
        self.disable_all_controls()
        
        # IMPORTANTE: Il pulsante di interruzione deve essere abilitato DOPO aver disabilitato tutti i controlli
        if hasattr(self, 'stop_button'):
            self.stop_button["state"] = "normal"
            # Forza aggiornamento immediato del bottone
            self.stop_button.update()
        
        self.progress_bar["value"] = 0
        self.status_label["text"] = "Ricerca in corso..."
        
        # Aggiorna il valore della profondità massima - FIX
        try:
            if hasattr(self, 'depth_spinbox') and self.depth_spinbox.winfo_exists():
                try:
                    self.max_depth = int(self.depth_spinbox.get())
                except (ValueError, tk.TclError) as e:
                    self.log_debug(f"Errore nell'accesso al depth_spinbox: {str(e)}")
                    self.max_depth = getattr(self, 'max_depth', 0)  # Usa il valore esistente o default
            else:
                # Se il widget non esiste, usa il valore memorizzato come attributo 
                # o il valore predefinito se non esiste
                self.max_depth = getattr(self, 'max_depth', 0)
                self.log_debug(f"Widget depth_spinbox non disponibile, usando valore: {self.max_depth}")
        except Exception as e:
            self.log_debug(f"Errore generale nell'aggiornamento della profondità: {str(e)}")
            self.max_depth = 0  # Valore predefinito sicuro
        
        # Ottieni le parole chiave di ricerca
        search_terms = [term.strip() for term in self.keywords.get().split(',') if term.strip()]
        
        # DEBUG: Verifica che la ricerca sia correttamente impostata
        self.log_debug(f"STATO RICERCA: is_searching={self.is_searching}, stop_search={self.stop_search}")
        self.log_debug(f"Stato pulsante interruzione: {self.stop_button['state']}")
            
        # Avvia la ricerca in un thread separato
        self.search_results = []  # Resetta i risultati
        
        # INIZIO MODIFICHE WINDOWS SEARCH
        # Verifica se è possibile utilizzare Windows Search
        use_windows_search = False
        
        # Inizializza WindowsSearchHelper se non già fatto
        if not hasattr(self, 'windows_search_helper'):
            self.windows_search_helper = WindowsSearchHelper(logger=self)
        
        # Determina se usare Windows Search
        if (self.windows_search_helper.available and 
                hasattr(self, 'use_windows_search_var') and 
                self.use_windows_search_var.get()):
            
            # Verifica se il percorso è indicizzato
            is_indexed = self.windows_search_helper.index_status(search_path)
            
            if is_indexed:
                use_windows_search = True
                self.log_debug(f"Il percorso {search_path} è indicizzato. Utilizzando Windows Search.")
            else:
                self.log_debug(f"Il percorso {search_path} non è indicizzato. Utilizzando ricerca standard.")
        
        # Avvia il thread di ricerca appropriato
        if use_windows_search:
            self.log_debug("INFO: Avvio ricerca utilizzando Windows Search (windows.edb)")
            search_thread = threading.Thread(target=self._windows_search_thread, 
                                        args=(search_path, search_terms, self.search_content.get()))
        else:
            self.log_debug("INFO: Avvio ricerca standard dei file")
            search_thread = threading.Thread(target=self._search_thread, 
                                        args=(search_path, search_terms, self.search_content.get()))
        # FINE MODIFICHE WINDOWS SEARCH
        
        search_thread.daemon = True
        search_thread.start()
        
        # Memorizza l'orario di inizio come oggetto datetime
        self.search_start_time = datetime.now()
        current_time = self.search_start_time.strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")
        self.total_time_label.config(text="--:--")  # Reset del tempo totale

        # Avvia l'aggiornamento della progress bar
        self.update_progress()
        
        # Salva original_params come attributo per ripristinarli dopo la ricerca
        if original_params:
            self.original_system_search_params = original_params

        # IMPORTANTE: Doppio controllo del pulsante di interruzione dopo aver iniziato la ricerca
        if hasattr(self, 'stop_button'):
            # Usa un approccio a ritardo doppio per garantire che sia abilitato
            self.root.after(100, lambda: self.stop_button.configure(state="normal"))
            self.root.after(500, lambda: self.stop_button.configure(state="normal"))
        
        self.log_debug("Ricerca avviata correttamente con pulsante interruzione abilitato")

    # Aggiungi questo metodo alla classe FileSearchApp
    def _windows_search_thread(self, path, keywords, search_content):
        """Thread di ricerca che utilizza Windows Search"""
        start_time = time.time()
        results = []
        
        try:
            self.log_info(f"Utilizzo Windows Search per la ricerca in {path}")
            
            # Prepara le estensioni filtrate se necessario
            extensions = None
            if not self.search_all_extensions_var.get():
                extensions = self.get_enabled_extensions()
            
            # Esegui la ricerca con Windows Search
            found_files = self.windows_search_helper.search_files(
                search_path=path,
                keywords=keywords,
                file_extensions=extensions
            )
            
            # Processa i risultati
            for file_path in found_files:
                if not self.search_running:
                    break
                    
                # Verifica che il file soddisfi i criteri
                if not self.should_skip_file(file_path):
                    # Anche se Windows Search ha già cercato nel contenuto,
                    # per coerenza con il resto dell'app, verificare ancora
                    # se search_content è attivo e il file è supportato
                    file_matches = True
                    
                    if search_content and self.should_search_content(file_path):
                        # Verificare ulteriormente le corrispondenze nel contenuto
                        # con la nostra logica personalizzata se necessario
                        try:
                            content = self.get_file_content(file_path)
                            file_matches = any(k.lower() in content.lower() for k in keywords)
                        except Exception as e:
                            self.log_error(f"Errore durante lettura contenuto: {file_path}", exception=e)
                    
                    if file_matches:
                        # Aggiungi ai risultati
                        file_info = self.create_file_info(file_path)
                        if file_info:
                            results.append(file_info)
                
                # Aggiorna il contatore
                self.total_files_processed += 1
                
        except Exception as e:
            self.log_error("Errore durante la ricerca con Windows Search", exception=e)
        
        # Completa la ricerca
        self.search_results = results
        self.search_running = False
        self.search_completed = True
        
        elapsed_time = time.time() - start_time
        self.total_search_time = elapsed_time
        
        self.log_info(f"Ricerca completata in {elapsed_time:.2f} secondi. "
                    f"Trovati {len(results)} risultati su {self.total_files_processed} file analizzati.")
        
        # Aggiorna l'interfaccia
        self.root.after(100, self.update_results_list)

    @error_handler
    def start_search_watchdog(self):
        """Avvia un timer di controllo per rilevare se la ricerca si è bloccata"""
        self.last_progress_time = time.time()
        self.watchdog_active = True
        self.check_search_progress()

    @error_handler
    def check_search_progress(self):
        """Controlla se la ricerca ha fatto progressi recentemente"""
        if not self.is_searching or not hasattr(self, 'watchdog_active') or not self.watchdog_active:
            return
            
        current_time = time.time()
        elapsed_since_progress = current_time - self.last_progress_time
        
        # Se nessun progresso per 3 minuti, considera la ricerca bloccata
        if elapsed_since_progress > 180:  # 3 minuti
            self.log_debug("La ricerca sembra bloccata - tentativo di recupero")
            
            # Prova a recuperare forzando la chiusura dell'executor e riavviandolo
            if hasattr(self, 'search_executor') and self.search_executor:
                try:
                    self.search_executor.shutdown(wait=False)
                    self.search_executor = concurrent.futures.ThreadPoolExecutor(
                        max_workers=max(1, self.worker_threads.get())
                    )
                    self.last_progress_time = time.time()  # Reset timer
                    self.status_label["text"] = "Recupero dalla ricerca bloccata..."
                except Exception as e:
                    self.log_debug(f"Errore durante il recupero della ricerca: {str(e)}")
            
        # Controlla di nuovo tra 30 secondi
        self.root.after(30000, self.check_search_progress)

    @error_handler # Aggiungi questo nuovo metodo per calcolare il tempo totale
    def update_total_time(self):
        """Calcola e visualizza il tempo totale trascorso tra inizio e fine ricerca"""
        if self.search_start_time:
            end_time = datetime.now()
            time_diff = end_time - self.search_start_time
            
            # Converti in minuti e secondi
            total_seconds = int(time_diff.total_seconds())
            minutes = total_seconds // 60
            seconds = total_seconds % 60
            
            # Formatta la stringa del tempo totale
            if minutes > 0:
                total_time_str = f"{minutes}min {seconds}sec"
            else:
                total_time_str = f"{seconds}sec"
            
            self.total_time_label.config(text=total_time_str) 

    @error_handler
    def calculate_block_priority(self, directory):
        """Calcola la priorità del blocco (numerica, più bassa = più alta priorità)"""
        # Se l'opzione è disabilitata, usa priorità standard per tutti
        if not self.prioritize_user_folders.get():
            return 1
            
        # Altrimenti, applica la prioritizzazione configurata
        # Dai priorità alle cartelle degli utenti
        if "users" in directory.lower() or "documenti" in directory.lower() or "desktop" in directory.lower():
            return 0  # Alta priorità
        # Dai priorità medio-alta alle cartelle di dati
        elif "data" in directory.lower() or "database" in directory.lower() or "downloads" in directory.lower():
            return 1  # Priorità media-alta
        # Dai priorità bassa a cartelle di sistema
        elif "windows" in directory.lower() or "program files" in directory.lower():
            return 3  # Priorità bassa
        # Priorità standard per altre cartelle
        return 2
    
    @error_handler
    def initialize_block_queue(self, root_path, block_queue, visited_dirs, files_checked, keywords, search_content, futures):
        """Inizializza la coda di blocchi con il percorso principale"""
        try:
            # Verifica se l'executor esiste, altrimenti crealo
            if not hasattr(self, 'executor'):
                import concurrent.futures
                self.executor = concurrent.futures.ThreadPoolExecutor(max_workers=os.cpu_count() or 4)
                self.log_debug("Executor creato dinamicamente durante l'inizializzazione della coda blocchi")
                
            # Verifica se il percorso è valido
            if not os.path.exists(root_path) or root_path in visited_dirs:
                return
                
            # Aggiungi il percorso alla lista delle cartelle visitate per evitare cicli
            try:
                visited_dirs.add(os.path.realpath(root_path))
            except:
                visited_dirs.add(root_path)
            
            # Ottieni la profondità massima dalle impostazioni
            max_depth = self.depth_var.get() if hasattr(self, 'depth_var') else self.max_depth
            
            # Log dell'inizio della ricerca con informazioni sulla profondità
            self.log_debug(f"Inizializzata ricerca in {root_path} con profondità {'illimitata' if max_depth == 0 else max_depth}")
            
            # CORREZIONE: Aggiungere sempre il percorso principale alla coda con profondità 0
            priority_root = self.calculate_block_priority(root_path)
            block_queue.put((priority_root, root_path, 0))  # Profondità 0 per la directory principale
            
            # Processa i file nella directory principale
            try:
                items = os.listdir(root_path)
                
                # Prima elabora le sottodirectory
                for item in items:
                    item_path = os.path.join(root_path, item)
                    try:
                        if os.path.isdir(item_path):
                            # Salta directory nascoste se richiesto
                            if self.ignore_hidden.get() and (item.startswith('.') or 
                                (os.name == 'nt' and os.path.exists(item_path) and 
                                os.stat(item_path).st_file_attributes & 2)):
                                continue
                            
                            # Salta directory escluse
                            if hasattr(self, 'excluded_paths') and any(
                                item_path.lower().startswith(excluded.lower()) 
                                for excluded in self.excluded_paths):
                                continue
                            
                            # Aggiungiamo sempre le sottocartelle di primo livello con profondità 1
                            priority = self.calculate_block_priority(item_path)
                            block_queue.put((priority, item_path, 1))  # Profondità 1 per le sottocartelle dirette
                            
                            try:
                                visited_dirs.add(os.path.realpath(item_path))
                            except:
                                visited_dirs.add(item_path)
                    except Exception as e:
                        self.log_debug(f"Errore nell'aggiunta del blocco {item_path}: {str(e)}")
                
                # Poi elabora i file
                for item in items:
                    if hasattr(self, 'stop_search') and self.stop_search:
                        return
                        
                    item_path = os.path.join(root_path, item)
                    try:
                        if not os.path.isdir(item_path):
                            if self.search_files.get():
                                # Verifica file nascosti
                                if self.ignore_hidden.get() and (item.startswith('.') or 
                                    (os.name == 'nt' and os.path.exists(item_path) and 
                                    os.stat(item_path).st_file_attributes & 2)):
                                    continue
                                
                                # Salta se il file dovrebbe essere ignorato
                                if self.should_skip_file(item_path):
                                    continue
                                
                                # Processa direttamente i file nella directory principale
                                files_checked[0] += 1
                                if files_checked[0] > self.max_files_to_check.get():
                                    if hasattr(self, 'stop_search'):
                                        self.stop_search = True
                                    return
                                
                                # Usa l'executor appropriato (search_executor o executor)
                                executor_to_use = self.search_executor if hasattr(self, 'search_executor') else self.executor
                                future = executor_to_use.submit(self.process_file, item_path, keywords, search_content)
                                futures.append(future)
                    except Exception as e:
                        self.log_debug(f"Errore nell'elaborazione del file {item_path}: {str(e)}")
                    
            except PermissionError:
                self.log_debug(f"Permesso negato per la directory {root_path}")
            except Exception as e:
                self.log_debug(f"Errore nell'inizializzazione dei blocchi da {root_path}: {str(e)}")
                
            # Aggiorna la UI con i progressi
            if hasattr(self, 'update_status_label'):
                self.update_status_label(f"Scansione directory: {root_path}")
                
        except Exception as e:
            self.log_error(f"Errore nell'inizializzazione della coda blocchi: {str(e)}")

    @error_handler
    def process_file_batch(self, file_batch, files_checked, last_update_time, 
                      calculation_enabled, keywords, search_content, futures):
        """Processa un batch di file in modo ottimizzato"""
        for file_path in file_batch:
            try:
                # Verifica limite file
                files_checked[0] += 1
                
                # Verifica ottimizzata del limite di file
                if files_checked[0] > self.max_files_to_check.get():
                    self.stop_search = True
                    self.progress_queue.put(("status", 
                        f"Limite di {self.max_files_to_check.get():,} file controllati raggiunto. "
                        f"Aumenta il limite nelle opzioni per cercare più file."))
                    return
                    
                # Calcola dimensione per statistiche se abilitato
                if calculation_enabled:
                    try:
                        if os.path.isfile(file_path) and not os.path.islink(file_path):
                            file_size = os.path.getsize(file_path)
                            self.current_search_size += file_size
                            
                            # Aggiorna la dimensione mostrata periodicamente
                            current_time = time.time()
                            if files_checked[0] % 1000 == 0 or (current_time - last_update_time[0]) > 5:
                                self.progress_queue.put(("update_dir_size", self.current_search_size))
                                last_update_time[0] = current_time
                    except:
                        pass  # Ignora errori durante il calcolo della dimensione
                
                # Verifica se il processo può ancora eseguire submit (non è stato chiuso)
                try:
                    if self.search_executor and not self.search_executor._shutdown:
                        # Usa la versione con timeout per evitare blocchi
                        future = self.search_executor.submit(
                            self.process_file_with_timeout, 
                            file_path, keywords, search_content
                        )
                        futures.append(future)
                    else:
                        # Fallback diretto se l'executor è chiuso
                        result = self.process_file(file_path, keywords, search_content)
                        if result:
                            self.search_results.append(result)
                except Exception as e:
                    self.log_debug(f"Errore nell'elaborazione parallela del file {file_path}: {str(e)}")
                    # Fallback se l'executor fallisce
                    try:
                        result = self.process_file(file_path, keywords, search_content)
                        if result:
                            self.search_results.append(result)
                    except:
                        pass
                    
            except Exception as e:
                self.log_debug(f"Errore nell'aggiunta del file {file_path} alla coda: {str(e)}")

    @error_handler
    def process_blocks(self, block_queue, visited_dirs, start_time, timeout, is_system_search, 
                  files_checked, dirs_checked, last_update_time, path, keywords, search_content, futures):
        """Elabora i blocchi dalla coda in base alla priorità"""
        # Determina il numero massimo di file per blocco in base alle impostazioni
        max_files_in_block = self.max_files_per_block.get()
        
        # CORREZIONE: Ottieni il valore della profondità massima in modo coerente
        max_depth = self.depth_var.get() if hasattr(self, 'depth_var') else self.max_depth
        
        # CORREZIONE: Semplifica la logica di tracking - true se c'è un limite, false se è illimitato (0)
        using_limited_depth = max_depth > 0
        
        # Ottimizza il calcolo della dimensione
        calculation_enabled = self.dir_size_calculation.get() != "disabilitato"
        
        # Log dell'inizio dell'elaborazione dei blocchi
        self.log_debug(f"Avvio elaborazione blocchi con profondità {'limitata a '+str(max_depth) if using_limited_depth else 'illimitata'}")
        
        # Adatta automaticamente la dimensione del blocco
        if self.block_size_auto_adjust.get():
            # Verifica se è un percorso di rete
            is_network_path = path.startswith('\\\\') or path.startswith('//')
            
            # Ottimizzazione per percorsi di rete
            if is_network_path:
                # Blocchi più grandi per la rete per ridurre l'overhead
                max_files_in_block = max(2000, max_files_in_block)
                self.log_debug(f"Dimensione blocco adattata per rete: {max_files_in_block}")
            # Ottimizzazione per ricerche di sistema
            elif is_system_search:
                max_files_in_block = max(1500, max_files_in_block)
                self.log_debug(f"Dimensione blocco adattata per sistema: {max_files_in_block}")
            # Ottimizzazione per grandi quantità di file
            elif files_checked[0] > 50000:
                max_files_in_block = max(3000, max_files_in_block)
                self.log_debug(f"Dimensione blocco adattata per ricerca grande: {max_files_in_block}")
        
        # Ottimizza il numero di blocchi paralleli in base al carico
        max_parallel = self.max_parallel_blocks.get()
        
        # Adatta il parallelismo in base al percorso
        if path.startswith('\\\\') or path.startswith('//'):
            # Limita il parallelismo su percorsi di rete
            max_parallel = min(max_parallel, 3)
            self.log_debug(f"Parallelismo limitato per percorso di rete: {max_parallel}")
        
        # Utilizziamo un set per tracciare i blocchi già processati
        processed_blocks = set()
        
        # Aggiungi contatori per controllo delle interruzioni e gestione memoria
        interrupt_check_counter = 0
        memory_check_counter = 0
        last_watchdog_update = time.time()
        
        # Tracciamento delle prestazioni
        last_performance_check = time.time()
        files_at_last_check = files_checked[0]
        
        while not block_queue.empty() and not self.stop_search:
            # Verifica timeout
            current_time = time.time()
            if timeout and current_time - start_time > timeout:
                self.progress_queue.put(("timeout", "Timeout raggiunto"))
                return
            
            # Aggiorna il watchdog più frequentemente
            if hasattr(self, 'last_progress_time') and current_time - last_watchdog_update > 10:
                self.last_progress_time = current_time
                last_watchdog_update = current_time
                
                # Tracciamento delle prestazioni
                if current_time - last_performance_check >= 30:  # Ogni 30 secondi
                    # Calcola velocità di elaborazione
                    elapsed = current_time - last_performance_check
                    files_processed = files_checked[0] - files_at_last_check
                    
                    if elapsed > 0:
                        speed = files_processed / elapsed
                        self.log_debug(f"Velocità di elaborazione: {speed:.1f} file/sec")
                        
                        # Adatta dinamicamente la dimensione del blocco
                        if self.block_size_auto_adjust.get():
                            if speed < 10:  # Molto lento
                                new_size = max(max_files_in_block // 2, 100)
                                if new_size != max_files_in_block:
                                    max_files_in_block = new_size
                                    self.log_debug(f"Riduzione dimensione blocco a {max_files_in_block} per bassa velocità")
                            elif speed > 1000:  # Molto veloce
                                new_size = min(max_files_in_block * 2, 5000)
                                if new_size != max_files_in_block:
                                    max_files_in_block = new_size
                                    self.log_debug(f"Aumento dimensione blocco a {max_files_in_block} per alta velocità")
                    
                    # Aggiorna i contatori per il prossimo controllo
                    last_performance_check = current_time
                    files_at_last_check = files_checked[0]
            
            try:
                # CORREZIONE: Estrai in modo sicuro dalla coda
                try:
                    # Controllo più robusto per la consistenza del formato della coda
                    queue_item = block_queue.get(block=False)
                    
                    if len(queue_item) >= 3:  # Nuovo formato con profondità
                        priority, current_block, current_depth = queue_item
                    else:  # Vecchio formato senza profondità
                        priority, current_block = queue_item
                        # Se non stiamo tracciando la profondità (max_depth = 0), imposta a 0
                        # altrimenti calcola la profondità basata sul percorso
                        current_depth = 0 if not using_limited_depth else current_block.count(os.path.sep) - path.count(os.path.sep)
                except (IndexError, queue.Empty):
                    # In caso di coda vuota, esci dal ciclo
                    break
                
                # CORREZIONE: Verifica la profondità massima in modo più chiaro
                # Se using_limited_depth è True (max_depth > 0) e la profondità corrente supera max_depth, salta
                if using_limited_depth and current_depth >= max_depth:
                    # AGGIUNTA: Log per debug quando una directory viene saltata per limiti di profondità
                    self.log_debug(f"Saltata directory {current_block} per limite di profondità (profondità: {current_depth}, max: {max_depth})")
                    continue
                
                # Salta blocchi già processati
                if current_block in processed_blocks:
                    continue
                    
                processed_blocks.add(current_block)
                
                # Aggiorna lo stato più frequentemente
                current_time = time.time()
                if current_time - last_update_time[0] >= 0.3:  # Ridotto da 0.5 a 0.3 secondi
                    elapsed_time = current_time - start_time
                    self.progress_queue.put(("status", 
                        f"Analisi blocco: {current_block} (Cartelle: {dirs_checked[0]}, File: {files_checked[0]}, Tempo: {int(elapsed_time)}s)"))
                    self.progress_queue.put(("progress", 
                        min(90, int((files_checked[0] / max(1, self.max_files_to_check.get())) * 100))))
                    last_update_time[0] = current_time
                
                # Implementa la verifica dei percorsi problematici prima di elaborare
                # Verifica se il percorso attuale è in una directory problematica o esclusa
                skip_block = False
                
                # Verifica più efficiente dei percorsi esclusi
                if hasattr(self, 'excluded_paths') and self.excluded_paths:
                    if any(current_block.lower().startswith(excluded.lower()) for excluded in self.excluded_paths):
                        self.log_debug(f"Salto blocco in percorso escluso: {current_block}")
                        skip_block = True
                
                # Verifica per directory problematiche note
                if not skip_block and hasattr(self, 'problematic_dirs'):
                    if any(problematic in current_block for problematic in self.problematic_dirs):
                        self.log_debug(f"Salto blocco in directory problematica: {current_block}")
                        skip_block = True
                
                if skip_block:
                    continue
                
                # Blocco attualmente in elaborazione
                dirs_checked[0] += 1
                
                # Usa try-except più granulare per gestire errori di accesso
                try:
                    items = os.listdir(current_block)
                except PermissionError:
                    if self.skip_permission_errors.get():
                        self.log_debug(f"Saltata directory con permesso negato: {current_block}")
                        continue
                    else:
                        # Gestione fallback per directory inaccessibili
                        dir_name = os.path.basename(current_block)
                        parent_dir = os.path.dirname(current_block)
                        is_user_folder = (parent_dir.lower() in ["c:/users", "c:\\users"] and 
                                        dir_name.lower() != getpass.getuser().lower())
                        if is_user_folder:
                            self.log_debug(f"Cartella di un altro utente inaccessibile: {current_block}")
                            self.progress_queue.put(("status", f"Saltata cartella utente protetta: {current_block}"))
                        else:
                            self.log_debug(f"Permesso negato per la directory {current_block}")
                            self.progress_queue.put(("status", f"Permesso negato: {current_block}"))
                        continue
                except Exception as e:
                    self.log_debug(f"Errore nell'accesso alla directory {current_block}: {str(e)}")
                    continue
                
                # Prima processa le sottocartelle (aggiungi nuovi blocchi)
                subfolders = []
                for item in items:
                    if self.stop_search:
                        return
                        
                    item_path = os.path.join(current_block, item)
                    
                    # Salta file/cartelle nascoste
                    try:
                        if self.ignore_hidden.get() and (item.startswith('.') or 
                            (os.name == 'nt' and os.path.exists(item_path) and 
                            os.stat(item_path).st_file_attributes & 2)):
                            continue
                    except Exception as e:
                        self.log_debug(f"Errore nel controllo hidden per {item_path}: {str(e)}")
                        continue
                        
                    # Gestione sottodirectory
                    try:
                        if os.path.isdir(item_path):
                            # Verifica se la cartella è già stata visitata
                            try:
                                real_path = os.path.realpath(item_path)
                                if real_path in visited_dirs:
                                    continue
                                visited_dirs.add(real_path)
                            except:
                                if item_path in visited_dirs:
                                    continue
                                visited_dirs.add(item_path)
                            
                            # Ottimizzazione verifica percorsi esclusi
                            # Verifica se il percorso deve essere escluso (più efficiente)
                            excluded = False
                            if hasattr(self, 'excluded_paths') and self.excluded_paths:
                                excluded = any(item_path.lower().startswith(excluded_path.lower()) 
                                            for excluded_path in self.excluded_paths)
                                
                            if excluded:
                                continue
                            
                            # CORREZIONE: Aggiungi sempre la sottocartella alla lista con profondità incrementata
                            # La verifica della profondità massima verrà fatta nel ciclo successivo
                            subfolders.append((item_path, current_depth + 1))
                            
                            # Verifica corrispondenza nome cartella
                            if self.search_folders.get():
                                # Usa match parola intera dove appropriato
                                matched = False
                                for keyword in keywords:
                                    if self.whole_word_search.get():
                                        if self.is_whole_word_match(keyword, item):
                                            matched = True
                                            break
                                    elif keyword.lower() in item.lower():
                                        matched = True
                                        break
                                        
                                if matched:
                                    folder_info = self.create_folder_info(item_path)
                                    self.search_results.append(folder_info)
                                    
                    except Exception as e:
                        self.log_debug(f"Errore nell'analisi della directory {item_path}: {str(e)}")
                
                # Ottimizza l'ordine di elaborazione delle sottocartelle
                if self.prioritize_user_folders.get():
                    # Estrai solo i percorsi per la funzione optimize_disk_search_order
                    subfolder_paths = [item[0] for item in subfolders]
                    optimized_paths = self.optimize_disk_search_order(path, subfolder_paths)
                    
                    # Ricostruisci la lista con le profondità
                    optimized_subfolders = []
                    path_to_depth = {item[0]: item[1] for item in subfolders}
                    for optimized_path in optimized_paths:
                        optimized_subfolders.append((optimized_path, path_to_depth.get(optimized_path, current_depth + 1)))
                    subfolders = optimized_subfolders
                
                # CORREZIONE: Aggiunta sottocartelle alla coda con priorità calcolata
                # Usa sempre il formato con profondità per maggiore coerenza
                for subfolder_info in subfolders:
                    subfolder, folder_depth = subfolder_info
                    
                    priority = self.calculate_block_priority(subfolder)
                    
                    # Aggiungi sempre la sottocartella alla coda, la verifica della profondità
                    # sarà fatta nel ciclo successivo
                    block_queue.put((priority, subfolder, folder_depth))
                    
                    # AGGIUNTA: Log per debug quando viene aggiunta una sottocartella
                    if folder_depth > current_depth:
                        self.log_debug(f"{subfolder} (profondità: {folder_depth})")
                
                # Ottimizza la gestione della memoria per blocchi grandi
                # Gestione della memoria più aggressiva
                memory_check_counter += 1
                if memory_check_counter >= 5:  # Controllo ogni 5 blocchi invece che per file
                    self.manage_memory()
                    memory_check_counter = 0
                
                # Processa i file nel blocco corrente con batch più piccoli
                file_batch = []
                for item in items:
                    if self.stop_search:
                        return
                        
                    item_path = os.path.join(current_block, item)
                    
                    # Salta i file nascosti e le directory (già processate)
                    try:
                        # Controllo file nascosto
                        if self.ignore_hidden.get() and (item.startswith('.') or 
                                    (os.name == 'nt' and os.path.exists(item_path) and 
                                        os.stat(item_path).st_file_attributes & 2)):
                            continue
                            
                        # Salta le directory (già processate)
                        if os.path.isdir(item_path):
                            continue
                    except Exception as e:
                        self.log_debug(f"Errore nel controllo del tipo per {item_path}: {str(e)}")
                        continue
                        
                    # Aggiungi file al batch se cerchiamo nei file
                    if self.search_files.get():
                        file_batch.append(item_path)
                        
                        # Quando il batch raggiunge dimensione massima, processalo
                        if len(file_batch) >= 100:  # Elabora a blocchi di 100 file
                            self.process_file_batch(file_batch, files_checked, last_update_time, 
                                                calculation_enabled, keywords, search_content, futures)
                            file_batch = []
                
                # Processa il batch finale di file
                if file_batch:
                    self.process_file_batch(file_batch, files_checked, last_update_time, 
                                        calculation_enabled, keywords, search_content, futures)
                    
            except queue.Empty:
                break

            # Controllo più frequente per interruzioni
            interrupt_check_counter += 1
            if interrupt_check_counter >= 3:  # Ridotto da 10 a 3 blocchi
                interrupt_check_counter = 0
                # Forza un aggiornamento dell'interfaccia più frequente
                try:
                    self.root.update_idletasks()
                except:
                    pass
    
    @error_handler
    # Funzione per calcolare il tempo rimanente stimato
    def calculate_remaining_time(self, files_processed, max_files, elapsed_time):
        """Calcola il tempo rimanente stimato in base al progresso attuale"""
        if files_processed == 0 or elapsed_time == 0:
            return "Calcolo..."
            
        try:
            # Calcola la velocità di elaborazione (files/secondo)
            files_per_second = files_processed / elapsed_time
            
            if files_per_second > 0:
                # Calcola il tempo rimanente stimato
                remaining_files = max_files - files_processed
                remaining_seconds = remaining_files / files_per_second
                
                # Formatta il risultato
                if remaining_seconds < 60:
                    return f"{int(remaining_seconds)} secondi"
                elif remaining_seconds < 3600:
                    minutes = int(remaining_seconds // 60)
                    seconds = int(remaining_seconds % 60)
                    return f"{minutes}m {seconds}s"
                else:
                    hours = int(remaining_seconds // 3600)
                    minutes = int((remaining_seconds % 3600) // 60)
                    return f"{hours}h {minutes}m"
            else:
                return "Calcolo..."
        except Exception as e:
            self.log_debug(f"Errore nel calcolo del tempo rimanente: {str(e)}")
            return "N/A"
        
    @error_handler
    def _search_thread(self, path, keywords, search_content):
        try:
            # Inizializza i contatori di file e directory esaminati e il tempo di inizio
            files_checked = [0]  # Uso una lista per poter modificare il valore nelle funzioni chiamate
            dirs_checked = [0]
            start_time = time.time()
            timeout = self.timeout_seconds.get() if self.timeout_enabled.get() else None
            
            # NUOVA RIGA: Verifica all'inizio se il calcolo della dimensione è abilitato
            calculation_enabled = self.dir_size_calculation.get() != "disabilitato"
            
            self.current_search_size = 0  # Dimensione totale dei file trovati
            last_size_update_time = time.time()  # Per aggiornamenti periodici

            # Determina se si tratta di una ricerca completa del sistema (C:/ o simile)
            is_system_search = path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]

            # Per ricerche di sistema, adatta automaticamente i parametri
            if is_system_search:
                # Informa l'utente che la ricerca potrebbe richiedere molto tempo
                self.progress_queue.put(("status", "Ricerca completa del sistema in corso - potrebbe richiedere molto tempo"))
                
                # Temporaneamente aumenta i limiti per la ricerca di sistema
                original_max_files = self.max_files_to_check.get()
                original_timeout = timeout
                
                # Imposta limiti più elevati per la ricerca completa
                self.max_files_to_check.set(5000000)  # 5 milioni di file
                
                # Disabilita temporaneamente il timeout o aumentalo significativamente
                if timeout and timeout < 3600:
                    timeout = 3600 * 8  # 8 ore
                
                # Avviso all'utente
                self.progress_queue.put(("status", "Ricerca completa avviata - parametri adattati per ricerca approfondita"))
            
            # Variabili per gestire l'aggiornamento ogni secondo
            last_update_time = [time.time()]
            
            # Crea un executor per processare i file in parallelo in modo sicuro
            try:
                max_workers = max(1, min(32, self.worker_threads.get()))  # Garantisce un valore valido
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
            except Exception as e:
                self.log_debug(f"Errore nella creazione del ThreadPoolExecutor: {str(e)}")
                # Fallback a un valore sicuro
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
            
            futures = []
            
            # Insieme per tenere traccia delle cartelle visitate (evita loop infiniti con symlink)
            visited_dirs = set()
            
            # Coda di priorità per i blocchi di ricerca
            block_queue = queue.PriorityQueue()

            # Aggiorna lo stato iniziale
            self.progress_queue.put(("status", f"Inizio ricerca a blocchi in: {path} (Profondità: {'illimitata' if self.max_depth == 0 else self.max_depth})"))
            
            # Assicurati che la lista di esclusioni sia inizializzata
            if not hasattr(self, 'excluded_paths'):
                self.excluded_paths = []
                
            # Avvia la ricerca a blocchi
            self.initialize_block_queue(path, block_queue, visited_dirs, files_checked, keywords, search_content, futures)
            self.process_blocks(block_queue, visited_dirs, start_time, timeout, is_system_search, 
                            files_checked, dirs_checked, last_update_time, path, keywords, search_content, futures)
            
            # Ripristina i parametri originali se erano stati modificati
            if is_system_search and self.max_depth == 0 and 'original_max_files' in locals():
                self.max_files_to_check.set(original_max_files)
            
            # Aggiorna lo stato finale di analisi
            self.progress_queue.put(("status", f"Elaborazione risultati... (analizzati {files_checked[0]} file in {dirs_checked[0]} cartelle)"))
            
            # Raccolta risultati dalle future con aggiornamento temporizzato
            completed = 0
            total_futures = len(futures)
            last_update_time[0] = time.time()
            
            # Gestione sicura delle future
            if self.search_executor and not self.search_executor._shutdown and futures:
                try:
                    for future in concurrent.futures.as_completed(futures):
                        if self.stop_search:
                            break
                            
                        try:
                            result = future.result()
                            if result:
                                self.search_results.append(result)
                            
                            completed += 1
                            
                            # Aggiorna il progresso e il tempo più frequentemente
                            current_time = time.time()
                            if completed % 20 == 0 or current_time - last_update_time[0] > 2:  # Aggiorna ogni 20 file o ogni 2 secondi
                                progress = 90 + min(10, int((completed / max(1, len(futures))) * 10))
                                self.progress_queue.put(("progress", progress))
                                
                                elapsed_time = current_time - start_time
                                self.progress_queue.put(("status", f"Elaborati {completed}/{len(futures)} file (tempo: {int(elapsed_time)}s)"))
                                
                                # Aggiorna anche il tempo totale durante la ricerca
                                if hasattr(self, 'search_start_time'):
                                    now = datetime.now()
                                    time_diff = now - self.search_start_time
                                    # Converti in minuti e secondi
                                    total_seconds = int(time_diff.total_seconds())
                                    minutes = total_seconds // 60
                                    seconds = total_seconds % 60
                                    
                                    if minutes > 0:
                                        total_time_str = f"{minutes}min {seconds}sec"
                                    else:
                                        total_time_str = f"{seconds}sec"
                                        
                                    self.progress_queue.put(("update_total_time", total_time_str))
                                
                                last_update_time[0] = current_time
                        except Exception as e:
                            self.log_debug(f"Errore nell'elaborazione di un risultato: {str(e)}")
                except Exception as e:
                    self.log_debug(f"Errore nella raccolta dei risultati: {str(e)}")
            
            # Completa la ricerca in modo sicuro
            try:
                if self.search_executor and not self.search_executor._shutdown:
                    self.search_executor.shutdown(wait=False)
                self.search_executor = None
            except Exception as e:
                self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                self.search_executor = None
            
            # Ripristina i parametri originali se erano stati modificati per la ricerca di sistema
            if hasattr(self, 'original_system_search_params'):
                try:
                    self.max_files_to_check.set(self.original_system_search_params["max_files"])
                    self.worker_threads.set(self.original_system_search_params["worker_threads"])
                    delattr(self, 'original_system_search_params')
                except:
                    pass

            # Riporta il risultato finale
            elapsed_time = time.time() - start_time
            self.progress_queue.put(("status", 
                f"Ricerca completata! Analizzati {files_checked[0]} file in {dirs_checked[0]} cartelle in {int(elapsed_time)} secondi."))
            
            # Ordina i risultati per tipo e nome
            self.search_results.sort(key=lambda x: (x[0], x[1]))
            
            self.log_debug(f"Ricerca completata. Trovati {len(self.search_results)} risultati")
            self.progress_queue.put(("complete", "Ricerca completata"))
            
        except Exception as e:
            error_msg = f"Si è verificato un errore durante la ricerca: {str(e)}\n{traceback.format_exc()}"
            self.log_debug(error_msg)
            self.progress_queue.put(("error", error_msg))

    @error_handler
    def is_whole_word_match(self, keyword, text):
        """Verifica se la keyword è presente nel testo come parola intera."""
        try:
            pattern = r'\b' + re.escape(keyword.lower()) + r'\b'
            return re.search(pattern, text.lower()) is not None
        except re.error as e:
            self.log_debug(f"Errore regex con il termine '{keyword}': {str(e)}")
            
            # Fallback manuale se la regex fallisce
            keyword_lower = keyword.lower()
            text_lower = text.lower()
            
            if keyword_lower not in text_lower:
                return False
                
            # Verifica manuale dei confini parola
            index = text_lower.find(keyword_lower)
            while index != -1:
                # Controlla il carattere prima della keyword
                has_char_before = index > 0 and text_lower[index-1].isalnum()
                
                # Controlla il carattere dopo la keyword
                end_pos = index + len(keyword_lower)
                has_char_after = end_pos < len(text_lower) and text_lower[end_pos].isalnum()
                
                # Se non ci sono caratteri alfanumerici prima e dopo, è una parola intera
                if not has_char_before and not has_char_after:
                    return True
                    
                # Cerca la prossima occorrenza
                index = text_lower.find(keyword_lower, index + 1)
            
            return False

    @error_handler  
    def search_current_user_only(self):
        """Imposta la ricerca solo nella cartella dell'utente corrente"""
        user_folder = os.path.join("C:/Users", getpass.getuser())
        if os.path.exists(user_folder):
            self.search_path.set(user_folder)
            
            # Mostra la conferma solo se l'utente non annulla il warning sulla ricerca nei contenuti
            if self.show_content_search_warning():
                messagebox.showinfo(
                    "Ricerca configurata", 
                    f"La ricerca è stata configurata per cercare solo nella tua cartella utente:\n{user_folder}\n\n"
                    "Questa modalità evita problemi di permesso con altre cartelle protette."
                )
                # Avvia direttamente la ricerca
                self.start_search()
        else:
            messagebox.showerror("Errore", "Impossibile trovare la tua cartella utente")

    @error_handler
    def create_file_info(self, file_path, from_attachment=False):
        """Crea le informazioni del file per la visualizzazione"""
        try:
            file_size = os.path.getsize(file_path)
            modified_time = datetime.fromtimestamp(os.path.getmtime(file_path))
            created_time = datetime.fromtimestamp(os.path.getctime(file_path))
            file_name = os.path.basename(file_path)
            file_extension = os.path.splitext(file_name)[1].lower()
            
            # Determina il tipo di file
            file_type = "File"  # Valore predefinito
            if os.path.isdir(file_path):
                file_type = "Directory"
            else:
                # Usa mimetypes per determinare il tipo
                import mimetypes
                mime_type, _ = mimetypes.guess_type(file_path)
                if mime_type:
                    file_type = mime_type.split('/')[0].capitalize()
                    if file_type == "Application":
                        if "pdf" in mime_type:
                            file_type = "PDF"
                        elif "word" in mime_type or file_extension == ".docx" or file_extension == ".doc":
                            file_type = "Word"
                        elif "excel" in mime_type or file_extension in [".xlsx", ".xls"]:
                            file_type = "Excel"
                        elif "powerpoint" in mime_type or file_extension in [".pptx", ".ppt"]:
                            file_type = "PowerPoint"
                        else:
                            file_type = "Documento"
                else:
                    # Fallback basato sull'estensione
                    if file_extension in ['.txt', '.md', '.rtf']:
                        file_type = "Testo"
                    elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                        file_type = "Immagine"
                    elif file_extension in ['.mp3', '.wav', '.ogg', '.flac']:
                        file_type = "Audio"
                    elif file_extension in ['.mp4', '.avi', '.mkv', '.mov']:
                        file_type = "Video"
                    elif file_extension in ['.exe', '.dll', '.bat']:
                        file_type = "Eseguibile"
                    elif file_extension == '.eml':
                        file_type = "Email"
            
            # Formatta dimensione file
            if file_size < 1024:
                size_str = f"{file_size} B"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.1f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.1f} MB"
            
            # Aggiungi il flag from_attachment alla tupla dei risultati
            return (
                file_type,
                file_name,
                size_str,
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                file_path,
                from_attachment  # Nuovo campo
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni del file {file_path}: {str(e)}")
            return (
                "File",
                os.path.basename(file_path),
                "N/A",
                "N/A",
                "N/A",
                file_path,
                False  # Default: non è un allegato
            )
    @error_handler
    def create_folder_info(self, folder_path):
        """Crea le informazioni della cartella per la visualizzazione"""
        try:
            modified_time = datetime.fromtimestamp(os.path.getmtime(folder_path))
            created_time = datetime.fromtimestamp(os.path.getctime(folder_path))
            folder_name = os.path.basename(folder_path)
            
            return (
                "Directory",
                folder_name,
                "",  # Le cartelle non hanno dimensione
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                folder_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni della cartella {folder_path}: {str(e)}")
            return (
                "Directory",
                os.path.basename(folder_path),
                "",
                "N/A",
                "N/A",
                folder_path
            )
    
    @error_handler
    def should_search_content(self, file_path):
        """Versione ottimizzata per determinare se analizzare il contenuto del file"""
        # Prima verifica le condizioni più veloci (principio fail-fast)
        if not self.search_content.get():
            return False
        
        # Controlla la dimensione del file usando il nuovo sistema di categorizzazione
        file_category = self.large_file_handler.get_file_size_category(file_path)
        
        # Verifica se il file è grande e se l'ottimizzazione per file grandi è disabilitata
        if file_category in ["large", "huge", "gigantic"] and not getattr(self, 'large_file_search_enabled', True):
            self.log_debug(f"File {file_category} ignorato (ottimizzazione disabilitata): {os.path.basename(file_path)}")
            return False
        
        # Aggiunge log dettagliato sulla dimensione del file per debug
        if file_category != "normal":
            self.log_debug(f"Elaborazione file di dimensione {file_category}: {os.path.basename(file_path)}")
        
        # Se il file è gigantesco, applica la logica speciale per file giganteschi
        if file_category == "gigantic":
            # Registra l'evento nei log
            self.log_debug(f"Applicazione strategia speciale per file gigantesco: {os.path.basename(file_path)}")
            
            # Imposta un flag per indicare alle altre funzioni che questo file richiede trattamento speciale
            file_size = os.path.getsize(file_path)
            file_extension = os.path.splitext(file_path)[1].lower()
            
            # Metodo 1: Verifica se il tipo di file è supportato per analisi parziale
            binary_types = ['.exe', '.dll', '.bin', '.iso', '.img', '.msi']
            archive_types = ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2']
            media_types = ['.mp4', '.avi', '.mkv', '.mov', '.mpg', '.wmv', '.mp3', '.wav', '.flac']
            
            # Per file binari molto grandi, probabilmente contengono dati non testuali
            if file_extension in binary_types and file_size > 5 * 1024 * 1024 * 1024:  # >5GB
                self.log_debug(f"File binario gigantesco, contenuto probabilmente non testuale: {os.path.basename(file_path)}")
                return False
                
            # Per archivi enormi, è generalmente meglio estrarre prima e cercare nei file estratti
            if file_extension in archive_types and file_size > 4 * 1024 * 1024 * 1024:  # >4GB
                # Qui potresti voler mostrare un messaggio all'utente suggerendo di estrarre prima l'archivio
                self.log_debug(f"Archivio gigantesco, si consiglia di estrarre e cercare nei file estratti: {os.path.basename(file_path)}")
                # Si potrebbe aggiungere un popup per l'utente qui
                return False
            
            # Per file multimediali enormi, la ricerca di testo è raramente utile
            if file_extension in media_types and file_size > 3 * 1024 * 1024 * 1024:  # >3GB
                self.log_debug(f"File multimediale gigantesco, contenuto probabilmente non testuale: {os.path.basename(file_path)}")
                return False
            
            # Metodo 2: Per database e file di log giganteschi, usa metodi speciali
            db_types = ['.db', '.sqlite', '.mdb', '.accdb', '.sql']
            log_types = ['.log', '.csv', '.tsv', '.txt']
            
            if file_extension in db_types or file_extension in log_types:
                # Imposta un attributo temporaneo per indicare che questo file necessita di elaborazione speciale
                # Questo verrà controllato dalla funzione di elaborazione
                self._mark_file_for_partial_analysis(file_path)
                self.log_debug(f"File dati gigantesco, verrà analizzato con metodi speciali: {os.path.basename(file_path)}")
                return True  # Continua con l'elaborazione, ma sarà gestito diversamente
            
            # Metodo 3: Per altri tipi di file giganteschi, offri all'utente la possibilità di scegliere
            # Questa è un'opzione più avanzata che potrebbe richiedere una UI specifica
            if file_size > 2.5 * 1024 * 1024 * 1024:  # >2.5GB
                # Implementa una logica di conferma tramite una variabile globale temporanea
                if not hasattr(self, '_gigantic_files_confirmed') or file_path not in self._gigantic_files_confirmed:
                    self.log_debug(f"File gigantesco richiede conferma per l'elaborazione: {os.path.basename(file_path)}")
                    # Restituisci False per ora, ma imposta un flag per ulteriore elaborazione
                    self._queue_gigantic_file_for_confirmation(file_path)
                    return False
        
        # Il resto del codice rimane invariato
        ext = os.path.splitext(file_path)[1].lower()
        
        # Attiva sempre la ricerca nei file Office, indipendentemente dal livello di ricerca
        if ext in ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']:
            self.log_debug(f"Analisi contenuto attivata per file Office: {os.path.basename(file_path)} ({ext})")
            return True
        
        # Seleziona il livello di ricerca attuale
        search_level = self.search_depth.get()
        
        # Ottieni le estensioni personalizzate dell'utente
        custom_extensions = self.get_extension_settings(search_level)
        
        # PRIORITÀ #1: Se l'estensione è stata aggiunta manualmente, cerca sempre il contenuto
        if ext in custom_extensions:
            self.log_debug(f"Ricerca contenuto in file con estensione personalizzata: {ext}")
            return True
        
        # PRIORITÀ #2: In modalità profonda senza estensioni personalizzate, cerca tutto
        if search_level == "profonda" and not custom_extensions:
            return True
        
        # Liste predefinite nel codice per ciascun livello
        base_extensions = ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.log', 
                        '.docx', '.doc', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls', '.rtf', '.odt', '.ods', '.odp',
                        '.csv','.eml', '.msg', '.emlx']

        advanced_extensions = base_extensions + ['.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', 
                                            '.vbs', '.js', '.config', '.ini', '.json', '.reg']
                                            
        deep_extensions = advanced_extensions + ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.mp3', '.mp4', 
                                            '.avi', '.mov', '.mkv', '.wav', '.flac', '.zip', '.rar', 
                                            '.7z', '.tar', '.gz', '.iso', '.psd', '.ai', '.svg']
        
        # Verifica il livello predefinito se non è nelle estensioni personalizzate
        if search_level == "base" and ext in base_extensions:
            return True
        elif search_level == "avanzata" and ext in advanced_extensions:
            return True
        elif search_level == "profonda" and ext in deep_extensions:
            return True
        
        return False

    @error_handler
    def _mark_file_for_partial_analysis(self, file_path):
        """Marca un file per l'analisi parziale invece che completa"""
        if not hasattr(self, '_partial_analysis_files'):
            self._partial_analysis_files = set()
        self._partial_analysis_files.add(file_path)

    @error_handler
    def _queue_gigantic_file_for_confirmation(self, file_path):
        """Accoda un file gigantesco per la conferma dell'utente"""
        if not hasattr(self, '_gigantic_files_queue'):
            self._gigantic_files_queue = []
        if not hasattr(self, '_gigantic_files_confirmed'):
            self._gigantic_files_confirmed = set()
        
        if file_path not in self._gigantic_files_queue and file_path not in self._gigantic_files_confirmed:
            self._gigantic_files_queue.append(file_path)
            
            # Schedula un prompt per l'utente se non già in corso
            if not hasattr(self, '_showing_gigantic_confirmation') or not self._showing_gigantic_confirmation:
                self._showing_gigantic_confirmation = True
                self.root.after(100, self._process_gigantic_file_queue)

    @error_handler
    def _process_gigantic_file_queue(self):
        """Processa la coda di file giganteschi richiedendo conferma all'utente"""
        if not hasattr(self, '_gigantic_files_queue') or not self._gigantic_files_queue:
            self._showing_gigantic_confirmation = False
            return
        
        file_path = self._gigantic_files_queue[0]
        file_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        formatted_size = self._format_size(file_size)
        
        # Calcola il tempo stimato di elaborazione (approssimativo)
        estimated_time = self._estimate_processing_time(file_size)
        
        # Mostra un dialog all'utente
        message = (f"Il file '{file_name}' è molto grande ({formatted_size}).\n\n"
                f"L'elaborazione potrebbe richiedere circa {estimated_time}.\n"
                "Vuoi procedere con l'analisi?")
        
        result = messagebox.askyesno("File gigantesco rilevato", message)
        
        if result:  # L'utente ha confermato
            self._gigantic_files_confirmed.add(file_path)
            # Riesamina il file nella ricerca
            if hasattr(self, 'redo_search_for_file'):
                self.redo_search_for_file(file_path)
        
        # Rimuovi il file dalla coda
        self._gigantic_files_queue.pop(0)
        
        # Se ci sono altri file nella coda, continua a processarli
        if self._gigantic_files_queue:
            self.root.after(100, self._process_gigantic_file_queue)
        else:
            self._showing_gigantic_confirmation = False

    @error_handler
    def _estimate_processing_time(self, file_size):
        """Stima il tempo di elaborazione per un file di grandi dimensioni"""
        # Valori empirici basati su test (questi andrebbero regolati in base alle prestazioni reali)
        bytes_per_second = 25 * 1024 * 1024  # ~25 MB/s per file di testo
        
        # Adatta la velocità in base alla dimensione (file più grandi sono più lenti da processare)
        if file_size > 5 * 1024 * 1024 * 1024:  # >5GB
            bytes_per_second = 15 * 1024 * 1024  # ~15 MB/s
        
        seconds = file_size / bytes_per_second
        
        # Formatta il tempo stimato
        if seconds < 60:
            return f"{int(seconds)} secondi"
        elif seconds < 3600:
            return f"{int(seconds / 60)} minuti"
        else:
            hours = int(seconds / 3600)
            minutes = int((seconds % 3600) / 60)
            return f"{hours} ore e {minutes} minuti"

    @error_handler
    def redo_search_for_file(self, file_path):
        """Riesamina un file specifico nella ricerca corrente"""
        if not hasattr(self, 'current_search_keywords') or not self.current_search_keywords:
            return
        
        # Esegui la ricerca solo su questo file specifico
        try:
            result = self.process_file(file_path, self.current_search_keywords, search_content=True)
            if result:
                self.search_results.append(result)
                self.update_results_list()
        except Exception as e:
            self.log_error(f"Errore durante l'elaborazione del file {file_path}", e)

    @error_handler
    def should_skip_file(self, file_path):
        """Verifica se un file deve essere saltato durante l'analisi del contenuto"""
        ext = os.path.splitext(file_path)[1].lower()
        skip_type = "File di sistema" if ext in self.system_file_extensions else "File"
        skip_filename = os.path.basename(file_path)
        
        # Ottieni le estensioni selezionate per il livello di ricerca corrente
        search_level = self.search_depth.get()
        selected_extensions = self.get_extension_settings(search_level)
        
        # CORREZIONE: Se l'estensione non è tra quelle selezionate dall'utente, salta il file
        # Ma procedi solo se l'estensione non è vuota (per includere file senza estensione)
        if ext and ext not in selected_extensions:
            self.log_debug(f"File saltato perché estensione deselezionata: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Estensione deselezionata")
            return True
        
        # Salta i file di Rights Management Services
        if "Rights Management Services" in file_path or "IRMProtectors" in file_path:
            self.log_debug(f"Saltato file protetto: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Rights Management Services")
            return True
                    
        # Salta file con estensioni problematiche
        problematic_extensions = [".msoprotector.doc", ".msoprotector.ppt", ".msoprotector.xls"]
        if any(ext in file_path for ext in problematic_extensions):
            self.log_debug(f"Saltato file con formato problematico: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Formato problematico")
            return True

        # Salta file di sistema solo se non sono nelle estensioni personalizzate
        if self.exclude_system_files.get() and ext in self.system_file_extensions:
            # Non saltare script files in ricerca avanzata/profonda
            script_files = ['.bat', '.cmd', '.ps1', '.vbs']
            if ext in script_files and search_level in ["avanzata", "profonda"]:
                return False
            
            self.log_debug(f"File di sistema escluso: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "File di sistema")
            return True
                    
        return False
    
    @error_handler
    def log_skipped_file(self, filepath, skiptype, filename, skipreason):
        """Registra i file saltati in un file di log"""
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_entry = f"{timestamp} - {skiptype} - {filename} - {filepath} - {skipreason}\n"
            
            with open(self.skipped_files_log_path, 'a', encoding='utf-8') as log_file:
                log_file.write(log_entry)
        except Exception as e:
            self.log_debug(f"Errore durante la scrittura del log dei file saltati: {str(e)}")

    @error_handler
    def export_skipped_files_log(self):
        """Esporta il log dei file saltati in un formato CSV"""
        try:
            
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da esportare.")
                return
                
            # Verifica che il file non sia vuoto
            if os.path.getsize(self.skipped_files_log_path) == 0:
                messagebox.showinfo("Informazione", "Il file di log esiste ma è vuoto. Non ci sono file esclusi da esportare.")
                return

            # Chiedi all'utente dove salvare il file esportato
            export_path = filedialog.asksaveasfilename(
                defaultextension=".csv",
                initialfile="file_esclusi_export.csv",
                filetypes=[("CSV files", "*.csv"), ("Text files", "*.txt"), ("All files", "*.*")],
                title="Salva il log dei file esclusi")
            
            if not export_path:  # L'utente ha annullato
                return
                
            # Leggi il file di log originale
            with open(self.skipped_files_log_path, 'r', encoding='utf-8') as log_file:
                log_content = log_file.readlines()
                
            # Crea il file CSV
            with open(export_path, 'w', newline='', encoding='utf-8') as csv_file:
                csv_writer = csv.writer(csv_file, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
                
                # Intestazione
                csv_writer.writerow(["Data e Ora", "Tipo", "Nome File", "Percorso Completo", "Motivo Esclusione"])
                
                # Analizza e scrive ogni riga del log
                for log_line in log_content:
                    try:
                        # Formato tipico: 2023-01-01 12:34:56 - File di sistema - nomefile.exe - C:/percorso/nomefile.exe - File di sistema
                        parts = log_line.strip().split(" - ", 4)
                        if len(parts) >= 5:
                            csv_writer.writerow(parts)
                    except Exception as e:
                        self.log_debug(f"Errore nell'elaborazione della riga di log: {str(e)}")
                        
                # Aggiunge statistiche alla fine
                csv_writer.writerow([])
                csv_writer.writerow([f"Totale file esclusi: {len(log_content)}"])
                # Utilizzo della data e utente forniti
                csv_writer.writerow([f"Esportazione eseguita il: 2025-04-19 19:25:10"])
                csv_writer.writerow([f"Utente: Nino19980"])
            
            # Aggiungi un link per aprire il file esportato
            open_export = messagebox.askyesno(
                "Esportazione completata", 
                f"Esportazione completata con successo!\n\nFile salvato in:\n{export_path}\n\nVuoi aprire il file?")
            
            if open_export:
                try:
                    # Solo Windows, rimossa la parte Linux/macOS
                    # Usare os.startfile è consigliato su Windows e non apre finestre CMD
                    os.startfile(export_path)
                except Exception as e:
                    self.log_debug(f"Errore nell'apertura del file esportato: {str(e)}")
                    messagebox.showinfo("Informazione", f"Il file è stato salvato in:\n{export_path}")
                    
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'esportazione: {str(e)}")
            self.log_debug(f"Errore nell'esportazione del log: {str(e)}")

    @error_handler
    def view_skipped_files_log(self):
        """Visualizza il log dei file saltati in una finestra separata"""
        try:
            if not hasattr(self, 'search_start_time') or self.search_start_time is None:
                messagebox.showinfo("Informazione", "Non è stata ancora effettuata una ricerca per visualizzare i file esclusi.")
                return
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da visualizzare.")
                return
                
            # Verifica che il file non sia vuoto
            if os.path.getsize(self.skipped_files_log_path) == 0:
                messagebox.showinfo("Informazione", "Il file di log esiste ma è vuoto. Non ci sono file esclusi da visualizzare.")
                return
                
            # Crea una nuova finestra
            log_window = ttk.Toplevel(self.root)
            log_window.title("Log dei file esclusi")
            log_window.geometry("1000x800")
            log_window.transient(self.root)
            
            # Frame principale
            main_frame = ttk.Frame(log_window, padding=10)
            main_frame.pack(fill=BOTH, expand=YES)
            
            # Intestazione
            ttk.Label(main_frame, text="File esclusi durante la ricerca", 
                    font=("", 12, "bold")).pack(anchor=W, pady=(0, 10))
            
            # Area di testo per visualizzare i log
            text_frame = ttk.Frame(main_frame)
            text_frame.pack(fill=BOTH, expand=YES)
            
            scrollbar_y = ttk.Scrollbar(text_frame)
            scrollbar_y.pack(side=RIGHT, fill=Y)
            
            scrollbar_x = ttk.Scrollbar(text_frame, orient="horizontal")
            scrollbar_x.pack(side=BOTTOM, fill=X)
            
            log_text = tk.Text(text_frame, wrap="none", 
                            xscrollcommand=scrollbar_x.set,
                            yscrollcommand=scrollbar_y.set)
            log_text.pack(fill=BOTH, expand=YES)
            
            scrollbar_y.config(command=log_text.yview)
            scrollbar_x.config(command=log_text.xview)
            
            # Leggi e inserisci il contenuto del log
            with open(self.skipped_files_log_path, 'r', encoding='utf-8') as log_file:
                log_content = log_file.read()
                log_text.insert("1.0", log_content)
                
            # Rendi il testo di sola lettura
            log_text.config(state="disabled")
            
            # Pulsanti
            btn_frame = ttk.Frame(main_frame)
            btn_frame.pack(fill=X, pady=(10, 0))
            
            ttk.Button(btn_frame, text="Esporta CSV", 
                    command=self.export_skipped_files_log).pack(side=LEFT)
            
            ttk.Button(btn_frame, text="Svuota Log", 
                    command=self.clear_skipped_files_log).pack(side=LEFT, padx=5)
            
            ttk.Button(btn_frame, text="Chiudi", 
                    command=log_window.destroy).pack(side=RIGHT)
                    
            # Centra la finestra
            log_window.update_idletasks()
            width = log_window.winfo_width()
            height = log_window.winfo_height()
            x = (log_window.winfo_screenwidth() // 2) - (width // 2)
            y = (log_window.winfo_screenheight() // 2) - (height // 2)
            log_window.geometry(f"{width}x{height}+{x}+{y}")
            
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante la visualizzazione del log: {str(e)}")
            self.log_debug(f"Errore nella visualizzazione del log: {str(e)}")

    @error_handler
    def clear_skipped_files_log(self):
        """Svuota il log dei file esclusi"""
        try:
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da cancellare.")
                return
                
            # Svuota il file di log
            open(self.skipped_files_log_path, 'w', encoding='utf-8').close()
            
            # Reinizializza la lista dei file saltati
            self.skipped_files = []
            
            # Fornisci feedback all'utente
            messagebox.showinfo("Operazione completata", "Il log dei file esclusi è stato svuotato con successo.")
            
            # Se la finestra di log è attualmente aperta, aggiornala per mostrare il log vuoto
            for widget in self.root.winfo_children():
                if isinstance(widget, tk.Toplevel) and widget.winfo_exists():
                    # Cerca il widget Text nel toplevel
                    for child in widget.winfo_children():
                        if isinstance(child, ttk.Frame):  # main_frame
                            for frame in child.winfo_children():
                                if isinstance(frame, ttk.Frame):  # text_frame
                                    for text_widget in frame.winfo_children():
                                        if isinstance(text_widget, tk.Text):
                                            text_widget.config(state="normal")
                                            text_widget.delete("1.0", tk.END)
                                            text_widget.config(state="disabled")
                                            return
            
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante la cancellazione del log: {str(e)}")
            self.log_debug(f"Errore nella cancellazione del log: {str(e)}")

    @error_handler
    def get_file_content(self, file_path):
        """Estrae il contenuto testuale dai vari formati di file - versione completa"""
        try:
            # Inizializza le variabili all'inizio per evitare problemi di scope
            content = ""
            text_content = ""
            result = ""
            extension = os.path.splitext(file_path.lower())[1]
        
            # Verifica se il file è un archivio compresso
            compressed_extensions = ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.iso', '.tgz', '.xz', '.cab', '.jar']
            if extension in compressed_extensions:
                # Utilizzo la funzione specializzata per gli archivi compressi
                return self.extract_archive_content(file_path)
                
            # Controlli preliminari
            if self.should_skip_file(file_path):
                return ""
                    
            ext = os.path.splitext(file_path)[1].lower()
            
            # Log per debug
            self.log_debug(f"Tentativo estrazione contenuto da: {os.path.basename(file_path)} ({os.path.splitext(file_path)[1]})")
            
            # Controllo dimensione file
            try:
                file_size = os.path.getsize(file_path)
                if file_size > self.max_file_size_mb.get() * 1024 * 1024:
                    self.log_debug(f"File {file_path} troppo grande per l'analisi del contenuto")
                    return ""
            except Exception as e:
                self.log_debug(f"Errore nel controllo dimensione del file {file_path}: {str(e)}")
                return ""
            
            if self.stop_search:
                return ""
            # Word DOCX
            if ext == '.docx':
                try:
                    import docx
                    self.log_debug(f"Processando file DOCX: {file_path}")
                    doc = docx.Document(file_path)
                    content = []
                    for paragraph in doc.paragraphs:
                        if paragraph.text.strip():
                            content.append(paragraph.text)
                    
                    # Estrai anche il testo dalle tabelle
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                content.append(" | ".join(row_text))
                    
                    result = '\n'.join(content)
                    self.log_debug(f"Estratti {len(result)} caratteri da DOCX")
                    return result
                except ImportError:
                    self.log_debug("Libreria python-docx non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file DOCX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Word DOC (vecchio formato)
            elif ext == '.doc':
                try:
                    # Prova prima con pywin32 (solo Windows)
                    if os.name == 'nt':
                        self.log_debug("Tentativo di estrazione da DOC con win32com...")
                        try:
                            import win32com.client
                            import pythoncom
                            
                            # Inizializzazione necessaria per i thread
                            pythoncom.CoInitialize()
                            
                            try:
                                word = win32com.client.Dispatch("Word.Application")
                                word.Visible = False
                                word.DisplayAlerts = False
                                
                                # Apri il documento in modalità sola lettura
                                doc = word.Documents.Open(os.path.abspath(file_path), ReadOnly=True)
                                text = doc.Content.Text
                                doc.Close(SaveChanges=False)
                                word.Quit()
                                
                                pythoncom.CoUninitialize()
                                
                                if text:
                                    self.log_debug(f"Estratti {len(text)} caratteri da DOC")
                                    return text
                                else:
                                    self.log_debug("Nessun testo estratto dal file DOC")
                                    return ""
                            except Exception as e:
                                self.log_debug(f"Errore nell'apertura del DOC con win32com: {str(e)}")
                                # Cleanup in caso di errore
                                try:
                                    if 'doc' in locals() and doc:
                                        doc.Close(SaveChanges=False)
                                    if 'word' in locals() and word:
                                        word.Quit()
                                except:
                                    pass
                                
                                pythoncom.CoUninitialize()
                                return ""
                        except ImportError:
                            self.log_debug("win32com non disponibile per i file DOC")
                            return ""
                    else:
                        self.log_debug("Estrazione da DOC non supportata su questa piattaforma")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore generale nell'elaborazione del file DOC {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Excel XLSX
            elif ext == '.xlsx':
                try:
                    import openpyxl
                    self.log_debug(f"Processando file XLSX: {file_path}")
                    
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    texts = []
                    
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        sheet_texts = []
                        
                        # Prima verifica il range utilizzato
                        max_row = min(sheet.max_row, 500) if sheet.max_row else 0
                        max_col = min(sheet.max_column, 50) if sheet.max_column else 0
                        
                        if max_row > 0 and max_col > 0:
                            # Utilizza openpyxl 2.6+ con la sintassi più efficiente
                            try:
                                # Estrazione di righe specifiche
                                for row_idx in range(1, max_row + 1):
                                    row_values = []
                                    for col_idx in range(1, max_col + 1):
                                        cell = sheet.cell(row=row_idx, column=col_idx)
                                        if cell.value:
                                            row_values.append(str(cell.value))
                                    
                                    if row_values:
                                        sheet_texts.append(" ".join(row_values))
                            except Exception as e:
                                self.log_debug(f"Errore nell'iterazione del foglio {sheet_name}: {str(e)}")
                        
                        if sheet_texts:
                            texts.append(f"--- Foglio: {sheet_name} ---")
                            texts.append("\n".join(sheet_texts))
                    
                    result = "\n".join(texts)
                    self.log_debug(f"Estratti {len(result)} caratteri da XLSX")
                    return result
                    
                except ImportError:
                    self.log_debug("Libreria openpyxl non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file XLSX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Excel XLS (vecchio formato)
            elif ext == '.xls':
                try:
                    # Prima prova con xlrd
                    try:
                        import xlrd
                        self.log_debug(f"Processando file XLS con xlrd: {file_path}")
                        
                        book = xlrd.open_workbook(file_path, on_demand=True)
                        texts = []
                        
                        for sheet_idx in range(book.nsheets):
                            sheet = book.sheet_by_index(sheet_idx)
                            sheet_texts = []
                            
                            # Limite a 500 righe per prestazioni
                            for row_idx in range(min(sheet.nrows, 500)):
                                row_values = sheet.row_values(row_idx)
                                row_texts = [str(value) for value in row_values if value]
                                if row_texts:
                                    sheet_texts.append(" ".join(row_texts))
                            
                            if sheet_texts:
                                texts.append(f"--- Foglio: {sheet.name} ---")
                                texts.append("\n".join(sheet_texts))
                        
                        book.release_resources()
                        result = "\n".join(texts)
                        self.log_debug(f"Estratti {len(result)} caratteri da XLS")
                        return result
                        
                    except ImportError:
                        # Fallback a pywin32 se disponibile e su Windows
                        if os.name == 'nt':
                            try:
                                import win32com.client
                                import pythoncom
                                
                                # Inizializzazione necessaria per i thread
                                pythoncom.CoInitialize()
                                
                                self.log_debug("Processando file XLS con win32com")
                                
                                try:
                                    excel = win32com.client.Dispatch("Excel.Application")
                                    excel.Visible = False
                                    excel.DisplayAlerts = False
                                    
                                    workbook = excel.Workbooks.Open(os.path.abspath(file_path), ReadOnly=True)
                                    texts = []
                                    
                                    for i in range(1, workbook.Sheets.Count + 1):
                                        sheet = workbook.Sheets(i)
                                        texts.append(f"--- Foglio: {sheet.Name} ---")
                                        
                                        # Verifica se c'è un range utilizzato
                                        if sheet.UsedRange and sheet.UsedRange.Cells.Count > 0:
                                            used_range = sheet.UsedRange
                                            row_count = min(used_range.Rows.Count, 500)
                                            col_count = min(used_range.Columns.Count, 50)
                                            
                                            for row_idx in range(1, row_count + 1):
                                                row_texts = []
                                                for col_idx in range(1, col_count + 1):
                                                    cell_value = used_range.Cells(row_idx, col_idx).Value
                                                    if cell_value:
                                                        row_texts.append(str(cell_value))
                                                
                                                if row_texts:
                                                    texts.append(" ".join(row_texts))
                                    
                                    workbook.Close(False)
                                    excel.Quit()
                                    
                                    pythoncom.CoUninitialize()
                                    
                                    result = "\n".join(texts)
                                    self.log_debug(f"Estratti {len(result)} caratteri da XLS con win32com")
                                    return result
                                    
                                except Exception as e:
                                    self.log_debug(f"Errore nell'apertura dell'XLS con win32com: {str(e)}")
                                    
                                    # Cleanup in caso di errore
                                    try:
                                        if 'workbook' in locals() and workbook:
                                            workbook.Close(False)
                                        if 'excel' in locals() and excel:
                                            excel.Quit()
                                    except:
                                        pass
                                    
                                    pythoncom.CoUninitialize()
                                    return ""
                                    
                            except ImportError:
                                self.log_debug("Nessuna libreria disponibile per i file XLS")
                                return ""
                        else:
                            self.log_debug("Nessuna libreria disponibile per i file XLS su questa piattaforma")
                            return ""
                            
                except Exception as e:
                    self.log_debug(f"Errore generale nell'elaborazione del file XLS {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # PowerPoint PPTX
            elif ext == '.pptx':
                try:
                    import pptx
                    self.log_debug(f"Processando file PPTX: {file_path}")
                    
                    presentation = pptx.Presentation(file_path)
                    texts = []
                    
                    for i, slide in enumerate(presentation.slides):
                        slide_text = []
                        texts.append(f"--- Diapositiva {i+1} ---")
                        
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text.append(shape.text)
                        
                        if slide_text:
                            texts.append("\n".join(slide_text))
                    
                    result = "\n".join(texts)
                    self.log_debug(f"Estratti {len(result)} caratteri da PPTX")
                    return result
                    
                except ImportError:
                    self.log_debug("Libreria python-pptx non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file PPTX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # PowerPoint PPT (vecchio formato)
            elif ext == '.ppt':
                # Su Windows, prova con pywin32
                if os.name == 'nt':
                    try:
                        import win32com.client
                        import pythoncom
                        
                        # Inizializzazione necessaria per i thread
                        pythoncom.CoInitialize()
                        
                        self.log_debug(f"Processando file PPT: {file_path}")
                        
                        try:
                            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                            powerpoint.Visible = False
                            
                            presentation = powerpoint.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
                            texts = []
                            
                            for slide_idx in range(1, presentation.Slides.Count + 1):
                                slide = presentation.Slides.Item(slide_idx)
                                texts.append(f"--- Diapositiva {slide_idx} ---")
                                slide_text = []
                                
                                for shape_idx in range(1, slide.Shapes.Count + 1):
                                    shape = slide.Shapes.Item(shape_idx)
                                    if shape.HasTextFrame:
                                        if shape.TextFrame.HasText:
                                            slide_text.append(shape.TextFrame.TextRange.Text)
                                
                                if slide_text:
                                    texts.append("\n".join(slide_text))
                            
                            presentation.Close()
                            powerpoint.Quit()
                            
                            pythoncom.CoUninitialize()
                            
                            result = "\n".join(texts)
                            self.log_debug(f"Estratti {len(result)} caratteri da PPT")
                            return result
                            
                        except Exception as e:
                            self.log_debug(f"Errore nell'apertura del PPT con win32com: {str(e)}")
                            
                            # Cleanup in caso di errore
                            try:
                                if 'presentation' in locals() and presentation:
                                    presentation.Close()
                                if 'powerpoint' in locals() and powerpoint:
                                    powerpoint.Quit()
                            except:
                                pass
                            
                            pythoncom.CoUninitialize()
                            return ""
                            
                    except ImportError:
                        self.log_debug("win32com non disponibile per i file PPT")
                        return ""
                else:
                    self.log_debug("Estrazione di testo dai file PPT non supportata su questa piattaforma")
                    return ""
            
            if self.stop_search:
                return ""
            # Rich Text Format (RTF) - NUOVO
            elif ext == '.rtf':
                try:
                    from striprtf.striprtf import rtf_to_text
                    self.log_debug(f"Processando file RTF: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        rtf_content = f.read()
                        content = rtf_to_text(rtf_content)
                        self.log_debug(f"Estratti {len(content)} caratteri da RTF")
                        return content
                except ImportError:
                    self.log_debug("Libreria striprtf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file RTF {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""        
            # OpenDocument Text (ODT) - NUOVO
            elif ext == '.odt':
                try:
                    from odf import opendocument, text
                    self.log_debug(f"Processando file ODT: {file_path}")
                    
                    doc = opendocument.load(file_path)
                    paragraphs = []
                    
                    # Estrai tutto il testo dai paragrafi
                    for element in doc.getElementsByType(text.P):
                        paragraphs.append(element.firstChild.data if element.firstChild else "")
                        
                    content = "\n".join(paragraphs)
                    self.log_debug(f"Estratti {len(content)} caratteri da ODT")
                    return content
                except ImportError:
                    self.log_debug("Libreria odf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODT {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""       
            # OpenDocument Spreadsheet (ODS) - NUOVO
            elif ext == '.ods':
                try:
                    from odf import opendocument, table, text
                    self.log_debug(f"Processando file ODS: {file_path}")
                    
                    doc = opendocument.load(file_path)
                    sheets = []
                    
                    # Estrai tutte le tabelle e celle
                    for sheet in doc.getElementsByType(table.Table):
                        sheet_rows = []
                        sheet_name = sheet.getAttribute('table:name')
                        sheets.append(f"--- Foglio: {sheet_name} ---")
                        
                        for row in sheet.getElementsByType(table.TableRow):
                            row_cells = []
                            for cell in row.getElementsByType(table.TableCell):
                                cell_text = ""
                                for p in cell.getElementsByType(text.P):
                                    cell_text += p.firstChild.data if p.firstChild else ""
                                row_cells.append(cell_text)
                            if row_cells:
                                sheet_rows.append(" | ".join(row_cells))
                                
                        sheets.append("\n".join(sheet_rows))
                        
                    content = "\n".join(sheets)
                    self.log_debug(f"Estratti {len(content)} caratteri da ODS")
                    return content
                except ImportError:
                    self.log_debug("Libreria odf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODS {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # OpenDocument Presentation (ODP) - NUOVO
            elif ext == '.odp':
                try:
                    from odf import opendocument, draw, text
                    self.log_debug(f"Processando file ODP: {file_path}")
                    
                    doc = opendocument.load(file_path)
                    slides = []
                    
                    # Estrai tutte le diapositive
                    for page in doc.getElementsByType(draw.Page):
                        slide_text = []
                        slide_name = page.getAttribute('draw:name')
                        slides.append(f"--- Diapositiva: {slide_name} ---")
                        
                        # Estrai il testo dai frame e forme
                        for element in page.getElementsByType(draw.Frame):
                            for p in element.getElementsByType(text.P):
                                if p.firstChild:
                                    slide_text.append(p.firstChild.data)
                                    
                        slides.append("\n".join(slide_text))
                        
                    content = "\n".join(slides)
                    self.log_debug(f"Estratti {len(content)} caratteri da ODP")
                    return content
                except ImportError:
                    self.log_debug("Libreria odf non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODP {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""        
            # EPUB - NUOVO
            elif ext == '.epub':
                try:
                    import ebooklib
                    from ebooklib import epub
                    from bs4 import BeautifulSoup
                    self.log_debug(f"Processando file EPUB: {file_path}")
                    
                    # Funzione per estrarre testo dall'HTML
                    def chapter_to_text(content):
                        soup = BeautifulSoup(content, 'html.parser')
                        return soup.get_text()
                    
                    book = epub.read_epub(file_path)
                    chapters = []
                    
                    # Estrai metadati
                    title = book.get_metadata('DC', 'title')
                    if title:
                        chapters.append(f"Titolo: {title[0][0]}")
                    
                    authors = book.get_metadata('DC', 'creator')
                    if authors:
                        chapters.append(f"Autore: {authors[0][0]}")
                        
                    # Estrai contenuto
                    for item in book.get_items():
                        if item.get_type() == ebooklib.ITEM_DOCUMENT:
                            chapters.append(chapter_to_text(item.get_content()))
                            
                    content = "\n".join(chapters)
                    self.log_debug(f"Estratti {len(content)} caratteri da EPUB")
                    return content
                except ImportError:
                    self.log_debug("Librerie ebooklib o bs4 non disponibili")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file EPUB {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # MOBI - NUOVO
            elif ext == '.mobi':
                try:
                    import mobi
                    import tempfile
                    import shutil
                    self.log_debug(f"Processando file MOBI: {file_path}")
                    
                    tempdir = tempfile.mkdtemp()
                    try:
                        # Estrai il contenuto del file MOBI
                        extractor = mobi.Mobi(file_path)
                        extractor.extract(tempdir)
                        
                        # Leggi il testo estratto
                        text_content = []
                        for filename in os.listdir(tempdir):
                            if filename.endswith('.txt'):
                                with open(os.path.join(tempdir, filename), 'r', encoding='utf-8', errors='replace') as f:
                                    text_content.append(f.read())
                        
                        content = "\n".join(text_content)
                        self.log_debug(f"Estratti {len(content)} caratteri da MOBI")
                        return content
                    finally:
                        # Pulisci i file temporanei
                        shutil.rmtree(tempdir)
                except ImportError:
                    self.log_debug("Libreria mobi non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file MOBI {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""       
            # LaTeX - NUOVO
            elif ext == '.tex':
                try:
                    self.log_debug(f"Processando file LaTeX: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        
                    # Rimuovi i comandi LaTeX più comuni
                    import re
                    # Rimuovi i comandi
                    content = re.sub(r'\\[a-zA-Z]+(\{[^}]*\}|\[[^\]]*\])*', ' ', content)
                    # Rimuovi gli ambienti
                    content = re.sub(r'\\begin\{[^}]*\}(.*?)\\end\{[^}]*\}', ' ', content, flags=re.DOTALL)
                    # Rimuovi i commenti
                    content = re.sub(r'%.*?(\n|$)', ' ', content)
                    # Rimuovi le graffe
                    content = re.sub(r'\{|\}', '', content)
                    # Sostituisci più spazi con uno solo
                    content = re.sub(r'\s+', ' ', content)
                    
                    self.log_debug(f"Estratti {len(content)} caratteri da LaTeX")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file LaTeX {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # reStructuredText - NUOVO
            elif ext == '.rst':
                try:
                    self.log_debug(f"Processando file reStructuredText: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        
                    # Rimuovi gli elementi di markup più comuni
                    import re
                    # Rimuovi i titoli
                    content = re.sub(r'(=+|-+|~+|\^+|"+)\n', ' ', content)
                    # Rimuovi i link
                    content = re.sub(r'`[^`]*`_', ' ', content)
                    # Rimuovi i riferimenti alle direttive
                    content = re.sub(r'\.\. [a-z]+::', ' ', content)
                    
                    self.log_debug(f"Estratti {len(content)} caratteri da RST")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file RST {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""        
            # SQLite database (.db, .sqlite, .sqlite3) - NUOVO
            elif ext in ['.db', '.sqlite', '.sqlite3']:
                try:
                    import sqlite3
                    self.log_debug(f"Processando file SQLite: {file_path}")
                    
                    # Verifica che sia un file SQLite valido
                    if not os.path.getsize(file_path) > 100:
                        return ""
                        
                    try:
                        # Connetti al database
                        conn = sqlite3.connect(file_path)
                        cursor = conn.cursor()
                        
                        # Ottieni la lista delle tabelle
                        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                        tables = cursor.fetchall()
                        
                        content_parts = []
                        content_parts.append(f"Database SQLite: {os.path.basename(file_path)}")
                        content_parts.append(f"Numero di tabelle: {len(tables)}")
                        
                        # Estrai struttura e campione di dati da ogni tabella
                        for table in tables:
                            table_name = table[0]
                            content_parts.append(f"\n--- Tabella: {table_name} ---")
                            
                            # Ottieni struttura della tabella
                            cursor.execute(f"PRAGMA table_info({table_name});")
                            columns = cursor.fetchall()
                            col_names = [col[1] for col in columns]
                            content_parts.append("Colonne: " + ", ".join(col_names))
                            
                            # Ottieni un campione di dati (massimo 10 righe)
                            try:
                                cursor.execute(f"SELECT * FROM {table_name} LIMIT 10;")
                                rows = cursor.fetchall()
                                if rows:
                                    content_parts.append(f"Campione dati ({len(rows)} righe):")
                                    for row in rows:
                                        content_parts.append(str(row))
                            except:
                                content_parts.append("Errore nell'estrazione del campione dati")
                        
                        conn.close()
                        content = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(content)} caratteri da SQLite")
                        return content
                        
                    except sqlite3.Error:
                        # Non è un database SQLite valido o è cifrato
                        return ""
                except ImportError:
                    self.log_debug("Libreria sqlite3 non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file SQLite {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # Microsoft Access (.mdb, .accdb)
            elif ext in ['.mdb', '.accdb']:
                try:
                    import pyodbc
                    self.log_debug(f"Processando file Access: {file_path}")
                    
                    # Verifica che siamo su Windows
                    if os.name != 'nt':
                        self.log_debug("L'accesso ai file MDB/ACCDB è supportato solo su Windows")
                        return ""
                    
                    # Connetti al database Access
                    driver = "Microsoft Access Driver (*.mdb, *.accdb)"
                    conn_str = f"Driver={{{driver}}};DBQ={file_path};"
                    
                    try:
                        conn = pyodbc.connect(conn_str)
                        cursor = conn.cursor()
                        
                        # Ottieni l'elenco delle tabelle
                        tables = []
                        for row in cursor.tables():
                            if row.table_type == 'TABLE':
                                tables.append(row.table_name)
                        
                        content_parts = []
                        content_parts.append(f"Database Access: {os.path.basename(file_path)}")
                        content_parts.append(f"Numero di tabelle: {len(tables)}")
                        
                        # Estrai struttura e dati da ogni tabella
                        for table in tables:
                            if self.stop_search:
                                return ""
                                
                            content_parts.append(f"\n--- Tabella: {table} ---")
                            
                            # Ottieni struttura della tabella
                            columns = cursor.columns(table=table)
                            col_names = [col.column_name for col in columns]
                            content_parts.append("Colonne: " + ", ".join(col_names))
                            
                            # Ricerca nei dati della tabella - MIGLIORIA PER RICERCA NEI DATI
                            try:
                                # Esegui query su tutti i dati della tabella
                                cursor.execute(f"SELECT * FROM [{table}]")
                                rows = cursor.fetchall()
                                if rows:
                                    # Aggiungi un campione di righe al contenuto per la ricerca
                                    content_parts.append(f"Dati ({len(rows)} righe):")
                                    # Memorizza tutte le righe come stringhe per permettere la ricerca
                                    for row in rows:
                                        row_str = " ".join([str(cell) for cell in row if cell is not None])
                                        content_parts.append(row_str)
                            except Exception as e:
                                content_parts.append(f"Errore nell'estrazione dei dati: {str(e)}")
                        
                        conn.close()
                        content = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(content)} caratteri da Access")
                        return content
                        
                    except pyodbc.Error as e:
                        self.log_debug(f"Errore di accesso al database: {str(e)}")
                        return ""
                except ImportError:
                    self.log_debug("Libreria pyodbc non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Access {file_path}: {str(e)}")
                    return ""

            # OpenDocument Database (.odb)
            elif ext == '.odb':
                try:
                    import zipfile
                    import xml.etree.ElementTree as ET
                    self.log_debug(f"Processando file ODB: {file_path}")
                    
                    # ODB è essenzialmente un file ZIP con file XML all'interno
                    if zipfile.is_zipfile(file_path):
                        with zipfile.ZipFile(file_path, 'r') as zip_ref:
                            content_parts = []
                            
                            # Estrai il file contenente lo schema
                            try:
                                with zip_ref.open('content.xml') as content_file:
                                    tree = ET.parse(content_file)
                                    root = tree.getroot()
                                    
                                    # Cerca namespace
                                    ns = {'db': 'urn:oasis:names:tc:opendocument:xmlns:database:1.0',
                                        'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'}
                                    
                                    # Estrai informazioni sulle tabelle
                                    tables = root.findall('.//db:table', ns)
                                    content_parts.append(f"Database ODB: {os.path.basename(file_path)}")
                                    content_parts.append(f"Numero di tabelle trovate: {len(tables)}")
                                    
                                    for table in tables:
                                        if self.stop_search:
                                            return ""
                                            
                                        if 'db:name' in table.attrib:
                                            table_name = table.get('{urn:oasis:names:tc:opendocument:xmlns:database:1.0}name')
                                            content_parts.append(f"\n--- Tabella: {table_name} ---")
                                            
                                            # Estrai colonne
                                            columns = table.findall('.//db:column', ns)
                                            col_names = []
                                            for column in columns:
                                                if 'db:name' in column.attrib:
                                                    col_name = column.get('{urn:oasis:names:tc:opendocument:xmlns:database:1.0}name')
                                                    col_names.append(col_name)
                                            
                                            content_parts.append("Colonne: " + ", ".join(col_names))
                                    
                                    # NUOVO: Cercare di estrarre anche i dati dalle tabelle
                                    # Controlla se esiste il file di dati tipicamente in format/database/data
                                    try:
                                        # In ODB, i dati possono essere in diversi formati/posizioni
                                        # Questo è solo un esempio, potrebbe richiedere ulteriori adattamenti
                                        data_files = [f for f in zip_ref.namelist() if f.startswith('database/') and f.endswith('.dbf')]
                                        if data_files:
                                            content_parts.append("\n--- Dati estratti ---")
                                            for data_file in data_files:
                                                with zip_ref.open(data_file) as df:
                                                    # Estrai contenuto raw per permettere la ricerca
                                                    content_parts.append(f"Contenuto del file dati: {data_file}")
                                                    # Leggi i primi KB del file per permettere la ricerca
                                                    content_parts.append(df.read(10240).decode('utf-8', errors='ignore'))
                                    except Exception as e:
                                        content_parts.append(f"Errore nell'estrazione dei dati: {str(e)}")
                            except Exception as e:
                                content_parts.append(f"Errore nell'estrazione dello schema: {str(e)}")
                            
                            content = "\n".join(content_parts)
                            self.log_debug(f"Estratti {len(content)} caratteri da ODB")
                            return content
                    else:
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file ODB {file_path}: {str(e)}")
                    return ""
            # Tab-Separated Values (.tsv) - NUOVO
            elif ext == '.tsv':
                try:
                    import csv
                    self.log_debug(f"Processando file TSV: {file_path}")
                    
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        reader = csv.reader(f, delimiter='\t')
                        rows = []
                        
                        # Limita a 1000 righe per file grandi
                        for i, row in enumerate(reader):
                            if i >= 1000:
                                rows.append("... (file troncato, troppe righe)")
                                break
                            rows.append("\t".join(row))
                            
                    content = "\n".join(rows)
                    self.log_debug(f"Estratti {len(content)} caratteri da TSV")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file TSV {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # dBase format (.dbf) - NUOVO
            elif ext == '.dbf':
                try:
                    import dbfread
                    self.log_debug(f"Processando file DBF: {file_path}")
                    
                    table = dbfread.DBF(file_path)
                    records = []
                    
                    # Ottieni nomi delle colonne
                    headers = table.field_names
                    records.append("Colonne: " + ", ".join(headers))
                    
                    # Ottieni un campione di dati
                    for i, record in enumerate(table):
                        if i >= 50:  # Limita a 50 record
                            records.append("... (file troncato, troppi record)")
                            break
                            
                        record_data = []
                        for field in headers:
                            record_data.append(f"{field}: {record[field]}")
                        records.append(" | ".join(record_data))
                        
                    content = "\n".join(records)
                    self.log_debug(f"Estratti {len(content)} caratteri da DBF")
                    return content
                except ImportError:
                    self.log_debug("Libreria dbfread non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file DBF {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # Data Interchange Format (.dif) - NUOVO
            elif ext == '.dif':
                try:
                    self.log_debug(f"Processando file DIF: {file_path}")
                    
                    # I file DIF hanno una struttura specifica
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        lines = f.readlines()
                        
                    if len(lines) < 3:
                        return ""
                        
                    content_lines = []
                    i = 0
                    
                    # Salta l'intestazione
                    while i < len(lines) and not lines[i].strip().startswith('DATA'):
                        i += 1
                        
                    # Estrai i dati
                    data_mode = False
                    current_row = []
                    
                    while i < len(lines):
                        line = lines[i].strip()
                        
                        if line.startswith('BOT'):  # Beginning of tuple
                            current_row = []
                        elif line.startswith('EOD'):  # End of data
                            break
                        elif not data_mode and line.startswith('1,0'):
                            data_mode = True
                        elif data_mode and line.startswith('V') or line.startswith('C'):
                            # La prossima riga contiene il valore
                            i += 1
                            if i < len(lines):
                                value = lines[i].strip().strip('"')
                                current_row.append(value)
                        elif line.startswith('EOT'):  # End of tuple
                            content_lines.append(",".join(current_row))
                            data_mode = False
                            
                        i += 1
                        
                    content = "\n".join(content_lines)
                    self.log_debug(f"Estratti {len(content)} caratteri da DIF")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file DIF {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # Testo semplice
            elif ext in ['.txt', '.csv', '.log', '.ini', '.xml', '.json', '.md', '.html', '.htm',
                        '.py', '.js', '.java', '.cpp', '.c', '.cs', '.php', '.rb', '.go', '.swift', 
                        '.sql', '.sh', '.bat', '.ps1', '.vbs', '.pl', '.ts', '.kt', '.scala',
                        '.h', '.hpp', '.vb', '.lua', '.rs', '.groovy', '.yml', '.yaml', '.toml',
                        '.properties', '.conf', '.config', '.cfg', '.reg']:
                try:
                    # Apri con diverse codifiche per essere robusto
                    encodings = ['utf-8', 'latin-1', 'windows-1252']
                    content = ""
                    
                    for encoding in encodings:
                        try:
                            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                                content = f.read(1024*1024)  # Leggi al massimo 1 MB
                                break
                        except UnicodeDecodeError:
                            continue
                        except Exception as e:
                            self.log_debug(f"Errore con codifica {encoding}: {str(e)}")
                    
                    if content:
                        self.log_debug(f"Estratti {len(content)} caratteri da file di testo")
                        return content
                    else:
                        self.log_debug("Nessun contenuto estratto dal file di testo")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file di testo {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # PDF (aggiunto per completezza)
            elif ext == '.pdf':
                try:
                    import PyPDF2
                    self.log_debug(f"Processando file PDF: {file_path}")
                    
                    content = []
                    with open(file_path, 'rb') as f:
                        try:
                            reader = PyPDF2.PdfReader(f)
                            num_pages = min(len(reader.pages), 50)  # Limita a 50 pagine
                            
                            for page_num in range(num_pages):
                                try:
                                    page_text = reader.pages[page_num].extract_text()
                                    if page_text and page_text.strip():
                                        content.append(f"--- Pagina {page_num+1} ---")
                                        content.append(page_text)
                                except Exception as e:
                                    self.log_debug(f"Errore nell'estrazione testo pagina {page_num}: {str(e)}")
                        except Exception as e:
                            self.log_debug(f"Errore nell'apertura del PDF: {str(e)}")
                    
                    result = "\n".join(content)
                    self.log_debug(f"Estratti {len(result)} caratteri da PDF")
                    return result
                    
                except ImportError:
                    self.log_debug("Libreria PyPDF2 non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nella lettura del file PDF: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # ===== EMAIL E CALENDARIO =====
            # File Email (.eml)
            elif ext == '.eml':
                try:
                    self.log_debug(f"Processando file Email: {file_path}")
                    import email
                    import base64
                    import quopri
                    from email.header import decode_header
                    
                    # Leggi il file come binario per gestire correttamente tutti i tipi di codifica
                    with open(file_path, 'rb') as f:
                        msg_data = f.read()
                        
                    # Verifica che abbiamo letto dei dati
                    if not msg_data:
                        self.log_debug(f"File EML vuoto: {file_path}")
                        return ""
                        
                    self.log_debug(f"Letti {len(msg_data)} bytes dal file EML")
                    
                    # Analizza il messaggio email
                    msg = email.message_from_bytes(msg_data)
                    content_parts = []
                    
                    # Estrai intestazioni
                    self.log_debug("Estrazione intestazioni email")
                    for header in ['From', 'To', 'Subject', 'Date']:
                        if msg[header]:
                            # Decodifica l'intestazione se necessario
                            try:
                                header_parts = decode_header(msg[header])
                                decoded_header = ""
                                for part, encoding in header_parts:
                                    if isinstance(part, bytes):
                                        decoded_header += part.decode(encoding or 'utf-8', errors='replace')
                                    else:
                                        decoded_header += part
                                content_parts.append(f"{header}: {decoded_header}")
                            except:
                                content_parts.append(f"{header}: {msg[header]}")
                    
                    # Aggiungi un separatore dopo le intestazioni
                    if content_parts:
                        content_parts.append("-" * 40)
                    
                    # Estrai corpo del messaggio e allegati
                    attachment_count = 0
                    self.log_debug("Inizio analisi parti dell'email")
                    
                    # Elabora tutte le parti del messaggio
                    for part in msg.walk():
                        content_type = part.get_content_type()
                        self.log_debug(f"Elaborazione parte email: {content_type}")
                        
                        # Estrai testo semplice dal corpo dell'email
                        if content_type == "text/plain" and not part.get_filename():
                            try:
                                payload = part.get_payload(decode=True)
                                if payload:
                                    charset = part.get_content_charset() or 'utf-8'
                                    text = payload.decode(charset, errors='replace')
                                    content_parts.append(text)
                                    self.log_debug(f"Estratti {len(text)} caratteri di testo")
                            except Exception as e:
                                self.log_debug(f"Errore nell'estrazione del testo: {str(e)}")
                                try:
                                    # Fallback a utf-8
                                    content_parts.append(payload.decode('utf-8', errors='replace'))
                                except:
                                    pass
                                    
                        # Estrai HTML dal corpo dell'email
                        elif content_type == "text/html" and not part.get_filename():
                            try:
                                payload = part.get_payload(decode=True)
                                if payload:
                                    charset = part.get_content_charset() or 'utf-8'
                                    html_content = payload.decode(charset, errors='replace')
                                    
                                    # Opzionalmente, prova a estrarre il testo dall'HTML
                                    try:
                                        from bs4 import BeautifulSoup
                                        soup = BeautifulSoup(html_content, 'html.parser')
                                        text_content = soup.get_text(separator=' ', strip=True)
                                        content_parts.append(text_content)
                                        self.log_debug(f"Estratti {len(text_content)} caratteri da HTML")
                                    except ImportError:
                                        # Se BeautifulSoup non è disponibile, usa l'HTML grezzo
                                        content_parts.append(html_content)
                                        self.log_debug(f"BeautifulSoup non disponibile, usato HTML grezzo ({len(html_content)} caratteri)")
                            except Exception as e:
                                self.log_debug(f"Errore nell'estrazione HTML: {str(e)}")
                        
                        # Gestisci gli allegati
                        elif part.get('Content-Disposition') or part.get_filename():
                            try:
                                # Ottieni il nome dell'allegato
                                filename = part.get_filename()
                                attachment_count += 1
                                
                                # Normalizza il nome dell'allegato
                                if not filename:
                                    filename = f"allegato_{attachment_count}"
                                
                                # Decodifica il nome dell'allegato se necessario
                                if isinstance(filename, bytes):
                                    try:
                                        filename = filename.decode('utf-8', errors='replace')
                                    except:
                                        filename = f"allegato_{attachment_count}"
                                
                                # Decodifica header encodati se necessario
                                try:
                                    decoded_parts = decode_header(filename)
                                    decoded_filename = ""
                                    for part_data, charset in decoded_parts:
                                        if isinstance(part_data, bytes):
                                            decoded_filename += part_data.decode(charset or 'utf-8', errors='replace')
                                        else:
                                            decoded_filename += part_data
                                    filename = decoded_filename
                                except:
                                    pass
                                
                                self.log_debug(f"Allegato trovato: {filename}")
                                
                                # Aggiungi informazioni sull'allegato
                                content_parts.append(f"\n--- ALLEGATO {attachment_count}: {filename} ---\n")
                                
                                # METODO 1: Usa get_payload(decode=True) che dovrebbe gestire automaticamente la decodifica
                                payload = part.get_payload(decode=True)
                                
                                # METODO 2: se il primo metodo fallisce, prova con decodifica manuale
                                if not payload:
                                    encoding = part.get('Content-Transfer-Encoding', '').lower()
                                    raw_payload = part.get_payload()
                                    
                                    if encoding == 'base64' and isinstance(raw_payload, str):
                                        try:
                                            payload = base64.b64decode(raw_payload)
                                            self.log_debug("Allegato decodificato con base64")
                                        except:
                                            self.log_debug("Errore nella decodifica base64")
                                    elif encoding == 'quoted-printable' and isinstance(raw_payload, str):
                                        try:
                                            payload = quopri.decodestring(raw_payload.encode('utf-8'))
                                            self.log_debug("Allegato decodificato con quoted-printable")
                                        except:
                                            self.log_debug("Errore nella decodifica quoted-printable")
                                    elif isinstance(raw_payload, str):
                                        try:
                                            payload = raw_payload.encode('utf-8', errors='replace')
                                            self.log_debug("Allegato convertito da stringa a bytes")
                                        except:
                                            self.log_debug("Errore nella conversione della stringa a bytes")
                                
                                if payload and len(payload) > 0:
                                    self.log_debug(f"Analisi contenuto allegato: {filename} ({len(payload)} bytes)")
                                    
                                    # IMPORTANTE: CHIAMA IL METODO PER ELABORARE L'ALLEGATO
                                    attachment_content = self.process_email_attachment(
                                        payload, filename, part.get_content_type())
                                    
                                    if attachment_content:
                                        content_parts.append(attachment_content)
                                        self.log_debug(f"Contenuto allegato {filename} aggiunto ({len(attachment_content)} caratteri)")
                                    else:
                                        content_parts.append(f"[Allegato {filename}: nessun contenuto estraibile]")
                                        self.log_debug(f"Nessun contenuto estraibile dall'allegato {filename}")
                                else:
                                    content_parts.append(f"[Allegato {filename}: vuoto o non decodificabile]")
                                    self.log_debug(f"Allegato vuoto o non decodificabile: {filename}")
                                    
                            except Exception as e:
                                self.log_debug(f"Errore nell'elaborazione dell'allegato: {str(e)}")
                                content_parts.append(f"[Errore nell'elaborazione dell'allegato: {str(e)}]")
                    
                    self.log_debug(f"Totale allegati trovati: {attachment_count}")
                    
                    # Se non è stato trovato alcun contenuto, prova un metodo più semplice
                    if not content_parts:
                        self.log_debug("Nessun contenuto trovato, tentativo con metodo alternativo")
                        try:
                            if msg.is_multipart():
                                for part in msg.walk():
                                    content_type = part.get_content_type()
                                    if content_type == "text/plain":
                                        payload = part.get_payload(decode=True)
                                        if payload:
                                            content_parts.append(payload.decode('utf-8', errors='replace'))
                            else:
                                payload = msg.get_payload(decode=True)
                                if payload:
                                    content_parts.append(payload.decode('utf-8', errors='replace'))
                        except Exception as e:
                            self.log_debug(f"Errore nel metodo alternativo: {str(e)}")
                    
                    content = "\n".join(content_parts)
                    self.log_debug(f"Estratti in totale {len(content)} caratteri da EML (inclusi allegati)")
                    return content
                except ImportError as e:
                    self.log_debug(f"Modulo necessario non disponibile: {str(e)}")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file EML {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File vCard (.vcf)
            elif ext == '.vcf':
                try:
                    self.log_debug(f"Processando file vCard: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        # Estrai campi più importanti per la ricerca
                        processed_content = []
                        for line in content.splitlines():
                            if line.startswith(('FN:', 'N:', 'EMAIL:', 'TEL:', 'ADR:', 'ORG:', 'TITLE:', 'NOTE:')):
                                processed_content.append(line)
                        
                        result = "\n".join(processed_content)
                        self.log_debug(f"Estratti {len(result)} caratteri da vCard")
                        return result
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file vCard {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File iCalendar (.ics)
            elif ext == '.ics':
                try:
                    self.log_debug(f"Processando file iCalendar: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        # Estrai campi più importanti per la ricerca
                        processed_content = []
                        current_event = []
                        in_event = False
                        
                        for line in content.splitlines():
                            line = line.strip()
                            if line == "BEGIN:VEVENT":
                                in_event = True
                                current_event = []
                            elif line == "END:VEVENT":
                                in_event = False
                                processed_content.append("\n".join(current_event))
                                processed_content.append("---")
                            elif in_event and line.startswith(('SUMMARY:', 'DESCRIPTION:', 'LOCATION:', 
                                                            'ORGANIZER:', 'ATTENDEE:', 'DTSTART:', 
                                                            'DTEND:', 'CATEGORIES:')):
                                current_event.append(line)
                        
                        result = "\n".join(processed_content)
                        self.log_debug(f"Estratti {len(result)} caratteri da iCalendar")
                        return result
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file iCalendar {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # ===== PRESENTAZIONI =====
            # File PowerPoint Show (.pps)
            elif ext == '.pps':
                try:
                    # Utilizziamo la stessa implementazione di PPT dato che hanno lo stesso formato
                    if os.name == 'nt':
                        import win32com.client
                        import pythoncom
                        
                        # Inizializzazione necessaria per i thread
                        pythoncom.CoInitialize()
                        
                        self.log_debug(f"Processando file PowerPoint Show: {file_path}")
                        
                        try:
                            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                            powerpoint.Visible = False
                            
                            presentation = powerpoint.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
                            texts = []
                            
                            for slide_idx in range(1, presentation.Slides.Count + 1):
                                slide = presentation.Slides.Item(slide_idx)
                                texts.append(f"--- Diapositiva {slide_idx} ---")
                                slide_text = []
                                
                                for shape_idx in range(1, slide.Shapes.Count + 1):
                                    shape = slide.Shapes.Item(shape_idx)
                                    if shape.HasTextFrame:
                                        if shape.TextFrame.HasText:
                                            slide_text.append(shape.TextFrame.TextRange.Text)
                                
                                if slide_text:
                                    texts.append("\n".join(slide_text))
                            
                            presentation.Close()
                            powerpoint.Quit()
                            
                            pythoncom.CoUninitialize()
                            
                            result = "\n".join(texts)
                            self.log_debug(f"Estratti {len(result)} caratteri da PPS")
                            return result
                            
                        except Exception as e:
                            self.log_debug(f"Errore nell'apertura del PPS con win32com: {str(e)}")
                            
                            # Cleanup in caso di errore
                            try:
                                if 'presentation' in locals() and presentation:
                                    presentation.Close()
                                if 'powerpoint' in locals() and powerpoint:
                                    powerpoint.Quit()
                            except:
                                pass
                            
                            pythoncom.CoUninitialize()
                            return ""
                    else:
                        self.log_debug("Estrazione di testo dai file PPS non supportata su questa piattaforma")
                        return ""
                except ImportError:
                    self.log_debug("win32com non disponibile per i file PPS")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore generale nell'elaborazione del file PPS {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File Keynote (.key)
            elif ext == '.key':
                try:
                    self.log_debug(f"Processando file Keynote: {file_path}")
                    # Keynote è essenzialmente un pacchetto compresso con file XML all'interno
                    if not zipfile.is_zipfile(file_path):
                        self.log_debug(f"File Keynote non valido (non è un file zip): {file_path}")
                        return ""
                        
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        content_parts = []
                        
                        # Cerca file di testo nel pacchetto Keynote
                        for file_info in zip_ref.infolist():
                            if file_info.filename.endswith(('.xml', '.txt')):
                                try:
                                    with zip_ref.open(file_info) as content_file:
                                        content = content_file.read().decode('utf-8', errors='replace')
                                        # Estrai solo il testo dalle presentazioni, rimuovendo i tag
                                        import re
                                        text_content = re.sub(r'<[^>]+>', ' ', content)
                                        text_content = re.sub(r'\s+', ' ', text_content).strip()
                                        if text_content:
                                            content_parts.append(text_content)
                                except:
                                    continue
                                    
                        result = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(result)} caratteri da Keynote")
                        return result
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Keynote {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # ===== FILE DI CONFIGURAZIONE =====
            # File YAML (.yml, .yaml)
            elif ext in ['.yml', '.yaml']:
                try:
                    self.log_debug(f"Processando file YAML: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da YAML")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file YAML {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File TOML (.toml)
            elif ext == '.toml':
                try:
                    self.log_debug(f"Processando file TOML: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da TOML")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file TOML {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File Registry Windows (.reg)
            elif ext == '.reg':
                try:
                    self.log_debug(f"Processando file Registry: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da Registry")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Registry {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File plist (.plist)
            elif ext == '.plist':
                try:
                    self.log_debug(f"Processando file Property List: {file_path}")
                    # Prova a leggere come XML
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        # Rimuovi i tag XML per estrarre solo il testo
                        import re
                        text_content = re.sub(r'<[^>]+>', ' ', content)
                        text_content = re.sub(r'\s+', ' ', text_content).strip()
                        self.log_debug(f"Estratti {len(text_content)} caratteri da plist")
                        return text_content
                except UnicodeDecodeError:
                    try:
                        # Potrebbe essere un file plist binario
                        self.log_debug("Tentativo di lettura come plist binario")
                        with open(file_path, 'rb') as f:
                            content = f.read().decode('latin-1', errors='replace')
                            # Estrai stringhe leggibili
                            printable = ''.join(c for c in content if c.isprintable() and len(c.strip()) > 0)
                            self.log_debug(f"Estratti {len(printable)} caratteri da plist binario")
                            return printable
                    except Exception as inner_e:
                        self.log_debug(f"Errore nella lettura del plist binario: {str(inner_e)}")
                        return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file plist {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File Properties (.properties)
            elif ext == '.properties':
                try:
                    self.log_debug(f"Processando file Properties: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da Properties")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file Properties {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File htaccess (.htaccess)
            elif file_path.endswith('.htaccess'):
                try:
                    self.log_debug(f"Processando file htaccess: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da htaccess")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file htaccess {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # ===== LINGUAGGI DI PROGRAMMAZIONE =====
            # I vari linguaggi di programmazione possono usare lo stesso parser di testo
            elif ext in ['.h', '.hpp', '.vb', '.lua', '.rs', '.groovy']:
                try:
                    self.log_debug(f"Processando file di codice {ext}: {file_path}")
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                        content = f.read()
                        self.log_debug(f"Estratti {len(content)} caratteri da file {ext}")
                        return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file {ext} {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File MSG (Outlook)
            elif ext == '.msg':
                try:
                    self.log_debug(f"Processando file MSG Outlook: {file_path}")
                    try:
                        import extract_msg
                        msg = extract_msg.Message(file_path)
                        
                        content_parts = []
                        
                        # Estrai intestazioni principali
                        from_address = msg.sender if hasattr(msg, 'sender') else "N/A"
                        to_address = msg.to if hasattr(msg, 'to') else "N/A"
                        subject = msg.subject if hasattr(msg, 'subject') else "N/A"
                        msg_date = msg.date if hasattr(msg, 'date') else "N/A"
                        
                        # Aggiungi intestazioni al contenuto
                        content_parts.append(f"Da: {from_address}")
                        content_parts.append(f"A: {to_address}")
                        content_parts.append(f"Oggetto: {subject}")
                        content_parts.append(f"Data: {msg_date}")
                        content_parts.append("-" * 40)  # Separatore
                        
                        # Estrai corpo del messaggio
                        msg_body = msg.body if hasattr(msg, 'body') else ""
                        if msg_body:
                            content_parts.append(msg_body)
                        
                        # Contatori per monitoraggio allegati
                        attachment_count = 0
                        
                        # Versione migliorata per gestire gli allegati di MSG
                        for attachment in msg.attachments:
                            try:
                                # Verifica se l'allegato è valido
                                if not attachment:
                                    continue
                                    
                                # Ottieni il nome del file 
                                filename = None
                                if hasattr(attachment, 'longFilename') and attachment.longFilename:
                                    filename = attachment.longFilename
                                elif hasattr(attachment, 'shortFilename') and attachment.shortFilename:
                                    filename = attachment.shortFilename
                                elif hasattr(attachment, 'filename') and attachment.filename:
                                    filename = attachment.filename
                                elif hasattr(attachment, 'name') and attachment.name:
                                    filename = attachment.name
                                
                                # Se ancora non abbiamo un nome, generiamone uno
                                if not filename:
                                    filename = f"allegato_{attachment_count + 1}"
                                
                                attachment_count += 1
                                content_parts.append(f"\n--- ALLEGATO {attachment_count}: {filename} ---\n")
                                
                                # Ottieni i dati binari dell'allegato - metodo migliorato
                                attachment_data = None
                                if hasattr(attachment, 'data'):
                                    attachment_data = attachment.data
                                elif hasattr(attachment, 'getBytes') and callable(attachment.getBytes):
                                    attachment_data = attachment.getBytes()
                                elif hasattr(attachment, 'getData') and callable(attachment.getData):
                                    attachment_data = attachment.getData()
                                # Per la nuova versione di extract_msg che usa 'content'
                                elif hasattr(attachment, 'content') and attachment.content:
                                    attachment_data = attachment.content
                                # Per la versione più recente che potrebbe usare 'bytes'
                                elif hasattr(attachment, 'bytes'):
                                    attachment_data = attachment.bytes
                                
                                if not attachment_data:
                                    content_parts.append(f"[Allegato vuoto o non leggibile]")
                                    continue
                                
                                # Ottieni il tipo di contenuto (MIME type)
                                content_type = ""
                                if hasattr(attachment, 'mimetype') and attachment.mimetype:
                                    content_type = attachment.mimetype
                                elif hasattr(attachment, 'contentType') and attachment.contentType:
                                    content_type = attachment.contentType
                                
                                # Logging dettagliato per debug
                                self.log_debug(f"Allegato trovato: {filename} ({len(attachment_data)} bytes, tipo: {content_type or 'non definito'})")
                                
                                # Processa l'allegato - più log per tracciare il processo
                                self.log_debug(f"Inizio elaborazione contenuto allegato {filename}")
                                attachment_content = self.process_email_attachment(
                                    attachment_data, filename, content_type)
                                
                                if attachment_content:
                                    self.log_debug(f"Contenuto estratto da allegato '{filename}': {len(attachment_content)} caratteri")
                                    content_parts.append(attachment_content)
                                else:
                                    self.log_debug(f"Nessun contenuto estratto da allegato '{filename}'")
                                    content_parts.append(f"[Allegato {filename}: nessun contenuto estraibile]")
                                    
                            except Exception as e:
                                self.log_debug(f"Errore nell'elaborazione dell'allegato {attachment_count}: {str(e)}")
                                content_parts.append(f"[Errore nell'elaborazione dell'allegato {attachment_count}: {str(e)}]")
                        
                        content = "\n".join(content_parts)
                        self.log_debug(f"Estratti {len(content)} caratteri da MSG (inclusi {attachment_count} allegati)")
                        return content
                        
                    except ImportError:
                        self.log_debug("Libreria extract_msg non disponibile")
                        # Metodo fallback - estrai contenuto usando regexp base
                        with open(file_path, 'rb') as f:
                            binary_content = f.read()
                            text_content = ""
                            # Cerca stringhe ASCII leggibili nel file binario
                            import re
                            text_chunks = re.findall(b'[\x20-\x7E\r\n]{4,}', binary_content)
                            for chunk in text_chunks:
                                try:
                                    text_content += chunk.decode('utf-8', errors='replace') + "\n"
                                except:
                                    pass
                            
                            if text_content:
                                self.log_debug(f"Estratti {len(text_content)} caratteri da MSG (metodo fallback)")
                                return text_content
                            return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file MSG {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File PST/OST (Outlook database) - versione semplificata senza dipendenze esterne
            elif ext in ['.pst', '.ost']:
                try:
                    self.log_debug(f"Processando file {ext} Outlook: {file_path}")
                    
                    # Estrai metadati di base
                    file_size = os.path.getsize(file_path)
                    content_parts = []
                    
                    content_parts.append(f"File {ext.upper()} Outlook")
                    content_parts.append(f"Percorso: {file_path}")
                    content_parts.append(f"Dimensione: {self._format_size(file_size)}")
                    content_parts.append(f"Data modifica: {datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M')}")
                    
                    # Estrai stringhe di testo leggibili (approccio base)
                    try:
                        with open(file_path, 'rb') as f:
                            binary_content = f.read(1024*1024)  # Leggi solo il primo MB
                            
                            # Usa espressione regolare per trovare stringhe ASCII leggibili
                            import re
                            # Trova stringhe alfanumeriche con spazi/punteggiatura di almeno 5 caratteri
                            pattern = re.compile(b'[a-zA-Z0-9\\s\\.,@\\-_:;\'"/]{5,}')
                            matches = pattern.findall(binary_content)
                            
                            # Filtra e converte le stringhe trovate
                            strings = []
                            for match in matches:
                                try:
                                    text = match.decode('utf-8', errors='replace')
                                    # Filtra stringhe che sembrano email valide o messaggi
                                    if ('@' in text or 
                                        text.startswith("To:") or 
                                        text.startswith("From:") or 
                                        text.startswith("Subject:")):
                                        strings.append(text)
                                    # O stringhe abbastanza lunghe da essere significative
                                    elif len(text) > 15:
                                        strings.append(text)
                                except:
                                    pass
                            
                            # Aggiungi i risultati all'output
                            if strings:
                                content_parts.append("\n--- Contenuto estratto ---\n")
                                content_parts.extend(strings)
                    except Exception as e:
                        self.log_debug(f"Errore nell'estrazione del testo: {str(e)}")
                    
                    content = "\n".join(content_parts)
                    self.log_debug(f"Estratti {len(content)} caratteri da {ext}")
                    return content
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file {ext} {file_path}: {str(e)}")
                    return ""

            if self.stop_search:
                return ""
            # File MBOX - usa il modulo mailbox standard
            elif ext == '.mbox':
                try:
                    self.log_debug(f"Processando file MBOX: {file_path}")
                    import mailbox
                    import email
                    
                    # Apri il file MBOX
                    mbox = mailbox.mbox(file_path)
                    content_parts = []
                    
                    # Limita il numero di messaggi da processare
                    max_messages = 50
                    processed = 0
                    
                    content_parts.append(f"File MBOX: {os.path.basename(file_path)}")
                    content_parts.append(f"Totale messaggi: {len(mbox)}")
                    content_parts.append("-" * 40)
                    
                    # Itera attraverso i messaggi
                    for key, msg in mbox.items():
                        if processed >= max_messages:
                            content_parts.append(f"[Limitato a {max_messages} messaggi]")
                            break
                        
                        try:
                            # Estrai le intestazioni principali
                            headers = []
                            for header in ['From', 'To', 'Subject', 'Date']:
                                if msg[header]:
                                    headers.append(f"{header}: {msg[header]}")
                            
                            content_parts.append(f"\n--- MESSAGGIO {processed+1} ---")
                            content_parts.extend(headers)
                            
                            # Estrai corpo del messaggio
                            if msg.is_multipart():
                                for part in msg.walk():
                                    content_type = part.get_content_type()
                                    if content_type == "text/plain" and not part.get_filename():
                                        try:
                                            payload = part.get_payload(decode=True)
                                            if payload:
                                                charset = part.get_content_charset() or 'utf-8'
                                                text = payload.decode(charset, errors='replace')
                                                content_parts.append(text)
                                        except:
                                            pass
                            else:
                                payload = msg.get_payload(decode=True)
                                if payload:
                                    charset = msg.get_content_charset() or 'utf-8'
                                    try:
                                        text = payload.decode(charset, errors='replace')
                                        content_parts.append(text)
                                    except:
                                        pass
                            
                            processed += 1
                        except Exception as e:
                            self.log_debug(f"Errore nel processare il messaggio MBOX: {str(e)}")
                    
                    mbox.close()
                    content = "\n".join(content_parts)
                    self.log_debug(f"Estratti {len(content)} caratteri da MBOX")
                    return content
                except ImportError:
                    self.log_debug("Libreria mailbox non disponibile")
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file MBOX {file_path}: {str(e)}")
                    return ""
            
            if self.stop_search:
                return ""
            # File EMLX (Apple Mail)
            elif ext == '.emlx':
                try:
                    self.log_debug(f"Processando file EMLX: {file_path}")
                    import email
                    
                    # I file EMLX hanno un formato particolare: prima riga è un numero, poi segue il messaggio email
                    with open(file_path, 'rb') as f:
                        content = f.read()
                        
                    # Separa il numero iniziale dal resto del contenuto email
                    try:
                        # La prima riga è un numero seguito da una nuova riga
                        parts = content.split(b'\n', 1)
                        if len(parts) > 1:
                            email_content = parts[1]
                            
                            # Parse dell'email con il modulo email standard
                            msg = email.message_from_bytes(email_content)
                            
                            # Da qui in poi possiamo usare lo stesso codice per gestire il messaggio
                            # come facciamo per i file EML standard
                            content_parts = []
                            
                            # Estrai intestazioni
                            for header in ['From', 'To', 'Subject', 'Date']:
                                if msg[header]:
                                    content_parts.append(f"{header}: {msg[header]}")
                            
                            # Estrai corpo del messaggio
                            if msg.is_multipart():
                                for part in msg.walk():
                                    content_type = part.get_content_type()
                                    if content_type == "text/plain" and not part.get_filename():
                                        payload = part.get_payload(decode=True)
                                        if payload:
                                            charset = part.get_content_charset() or 'utf-8'
                                            try:
                                                text = payload.decode(charset, errors='replace')
                                                content_parts.append(text)
                                            except:
                                                pass
                            else:
                                payload = msg.get_payload(decode=True)
                                if payload:
                                    try:
                                        text = payload.decode('utf-8', errors='replace')
                                        content_parts.append(text)
                                    except:
                                        pass
                            
                            result = "\n".join(content_parts)
                            self.log_debug(f"Estratti {len(result)} caratteri da EMLX")
                            return result
                    except Exception as e:
                        self.log_debug(f"Errore nell'analisi del file EMLX: {str(e)}")
                    
                    return ""
                except Exception as e:
                    self.log_debug(f"Errore nell'analisi del file EMLX {file_path}: {str(e)}")
                    return ""
            
        except Exception as e:
            self.log_debug(f"Errore generale nella lettura del file {file_path}: {str(e)}")
            return ""
    
    @error_handler
    def extract_archive_content(self, file_path):
        """Estrae e cerca all'interno dei file compressi.
        Supporta ZIP, RAR, 7z, TAR, GZ, BZ2, ecc."""
        import subprocess
        import tempfile
        import os
        import shutil
        
        try:
            self.log_debug(f"Tentativo di estrazione del contenuto da {file_path}")
            extension = os.path.splitext(file_path.lower())[1]
            extracted_contents = {}
            temp_dir = None
            
            # Gestione file ZIP
            if extension == '.zip':
                try:
                    import zipfile
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        # Elenco dei file nell'archivio
                        file_list = [f for f in zip_ref.namelist() if not f.endswith('/')]
                        
                        # Limita il numero di file da estrarre per prestazioni
                        max_files = 100  # Limita a 100 file per archivio
                        if len(file_list) > max_files:
                            self.log_debug(f"Archivio con troppi file ({len(file_list)}), limitato a {max_files}")
                            file_list = file_list[:max_files]
                        
                        for file_in_zip in file_list:
                            try:
                                # Verifica se il file dovrebbe essere processato in base all'estensione
                                if self.should_search_content(file_in_zip):
                                    with zip_ref.open(file_in_zip) as f:
                                        # Limita la dimensione del file da estrarre
                                        content = f.read(10 * 1024 * 1024)  # Max 10MB per file
                                        try:
                                            # Tenta di decodificare come testo
                                            text_content = content.decode('utf-8', errors='ignore')
                                            extracted_contents[file_in_zip] = text_content
                                        except:
                                            # Se fallisce, tratta come file binario
                                            extracted_contents[file_in_zip] = f"[Contenuto binario: {file_in_zip}]"
                            except Exception as e:
                                self.log_error(f"Errore nell'estrazione del file {file_in_zip} dall'archivio", exception=e)
                except ImportError:
                    self.log_error("Modulo zipfile non disponibile")
                    if "zipfile" not in missing_libraries:
                        missing_libraries.append("zipfile")
            
            # Gestione file RAR
            elif extension == '.rar':
                try:
                    import rarfile
                    with rarfile.RarFile(file_path) as rar_ref:
                        file_list = [f for f in rar_ref.namelist() if not rar_ref.getinfo(f).isdir()]
                        
                        # Limita il numero di file
                        max_files = 100
                        if len(file_list) > max_files:
                            self.log_debug(f"Archivio con troppi file ({len(file_list)}), limitato a {max_files}")
                            file_list = file_list[:max_files]
                        
                        for file_in_rar in file_list:
                            try:
                                if self.should_search_content(file_in_rar):
                                    with rar_ref.open(file_in_rar) as f:
                                        content = f.read(10 * 1024 * 1024)  # Max 10MB
                                        try:
                                            text_content = content.decode('utf-8', errors='ignore')
                                            extracted_contents[file_in_rar] = text_content
                                        except:
                                            extracted_contents[file_in_rar] = f"[Contenuto binario: {file_in_rar}]"
                            except Exception as e:
                                self.log_error(f"Errore nell'estrazione del file {file_in_rar} dall'archivio", exception=e)
                except ImportError:
                    self.log_error("Modulo rarfile non disponibile")
                    if "rarfile" not in missing_libraries:
                        missing_libraries.append("rarfile")
            
            # Gestione file 7Z - MODIFICATO PER GESTIRE L'ERRORE
            elif extension == '.7z':
                try:
                    # Crea una directory temporanea per l'estrazione
                    temp_dir = tempfile.mkdtemp()
                    self.log_debug(f"Directory temporanea creata per 7z: {temp_dir}")
                    
                    # Verifica se 7-Zip è installato (percorso comune su Windows)
                    seven_zip_paths = [
                        r"C:\Program Files\7-Zip\7z.exe",
                        r"C:\Program Files (x86)\7-Zip\7z.exe",
                        "7z"  # Se è nel PATH di sistema
                    ]
                    
                    seven_zip_exe = None
                    for path in seven_zip_paths:
                        try:
                            if os.path.isfile(path):
                                seven_zip_exe = path
                                self.log_debug(f"7-Zip trovato in: {path}")
                                break
                            elif shutil.which(path):
                                seven_zip_exe = shutil.which(path)
                                self.log_debug(f"7-Zip trovato nel PATH come: {seven_zip_exe}")
                                break
                        except Exception as e:
                            self.log_debug(f"Errore durante la verifica del percorso {path}: {str(e)}")
                    
                    if seven_zip_exe:
                        # Usa il parametro -aou per sovrascrivere automaticamente i file senza chiedere
                        command = [seven_zip_exe, "x", file_path, f"-o{temp_dir}", "-y", "-aou"]
                        self.log_debug(f"Esecuzione comando 7z con argomenti: {command}")
                        
                        import platform
                        CREATE_NO_WINDOW = 0x08000000  # Per versioni di Python precedenti alla 3.7
                        
                        # Esegui il comando con i parametri appropriati per nascondere la finestra CMD in Windows
                        if platform.system() == "Windows":
                            process = subprocess.run(
                                command,
                                stdout=subprocess.PIPE,
                                stderr=subprocess.PIPE,
                                text=True,
                                check=False,  # Non sollevare eccezioni se il codice di ritorno è diverso da zero
                                creationflags=CREATE_NO_WINDOW)
                        else:
                            # Per sistemi non-Windows
                            process = subprocess.run(
                                command,
                                stdout=subprocess.PIPE,
                                stderr=subprocess.PIPE,
                                text=True,
                                check=False)
                        
                        self.log_debug(f"Codice di uscita del processo 7z: {process.returncode}")
                        if process.stdout:
                            self.log_debug(f"Output di 7z: {process.stdout[:200]}...")  # Primi 200 caratteri
                        if process.stderr:
                            self.log_debug(f"Errori di 7z: {process.stderr}")
                        
                        if process.returncode == 0:
                            # L'estrazione è riuscita, verifichiamo i file
                            all_files = []
                            for root, _, files in os.walk(temp_dir):
                                for file in files:
                                    file_path_full = os.path.join(root, file)
                                    rel_path = os.path.relpath(file_path_full, temp_dir)
                                    all_files.append(rel_path)
                            
                            self.log_debug(f"Trovati {len(all_files)} file estratti: {all_files[:10]}...")
                            
                            # IMPORTANTE: Ottieni le parole chiave di ricerca dalla variabile corrente
                            current_keywords = self.current_search_keywords if hasattr(self, 'current_search_keywords') else []
                            if not current_keywords and hasattr(self, 'keywords_entry'):
                                try:
                                    current_keywords = self.keywords_entry.get().strip().split()
                                except:
                                    pass
                            
                            self.log_debug(f"Parole chiave per la ricerca nei file estratti: {current_keywords}")
                            
                            # Processa i file estratti (limitato a 100)
                            processed_files_count = 0
                            matched_files_count = 0
                            
                            for root, _, files in os.walk(temp_dir):
                                if processed_files_count >= 100:
                                    break
                                    
                                for file in files:
                                    if processed_files_count >= 100:
                                        break
                                        
                                    file_path_full = os.path.join(root, file)
                                    rel_path = os.path.relpath(file_path_full, temp_dir)
                                    processed_files_count += 1
                                    
                                    # Verifica se il file dovrebbe essere cercato in base all'estensione
                                    if self.should_search_content(file_path_full):
                                        self.log_debug(f"Analisi del file: {rel_path}")
                                        try:
                                            with open(file_path_full, 'rb') as f:
                                                try:
                                                    content = f.read(10 * 1024 * 1024)  # Max 10MB
                                                    text_content = content.decode('utf-8', errors='ignore')
                                                    
                                                    # Verifica esplicita della presenza delle parole chiave
                                                    found_match = False
                                                    if current_keywords:
                                                        for keyword in current_keywords:
                                                            if keyword.lower() in text_content.lower():
                                                                self.log_debug(f"Trovata corrispondenza per '{keyword}' in {rel_path}")
                                                                found_match = True
                                                                break
                                                    else:
                                                        # Se non ci sono parole chiave, includi tutti i file
                                                        found_match = True
                                                    
                                                    if found_match:
                                                        extracted_contents[rel_path] = text_content
                                                        matched_files_count += 1
                                                except UnicodeDecodeError:
                                                    # File binario - registriamo solo se corrisponde esattamente al nome cercato
                                                    if any(keyword.lower() in file.lower() for keyword in current_keywords):
                                                        extracted_contents[rel_path] = f"[Contenuto binario: {rel_path}]"
                                                        matched_files_count += 1
                                                        self.log_debug(f"Trovata corrispondenza nel nome del file binario: {rel_path}")
                                        except Exception as e:
                                            self.log_error(f"Errore nella lettura del file {rel_path}: {str(e)}")
                                    else:
                                        self.log_debug(f"File {rel_path} saltato per estensione non supportata")
                            
                            self.log_debug(f"File processati: {processed_files_count}, File con corrispondenze: {matched_files_count}")
                        else:
                            error_msg = process.stderr if process.stderr else "Errore sconosciuto"
                            self.log_error(f"Errore nell'estrazione dell'archivio 7z (codice {process.returncode}): {error_msg}")
                    else:
                        self.log_error("7-Zip non trovato sul sistema. Installare 7-Zip per supportare gli archivi .7z")
                        if "7-Zip" not in missing_libraries:
                            missing_libraries.append("7-Zip")
                except Exception as e:
                    self.log_error(f"Errore durante l'elaborazione dell'archivio 7z: {str(e)}", exception=e)
            
            # Gestione file TAR (inclusi .tar.gz, .tar.bz2)
            elif extension in ['.tar', '.tgz', '.gz', '.bz2', '.xz']:
                try:
                    import tarfile
                    mode = 'r'
                    if file_path.endswith('.gz') or file_path.endswith('.tgz'):
                        mode = 'r:gz'
                    elif file_path.endswith('.bz2'):
                        mode = 'r:bz2'
                    elif file_path.endswith('.xz'):
                        mode = 'r:xz'
                    
                    with tarfile.open(file_path, mode) as tar_ref:
                        members = tar_ref.getmembers()
                        file_members = [m for m in members if m.isfile()][:100]  # Limita a 100 file
                        
                        for member in file_members:
                            try:
                                if self.should_search_content(member.name):
                                    f = tar_ref.extractfile(member)
                                    if f:
                                        content = f.read(10 * 1024 * 1024)  # Max 10MB
                                        try:
                                            text_content = content.decode('utf-8', errors='ignore')
                                            extracted_contents[member.name] = text_content
                                        except:
                                            extracted_contents[member.name] = f"[Contenuto binario: {member.name}]"
                            except Exception as e:
                                self.log_error(f"Errore nell'estrazione del file {member.name}", exception=e)
                except ImportError:
                    self.log_error("Modulo tarfile non disponibile")
                    if "tarfile" not in missing_libraries:
                        missing_libraries.append("tarfile")
                        
            # JAR files (sono essenzialmente file ZIP)
            elif extension == '.jar':
                try:
                    import zipfile
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        file_list = [f for f in zip_ref.namelist() if not f.endswith('/')][:100]
                        
                        for file_in_jar in file_list:
                            try:
                                if self.should_search_content(file_in_jar):
                                    with zip_ref.open(file_in_jar) as f:
                                        content = f.read(10 * 1024 * 1024)  # Max 10MB
                                        try:
                                            text_content = content.decode('utf-8', errors='ignore')
                                            extracted_contents[file_in_jar] = text_content
                                        except:
                                            extracted_contents[file_in_jar] = f"[Contenuto binario: {file_in_jar}]"
                            except Exception as e:
                                self.log_error(f"Errore nell'estrazione del file {file_in_jar} dall'archivio JAR", exception=e)
                except ImportError:
                    self.log_error("Modulo zipfile non disponibile")
                    if "zipfile" not in missing_libraries:
                        missing_libraries.append("zipfile")
                
                # Gestione di altri formati compressi
            elif extension in ['.iso', '.cab']:
                self.log_debug(f"Formato compresso {extension} non supportato attualmente")
                extracted_contents["info"] = f"Il formato {extension} richiede librerie aggiuntive non installate"
            
            # Pulisci le directory temporanee
            if temp_dir and os.path.exists(temp_dir):
                import shutil
                shutil.rmtree(temp_dir, ignore_errors=True)
                    
                self.log_debug(f"Estratti {len(extracted_contents)} file con corrispondenze dall'archivio {file_path}")
            return extracted_contents
        except Exception as e:
            self.log_error(f"Errore generale nell'estrazione dell'archivio {file_path}", exception=e)
            return {}
        finally:
            # Pulisci le directory temporanee
            if temp_dir and os.path.exists(temp_dir):
                try:
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    self.log_debug(f"Directory temporanea rimossa: {temp_dir}")
                except Exception as e:
                    self.log_debug(f"Errore durante la rimozione della directory temporanea: {str(e)}")
    
    @error_handler
    def process_email_attachment(self, attachment_data, attachment_name, content_type):
        """Process an email attachment and extract its content"""
        temp_dir = None
        temp_file_path = None
        
        try:
            import tempfile
            import os

            # Log dettagliato per il debug
            self.log_debug(f"Inizio elaborazione allegato: {attachment_name} ({len(attachment_data)} bytes)")
            
            # Sanitizza il nome del file per evitare caratteri problematici
            safe_name = ''.join(c for c in attachment_name if c.isalnum() or c in '._- ')
            if not safe_name:
                safe_name = f"attachment_{hash(attachment_data) % 10000}.bin"
            
            # Se il nome file non ha estensione ma abbiamo un content_type, aggiungila
            if '.' not in safe_name and content_type:
                if content_type == 'application/pdf':
                    safe_name += '.pdf'
                elif content_type == 'application/msword':
                    safe_name += '.doc'
                elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    safe_name += '.docx'
                elif content_type == 'application/vnd.ms-excel':
                    safe_name += '.xls'
                elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                    safe_name += '.xlsx'
                elif content_type == 'text/csv':
                    safe_name += '.csv'
                elif content_type == 'application/vnd.oasis.opendocument.text':
                    safe_name += '.odt'
                elif content_type == 'application/rtf' or content_type == 'text/rtf':
                    safe_name += '.rtf'
                elif content_type == 'application/xml' or content_type == 'text/xml':
                    safe_name += '.xml'
                elif content_type == 'text/plain':
                    safe_name += '.txt'
            
            # Crea il file temporaneo
            temp_dir = tempfile.mkdtemp(prefix="email_att_")
            temp_file_path = os.path.join(temp_dir, safe_name)
            
            # Salva l'allegato su file temporaneo
            with open(temp_file_path, 'wb') as f:
                f.write(attachment_data)
            
            self.log_debug(f"File temporaneo creato: {temp_file_path}")
            
            # Ottieni l'estensione del file
            _, ext = os.path.splitext(temp_file_path)
            ext = ext.lower()
            
            # IMPORTANTE: Verifica se l'estensione è supportata - log dettagliato
            supports_content_search = self.should_search_content(temp_file_path)
            self.log_debug(f"L'estensione {ext} supporta la ricerca nei contenuti? {supports_content_search}")
            
            # Force include common document types regardless of extension settings
            force_include = ext in ['.pdf', '.doc', '.docx', '.xls', '.xlsx', '.csv', '.odt', '.rtf', '.xml', '.txt']
            
            # Check if we should process this type of file
            if supports_content_search or force_include:
                # Skip nested email files to prevent recursion
                if ext == '.eml':
                    self.log_debug(f"Allegato email rilevato, evito ricorsione: {temp_file_path}")
                    return "[Contenuto allegato email non elaborato per evitare ricorsione]"
                
                # For supported formats, use the existing content extraction
                try:
                    self.log_debug(f"Inizio estrazione contenuto da {safe_name} ({ext})")
                    
                    # Gestione specifica per tipo di file
                    content = ""
                    
                    # File di testo semplice
                    if ext in ['.txt', '.csv', '.log', '.ini', '.xml', '.json', '.md']:
                        try:
                            with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                content = f.read()
                            self.log_debug(f"Letti {len(content)} caratteri da file di testo {ext}")
                        except Exception as text_err:
                            self.log_debug(f"Errore lettura testo: {str(text_err)}")
                    
                    # PDF
                    elif ext == '.pdf':
                        try:
                            import PyPDF2
                            with open(temp_file_path, 'rb') as f:
                                try:
                                    reader = PyPDF2.PdfReader(f)
                                    pdf_text = []
                                    for page_num in range(len(reader.pages)):
                                        page = reader.pages[page_num]
                                        pdf_text.append(page.extract_text())
                                    content = "\n".join(pdf_text)
                                    self.log_debug(f"Estratti {len(content)} caratteri da PDF")
                                except Exception as pdf_err:
                                    self.log_debug(f"Errore PDF: {str(pdf_err)}")
                        except ImportError:
                            self.log_debug("PyPDF2 non disponibile")
                    
                    # Word DOCX
                    elif ext == '.docx':
                        try:
                            import docx
                            doc = docx.Document(temp_file_path)
                            doc_text = []
                            for para in doc.paragraphs:
                                if para.text:
                                    doc_text.append(para.text)
                            content = "\n".join(doc_text)
                            self.log_debug(f"Estratti {len(content)} caratteri da DOCX")
                        except ImportError:
                            self.log_debug("python-docx non disponibile")
                        except Exception as docx_err:
                            self.log_debug(f"Errore DOCX: {str(docx_err)}")
                    
                    # Excel XLSX - NUOVO
                    elif ext == '.xlsx':
                        try:
                            import openpyxl
                            wb = openpyxl.load_workbook(temp_file_path, read_only=True, data_only=True)
                            xlsx_texts = []
                            
                            for sheet_name in wb.sheetnames:
                                sheet = wb[sheet_name]
                                xlsx_texts.append(f"--- Foglio: {sheet_name} ---")
                                
                                # Max 500 righe e 50 colonne per evitare problemi con file molto grandi
                                max_row = min(sheet.max_row, 500) if sheet.max_row else 0
                                max_col = min(sheet.max_column, 50) if sheet.max_column else 0
                                
                                for row_idx in range(1, max_row + 1):
                                    row_values = []
                                    for col_idx in range(1, max_col + 1):
                                        cell = sheet.cell(row=row_idx, column=col_idx)
                                        if cell.value:
                                            row_values.append(str(cell.value))
                                    if row_values:
                                        xlsx_texts.append(" ".join(row_values))
                                        
                            content = "\n".join(xlsx_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da XLSX")
                        except ImportError:
                            self.log_debug("Libreria openpyxl non disponibile")
                        except Exception as xlsx_err:
                            self.log_debug(f"Errore XLSX: {str(xlsx_err)}")
                    
                    # Excel XLS (vecchio formato) - NUOVO
                    elif ext == '.xls':
                        try:
                            import xlrd
                            book = xlrd.open_workbook(temp_file_path)
                            xls_texts = []
                            
                            for sheet_idx in range(book.nsheets):
                                sheet = book.sheet_by_index(sheet_idx)
                                xls_texts.append(f"--- Foglio: {sheet.name} ---")
                                
                                # Limita il numero di righe per prestazioni
                                for row_idx in range(min(sheet.nrows, 500)):
                                    row_values = sheet.row_values(row_idx)
                                    row_texts = [str(value) for value in row_values if value]
                                    if row_texts:
                                        xls_texts.append(" ".join(row_texts))
                                        
                            content = "\n".join(xls_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da XLS")
                        except ImportError:
                            self.log_debug("Libreria xlrd non disponibile")
                        except Exception as xls_err:
                            self.log_debug(f"Errore XLS: {str(xls_err)}")
                    
                    # CSV - NUOVO
                    elif ext == '.csv':
                        try:
                            import csv
                            csv_texts = []
                            
                            with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                reader = csv.reader(f)
                                # Limita a 1000 righe per file grandi
                                for i, row in enumerate(reader):
                                    if i >= 1000:
                                        csv_texts.append("... (file troncato, troppe righe)")
                                        break
                                    csv_texts.append("\t".join(row))
                                    
                            content = "\n".join(csv_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da CSV")
                        except Exception as csv_err:
                            self.log_debug(f"Errore CSV: {str(csv_err)}")
                    
                    # OpenDocument Text (ODT) - NUOVO
                    elif ext == '.odt':
                        try:
                            from odf import opendocument, text
                            
                            doc = opendocument.load(temp_file_path)
                            paragraphs = []
                            
                            # Estrai tutto il testo dai paragrafi
                            for element in doc.getElementsByType(text.P):
                                if element.firstChild:
                                    paragraphs.append(element.firstChild.data)
                                else:
                                    paragraphs.append("")
                                    
                            content = "\n".join(paragraphs)
                            self.log_debug(f"Estratti {len(content)} caratteri da ODT")
                        except ImportError:
                            self.log_debug("Libreria odf non disponibile")
                        except Exception as odt_err:
                            self.log_debug(f"Errore ODT: {str(odt_err)}")
                    
                    # Rich Text Format (RTF) - NUOVO
                    elif ext == '.rtf':
                        try:
                            from striprtf.striprtf import rtf_to_text
                            
                            with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                rtf_content = f.read()
                                content = rtf_to_text(rtf_content)
                                self.log_debug(f"Estratti {len(content)} caratteri da RTF")
                        except ImportError:
                            self.log_debug("Libreria striprtf non disponibile")
                        except Exception as rtf_err:
                            self.log_debug(f"Errore RTF: {str(rtf_err)}")
                    
                    # XML - NUOVO
                    elif ext == '.xml':
                        try:
                            import xml.etree.ElementTree as ET
                            
                            tree = ET.parse(temp_file_path)
                            root = tree.getroot()
                            
                            # Funzione ricorsiva per estrarre tutti i testi
                            def extract_text_from_element(element):
                                result = []
                                if element.text and element.text.strip():
                                    result.append(element.text.strip())
                                for child in element:
                                    result.extend(extract_text_from_element(child))
                                if element.tail and element.tail.strip():
                                    result.append(element.tail.strip())
                                return result
                            
                            xml_texts = extract_text_from_element(root)
                            content = "\n".join(xml_texts)
                            self.log_debug(f"Estratti {len(content)} caratteri da XML")
                        except Exception as xml_err:
                            self.log_debug(f"Errore XML: {str(xml_err)}")
                            # Fallback a lettura semplice del file
                            try:
                                with open(temp_file_path, 'r', encoding='utf-8', errors='replace') as f:
                                    content = f.read()
                                    self.log_debug(f"XML letto come testo semplice: {len(content)} caratteri")
                            except:
                                pass
                    
                    # Se è un altro tipo di file, usa get_file_content
                    else:
                        try:
                            content = self.get_file_content(temp_file_path)
                            self.log_debug(f"Contenuto estratto via get_file_content: {len(content)} caratteri")
                        except Exception as get_file_err:
                            self.log_debug(f"Errore get_file_content: {str(get_file_err)}")
                    
                    if content and len(content) > 0:
                        self.log_debug(f"Estratti con successo {len(content)} caratteri dall'allegato {safe_name}")
                        return content
                    else:
                        self.log_debug(f"Nessun contenuto estratto dall'allegato {safe_name}")
                        return f"[Allegato {safe_name}: nessun contenuto estratto]"
                except Exception as e:
                    self.log_debug(f"Errore nell'estrazione del contenuto dell'allegato {safe_name}: {str(e)}")
                    return f"[Errore nell'estrazione: {str(e)}]"
            else:
                self.log_debug(f"Allegato {safe_name} non elaborabile secondo le impostazioni attuali")
                return f"[Allegato {safe_name} ({ext}) non elaborabile]"
            
        except Exception as e:
            self.log_debug(f"Errore generale nell'elaborazione dell'allegato {attachment_name}: {str(e)}")
            return ""
        finally:
            # Clean up temp directory
            if temp_dir and os.path.exists(temp_dir):
                try:
                    import shutil
                    shutil.rmtree(temp_dir)
                    self.log_debug(f"Directory temporanea rimossa: {temp_dir}")
                except Exception as e:
                    self.log_debug(f"Errore nella pulizia dei file temporanei: {str(e)}")

    @error_handler
    def _partial_content_search(self, file_path, keywords):
        """Esegue una ricerca parziale in un file molto grande"""
        try:
            self.log_debug(f"Inizio analisi parziale per file gigantesco: {os.path.basename(file_path)}")
            file_size = os.path.getsize(file_path)
            
            # Definisci quanto analizzare all'inizio e alla fine
            head_size = 20 * 1024 * 1024  # Primi 20 MB
            tail_size = 10 * 1024 * 1024  # Ultimi 10 MB
            
            # Adatta la dimensione all'analisi in base al file
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # Per file di log, analizza più contenuto all'inizio (eventi recenti)
            if file_ext in ['.log', '.txt']:
                head_size = 30 * 1024 * 1024  # 30 MB all'inizio
                tail_size = 5 * 1024 * 1024   # 5 MB alla fine
            
            # Per database, analizza più alla fine (record più recenti)
            if file_ext in ['.db', '.sqlite', '.mdb']:
                head_size = 10 * 1024 * 1024  # 10 MB all'inizio
                tail_size = 25 * 1024 * 1024  # 25 MB alla fine
            
            # Crea buffer vuoti per contenere i dati
            head_data = b''
            tail_data = b''
            
            with open(file_path, 'rb') as f:
                # Leggi l'inizio del file
                head_data = f.read(min(head_size, file_size))
                
                # Vai alla fine meno tail_size se il file è abbastanza grande
                if file_size > head_size + tail_size:
                    f.seek(max(head_size, file_size - tail_size))
                    # Leggi la fine del file
                    tail_data = f.read(min(tail_size, file_size - head_size))
                
            # Converti i dati binari in testo (gestisce errori di decodifica)
            encodings = ['utf-8', 'latin-1', 'windows-1252']
            head_text = ""
            tail_text = ""
            
            # Prova diverse codifiche per decodificare il contenuto
            for encoding in encodings:
                try:
                    head_text = head_data.decode(encoding, errors='ignore')
                    if tail_data:
                        tail_text = tail_data.decode(encoding, errors='ignore')
                    break
                except:
                    continue
            
            # Se nessuna codifica ha funzionato, usa l'ultima con ignore
            if not head_text:
                head_text = head_data.decode('utf-8', errors='ignore')
            if tail_data and not tail_text:
                tail_text = tail_data.decode('utf-8', errors='ignore')
            
            # Se il file è piccolo o la seconda lettura non ha ottenuto dati
            if not tail_text or file_size <= head_size + tail_size:
                combined_text = head_text
            else:
                # Unisci i testi con un indicatore che mostra che è un'analisi parziale
                combined_text = head_text + "\n[...CONTENUTO INTERMEDIO NON ANALIZZATO...]\n" + tail_text
            
            # Cerca le keywords nel testo combinato
            for keyword in keywords:
                # Verifica corrispondenza intera parola se richiesto
                if self.whole_word_search.get():
                    if self.is_whole_word_match(keyword, combined_text):
                        self.log_debug(f"Match trovato in analisi parziale di file gigantesco: {os.path.basename(file_path)}")
                        return True
                # Altrimenti cerca match normale
                elif keyword.lower() in combined_text.lower():
                    self.log_debug(f"Match trovato in analisi parziale di file gigantesco: {os.path.basename(file_path)}")
                    return True
            
            self.log_debug(f"Nessun match trovato in analisi parziale di file gigantesco: {os.path.basename(file_path)}")
            return False
            
        except Exception as e:
            self.log_error(f"Errore durante l'analisi parziale del file {file_path}", e)
            return False

    @error_handler
    def _mark_file_for_partial_analysis(self, file_path):
        """Marca un file per l'analisi parziale invece che completa"""
        if not hasattr(self, '_partial_analysis_files'):
            self._partial_analysis_files = set()
        self._partial_analysis_files.add(file_path)
        self.log_debug(f"File marcato per analisi parziale: {os.path.basename(file_path)}")

    @error_handler
    def _queue_gigantic_file_for_confirmation(self, file_path):
        """Accoda un file gigantesco per la conferma dell'utente"""
        if not hasattr(self, '_gigantic_files_queue'):
            self._gigantic_files_queue = []
        if not hasattr(self, '_gigantic_files_confirmed'):
            self._gigantic_files_confirmed = set()
        
        if file_path not in self._gigantic_files_queue and file_path not in self._gigantic_files_confirmed:
            self._gigantic_files_queue.append(file_path)
            self.log_debug(f"File gigantesco accodato per conferma: {os.path.basename(file_path)}")
            
            # Schedula un prompt per l'utente se non già in corso
            if not hasattr(self, '_showing_gigantic_confirmation') or not self._showing_gigantic_confirmation:
                self._showing_gigantic_confirmation = True
                if self.root:  # Aggiungi controllo per evitare errori
                    self.root.after(100, self._process_gigantic_file_queue)

    @error_handler
    def _process_gigantic_file_queue(self):
        """Processa la coda di file giganteschi richiedendo conferma all'utente"""
        if not hasattr(self, '_gigantic_files_queue') or not self._gigantic_files_queue:
            self._showing_gigantic_confirmation = False
            return
        
        # Verifica se l'interfaccia è disponibile
        if not self.root or not self.root.winfo_exists():
            self._showing_gigantic_confirmation = False
            return
        
        file_path = self._gigantic_files_queue[0]
        file_name = os.path.basename(file_path)
        
        try:
            file_size = os.path.getsize(file_path)
            formatted_size = self._format_size(file_size)
            
            # Calcola il tempo stimato di elaborazione (approssimativo)
            estimated_time = self._estimate_processing_time(file_size)
            
            # Mostra un dialog all'utente
            message = (f"Il file '{file_name}' è molto grande ({formatted_size}).\n\n"
                    f"L'elaborazione potrebbe richiedere circa {estimated_time}.\n"
                    "Vuoi procedere con l'analisi?")
            
            from tkinter import messagebox
            result = messagebox.askyesno("File gigantesco rilevato", message)
            
            if result:  # L'utente ha confermato
                self._gigantic_files_confirmed.add(file_path)
                self.log_debug(f"Utente ha confermato l'analisi del file gigantesco: {file_name}")
                # Riesamina il file nella ricerca
                if hasattr(self, 'current_search_keywords') and self.current_search_keywords:
                    self.root.after(100, lambda: self.redo_search_for_file(file_path))
        except Exception as e:
            self.log_error(f"Errore nel processare la conferma per file gigantesco: {file_path}", e)
        
        # Rimuovi il file dalla coda
        if self._gigantic_files_queue:
            self._gigantic_files_queue.pop(0)
        
        # Se ci sono altri file nella coda, continua a processarli
        if self._gigantic_files_queue:
            self.root.after(100, self._process_gigantic_file_queue)
        else:
            self._showing_gigantic_confirmation = False

    @error_handler
    def _estimate_processing_time(self, file_size):
        """Stima il tempo di elaborazione per un file di grandi dimensioni"""
        # Valori empirici basati su test (questi andrebbero regolati in base alle prestazioni reali)
        bytes_per_second = 25 * 1024 * 1024  # ~25 MB/s per file di testo
        
        # Adatta la velocità in base alla dimensione (file più grandi sono più lenti da processare)
        if file_size > 5 * 1024 * 1024 * 1024:  # >5GB
            bytes_per_second = 15 * 1024 * 1024  # ~15 MB/s
        elif file_size > 3 * 1024 * 1024 * 1024:  # >3GB
            bytes_per_second = 20 * 1024 * 1024  # ~20 MB/s
        
        seconds = file_size / bytes_per_second
        
        # Formatta il tempo stimato
        if seconds < 60:
            return f"{int(seconds)} secondi"
        elif seconds < 3600:
            return f"{int(seconds / 60)} minuti"
        else:
            hours = int(seconds / 3600)
            minutes = int((seconds % 3600) / 60)
            return f"{hours} ore e {minutes} minuti"

    @error_handler
    def redo_search_for_file(self, file_path):
        """Riesamina un file specifico nella ricerca corrente"""
        if not hasattr(self, 'current_search_keywords') or not self.current_search_keywords:
            self.log_debug(f"Impossibile eseguire ricerca su file {os.path.basename(file_path)}: nessuna keyword impostata")
            return
        
        self.log_debug(f"Esecuzione ricerca sul file confermato: {os.path.basename(file_path)}")
        
        # Esegui la ricerca solo su questo file specifico
        try:
            # Rimuovi il file dai processati, se presente
            normalized_path = os.path.normpath(os.path.abspath(file_path))
            if hasattr(self, 'processed_files') and normalized_path in self.processed_files:
                self.processed_files.remove(normalized_path)
            
            # Process file può restituire None se non trova match
            result = self.process_file(file_path, self.current_search_keywords, search_content=True)
            if result:
                if not hasattr(self, 'search_results'):
                    self.search_results = []
                self.search_results.append(result)
                self.log_debug(f"Match trovato nel file gigantesco confermato: {os.path.basename(file_path)}")
                if hasattr(self, 'update_results_list'):
                    self.update_results_list()
        except Exception as e:
            self.log_error(f"Errore durante l'elaborazione del file {file_path}", e)

    @error_handler
    def update_progress(self):
        if self.is_searching:
            try:
                # Process only a limited number of messages per cycle
                max_messages_per_cycle = 15  # Reduced from 20 to 15 for better responsiveness
                messages_processed = 0
                
                start_time = time.time()
                # Prevent processing for too long
                max_processing_time = 0.03  # Reduced from 0.05 to 0.03 seconds for more responsive UI
                
                while messages_processed < max_messages_per_cycle:
                    try:
                        # Processa tutti i messaggi nella coda
                        progress_type, value = self.progress_queue.get_nowait()
                        
                        # CORREZIONE: Ottimizza l'elaborazione dei messaggi di stato
                        if progress_type == "update_total_time":
                            if hasattr(self, 'total_time_label') and self.total_time_label.winfo_exists():
                                self.total_time_label.config(text=value)
                        elif progress_type == "progress":
                            if hasattr(self, 'progress_bar') and self.progress_bar.winfo_exists():
                                current_value = self.progress_bar["value"]
                                # Evita aggiornamenti inutili se il valore non è cambiato significativamente
                                if abs(current_value - value) >= 1:
                                    self.progress_bar["value"] = value
                        elif progress_type == "status":
                            # CORREZIONE: Migliora la gestione dei messaggi di stato
                            if ("analizzati" in value.lower() or "cartelle:" in value.lower() or 
                                "file:" in value.lower()) and hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                # Estrai il percorso se presente
                                if "analisi:" in value.lower():
                                    parts = value.split("(", 1)
                                    if len(parts) > 1:
                                        path_part = parts[0].strip()
                                        counts_part = "(" + parts[1]
                                        if hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                            self.status_label["text"] = path_part
                                        if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                            self.analyzed_files_label["text"] = counts_part
                                    else:
                                        if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                            self.analyzed_files_label["text"] = value
                                else:
                                    # È solo un messaggio di stato semplice
                                    if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                        self.analyzed_files_label["text"] = value
                            elif hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                # CORREZIONE: Evita aggiornamenti ripetitivi dello stesso messaggio
                                if not hasattr(self, 'last_status_message') or self.last_status_message != value:
                                    self.status_label["text"] = value
                                    self.last_status_message = value
                        elif progress_type == "update_dir_size":
                            # CORREZIONE: Aggiorna solo se abilitato e necessario
                            calculation_mode = self.dir_size_calculation.get()
                            if hasattr(self, 'dir_size_var') and calculation_mode != "disabilitato":
                                current_text = self.dir_size_var.get()
                                new_text = self._format_size(value)
                                # Aggiorna solo se il valore è cambiato significativamente
                                if current_text != new_text:
                                    self.dir_size_var.set(new_text)
                            elif hasattr(self, 'dir_size_var') and calculation_mode == "disabilitato":
                                self.dir_size_var.set("Calcolo disattivato")
                        elif progress_type == "complete":
                            self.is_searching = False
                            self.enable_all_controls()
                            if hasattr(self, 'stop_button') and self.stop_button.winfo_exists():
                                self.stop_button["state"] = "disabled"
                            
                            # CORREZIONE: Aggiungi un breve ritardo prima di aggiornare i risultati
                            # per garantire che l'interfaccia sia reattiva
                            self.root.after(100, self.update_results_list)
                            
                            # Aggiorna il tempo finale
                            if hasattr(self, 'search_start_time') and self.search_start_time:
                                # Imposta il timestamp di fine ricerca
                                self.search_end_time = datetime.now()
                                current_time = self.search_end_time.strftime('%H:%M')
                                
                                if hasattr(self, 'end_time_label') and self.end_time_label.winfo_exists():
                                    self.end_time_label.config(text=current_time)
                                    
                                # Aggiorna il tempo totale
                                self.update_total_time()
                                
                                # Aggiorna il debug log
                                self.log_debug(f"Ricerca completata. Trovati {len(self.search_results)} risultati")
                        
                        messages_processed += 1
                        
                        # Check if we've been processing for too long
                        if time.time() - start_time > max_processing_time:
                            break
                            
                    except queue.Empty:
                        break
                    except tk.TclError as e:
                        # Widget non più esistente o distrutto - non grave, ignora e continua
                        self.log_debug(f"Widget non più disponibile: {str(e)}")
                        continue
                    except Exception as e:
                        # Log altri errori ma non interrompere l'aggiornamento
                        self.log_debug(f"Errore durante l'elaborazione messaggio: {str(e)}")
                        continue
                
                # Force UI update after processing messages
                try:
                    self.root.update_idletasks()
                except:
                    pass
                
                # CORREZIONE: Adatta meglio la frequenza di aggiornamento
                if hasattr(self.progress_queue, 'qsize'):
                    queue_size = self.progress_queue.qsize()
                    if queue_size > 100:
                        # Many messages in queue, update more frequently
                        self.root.after(20, self.update_progress)  # 20ms (50Hz)
                    elif queue_size > 20:
                        # Moderate queue, moderate update rate
                        self.root.after(50, self.update_progress)  # 50ms (20Hz)
                    else:
                        # Few messages, standard update rate
                        self.root.after(100, self.update_progress)  # 100ms (10Hz)
                else:
                    # No qsize method, use standard rate
                    self.root.after(100, self.update_progress)
                    
            except tk.TclError as e:
                self.log_debug(f"TclError nell'aggiornamento del progresso: {str(e)}")
                self.root.after(500, self.update_progress)
            except Exception as e:
                self.log_debug(f"Errore nell'aggiornamento del progresso: {str(e)}")
                self.root.after(500, self.update_progress)
    
    @error_handler
    def reset_search_state(self):
        """Reimpostazione completa dello stato di ricerca per garantire coerenza"""
        self.log_debug("Reimpostazione dello stato di ricerca")
        
        # Variabili di controllo principali
        self.is_searching = False
        self.stop_search = False
        
        # Ferma il monitoraggio della memoria
        self.stop_memory_monitoring()

        # Resetta l'elenco dei file processati
        self.processed_files = set()
        # Elimina flag temporanei di interruzione se presenti
        if hasattr(self, '_stopping_in_progress'):
            delattr(self, '_stopping_in_progress')
        
        # Chiudi executor se ancora esistente
        if hasattr(self, 'search_executor') and self.search_executor:
            try:
                self.search_executor.shutdown(wait=False)
            except:
                pass
            self.search_executor = None
        
        # Ferma il watchdog
        self.watchdog_active = False
        
        # Svuota le code
        try:
            while not self.progress_queue.empty():
                try:
                    self.progress_queue.get_nowait()
                    self.progress_queue.task_done()
                except:
                    break
        except:
            pass
        
        # Aggiorna UI per riflettere lo stato corretto
        if hasattr(self, 'stop_button'):
            self.stop_button["state"] = "disabled"
        
        self.log_debug("Stato di ricerca reimpostato correttamente")

    @error_handler
    def stop_search_process(self):
        """Ferma il processo di ricerca in corso in modo aggressivo"""
        # Verificare se una ricerca è effettivamente in corso
        if not self.is_searching:
            self.log_debug("Tentativo di interruzione ma nessuna ricerca in corso")
            return
        
        self.log_debug("INTERRUZIONE: Stop ricerca richiesto dall'utente")
        
        # 1. Imposta i flag di interruzione
        self.stop_search = True
        self.is_searching = False
        
        # Ferma il monitoraggio della memoria
        self.stop_memory_monitoring()

        # 2. Aggiorna l'interfaccia
        self.status_label["text"] = "Interruzione ricerca in corso... attendere"
        self.analyzed_files_label["text"] = "Ricerca interrotta dall'utente"
        self.progress_bar["value"] = 100
        
        # 3. Cattura il tempo di fine
        self.search_end_time = datetime.now()
        current_time = self.search_end_time.strftime('%H:%M')
        self.end_time_label.config(text=current_time)
        self.update_total_time()  # Aggiorna il tempo totale
        
        # 4. CRITICO: Ferma tutti i thread in modo aggressivo
        if hasattr(self, 'search_executor') and self.search_executor:
            try:
                self.log_debug("INTERRUZIONE: Chiusura dell'executor...")
                
                # Ferma l'executor nel modo più aggressivo possibile
                try:
                    import sys
                    if sys.version_info >= (3, 9):
                        self.search_executor.shutdown(wait=False, cancel_futures=True)
                    else:
                        self.search_executor.shutdown(wait=False)
                        
                    # Crea un nuovo executor per le future ricerche
                    self.search_executor = None
                    
                except Exception as e:
                    self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                    self.search_executor = None  # Forza nullificazione
                
            except Exception as e:
                self.log_debug(f"Errore generale nell'interruzione: {str(e)}")
                self.search_executor = None
        
        # 5. Visualizza risultati parziali trovati
        self.root.after(500, self.update_results_list)
        
        # 6. Riabilita l'interfaccia utente dopo l'interruzione
        self.root.after(1000, self._complete_interrupt_process)
        
        # 7. Forza aggiornamento GUI
        try:
            self.root.update_idletasks()
        except:
            pass
        
        self.log_debug("INTERRUZIONE: Processo di interruzione ricerca completato")
    
    @error_handler
    def _complete_interrupt_process(self):
        """Completa il processo di interruzione ripristinando l'interfaccia"""
        try:
            # Riabilita tutti i controlli
            self.enable_all_controls()
            
            # Reimpostazione completa dello stato
            self.reset_search_state()
            
            # Calcola la dimensione totale dei file trovati fino all'interruzione
            self.update_total_files_size()

            self.status_label["text"] = "Ricerca interrotta dall'utente"
            self.log_debug("Interfaccia ripristinata dopo interruzione")
        except Exception as e:
            self.log_debug(f"Errore nel completamento dell'interruzione: {str(e)}")

    @error_handler
    def update_results_list(self):
        """Aggiorna la lista dei risultati con i risultati trovati"""
        # Aggiorna i colori del tema prima di aggiornare la lista (garantisce i colori corretti)
        self.update_theme_colors()
        
        # Pulisci la lista attuale
        for item in self.results_list.get_children():
            self.results_list.delete(item)
        
        attachment_count = 0  # Counter per debug
        non_attachment_items = []  # Lista per tenere traccia di elementi che non sono allegati
        
        # Aggiungi i risultati alla lista
        for result in self.search_results:
            # Verifica se il risultato è un dizionario o una tupla/lista
            if isinstance(result, dict):
                # Formato dizionario (nuovo formato)
                item_type = result.get("type", "File")
                author = result.get("author", "")
                size = result.get("size", "0 B")
                modified = result.get("modified", "")
                created = result.get("created", "")
                path = result.get("path", "")
                from_attachment = result.get("is_attachment", False)
            else:
                # Formato tupla/lista (vecchio formato)
                if len(result) >= 7:
                    item_type, author, size, modified, created, path, from_attachment = result
                else:
                    item_type, author, size, modified, created, path = result
                    from_attachment = False
            
            # Imposta l'icona della graffetta nella colonna dedicata
            attachment_icon = "📎" if from_attachment else ""
            
            # Applica stile in base al tipo di elemento con priorità per gli allegati
            if from_attachment:
                tags = ("attachment",)  # Tag speciale per gli allegati
                attachment_count += 1
                self.log_debug(f"Allegato trovato ({attachment_count}): {path}")
            elif item_type == "Directory":
                tags = ("directory",)
            else:
                tags = ("file",)
            
            display_values = (
                item_type,        # Tipo
                attachment_icon,  # Icona allegato
                size,             # Dimensione
                modified,         # Data modifica
                created,          # Data creazione
                author,           # Nome/Autore
                path              # Percorso
            )
                
            # Inserisci l'elemento nella TreeView con i tag appropriati
            item_id = self.results_list.insert("", "end", values=display_values, tags=tags)
            
            # Memorizza solo gli elementi che NON sono allegati
            if not from_attachment:
                non_attachment_items.append(item_id)
        
        # Aggiorna lo stato
        self.status_label["text"] = f"Trovati {len(self.search_results)} risultati"
        if attachment_count > 0:
            self.status_label["text"] += f" (inclusi {attachment_count} allegati)"
        
        # Aggiorna la dimensione totale dei file trovati
        self.update_total_files_size()
        
        # Se ci sono risultati NON allegati, seleziona il primo di questi
        if non_attachment_items:
            self.results_list.selection_set(non_attachment_items[0])
            self.results_list.focus(non_attachment_items[0])

    
    @error_handler
    def update_total_files_size(self):
        """Calcola e aggiorna la dimensione totale dei file trovati"""
        try:
            total_size = 0
            file_count = 0
            
            # Output debug info
            self.log_debug(f"Calculating total size for {len(self.search_results)} results")
            
            # Calcola la dimensione totale dai risultati
            for result in self.search_results:
                try:
                    # Handle both 6-element and 7-element results
                    if len(result) >= 7:  # If it includes from_attachment flag
                        item_type, _, size_str, _, _, _, _ = result
                    else:
                        item_type, _, size_str, _, _, _ = result
                    
                    # Conta solo i file, non le directory
                    if item_type != "Directory":
                        file_count += 1
                        
                        # Estrai il valore numerico dalla stringa della dimensione
                        if size_str and isinstance(size_str, str):
                            if 'KB' in size_str:
                                size_value = float(size_str.split()[0]) * 1024
                            elif 'MB' in size_str:
                                size_value = float(size_str.split()[0]) * 1024 * 1024
                            elif 'GB' in size_str:
                                size_value = float(size_str.split()[0]) * 1024 * 1024 * 1024
                            else:
                                # Assume B or other unit
                                size_value = float(size_str.split()[0])
                            
                            total_size += size_value
                except Exception as e:
                    self.log_debug(f"Error processing file size: {str(e)}")
                    continue
            
            # Formatta la dimensione totale
            formatted_size = self._format_size(total_size)
            
            # Aggiorna il label con la dimensione totale e il numero di file
            if hasattr(self, 'total_files_size_label') and self.total_files_size_label.winfo_exists():
                new_text = f"Dimensione totale: {formatted_size} ({file_count} file)"
                self.log_debug(f"Updating size label to: {new_text}")
                self.total_files_size_label.config(text=new_text)
                
                # Forza aggiornamento dell'interfaccia
                self.root.update_idletasks()
        except Exception as e:
            self.log_debug(f"Error in update_total_files_size: {str(e)}")

    @error_handler
    def update_selected_files_size(self, event=None):
        """Calcola e visualizza la dimensione totale dei file selezionati"""
        # Ottieni SOLO gli elementi esplicitamente selezionati dall'utente
        selected_items = self.results_list.selection()
        
        total_size = 0
        file_count = 0
        
        for item in selected_items:
            values = self.results_list.item(item)['values']
            # Verifica che sia un file (non una directory)
            if values and values[0] != "Directory":
                file_count += 1
                size_str = values[2]  # La dimensione è nella terza colonna (indice 2)
                
                # Estrai il valore numerico dalla stringa della dimensione
                if size_str and isinstance(size_str, str):
                    try:
                        if 'TB' in size_str:
                            size_value = float(size_str.split()[0]) * 1024 * 1024 * 1024 * 1024
                        elif 'GB' in size_str:
                            size_value = float(size_str.split()[0]) * 1024 * 1024 * 1024
                        elif 'MB' in size_str:
                            size_value = float(size_str.split()[0]) * 1024 * 1024
                        elif 'KB' in size_str:
                            size_value = float(size_str.split()[0]) * 1024
                        else:
                            # Assume B o altra unità
                            size_value = float(size_str.split()[0])
                        
                        total_size += size_value
                    except:
                        pass  # Ignora errori di parsing
        
        # Formatta la dimensione e aggiorna l'etichetta
        if file_count > 0:
            formatted_size = self._format_size(total_size)
            self.selected_files_size_label.config(text=f"Selezionati: {formatted_size} ({file_count} file)")
        else:
            self.selected_files_size_label.config(text="Selezionati: 0 (0 file)")

    @error_handler
    def update_theme_colors(self, theme=None):
        """Aggiorna i colori del tema per evidenziare cartelle e file"""
        if theme is None:
            # Se non viene specificato un tema, usa quello attualmente selezionato
            if hasattr(self, 'theme_combobox'):
                theme = self.theme_combobox.get()
            else:
                theme = "darkly"  # Predefinito
        
        style = ttk.Style()
        
        # Definisci i colori di sfondo per ogni tema
        bg_colors = {
            "minty": "#f8f9fa",
            "cosmo": "#f8f9fa",
            "darkly": "#222",
            "cyborg": "#060606"
        }
        
        # Ottieni il colore di sfondo corrente
        bg_color = bg_colors.get(theme, "#222")  # Default se il tema non è nella lista
        
        # Configura i colori in base al tema
        if theme in ["minty", "cosmo"]:  # Temi chiari
            # Usa lo sfondo dell'applicazione anche per la treeview
            style.configure("Treeview", background=bg_color, foreground="#000000", fieldbackground=bg_color)
            if hasattr(self, 'results_list'):
                self.results_list.tag_configure("directory", background="#e6f2ff", foreground="#000000")
                self.results_list.tag_configure("file", background=bg_color, foreground="#000000")
                self.results_list.tag_configure("attachment", background="#f8f8e0", foreground="#000000")
            
            # Configura la finestra di debug se esiste
            if hasattr(self, 'debug_log_text'):
                self.debug_log_text.configure(background=bg_color, foreground="#000000")
        
        elif theme in ["darkly", "cyborg"]:  # Temi scuri
            # Usa lo sfondo dell'applicazione anche per la treeview
            style.configure("Treeview", background=bg_color, foreground="#ffffff", fieldbackground=bg_color)
            
            # Configura i colori specifici per ogni tema
            if theme == "darkly":
                if hasattr(self, 'results_list'):
                    self.results_list.tag_configure("directory", background="#303030", foreground="#ffffff")
                    self.results_list.tag_configure("file", background=bg_color, foreground="#ffffff")
                    self.results_list.tag_configure("attachment", background="#a0a080", foreground="#ffffff")
                    
            elif theme == "cyborg":
                if hasattr(self, 'results_list'):
                    self.results_list.tag_configure("directory", background="#181818", foreground="#2a9fd6")
                    self.results_list.tag_configure("file", background=bg_color, foreground="#ffffff")
                    self.results_list.tag_configure("attachment", background="#a0a080", foreground="#ffffff")
            
            # Configura la finestra di debug se esiste
            if hasattr(self, 'debug_log_text'):
                self.debug_log_text.configure(background=bg_color, foreground="#ffffff")
        
        # Aggiorna anche la selezione nella treeview
        style.map("Treeview",
                background=[('selected', '#375a7f' if theme in ["darkly", "cyborg"] else '#2780e3')],
                foreground=[('selected', '#ffffff')])

    @error_handler
    def copy_selected(self):
        """Copia i file selezionati in una directory di destinazione"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da copiare")
            return
            
        # Chiedi la directory di destinazione
        dest_dir = filedialog.askdirectory(title="Seleziona la cartella di destinazione")
        if not dest_dir:
            return
            
        # Prepara variabili per tracciare l'avanzamento
        total = len(selected_items)
        copied = 0
        failed = 0
        skipped = 0
        
        try:
            for item in selected_items:
                values = self.results_list.item(item)['values']
                item_type, _, _, _, _, source_path = values
                
                # Ottieni il nome dell'elemento senza il percorso completo
                basename = os.path.basename(source_path)
                dest_path = os.path.join(dest_dir, basename)
                
                try:
                    if item_type == "Directory":
                        # Se la cartella esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("Cartella esistente", 
                                                f"La cartella {basename} esiste già. Vuoi sovrascriverla?"):
                                # Elimina la cartella esistente
                                shutil.rmtree(dest_path)
                            else:
                                skipped += 1
                                continue
                                
                        # Copia ricorsiva della cartella
                        shutil.copytree(source_path, dest_path)
                        copied += 1
                    else:
                        # Se il file esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("File esistente", 
                                                f"Il file {basename} esiste già. Vuoi sovrascriverlo?"):
                                # Continua con la sovrascrittura
                                pass
                            else:
                                skipped += 1
                                continue
                                
                        # Copia il file
                        shutil.copy2(source_path, dest_path)
                        copied += 1
                        
                except Exception as e:
                    failed += 1
                    self.log_debug(f"Errore durante la copia di {source_path}: {str(e)}")
                    
                # Aggiorna la barra di avanzamento
                progress = ((copied + failed + skipped) / total) * 100
                self.progress_bar["value"] = progress
                self.status_label["text"] = f"Copiati {copied}/{total} elementi..."
                self.root.update()
                
            # Messaggio di completamento
            if failed > 0:
                messagebox.showwarning("Copia completata con errori", 
                                    f"Copiati: {copied}\nSaltati: {skipped}\nFalliti: {failed}")
            else:
                messagebox.showinfo("Copia completata", 
                                 f"Copiati con successo: {copied}\nSaltati: {skipped}")
                
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'operazione: {str(e)}")
            
        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    @error_handler
    def get_zip_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome del file ZIP"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Nome file ZIP")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome del file ZIP (senza estensione):", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = StringVar(value="archivio")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            zip_name = name_var.get().strip()
            if zip_name:
                result["name"] = zip_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per il file ZIP", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea ZIP", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]

    @error_handler
    def compress_selected(self):
        """Versione avanzata della compressione che mantiene la struttura originale delle directory"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da comprimere")
            return

        # Ottieni il nome del file ZIP
        zip_name = self.get_zip_name()
        if not zip_name:
            return

        # Ottieni il nome della cartella principale
        main_folder_name = self.get_main_folder_name()
        if main_folder_name is None:
            return  # L'utente ha annullato
        if not main_folder_name:
            main_folder_name = "files"  # Default se l'utente inserisce valore vuoto

        # Opzioni di compressione più dettagliate
        compression_dialog = tk.Toplevel(self.root)
        compression_dialog.title("Opzioni Compressione")
        compression_dialog.transient(self.root)
        compression_dialog.grab_set()

        frame = ttk.Frame(compression_dialog, padding=20)
        frame.pack(fill=tk.BOTH, expand=True)

        ttk.Label(frame, text="Scegli il tipo di compressione:").pack(anchor=tk.W, pady=(0,10))

        comp_var = tk.StringVar(value="nessuna")
        preserve_var = tk.BooleanVar(value=True)  # Nuova opzione per preservare la struttura

        ttk.Radiobutton(frame, text="Nessuna (solo archiviazione)", 
            variable=comp_var, value="nessuna").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Minima (massima velocità) Liv.1", 
                    variable=comp_var, value="minima").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Standard (buon equilibrio) Liv.6", 
                    variable=comp_var, value="standard").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Massima (compressione ottimale, più lenta) Liv.9", 
                    variable=comp_var, value="massima").pack(anchor=tk.W)

        # Frame per le opzioni di hash
        hash_frame = ttk.LabelFrame(frame, text="Calcolo hash del file ZIP")
        hash_frame.pack(fill=tk.X, pady=(10,5))

        ttk.Label(hash_frame, text="Seleziona gli algoritmi di hash da calcolare sul file compresso:").pack(anchor=tk.W, pady=(5,5))

        # Variabili per gli algoritmi di hash
        hash_md5 = tk.BooleanVar(value=True)
        hash_sha1 = tk.BooleanVar(value=False)
        hash_sha256 = tk.BooleanVar(value=False)

        # Checkbox per gli algoritmi
        ttk.Checkbutton(hash_frame, text="MD5 (più veloce)", variable=hash_md5).pack(anchor=tk.W, padx=20)
        ttk.Checkbutton(hash_frame, text="SHA1", variable=hash_sha1).pack(anchor=tk.W, padx=20)
        ttk.Checkbutton(hash_frame, text="SHA256 (più sicuro)", variable=hash_sha256).pack(anchor=tk.W, padx=20)

        ttk.Label(frame, text="\nOpzioni struttura:", font=("", 9, "bold")).pack(anchor=tk.W, pady=(10,5))

        ttk.Checkbutton(frame, text="Preserva la struttura delle directory", 
                    variable=preserve_var).pack(anchor=tk.W)
        ttk.Label(frame, text="Se attivato, mantiene l'alberatura originale delle cartelle", 
                    font=("", 8), foreground="gray").pack(anchor=tk.W, padx=(20, 5))

        ttk.Label(frame, text="\nOpzioni per file di grandi dimensioni:").pack(anchor=tk.W, pady=(10,5))

        use_chunks = tk.BooleanVar(value=False)
        ttk.Checkbutton(frame, text="Elabora in blocchi (per file molto grandi)", 
                    variable=use_chunks).pack(anchor=tk.W)

        btn_frame = ttk.Frame(frame)
        btn_frame.pack(fill=tk.X, pady=(15,0))

        result = {"option": None, "chunks": False, "preserve": True}

        def on_cancel():
            result["option"] = None
            compression_dialog.destroy()

        def on_ok():
            result["option"] = comp_var.get()
            result["chunks"] = use_chunks.get()
            result["preserve"] = preserve_var.get()
            
            # Aggiungi gli algoritmi selezionati
            result["hash_algorithms"] = []
            if hash_md5.get():
                result["hash_algorithms"].append("md5")
            if hash_sha1.get():
                result["hash_algorithms"].append("sha1")
            if hash_sha256.get():
                result["hash_algorithms"].append("sha256")
                
            compression_dialog.destroy()

        ttk.Button(btn_frame, text="Annulla", command=on_cancel).pack(side=tk.LEFT)
        ttk.Button(btn_frame, text="OK", command=on_ok).pack(side=tk.RIGHT)

        # Centra la finestra
        compression_dialog.update_idletasks()
        width = compression_dialog.winfo_reqwidth() + 50
        height = compression_dialog.winfo_reqheight() + 20
        x = (compression_dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (compression_dialog.winfo_screenheight() // 2) - (height // 2)
        compression_dialog.geometry(f"{width}x{height}+{x}+{y}")

        compression_dialog.wait_window()

        if result["option"] is None:
            return  # Utente ha annullato

        # Determina il metodo di compressione
        compression_level = 0  # Valore predefinito

        if result["option"] == "nessuna":
            compression_method = zipfile.ZIP_STORED
            compression_text = "nessuna"
        elif result["option"] == "minima":
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "minima"
            compression_level = 1  # Compressione veloce, riduzione minima
        elif result["option"] == "massima":
            # Prova LZMA se disponibile, altrimenti BZIP2, altrimenti fallback su DEFLATED con livello massimo
            if hasattr(zipfile, 'ZIP_LZMA'):
                compression_method = zipfile.ZIP_LZMA
                compression_text = "massima (LZMA)"
                compression_level = 9
            elif hasattr(zipfile, 'ZIP_BZIP2'):
                compression_method = zipfile.ZIP_BZIP2
                compression_text = "massima (BZIP2)"
                compression_level = 9
            else:
                compression_method = zipfile.ZIP_DEFLATED
                compression_text = "massima (DEFLATED)"
                compression_level = 9
        else:  # standard
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "standard"
            compression_level = 6  # Compressione bilanciata
            
        # Log della scelta di compressione
        self.log_debug(f"Utilizzo compressione {compression_text} (metodo: {compression_method}, livello: {compression_level})")
        self.log_debug(f"Preserva struttura directory: {result['preserve']}")

        zip_path = filedialog.asksaveasfilename(
            defaultextension=".zip",
            initialfile=f"{zip_name}.zip",
            filetypes=[("ZIP files", "*.zip")],
            title="Salva file ZIP"
        )

        if not zip_path:
            return

        # Determinare il percorso base per i calcoli relativi
        base_path = self._find_common_base_path(selected_items)
        self.log_debug(f"Percorso base per la struttura: {base_path}")

        # Raccogli tutti i file delle cartelle selezionate
        files_in_folders = set()
        folder_paths = []
        single_files = []

        # Prima fase: raccogli informazioni su cartelle e file
        for item in selected_items:
            values = self.results_list.item(item)['values']
            item_type, _, _, _, _, _, source_path = values

            if item_type == "Directory":
                folder_paths.append(source_path)
                # Raccogli tutti i file nelle cartelle selezionate
                for root, _, files in os.walk(source_path):
                    for file in files:
                        full_path = os.path.join(root, file)
                        files_in_folders.add(os.path.abspath(full_path))
            else:
                single_files.append(source_path)

        # Filtra i file singoli che sono già presenti nelle cartelle
        filtered_single_files = [f for f in single_files if os.path.abspath(f) not in files_in_folders]

        total_items = len(folder_paths) + len(filtered_single_files)
        processed = 0

        # Lista per tenere traccia dei file (tutti, non solo omonimi)
        all_files_log = []
        # Lista per tenere traccia dei file omonimi (come prima)
        omonimi_log = []

        # Ottieni informazioni sulla ricerca corrente per i log
        search_directory = self.search_path.get() if hasattr(self, 'search_path') else "N/A"
        search_keywords = self.keywords.get() if hasattr(self, 'keywords') else "N/A"

        try:
            # Primo passaggio: analizza tutti i file e identifica gli omonimi
            self.log_debug(f"Analisi dei file da comprimere nella cartella {main_folder_name}...")

            # Dizionario {nome_file: [lista_percorsi]}
            file_names_map = {}

            # Raccogli i nomi dei file dalle cartelle
            for folder_path in folder_paths:
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        file_name = os.path.basename(file_path)

                        # Aggiungi al log di tutti i file
                        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                        modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                        
                        all_files_log.append({
                            "nome_file": file_name,
                            "percorso_originale": file_path,
                            "dimensione": self._format_size(file_size),
                            "ultima_modifica": modified_time,
                            "tipo": self._get_file_type(file_path)
                        })

                        if file_name not in file_names_map:
                            file_names_map[file_name] = []
                        file_names_map[file_name].append(file_path)

            # Raccogli i nomi dei file singoli
            for file_path in filtered_single_files:
                if os.path.exists(file_path):
                    file_name = os.path.basename(file_path)
                    
                    # Aggiungi al log di tutti i file
                    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                    modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                    
                    all_files_log.append({
                        "nome_file": file_name,
                        "percorso_originale": file_path,
                        "dimensione": self._format_size(file_size),
                        "ultima_modifica": modified_time,
                        "tipo": self._get_file_type(file_path)
                    })
                    
                    if file_name not in file_names_map:
                        file_names_map[file_name] = []
                    file_names_map[file_name].append(file_path)

            # Identifica i file omonimi (con lo stesso nome ma percorsi diversi)
            omonimi_files = {name: paths for name, paths in file_names_map.items() if len(paths) > 1}
            self.log_debug(f"Trovati {len(omonimi_files)} file con nomi duplicati")
            
            # Per l'elaborazione a blocchi
            if result["chunks"] and total_items > 50:
                chunk_size = 10  # Elabora 10 file alla volta
                self.log_debug(f"Elaborazione a blocchi attivata: {chunk_size} file per blocco")

            # Secondo passaggio: crea l'archivio ZIP
            with zipfile.ZipFile(zip_path, 'w', compression_method, compresslevel=compression_level) as zipf:
                # Tiene traccia dei percorsi già aggiunti al file ZIP
                added_zip_paths = set()

                # Comprimi le cartelle
                for folder_path in folder_paths:
                    # Utilizzo la struttura originale o piatta in base alla scelta dell'utente
                    for root, _, files in os.walk(folder_path):
                        for file in files:
                            file_path = os.path.join(root, file)
                            file_name = os.path.basename(file_path)

                            try:
                                # Salta file non leggibili
                                if not os.access(file_path, os.R_OK):
                                    self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                    continue

                                # Calcola percorso relativo in base all'opzione selezionata
                                if result["preserve"]:
                                    # Calcola il percorso relativo rispetto alla base
                                    rel_path = os.path.relpath(file_path, base_path)
                                    zip_path_to_use = os.path.join(main_folder_name, rel_path)
                                else:
                                    # Comportamento originale senza preservare la struttura
                                    rel_path = os.path.relpath(file_path, os.path.dirname(folder_path))
                                    
                                    # Verifica se è un file omonimo
                                    if file_name in omonimi_files:
                                        # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                        if omonimi_files[file_name][0] == file_path:
                                            # Usa il percorso all'interno della cartella principale
                                            zip_path_to_use = os.path.join(main_folder_name, rel_path)
                                        else:
                                            # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                            parent_folder = os.path.basename(os.path.dirname(file_path))
                                            unique_name = f"{parent_folder}_{file_name}"

                                            # Se anche questo nome è duplicato, aggiungi un contatore
                                            counter = 1
                                            while os.path.join("omonimi", unique_name) in added_zip_paths:
                                                unique_name = f"{parent_folder}_{counter}_{file_name}"
                                                counter += 1

                                            zip_path_to_use = os.path.join("omonimi", unique_name)

                                            # Registra questo file nel log degli omonimi
                                            omonimi_log.append({
                                                "nome_file": file_name,
                                                "percorso_originale": file_path,
                                                "primo_percorso": omonimi_files[file_name][0],
                                                "posizione_zip": zip_path_to_use
                                            })
                                    else:
                                        # Non è un omonimo, inseriscilo nella cartella principale
                                        zip_path_to_use = os.path.join(main_folder_name, rel_path)

                                # Verifica se questo percorso ZIP è già stato usato (potrebbe accadere anche con la struttura preservata)
                                if zip_path_to_use in added_zip_paths:
                                    self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                    # Crea un nome alternativo
                                    alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)

                                    # Se anche questo nome è già usato, aggiungi un contatore
                                    counter = 1
                                    while zip_path_to_use in added_zip_paths:
                                        alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                        zip_path_to_use = os.path.join("omonimi", alt_name)
                                        counter += 1

                                # Aggiungi il file al ZIP e registra il percorso
                                zipf.write(file_path, zip_path_to_use)
                                added_zip_paths.add(zip_path_to_use)

                            except Exception as e:
                                self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Comprimi i file singoli
                for file_path in filtered_single_files:
                    if os.path.exists(file_path):
                        file_name = os.path.basename(file_path)

                        try:
                            # Salta file non leggibili
                            if not os.access(file_path, os.R_OK):
                                self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                continue

                            # Calcola percorso relativo in base all'opzione selezionata
                            if result["preserve"]:
                                # Calcola il percorso relativo rispetto alla base
                                rel_path = os.path.relpath(file_path, base_path)
                                zip_path_to_use = os.path.join(main_folder_name, rel_path)
                            else:
                                # Comportamento originale senza preservare la struttura
                                # Verifica se è un file omonimo
                                if file_name in omonimi_files:
                                    # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                    if omonimi_files[file_name][0] == file_path:
                                        # Usa il nome del file all'interno della cartella principale
                                        zip_path_to_use = os.path.join(main_folder_name, file_name)
                                    else:
                                        # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                        parent_folder = os.path.basename(os.path.dirname(file_path))
                                        unique_name = f"{parent_folder}_{file_name}"

                                        # Se anche questo nome è duplicato, aggiungi un contatore
                                        counter = 1
                                        while os.path.join("omonimi", unique_name) in added_zip_paths:
                                            unique_name = f"{parent_folder}_{counter}_{file_name}"
                                            counter += 1

                                        zip_path_to_use = os.path.join("omonimi", unique_name)

                                        # Registra questo file nel log degli omonimi
                                        omonimi_log.append({
                                            "nome_file": file_name,
                                            "percorso_originale": file_path,
                                            "primo_percorso": omonimi_files[file_name][0],
                                            "posizione_zip": zip_path_to_use
                                        })
                                else:
                                    # Non è un omonimo, inseriscilo nella cartella principale
                                    zip_path_to_use = os.path.join(main_folder_name, file_name)

                            # Verifica se questo percorso ZIP è già stato usato
                            if zip_path_to_use in added_zip_paths:
                                self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                # Crea un nome alternativo
                                alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                zip_path_to_use = os.path.join("omonimi", alt_name)

                                # Se anche questo nome è già usato, aggiungi un contatore
                                counter = 1
                                while zip_path_to_use in added_zip_paths:
                                    alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)
                                    counter += 1

                            # Aggiungi il file al ZIP e registra il percorso
                            zipf.write(file_path, zip_path_to_use)
                            added_zip_paths.add(zip_path_to_use)

                        except Exception as e:
                            self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Crea log in formato CSV con collegamenti ipertestuali
                current_time_local = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                current_user = getpass.getuser()

                # Identifica quali file sono omonimi per poterli escludere dal log generale
                omonimo_paths = set()
                for entry in omonimi_log:
                    omonimo_paths.add(entry['percorso_originale'])

                # Lista di file normali (esclude gli omonimi)
                normal_files_log = [f for f in all_files_log if f['percorso_originale'] not in omonimo_paths]

                # ---------- CREAZIONE DEL CSV DEI FILE NORMALI ----------
                # Creazione intestazione per il CSV con informazioni sulla ricerca
                csv_header = f"LOG DEI FILE TROVATI - {current_time_local}\n"
                csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                csv_header += f"Directory di ricerca: {search_directory}\n"
                csv_header += f"Parole chiave ricercate: {search_keywords}\n"
                csv_header += f"Opzione struttura: {'Preservata' if result['preserve'] else 'Piatta'}\n"
                csv_header += f"Totale file: {len(all_files_log)}\n\n"

                # Prepara l'output CSV
                csv_content = io.StringIO()
                csv_writer = csv.writer(csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)

                # Intestazione colonne
                csv_writer.writerow(["Nome file", "Tipo", "Dimensione", "Ultima modifica", "Percorso completo", "Link"])

                # Aggiungi i file normali (non omonimi)
                for file_entry in sorted(all_files_log, key=lambda x: x["nome_file"]):
                    # Ottieni il percorso della directory che contiene il file
                    file_path = file_entry['percorso_originale']
                    directory_path = os.path.dirname(file_path).replace('\\', '/')
                    # Formatta il percorso della directory per il collegamento ipertestuale
                    folder_uri = f"file:///{directory_path}"
                    
                    # Aggiungi riga al CSV con la formula italiana per i collegamenti alla CARTELLA
                    csv_writer.writerow([
                        file_entry['nome_file'],
                        file_entry['tipo'],
                        file_entry['dimensione'],
                        file_entry['ultima_modifica'],
                        file_entry['percorso_originale'],
                        f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")'
                    ])

                # Ottieni il contenuto del CSV come stringa
                csv_data = csv_header + csv_content.getvalue()

                # Aggiungi il file CSV dei file normali all'archivio
                zipf.writestr("Log_file_Trovati.csv", csv_data)

                # ---------- CREAZIONE DEL CSV DEGLI OMONIMI (SOLO SE CE NE SONO) ----------
                if omonimi_log:
                    # Creazione intestazione per il CSV degli omonimi con info di ricerca
                    omonimi_csv_header = f"LOG DEI FILE OMONIMI - {current_time_local}\n"
                    omonimi_csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                    omonimi_csv_header += f"Directory di ricerca: {search_directory}\n"
                    omonimi_csv_header += f"Parole chiave ricercate: {search_keywords}\n"
                    omonimi_csv_header += f"Totale file omonimi trovati: {len(omonimi_log)}\n\n"
                    
                    # Prepara l'output CSV per gli omonimi
                    omonimi_csv_content = io.StringIO()
                    omonimi_csv_writer = csv.writer(omonimi_csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
                    
                    # Intestazione colonne specifiche per gli omonimi
                    omonimi_csv_writer.writerow([
                        "Nome file", "Tipo", "Dimensione", "Ultima modifica", 
                        "Percorso originale", "Link", "Posizione nello ZIP", "In conflitto con"
                    ])
                    
                    # Estrai dai log completi solo i file omonimi
                    omonimi_files_details = [f for f in all_files_log if f['percorso_originale'] in omonimo_paths]
                    
                    # Crea un dizionario per unire le informazioni di omonimi_log e omonimi_files_details
                    omonimi_details_dict = {}
                    for entry in omonimi_log:
                        omonimi_details_dict[entry['percorso_originale']] = entry
                    
                    # Aggiungi i file omonimi al CSV
                    for file_entry in sorted(omonimi_files_details, key=lambda x: x["nome_file"]):
                        # Ottieni il percorso della directory che contiene il file
                        file_path = file_entry['percorso_originale']
                        directory_path = os.path.dirname(file_path).replace('\\', '/')
                        # Formatta il percorso della directory per il collegamento ipertestuale
                        folder_uri = f"file:///{directory_path}"
                        
                        # Trova le informazioni aggiuntive sugli omonimi
                        omonimo_info = omonimi_details_dict.get(file_entry['percorso_originale'], {})
                        posizione_zip = omonimo_info.get('posizione_zip', 'N/A')
                        in_conflitto = omonimo_info.get('primo_percorso', 'N/A')
                        
                        # Aggiungi riga al CSV degli omonimi con la formula italiana per i collegamenti alla CARTELLA
                        omonimi_csv_writer.writerow([
                            file_entry['nome_file'],
                            file_entry['tipo'],
                            file_entry['dimensione'],
                            file_entry['ultima_modifica'],
                            file_entry['percorso_originale'],
                            f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")',
                            posizione_zip,
                            in_conflitto
                        ])
                    
                    # Ottieni il contenuto del CSV come stringa
                    omonimi_csv_data = omonimi_csv_header + omonimi_csv_content.getvalue()
                    
                    # Aggiungi il file CSV degli omonimi all'archivio
                    zipf.writestr("omonimi_log.csv", omonimi_csv_data)

            # Verifica se ci sono algoritmi di hash selezionati
            hash_algorithms = result.get("hash_algorithms", [])
            hash_results = {}
            hash_success = False

            # Calcola gli hash del file ZIP se richiesto
            if hash_algorithms:
                try:
                    self.status_label["text"] = f"Calcolo hash del file ZIP..."
                    self.root.update_idletasks()
                    
                    # Calcola gli hash sul file ZIP
                    hash_results = self.calculate_file_hash(zip_path, hash_algorithms)
                    hash_success = True
                    
                    # Crea un file di report per gli hash
                    hash_report_path = os.path.splitext(zip_path)[0] + "_hash.txt"
                    with open(hash_report_path, 'w', encoding='utf-8') as hash_file:
                        hash_file.write(f"REPORT HASH PER IL FILE: {os.path.basename(zip_path)}\n")
                        hash_file.write(f"Data e ora: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n")
                        hash_file.write(f"Utente: {self.current_user}\n\n")
                        
                        for algorithm in hash_algorithms:
                            hash_value = hash_results.get(algorithm, "Errore")
                            hash_file.write(f"{algorithm.upper()}: {hash_value}\n")
                
                except Exception as e:
                    self.log_debug(f"Errore nel calcolo hash del file ZIP: {str(e)}")

            # Prepara il messaggio di completamento
            skipped_files = len(single_files) - len(filtered_single_files)
            message = f"Compressione completata!\nFile salvato in: {zip_path}\n"
            message += f"File organizzati nella cartella '{main_folder_name}'\n"
            message += f"Tipo di compressione utilizzata: {compression_text}\n"
            
            # Aggiungi informazioni sugli hash calcolati
            if hash_algorithms:
                if hash_success:
                    message += f"\nHash calcolati sul file ZIP:\n"
                    for algorithm in hash_algorithms:
                        hash_value = hash_results.get(algorithm, "Errore")
                        message += f"{algorithm.upper()}: {hash_value}\n"
                    message += f"\nI dettagli sono stati salvati in: {os.path.basename(hash_report_path)}"
                else:
                    message += f"\nErrore nel calcolo degli hash richiesti."

            if result["preserve"]:
                message += f"Struttura delle directory originali: Preservata\n"
            else:
                message += f"Struttura delle directory originali: Non preservata\n"
                
            message += f"\nCreato file di log completo ('Log_file_Trovati.csv')"

            if result["chunks"]:
                message += "\nElaborazione a blocchi: Attiva"

            if omonimi_log:
                message += f"\nTrovati {len(omonimi_log)} file omonimi."
                message += f"\nCreato log dettagliato degli omonimi ('omonimi_log.csv')"

            if skipped_files > 0:
                message += f"\n{skipped_files} file saltati perché già presenti nelle cartelle"

            messagebox.showinfo("Completato", message)

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            self.log_debug(f"Errore dettagliato durante la compressione:\n{error_details}")

            messagebox.showerror("Errore", 
                f"Errore durante la compressione: {str(e)}\n\n"
                f"Tipo di errore: {type(e).__name__}")

        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    @error_handler
    def _find_common_base_path(self, selected_items):
        """Trova il percorso base comune a tutti i file selezionati"""
        if not selected_items:
            return ""
        
        paths = []
        
        # Raccogli tutti i percorsi
        for item in selected_items:
            values = self.results_list.item(item)['values']
            paths.append(values[5])  # Il percorso è nella sesta colonna
        
        # Trova il percorso comune più lungo
        def common_path(paths):
            if not paths:
                return ""
            
            # Normalizza i percorsi per gestire i diversi separatori di directory
            paths = [os.path.normpath(p) for p in paths]
            
            # Dividi ogni percorso in componenti
            components = [p.split(os.path.sep) for p in paths]
            
            # Se ci sono percorsi di unità diverse (es. C:\ e D:\) su Windows
            if os.name == 'nt' and len(set(c[0] for c in components if c)) > 1:
                return ""  # Non c'è un percorso comune tra unità diverse
            
            common = []
            for i in range(min(len(c) for c in components)):
                if len(set(c[i] for c in components)) == 1:
                    common.append(components[0][i])
                else:
                    break
            
            # Se non c'è un elemento comune, restituisci la radice/unità
            if not common and os.name == 'nt' and components and components[0]:
                return components[0][0] + os.path.sep  # es. "C:\"
            elif not common:
                return os.path.sep  # radice Unix "/"
                
            return os.path.sep.join(common)
        
        base_path = common_path(paths)
        
        # Se il percorso base finisce con un separatore, va bene
        # altrimenti dobbiamo aggiungere il separatore se base_path non è vuoto
        if base_path and not base_path.endswith(os.path.sep):
            base_path = os.path.dirname(base_path)
        
        # Se non abbiamo trovato un percorso comune, usa il percorso di ricerca corrente
        if not base_path and hasattr(self, 'search_path'):
            base_path = self.search_path.get()
        
        return base_path
    
    @error_handler
    def get_directory_size(self, path):
        """Calculate the total size of a directory"""
        # Implementazione di una cache in memoria per risultati recenti
        if not hasattr(self, '_dir_size_cache'):
            self._dir_size_cache = {}
            self._dir_size_cache_timestamp = {}
        
        # Controllo rapido in cache
        cache_valid = False
        if path in self._dir_size_cache:
            # Verifica se il risultato in cache è ancora valido (max 60 secondi)
            if time.time() - self._dir_size_cache_timestamp.get(path, 0) < 60:
                cache_valid = True
                # Ulteriore verifica: controlla se la directory è stata modificata
                try:
                    if os.path.exists(path):
                        last_modified = os.path.getmtime(path)
                        if last_modified <= self._dir_size_cache_timestamp.get(f"{path}_mtime", 0):
                            return self._dir_size_cache[path]
                except:
                    pass
        
        if not os.path.exists(path):
            return 0
            
        try:
            total_size = 0
            if os.path.isfile(path):
                return os.path.getsize(path)
                
            # Per evitare il blocco dell'interfaccia, imposta un timeout massimo
            start_time = time.time()
            max_time = 30  # massimo 30 secondi
            files_count = 0
            error_count = 0
            
            # Ottimizzazione 1: Usa set per tracciare i percorsi già visitati
            visited = set()
            
            # Ottimizzazione 2: Buffer per file da processare e aggiornamenti UI più rari
            file_batch = []
            batch_size = 0
            update_interval = 0.5  # aggiorna UI ogni 0.5 secondi
            last_update = time.time()
                
            # For directories, walk through all files and subdirectories
            for dirpath, dirnames, filenames in os.walk(path):
                # Ottimizzazione 3: Salta directory già visitate (simbolici o hardlink)
                if dirpath in visited:
                    continue
                visited.add(dirpath)
                
                # Ottimizzazione 4: Processa i file in batch
                for f in filenames:
                    # Verifica se il timeout è scaduto
                    if time.time() - start_time > max_time:
                        self.log_debug(f"Timeout nel calcolo della dimensione per {path}")
                        # Salva il risultato parziale in cache
                        self._dir_size_cache[path] = total_size
                        self._dir_size_cache_timestamp[path] = time.time()
                        self._dir_size_cache_timestamp[f"{path}_mtime"] = os.path.getmtime(path) if os.path.exists(path) else 0
                        return total_size
                    
                    fp = os.path.join(dirpath, f)
                    file_batch.append(fp)
                    batch_size += 1
                    
                    # Processa il batch quando raggiunge una dimensione appropriata
                    if batch_size >= 100:
                        total_size += self._process_file_batch(file_batch, error_count)
                        files_count += batch_size - error_count
                        batch_size = 0
                        file_batch = []
                    
                    # Update the status periodically to show progress
                    current_time = time.time()
                    if current_time - last_update > update_interval:
                        self.root.after(0, lambda size=total_size: 
                            self.status_label.config(text=f"Calcolando dimensione: {self._format_size(size)}...") 
                            if hasattr(self, 'status_label') and self.status_label.winfo_exists() else None)
                        last_update = current_time
            
            # Processa gli ultimi file rimasti
            if file_batch:
                total_size += self._process_file_batch(file_batch, error_count)
            
            # Salva il risultato in cache
            self._dir_size_cache[path] = total_size
            self._dir_size_cache_timestamp[path] = time.time()
            self._dir_size_cache_timestamp[f"{path}_mtime"] = os.path.getmtime(path) if os.path.exists(path) else 0
                
            return total_size
        except Exception as e:
            self.log_debug(f"Error calculating directory size for {path}: {str(e)}")
            return 0

    @error_handler
    def _process_file_batch(self, file_batch, error_count):
        """Processo un batch di file per calcolare la dimensione totale"""
        batch_size = 0
        for fp in file_batch:
            try:
                # Verifica esplicita che il file esiste ancora prima di tentare di leggerne la dimensione
                if os.path.exists(fp) and not os.path.islink(fp):
                    batch_size += os.path.getsize(fp)
            except FileNotFoundError:
                # Ignora silenziosamente i file che non esistono più
                pass
            except (OSError, PermissionError) as e:
                error_count += 1
                # Limita il numero di errori da registrare per evitare spam nel log
                if error_count < 100:  
                    self.log_debug(f"Error getting size of {fp}: {str(e)}")
        return batch_size

    @error_handler
    def get_directory_size_system(self, path):
        """Utilizza comandi di sistema per ottenere dimensioni di directory molto grandi"""
        # Cache system per comandi di sistema
        if not hasattr(self, '_system_size_cache'):
            self._system_size_cache = {}
            self._system_size_timestamp = {}
        
        # Verifica cache per risultati recenti
        if path in self._system_size_cache:
            if time.time() - self._system_size_timestamp.get(path, 0) < 120:  # 2 minuti
                return self._system_size_cache[path]
        
        # Verifica se è un'unità disco completa o una directory troppo grande
        is_root_drive = False
        try:
            if os.name == 'nt':
                # Verifica se è una radice di unità come "C:/" o "C:\"
                if path.endswith(':\\') or path.endswith(':/'):
                    is_root_drive = True
                elif len(path) <= 3 and path[1:] == ':\\':  # come "C:\"
                    is_root_drive = True
        except:
            pass
        
        # Se è una radice di unità, fallback al metodo di stima
        if is_root_drive:
            self.log_debug(f"Rilevata richiesta per unità disco completa: {path}. Uso metodo di stima.")
            return self.estimate_directory_size(path)
        
        # Tenta di ottenere il numero di file per verificare se è una directory molto grande
        try:
            file_count = 0
            dir_count = 0
            for _, dirs, files in os.walk(path, topdown=True):
                file_count += len(files)
                dir_count += len(dirs)
                # Se la directory contiene troppi file/cartelle, usa il metodo di stima
                if file_count + dir_count > 50000:
                    self.log_debug(f"Directory troppo grande per comando di sistema: {path} ({file_count} file). Uso metodo di stima.")
                    return self.estimate_directory_size(path)
                # Limita la scansione iniziale
                if file_count > 1000:
                    break
        except:
            pass
        
        try:
            if os.name == 'nt':  # Windows
                import platform
                CREATE_NO_WINDOW = 0x08000000
                
                # Imposta un timeout più breve per evitare blocchi
                timeout = 30  # 30 secondi massimo
                
                # Usa robocopy che è più affidabile per questa operazione
                cmd = f'robocopy "{path}" NULL /L /S /NJH /BYTES /NC /NFL /NDL /XJ'
                
                startupinfo = None
                if hasattr(subprocess, 'STARTUPINFO'):
                    startupinfo = subprocess.STARTUPINFO()
                    startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                    startupinfo.wShowWindow = 0  # SW_HIDE
                
                try:
                    result = subprocess.check_output(
                        cmd, 
                        stderr=subprocess.STDOUT,
                        timeout=timeout,
                        creationflags=CREATE_NO_WINDOW if os.name == 'nt' else 0,
                        startupinfo=startupinfo)
                    
                    # Estrazione della dimensione dal risultato
                    output = result.decode('utf-8', errors='ignore')
                    size = 0
                    
                    # Cerca la riga con il totale bytes
                    for line in output.splitlines():
                        if "Bytes :" in line:
                            # Trova numeri nella stringa
                            size_str = line.split("Bytes :")[1].strip()
                            try:
                                # Rimuove eventuali separatori di migliaia
                                size_str = size_str.replace(',', '').replace('.', '')
                                size = int(size_str)
                                break
                            except ValueError:
                                pass
                    
                    # Memorizza in cache
                    self._system_size_cache[path] = size
                    self._system_size_timestamp[path] = time.time()
                    
                    return size
                    
                except subprocess.TimeoutExpired:
                    self.log_debug(f"Timeout durante il calcolo della dimensione per {path}")
                    return self.estimate_directory_size(path)
                    
                except Exception as e:
                    self.log_debug(f"Fallito metodo robocopy: {str(e)}. Provo con PowerShell.")
                    
                    # Prova con PowerShell come backup se robocopy fallisce
                    try:
                        ps_cmd = f'powershell -command "Get-ChildItem -Path \'{path}\' -Recurse -Force -ErrorAction SilentlyContinue | Measure-Object -Property Length -Sum | Select-Object -ExpandProperty Sum"'
                        
                        result = subprocess.check_output(
                            ps_cmd, 
                            shell=True, 
                            stderr=subprocess.STDOUT,
                            timeout=timeout,
                            creationflags=CREATE_NO_WINDOW if os.name == 'nt' else 0,
                            startupinfo=startupinfo)
                        
                        size = int(result.strip())
                        
                        # Memorizza in cache
                        self._system_size_cache[path] = size
                        self._system_size_timestamp[path] = time.time()
                        
                        return size
                    except:
                        # Fallback finale al metodo di stima
                        return self.estimate_directory_size(path)
        except Exception as e:
            self.log_error(f"Errore nel calcolo della dimensione della directory: {str(e)}")
            # Se fallisce il metodo system, prova con un metodo alternativo
            return self.estimate_directory_size(path)

    @error_handler
    def _calculate_dir_size_thread(self, path):
        """Thread function to calculate directory size with improved disk handling"""
        # Limita il carico di lavoro per evitare blocchi del sistema
        import threading
        try:
            # Imposta una priorità più bassa per questo thread
            if hasattr(threading.current_thread(), "setName"):
                threading.current_thread().name = "LowPriority_DirSize"
        except:
            pass
            
        calculation_mode = self.dir_size_calculation.get()
        dir_size = 0
        
        try:
            # Verifica se è un'unità disco
            is_drive_root = False
            if os.name == 'nt':
                if path.endswith(':\\') or path.endswith(':/'):
                    is_drive_root = True
                elif len(path) <= 3 and path[1:] == ':':
                    is_drive_root = True
            
            # Per unità disco, usa sempre il metodo accurato
            if is_drive_root:
                self.log_debug(f"Rilevata unità disco: {path}, utilizzo metodo accurato")
                accurate_size = self.get_disk_accurate_size(path)
                if accurate_size is not None:
                    dir_size = accurate_size
                    # Aggiorna l'interfaccia
                    if self.root and self.root.winfo_exists():
                        self.root.after(0, lambda: self.dir_size_var.set(self._format_size(dir_size)))
                        self.root.after(0, lambda: self.status_label.config(text="In attesa...") 
                                    if hasattr(self, 'status_label') and self.status_label.winfo_exists() else None)
                    return
            
            # Per le directory normali o se il metodo accurato fallisce, usa il metodo selezionato
            if calculation_mode == "preciso":
                dir_size = self.get_directory_size(path)
            elif calculation_mode == "stimato":
                dir_size = self.estimate_directory_size(path)
            elif calculation_mode == "sistema":
                dir_size = self.get_directory_size_system(path)
            else:  # incrementale o fallback
                dir_size = self.get_directory_size(path)
                
            # Update the UI from the main thread - CORREZIONE con check
            if self.root and self.root.winfo_exists():
                self.root.after(0, lambda: self.dir_size_var.set(self._format_size(dir_size)))
                self.root.after(0, lambda: self.status_label.config(text="In attesa...") 
                            if hasattr(self, 'status_label') and self.status_label.winfo_exists() else None)
        
        except Exception as e:
            self.log_debug(f"Errore nel calcolo della dimensione: {str(e)}")
            if self.root and self.root.winfo_exists():
                self.root.after(0, lambda: self.dir_size_var.set("Errore"))
                self.root.after(0, lambda: self.status_label.config(text="In attesa...") 
                            if hasattr(self, 'status_label') and self.status_label.winfo_exists() else None)

    @error_handler
    def get_disk_accurate_size(self, drive_path):
        """Calcola la dimensione reale dei dati su un'unità disco completa
        utilizzando le API di Windows che forniscono dati più accurati."""
        # Verifica se il percorso è una radice di unità
        is_drive_root = False
        drive_letter = None
        
        try:
            if os.name == 'nt':
                # Normalizza il percorso
                if drive_path.endswith(':\\') or drive_path.endswith(':/'):
                    is_drive_root = True
                    drive_letter = drive_path[0]
                elif len(drive_path) <= 3 and drive_path[1:] == ':':
                    is_drive_root = True
                    drive_letter = drive_path[0]
        except:
            pass
        
        if not is_drive_root or not drive_letter:
            # Non è un'unità disco, usa il metodo standard
            self.log_debug(f"Il percorso {drive_path} non è un'unità disco, utilizzo metodo standard")
            return None
        
        try:
            # Usa WMI per ottenere informazioni accurate sul disco
            try:
                import pythoncom
                import wmi
                import sys
                
            except ImportError:
                self.log_debug("Libreria WMI non disponibile, utilizzo statistiche del sistema operativo")
                # Fallback a shutil.disk_usage che è più accurato di get_directory_size per unità intere
                total, used, free = shutil.disk_usage(drive_letter + ":\\")
                return used
            
            # Inizializza COM in questo thread
            pythoncom.CoInitialize()
            
            c = wmi.WMI()
            
            # Ottieni un'istanza più affidabile di DriveType=3 (disco locale)
            for logical_disk in c.Win32_LogicalDisk(DriveType=3):
                if logical_disk.DeviceID[0].lower() == drive_letter.lower():
                    # Ottieni le dimensioni dal sistema operativo in modo affidabile
                    used_space = int(logical_disk.Size) - int(logical_disk.FreeSpace)
                    
                    self.log_debug(f"Dimensione unità {drive_letter}: totale={self._format_size(int(logical_disk.Size))}, "
                                f"usato={self._format_size(used_space)}, "
                                f"libero={self._format_size(int(logical_disk.FreeSpace))}")
                    
                    # Restituisci lo spazio utilizzato, che è il dato più rilevante
                    return used_space
            
            # Se non troviamo il disco specifico, usiamo shutil come fallback
            self.log_debug(f"Disco {drive_letter} non trovato via WMI, uso metodo standard")
            total, used, free = shutil.disk_usage(drive_letter + ":\\")
            return used
        
        except Exception as e:
            self.log_debug(f"Errore nel calcolo accurato della dimensione del disco {drive_letter}: {str(e)}")
            return None
        finally:
            # Rilascia COM
            try:
                if 'pythoncom' in sys.modules:
                    pythoncom.CoUninitialize()
            except:
                pass

    @error_handler
    def estimate_directory_size(self, path, sample_size=100):
        """Stima la dimensione di una directory campionando alcuni file - Versione ottimizzata"""
        import random
        
        # Usa la cache quando disponibile
        if not hasattr(self, '_estimate_size_cache'):
            self._estimate_size_cache = {}
            self._estimate_timestamp = {}
        
        # Verifica cache validità 5 minuti
        if path in self._estimate_size_cache:
            if time.time() - self._estimate_timestamp.get(path, 0) < 300:
                return self._estimate_size_cache[path]
        
        if not os.path.exists(path) or os.path.isfile(path):
            return self.get_directory_size(path)  # Usa il metodo esatto per file o percorsi non validi
        
        try:
            # Ottieni un conteggio rapido dei file (questo è veloce)
            total_files = 0
            sampled_files = 0
            total_sampled_size = 0
            
            # Ottimizzazione: campionamento progressivo durante la scansione
            file_paths = []
            
            # Prima passata veloce combinata con campionamento
            for root, _, files in os.walk(path, topdown=True):
                # Aggiungi conteggio file
                file_count = len(files)
                total_files += file_count
                
                # Seleziona alcuni file da questo batch per il campionamento
                # Campiona in modo uniforme ma non randomico (più efficiente)
                if files and total_files < 50000:  # Limita la raccolta per grandi directory
                    # Seleziona alcuni file da questo batch in modo più efficiente
                    step = max(1, len(files) // min(10, sample_size // 10 + 1))
                    for i in range(0, len(files), step):
                        if len(file_paths) < sample_size * 2:  # Raccogliamo più file del necessario
                            file_paths.append(os.path.join(root, files[i]))
                
                # Limita il tempo della prima passata ma assicura un campione minimo
                if total_files > 10000 and len(file_paths) >= sample_size:
                    break
                    
            # Se pochi file, usare metodo preciso
            if total_files < 500:
                return self.get_directory_size(path)
                
            # Seconda fase: calcola dimensione dei file campionati
            # Usa un sottoinsieme casuale dei percorsi dei file per un campione rappresentativo
            if len(file_paths) > sample_size:
                random.shuffle(file_paths)
                file_paths = file_paths[:sample_size]
                
            for file_path in file_paths:
                try:
                    if os.path.exists(file_path) and not os.path.islink(file_path):
                        file_size = os.path.getsize(file_path)
                        total_sampled_size += file_size
                        sampled_files += 1
                        
                        # Feedback progressivo
                        if sampled_files % 10 == 0:
                            avg_so_far = total_sampled_size / sampled_files
                            est_so_far = avg_so_far * total_files
                            if hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                self.root.after(0, lambda s=self._format_size(est_so_far): 
                                        self.status_label.config(text=f"Stima dimensione: {s}..."))
                except:
                    pass
            
            # Calcola la stima finale
            if sampled_files > 0:
                avg_file_size = total_sampled_size / sampled_files
                estimated_size = avg_file_size * total_files
                
                # Salva in cache
                self._estimate_size_cache[path] = estimated_size
                self._estimate_timestamp[path] = time.time()
                
                self.log_debug(f"Dimensione stimata per {path}: {self._format_size(estimated_size)} (basata su {sampled_files} campioni)")
                return estimated_size
            else:
                return self.get_directory_size(path)  # Fallback al metodo standard
                
        except Exception as e:
            self.log_debug(f"Errore nella stima della dimensione: {str(e)}")
            return 0

    @error_handler
    def refresh_directory_size(self):
        """Aggiorna manualmente il calcolo della dimensione della directory - Versione ottimizzata"""
        # Ottieni il percorso corrente
        path = self.search_path.get()
        
        # Verifica che il percorso esista
        if not path or not os.path.exists(path):
            messagebox.showinfo("Informazione", "Seleziona prima un percorso valido")
            return
            
        # Verifica la modalità di calcolo
        calculation_mode = self.dir_size_calculation.get()
        if calculation_mode == "disabilitato":
            response = messagebox.askyesno("Calcolo disabilitato", 
                                        "Il calcolo della dimensione è attualmente disabilitato.\n\n" +
                                        "Vuoi attivarlo e procedere con il calcolo?")
            if response:
                # Seleziona la modalità "preciso"
                self.dir_size_calculation.set("preciso")
            else:
                return
        
        # Aggiorna lo stato
        self.dir_size_var.set("Calcolo in corso...")
        if hasattr(self, 'status_label') and self.status_label.winfo_exists():
            self.status_label.config(text="Calcolo dimensione directory...")
        
        # Interrompi eventuali thread in esecuzione
        if hasattr(self, '_dir_size_thread') and self._dir_size_thread is not None:
            # Segnala interruzione
            if hasattr(self, '_stop_calculation'):
                self._stop_calculation = True
        
        # Flag per eventuali interruzioni
        self._stop_calculation = False
        
        # Esegui il calcolo in un thread separato
        self._dir_size_thread = threading.Thread(target=self._calculate_dir_size_thread, args=(path,), daemon=True)
        self._dir_size_thread.start()

    @error_handler
    def get_disk_space(self, path):
        """Get disk space information for the partition containing the path"""
        if not os.path.exists(path):
            return (0, 0, 0)
        
        try:
            if os.name == 'nt':  # Windows
                total, used, free = shutil.disk_usage(os.path.splitdrive(path)[0] + '\\')
            else:  # Linux/Mac
                total, used, free = shutil.disk_usage(os.path.abspath(os.path.join(path, os.pardir)))
            return (total, used, free)
        except Exception as e:
            self.log_debug(f"Error getting disk space for {path}: {str(e)}")
            return (0, 0, 0)

    @error_handler
    def update_disk_info(self, path=None, calculate_dir_size=True):
        """Aggiorna le informazioni sul disco e sulla directory"""
        if path is None:
            path = self.search_path.get()
        
        if not path or not os.path.exists(path):
            return

        # Aggiorna la variabile di controllo globale in base al parametro
        self.directory_calculation_enabled = calculate_dir_size
        
        # Solo per debug
        self.log_debug(f"Aggiornamento info disco: path={path}, calculation_enabled={self.directory_calculation_enabled}")
        
        # Se il calcolo è disabilitato, aggiorna solo le informazioni del disco
        if not self.directory_calculation_enabled:
            self.log_debug("Calcolo dimensione directory disabilitato nelle impostazioni")
            
            # Aggiorniamo solo le informazioni del disco senza calcolare la dimensione della directory
            try:
                total, used, free = self.get_disk_space(path)
                percent_used = (used / total) * 100 if total > 0 else 0
                
                total_formatted = self._format_size(total)
                free_formatted = self._format_size(free)
                
                # CORREZIONE: Verifica se il widget esiste prima di usarlo
                if hasattr(self, 'disk_info_label') and self.disk_info_label.winfo_exists():
                    self.disk_info_label.config(text=f"Disco: {total_formatted} totali, {free_formatted} liberi ({percent_used:.1f}% usato)")
                else:
                    self.log_debug("Widget disk_info_label non disponibile")
                
                # Importante: NON avviare il thread di calcolo della directory
                # CORREZIONE: Verifica se il widget esiste prima di usarlo
                if hasattr(self, 'dir_size_label') and self.dir_size_label.winfo_exists():
                    self.dir_size_label.config(text="Dimensione directory: calcolo disabilitato")
                else:
                    self.log_debug("Widget dir_size_label non disponibile")
            except Exception as e:
                self.log_debug(f"Errore nell'aggiornamento info disco: {str(e)}")
            
            return
        
        # Se arriviamo qui, il calcolo è abilitato e dobbiamo procedere
        self._async_update_disk_info(path, calculate_dir_size)
    
    @error_handler
    def _update_disk_info_thread(self, path, calculate_dir_size):
        """Thread per calcolare le informazioni del disco"""
        try:
            # Get disk space info
            total, used, free = self.get_disk_space(path)
            
            # Aggiorna l'UI dal thread principale
            self.root.after(0, lambda: self.total_disk_var.set(self._format_size(total)))
            self.root.after(0, lambda: self.used_disk_var.set(self._format_size(used)))
            self.root.after(0, lambda: self.free_disk_var.set(self._format_size(free)))
            self.root.after(0, lambda: self.status_label.config(text="In attesa..."))
            
            # Log per debug
            self.log_debug(f"Info disco aggiornate: Totale={self._format_size(total)}, Libero={self._format_size(free)}")
            
        except Exception as e:
            self.log_debug(f"Errore nel calcolo dello spazio disco: {str(e)}")
            self.root.after(0, lambda: self.total_disk_var.set("Errore"))
            self.root.after(0, lambda: self.used_disk_var.set("Errore"))
            self.root.after(0, lambda: self.free_disk_var.set("Errore"))
        
        # Controlla esplicitamente le impostazioni per il calcolo della directory
        calculation_mode = self.dir_size_calculation.get()
        
        # Gestisce il caso in cui non si deve calcolare la dimensione
        if not calculate_dir_size or calculation_mode == "disabilitato" or self.is_searching:
            # Se il calcolo è disabilitato o stiamo cercando, mostra "N/D" invece di calcolare
            self.log_debug("Calcolo dimensione directory saltato: " + 
                        ("parametro disabilitato" if not calculate_dir_size else 
                        "impostazione disabilitata" if calculation_mode == "disabilitato" else 
                        "ricerca in corso"))
            
            self.root.after(0, lambda: self.dir_size_var.set("N/D"))
            return
        
        # Se arriviamo qui, dobbiamo calcolare la dimensione
        try:
            self.log_debug(f"Avvio calcolo dimensione directory: {path} (modalità: {calculation_mode})")
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            self._calculate_dir_size_thread(path)
        except Exception as e:
            self.log_debug(f"Errore nell'avvio del calcolo dimensione directory: {str(e)}")
            self.root.after(0, lambda: self.dir_size_var.set("Errore"))

    @error_handler
    def _async_update_disk_info(self, path, calculate_dir_size):
        """Esegue il calcolo delle informazioni del disco in background"""
        try:
            # Get disk space info
            total, used, free = self.get_disk_space(path)
            
            # Aggiorna l'UI dal thread principale
            self.root.after(0, lambda: self.total_disk_var.set(self._format_size(total)))
            self.root.after(0, lambda: self.used_disk_var.set(self._format_size(used)))
            self.root.after(0, lambda: self.free_disk_var.set(self._format_size(free)))
        except Exception as e:
            self.log_debug(f"Errore nel calcolo dello spazio disco: {str(e)}")
            self.root.after(0, lambda: self.total_disk_var.set("Errore"))
            self.root.after(0, lambda: self.used_disk_var.set("Errore"))
            self.root.after(0, lambda: self.free_disk_var.set("Errore"))
        
        # Verifica modalità di calcolo dimensione
        calculation_mode = self.dir_size_calculation.get()
        
        # Rispetta sempre l'impostazione disabilitato, indipendentemente dal valore di calculate_dir_size
        if calculation_mode == "disabilitato":
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo disattivato"))
            self.log_debug("Calcolo dimensione directory disattivato nelle impostazioni")
            return
        
        # Se è incrementale e siamo in fase di ricerca, non fare niente qui
        if calculation_mode == "incrementale" and self.is_searching:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            return
                
        # Per le altre modalità o se non siamo in ricerca
        if calculate_dir_size:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            # Il calcolo avviene già in un thread separato
            self._calculate_dir_size_thread(path)
        else:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo disattivato"))
    

    @error_handler # Funzione helper per formattare la dimensione del file
    def _format_size(self, size_bytes):
        """Formatta la dimensione del file in modo leggibile"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.2f} KB"
        elif size_bytes < 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024):.2f} MB"
        elif size_bytes < 1024 * 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"
        else:
            return f"{size_bytes / (1024 * 1024 * 1024 * 1024):.2f} TB"
        
    @error_handler
    def calculate_file_hash(self, file_path, algorithms=None):
        """Calcola gli hash di un file usando gli algoritmi specificati :param file_path: Percorso del file
        :param algorithms: Lista degli algoritmi da usare ('md5', 'sha1', 'sha256'):return: Dizionario con gli hash calcolati"""
        import hashlib
        
        if algorithms is None:
            algorithms = ['md5']
        
        results = {}
        
        try:
            for algorithm in algorithms:
                if algorithm == 'md5':
                    hash_obj = hashlib.md5()
                elif algorithm == 'sha1':
                    hash_obj = hashlib.sha1()
                elif algorithm == 'sha256':
                    hash_obj = hashlib.sha256()
                else:
                    continue
                    
                # Leggi il file in blocchi per gestire file di grandi dimensioni
                with open(file_path, 'rb') as f:
                    for chunk in iter(lambda: f.read(4096), b''):
                        hash_obj.update(chunk)
                    
                results[algorithm] = hash_obj.hexdigest()
                
                # Log per debug
                self.log_debug(f"Calcolato {algorithm} per {os.path.basename(file_path)}: {results[algorithm]}")
                
        except Exception as e:
            self.log_debug(f"Errore nel calcolo hash {algorithm} per {file_path}: {str(e)}")
            # In caso di errore, restituisci None per quell'algoritmo
            for algorithm in algorithms:
                if algorithm not in results:
                    results[algorithm] = "Errore"
        
        return results

    @error_handler # Funzione helper per determinare il tipo di file
    def _get_file_type(self, file_path):
        """Determina il tipo di file in base all'estensione"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.txt', '.md', '.rtf']:
            return "Documento di testo"
        elif ext in ['.doc', '.docx', '.odt']:
            return "Documento Word"
        elif ext in ['.xls', '.xlsx', '.ods']:
            return "Foglio di calcolo"
        elif ext in ['.pdf']:
            return "PDF"
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            return "Immagine"
        elif ext in ['.mp3', '.wav', '.ogg', '.flac']:
            return "Audio"
        elif ext in ['.mp4', '.avi', '.mkv', '.mov']:
            return "Video"
        elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
            return "Archivio"
        elif ext in ['.exe', '.dll', '.bat', '.cmd']:
            return "Eseguibile"
        else:
            return "File"
    
    @error_handler
    def get_main_folder_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome della cartella principale"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Cartella Principale")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=tk.BOTH, expand=tk.YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome della cartella principale nell'archivio:", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = tk.StringVar(value="files")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=tk.X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=tk.X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            folder_name = name_var.get().strip()
            if folder_name:
                result["name"] = folder_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per la cartella", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea Cartella", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=tk.RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]
    
    @error_handler
    def open_file_location(self, event=None):
        """Apre il percorso del file selezionato nel file explorer"""
            
        selected_items = self.results_list.selection()
        if not selected_items:
            return
            
        selected_item = selected_items[0]  # Prendi il primo elemento selezionato
        file_path = self.results_list.item(selected_item, "values")[6]  # Modificato da 5 a 6
        
        try:
            if os.path.exists(file_path):
                # Ottieni la directory contenente il file
                directory = os.path.dirname(file_path)
                
                if os.name == 'nt':  # Windows
                    # Converti eventuali forward slash in backslash per Windows
                    file_path = os.path.normpath(file_path)
                    
                    # Usa il flag CREATE_NO_WINDOW per nascondere la finestra del CMD
                    CREATE_NO_WINDOW = 0x08000000  # Per Python < 3.7
                    
                    # Importante: impostare shell=False e usare creationflags per nascondere la finestra
                    subprocess.run(
                    ['explorer', f'/select,{file_path}'],
                    shell=False,
                    creationflags=CREATE_NO_WINDOW)
                
                self.log_debug(f"Apertura percorso: {file_path}")
            else:
                messagebox.showinfo("Informazione", f"Il percorso non esiste: {file_path}")
                self.log_debug(f"Percorso non esistente: {file_path}")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile aprire il percorso: {str(e)}")
            self.log_debug(f"Errore nell'apertura del percorso: {str(e)}")

    @error_handler
    def configure_extensions(self, mode="base"):
        """Dialog to configure file extensions for different search modes"""
        dialog = ttk.Toplevel(self.root)
        dialog.title(f"Configura estensioni - Modalità {mode.capitalize()}")
        dialog.geometry("1100x450")
        dialog.transient(self.root)
        dialog.grab_set()
        
        main_frame = ttk.Frame(dialog, padding=15)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Heading text
        ttk.Label(main_frame, text=f"Seleziona le estensioni di file da includere nella ricerca per la modalità {mode}", 
                font=("", 10)).pack(anchor=W, pady=(0, 15))
        
        # Category tabs
        notebook = ttk.Notebook(main_frame)
        notebook.pack(fill=BOTH, expand=YES, pady=(0, 15))
        
        # Initialize dictionary to store checkbutton variables
        ext_vars = {}
        
        # Define extensions by category - EXPANDED with more extensions
        extension_categories = {
        "Documenti": [
            (".txt", "File di testo"),
            (".doc", "Word vecchio"),
            (".docx", "Word"),
            (".pdf", "PDF"),
            (".rtf", "Rich Text"),
            (".odt", "OpenDoc Text"),
            (".md", "Markdown"),
            (".xml", "XML"),
            (".html", "HTML"),
            (".htm", "HTM"),
            (".log", "Log file"),
            (".tex", "LaTeX"),
            (".rst", "reStructuredText"),
            (".epub", "E-book EPUB"),
            (".mobi", "E-book Mobi"),
            (".vcf", "vCard"),
            (".ics", "iCalendar"),
        ],
        "Email": [
            (".eml", "Email standard"),
            (".msg", "Email formato Outlook"),
            (".pst", "Archivio Outlook"),
            (".ost", "Archivio Outlook Offline"),
            (".mbox", "Mailbox Unix/Linux"),
            (".emlx", "Email formato Apple Mail"),
            (".mbx", "Mailbox formato Eudora/Thunderbird"),
            (".dbx", "Archivio Outlook Express"),
            (".wab", "Windows Address Book"),
            (".nws", "Email formato Outlook Express"),
            (".mht", "MIME HTML Archive"),
            (".mhtml", "MIME HTML Archive"),
            (".imapmbox", "IMAP Mailbox"),
            (".email", "Email generica")
        ],
        "Fogli calcolo": [
            (".xls", "Excel vecchio"),
            (".xlsx", "Excel"),
            (".ods", "OpenCalc"),
            (".csv", "CSV (valori separati da virgola)"),
            (".tsv", "TSV (valori separati da tab)"),
            (".dbf", "Database File"),
            (".dif", "Data Interchange Format")
        ],
        "Presentazioni": [
            (".ppt", "PowerPoint vecchio"),
            (".pptx", "PowerPoint"),
            (".odp", "OpenImpress"),
            (".key", "Keynote"),
            (".pps", "PowerPoint Show")
        ],
        "Database": [
            (".db", "Database generico"),
            (".sqlite", "SQLite"),
            (".sqlite3", "SQLite3"),
            (".mdb", "Access DB"),
            (".accdb", "Access DB nuovo"),
            (".odb", "OpenOffice DB"),
            (".sql", "SQL script")
        ],
        "Immagini": [
            (".jpg", "JPEG"),
            (".jpeg", "JPEG"),
            (".png", "PNG"),
            (".gif", "GIF"),
            (".bmp", "Bitmap"),
            (".tiff", "TIFF"),
            (".tif", "TIF"),
            (".svg", "SVG"),
            (".webp", "WebP"),
            (".ico", "Icon"),
            (".raw", "Raw"),
            (".psd", "Photoshop"),
            (".ai", "Illustrator"),
            (".odg", "OpenOffice Draw"),
            (".xcf", "GIMP"),
            (".heic", "HEIC")
        ],
        "Audio": [
            (".mp3", "MP3"),
            (".wav", "WAV"),
            (".ogg", "OGG"),
            (".flac", "FLAC"),
            (".aac", "AAC"),
            (".m4a", "M4A"),
            (".wma", "WMA"),
            (".mid", "MIDI"),
            (".midi", "MIDI"),
            (".aiff", "AIFF"),
            (".opus", "Opus")
        ],
        "Video": [
            (".mp4", "MP4"),
            (".avi", "AVI"),
            (".mkv", "MKV"),
            (".mov", "MOV"),
            (".wmv", "WMV"),
            (".flv", "FLV"),
            (".webm", "WebM"),
            (".m4v", "M4V"),
            (".mpg", "MPEG"),
            (".mpeg", "MPEG"),
            (".3gp", "3GP"),
            (".ogv", "OGV"),
            (".ts", "TS (Transport Stream)")
        ],
        "Archivi": [
            (".zip", "ZIP"),
            (".rar", "RAR"),
            (".7z", "7-Zip"),
            (".tar", "TAR"),
            (".gz", "GZip"),
            (".bz2", "BZip2"),
            (".tgz", "Tar GZipped"),
            (".xz", "XZ"),
            (".jar", "Java Archive")
        ],
        "Eseguibili": [
            (".exe", "Eseguibile"),
            (".dll", "Libreria"),
            (".bat", "Batch"),
            (".cmd", "Command"),
            (".ps1", "PowerShell"),
            (".vbs", "VBScript"),
            (".sh", "Shell script"),
            (".msi", "Installer Windows"),
            (".app", "Applicazione macOS"),
            (".deb", "Pacchetto Debian"),
            (".rpm", "Red Hat Package"),
            (".apk", "Android Package")
        ],
        "Configurazione": [
            (".ini", "INI"),
            (".config", "Config"),
            (".conf", "Config"),
            (".reg", "Registry"),
            (".cfg", "Config"),
            (".properties", "Properties"),
            (".yml", "YAML"),
            (".yaml", "YAML"),
            (".json", "JSON Config"),
            (".toml", "TOML"),
            (".env", "Environment"),
            (".htaccess", "Apache Config"),
            (".plist", "macOS Property List")
        ],
        "Programmazione": [
            (".c", "C"),
            (".cpp", "C++"),
            (".cs", "C#"),
            (".java", "Java"),
            (".py", "Python"),
            (".js", "JavaScript"),
            (".php", "PHP"),
            (".rb", "Ruby"),
            (".go", "Golang"),
            (".rs", "Rust"),
            (".swift", "Swift"),
            (".pl", "Perl"),
            (".lua", "Lua"),
            (".h", "Header C"),
            (".hpp", "Header C++"),
            (".vb", "Visual Basic"),
            (".ts", "TypeScript"),
            (".scala", "Scala"),
            (".groovy", "Groovy"),
            (".kt", "Kotlin")
        ]
    }
        
        # Create a list of all extensions for "profonda" mode
        all_extensions = []
        for category_extensions in extension_categories.values():
            for ext, _ in category_extensions:
                all_extensions.append(ext.lower())
        
        # Load current settings (this would load from your saved settings)
        current_settings = self.get_extension_settings(mode)
        
        # Extensions that should be included in each search level
        base_extensions = ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.log', 
                        '.docx', '.doc', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls', '.rtf', '.odt', '.ods', '.odp',
                        '.csv','.eml', '.msg', '.emlx']
                        
        advanced_extensions = base_extensions + ['.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', 
                                            '.vbs', '.js', '.config', '.ini', '.reg',
                                            '.pst', '.ost', '.mbox']
        
        # Per modalità profonda, usa tutte le estensioni definite
        if mode == "profonda":
            # Sovrascrive le impostazioni correnti per far sì che tutte le estensioni siano selezionate
            current_settings = all_extensions
        
        # Dizionario per tenere traccia dei pulsanti per categoria
        category_buttons = {}
        
        # Create tabs for each category
        for category, extensions in extension_categories.items():
            # Create a frame for this category
            category_frame = ttk.Frame(notebook, padding=10)
            notebook.add(category_frame, text=category)
            
            # Aggiungi pulsante "Seleziona tutto" in alto
            button_frame = ttk.Frame(category_frame)
            button_frame.pack(fill=X, pady=(0, 10))
            
            # Funzione per selezionare tutte le estensioni in questa categoria
            def select_all_category(cat=category):
                for ext, desc in extension_categories[cat]:
                    if ext.lower() in ext_vars:
                        ext_vars[ext.lower()].set(True)
            
            # Funzione per deselezionare tutte le estensioni in questa categoria
            def deselect_all_category(cat=category):
                for ext, desc in extension_categories[cat]:
                    if ext.lower() in ext_vars:
                        ext_vars[ext.lower()].set(False)
            
            # Aggiungi i pulsanti
            select_all_btn = ttk.Button(button_frame, text=f"Seleziona tutto", 
                                    command=select_all_category)
            select_all_btn.pack(side=LEFT, padx=(0, 5))
            
            deselect_all_btn = ttk.Button(button_frame, text=f"Deseleziona tutti", 
                                        command=deselect_all_category)
            deselect_all_btn.pack(side=LEFT)
            
            # Salva i pulsanti nel dizionario per riferimento futuro
            category_buttons[category] = (select_all_btn, deselect_all_btn)
            
            # Create a grid layout
            content_frame = ttk.Frame(category_frame)
            content_frame.pack(fill=BOTH, expand=YES)
            row, col = 0, 0
            
            # Add checkboxes for each extension in this category
            for ext, desc in extensions:
                # Check if this extension should be selected based on mode
                is_selected = ext.lower() in current_settings
                
                # Create a variable for this checkbox
                var = BooleanVar(value=is_selected)
                ext_vars[ext.lower()] = var
                
                # Create the checkbox
                cb = ttk.Checkbutton(content_frame, text=f"{ext} ({desc})", variable=var)
                cb.grid(row=row, column=col, sticky=W, padx=5, pady=3)
                
                # Update grid position
                col += 1
                if col > 2:  # 3 columns per row
                    col = 0
                    row += 1
        
        # Button frame
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        # Helper functions
        def select_all():
            for var in ext_vars.values():
                var.set(True)
                
        def deselect_all():
            for var in ext_vars.values():
                var.set(False)
                
        def restore_defaults():
            # Modified restore_defaults function
            if mode == "profonda":
                # For deep mode, select ALL extensions
                for var in ext_vars.values():
                    var.set(True)
            else:
                # Reset to default extensions for base/avanzata modes
                default_list = base_extensions if mode == "base" else advanced_extensions
                for ext, var in ext_vars.items():
                    var.set(ext in default_list)
        
        def save_settings():
            # Save the selected extensions
            selected_extensions = [ext for ext, var in ext_vars.items() if var.get()]
            self.save_extension_settings(mode, selected_extensions)
            dialog.destroy()
            
            # Update the search_depth combo if needed
            current_depth = self.search_depth.get()
            if current_depth == mode:
                # Refresh the UI to reflect the new settings
                self.log_debug(f"Aggiornate le estensioni per la modalità {mode} ({len(selected_extensions)} estensioni)")
        
        # Buttons
        ttk.Button(btn_frame, text="Seleziona tutti", command=select_all).pack(side=LEFT, padx=5)
        ttk.Button(btn_frame, text="Deseleziona tutti", command=deselect_all).pack(side=LEFT, padx=5)
        ttk.Button(btn_frame, text="Ripristina default", command=restore_defaults).pack(side=LEFT, padx=5)
        
        ttk.Button(btn_frame, text="Annulla", command=dialog.destroy).pack(side=RIGHT, padx=5)
        ttk.Button(btn_frame, text="Salva", command=save_settings).pack(side=RIGHT, padx=5)
        
        # Center the dialog on screen
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Initial focus
        notebook.select(0)  # Focus first tab
    
    @error_handler
    def get_default_extensions(self, mode="base"):
        """Get default extensions for the specified search mode"""
        if mode == "base":
            return [
                # Documenti di testo
                '.txt', '.md', '.html', '.htm', '.xml', '.json', '.log',
                # Documenti Office
                '.docx', '.doc', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls',
                # Formati OpenDocument
                '.rtf', '.odt', '.ods', '.odp', 
                # Fogli di calcolo
                '.csv',
                # Formati email di base
                '.eml', '.msg', '.emlx'
            ]
        elif mode == "avanzata":
            # Il resto del codice rimane uguale
            base_exts = self.get_default_extensions("base")
            advanced_only = [
                # Estensioni avanzate...
                '.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', '.vbs', 
                '.config', '.ini', '.reg',
                '.py', '.java', '.php', '.cs', '.cpp', '.c', '.h', '.rb', '.js',
                '.db', '.sqlite', '.sqlite3',
                '.pst', '.ost', '.mbox', '.mbx', '.dbx',
                '.env', '.yml', '.yaml', '.toml', '.json5',
                '.bak', '.old', '.tmp', '.temp',
                '.epub', '.tex', '.rst',
                '.css', '.less', '.scss', '.jsp', '.asp', '.aspx',
                '.dot', '.dotx', '.xlt', '.xltx', '.pot', '.potx', '.ppsx'
            ]
            return base_exts + [ext for ext in advanced_only if ext not in base_exts]
        else:  # profonda
            return []  # Il valore effettivo viene determinato in configure_extensions
    
    @error_handler
    def get_extension_settings(self, mode="base"):
        """Load saved extension settings for the specified search mode"""
        if not hasattr(self, 'extension_settings'):
            # Initialize with defaults
            self.extension_settings = {
                "base": self.get_default_extensions("base"),
                "avanzata": self.get_default_extensions("avanzata"),
                "profonda": self.get_default_extensions("profonda")
            }
        
        extensions = self.extension_settings.get(mode)
        
        # Se non abbiamo estensioni per questa modalità o la lista è vuota, 
        # usa i valori predefiniti
        if extensions is None or len(extensions) == 0:
            extensions = self.get_default_extensions(mode)
            # Salva queste estensioni predefinite per usi futuri
            if not hasattr(self, 'extension_settings'):
                self.extension_settings = {}
            self.extension_settings[mode] = extensions
        
        self.log_debug(f"Estensioni caricate per modalità {mode}: {', '.join(extensions)}")
        return extensions
        
    @error_handler
    def save_extension_settings(self, mode, extensions):
        """Save extension settings for the specified search mode"""
        if not hasattr(self, 'extension_settings'):
            self.extension_settings = {}
        
        # Ensure all extensions have a leading dot and are lowercase
        normalized_extensions = []
        for ext in extensions:
            if not ext.startswith('.'):
                ext = '.' + ext
            normalized_extensions.append(ext.lower())
        
        self.extension_settings[mode] = normalized_extensions
            
        # Qui potresti aggiungere codice per salvare le impostazioni su file
        self.save_settings_to_file()
    
    @error_handler
    def save_settings_to_file(self):
        """Salva le impostazioni dell'applicazione su un file JSON"""
        try:
            import json
            import os

            # Cartella per le impostazioni
            settings_dir = os.path.join(os.path.expanduser("~"), ".file_search_tool")
            if not os.path.exists(settings_dir):
                os.makedirs(settings_dir)
            
            # File per salvare le impostazioni
            settings_file = os.path.join(settings_dir, "application_settings.json")
            
            # Creazione del dizionario delle impostazioni
            settings = {
                # Impostazioni esistenti
                "extension_settings": getattr(self, 'extension_settings', {}),
                "dir_size_calculation": self.dir_size_calculation.get(),
                "timeout_enabled": self.timeout_enabled.get(),
                "timeout_seconds": self.timeout_seconds.get(),
                "max_files_to_check": self.max_files_to_check.get(),
                "max_results": self.max_results.get(),
                "worker_threads": self.worker_threads.get(),
                "max_file_size_mb": self.max_file_size_mb.get(),
                "use_indexing": self.use_indexing.get(),
                "skip_permission_errors": self.skip_permission_errors.get(),
                
                # Nuove impostazioni per la gestione della RAM
                "auto_memory_management": getattr(self, 'auto_memory_management', True),
                "memory_usage_percent": getattr(self, 'memory_usage_percent', 75),
                
                # Impostazioni di aggiornamento
                "update_settings": getattr(self, 'update_settings', {
                    "auto_update": True,
                    "update_frequency": "All'avvio"
                })
            }
            
            # Salva le impostazioni su file
            with open(settings_file, 'w') as f:
                json.dump(settings, f, indent=4)
            
            self.log_debug(f"Impostazioni dell'applicazione salvate in {settings_file}")
        except Exception as e:
            self.log_error("Errore durante il salvataggio delle impostazioni", e)
    
    @error_handler
    def load_settings_from_file(self):
        """Carica le impostazioni dell'applicazione da un file JSON"""
        try:
            import json
            import os

            # File per le impostazioni
            settings_file = os.path.join(os.path.expanduser("~"), ".file_search_tool", "application_settings.json")

            # Controlla se il file esiste
            if os.path.exists(settings_file):
                with open(settings_file, 'r') as f:
                    settings = json.load(f)
                    self.log_debug(f"Impostazioni dell'applicazione caricate da {settings_file}")

                    # Carica le impostazioni esistenti
                    self.extension_settings = settings.get("extension_settings", {
                        "base": self.get_default_extensions("base"),
                        "avanzata": self.get_default_extensions("avanzata"),
                        "profonda": self.get_default_extensions("profonda")
                    })
                    self.dir_size_calculation.set(settings.get("dir_size_calculation", "disabilitato"))
                    self.timeout_enabled.set(settings.get("timeout_enabled", False))
                    self.timeout_seconds.set(settings.get("timeout_seconds", 3600))
                    self.max_files_to_check.set(settings.get("max_files_to_check", 100000))
                    self.max_results.set(settings.get("max_results", 50000))
                    self.worker_threads.set(settings.get("worker_threads", min(8, os.cpu_count() or 4)))
                    self.max_file_size_mb.set(settings.get("max_file_size_mb", 100))
                    self.use_indexing.set(settings.get("use_indexing", True))
                    self.skip_permission_errors.set(settings.get("skip_permission_errors", True))

                    # Carica le impostazioni per la gestione della memoria
                    self.auto_memory_management = settings.get("auto_memory_management", True)
                    self.memory_usage_percent = settings.get("memory_usage_percent", 75)
                    
                    # Carica le impostazioni per l'aggiornamento
                    self.update_settings = settings.get("update_settings", {
                        "auto_update": True,
                        "update_frequency": "All'avvio",
                        "last_update_check": "Mai"
                    })
                    
                    return settings
            else:
                # Inizializza con valori predefiniti se il file non esiste
                self.log_debug(f"File delle impostazioni non trovato: {settings_file}. Inizializzo con i valori predefiniti.")
                self.extension_settings = {
                    "base": self.get_default_extensions("base"),
                    "avanzata": self.get_default_extensions("avanzata"),
                    "profonda": self.get_default_extensions("profonda")
                }
                self.dir_size_calculation.set("disabilitato")
                self.timeout_enabled.set(False)
                self.timeout_seconds.set(3600)
                self.max_files_to_check.set(100000)
                self.max_results.set(50000)
                self.worker_threads.set(min(8, os.cpu_count() or 4))
                self.max_file_size_mb.set(100)
                self.use_indexing.set(True)
                self.skip_permission_errors.set(True)
                self.auto_memory_management = True
                self.memory_usage_percent = 75
                
                # Inizializza le impostazioni di aggiornamento
                self.update_settings = {
                    "auto_update": True,
                    "update_frequency": "All'avvio",
                    "last_update_check": "Mai"
                }
                
                return {}
        except Exception as e:
            self.log_error("Errore durante il caricamento delle impostazioni", e)
            # Fallback ai valori predefiniti in caso di errore
            self.extension_settings = {
                "base": self.get_default_extensions("base"),
                "avanzata": self.get_default_extensions("avanzata"),
                "profonda": self.get_default_extensions("profonda")
            }
            self.dir_size_calculation.set("disabilitato")
            self.timeout_enabled.set(False)
            self.timeout_seconds.set(3600)
            self.max_files_to_check.set(100000)
            self.max_results.set(50000)
            self.worker_threads.set(min(8, os.cpu_count() or 4))
            self.max_file_size_mb.set(100)
            self.use_indexing.set(True)
            self.skip_permission_errors.set(True)
            self.auto_memory_management = True
            self.memory_usage_percent = 75
            
            # Inizializza le impostazioni di aggiornamento
            self.update_settings = {
                "auto_update": True,
                "update_frequency": "All'avvio",
                "last_update_check": "Mai",
                "include_beta": False
            }
            
            return {}

    @error_handler
    def create_widgets(self):
        # Frame principale che conterrà tutto
        main_container = ttk.Frame(self.root)
        main_container.pack(fill=BOTH, expand=YES)
        
        # Intestazione (titolo e informazioni)
        header_frame = ttk.Frame(main_container, padding="10")
        header_frame.pack(fill=X)

        # Layout a tre colonne in una singola riga
        # 1. Titolo e icona a sinistra
        title_frame = ttk.Frame(header_frame)
        title_frame.pack(side=LEFT, fill=Y)

        try:
            # Usa PIL per il supporto di più formati e il ridimensionamento
            from PIL import Image, ImageTk
            
            # Sostituisci 'logo.png' con il percorso della tua immagine
            current_dir = os.path.dirname(os.path.abspath(__file__))
            image_path = os.path.join(current_dir, "logo.png")
            
            # Carica e ridimensiona l'immagine
            original_image = Image.open(image_path)
            resized_image = original_image.resize((84, 84))
            self.logo_image = ImageTk.PhotoImage(resized_image)
            
            # Crea un label per l'immagine
            logo_label = ttk.Label(title_frame, image=self.logo_image)
            logo_label.pack(side=LEFT, padx=(0, 10))
        except Exception as e:
            self.log_debug(f"Impossibile caricare l'immagine del logo: {str(e)}")
            logo_label = None

        # Container per titolo e nome impilati verticalmente
        title_text_container = ttk.Frame(title_frame)
        title_text_container.pack(side=LEFT)

        title_label = ttk.Label(title_text_container, text=APP_TITLE, 
                            font=("Helvetica", 14, "bold"))
        title_label.pack(anchor=W)

        # Aggiunta del testo "Antonino" sotto il titolo
        antonino_label = ttk.Label(title_text_container, text="APS.QS Antonino Tessio", 
                            font=("Helvetica", 12))
        antonino_label.pack(anchor=W)

        # 2. Tema al centro
        theme_frame = ttk.Frame(header_frame)
        theme_frame.pack(side=LEFT, expand=True, fill=Y)

        ttk.Label(theme_frame, text="Tema:").pack(side=LEFT)
        # Limitiamo i temi disponibili a quelli specificati
        available_themes = ["minty", "cosmo", "darkly", "cyborg"]
        self.theme_combobox = ttk.Combobox(theme_frame, values=available_themes, width=15)
        self.theme_combobox.pack(side=LEFT, padx=5)

        # Carica il tema salvato o usa darkly come predefinito
        saved_theme = self.load_saved_theme()
        if saved_theme in available_themes:
            self.theme_combobox.set(saved_theme)
        else:
            self.theme_combobox.set("darkly")

        self.theme_combobox.bind("<<ComboboxSelected>>", lambda e: [
            ttk.Style().theme_use(self.theme_combobox.get()),
            self.update_theme_colors(self.theme_combobox.get()),
            self.save_theme_preference(self.theme_combobox.get())
        ])
        # 3. Data/ora e utente a destra - rimane invariato
        datetime_frame = ttk.Frame(header_frame)
        datetime_frame.pack(side=RIGHT, fill=Y)

        datetime_label = ttk.Label(datetime_frame, textvariable=self.datetime_var, font=("Helvetica", 9))
        datetime_label.pack(side=RIGHT)

        # ==========================================================
        # SEZIONE DEI CONTROLLI DI RICERCA (organizzati per righe)
        # ==========================================================
        controls_frame = ttk.LabelFrame(main_container, text="Parametri di ricerca", padding=10)
        controls_frame.pack(fill=X, padx=10, pady=5)
        
        # ------------------------------------------------------
        # RIGA 1: Directory di ricerca
        # ------------------------------------------------------
        path_frame = ttk.Frame(controls_frame)
        path_frame.pack(fill=X, pady=5)
        
        path_label = ttk.Label(path_frame, text="Directory:", width=12, anchor=W)
        path_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(path_label, "Seleziona la directory per effettuare la ricerca dei file")
        
        self.path_entry = ttk.Entry(path_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        self.browse_btn = ttk.Button(path_frame, text="📁 Sfoglia", command=self.browse_directory, width=10)
        self.browse_btn.pack(side=LEFT)
        
        # ------------------------------------------------------
        # RIGA 2: Parole chiave
        # ------------------------------------------------------
        keyword_frame = ttk.Frame(controls_frame)
        keyword_frame.pack(fill=X, pady=5)
        
        keyword_label = ttk.Label(keyword_frame, text="Parole chiave:", width=12, anchor=W)
        keyword_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(keyword_label, "Per la ricerca di più parole usa la virgola. Esempio: documento, fattura, contratto\n"
                                        "Attiva 'Parola intera' per cercare 'log' senza trovare 'login' o 'logo'")
        
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)

        # ------------------------------------------------------
        # RIGA 3: Opzioni base di ricerca (checkbox)
        # ------------------------------------------------------
        options_frame = ttk.Frame(controls_frame)
        options_frame.pack(fill=X, pady=5)

        # Prima parte: Livelli di ricerca
        search_options = ttk.Frame(options_frame)
        search_options.pack(side=LEFT, fill=Y)

        ttk.Label(search_options, text="Livelli di ricerca:").pack(side=LEFT, padx=(0, 5))
        search_depth_combo = ttk.Combobox(search_options, textvariable=self.search_depth, 
                                        values=["base", "avanzata", "profonda"], 
                                        width=10, state="readonly")
        search_depth_combo.pack(side=LEFT, padx=5)
        search_depth_combo.current(0)

        extensions_btn = ttk.Button(search_options, text="📄 Configura estensioni", 
                            command=lambda: self.configure_extensions(self.search_depth.get()))
        extensions_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(extensions_btn, "Configura quali estensioni di file includere nella ricerca")

        # Pulsante impostazioni avanzate
        advanced_options_btn = ttk.Button(search_options, text="⚙ Impostazioni avanzate", 
                                        command=self.show_advanced_options)
        advanced_options_btn.pack(side=LEFT, padx=10)
        self.create_tooltip(advanced_options_btn, "Configura tutte le impostazioni avanzate (profondità, filtri, esclusioni, performance)")

        # Separatore per creare più spazio
        separator = ttk.Frame(options_frame, width=30)
        separator.pack(side=LEFT)

        # Pulsanti di azione (spostati qui dopo lo spazio)
        action_buttons = ttk.Frame(options_frame)
        action_buttons.pack(side=LEFT, fill=Y)

        # Pulsante di ricerca (principale)
        self.search_button = ttk.Button(action_buttons, text="🔍 CERCA", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=15)
        self.search_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.search_button, "Avvia la ricerca con i criteri specificati")

        # Pulsante per interrompere la ricerca
        self.stop_button = ttk.Button(action_buttons, text="⏹️ Interrompi ricerca",
                                    command=self.stop_search_process,
                                    style="danger.TButton", width=20,
                                    state="disabled")
        self.stop_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.stop_button, "Ferma immediatamente la ricerca in corso")

        # Pulsante per pulire i campi di ricerca
        self.clear_btn = ttk.Button(action_buttons, text="🧹 Pulisci campi", 
                    command=lambda: [self.search_path.set(""), self.keywords.set("")],
                    style="secondary.Outline.TButton", width=15)
        self.clear_btn.pack(side=LEFT, padx=10)
        self.create_tooltip(self.clear_btn, "Cancella i campi di ricerca")

        # Pulsante admin solo su Windows
        if os.name == 'nt':
            self.admin_button = ttk.Button(action_buttons, text="🛡️Avvia Admin", 
                                    command=self.restart_as_admin,
                                    style="info.Outline.TButton", width=20)
            self.admin_button.pack(side=LEFT, padx=0)
            
            # Disabilita il pulsante se l'app è già avviata come amministratore
            if self.is_admin:
                self.admin_button.config(state="disabled")
                self.create_tooltip(self.admin_button, "L'applicazione è già in esecuzione come amministratore")
            else:
                self.create_tooltip(self.admin_button, "Riavvia l'applicazione con privilegi di amministratore per accedere a tutte le cartelle")
        
        # ==========================================================
        # SEZIONE INFORMAZIONI TEMPORALI E SPAZIO DISCO (SOPRA LA LISTA RISULTATI)
        # ==========================================================
        info_bar = ttk.Frame(main_container)
        info_bar.pack(fill=X, padx=10, pady=5)
        
        # Frame per le informazioni temporali a sinistra
        time_frame = ttk.LabelFrame(info_bar, text="Informazioni temporali", padding=5)
        time_frame.pack(side=LEFT, fill=X, padx=(0, 5), expand=YES)
        
        time_grid = ttk.Frame(time_frame)
        time_grid.pack(fill=X, pady=2)
        
        ttk.Label(time_grid, text="Avvio:").grid(row=0, column=0, sticky=W, padx=5)
        self.start_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.start_time_label.grid(row=0, column=1, sticky=W, padx=5)
        
        ttk.Label(time_grid, text="Fine:").grid(row=0, column=2, sticky=W, padx=5)
        self.end_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.end_time_label.grid(row=0, column=3, sticky=W, padx=5)
        
        ttk.Label(time_grid, text="Durata:").grid(row=0, column=4, sticky=W, padx=5)
        self.total_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.total_time_label.grid(row=0, column=5, sticky=W, padx=5)
        
        # Frame per info disco a destra
        disk_info_frame = ttk.LabelFrame(info_bar, text="Spazio disco", padding=5)
        disk_info_frame.pack(side=LEFT, fill=X, expand=YES)

        # Uso layout semplice con pack per ordinare gli elementi in linea
        disk_grid = ttk.Frame(disk_info_frame)
        disk_grid.pack(fill=X, pady=2)

        # 1. Usato
        ttk.Label(disk_grid, text="Usato:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.used_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 2. Libero
        ttk.Label(disk_grid, text="Libero:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.free_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 3. Totale
        ttk.Label(disk_grid, text="Totale:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.total_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 4. Directory
        ttk.Label(disk_grid, text="Directory:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.dir_size_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 5. Pulsante aggiorna
        refresh_size_btn = ttk.Button(disk_grid, text="🔄 Aggiorna", command=self.refresh_directory_size, 
                                width=12, style="info.TButton")
        refresh_size_btn.pack(side=LEFT, padx=5)
        
        # ==========================================================
        # SEZIONE STATO RICERCA
        # ==========================================================
        # Contenitore per le informazioni di stato sopra i risultati
        status_frame = ttk.LabelFrame(main_container, text="Stato ricerca", padding=5)
        status_frame.pack(fill=X, padx=10, pady=5)

        # Progress bar
        self.progress_bar = ttk.Progressbar(status_frame, mode='determinate')
        self.progress_bar.pack(fill=X, pady=5)

        # Status grid with improved layout
        status_grid = ttk.Frame(status_frame)
        status_grid.pack(fill=X)

        # Row 1: Status with debug button
        status_row = ttk.Frame(status_grid)
        status_row.pack(fill=X, pady=2)

        ttk.Label(status_row, text="Analisi:", width=12, anchor=W, font=("", 9, "bold")).pack(side=LEFT, padx=5)
        self.status_label = ttk.Label(status_row, text="In attesa...", wraplength=1000)
        self.status_label.pack(side=LEFT, fill=X, expand=YES, padx=5)

        # Add Debug Button on the right side of the status row
        self.debug_button = ttk.Button(status_row, text="📊 Debug Log", command=self.show_debug_log, 
                                    style="info.Outline.TButton", width=12)
        self.debug_button.pack(side=RIGHT, padx=5)
        self.create_tooltip(self.debug_button, "Mostra la finestra di debug con log dettagliati sulla ricerca in corso")

        # Row 2: File analyzed
        files_row = ttk.Frame(status_grid)
        files_row.pack(fill=X, pady=2)

        ttk.Label(files_row, text="File analizzati:", width=12, anchor=W, font=("", 9, "bold")).pack(side=LEFT, padx=5)
        self.analyzed_files_label = ttk.Label(files_row, text="Nessuna ricerca avviata", wraplength=2000)
        self.analyzed_files_label.pack(side=LEFT, fill=X, expand=YES, padx=5)
        # ==========================================================
        # SEZIONE RISULTATI 
        # ==========================================================
        results_container = ttk.LabelFrame(main_container, text="Risultati di ricerca", padding=10)
        results_container.pack(fill=BOTH, expand=YES, padx=10, pady=(5, 10))
        
        # Frame per i pulsanti di azione sui risultati
        actions_frame = ttk.Frame(results_container)
        actions_frame.pack(fill=X, pady=(0, 5))

        # Pulsanti per la selezione
        selection_frame = ttk.Frame(actions_frame)
        selection_frame.pack(side=LEFT)

        select_all_btn = ttk.Button(selection_frame, text="☑ Seleziona tutto", command=self.select_all)
        select_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(select_all_btn, "Seleziona tutti i risultati nella lista")

        deselect_all_btn = ttk.Button(selection_frame, text="☐ Deseleziona tutto", command=self.deselect_all)
        deselect_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(deselect_all_btn, "Deseleziona tutti i risultati")

        # Crea frame centrale che si espande per riempire lo spazio disponibile
        center_frame = ttk.Frame(actions_frame)
        center_frame.pack(side=LEFT, fill=X, expand=YES)

        # Creiamo un sotto-frame per contenere entrambi i label affiancati
        labels_frame = ttk.Frame(center_frame)
        labels_frame.pack(anchor=CENTER)

        # Label per la dimensione totale (esistente)
        self.total_files_size_label = ttk.Label(labels_frame, text="Dimensione totale: 0  (0 file)", font=("", 9, "bold"))
        self.total_files_size_label.pack(side=LEFT, padx=10)

        # NUOVO: Label per la dimensione dei file selezionati
        self.selected_files_size_label = ttk.Label(labels_frame, text="Selezionati: 0  (0 file)", font=("", 9, "bold"))
        self.selected_files_size_label.pack(side=LEFT, padx=10)

        # Pulsanti per le azioni
        action_frame = ttk.Frame(actions_frame)
        action_frame.pack(side=RIGHT)

        self.copy_button = ttk.Button(action_frame, text="📋 Copia selezionati",
                                    command=self.copy_selected,
                                    style="TButton")
        self.copy_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.copy_button, "Copia i file selezionati nella directory specificata")

        self.compress_button = ttk.Button(action_frame, text="📦 Comprimi selezionati",
                                        command=self.compress_selected,
                                        style="TButton")
        self.compress_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.compress_button, "Comprimi i file selezionati in un archivio ZIP")

        self.view_log_button = ttk.Button(action_frame, text="🚫 Visualizza file esclusi",
                                        command=self.view_skipped_files_log,
                                        style="secondary.TButton")
        self.view_log_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.view_log_button, "Visualizza il log dei file esclusi dalla ricerca")

        # TreeView con scrollbar
        treeview_container = ttk.Frame(results_container)
        treeview_container.pack(fill=BOTH, expand=True)

        # Creazione della TreeView
        self.results_list = ttk.Treeview(treeview_container, selectmode="extended",
                                columns=("type", "attachment", "size", "modified", "created", "author", "path"), 
                                show="headings")

        # Creazione delle scrollbar
        vsb = ttk.Scrollbar(treeview_container, orient="vertical", command=self.results_list.yview)
        hsb = ttk.Scrollbar(treeview_container, orient="horizontal", command=self.results_list.xview)

        # Configurazione delle scrollbar nella TreeView
        self.results_list.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

        # Posizionamento con grid
        self.results_list.grid(column=0, row=0, sticky='nsew')
        vsb.grid(column=1, row=0, sticky='ns')
        hsb.grid(column=0, row=1, sticky='ew')

        # Configura grid layout per l'espansione
        treeview_container.grid_columnconfigure(0, weight=1)
        treeview_container.grid_rowconfigure(0, weight=1)

        # Impostazione delle intestazioni delle colonne
        self.results_list.heading("type", text="Tipo")
        self.results_list.heading("attachment", text="Allegato")
        self.results_list.heading("size", text="Dimensione")
        self.results_list.heading("modified", text="Modificato")
        self.results_list.heading("created", text="Creato")
        self.results_list.heading("author", text="Nome")
        self.results_list.heading("path", text="Percorso")

        # Impostazione delle intestazioni delle colonne con funzione di ordinamento
        self.results_list.heading("type", text="Tipo", 
                            command=lambda: self.treeview_sort_column(self.results_list, "type", False))
        self.results_list.heading("size", text="Dimensione", 
                            command=lambda: self.treeview_sort_column(self.results_list, "size", False))
        self.results_list.heading("modified", text="Modificato", 
                            command=lambda: self.treeview_sort_column(self.results_list, "modified", False))
        self.results_list.heading("created", text="Creato", 
                            command=lambda: self.treeview_sort_column(self.results_list, "created", False))
        self.results_list.heading("author", text="Nome", 
                            command=lambda: self.treeview_sort_column(self.results_list, "author", False))
        self.results_list.heading("path", text="Percorso", 
                            command=lambda: self.treeview_sort_column(self.results_list, "path", False))

        # Impostazione delle larghezze fisse delle colonne
        self.results_list.column("type", width=120, minwidth=50, stretch=NO, anchor="center")
        self.results_list.column("attachment", width=70, minwidth=30, stretch=NO, anchor="center")
        self.results_list.column("size", width=100, minwidth=80, stretch=NO, anchor="center")
        self.results_list.column("modified", width=150, minwidth=120, stretch=NO, anchor="center")
        self.results_list.column("created", width=150, minwidth=120, stretch=NO, anchor="center")
        self.results_list.column("author", width=400, minwidth=80, stretch=NO, anchor="w")
        self.results_list.column("path", width=600, minwidth=200, stretch=YES, anchor="w")

        # Aggiungi binding per l'evento di doppio clic
        self.results_list.bind("<Double-1>", self.open_file_location)

        # Aggiungi binding per l'evento di selezione
        self.results_list.bind("<<TreeviewSelect>>", self.update_selected_files_size)

        # Auto-focus sull'entry del percorso all'avvio
        self.path_entry.focus_set()

    @error_handler
    def change_theme(self, theme):
        """Cambia il tema dell'interfaccia e salva la preferenza"""
        # Verifica che il tema sia valido
        valid_themes = ["minty", "cosmo", "darkly", "cyborg"]
        if theme not in valid_themes:
            self.log_debug(f"Tema '{theme}' non valido, uso darkly come predefinito")
            theme = "darkly"
        
        # Applica il tema a livello di stile ttk
        ttk.Style().theme_use(theme)
        
        # Aggiorna i colori personalizzati
        self.update_theme_colors(theme)
        
        # Aggiorna la selezione nella combobox se esiste
        if hasattr(self, 'theme_combobox') and self.theme_combobox is not None:
            self.theme_combobox.set(theme)
        
        # Salva la preferenza
        self.save_theme_preference(theme)

    @error_handler
    def save_theme_preference(self, theme):
        """Salva la preferenza del tema nelle impostazioni"""
        try:
            settings_file = os.path.join(os.path.expanduser("~"), ".file_search_settings.json")
            
            # Carica le impostazioni esistenti o crea un nuovo dizionario
            if os.path.exists(settings_file):
                with open(settings_file, 'r', encoding='utf-8') as f:
                    settings = json.load(f)
            else:
                settings = {}
            
            # Aggiorna o aggiungi la preferenza del tema
            settings["theme"] = theme
            
            # Salva le impostazioni aggiornate
            with open(settings_file, 'w', encoding='utf-8') as f:
                json.dump(settings, f, indent=2)
            
            self.log_debug(f"Tema '{theme}' salvato nelle impostazioni")
        except Exception as e:
            self.log_error(f"Impossibile salvare la preferenza del tema: {str(e)}")

    @error_handler
    def load_saved_theme(self):
        """Carica il tema salvato dalle impostazioni"""
        try:
            settings_file = os.path.join(os.path.expanduser("~"), ".file_search_settings.json")
            
            if os.path.exists(settings_file):
                with open(settings_file, 'r', encoding='utf-8') as f:
                    settings = json.load(f)
                
                if "theme" in settings:
                    self.log_debug(f"Tema caricato dalle impostazioni: {settings['theme']}")
                    return settings["theme"]
            
            # Se non è stato trovato alcun tema salvato, restituisci None
            return None
        except Exception as e:
            self.log_error(f"Impossibile caricare la preferenza del tema: {str(e)}")
            return None
    
    @error_handler
    def show_advanced_options(self):
        """Mostra una finestra di dialogo unificata per tutte le opzioni avanzate"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Impostazioni avanzate")
        dialog.withdraw()
        dialog.geometry("800x800")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Crea un notebook con schede per le diverse categorie di opzioni
        notebook = ttk.Notebook(main_frame)
        notebook.pack(fill=BOTH, expand=YES, pady=10)
        
        # ================= Scheda 1: Opzioni di ricerca =================
        search_options_frame = ttk.Frame(notebook, padding=15)
        notebook.add(search_options_frame, text="Opzioni di ricerca")
        
        # Profondità di ricerca
        depth_frame = ttk.LabelFrame(search_options_frame, text="Profondità di ricerca", padding=10)
        depth_frame.pack(fill=X, pady=10)
        
        depth_desc = ttk.Label(depth_frame, 
                        text="Imposta il numero massimo di livelli di cartella da esplorare.\n"
                            "Il valore 0 indica profondità illimitata.",
                        wraplength=700)
        depth_desc.pack(anchor=W, pady=(0, 10))
        
        depth_control = ttk.Frame(depth_frame)
        depth_control.pack(fill=X)
        
        ttk.Label(depth_control, text="Profondità:").pack(side=LEFT)
        
        # Usa variabili per i controlli
        depth_var = IntVar(value=self.max_depth)
        depth_spinbox = ttk.Spinbox(depth_control, from_=0, to=20, width=3, textvariable=depth_var)
        depth_spinbox.pack(side=LEFT, padx=5)
        ttk.Label(depth_control, text="(0 = illimitata)", foreground="gray").pack(side=LEFT)
        
        # AGGIUNTO TOOLTIP
        self.create_tooltip(depth_spinbox, 
                    "Controlla quanto a fondo cercare nelle sottocartelle.\n"
                    "Un valore più alto aumenta il tempo di ricerca ma trova file in cartelle più profonde.\n"
                    "Il valore 0 cerca in tutte le sottocartelle senza limiti di profondità.")
        
        # Contenuti da cercare
        content_frame = ttk.LabelFrame(search_options_frame, text="Contenuti da cercare", padding=10)
        content_frame.pack(fill=X, pady=10)

        search_files_var = BooleanVar(value=self.search_files.get())
        search_files_cb = ttk.Checkbutton(content_frame, text="Cerca nei file", variable=search_files_var)
        search_files_cb.pack(anchor=W, pady=2)
        self.create_tooltip(search_files_cb, "Attiva per includere i file nei risultati di ricerca")
        
        search_folders_var = BooleanVar(value=self.search_folders.get())
        search_folders_cb = ttk.Checkbutton(content_frame, text="Cerca nelle cartelle", variable=search_folders_var)
        search_folders_cb.pack(anchor=W, pady=2)
        self.create_tooltip(search_folders_cb, "Attiva per includere le cartelle nei risultati di ricerca")
        
        search_content_var = BooleanVar(value=self.search_content.get())
        search_content_cb = ttk.Checkbutton(content_frame, text="Cerca nei contenuti dei file", variable=search_content_var)
        search_content_cb.pack(anchor=W, pady=2)
        self.create_tooltip(search_content_cb, 
                        "Attiva per cercare le parole chiave all'interno dei file.\n"
                        "Questo rallenta la ricerca ma permette di trovare i file\n"
                        "che contengono il testo cercato anche se non è nel nome.")
        
        whole_word_var = BooleanVar(value=self.whole_word_search.get())
        whole_word_cb = ttk.Checkbutton(content_frame, text="Cerca parole intere", variable=whole_word_var)
        whole_word_cb.pack(anchor=W, pady=2)
        self.create_tooltip(whole_word_cb, 
                        "Quando attivato, cerca solo corrispondenze esatte delle parole.\n"
                        "Esempio: cercando 'log' NON troverà 'login' o 'catalogo'.")
        
        # Windows Search Frame - Da aggiungere dopo content_frame e prima della definizione della scheda "Filtri avanzati"
        windows_search_frame = ttk.LabelFrame(search_options_frame, text="Windows Search (windows.edb)", padding=10)
        windows_search_frame.pack(fill=X, pady=10)

        # Verifica se WindowsSearchHelper è disponibile
        has_windows_search = False
        try:
            import win32com.client
            import pythoncom
            has_windows_search = True
        except ImportError:
            has_windows_search = False

        # Variabile per l'utilizzo di Windows Search
        if not hasattr(self, 'use_windows_search_var'):
            self.use_windows_search_var = tk.BooleanVar(value=True)

        # Checkbox per abilitare/disabilitare Windows Search
        use_windows_search_cb = ttk.Checkbutton(
            windows_search_frame, 
            text="Utilizza Windows Search quando disponibile (accelera la ricerca)",
            variable=self.use_windows_search_var,
            state="normal" if has_windows_search else "disabled"
        )
        use_windows_search_cb.pack(anchor=W, padx=10, pady=5)
        self.create_tooltip(use_windows_search_cb, 
                        "Windows Search utilizza l'indice del sistema per accelerare significativamente\n"
                        "la ricerca di file in cartelle indicizzate.\n"
                        "Funziona solo per percorsi che sono già indicizzati da Windows Search.")

        # Informazioni sullo stato del servizio Windows Search
        windows_search_status = ttk.Label(windows_search_frame, text="Verifica stato di Windows Search...")
        windows_search_status.pack(anchor=W, padx=10, pady=5)

        # Verifica lo stato in background
        def check_windows_search_status():
            if not has_windows_search:
                windows_search_status.config(
                    text="❌ Libreria pywin32 non installata. Usa 'pip install pywin32' per abilitare questa funzionalità.",
                    foreground="red"
                )
                return
                
            # Verifica se il servizio Windows Search è disponibile
            try:
                pythoncom.CoInitialize()
                connection = win32com.client.Dispatch("ADODB.Connection")
                connection.Open("Provider=Search.CollatorDSO;Extended Properties='Application=Windows';")
                connection.Close()
                pythoncom.CoUninitialize()
                
                windows_search_status.config(
                    text="✅ Servizio Windows Search disponibile",
                    foreground="green"
                )
            except Exception:
                windows_search_status.config(
                    text="❌ Servizio Windows Search non disponibile o non attivo",
                    foreground="red"
                )
                pythoncom.CoUninitialize()

        # Chiamata asincrona per non bloccare l'interfaccia
        if has_windows_search:
            dialog.after(100, check_windows_search_status)
        else:
            check_windows_search_status()

        # Informazioni aggiuntive
        windows_search_info = ttk.Label(
            windows_search_frame,
            text=(
                "Nota: Windows Search funziona solo su cartelle già indicizzate dal sistema.\n"
                "La ricerca può essere fino a 10-20 volte più veloce quando questa opzione è attiva."
            ),
            wraplength=650
        )
        windows_search_info.pack(anchor=W, padx=10, pady=5)

        # Pulsante per aprire le impostazioni di indicizzazione di Windows
        def open_indexing_settings():
            os.system("control.exe srchadmin.dll")

        indexing_btn = ttk.Button(
            windows_search_frame,
            text="Apri impostazioni di indicizzazione di Windows",
            command=open_indexing_settings
        )
        indexing_btn.pack(anchor=W, padx=10, pady=10)

        # ================= Scheda 2: Filtri avanzati =================
        filters_frame = ttk.Frame(notebook, padding=15)
        notebook.add(filters_frame, text="Filtri avanzati")
        
        # Filtri dimensione
        size_frame = ttk.LabelFrame(filters_frame, text="Dimensione file", padding=10)
        size_frame.pack(fill=X, pady=10)
        
        ttk.Label(size_frame, text="Filtra i file in base alla dimensione:").pack(anchor=W, pady=(0, 10))
        
        size_grid = ttk.Frame(size_frame)
        size_grid.pack(fill=X)
        
        size_min_label = ttk.Label(size_grid, text="Dimensione minima (KB):")
        size_min_label.grid(row=0, column=0, padx=5, pady=5, sticky=W)
        min_size_var = StringVar(value=str(self.advanced_filters["size_min"] // 1024))
        min_size = ttk.Entry(size_grid, width=10, textvariable=min_size_var)
        min_size.grid(row=0, column=1, padx=5, pady=5)
        self.create_tooltip(min_size, 
                    "Filtra i file più piccoli della dimensione specificata in KB.\n"
                    "Esempio: impostando 100 verranno ignorati i file minori di 100KB.")
        
        size_max_label = ttk.Label(size_grid, text="Dimensione massima (KB):")
        size_max_label.grid(row=1, column=0, padx=5, pady=5, sticky=W)
        max_size_var = StringVar(value=str(self.advanced_filters["size_max"] // 1024 if self.advanced_filters["size_max"] else 0))
        max_size = ttk.Entry(size_grid, width=10, textvariable=max_size_var)
        max_size.grid(row=1, column=1, padx=5, pady=5)
        self.create_tooltip(max_size, 
                    "Filtra i file più grandi della dimensione specificata in KB.\n"
                    "Esempio: impostando 1000 verranno ignorati i file maggiori di 1000KB (1MB).\n"
                    "Impostare a 0 per nessun limite.")
        
        # Filtri data
        date_frame = ttk.LabelFrame(filters_frame, text="Data di modifica", padding=10)
        date_frame.pack(fill=X, pady=10)
        
        ttk.Label(date_frame, text="Filtra i file in base alla data di modifica:").pack(anchor=W, pady=(0, 10))
        
        date_grid = ttk.Frame(date_frame)
        date_grid.pack(fill=X)
        
        date_min_label = ttk.Label(date_grid, text="Data inizio (DD-MM-YYYY):")
        date_min_label.grid(row=0, column=0, padx=5, pady=5, sticky=W)
        min_date = ttk.DateEntry(date_grid, dateformat="%d-%m-%Y")
        min_date.grid(row=0, column=1, padx=5, pady=5)
        min_date.entry.delete(0, 'end')
        if self.advanced_filters["date_min"]:
            min_date.entry.insert(0, self.advanced_filters["date_min"])
        self.create_tooltip(min_date, 
                    "Includi solo file modificati dopo questa data.\n"
                    "Lascia vuoto per nessun limite di data minima.")
        
        date_max_label = ttk.Label(date_grid, text="Data fine (DD-MM-YYYY):")
        date_max_label.grid(row=1, column=0, padx=5, pady=5, sticky=W)
        max_date = ttk.DateEntry(date_grid, dateformat="%d-%m-%Y")
        max_date.grid(row=1, column=1, padx=5, pady=5)
        max_date.entry.delete(0, 'end')
        if self.advanced_filters["date_max"]:
            max_date.entry.insert(0, self.advanced_filters["date_max"])
        self.create_tooltip(max_date, 
                    "Includi solo file modificati prima di questa data.\n"
                    "Lascia vuoto per nessun limite di data massima.")
        
        # ================= Scheda 3: Gestione esclusioni =================
        exclusions_frame = ttk.Frame(notebook, padding=15)
        notebook.add(exclusions_frame, text="Esclusioni")

        ttk.Label(exclusions_frame, text="Aggiungi cartelle da escludere dalla ricerca:", wraplength=700).pack(anchor=W, pady=(0, 10))

        system_exclusions_frame = ttk.Frame(exclusions_frame)
        system_exclusions_frame.pack(fill=X, pady=5)

        # Definizione dei percorsi di sistema
        system_paths = [
            "C:/Windows", 
            "C:/Program Files", 
            "C:/Program Files (x86)",
            "C:/ProgramData",
            "C:/Users/All Users",
            "C:/Program Files (x86)/Client Active Directory Rights Management Services"
        ]

        system_exclusions_var = tk.BooleanVar(value=False)

        # Imposta il valore iniziale della checkbox in base alle esclusioni attuali
        if hasattr(self, 'excluded_paths') and all(path in self.excluded_paths for path in system_paths):
            system_exclusions_var.set(True)

        # Funzione per gestire il toggle delle esclusioni di sistema
        def toggle_system_exclusions():
            if not hasattr(self, 'excluded_paths'):
                self.excluded_paths = []
                
            if system_exclusions_var.get():
                # Aggiungi percorsi di sistema
                added = 0
                for path in system_paths:
                    if path not in self.excluded_paths:
                        self.excluded_paths.append(path)
                        added += 1
                
                # Aggiorna la lista visualizzata
                if added > 0:
                    for item in excluded_list.get_children():
                        excluded_list.delete(item)
                    for path in self.excluded_paths:
                        excluded_list.insert("", "end", values=(path,))
                    messagebox.showinfo("Esclusioni sistema", f"Aggiunti {added} percorsi di sistema alle esclusioni")
            else:
                # Rimuovi percorsi di sistema
                removed = 0
                for path in system_paths[:]:
                    if path in self.excluded_paths:
                        self.excluded_paths.remove(path)
                        removed += 1
                
                # Aggiorna la lista visualizzata
                if removed > 0:
                    for item in excluded_list.get_children():
                        excluded_list.delete(item)
                    for path in self.excluded_paths:
                        excluded_list.insert("", "end", values=(path,))
                    messagebox.showinfo("Esclusioni sistema", f"Rimossi {removed} percorsi di sistema dalle esclusioni")

        # Checkbox per esclusioni di sistema
        system_check = ttk.Checkbutton(system_exclusions_frame, text="Escludi cartelle di sistema (Windows, Program Files, ecc.)", 
            variable=system_exclusions_var, command=toggle_system_exclusions)
        system_check.pack(anchor=W)
        self.create_tooltip(system_check, 
                            "Quando selezionato, esclude automaticamente cartelle di sistema come:\n"
                            "- Windows\n- Program Files\n- ProgramData\n"
                            "Questo velocizzerà notevolmente le ricerche sui dischi di sistema.")

        # Lista dei percorsi esclusi
        list_frame = ttk.Frame(exclusions_frame)
        list_frame.pack(fill=BOTH, expand=YES, pady=5)  # Cambiato expand da NO a YES

        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side=RIGHT, fill=Y)

        # Aggiungi il parametro height per specificare il numero di righe visibili
        excluded_list = ttk.Treeview(list_frame, columns=("path",), show="headings", 
                                yscrollcommand=scrollbar.set, selectmode="extended", height=15)  # Aggiunto height=15
        excluded_list.heading("path", text="Percorso")
        excluded_list.column("path", width=450)
        excluded_list.pack(fill=BOTH, expand=YES)
        self.create_tooltip(excluded_list, 
                    "Elenco di cartelle escluse dalla ricerca.\n"
                    "Le sottocartelle di questi percorsi non verranno analizzate.")
        
        scrollbar.config(command=excluded_list.yview)
        
        # Aggiungi i percorsi esclusi alla lista
        if hasattr(self, 'excluded_paths'):
            for path in self.excluded_paths:
                excluded_list.insert("", "end", values=(path,))
        
        # Frame per aggiungere nuovi percorsi
        add_frame = ttk.Frame(exclusions_frame)
        add_frame.pack(fill=X, pady=10)
        
        path_var = StringVar()
        path_entry = ttk.Entry(add_frame, textvariable=path_var, width=50)
        path_entry.pack(side=LEFT, padx=(0, 5), fill=X, expand=YES)
        self.create_tooltip(path_entry, 
                    "Inserisci il percorso completo della cartella da escludere.\n"
                    "Esempio: C:/Windows o C:/Program Files")
        
        def browse_exclude():
            directory = filedialog.askdirectory()
            if directory:
                path_var.set(directory)
        
        browse_btn = ttk.Button(add_frame, text="Sfoglia", command=browse_exclude)
        browse_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(browse_btn, "Seleziona la cartella da escludere usando una finestra di dialogo")
        
        def add_exclusion():
            path = path_var.get().strip()
            if path:
                excluded_list.insert("", "end", values=(path,))
                path_var.set("")
        
        add_btn = ttk.Button(add_frame, text="Aggiungi", command=add_exclusion)
        add_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(add_btn, "Aggiungi il percorso specificato alla lista delle esclusioni")
        
        # Pulsante per rimuovere elementi selezionati
        def remove_selected():
            selected = excluded_list.selection()
            if selected:
                for item in selected:
                    excluded_list.delete(item)
        
        remove_btn = ttk.Button(exclusions_frame, text="Rimuovi selezionati", command=remove_selected)
        remove_btn.pack(anchor=W, pady=5)
        self.create_tooltip(remove_btn, "Rimuovi i percorsi selezionati dalla lista delle esclusioni")

        
        # ================= Scheda 4: Opzioni a blocchi =================
        blocks_frame = ttk.Frame(notebook, padding=15)
        notebook.add(blocks_frame, text="Opzioni a blocchi")
        
        desc_label = ttk.Label(blocks_frame, text="La ricerca a blocchi divide le cartelle in unità di lavoro più piccole per migliorare le prestazioni e la reattività.",
                        wraplength=700)
        desc_label.pack(fill=X, pady=(0, 15))
        
        # Dimensione blocchi
        size_frame = ttk.LabelFrame(blocks_frame, text="Dimensione e parallelismo", padding=10)
        size_frame.pack(fill=X, pady=10)
        
        size_grid = ttk.Frame(size_frame)
        size_grid.pack(fill=X)
        
        files_block_label = ttk.Label(size_grid, text="Max file per blocco:")
        files_block_label.grid(row=0, column=0, sticky=W, padx=5, pady=5)
        files_block_var = IntVar(value=self.max_files_per_block.get())
        files_block = ttk.Spinbox(size_grid, from_=100, to=10000, increment=100, width=7, textvariable=files_block_var)
        files_block.grid(row=0, column=1, padx=5, pady=5)
        self.create_tooltip(files_block, 
                    "Numero massimo di file da processare in un singolo blocco.\n"
                    "Valori più bassi aumentano la reattività dell'interfaccia ma\n"
                    "possono ridurre leggermente la velocità di ricerca complessiva.")
        
        parallel_label = ttk.Label(size_grid, text="Blocchi paralleli:")
        parallel_label.grid(row=0, column=2, sticky=W, padx=(20, 5), pady=5)
        parallel_var = IntVar(value=self.max_parallel_blocks.get())
        parallel = ttk.Spinbox(size_grid, from_=1, to=16, increment=1, width=5, textvariable=parallel_var)
        parallel.grid(row=0, column=3, padx=5, pady=5)
        self.create_tooltip(parallel, 
                    "Numero di blocchi da elaborare contemporaneamente.\n"
                    "Aumentare questo valore su sistemi con molte CPU può\n"
                    "velocizzare la ricerca, ma consuma più risorse.")
        
        # Opzioni aggiuntive
        options_frame = ttk.LabelFrame(blocks_frame, text="Ottimizzazioni", padding=10)
        options_frame.pack(fill=X, pady=10)
        
        auto_adjust_var = BooleanVar(value=self.block_size_auto_adjust.get())
        auto_adjust_cb = ttk.Checkbutton(options_frame, text="Adatta automaticamente la dimensione dei blocchi", 
                    variable=auto_adjust_var)
        auto_adjust_cb.pack(anchor=W, pady=2)
        self.create_tooltip(auto_adjust_cb, 
                    "Quando attivato, la dimensione dei blocchi viene adattata\n"
                    "automaticamente in base al tipo di ricerca e alla velocità del sistema.\n"
                    "Utile per bilanciare prestazioni e reattività.")
        
        prioritize_var = BooleanVar(value=self.prioritize_user_folders.get())
        prioritize_cb = ttk.Checkbutton(options_frame, text="Dare priorità alle cartelle utente", 
                    variable=prioritize_var)
        prioritize_cb.pack(anchor=W, pady=2)
        self.create_tooltip(prioritize_cb, 
                    "Analizza prima le cartelle più importanti come Documenti, Desktop,\n"
                    "per trovare rapidamente i file più rilevanti.\n"
                    "Utile quando si cercano file personali.")
        
        # ================= Scheda 5: Performance =================
        performance_container = ttk.Frame(notebook)
        notebook.add(performance_container, text="Performance")
        
        # Creiamo il canvas con scrollbar
        perf_canvas = tk.Canvas(performance_container)
        perf_scrollbar = ttk.Scrollbar(performance_container, orient="vertical", command=perf_canvas.yview)
        perf_canvas.configure(yscrollcommand=perf_scrollbar.set)
        
        # Posizionamento del canvas e della scrollbar
        perf_scrollbar.pack(side=RIGHT, fill=Y)
        perf_canvas.pack(side=LEFT, fill=BOTH, expand=YES)
        
        # Frame interno che conterrà tutti i controlli
        performance_frame = ttk.Frame(perf_canvas, padding=15)
        
        # Crea un ID per il frame all'interno del canvas
        perf_canvas_frame = perf_canvas.create_window((0, 0), window=performance_frame, anchor="nw")
        
        # Funzione per aggiornare la regione di scorrimento
        def configure_perf_scroll_region(event):
            perf_canvas.configure(scrollregion=perf_canvas.bbox("all"))
            
        # Funzione per aggiornare la larghezza del frame nel canvas
        def configure_perf_canvas(event):
            perf_canvas.itemconfig(perf_canvas_frame, width=event.width)
            
        # Configura eventi per il ridimensionamento
        performance_frame.bind("<Configure>", configure_perf_scroll_region)
        perf_canvas.bind("<Configure>", configure_perf_canvas)

        # Timeout e limiti
        timeout_frame = ttk.LabelFrame(performance_frame, text="Timeout e limiti", padding=10)
        timeout_frame.pack(fill=X, pady=10)

        # Rimuoviamo il posizionamento separato del checkbox e lo inseriamo direttamente nella griglia
        timeout_grid = ttk.Frame(timeout_frame)
        timeout_grid.pack(fill=X, pady=5)

        # Riga 0: Checkbox e secondi nella stessa riga
        timeout_enabled_var = BooleanVar(value=self.timeout_enabled.get())
        timeout_check = ttk.Checkbutton(timeout_grid, text="Attiva timeout ricerca", variable=timeout_enabled_var)
        timeout_check.grid(row=0, column=0, sticky=W, padx=5, pady=2)
        self.create_tooltip(timeout_check, 
                    "Interrompe automaticamente la ricerca dopo il tempo specificato.\n"
                    "Utile per evitare ricerche troppo lunghe su grandi volumi di dati.")

        timeout_label = ttk.Label(timeout_grid, text="Secondi:")
        timeout_label.grid(row=0, column=1, sticky=W, padx=(55, 5), pady=2)
        timeout_seconds_var = IntVar(value=self.timeout_seconds.get())
        timeout_spin = ttk.Spinbox(timeout_grid, from_=10, to=3600, width=5, textvariable=timeout_seconds_var)
        timeout_spin.grid(row=0, column=2, padx=5, pady=2, sticky=W)
        self.create_tooltip(timeout_spin, 
                    "Durata massima della ricerca in secondi prima dell'interruzione automatica.\n"
                    "Es. 300 = 5 minuti, 3600 = 1 ora")
        
        max_files_label = ttk.Label(timeout_grid, text="Max file da controllare:")
        max_files_label.grid(row=1, column=0, sticky=W, padx=5, pady=5)
        max_files_var = IntVar(value=self.max_files_to_check.get())
        max_files = ttk.Spinbox(timeout_grid, from_=1000, to=10000000, width=8, textvariable=max_files_var)
        max_files.grid(row=1, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_files, 
                    "Numero massimo di file da controllare prima di terminare la ricerca.\n"
                    "Limita le ricerche molto estese per evitare tempi di attesa eccessivi.\n"
                    "Valori consigliati: 10000 per ricerche rapide, 100000 per ricerche approfondite.")
        
        max_results_label = ttk.Label(timeout_grid, text="Max risultati:")
        max_results_label.grid(row=1, column=2, sticky=W, padx=5, pady=5)
        max_results_var = IntVar(value=self.max_results.get())
        max_results = ttk.Spinbox(timeout_grid, from_=500, to=100000, width=8, textvariable=max_results_var)
        max_results.grid(row=1, column=3, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_results, 
                    "Numero massimo di risultati da mostrare.\n"
                    "Limita il numero di risultati per migliorare le prestazioni\n"
                    "dell'interfaccia in caso di ricerche con molte corrispondenze.")
        
        # Processamento
        process_frame = ttk.LabelFrame(performance_frame, text="Processamento", padding=10)
        process_frame.pack(fill=X, pady=10)
        
        process_grid = ttk.Frame(process_frame)
        process_grid.pack(fill=X)
        
        threads_label = ttk.Label(process_grid, text="Thread paralleli:")
        threads_label.grid(row=0, column=0, sticky=W, padx=5, pady=5)
        threads_var = IntVar(value=self.worker_threads.get())
        threads = ttk.Spinbox(process_grid, from_=1, to=16, width=3, textvariable=threads_var)
        threads.grid(row=0, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(threads, 
                    "Numero di thread paralleli per la ricerca.\n"
                    "Più thread accelerano la ricerca ma usano più CPU.\n"
                    "Consigliato: 4-8 thread su PC moderni.")
        
        max_size_label = ttk.Label(process_grid, text="Dimensione max file (MB):")
        max_size_label.grid(row=0, column=2, sticky=W, padx=5, pady=5)
        max_size_mb_var = IntVar(value=self.max_file_size_mb.get())
        max_size_mb = ttk.Spinbox(process_grid, from_=1, to=1000, width=5, textvariable=max_size_mb_var)
        max_size_mb.grid(row=0, column=3, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_size_mb, 
                    "Dimensione massima in MB dei file di cui analizzare il contenuto.\n"
                    "I file più grandi saranno considerati solo in base al nome.\n"
                    "Un valore più basso aumenta la velocità di ricerca.")
        
        # Calcolo dimensioni
        calc_frame = ttk.LabelFrame(performance_frame, text="Calcolo dimensioni", padding=10)
        calc_frame.pack(fill=X, pady=10)
        
        calc_label = ttk.Label(calc_frame, text="Modalità di calcolo:")
        calc_label.pack(side=LEFT, padx=5)
        dir_size_calc_var = StringVar(value=self.dir_size_calculation.get())
        calc_combo = ttk.Combobox(calc_frame, textvariable=dir_size_calc_var, 
                            values=["incrementale", "preciso", "stimato", "sistema", "disabilitato"], 
                            width=12, state="readonly")
        calc_combo.pack(side=LEFT, padx=5)
        self.create_tooltip(calc_combo, 
                    "Modalità di calcolo della dimensione delle directory:\n"
                    "• incrementale: aggiorna durante la ricerca\n"
                    "• preciso: calcolo completo ma più lento\n"
                    "• stimato: più veloce ma approssimato\n"
                    "• sistema: usa comandi di sistema esterni\n"
                    "• disabilitato: non calcolare le dimensioni")
        
        # Gestione memoria (Aggiunta)
        memory_frame = ttk.LabelFrame(performance_frame, text="Gestione memoria", padding=10)
        memory_frame.pack(fill=tk.X, pady=10)

        auto_memory_var = BooleanVar(value=self.auto_memory_management)
        memory_percent_var = IntVar(value=self.memory_usage_percent)

        auto_memory_check = ttk.Checkbutton(memory_frame, text="Gestione automatica della memoria", variable=auto_memory_var)
        auto_memory_check.pack(anchor=tk.W, pady=5)

        # Creiamo un frame contenitore per le due label affiancate
        labels_container = ttk.Frame(memory_frame)
        labels_container.pack(fill=tk.X, padx=5, pady=2)

        # Creiamo la label che mostrerà la percentuale attuale (a sinistra)
        memory_slider_label = ttk.Label(labels_container, text=f"Soglia utilizzo RAM: {memory_percent_var.get()}%")
        memory_slider_label.pack(side=tk.LEFT, padx=(0, 10))

        # Nuova label per visualizzare la quantità di RAM corrispondente (a destra)
        memory_ram_value_label = ttk.Label(labels_container, text="")
        memory_ram_value_label.pack(side=tk.LEFT)

        # Funzione per calcolare e aggiornare la RAM corrispondente alla percentuale
        def update_ram_value_display():
            percent = memory_percent_var.get()
            try:
                total_ram = psutil.virtual_memory().total / (1024 ** 3)
                ram_threshold = total_ram * (percent / 100)
                memory_ram_value_label.config(
                    text=f"(Equivalente a {ram_threshold:.2f} GB di {total_ram:.2f} GB totali)"
                )
            except Exception as e:
                memory_ram_value_label.config(
                    text="(Errore nel calcolo della RAM)"
                )

        # Funzione per aggiornare il testo della label con la percentuale corrente
        def update_slider_label(*args):
            memory_slider_label.config(text=f"Soglia utilizzo RAM: {memory_percent_var.get()}%")
            update_ram_value_display()  # Aggiorna anche la visualizzazione della RAM

        # Utilizziamo trace_add invece di trace (che è deprecato)
        memory_percent_var.trace_add("write", update_slider_label)

        def toggle_memory_slider():
            if auto_memory_var.get():
                # When auto memory is enabled, reset slider to default 75%
                memory_slider.set(75)
                memory_slider.config(state="disabled")
                # Le label si aggiorneranno automaticamente grazie al trace
            else:
                # When manual control is enabled, enable slider
                memory_slider.config(state="normal")

        auto_memory_check.config(command=toggle_memory_slider)

        memory_slider = ttk.Scale(memory_frame, from_=10, to=95, orient=tk.HORIZONTAL,
                                variable=memory_percent_var, command=lambda v: memory_percent_var.set(int(float(v))))
        memory_slider.pack(fill=tk.X, padx=5, pady=5)

        def update_memory_details():
            percent = memory_percent_var.get()
            try:
                total_ram = psutil.virtual_memory().total / (1024 ** 3)
                ram_usage = total_ram * (percent / 100)
                memory_details_label.config(
                    text=f"RAM utilizzabile: {total_ram:.2f} GB totali"
                )
            except Exception as e:
                memory_details_label.config(
                    text="Errore nel calcolo della RAM. Assicurati che psutil sia installato."
                )
                print(f"Errore nel calcolo della RAM: {e}")

        memory_details_label = ttk.Label(memory_frame, text="")
        memory_details_label.pack(anchor=tk.W, padx=5, pady=5)
        update_memory_details()
        update_ram_value_display()  # Inizializza la visualizzazione della RAM

        self.create_tooltip(auto_memory_check, 
                        "Quando attivata, l'applicazione gestisce automaticamente la memoria\n"
                        "regolando la soglia al 75% e applicando ottimizzazioni quando necessario.\n"
                        "Consigliato per la maggior parte degli utenti.")

        self.create_tooltip(memory_slider, 
                            "Imposta la soglia percentuale di utilizzo della RAM oltre la quale\n"
                            "l'applicazione attiva automaticamente meccanismi di ottimizzazione.\n\n"
                            "Quando l'utilizzo della memoria supera questa soglia, vengono attivate:\n"
                            "• Pulizia della cache di ricerca\n"
                            "• Raccolta automatica dei rifiuti (garbage collection)\n"
                            "• Limitazione dei processi paralleli\n"
                            "• Ottimizzazione delle strutture dati\n\n"
                            "Valori più bassi (60-70%) riducono il rischio di rallentamenti ma\n"
                            "potrebbero causare frequenti ottimizzazioni. Valori più alti (80-90%)\n"
                            "massimizzano le prestazioni ma possono causare picchi di utilizzo RAM.")

        # Aggiorniamo il tooltip per coprire entrambe le label nel container
        self.create_tooltip(labels_container, 
                            "Mostra la soglia percentuale di memoria RAM e il corrispondente\n"
                            "valore in GB oltre il quale l'applicazione avvia procedure\n"
                            "di ottimizzazione della memoria.")

        self.create_tooltip(memory_details_label, 
                            "Riepilogo delle impostazioni di gestione della memoria\n"
                            "e del loro impatto sulle risorse del sistema.")

        # Ottimizzazione Rete
        network_frame = ttk.LabelFrame(performance_frame, text="Ottimizzazione Rete", padding=10)
        network_frame.pack(fill=X, pady=10)

        network_grid = ttk.Frame(network_frame)
        network_grid.pack(fill=X)

        # Checkbox per abilitare ottimizzazione rete
        network_search_var = BooleanVar(value=getattr(self, 'network_search_enabled', True))
        network_check = ttk.Checkbutton(network_grid, text="Ottimizza ricerca su percorsi di rete", 
                                    variable=network_search_var)
        network_check.grid(row=0, column=0, columnspan=2, sticky=W, padx=5, pady=5)
        self.create_tooltip(network_check, 
                        "Abilita ottimizzazioni specifiche per migliorare le prestazioni\n"
                        "durante la ricerca su unità di rete o percorsi UNC.\n"
                        "Migliora la stabilità e riduce i timeout su connessioni lente.")

        # Tentativi di connessione
        retry_label = ttk.Label(network_grid, text="Tentativi di connessione:")
        retry_label.grid(row=1, column=0, sticky=W, padx=5, pady=5)
        network_retry_var = IntVar(value=getattr(self, 'network_retry_count', 3))
        network_retry = ttk.Spinbox(network_grid, from_=1, to=10, width=5, textvariable=network_retry_var)
        network_retry.grid(row=1, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(network_retry, 
                        "Numero di tentativi prima di considerare fallita una\n"
                        "connessione di rete. Valori più alti migliorano l'affidabilità\n"
                        "su reti instabili ma possono rallentare la ricerca.")

        # Ricerche parallele su rete
        parallel_net_label = ttk.Label(network_grid, text="Ricerche parallele:")
        parallel_net_label.grid(row=1, column=2, sticky=W, padx=20, pady=5)
        network_parallel_var = IntVar(value=getattr(self, 'network_parallel_searches', 4))
        network_parallel = ttk.Spinbox(network_grid, from_=1, to=16, width=5, textvariable=network_parallel_var)
        network_parallel.grid(row=1, column=3, padx=5, pady=5, sticky=W)
        self.create_tooltip(network_parallel, 
                        "Numero di operazioni di rete da eseguire in parallelo.\n"
                        "Valori più alti possono accelerare la ricerca ma\n"
                        "potrebbero sovraccaricare connessioni lente.")

        # Ottimizzazione File Grandi
        large_file_frame = ttk.LabelFrame(performance_frame, text="Ricerca su File Grandi", padding=10)
        large_file_frame.pack(fill=X, pady=10)

        large_file_grid = ttk.Frame(large_file_frame)
        large_file_grid.pack(fill=X)

        # Checkbox per abilitare ricerca in file grandi
        large_file_var = BooleanVar(value=getattr(self, 'large_file_search_enabled', True))
        large_file_check = ttk.Checkbutton(large_file_grid, text="Includi file di grandi dimensioni nella ricerca", 
                                        variable=large_file_var, command=lambda: toggle_size_controls())
        large_file_check.grid(row=0, column=0, columnspan=4, sticky=W, padx=5, pady=5)
        self.create_tooltip(large_file_check, 
                            "Quando attivata, l'applicazione analizzerà anche i file di grandi dimensioni.\n"
                            "Per file giganteschi (>2GB), verranno applicate tecniche di analisi parziale.\n"
                            "Disattivare questa opzione accelera la ricerca escludendo file di grandi dimensioni.")

        # Soglia file grandi
        large_threshold_label = ttk.Label(large_file_grid, text="Soglia file grandi (MB):")
        large_threshold_label.grid(row=1, column=0, sticky=W, padx=5, pady=5)
        large_file_threshold_var = IntVar(value=getattr(self, 'large_file_threshold', 50 * 1024 * 1024) // (1024 * 1024))
        large_file_threshold = ttk.Spinbox(large_file_grid, from_=10, to=1000, width=5, textvariable=large_file_threshold_var)
        large_file_threshold.grid(row=1, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(large_file_threshold, 
                        "Dimensione in MB oltre la quale un file viene considerato 'grande'.\n"
                        "I file che superano questa soglia verranno elaborati utilizzando\n"
                        "tecniche di lettura a blocchi per ridurre il consumo di memoria.")

        # Soglia file enormi
        huge_threshold_label = ttk.Label(large_file_grid, text="Soglia file enormi (MB):")
        huge_threshold_label.grid(row=1, column=2, sticky=W, padx=20, pady=5)
        huge_file_threshold_var = IntVar(value=getattr(self, 'huge_file_threshold', 500 * 1024 * 1024) // (1024 * 1024))
        huge_file_threshold = ttk.Spinbox(large_file_grid, from_=100, to=5000, width=5, textvariable=huge_file_threshold_var)
        huge_file_threshold.grid(row=1, column=3, padx=5, pady=5, sticky=W)
        self.create_tooltip(huge_file_threshold, 
                        "Dimensione in MB oltre la quale un file viene considerato 'enorme'.\n"
                        "Per questi file verranno applicate ulteriori ottimizzazioni e\n"
                        "potrebbero essere analizzati solo parzialmente per garantire prestazioni.")

        # Soglia file medi
        medium_threshold_label = ttk.Label(large_file_grid, text="Soglia file medi (MB):")
        medium_threshold_label.grid(row=2, column=0, sticky=W, padx=5, pady=5)
        medium_file_threshold_var = IntVar(value=getattr(self, 'medium_file_threshold', 10 * 1024 * 1024) // (1024 * 1024))
        medium_file_threshold = ttk.Spinbox(large_file_grid, from_=1, to=100, width=5, textvariable=medium_file_threshold_var)
        medium_file_threshold.grid(row=2, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(medium_file_threshold, 
                        "Dimensione in MB oltre la quale un file viene considerato 'medio'.\n"
                        "I file che superano questa soglia potrebbero richiedere ottimizzazioni leggere.")

        # Soglia file giganteschi
        gigantic_threshold_label = ttk.Label(large_file_grid, text="Soglia file giganteschi (MB):")
        gigantic_threshold_label.grid(row=2, column=2, sticky=W, padx=20, pady=5)
        gigantic_file_threshold_var = IntVar(value=getattr(self, 'gigantic_file_threshold', 2048 * 1024 * 1024) // (1024 * 1024))
        gigantic_file_threshold = ttk.Spinbox(large_file_grid, from_=1000, to=10000, width=5, textvariable=gigantic_file_threshold_var)
        gigantic_file_threshold.grid(row=2, column=3, padx=5, pady=5, sticky=W)
        self.create_tooltip(gigantic_file_threshold, 
                        "Dimensione in MB oltre la quale un file viene considerato 'gigantesco'.\n"
                        "Per questi file verranno richieste conferme aggiuntive e\n"
                        "saranno analizzati con tecniche speciali per evitare problemi di memoria.")
        # Funzione per abilitare/disabilitare i controlli delle soglie
        def toggle_size_controls():
            enabled = large_file_var.get()
            state = "normal" if enabled else "disabled"
            
            # Aggiorna stato dei controlli
            medium_threshold_label.configure(state=state)
            medium_file_threshold.configure(state=state)
            large_threshold_label.configure(state=state)
            large_file_threshold.configure(state=state)
            huge_threshold_label.configure(state=state)
            huge_file_threshold.configure(state=state)
            gigantic_threshold_label.configure(state=state)
            gigantic_file_threshold.configure(state=state)

        # Imposta lo stato iniziale dei controlli
        toggle_size_controls()

        toggle_memory_slider()
        # Pulsanti finali per la finestra
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(15, 0))

        # Aggiungi questa funzione prima della definizione dei pulsanti finali
        def restore_defaults():
            # Valori predefiniti per le opzioni di ricerca
            depth_var.set(0)  # Profondità illimitata
            search_files_var.set(True)
            search_folders_var.set(True)
            search_content_var.set(True)  # Per default non cerchiamo nei contenuti
            whole_word_var.set(False)
            
            # Valori predefiniti per i filtri avanzati
            min_size_var.set("0")
            max_size_var.set("0")
            min_date.entry.delete(0, 'end')
            max_date.entry.delete(0, 'end')
            
            # Ripristino esclusioni
            system_exclusions_var.set(True)
            for item in excluded_list.get_children():
                excluded_list.delete(item)
            
            # Percorsi di sistema standard
            system_paths = [
                "C:/Windows", 
                "C:/Program Files", 
                "C:/Program Files (x86)",
                "C:/ProgramData",
                "C:/Users/All Users",
                "C:/Program Files (x86)/Client Active Directory Rights Management Services"
            ]
            for path in system_paths:
                excluded_list.insert("", "end", values=(path,))
            
            # Opzioni a blocchi predefinite
            files_block_var.set(500)  # Valore predefinito
            parallel_var.set(2)      # Valore predefinito
            auto_adjust_var.set(True)
            prioritize_var.set(True)
            
            # Opzioni di performance predefinite
            timeout_enabled_var.set(True)
            timeout_seconds_var.set(3600)
            max_files_var.set(100000)
            max_results_var.set(50000)
            threads_var.set(4)
            max_size_mb_var.set(50)
            dir_size_calc_var.set("disabilitato")
            
            # Opzioni memoria predefinite
            auto_memory_var.set(True)
            memory_percent_var.set(75)
            update_memory_details()
            toggle_memory_slider()
            
            # Aggiungi il ripristino delle impostazioni di rete e file grandi
            network_search_var.set(True)
            network_retry_var.set(3)
            network_parallel_var.set(4)
            
            
            # Ripristina i valori predefiniti per le soglie dei file
            large_file_var.set(True)  # Abilita ricerca file grandi di default
            # Ripristina valori predefiniti in MB (divisi per 1024*1024 per ottenere MB)
            medium_file_threshold_var.set(10)  # 10 MB
            large_file_threshold_var.set(50)   # 50 MB
            huge_file_threshold_var.set(500)   # 500 MB
            gigantic_file_threshold_var.set(2048)  # 2 GB = 2048 MB
            
            # Riattiva i controlli delle soglie se erano disabilitati
            toggle_size_controls()

            # Mostra conferma all'utente
            messagebox.showinfo("Ripristino", "I valori predefiniti sono stati ripristinati.")

        # E poi, modifica il pulsante per usare la funzione appena creata:
        defaults_btn = ttk.Button(btn_frame, text="Ripristina valori predefiniti", command=restore_defaults)
        defaults_btn.pack(side=LEFT)
        self.create_tooltip(defaults_btn, "Ripristina tutte le impostazioni ai valori predefiniti")
        
        # ================= Scheda 6: Aggiornamento =================
        update_frame = ttk.Frame(notebook, padding=15)
        notebook.add(update_frame, text="Aggiornamento")

        # Descrizione generale
        update_desc = ttk.Label(update_frame, 
                            text="Gestisci le impostazioni per l'aggiornamento automatico dell'applicazione.",
                            wraplength=700)
        update_desc.pack(anchor=W, pady=(0, 15))

        # Versione attuale
        current_version_frame = ttk.Frame(update_frame)
        current_version_frame.pack(fill=X, pady=5)

        ttk.Label(current_version_frame, text="Versione attuale:").pack(side=LEFT, padx=5)
        version_value = ttk.Label(current_version_frame, text=f"{APP_VERSION} {APP_STAGE}")
        version_value.pack(side=LEFT, padx=5)

        # Inizializza le variabili delle impostazioni di aggiornamento
        if not hasattr(self, 'update_settings'):
            self.update_settings = {
                "auto_update": True,
                "update_frequency": "All'avvio"
            }

        # Checkbox per aggiornamento automatico
        auto_update_var = BooleanVar(value=self.update_settings.get("auto_update", True))
        auto_update_check = ttk.Checkbutton(update_frame, 
                                        text="Controlla aggiornamenti all'avvio", 
                                        variable=auto_update_var)
        auto_update_check.pack(anchor=W, padx=10, pady=10)
        self.create_tooltip(auto_update_check, 
                        "Attiva per controllare automaticamente la disponibilità\n"
                        "di nuovi aggiornamenti all'avvio dell'applicazione.")

        # Ultimo controllo
        last_check_frame = ttk.Frame(update_frame)
        last_check_frame.pack(fill=X, padx=10, pady=5)

        ttk.Label(last_check_frame, text="Ultimo controllo:").pack(side=LEFT, padx=5)

        last_check_value = ttk.Label(last_check_frame, 
                                text=self.update_settings.get("last_update_check", "Mai"))
        last_check_value.pack(side=LEFT, padx=5)

        # Separatore
        ttk.Separator(update_frame, orient=HORIZONTAL).pack(fill=X, padx=10, pady=15)

        # Frame per controllo manuale e status
        manual_check_frame = ttk.Frame(update_frame)
        manual_check_frame.pack(fill=X, padx=10, pady=5)

        # Status dell'aggiornamento
        update_status_var = StringVar(value="Nessun controllo effettuato")
        update_status_label = ttk.Label(manual_check_frame, 
                                    textvariable=update_status_var,
                                    wraplength=400)
        update_status_label.pack(side=LEFT, padx=10, fill=X, expand=YES)

        # Pulsante di controllo
        check_button = ttk.Button(manual_check_frame, 
                            text="Controlla ora", 
                            command=lambda: self.check_for_updates(update_status_var, last_check_value))
        check_button.pack(side=RIGHT, padx=5)
        self.create_tooltip(check_button, 
                        "Verifica immediatamente se è disponibile\n"
                        "una nuova versione dell'applicazione.")

        # Pulsante per scaricare l'aggiornamento 
        self.download_button = ttk.Button(
            manual_check_frame,
            text="Scarica aggiornamento",
            command=lambda: self.download_update(self.latest_release_url))
        self.download_button.pack(side=RIGHT, padx=0)
        self.download_button.pack_forget()  # Nasconde il pulsante

        release_notes_frame = ttk.LabelFrame(update_frame, text="Note di rilascio")
        release_notes_frame.pack(fill="both", expand=True, padx=5, pady=5)
        
        # Utilizziamo un widget Text con scrollbar invece di una semplice Label
        # per gestire meglio il testo potenzialmente lungo
        release_notes_text = tk.Text(release_notes_frame, height=8, width=50, wrap="word")
        release_notes_text.pack(side="left", fill="both", expand=True, padx=5, pady=5)
        
        # Aggiungi una scrollbar
        notes_scrollbar = ttk.Scrollbar(release_notes_frame, orient="vertical", command=release_notes_text.yview)
        notes_scrollbar.pack(side="right", fill="y")
        release_notes_text.config(yscrollcommand=notes_scrollbar.set)
        
        # Configura il Text widget come di sola lettura
        release_notes_text.config(state="disabled")
        
        # Salva il riferimento al widget per poterlo aggiornare in seguito
        self.release_notes_text = release_notes_text
        
        # Inizialmente il widget è vuoto o mostra un messaggio predefinito
        self.update_release_notes("Verifica gli aggiornamenti per visualizzare le note di rilascio.")
        
        def save_options():
            try:
                               
                # Salva le opzioni di ricerca
                self.max_depth = int(depth_var.get())
                self.search_files.set(search_files_var.get())
                self.search_folders.set(search_folders_var.get())
                self.search_content.set(search_content_var.get())
                self.whole_word_search.set(whole_word_var.get())
                if not hasattr(self, 'search_settings'):
                    self.search_settings = {}
                    self.search_settings["use_windows_search"] = self.use_windows_search_var.get()

                # Salva i filtri avanzati
                min_kb = int(min_size_var.get() or 0)
                max_kb = int(max_size_var.get() or 0)
                self.advanced_filters["size_min"] = min_kb * 1024
                self.advanced_filters["size_max"] = max_kb * 1024 if max_kb > 0 else 0
                
                self.advanced_filters["date_min"] = min_date.entry.get().strip()
                self.advanced_filters["date_max"] = max_date.entry.get().strip()
                
                # Salva i percorsi esclusi
                new_excluded_paths = []
                for item in excluded_list.get_children():
                    values = excluded_list.item(item)["values"]
                    if values:
                        new_excluded_paths.append(values[0])
                self.excluded_paths = new_excluded_paths
                
                # Salva le opzioni a blocchi
                self.max_files_per_block.set(files_block_var.get())
                self.max_parallel_blocks.set(parallel_var.get())
                self.block_size_auto_adjust.set(auto_adjust_var.get())
                self.prioritize_user_folders.set(prioritize_var.get())
                
                # Salva le opzioni di performance
                self.timeout_enabled.set(timeout_enabled_var.get())
                self.timeout_seconds.set(timeout_seconds_var.get())
                self.max_files_to_check.set(max_files_var.get())
                self.max_results.set(max_results_var.get())
                self.worker_threads.set(threads_var.get())
                self.max_file_size_mb.set(max_size_mb_var.get())
                self.dir_size_calculation.set(dir_size_calc_var.get())
                
                # Salva le opzioni di gestione della memoria
                self.auto_memory_management = auto_memory_var.get()
                self.memory_usage_percent = memory_percent_var.get()

                # Salva le opzioni di rete
                self.network_search_enabled = network_search_var.get()
                self.network_retry_count = network_retry_var.get()
                self.network_parallel_searches = network_parallel_var.get()

                # Aggiorna il NetworkSearchOptimizer se è stato inizializzato
                if hasattr(self, 'network_optimizer'):
                    self.network_optimizer.retry_count = self.network_retry_count

                # Salva le opzioni per file di grandi dimensioni
                self.large_file_search_enabled = large_file_var.get()
                self.large_file_threshold = large_file_threshold_var.get() * 1024 * 1024
                self.huge_file_threshold = huge_file_threshold_var.get() * 1024 * 1024
                self.medium_file_threshold = medium_file_threshold_var.get() * 1024 * 1024
                self.gigantic_file_threshold = gigantic_file_threshold_var.get() * 1024 * 1024

                # Aggiorna il LargeFileHandler se è stato inizializzato
                if hasattr(self, 'large_file_handler'):
                    self.large_file_handler.large_file_threshold = self.large_file_threshold
                    self.large_file_handler.huge_file_threshold = self.huge_file_threshold

                # Aggiorna le variabili dell'update settings nel metodo di salvataggio
                self.update_settings["auto_update"] = auto_update_var.get()
                
                #  Salva effettivamente le impostazioni su file
                if hasattr(self, 'save_settings_to_file'):
                    self.save_settings_to_file()
                    self.log_debug("Impostazioni avanzate salvate su file permanente")
                else:
                    self.log_debug("AVVISO: Metodo save_settings_to_file non disponibile, salvataggio permanente non effettuato")
            
                messagebox.showinfo("Impostazioni", "Opzioni salvate con successo!")
                dialog.destroy()
                
            except ValueError as e:
                error_msg = f"Errore di valore non valido: {str(e)}"
                self.log_debug(f"ERRORE SALVATAGGIO: {error_msg}")
                messagebox.showerror("Errore", error_msg)
                return
            except Exception as e:
                error_msg = f"Errore durante il salvataggio: {str(e)}"
                self.log_debug(f"ERRORE SALVATAGGIO: {error_msg}")
                import traceback
                self.log_debug(traceback.format_exc())  # Registra il traceback completo dell'errore
                messagebox.showerror("Errore", error_msg)
                return
        
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=dialog.destroy)
        cancel_btn.pack(side=RIGHT, padx=5)
        self.create_tooltip(cancel_btn, "Chiudi la finestra senza salvare le modifiche")
        
        save_btn = ttk.Button(btn_frame, text="Salva", command=save_options)
        save_btn.pack(side=RIGHT, padx=5)
        self.create_tooltip(save_btn, "Salva tutte le impostazioni e chiudi la finestra")
        
        # Centra la finestra
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Mostra la finestra solo quando è pronta e posizionata
        dialog.deiconify()
        dialog.grab_set()  # Imposta il focus sulla finestra

    def update_release_notes(self, text):
        """Aggiorna il widget delle note di rilascio con il testo fornito"""
        if hasattr(self, 'release_notes_text') and self.release_notes_text.winfo_exists():
            # Abilita temporaneamente la modifica
            self.release_notes_text.config(state="normal")
            # Cancella il contenuto attuale
            self.release_notes_text.delete(1.0, tk.END)
            # Inserisci il nuovo testo
            self.release_notes_text.insert(tk.END, text)
            # Disabilita di nuovo la modifica
            self.release_notes_text.config(state="disabled")
            
    def check_for_updates(self, status_var=None, last_check_label=None):
        """Controlla se sono disponibili aggiornamenti sulla repository GitHub"""
        try:
            # Nascondi il pulsante di download all'inizio del controllo
            if hasattr(self, 'download_button') and self.download_button.winfo_exists():
                self.download_button.pack_forget()
            
            # Aggiorna il widget delle note di rilascio con un messaggio iniziale
            if hasattr(self, 'release_notes_text') and self.release_notes_text.winfo_exists():
                self.update_release_notes("Verifica aggiornamenti in corso...")
            
            # Aggiorna variabili UI
            if status_var:
                status_var.set("Verifica aggiornamenti in corso...")
            self.log_debug("Verifica aggiornamenti in corso...")
            
            # Importazioni necessarie per HTTP e JSON
            import urllib.request
            import urllib.error
            import json
            
            # Ottieni la versione attuale (rimuovi "V" o "v" iniziale se presente)
            current_version = APP_VERSION
            current_stage = APP_STAGE
            if current_version.lower().startswith("v"):
                current_version = current_version[1:]
                
            # Converti la versione in componenti numerici
            try:
                # Estrai solo la parte numerica dalla versione corrente
                import re
                current_version_match = re.match(r"(\d+(?:\.\d+)*)", current_version)
                if not current_version_match:
                    self.log_debug(f"Errore nella conversione della versione corrente: {current_version}")
                    if status_var:
                        status_var.set("Errore nel controllo: formato versione non valido")
                    self.update_release_notes("Errore nel controllo: formato versione non valido")
                    return False, None
                    
                # Ottieni solo la parte numerica
                current_numeric_version = current_version_match.group(1)
                current_parts = current_numeric_version.split(".")
                current_v_tuple = tuple(map(int, current_parts))
                
            except ValueError:
                self.log_debug(f"Errore nella conversione della versione corrente: {current_version}")
                if status_var:
                    status_var.set("Errore nel controllo: formato versione non valido")
                self.update_release_notes("Errore nella conversione della versione corrente")
                return False, None
                
            try:
                # Chiama l'API GitHub per le release
                # Verifico prima se la repository esiste
                try:
                    # Verifica che la repository esista
                    repo_url = "https://api.github.com/repos/Nino19980/File-search-tools"
                    
                    req = urllib.request.Request(repo_url)
                    req.add_header('User-Agent', 'FileSearchApp/1.0')  # GitHub richiede uno User-Agent
                    
                    try:
                        repo_response = urllib.request.urlopen(req, timeout=10)
                        # La connessione è riuscita se arriviamo qui
                    except urllib.error.HTTPError as e:
                        if e.code == 404:
                            self.log_debug("Repository non trovata: verificare il nome utente e il nome repository")
                            if status_var:
                                status_var.set("Repository GitHub non trovata. Verificare la connessione Internet o le impostazioni.")
                            self.update_release_notes("Repository GitHub non trovata. Verificare la connessione Internet o le impostazioni.")
                            return False, None
                        else:
                            # Gestione di altri errori HTTP
                            self.log_debug(f"Errore nella risposta API GitHub: {e.code}")
                            
                            # Messaggio più informativo per l'utente
                            error_message = f"Errore nel controllo: server GitHub ha risposto con codice {e.code}"
                            
                            if e.code == 403:
                                error_message = "Limite di richieste API GitHub superato. Riprova più tardi."
                            elif e.code >= 500:
                                error_message = "Errore del server GitHub. Riprova più tardi."
                            
                            if status_var:
                                status_var.set(error_message)
                            self.update_release_notes(error_message)
                            return False, None
                    except urllib.error.URLError as e:
                        self.log_debug(f"Errore di connessione durante la verifica della repository: {str(e)}")
                        if status_var:
                            status_var.set(f"Errore di connessione: {str(e.reason)}")
                        self.update_release_notes(f"Errore di connessione: {str(e.reason)}")
                        return False, None
                        
                except Exception as e:
                    self.log_debug(f"Errore nella verifica della repository: {str(e)}")
                    if status_var:
                        status_var.set(f"Errore durante il controllo: {str(e)}")
                    self.update_release_notes(f"Errore durante il controllo: {str(e)}")
                    return False, None
                        
                # Procedi con la richiesta delle release
                api_url = "https://api.github.com/repos/Nino19980/File-search-tools/releases"
                self.log_debug(f"Richiesta API GitHub a: {api_url}")
                
                req = urllib.request.Request(api_url)
                req.add_header('User-Agent', 'FileSearchApp/1.0')  # GitHub richiede uno User-Agent
                
                try:
                    response = urllib.request.urlopen(req, timeout=10)
                    response_data = response.read()
                    
                    # Log della risposta per debug
                    self.log_debug(f"Risposta API GitHub: status={response.getcode()}, content-length={len(response_data)}")
                    
                    # Analizza la risposta JSON
                    releases = json.loads(response_data.decode('utf-8'))
                    
                except urllib.error.HTTPError as e:
                    self.log_debug(f"Errore nella risposta API GitHub: {e.code}")
                    
                    # Messaggio più informativo per l'utente
                    error_message = f"Errore nel controllo: server GitHub ha risposto con codice {e.code}"
                    
                    if e.code == 404:
                        error_message = "Repository non trovata su GitHub. Verificare il nome repository o la connessione Internet."
                    elif e.code == 403:
                        error_message = "Limite di richieste API GitHub superato. Riprova più tardi."
                    elif e.code >= 500:
                        error_message = "Errore del server GitHub. Riprova più tardi."
                    
                    if status_var:
                        status_var.set(error_message)
                    self.update_release_notes(error_message)
                    return False, None
                except urllib.error.URLError as e:
                    error_msg = f"Errore di connessione durante il controllo: {str(e.reason)}"
                    self.log_debug(error_msg)
                    if status_var:
                        status_var.set(error_msg)
                    self.update_release_notes(f"Impossibile recuperare le note di rilascio: {error_msg}")
                    return False, None
                except json.JSONDecodeError as e:
                    error_msg = f"Errore nell'analisi della risposta JSON: {str(e)}"
                    self.log_debug(error_msg)
                    if status_var:
                        status_var.set(error_msg)
                    self.update_release_notes(f"Impossibile recuperare le note di rilascio: {error_msg}")
                    return False, None
                
                if not releases:
                    self.log_debug("Nessuna release trovata su GitHub")
                    if status_var:
                        status_var.set("Nessuna release trovata su GitHub")
                    self.update_release_notes("Nessuna release trovata su GitHub")
                    return False, None
                    
                # Escludiamo sempre le versioni beta (prerelease)
                releases = [r for r in releases if not r.get("prerelease", False)]
                    
                if not releases:
                    self.log_debug("Nessuna release stabile trovata (versioni beta escluse)")
                    if status_var:
                        status_var.set("Nessuna release stabile disponibile")
                    self.update_release_notes("Nessuna release stabile disponibile (versioni beta escluse)")
                    return False, None
                    
                # Trova la release più recente
                latest_release = releases[0]  # Le release sono già ordinate cronologicamente
                latest_version = latest_release.get("tag_name", "")

                self.log_debug(f"Release più recente trovata: {latest_version}")

                # Verifica che il tag non sia solo 'v' o 'V' e contenga effettivamente numeri di versione
                if latest_version.lower() == "v" or latest_version.strip() == "":
                    self.log_debug(f"Tag di versione non valido: {latest_version}. Saltando controllo aggiornamenti.")
                    if status_var:
                        status_var.set("Tag di versione non valido. Impossibile verificare aggiornamenti.")
                    self.update_release_notes(f"Tag di versione non valido: {latest_version}. Impossibile verificare aggiornamenti.")
                    return False, None

                # Rimuovi 'v' o 'V' iniziale se presente
                if latest_version.lower().startswith("v"):
                    latest_version = latest_version[1:]

                # Estrai solo la parte numerica dalla versione GitHub
                latest_version_match = re.match(r"(\d+(?:\.\d+)*)", latest_version)
                if not latest_version_match:
                    self.log_debug(f"Formato di versione non valido: {latest_version}")
                    if status_var:
                        status_var.set(f"Formato versione GitHub non valido: {latest_version}")
                    self.update_release_notes(f"Formato di versione non valido: {latest_version}")
                    return False, None
                    
                # Ottieni solo la parte numerica
                latest_numeric_version = latest_version_match.group(1)
                
                # Converti in componenti numerici per confronto
                try:
                    latest_parts = latest_numeric_version.split(".")
                    # Verifica che ci siano parti da convertire e che non siano vuote
                    if not latest_parts or any(part.strip() == "" for part in latest_parts):
                        self.log_debug(f"Formato di versione non valido: {latest_numeric_version}")
                        if status_var:
                            status_var.set(f"Formato versione GitHub non valido: {latest_numeric_version}")
                        self.update_release_notes(f"Formato di versione non valido: {latest_numeric_version}")
                        return False, None
                    
                    latest_v_tuple = tuple(map(int, latest_parts))
                    
                    self.log_debug(f"Confronto versioni: attuale={current_v_tuple}, più recente={latest_v_tuple}")
                except ValueError as e:
                    self.log_debug(f"Errore nella conversione della versione GitHub: {latest_numeric_version} - {str(e)}")
                    if status_var:
                        status_var.set(f"Errore nel confronto versioni: formato non valido ({latest_numeric_version})")
                    self.update_release_notes(f"Errore nella conversione della versione GitHub: {latest_numeric_version}")
                    return False, None
                    
                # Confronta le versioni
                is_update_available = latest_v_tuple > current_v_tuple
                
                # Aggiorna l'ultimo controllo
                now = datetime.now().strftime("%d/%m/%Y %H:%M")
                self.update_settings["last_update_check"] = now
                
                if last_check_label:
                    last_check_label.config(text=now)
                    
                # Salva le impostazioni aggiornate
                self.save_settings_to_file()
                
                # Aggiorna lo stato dell'interfaccia
                if is_update_available:
                    # Usa la versione originale con suffisso per il messaggio
                    message = f"Aggiornamento disponibile: {latest_release.get('tag_name')} (attuale: {APP_VERSION})"
                    self.log_debug(message)
                    
                    # Memorizza l'URL della release per il pulsante di download
                    self.latest_release_url = latest_release.get("html_url")
                    
                    if status_var:
                        download_message = (f"È disponibile la versione {latest_release.get('tag_name')}")
                        status_var.set(download_message)
                        
                    # Mostra il pulsante di download se esiste
                    if hasattr(self, 'download_button') and self.download_button.winfo_exists():
                        self.download_button.pack(side=RIGHT, padx=0)
                    
                    # Rendi cliccabile il messaggio di stato
                    if hasattr(self, 'update_status_label') and self.update_status_label.winfo_exists():
                        self.update_status_label.bind("<Button-1>", 
                                                lambda e: self.download_update(latest_release.get("html_url")))
                        self.update_status_label.config(cursor="hand2", foreground="blue")
                    
                    # Estrai e visualizza le note di rilascio
                    release_body = latest_release.get("body", "Nessuna nota di rilascio disponibile.")
                    release_notes = f"Novità nella versione {latest_release.get('tag_name')}:\n\n{release_body}"
                    self.update_release_notes(release_notes)
                        
                    return True, latest_release
                else:
                    message = f"Sei già alla versione più recente: {APP_VERSION}"
                    self.log_debug(message)
                    if status_var:
                        status_var.set(message)
                    
                    # Nascondi il pulsante di download se non ci sono aggiornamenti
                    if hasattr(self, 'download_button') and self.download_button.winfo_exists():
                        self.download_button.pack_forget()
                    
                    # Aggiorna il testo delle note di rilascio
                    self.update_release_notes("Sei già alla versione più recente. Non ci sono nuove note di rilascio da visualizzare.")
                        
                    return False, None
                    
            except Exception as e:
                error_msg = f"Errore nel controllo aggiornamenti: {str(e)}"
                self.log_debug(error_msg)
                if status_var:
                    status_var.set(error_msg)
                self.update_release_notes(f"Impossibile recuperare le note di rilascio: {error_msg}")
                return False, None
        except Exception as e:
            error_msg = f"Errore nel controllo aggiornamenti: {str(e)}"
            self.log_debug(error_msg)
            if status_var:
                status_var.set(error_msg)
            self.update_release_notes(f"Impossibile recuperare le note di rilascio: {error_msg}")
            return False, None

    @error_handler
    def download_update(self, release_url):
        """Apre il browser per scaricare l'aggiornamento"""
        try:
            webbrowser.open(release_url)
            self.log_debug(f"Apertura pagina di download: {release_url}")
        except Exception as e:
            self.log_debug(f"Errore nell'apertura del browser: {str(e)}")
            messagebox.showerror("Errore", 
                            f"Impossibile aprire il browser automaticamente.\n\n"
                            f"Visita manualmente: {release_url}")

    @error_handler
    def check_for_updates_on_startup(self):
        """Controlla aggiornamenti all'avvio se abilitato"""
        if not hasattr(self, 'update_settings'):
            return
            
        auto_update = self.update_settings.get("auto_update", True)
        if not auto_update:
            return
            
        # Verifica la frequenza impostata
        frequency = self.update_settings.get("update_frequency", "All'avvio")
        last_check = self.update_settings.get("last_update_check", "Mai")
        
        should_check = False
        
        if frequency == "All'avvio" or last_check == "Mai":
            should_check = True
        elif frequency in ["Giornaliera", "Settimanale", "Mensile"] and last_check != "Mai":
            try:
                last_date = datetime.strptime(last_check, "%d/%m/%Y %H:%M")
                now = datetime.now()
                delta = now - last_date
                
                if frequency == "Giornaliera" and delta.days >= 1:
                    should_check = True
                elif frequency == "Settimanale" and delta.days >= 7:
                    should_check = True
                elif frequency == "Mensile" and delta.days >= 30:
                    should_check = True
            except ValueError:
                # Se c'è un errore nel formato della data, forza il controllo
                should_check = True
        
        if should_check:
            # Esegui il controllo in background dopo un breve ritardo
            self.root.after(3000, self._check_updates_background)

    @error_handler
    def _check_updates_background(self):
        """Esegue il controllo aggiornamenti in background"""
        # Avvia un thread separato per non bloccare l'interfaccia
        threading.Thread(target=self._async_check_updates, daemon=True).start()

    @error_handler
    def _async_check_updates(self):
        """Esegue il controllo aggiornamenti in un thread separato"""
        try:
            is_update_available, latest_release = self.check_for_updates()
            
            if is_update_available and latest_release:
                # Notifica l'utente solo se l'aggiornamento è disponibile
                self.root.after(0, lambda: self._show_update_notification(latest_release))
        except Exception as e:
            self.log_debug(f"Errore nel controllo aggiornamenti in background: {str(e)}")

    @error_handler
    def _show_update_notification(self, release):
        """Mostra una notifica quando è disponibile un aggiornamento"""
        version = release.get("tag_name", "")
        url = release.get("html_url", "")
        
        response = messagebox.askyesno(
            "Aggiornamento disponibile",
            f"È disponibile una nuova versione di {APP_NAME}: {version}\n\n"
            f"Vuoi aprire la pagina di download?",
            icon="info"
        )
        
        if response:
            self.download_update(url)

    @error_handler
    def create_tooltip(self, widget, text, delay=500, fade=True):
        """Crea tooltip con ritardo, effetti di dissolvenza e larghezza automatica"""
        
        tooltip_timer = None        
        
        def show_tooltip():
            x = widget.winfo_rootx() + 25
            y = widget.winfo_rooty() + 25
            
            tooltip = ttk.Toplevel(widget)
            tooltip.wm_overrideredirect(True)
            tooltip.attributes("-alpha", 0.0)  # Inizia invisibile
            tooltip.attributes("-topmost", True)  # Mantiene sopra altre finestre
            
            # Crea un frame con bordo e padding
            frame = ttk.Frame(tooltip, borderwidth=1, relief="solid", padding=10)
            frame.pack(fill="both", expand=True)
            
            # Calcolo della larghezza del testo
            font = ("Segoe UI", 9)  # Puoi modificare questo per usare un altro font
            
            # Calcola la lunghezza approssimativa del testo
            temp_label = tk.Label(font=font)
            
            # Determina se il testo necessita di essere diviso in più righe
            lines = text.split("\n")
            max_line_width = 0
            
            for line in lines:
                temp_label.configure(text=line)
                line_width = temp_label.winfo_reqwidth()
                max_line_width = max(max_line_width, line_width)
            
            # Limita la larghezza massima a 500 pixel
            max_width = min(max_line_width, 500)
            
            wraplength = max_width
                
            # Crea l'etichetta con il testo
            label = ttk.Label(frame, text=text, justify="left", 
                            wraplength=wraplength, font=font)
            label.pack(fill="both", expand=True)
            
            # Distruggi il label temporaneo
            temp_label.destroy()
            
            # Aggiorna immediatamente per calcolare le dimensioni
            tooltip.update_idletasks()
            
            # Posiziona il tooltip in modo che non vada fuori dallo schermo
            tooltip_width = tooltip.winfo_reqwidth()
            tooltip_height = tooltip.winfo_reqheight()
            
            screen_width = tooltip.winfo_screenwidth()
            screen_height = tooltip.winfo_screenheight()
            
            # Aggiusta la posizione se il tooltip esce dallo schermo
            if x + tooltip_width > screen_width:
                x = screen_width - tooltip_width - 10
            
            if y + tooltip_height > screen_height:
                y = screen_height - tooltip_height - 10
                
            # Imposta la posizione finale
            tooltip.wm_geometry(f"+{x}+{y}")
            
            widget._tooltip = tooltip
            
            if fade:
                # Avvia il fade in utilizzando la nuova funzione generica
                self.fade(tooltip, 0.0, 1.0, 0.1)
        
        def enter(event):
            nonlocal tooltip_timer
            # Avvia il timer per mostrare il tooltip dopo un certo ritardo
            tooltip_timer = widget.after(delay, show_tooltip)
        
        def leave(event):
            nonlocal tooltip_timer
            
            # Cancella il timer se esiste
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
                tooltip_timer = None
            
            # Rimuovi il tooltip se esiste
            if hasattr(widget, "_tooltip"):
                tooltip = widget._tooltip
                
                def destroy_tooltip():
                    try:
                        if hasattr(widget, "_tooltip"):
                            del widget._tooltip
                        tooltip.destroy()
                    except Exception:
                        pass  # Ignora errori durante la distruzione
                        
                if fade:
                    # Tenta di avviare il fade out
                    try:
                        self.fade(tooltip, 1.0, 0.0, -0.1, destroy_tooltip)
                    except Exception:
                        # Se il fade fallisce, distruggi immediatamente
                        destroy_tooltip()
                else:
                    destroy_tooltip()
        
        # Gestisce anche il caso in cui il widget venga distrutto
        def on_destroy(event):
            nonlocal tooltip_timer
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
            if hasattr(widget, "_tooltip"):
                widget._tooltip.destroy()
        
        widget.bind("<Enter>", enter)
        widget.bind("<Leave>", leave)
        widget.bind("<Destroy>", on_destroy)

    @error_handler
    def fade(self, widget, from_alpha, to_alpha, step, callback=None):
        """Anima la trasparenza di un widget da un valore alpha a un altro."""
        # Memorizza l'ID del widget per verificare se è lo stesso nelle chiamate successive
        widget_id = str(widget)
        
        # Imposta l'alpha iniziale
        try:
            if not widget.winfo_exists():
                return  # Il widget non esiste più
            widget.attributes("-alpha", from_alpha)
        except (tk.TclError, RuntimeError, Exception) as e:
            self.log_error(f"Errore nell'inizializzazione dell'animazione fade: {e}")
            return
        
        # Funzione interna per gestire l'animazione frame per frame
        def update_alpha(current_alpha=from_alpha):
            # Verifica se il widget esiste ancora e ha lo stesso ID
            try:
                if not widget.winfo_exists() or str(widget) != widget_id:
                    return  # Il widget non esiste più o è cambiato
                    
                # Calcola il nuovo valore alpha
                new_alpha = current_alpha + step
                
                # Controlla se l'animazione è completa
                if (step > 0 and new_alpha >= to_alpha) or (step < 0 and new_alpha <= to_alpha):
                    # Imposta il valore finale esatto
                    widget.attributes("-alpha", to_alpha)
                    # Esegui il callback se fornito
                    if callback:
                        try:
                            callback()
                        except Exception as e:
                            self.log_error(f"Errore nell'esecuzione del callback di fade: {e}")
                    return
                
                # Altrimenti, continua l'animazione
                widget.attributes("-alpha", new_alpha)
                widget.after(20, lambda: update_alpha(new_alpha))
                    
            except (tk.TclError, RuntimeError, Exception) as e:
                # Log dell'errore ma senza interrompere l'applicazione
                self.log_debug(f"Animazione fade interrotta: {e}")
                return
        
        # Avvia l'animazione
        try:
            widget.after(20, update_alpha)
        except (tk.TclError, RuntimeError, Exception) as e:
            self.log_debug(f"Impossibile avviare l'animazione fade: {e}")

    def select_all(self):
        self.results_list.selection_set(self.results_list.get_children())
        
    def deselect_all(self):
        self.results_list.selection_remove(self.results_list.get_children())

    @error_handler   
    def invert_selection(self):
        all_items = self.results_list.get_children()
        selected_items = self.results_list.selection()
        self.results_list.selection_remove(selected_items)
        to_select = set(all_items) - set(selected_items)
        for item in to_select:
            self.results_list.selection_add(item)
    
    @error_handler
    def treeview_sort_column(self, tv, col, reverse):
        """Ordina il TreeView in base alla colonna cliccata"""
        # Ottieni la lista di item con i loro valori
        l = [(tv.set(k, col), k) for k in tv.get_children('')]
        
        # Determina il tipo di ordinamento in base alla colonna
        try:
            if col == "size":
                # Gestione speciale per le dimensioni (KB, MB, GB)
                def extract_size(size_str):
                    if not size_str:
                        return 0
                    try:
                        # Estrai il valore numerico e converti tutto in byte per confrontare correttamente
                        value = float(size_str.split()[0])
                        if "KB" in size_str:
                            return value * 1024
                        elif "MB" in size_str:
                            return value * 1024 * 1024
                        elif "GB" in size_str:
                            return value * 1024 * 1024 * 1024
                        else:
                            return value  # Assumiamo che sia in byte
                    except:
                        return 0
                        
                l.sort(key=lambda x: extract_size(x[0]), reverse=reverse)
            elif col in ("modified", "created"):
                # Gestione per le date (assumendo formato DD/MM/YYYY HH:MM)
                from datetime import datetime
                def parse_date(date_str):
                    try:
                        if date_str and date_str != "N/A":
                            return datetime.strptime(date_str, "%d/%m/%Y %H:%M")
                        return datetime(1900, 1, 1)  # Data di default per valori vuoti
                    except:
                        return datetime(1900, 1, 1)
                        
                l.sort(key=lambda x: parse_date(x[0]), reverse=reverse)
            else:
                # Ordinamento standard alfanumerico
                l.sort(reverse=reverse)
        except Exception as e:
            self.log_debug(f"Errore durante l'ordinamento: {str(e)}")
            # Fallback all'ordinamento alfanumerico semplice
            l.sort(reverse=reverse)
            
        # Riorganizza gli elementi nel TreeView
        for index, (val, k) in enumerate(l):
            tv.move(k, '', index)

        # Memorizza la colonna ordinata e la direzione per il prossimo click
        tv.heading(col, command=lambda _col=col: self.treeview_sort_column(tv, _col, not reverse))
        
        # Aggiorna l'indicatore visivo di ordinamento (aggiunge freccia all'intestazione)
        for c in tv["columns"]:
            if c == col:
                tv.heading(c, text=f"{tv.heading(c, 'text').split(' ')[0]} {'▼' if reverse else '▲'}")
            else:
                # Rimuovi eventuali indicatori di ordinamento da altre colonne
                tv.heading(c, text=tv.heading(c, 'text').split(' ')[0])
    
    @error_handler
    def fade(self, widget, from_alpha, to_alpha, step, callback=None):
        """Anima la trasparenza di un widget con gestione robusta degli errori."""
        # Aggiungiamo un ID di animazione univoco al widget
        animation_id = str(uuid.uuid4())
        
        try:
            # Assegna l'ID dell'animazione al widget
            widget._fade_animation_id = animation_id
            widget.attributes("-alpha", from_alpha)
        except Exception:
            return  # Widget non più valido, abbandona l'animazione
        
        def update_alpha(current_alpha=from_alpha):
            try:
                # Verifica che il widget esista ancora e che l'animazione non sia stata sostituita
                if not hasattr(widget, '_fade_animation_id') or widget._fade_animation_id != animation_id:
                    return  # Animazione cancellata o sostituita
                    
                # Calcola il nuovo valore alpha
                new_alpha = current_alpha + step
                
                # Controlla se l'animazione è completa
                if (step > 0 and new_alpha >= to_alpha) or (step < 0 and new_alpha <= to_alpha):
                    try:
                        widget.attributes("-alpha", to_alpha)
                        # Rimuovi l'ID dell'animazione
                        if hasattr(widget, '_fade_animation_id'):
                            del widget._fade_animation_id
                        # Esegui il callback in modo sicuro
                        if callback:
                            try:
                                callback()
                            except Exception:
                                pass
                    except Exception:
                        pass  # Ignora gli errori finali
                    return
                
                # Continua l'animazione
                try:
                    widget.attributes("-alpha", new_alpha)
                    widget.after(20, lambda: update_alpha(new_alpha))
                except Exception:
                    pass  # Ignora gli errori e termina l'animazione
                    
            except Exception:
                pass  # Ignora qualsiasi errore e termina l'animazione
        
        # Avvia l'animazione
        try:
            widget.after(20, update_alpha)
        except Exception:
            pass  # Avvio animazione fallito, abbandona

    @error_handler # Show debug log window with export functionality
    def show_debug_log(self):
        """Mostra una finestra con i log di debug e funzionalità di filtraggio, rispettando il tema corrente"""
        # Ottieni il tema corrente
        current_theme = "darkly"  # Valore predefinito
        if hasattr(self, 'theme_combobox') and self.theme_combobox is not None:
            current_theme = self.theme_combobox.get()
        
        # Determina i colori in base al tema
        is_dark_theme = current_theme in ["darkly", "cyborg"]
        
        # Colori per il testo di debug
        bg_color = "#1e1e1e" if is_dark_theme else "#ffffff"
        fg_color = "#e0e0e0" if is_dark_theme else "#000000"
        
        # Colori per i tipi di messaggi
        error_color = "#ff6b6b" if is_dark_theme else "#cc0000"
        warning_color = "#ffb86c" if is_dark_theme else "#ff8800"
        info_color = "#ffffff" if is_dark_theme else "#000000"
        
        # Personalizzazioni specifiche per tema
        if current_theme == "darkly":
            bg_color = "#121212"
        elif current_theme == "cyborg":
            bg_color = "#060606"
            info_color = "#2a9fd6"
        elif current_theme == "minty":
            info_color = "#78c2ad"
        elif current_theme == "cosmo":
            info_color = "#2780e3"
        
        if not hasattr(self, 'debug_window') or not self.debug_window.winfo_exists():
            # Crea la finestra di debug
            self.debug_window = tk.Toplevel(self.root)
            self.debug_window.title("Debug Log")
            self.debug_window.geometry("1400x900")
            
            frame = ttk.Frame(self.debug_window)
            frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
            
            # Header frame con conteggio messaggi e controlli di filtraggio
            header_frame = ttk.Frame(frame)
            header_frame.pack(fill=tk.X, pady=(0, 10))
            
            # Verifica che il debug log sia inizializzato
            if not hasattr(self, 'debug_log'):
                self.debug_log = []
            
            # Etichetta informativa con conteggio (spostata a sinistra)
            self.log_count_label = ttk.Label(header_frame, text=f"Registro di debug dell'applicazione: {len(self.debug_log)} messaggi")
            self.log_count_label.pack(side=tk.LEFT, padx=(0, 20))
            
            # Frame centrale per i controlli di filtraggio
            filter_frame = ttk.Frame(header_frame)
            filter_frame.pack(side=tk.LEFT, fill=tk.X, expand=True)
            
            # Etichetta e combobox per il filtraggio (centrati)
            filter_inner_frame = ttk.Frame(filter_frame)
            filter_inner_frame.pack(side=tk.TOP, anchor=tk.CENTER)
            
            filter_label = ttk.Label(filter_inner_frame, text="Filtra per tipo:")
            filter_label.pack(side=tk.LEFT, padx=(0, 5))
            
            # Combobox per il filtraggio
            self.filter_var = tk.StringVar(value="Tutti")
            self.filter_combo = ttk.Combobox(filter_inner_frame, textvariable=self.filter_var, 
                                            values=["Tutti", "Errore", "Avviso", "Info"],
                                            width=10, state="readonly")
            self.filter_combo.pack(side=tk.LEFT, padx=5)
            self.filter_combo.bind("<<ComboboxSelected>>", self.filter_log_messages)
            
            # Opzione auto-scroll (spostata a destra)
            self.autoscroll_var = tk.BooleanVar(value=True)
            ttk.Checkbutton(header_frame, text="Auto scorrimento", variable=self.autoscroll_var).pack(side=tk.RIGHT, padx=5)
            
            # Crea un text widget con scrollbar
            text_frame = ttk.Frame(frame)
            text_frame.pack(fill=tk.BOTH, expand=True)
            
            # Scrollbar verticale
            v_scrollbar = ttk.Scrollbar(text_frame)
            v_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
            
            # Scrollbar orizzontale
            h_scrollbar = ttk.Scrollbar(text_frame, orient=tk.HORIZONTAL)
            h_scrollbar.pack(side=tk.BOTTOM, fill=tk.X)
            
            # Text widget con colori basati sul tema
            self.debug_text = tk.Text(
                text_frame, 
                wrap=tk.NONE,  # Permette scroll orizzontale
                width=80, 
                height=20,
                bg=bg_color,  # Sfondo basato sul tema
                fg=fg_color,  # Testo basato sul tema
                font=("Consolas", 10),  # Font a larghezza fissa
                xscrollcommand=h_scrollbar.set,
                yscrollcommand=v_scrollbar.set
            )
            
            self.debug_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
            
            # Configura le scrollbar
            v_scrollbar.config(command=self.debug_text.yview)
            h_scrollbar.config(command=self.debug_text.xview)
            
            # Configura i tag per i colori in base al tema
            self.debug_text.tag_configure("error", foreground=error_color)
            self.debug_text.tag_configure("warning", foreground=warning_color)
            self.debug_text.tag_configure("info", foreground=info_color)
            
            # Aggiungi pulsanti di utilità
            btn_frame = ttk.Frame(self.debug_window)
            btn_frame.pack(fill=tk.X, padx=10, pady=5)
            
            ttk.Button(btn_frame, text="Pulisci Log", command=self.clear_log).pack(side=tk.LEFT, padx=5)
            ttk.Button(btn_frame, text="Esporta in TXT", command=self.export_log_to_txt).pack(side=tk.LEFT, padx=5)
            
            # Posiziona la finestra al centro
            self.debug_window.update_idletasks()
            width = self.debug_window.winfo_width()
            height = self.debug_window.winfo_height()
            x = (self.debug_window.winfo_screenwidth() // 2) - (width // 2)
            y = (self.debug_window.winfo_screenheight() // 2) - (height // 2)
            self.debug_window.geometry(f"{width}x{height}+{x}+{y}")
            
            # Imposta una dimensione minima ragionevole
            self.debug_window.minsize(800, 500)
            
            # Memorizza la configurazione di filtraggio iniziale
            self.current_filter = "Tutti"
            
            # Memorizza tutti i messaggi originali per il filtraggio
            self.all_log_messages = []
            if hasattr(self, 'debug_log'):
                self.all_log_messages = self.debug_log.copy()
            
            # Mostra il log corrente
            self.update_log_display()
        else:
            # Se la finestra esiste già, portala in primo piano
            self.debug_window.lift()
            self.debug_window.focus_force()
            
            # Aggiorna il tema se il widget debug_text esiste
            if hasattr(self, 'debug_text'):
                # Applica i colori del tema corrente
                self.debug_text.configure(bg=bg_color, fg=fg_color)
                
                # Aggiorna i tag per i colori
                self.debug_text.tag_configure("error", foreground=error_color)
                self.debug_text.tag_configure("warning", foreground=warning_color)
                self.debug_text.tag_configure("info", foreground=info_color)
            
            # Aggiorna il contenuto
            self.update_log_display()

    @error_handler
    def filter_log_messages(self, event=None):
        """Filtra i messaggi di log in base al tipo selezionato"""
        # Ottieni il filtro selezionato
        selected_filter = self.filter_var.get()
        self.current_filter = selected_filter
        
        # Se non ci sono log o la finestra non esiste, esci
        if not hasattr(self, 'debug_log') or not hasattr(self, 'debug_text') or not self.debug_window.winfo_exists():
            return
        
        # Salva tutti i messaggi se non l'abbiamo già fatto
        if not hasattr(self, 'all_log_messages') or not self.all_log_messages:
            self.all_log_messages = self.debug_log.copy()
        
        # Pulisci il testo esistente
        self.debug_text.config(state=tk.NORMAL)
        self.debug_text.delete("1.0", tk.END)
        
        # Applica il filtro
        filtered_messages = []
        
        if selected_filter == "Tutti":
            filtered_messages = self.all_log_messages
        else:
            # Mappa la selezione al prefisso corrispondente 
            # (supporta sia prefissi in italiano che in inglese durante la transizione)
            filter_map = {
                "Errore": ["[ERRORE]", "[ERROR]"],
                "Avviso": ["[AVVISO]", "[WARNING]"],
                "Info": ["[INFO]"]
            }
            
            prefixes = filter_map.get(selected_filter, [])
            if prefixes:
                filtered_messages = [msg for msg in self.all_log_messages 
                                    if any(prefix in msg for prefix in prefixes)]
        
        # Aggiorna l'etichetta con il conteggio dei messaggi filtrati
        if hasattr(self, 'log_count_label'):
            if selected_filter == "Tutti":
                self.log_count_label.config(
                    text=f"Registro di debug dell'applicazione: {len(filtered_messages)} messaggi"
                )
            else:
                self.log_count_label.config(
                    text=f"Registro di debug dell'applicazione: {len(filtered_messages)} messaggi ({selected_filter})"
                )
        
        # Inserisci i messaggi filtrati
        for message in filtered_messages:
            self.debug_text.insert(tk.END, message + "\n")
        
        # Applica la colorazione
        self.highlight_errors()
        
        # Scorri alla fine se richiesto
        if hasattr(self, 'autoscroll_var') and self.autoscroll_var.get():
            self.debug_text.see(tk.END)
        
        # Rendi il testo di nuovo sola lettura
        self.debug_text.config(state=tk.DISABLED)

    def reset_log_filter(self):
        """Resetta il filtro e mostra tutti i messaggi"""
        if hasattr(self, 'filter_combo'):
            self.filter_var.set("Tutti")
            self.filter_log_messages()

    @error_handler
    def update_log_display(self):
        """Aggiorna completamente la visualizzazione dei log nella finestra di debug"""
        if not hasattr(self, 'debug_window') or not self.debug_window.winfo_exists():
            return
            
        if not hasattr(self, 'debug_text'):
            return
        
        # Reset del contatore dei log visualizzati
        self.last_displayed_log_index = 0
            
        # Cancella il contenuto attuale
        self.debug_text.config(state=tk.NORMAL)
        self.debug_text.delete(1.0, tk.END)
        
        # Verifica che il debug log sia inizializzato
        if not hasattr(self, 'debug_log'):
            self.debug_log = []
            
        # Aggiorna l'etichetta con il conteggio dei messaggi
        if hasattr(self, 'log_count_label'):
            self.log_count_label.config(text=f"Registro di debug dell'applicazione: {len(self.debug_log)} messaggi")
            
        if self.debug_log:
            # Limita la visualizzazione a 5000 messaggi per non rallentare l'interfaccia
            max_display = 5000
            if len(self.debug_log) > max_display:
                self.debug_text.insert(tk.END, f"[Mostrando solo gli ultimi {max_display} di {len(self.debug_log)} messaggi...]\n\n")
                log_entries = self.debug_log[-max_display:]
                self.last_displayed_log_index = len(self.debug_log) - max_display
            else:
                log_entries = self.debug_log
                self.last_displayed_log_index = 0
                
            # Inserisci i log
            log_text = "\n".join(log_entries)
            self.debug_text.insert(tk.END, log_text)
            
            # Aggiorna l'indice dell'ultimo messaggio visualizzato
            self.last_displayed_log_index = len(self.debug_log)
            
            # Evidenzia gli errori con colore rosso
            self.highlight_errors()
            
            # Forza un aggiornamento dell'interfaccia
            self.debug_window.update_idletasks()
            
            # Prima imposta l'indice orizzontale all'inizio (più a sinistra)
            self.debug_text.xview_moveto(0.0)
            
            # Poi imposta anche il primo carattere visibile nella prima colonna
            self.debug_text.mark_set("insert", "1.0")
            
            # Decidi se mantenere la posizione verticale precedente o scorrere alla fine
            if hasattr(self, 'autoscroll_var') and self.autoscroll_var.get():
                self.debug_text.see(tk.END)  # Scorri alla fine verticalmente
        else:
            self.debug_text.insert(tk.END, "Nessun messaggio di debug disponibile.")
            
        # Rendi il testo di nuovo sola lettura
        self.debug_text.config(state=tk.DISABLED)

    @error_handler
    def clear_log(self):
        """Pulisce completamente i log di debug"""
        # Pulisce l'array dei log
        self.debug_log = []
        
        # Pulisce anche l'array della cronologia completa
        if hasattr(self, 'complete_debug_log_history'):
            self.complete_debug_log_history = []
        
        # Resetta il contatore dei log visualizzati
        if hasattr(self, 'last_displayed_log_index'):
            self.last_displayed_log_index = 0
        
        # Pulisce la coda dei log
        if hasattr(self, 'debug_logs_queue'):
            try:
                while not self.debug_logs_queue.empty():
                    self.debug_logs_queue.get_nowait()
            except:
                pass
        
        # Pulisce il text widget se la finestra è aperta
        if hasattr(self, 'debug_window') and self.debug_window.winfo_exists() and hasattr(self, 'debug_text'):
            self.debug_text.config(state=tk.NORMAL)
            self.debug_text.delete('1.0', tk.END)
            self.debug_text.insert(tk.END, "Log pulito. Nessun messaggio di debug disponibile.")
            self.debug_text.config(state=tk.DISABLED)
            
            # Aggiorna l'etichetta con il conteggio
            if hasattr(self, 'log_count_label'):
                self.log_count_label.config(text=f"Registro di debug dell'applicazione: 0 messaggi")
        
        # Aggiungi un messaggio che indica che i log sono stati puliti
        self.log_debug("I log sono stati puliti manualmente dall'utente")

    @error_handler
    def export_log_to_txt(self):
        """Esporta i log in un file di testo in base al filtro selezionato e registra l'operazione"""
        # Ottieni il filtro corrente dalla combobox
        filtro_attuale = self.filter_var.get()
        
        # Assicurati che all_log_messages sia popolato
        if not hasattr(self, 'all_log_messages') or not self.all_log_messages:
            if hasattr(self, 'debug_log'):
                self.all_log_messages = self.debug_log.copy()
            else:
                self.all_log_messages = []
        
        # Applica il filtro con una logica robusta
        if filtro_attuale == "Tutti":
            # Utilizza tutti i messaggi di log
            log_filtrati = self.all_log_messages
        elif filtro_attuale == "Errore":
            # Cerca diverse varianti di errore
            log_filtrati = [msg for msg in self.all_log_messages 
                        if "[Errore]" in msg or "[ERRORE]" in msg or "[errore]" in msg
                        or "ERROR" in msg.upper() or "ERRORE" in msg.upper()]
        elif filtro_attuale == "Avviso":
            # Cerca diverse varianti di avviso
            log_filtrati = [msg for msg in self.all_log_messages 
                        if "[Avviso]" in msg or "[AVVISO]" in msg or "[avviso]" in msg
                        or "WARNING" in msg.upper() or "AVVISO" in msg.upper()]
        elif filtro_attuale == "Info":
            # Cerca diverse varianti di info
            log_filtrati = [msg for msg in self.all_log_messages 
                        if "[Info]" in msg or "[INFO]" in msg or "[info]" in msg
                        or "INFO" in msg.upper()]
        else:
            # Fallback a tutti i messaggi
            log_filtrati = self.all_log_messages
        
        # Verifica se ci sono log da esportare
        if not log_filtrati:
            tk.messagebox.showinfo("Nessun log da esportare", 
                                f"Non ci sono messaggi di log del tipo '{filtro_attuale}' da esportare.")
            return
        
        # Chiedi all'utente dove salvare il file
        file_path = tk.filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt"), ("All files", "*.*")],
            title=f"Salva log {filtro_attuale} come file di testo"
        )
        
        if not file_path:  # Utente ha annullato
            return
            
        try:
            # Ottieni username e timestamp corrente
            import getpass
            username = getpass.getuser()
            current_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            
            log_content = "\n".join(log_filtrati)
            
            # Aggiungi intestazione con timestamp e info utente
            header = f"Debug Log ({filtro_attuale}) - Esportato il {current_time}\n"
            header += f"Utente: {username}\n"
            header += f"Applicazione: {APP_FULL_NAME}\n"
            header += f"Numero totale messaggi: {len(log_filtrati)}\n"
            header += "-" * 80 + "\n\n"
            
            # Scrivi su file
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(header + log_content)

            # Messaggio di successo
            tk.messagebox.showinfo("Esportazione completata", 
                            f"Log di tipo '{filtro_attuale}' salvati con successo in:\n{file_path}\n\n"
                            f"Il file contiene {len(log_filtrati)} messaggi di log.")
            
            # Registra l'operazione di esportazione nei log dell'applicazione
            export_log_message = f"[Info] {current_time} - Utente {username} ha esportato {len(log_filtrati)} messaggi di log di tipo '{filtro_attuale}' in: {file_path}"
            
            # Aggiungi il messaggio ai log
            if hasattr(self, 'debug_log'):
                self.debug_log.append(export_log_message)
            
            # Aggiorna anche all_log_messages
            if hasattr(self, 'all_log_messages'):
                self.all_log_messages.append(export_log_message)
            
            # Aggiorna la visualizzazione dei log se la finestra è aperta
            if hasattr(self, 'debug_window') and self.debug_window.winfo_exists():
                # Aggiungi il nuovo messaggio alla visualizzazione
                self.debug_text.config(state=tk.NORMAL)
                self.debug_text.insert(tk.END, export_log_message + "\n", "info")
                
                # Se auto-scroll è attivo, scorri alla fine
                if hasattr(self, 'autoscroll_var') and self.autoscroll_var.get():
                    self.debug_text.see(tk.END)
                
                self.debug_text.config(state=tk.DISABLED)
                
                # Aggiorna il contatore dei messaggi
                if hasattr(self, 'log_count_label'):
                    self.log_count_label.config(text=f"Registro di debug dell'applicazione: {len(self.debug_log)} messaggi")
                
        except Exception as e:
            tk.messagebox.showerror("Errore esportazione", 
                            f"Si è verificato un errore durante l'esportazione:\n{str(e)}")

# Funzione principale per eseguire l'applicazione
def main():
    import sys
    
    # Controlla se ci sono impostazioni salvate per il tema
    settings_file = os.path.join(os.path.expanduser("~"), ".file_search_settings.json")
    initial_theme = "darkly"  # Tema predefinito
    
    try:
        if os.path.exists(settings_file):
            with open(settings_file, 'r', encoding='utf-8') as f:
                settings = json.load(f)
            
            if "theme" in settings and settings["theme"] in ["minty", "cosmo", "darkly", "cyborg"]:
                initial_theme = settings["theme"]
    except Exception:
        pass  # In caso di errore, usa il tema predefinito
    
    # Crea la finestra principale con il tema caricato
    root = ttk.Window(themename=initial_theme)
    root.withdraw()  # Nascondi completamente la finestra durante l'inizializzazione
    
    # Crea la schermata di splash con dimensioni più piccole per caricamento più veloce
    splash = create_splash_screen(root)
    
    # Inizializza l'applicazione (ma l'interfaccia rimane nascosta)
    app = FileSearchApp(root)
    
    # Controlla se ci sono argomenti da linea di comando
    if len(sys.argv) > 1:
        app.search_path.set(sys.argv[1])
    
    # Funzione per completare l'avvio e mostrare la finestra principale
    def finish_startup():
        splash.destroy()
        
        # Dimensioni desiderate per la finestra principale
        width, height = 1400, 850
        
        # Calcola la posizione centrale
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        x = (screen_width // 2) - (width // 2)
        y = (screen_height // 2) - (height // 2)
        
        # Imposta la geometria per centrarla
        root.geometry(f"{width}x{height}+{x}+{y}")
        
        # Mostra la finestra completamente costruita
        root.deiconify()
        
        # Avvia il caricamento di componenti non essenziali in background
        root.after(100, app._delayed_startup_tasks)
    
    # Ridotto il delay da 1500 a 500 ms
    root.after(500, finish_startup)
    
    root.mainloop()

def create_splash_screen(parent):
    splash_win = tk.Toplevel(parent)
    splash_win.title("")
    splash_win.overrideredirect(True)
    splash_win.attributes("-topmost", True)
    
    # Ridotte le dimensioni dello splash per caricamento più veloce
    width, height = 500, 250
    screen_width = splash_win.winfo_screenwidth()
    screen_height = splash_win.winfo_screenheight()
    x = (screen_width // 2) - (width // 2)
    y = (screen_height // 2) - (height // 2)
    splash_win.geometry(f"{width}x{height}+{x}+{y}")
    
    # Contenuto semplificato dello splash
    frame = ttk.Frame(splash_win, padding=10)
    frame.pack(fill=tk.BOTH, expand=tk.YES)
    
    ttk.Label(frame, text=APP_FULL_NAME, font=("Helvetica", 18, "bold")).pack(pady=(10, 5))
    ttk.Label(frame, text=APP_NAME, font=("Helvetica", 14)).pack(pady=(0, 20))
    ttk.Label(frame, text="Caricamento applicazione in corso...").pack()
    
    progress = ttk.Progressbar(frame, mode="indeterminate")
    progress.pack(fill=tk.X, pady=10)
    progress.start(15)
    
    return splash_win

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        import traceback
        with open("error_log.txt", "w", encoding="utf-8") as f:
            f.write(f"Si è verificato un errore: {str(e)}\n")
            f.write(traceback.format_exc())
        print(f"Si è verificato un errore: {str(e)}")
        
        # Se stai usando tkinter, puoi anche mostrare un messaggio all'utente
        try:
            import tkinter as tk
            from tkinter import messagebox
            root = tk.Tk()
            root.withdraw()
            messagebox.showerror("Errore durante l'avvio", 
                f"Si è verificato un errore durante l'avvio dell'applicazione.\n\n"
                f"Errore: {str(e)}\n\n"
                f"I dettagli sono stati salvati nel file error_log.txt")
        except:
            pass  # Se anche la visualizzazione del messaggio fallisce, continua

