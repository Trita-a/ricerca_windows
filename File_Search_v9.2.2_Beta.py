# Importazioni essenziali per l'avvio
import io
import os
import tkinter as tk
import ttkbootstrap as ttk
from ttkbootstrap.constants import *
from tkinter import filedialog, messagebox, BooleanVar, StringVar, IntVar
import threading
import queue
from datetime import datetime
import getpass
import csv
# Importazioni che erano nel metodo import_non_essential_modules
import shutil
import zipfile
import traceback
import time
import concurrent.futures
import mimetypes
import signal
import re
import subprocess
import odfdo

# Dizionario per tracciare il supporto alle librerie - sarà popolato in seguito
file_format_support = {
    "docx": False, "pdf": False, "pptx": False, "excel": False,
    "odt": False, "rtf": False, "xls": False, "doc": False
}

# Elenco di librerie da installare se mancanti
missing_libraries = []

class FileSearchApp:
    def __init__(self, root):
        self.root = root
        self.root.title("File Search Tool V9.2 Beta Forensics G.di F.")
        
        # Imposta subito il debug mode per poter loggare
        self.debug_mode = True
        
        # Aggiungi questa riga per inizializzare current_user
        self.current_user = getpass.getuser()
        
        # Inizializza tutte le variabili in un passaggio
        self._init_essential_variables()
        self._init_remaining_variables()
        
        # Crea l'intera interfaccia in una volta sola
        self.create_widgets()
        
        # Applica il tema una sola volta
        self.update_theme_colors("dark")
        
        # Esegui attività di background dopo un breve ritardo
        self.root.after(500, self._delayed_startup_tasks)

    def create_base_interface(self):
        """Crea solo l'interfaccia essenziale per un avvio veloce"""
        # Forza l'aggiornamento dell'interfaccia
        self.root.update_idletasks()

    def complete_initialization(self):
        """Completa l'inizializzazione dell'applicazione"""
        try:
            # Rimuovi i componenti temporanei
            if hasattr(self, 'status_label'):
                self.status_label.destroy()
            if hasattr(self, 'progress_bar'):
                self.progress_bar.stop()
                self.progress_bar.destroy()
            
            # Inizializza le variabili rimanenti
            self._init_remaining_variables()
            
            # Crea l'interfaccia completa
            self.create_widgets()
            
            # Esegui le altre operazioni iniziali
            self._delayed_startup_tasks()
            
        except Exception as e:
            # In caso di errore, mostra un messaggio e prova a ripristinare l'applicazione
            messagebox.showerror("Errore di inizializzazione", 
                            f"Si è verificato un errore durante l'avvio: {str(e)}")
            # Tenta comunque di creare l'interfaccia
            try:
                self.create_widgets()
            except:
                pass
            
    def _init_essential_variables(self):
        """Inizializza solo le variabili essenziali per l'avvio"""
        # Variabili principali per la ricerca
        self.search_content = BooleanVar(value=True)
        self.search_path = StringVar()
        self.keywords = StringVar()
        self.search_results = []
        self.search_files = BooleanVar(value=True)
        self.search_folders = BooleanVar(value=True)
        self.is_searching = False
        self.progress_queue = queue.Queue()
        self.search_depth = StringVar(value="base") 

        # Variabili per data/ora e utente
        self.datetime_var = StringVar()
        self.user_var = StringVar(value=getpass.getuser())
        
        # Variabili essenziali per l'interfaccia
        self.ignore_hidden = BooleanVar(value=True)
        self.search_executor = None
        self.exclude_system_files = BooleanVar(value=True)
        self.whole_word_search = BooleanVar(value=False)
        self.dir_size_calculation = StringVar(value="disabilitato")
        
        # Variabili per la visualizzazione
        self.dir_size_var = StringVar(value="")
        self.total_disk_var = StringVar(value="")
        self.used_disk_var = StringVar(value="")
        self.free_disk_var = StringVar(value="")

    def _init_remaining_variables(self):
        """Inizializza le variabili non essenziali per l'avvio"""
        # Variabili per la ricerca a blocchi
        self.max_files_per_block = IntVar(value=1000)
        self.max_parallel_blocks = IntVar(value=4)
        self.prioritize_user_folders = BooleanVar(value=True)
        self.block_size_auto_adjust = BooleanVar(value=True)
        
        # Variabili aggiuntive
        self.search_start_time = None
        self.stop_search = False
        self.max_depth = 0  # 0 = illimitato
        self.last_search_params = {}
        self.search_history = []
        
        # Filtri avanzati
        self.advanced_filters = {
            "size_min": 0,
            "size_max": 0,
            "date_min": None,
            "date_max": None,
            "extensions": []
        }
        
        # Opzioni per la gestione dei permessi
        self.skip_permission_errors = BooleanVar(value=True)
        self.excluded_paths = []
        self.load_settings_from_file()
        
        # Directory problematiche da escludere automaticamente
        self.problematic_dirs = [
            "Client Active Directory Rights Management Services",
            "Windows Resource Protection",
            "Windows Defender",
            "Microsoft Office",
            "System Volume Information",
            "$Recycle.Bin",
            "$WINDOWS.~BT",
            "$Windows.~WS"
        ]

        # Verifica dei privilegi di amministratore
        self.is_admin = False
        try:
            import ctypes
            self.is_admin = ctypes.windll.shell32.IsUserAnAdmin() != 0
        except:
            self.is_admin = False

        # Variabili per limiti di tempo e altre ottimizzazioni
        self.timeout_enabled = BooleanVar(value=False)
        self.timeout_seconds = IntVar(value=3600)
        self.max_files_to_check = IntVar(value=100000)
        self.max_results = IntVar(value=50000)
        self.chunk_size = 8192
        self.max_file_size_mb = IntVar(value=100)
        self.worker_threads = IntVar(value=min(8, os.cpu_count() or 4))
        self.use_indexing = BooleanVar(value=True)
        self.search_index = {}
        
        # Lista di estensioni di file di sistema da escludere dalla ricerca nei contenuti
        self.system_file_extensions = [
            # File eseguibili e librerie
            '.exe', '.dll', '.sys', '.drv', '.ocx', '.vxd', '.com', '.bat', '.cmd', '.scr', '.app', '.dylib', '.exp', '.bpl',
            # ... resto delle estensioni ... 
        ]

        # Percorso del log dei file saltati
        self.skipped_files_log_path = os.path.join(os.path.expanduser("~"), "skipped_files_log.txt")
        
        # Aggiorna l'orario
        self.update_datetime()

    def _delayed_startup_tasks(self):
        """Esegue attività non essenziali all'avvio"""
        # Esegui queste operazioni in background
        threading.Thread(target=self._background_tasks, daemon=True).start()

    def _background_tasks(self):
        """Esegue operazioni in background"""
        try:
            # Verifica librerie disponibili
            global file_format_support, missing_libraries
            
            # Controlla ogni libreria
            for module_name, format_key, import_name in [
                ("docx", "docx", "python-docx"),
                ("PyPDF2", "pdf", "PyPDF2"),
                # ... e così via per le altre librerie
            ]:
                try:
                    __import__(module_name)
                    file_format_support[format_key] = True
                except ImportError:
                    missing_libraries.append(import_name)
            
            # Mostra notifica dopo un ritardo
            if missing_libraries:
                self.root.after(2000, self.check_and_notify_missing_libraries)
            
            # Aggiorna informazioni disco senza calcolare dimensione directory
            if hasattr(self, 'search_path'):
                path = self.search_path.get()
                if path and os.path.exists(path):
                    self.update_disk_info(path, calculate_dir_size=False)
            
        except Exception as e:
            self.log_debug(f"Errore nelle operazioni in background: {str(e)}")

    def _create_minimal_interface(self):
        """Crea solo i componenti essenziali dell'interfaccia"""
        # Frame principale che conterrà tutto
        self.main_container = ttk.Frame(self.root)
        self.main_container.pack(fill=BOTH, expand=YES)
                
        # Controlli minimi per la ricerca
        self.controls_frame = ttk.LabelFrame(self.main_container, text="Parametri di ricerca", padding=10)
        self.controls_frame.pack(fill=X, padx=10, pady=5)
        
        # Solo i controlli essenziali di ricerca
        self._create_essential_search_controls()
        
        # Barra di stato per informazioni
        status_frame = ttk.Frame(self.main_container, padding=5)
        status_frame.pack(fill=X, padx=10)
        
        self.status_label = ttk.Label(status_frame, text="Inizializzazione in corso...", font=("", 9))
        self.status_label.pack(side=LEFT)
        
        self.progress_bar = ttk.Progressbar(status_frame, mode='indeterminate')
        self.progress_bar.pack(side=LEFT, padx=10, fill=X, expand=YES)
        self.progress_bar.start(10)  # Avvia l'animazione di caricamento
        
        # Placeholder per i risultati
        self.results_container = ttk.LabelFrame(self.main_container, text="Risultati di ricerca", padding=10)
        self.results_container.pack(fill=BOTH, expand=YES, padx=10, pady=(5, 10))
        
        # Messaggio di caricamento
        self.loading_label = ttk.Label(self.results_container, text="Caricamento interfaccia in corso...", font=("", 12))
        self.loading_label.pack(expand=YES, pady=50)

    def _create_essential_search_controls(self):
        """Crea solo i controlli essenziali per la ricerca"""
        # RIGA 1: Directory di ricerca
        path_frame = ttk.Frame(self.controls_frame)
        path_frame.pack(fill=X, pady=5)
        
        path_label = ttk.Label(path_frame, text="Directory:", width=12, anchor=W)
        path_label.pack(side=LEFT, padx=(0, 5))
        
        self.path_entry = ttk.Entry(path_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        self.browse_btn = ttk.Button(path_frame, text="Sfoglia", command=self.browse_directory, width=10)
        self.browse_btn.pack(side=LEFT)
        
        # RIGA 2: Parole chiave
        keyword_frame = ttk.Frame(self.controls_frame)
        keyword_frame.pack(fill=X, pady=5)
        
        keyword_label = ttk.Label(keyword_frame, text="Parole chiave:", width=12, anchor=W)
        keyword_label.pack(side=LEFT, padx=(0, 5))
        
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        # Solo il pulsante di ricerca
        button_frame = ttk.Frame(self.controls_frame)
        button_frame.pack(fill=X, pady=(10, 5))
        
        center_frame = ttk.Frame(button_frame)
        center_frame.pack(side=TOP)
        
        self.search_button = ttk.Button(center_frame, text="CERCA", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=15)
        self.search_button.pack(side=LEFT, padx=10)
        
        # Disabilita il pulsante fino al caricamento completo
        self.search_button["state"] = "disabled"

    def _create_remaining_widgets(self):
        """Crea il resto dell'interfaccia in background"""
        try:
            # Aggiorna il messaggio di stato
            self.status_label["text"] = "Caricamento componenti avanzati..."
            
            # Rimuovi gli oggetti temporanei
            if hasattr(self, 'loading_label'):
                self.loading_label.destroy()
            
            # Ferma l'animazione della barra di progresso
            self.progress_bar.stop()
            self.progress_bar["mode"] = "determinate"
            
            # Crea tutti i componenti dell'interfaccia completa
            self.create_widgets()
            
            # Aggiorna lo stato
            self.status_label["text"] = "Applicazione pronta"
            
            # Riabilita il pulsante di ricerca
            self.search_button["state"] = "normal"
            
            # Imposta il focus sul campo di ricerca
            self.path_entry.focus_set()
            
            # Imposta il tema iniziale
            self.update_theme_colors("dark")
            
            self.log_debug("Interfaccia utente completamente caricata")
            
        except Exception as e:
            self.log_debug(f"Errore nel caricamento dell'interfaccia: {str(e)}")
            self.status_label["text"] = "Errore nel caricamento dell'interfaccia"

    def _check_available_libraries(self):
        """Verifica la disponibilità delle librerie in background senza bloccare l'avvio"""
        # Esegui il controllo in un thread separato
        threading.Thread(target=self._async_check_libraries, daemon=True).start()

    def _async_check_libraries(self):
        """Controlla le librerie in background senza bloccare l'interfaccia"""
        global file_format_support, missing_libraries
        
        # Funzione per verificare un singolo modulo
        def check_module(module_name, format_key, import_name=None):
            try:
                # Usa importlib per evitare bloccaggi
                import importlib
                importlib.import_module(module_name)
                file_format_support[format_key] = True
                self.log_debug(f"Supporto {format_key} attivato")
            except ImportError:
                missing_libraries.append(import_name or module_name)
                self.log_debug(f"Supporto {format_key} non disponibile")
        
        # Controlla ogni libreria
        check_module("docx", "docx", "python-docx")
        check_module("PyPDF2", "pdf", "PyPDF2")
        check_module("pptx", "pptx", "python-pptx")
        check_module("openpyxl", "excel", "openpyxl")
        check_module("striprtf.striprtf", "rtf", "striprtf")
        check_module("win32com.client", "doc", "pywin32")
        check_module("win32com.client", "xls", "pywin32")
        check_module("xlrd", "xls_native", "xlrd")
        
        # Controlla supporto ODT con gestione di alternative
        try:
            import importlib
            importlib.import_module("odfdo")
            file_format_support["odt"] = True
            self.log_debug("Supporto ODT attivato (libreria odfdo)")
        except ImportError:
            try:
                from odf import opendocument
                file_format_support["odt"] = True
                self.log_debug("Supporto ODT attivato (libreria odf)")
            except ImportError:
                missing_libraries.append("odfdo")
                self.log_debug("Supporto ODT non disponibile")
        
        # Mostra notifica dopo un ritardo più lungo per non disturbare l'avvio
        if missing_libraries:
            self.root.after(2000, self.check_and_notify_missing_libraries)

    def process_file(self, file_path, keywords, search_content=True):
        if self.stop_search:
            return None
            
        # Salta direttamente i file problematici
        if self.should_skip_file(file_path):
            return None
        
        try:
            # Implementazione timeout cross-platform tramite threading
            result = [None]
            exception = [None]
            processing_completed = [False]
            
            # Avvia il thread con un nome identificativo
            import threading
            worker_thread = threading.Thread(
                target=lambda: self.process_with_timeout(file_path, keywords, result, exception, processing_completed, search_content), 
                name=f"Process-{os.path.basename(file_path)}"
            )
            worker_thread.daemon = True
            worker_thread.start()
            
            # Attendi il completamento del thread con timeout ridotto
            if os.path.splitext(file_path)[1].lower() in ['.doc', '.xls']:
                # Timeout ridotto per file problematici
                timeout_seconds = 5
                self.log_debug(f"Utilizzo timeout ridotto (5s) per il file: {file_path}")
            else:
                timeout_seconds = 20
            worker_thread.join(timeout_seconds)
            
            # Verifica se il thread è ancora in esecuzione (timeout raggiunto)
            if not processing_completed[0]:
                self.log_debug(f"Timeout nella elaborazione del file {file_path}")
                # Non attendere più il thread, continua semplicemente
                return None
                
            # Verifica se si è verificata un'eccezione
            if exception[0]:
                self.log_debug(f"Eccezione nell'elaborazione del file {file_path}: {str(exception[0])}")
                return None
                
            # Restituisci il risultato dal thread
            return result[0]
            
        except Exception as e:
            self.log_debug(f"Errore generale nella elaborazione del file {file_path}: {str(e)}")
            return None
            
    def process_with_timeout(self, file_path, keywords, result, exception, processing_completed, search_content=True):
        try:
            # Verifica filtri di dimensione
            file_size = os.path.getsize(file_path)
            if (self.advanced_filters["size_min"] > 0 and file_size < self.advanced_filters["size_min"]) or \
            (self.advanced_filters["size_max"] > 0 and file_size > self.advanced_filters["size_max"]):
                result[0] = None
                return
            
            # Verifica filtri di data
            if self.advanced_filters["date_min"] or self.advanced_filters["date_max"]:
                mod_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                
                if self.advanced_filters["date_min"]:
                    min_date = datetime.strptime(self.advanced_filters["date_min"], "%d-%m-%Y")
                    if mod_time < min_date:
                        result[0] = None
                        return
                        
                if self.advanced_filters["date_max"]:
                    max_date = datetime.strptime(self.advanced_filters["date_max"], "%d-%m-%Y")
                    if mod_time > max_date:
                        result[0] = None
                        return
            
            # Verifica filtri estensione
            if self.advanced_filters["extensions"] and not any(file_path.lower().endswith(ext.lower()) 
                                                        for ext in self.advanced_filters["extensions"]):
                result[0] = None
                return
                
            if os.path.splitext(file_path)[1].lower() in ['.doc', '.xls']:
                self.progress_queue.put(("progress", self.progress_bar["value"]))  # Forza aggiornamento
                
            # Verifica corrispondenza nel nome file
            filename = os.path.basename(file_path)

            matched = False
            for keyword in keywords:
                # Verifica se la ricerca di parole intere è attivata
                if self.whole_word_search.get():
                    # Gestisce anche frasi con spazi usando la nuova funzione helper
                    if self.is_whole_word_match(keyword, filename):
                        matched = True
                        break
                # Se è un termine con spazi (contesto specifico)
                elif ' ' in keyword:
                    if keyword.lower() in filename.lower():
                        matched = True
                        break
                # Ricerca normale
                else:
                    if keyword.lower() in filename.lower():
                        matched = True
                        break
            
            if matched:
                if self.debug_mode and self.whole_word_search.get():
                    self.log_debug(f"Trovata corrispondenza per parola intera: '{keyword}' nel nome del file {os.path.basename(file_path)}")
                result[0] = self.create_file_info(file_path)
                return
            
            # Verifica contenuto se richiesto
            if search_content and self.should_search_content(file_path):
                max_size_bytes = self.max_file_size_mb.get() * 1024 * 1024
                
                # Salta file troppo grandi
                if file_size > max_size_bytes:
                    self.log_debug(f"File {file_path} troppo grande per l'analisi del contenuto")
                    result[0] = None
                    return
                    
                content = self.get_file_content(file_path)
                if content:
                    matched = False
                    for keyword in keywords:
                        # Verifica se la ricerca di parole intere è attivata
                        if self.whole_word_search.get():
                            # Gestisce anche frasi con spazi usando la nuova funzione helper
                            if self.is_whole_word_match(keyword, content):
                                matched = True
                                break
                        # Se è un termine con spazi (contesto specifico)
                        elif ' ' in keyword:
                            if keyword.lower() in content.lower():
                                matched = True
                                break
                        # Ricerca normale
                        else:
                            if keyword.lower() in content.lower():
                                matched = True
                                break
                                
                    if matched:
                        if self.debug_mode and self.whole_word_search.get():
                            self.log_debug(f"Trovata corrispondenza per parola intera: '{keyword}' nel contenuto del file {os.path.basename(file_path)}")
                        result[0] = self.create_file_info(file_path)
                        return
            else:
                # Log per i file di sistema esclusi
                if search_content and os.path.splitext(file_path)[1].lower() in self.system_file_extensions:
                    self.log_debug(f"File di sistema escluso dall'analisi del contenuto: {file_path}")
                    
            result[0] = None
            
        except Exception as e:
            exception[0] = e
            self.log_debug(f"Errore durante l'elaborazione del file {file_path}: {str(e)}")
            result[0] = None
        finally:
            processing_completed[0] = True
    def manage_memory(self):
        """ Gestisce l'utilizzo della memoria durante la ricerca per evitare problemi"""
        import psutil
        import gc
        
        process = psutil.Process(os.getpid())
        memory_usage = process.memory_info().rss / 1024 / 1024  # MB
        
        # Se l'uso della memoria è troppo alto, esegui la garbage collection forzata
        if memory_usage > 500:  # 500 MB
            self.log_debug(f"Utilizzo memoria elevato: {memory_usage:.2f} MB. Esecuzione garbage collection.")
            gc.collect()
            
            # Dopo la GC, ricalcola l'utilizzo della memoria
            memory_usage = process.memory_info().rss / 1024 / 1024
            self.log_debug(f"Utilizzo memoria dopo GC: {memory_usage:.2f} MB")    

    def run_with_timeout(self, func, args=(), kwargs={}, timeout_sec=10):
        """ Esegue una funzione con un timeout, evitando blocchi indefiniti Returns: (result, completed)
            - result: risultato della funzione o None
            - completed: True se completata, False se interrotta per timeout
        """
        result = [None]
        completed = [False]
        
        def target():
            try:
                result[0] = func(*args, **kwargs)
            except Exception as e:
                self.log_debug(f"Eccezione nella funzione con timeout: {str(e)}")
            finally:
                completed[0] = True
        
        thread = threading.Thread(target=target)
        thread.daemon = True
        thread.start()
        thread.join(timeout_sec)
        
        return result[0], completed[0]
    
    def _get_all_descendants(self, widget):
        """Ottiene ricorsivamente tutti i widget discendenti"""
        descendants = []
        for child in widget.winfo_children():
            descendants.append(child)
            descendants.extend(self._get_all_descendants(child))
        return descendants
    
    def disable_all_controls(self):
        """Disabilita tutti i controlli UI durante la ricerca"""
        try:
            # Input e campi di testo
            self.path_entry["state"] = "disabled"
            self.keyword_entry["state"] = "disabled"
            
            # Pulsante Sfoglia
            if hasattr(self, 'browse_btn'):
                self.browse_btn["state"] = "disabled"
            
            # Pulsante Pulisci campi
            if hasattr(self, 'clear_btn'):
                self.clear_btn["state"] = "disabled"
            
            # Combobox e spinbox
            if hasattr(self, 'theme_combobox'):
                self.theme_combobox["state"] = "disabled"
            
            if hasattr(self, 'depth_spinbox'):
                self.depth_spinbox["state"] = "disabled"
            
            # Pulsanti principali
            if hasattr(self, 'search_button'):
                self.search_button["state"] = "disabled"
            
            # Gestisci i pulsanti dei filtri
            for btn_name in ['advanced_filters_btn', 'exclusions_btn', 'block_options_btn']:
                if hasattr(self, btn_name):
                    getattr(self, btn_name)["state"] = "disabled"
            
            # Gestisci specificamente il pulsante admin
            if hasattr(self, 'admin_button') and not self.is_admin:
                self.admin_button["state"] = "disabled"

            # Disabilita il pulsante delle opzioni di performance
            if hasattr(self, 'perf_options_btn'):
                self.perf_options_btn["state"] = "disabled"
            
            # Pulsanti di azione sui risultati
            if hasattr(self, 'copy_button'):
                self.copy_button["state"] = "disabled"
            if hasattr(self, 'compress_button'):
                self.compress_button["state"] = "disabled"
                
        except Exception as e:
            # Registra l'errore senza interrompere l'esecuzione
            self.log_debug(f"Errore nella disabilitazione dei controlli: {str(e)}")

    def enable_all_controls(self):
        """Riabilita tutti i controlli UI dopo la ricerca"""
        # Input e campi di testo
        self.path_entry["state"] = "normal"
        self.keyword_entry["state"] = "normal"
        
        # Pulsante Sfoglia
        if hasattr(self, 'browse_btn'):
            self.browse_btn["state"] = "normal"
        
        # Pulsante Pulisci campi
        if hasattr(self, 'clear_btn'):
            self.clear_btn["state"] = "normal"
        
        # Combobox e spinbox
        if hasattr(self, 'theme_combobox'):
            self.theme_combobox["state"] = "readonly"
        if hasattr(self, 'depth_spinbox'):
            self.depth_spinbox["state"] = "normal"
        
        # Pulsanti principali
        self.search_button["state"] = "normal"
        
        # Gestisci i pulsanti dei filtri
        for btn_name in ['advanced_filters_btn', 'exclusions_btn']:
            if hasattr(self, btn_name):
                getattr(self, btn_name)["state"] = "normal"

        # Gestisci i pulsanti  a blocchi
        if hasattr(self, 'block_options_btn'):
            self.block_options_btn["state"] = "normal"

        # Gestisci specificamente il pulsante admin
        if hasattr(self, 'admin_button') and not self.is_admin:
            self.admin_button["state"] = "normal"
        
        # Checkbox
        for widget in self.root.winfo_children():
            self._enable_checkbuttons_recursive(widget)
        
        # Pulsanti di azione sui risultati
        if hasattr(self, 'copy_button'):
            self.copy_button["state"] = "normal"
        if hasattr(self, 'compress_button'):
            self.compress_button["state"] = "normal"
        
        # Abilita il pulsante delle opzioni di performance
        if hasattr(self, 'perf_options_btn'):
            self.perf_options_btn["state"] = "normal"
                
    def _disable_checkbuttons_recursive(self, widget):
        """Disabilita ricorsivamente tutte le checkbox nei widget"""
        if isinstance(widget, ttk.Checkbutton):
            widget["state"] = "disabled"
        
        # Processa i widget figli ricorsivamente
        if hasattr(widget, 'winfo_children'):
            for child in widget.winfo_children():
                self._disable_checkbuttons_recursive(child)
                
    def _enable_checkbuttons_recursive(self, widget):
        """Riabilita ricorsivamente tutte le checkbox nei widget"""
        if isinstance(widget, ttk.Checkbutton):
            widget["state"] = "normal"
        
        # Processa i widget figli ricorsivamente
        if hasattr(widget, 'winfo_children'):
            for child in widget.winfo_children():
                self._enable_checkbuttons_recursive(child)
                
    def show_content_search_warning(self):
        """Mostra un avviso se la ricerca nei contenuti è attivata"""
        if self.search_content.get():  # Se la checkbox è attivata
            message = "Hai attivato la ricerca nei contenuti dei file."
            
            # MODIFICA QUI: Aggiungi un log di debug per verificare il valore effettivo
            self.log_debug(f"Livello di ricerca attuale: {self.search_depth.get()}")
            
            # Avviso specifico per la ricerca profonda
            if hasattr(self, 'search_depth') and self.search_depth.get() == "profonda":
                message += "\n\nATTENZIONE: Hai selezionato la ricerca PROFONDA che include tutti i tipi di file!\n" + \
                        "Questo potrebbe rallentare notevolmente la ricerca, soprattutto con grandi quantità di file."
            else:
                message += "\n\nQuesta operazione può richiedere molto più tempo, soprattutto con grandi quantità di file."
                
            message += "\n\nVuoi procedere con la ricerca nei contenuti?"
            
            return messagebox.askyesno(
                "Attenzione - Ricerca nei contenuti", 
                message,
                icon="warning"
            )
        return True   # Se la ricerca nei contenuti non è attivata, procedi senza avviso   
    
    def register_interrupt_handler(self):
        """Registra il gestore degli interrupt (CTRL+C)"""
        def handle_interrupt(sig, frame):
            if self.is_searching:
                self.stop_search_process()
            else:
                self.root.quit()
        
        signal.signal(signal.SIGINT, handle_interrupt)

    def log_debug(self, message):
        """Funzione per logging, stampa solo quando debug_mode è True"""
        if self.debug_mode:
            print(f"[DEBUG] {message}")

    def check_and_notify_missing_libraries(self):
        """Verifica e notifica l'utente di eventuali librerie mancanti"""
        missing = []
        
        if not file_format_support["docx"]:
            missing.append("python-docx (per file Word)")
        if not file_format_support["pdf"]: 
            missing.append("PyPDF2 (per file PDF)")
        if not file_format_support["pptx"]:
            missing.append("python-pptx (per file PowerPoint)")
        if not file_format_support["excel"]:
            missing.append("openpyxl (per file Excel)")
        if not file_format_support["rtf"]:
            missing.append("striprtf (per file RTF)")
        if not file_format_support["odt"]:
            missing.append("odfdo (per file OpenDocument)")
        
        if missing:
            message = "Alcune funzionalità di ricerca nei contenuti sono disabilitate.\n\n"
            message += "Per abilitare il supporto completo ai vari formati di file, installa le seguenti librerie:\n\n"
            
            for lib in missing:
                message += f"- {lib}\n"
            
            message += "\nPuoi installarle con il comando:\n"
            message += "pip install " + " ".join([lib.split(" ")[0] for lib in missing])
            
            # Mostra la notifica dopo un breve ritardo per permettere all'UI di caricarsi
            self.root.after(1000, lambda: messagebox.showinfo("Librerie opzionali mancanti", message))

    def update_datetime(self):
        current_time = datetime.now().strftime('%d/%m/%Y %H:%M:%S')
        self.datetime_var.set(f"Data: {current_time} | Utente: {self.user_var.get()}")
        self.root.after(1000, self.update_datetime)

    def browse_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.search_path.set(directory)
            self.update_disk_info(directory)

    def optimize_system_search(self, path):
        """Ottimizza la ricerca per percorsi di sistema come C:/ impostando parametri appropriati"""
        is_system_path = path.lower() in ["c:/", "c:\\", "d:/", "d:\\"]
        
        # Verifica se è un percorso di rete o server (inizia con \\ o //)
        is_network_path = path.startswith('\\\\') or path.startswith('//')
        
        # Verifica se è un percorso Unix root o home
        is_unix_root = path in ['/', '/home', '/usr', '/var', '/etc']
        
        if is_system_path or is_network_path or is_unix_root:
            # Salva i parametri attuali
            original_params = {
                "max_files": self.max_files_to_check.get(),
                "worker_threads": self.worker_threads.get(),
                "timeout": self.timeout_seconds.get() if self.timeout_enabled.get() else None
            }
            
            # Applica parametri ottimizzati
            self.max_files_to_check.set(1000000)  # Un milione di file
            self.worker_threads.set(min(12, os.cpu_count() or 4))  # Aumenta thread per server
            
            if self.timeout_enabled.get():
                self.timeout_seconds.set(max(3600, self.timeout_seconds.get()))  # Minimo 1 ora
            
            # Notifica l'utente
            messagebox.showinfo(
                "Ricerca su sistema o server",
                "Stai avviando una ricerca su un percorso di sistema o server.\n\n"
                "Per ottimizzare la ricerca, sono stati temporaneamente modificati alcuni parametri "
                "per consentire una scansione più approfondita.\n\n"
                "La ricerca potrebbe richiedere tempo e risorse significative."
            )
            
            return original_params
        return None
    
    def show_optimization_tips(self, path):
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]:
            # Determina la lettera del disco corrente
            if os.name == 'nt':  # Windows
                if path.lower() in ["c:/", "c:\\"]:
                    drive_letter = "C:"
                elif path.lower() in ["d:/", "d:\\"]:
                    drive_letter = "D:"
                elif path.lower() in ["e:/", "e:\\"]:
                    drive_letter = "E:"
                else:
                    drive_letter = path[:2] if len(path) >= 2 else "Disco"
            else:
                drive_letter = "Disco"  # Per sistemi non Windows
            
            optimization_done = False
            
            # Verifica se l'utente ha già implementato le ottimizzazioni
            if hasattr(self, 'excluded_paths'):
                windows_excluded = any("windows" in p.lower() for p in self.excluded_paths)
                program_files_excluded = any("program files (x86)" in p.lower() for p in self.excluded_paths)
                optimization_done = windows_excluded and program_files_excluded
            
            # Verifica la profondità di ricerca
            depth_optimized = self.max_depth >= 1 and self.max_depth <= 10
            
            # Se le ottimizzazioni non sono già state applicate, mostra il messaggio
            if not (optimization_done and depth_optimized):
                response = messagebox.askyesno(
                    "Ottimizza la ricerca",
                    f"Stai per avviare una ricerca sull'intero disco {drive_letter}\\\n\n"
                    "Per migliorare notevolmente le prestazioni, è consigliato:\n\n"
                    "1. Escludere le cartelle di sistema (Windows, Program Files)\n"
                    "2. Escludere le cartelle di altri utenti\n"
                    "3. Limitare la profondità di ricerca a 5-10 livelli\n\n"
                    "Vuoi applicare automaticamente queste ottimizzazioni?",
                    icon="question"
                )
                
                if response:
                    # Applica le ottimizzazioni
                    if not hasattr(self, 'excluded_paths'):
                        self.excluded_paths = []
                    
                    # Aggiungi esclusioni di sistema
                    system_paths = ["C:/Windows", "C:/Program Files", "C:/Program Files (x86)"]
                    for sys_path in system_paths:
                        if sys_path not in self.excluded_paths:
                            self.excluded_paths.append(sys_path)
                    
                    # Escludi altri utenti
                    current_user = getpass.getuser()
                    users_dir = "C:/Users"
                    if os.path.exists(users_dir):
                        try:
                            for user in os.listdir(users_dir):
                                user_path = os.path.join(users_dir, user)
                                if user.lower() != current_user.lower() and user.lower() not in ["public", "default", "all users"]:
                                    if user_path not in self.excluded_paths:
                                        self.excluded_paths.append(user_path)
                        except Exception as e:
                            self.log_debug(f"Errore nell'esclusione degli altri utenti: {str(e)}")
                    
                    # Imposta la profondità di ricerca a 7 (valore ragionevole)
                    self.depth_spinbox.set("7")
                    self.max_depth = 7
                    
                    # Notifica l'utente
                    messagebox.showinfo(
                        "Ottimizzazioni applicate",
                        f"Le ottimizzazioni sono state applicate:\n\n"
                        f"• Escluse {len(system_paths)} cartelle di sistema\n"
                        f"• Escluse le cartelle di altri utenti\n"
                        f"• Profondità di ricerca impostata a 7 livelli\n\n"
                        f"La ricerca sarà ora più veloce e stabile."
                    )
                    
                    return True
                else:
                    # NUOVA PARTE: Se l'utente risponde "No", imposta esplicitamente la profondità a 0
                    self.depth_spinbox.set("0")
                    self.max_depth = 0
                    
                    # Notifica l'utente che sta per iniziare una ricerca illimitata
                    self.log_debug("Utente ha rifiutato le ottimizzazioni, impostata profondità illimitata (0)")
                    
                    # Opzionalmente, puoi mostrare un messaggio di avviso
                    messagebox.showinfo(
                        "Ricerca illimitata",
                        "Stai per avviare una ricerca illimitata sull'intero disco.\n\n"
                        "La profondità di ricerca è stata impostata a illimitata (0).\n"
                        "Questa operazione potrebbe richiedere molto tempo e risorse di sistema."
                    )
                
                return False
            
            return False  # Per percorsi non di sistema, non fare nulla
    def debug_exclusions(self):
        """Visualizza lo stato corrente delle esclusioni per debug"""
        if hasattr(self, 'excluded_paths'):
            paths = '\n'.join(self.excluded_paths) if self.excluded_paths else "Nessun percorso escluso"
            messagebox.showinfo("Debug esclusioni", 
                            f"Stato esclusioni:\n{paths}\n\n"
                            f"Totale percorsi: {len(self.excluded_paths)}")
        else:
            messagebox.showinfo("Debug esclusioni", "Lista esclusioni non inizializzata")

    def manage_exclusions(self):
        """Apre una finestra di dialogo per gestire i percorsi esclusi dalla ricerca"""
        if not hasattr(self, 'excluded_paths'):
            self.excluded_paths = []
        dialog = ttk.Toplevel(self.root)
        dialog.title("Gestione percorsi esclusi")
        dialog.geometry("500x400")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Istruzioni
        ttk.Label(main_frame, text="Aggiungi cartelle da escludere dalla ricerca (es. C:/Windows)", 
                wraplength=480).pack(anchor=W, pady=(0, 10))
        
        # Frame per aggiungere nuovi percorsi
        add_frame = ttk.Frame(main_frame)
        add_frame.pack(fill=X, pady=5)
        
        path_var = StringVar()
        entry = ttk.Entry(add_frame, textvariable=path_var, width=40)
        entry.pack(side=LEFT, padx=(0, 5), fill=X, expand=YES)
        
        def browse_exclude_dir():
            directory = filedialog.askdirectory()
            if directory:
                path_var.set(directory)
        
        def add_exclusion():
            path = path_var.get().strip()
            if path and path not in self.excluded_paths:
                self.excluded_paths.append(path)
                update_list()
                path_var.set("")
        
        ttk.Button(add_frame, text="Sfoglia", command=browse_exclude_dir).pack(side=LEFT, padx=2)
        ttk.Button(add_frame, text="Aggiungi", command=add_exclusion).pack(side=LEFT, padx=2)
        
        # Lista dei percorsi esclusi
        ttk.Label(main_frame, text="Percorsi attualmente esclusi:").pack(anchor=W, pady=(10, 5))
        
        list_frame = ttk.Frame(main_frame)
        list_frame.pack(fill=BOTH, expand=YES, pady=5)
        
        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side=RIGHT, fill=Y)
        
        excluded_list = ttk.Treeview(list_frame, columns=("path",), show="headings", 
                                yscrollcommand=scrollbar.set, selectmode="extended")
        excluded_list.heading("path", text="Percorso")
        excluded_list.column("path", width=450)
        excluded_list.pack(fill=BOTH, expand=YES)
        
        scrollbar.config(command=excluded_list.yview)
        
        def update_list():
            # Pulisci la lista
            for item in excluded_list.get_children():
                excluded_list.delete(item)
            
            # Aggiungi i percorsi esclusi
            for path in self.excluded_paths:
                excluded_list.insert("", "end", values=(path,))
        
        def remove_selected():
            selected = excluded_list.selection()
            if selected:
                for item in selected:
                    values = excluded_list.item(item)["values"]
                    if values and values[0] in self.excluded_paths:
                        self.excluded_paths.remove(values[0])
                update_list()
        
        # Pulsanti per gestire la lista
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=10)
        
        ttk.Button(btn_frame, text="Rimuovi selezionati", command=remove_selected).pack(side=LEFT)
        
        # Preset di esclusioni comuni
        def add_common_exclusions():
            common_paths = [
                "C:/Windows", 
                "C:/Program Files", 
                "C:/Program Files (x86)",
                "C:/ProgramData",
                "C:/Users/All Users",
                "C:/Program Files (x86)/Client Active Directory Rights Management Services"
            ]
            added = 0
            for path in common_paths:
                if path not in self.excluded_paths:
                    self.excluded_paths.append(path)
                    added += 1
            
            if added > 0:
                update_list()
                messagebox.showinfo("Aggiunto", f"Aggiunti {added} percorsi comuni alle esclusioni")
        
        ttk.Button(btn_frame, text="Aggiungi esclusioni comuni", 
                command=add_common_exclusions).pack(side=LEFT, padx=5)
        
        # Aggiungi automaticamente altre cartelle utente
        def exclude_other_users():
            current_user = getpass.getuser()
            users_dir = "C:/Users"
            added = 0
            
            if os.path.exists(users_dir):
                try:
                    for user in os.listdir(users_dir):
                        user_path = os.path.join(users_dir, user)
                        # Escludi tutte le cartelle utente tranne quella dell'utente corrente
                        if user.lower() != current_user.lower() and user.lower() not in ["public", "default", "all users"]:
                            if user_path not in self.excluded_paths:
                                self.excluded_paths.append(user_path)
                                added += 1
                except Exception as e:
                    self.log_debug(f"Errore nell'esclusione degli altri utenti: {str(e)}")
            
            if added > 0:
                update_list()
                messagebox.showinfo("Aggiunto", f"Escluse {added} cartelle di altri utenti")
        
        ttk.Button(btn_frame, text="Escludi altri utenti", 
                command=exclude_other_users).pack(side=LEFT, padx=5)
        
        # Pulsanti OK/Annulla
        dialog_btn_frame = ttk.Frame(main_frame)
        dialog_btn_frame.pack(fill=X, pady=(10, 0))
        
        ttk.Button(dialog_btn_frame, text="OK", command=dialog.destroy).pack(side=RIGHT)
        
        # Inizializza la lista
        update_list()

        # Aggiorna la finestra per ottenere le dimensioni corrette
        dialog.update_idletasks()
        
        # Ottieni le dimensioni della finestra
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        
        # Calcola la posizione x,y per centrare la finestra
        screen_width = dialog.winfo_screenwidth()
        screen_height = dialog.winfo_screenheight()
        x = (screen_width // 2) - (width // 1)
        y = (screen_height // 2) - (height // 1)
        
        # Imposta la geometria con la posizione calcolata
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima per la finestra
        dialog.minsize(800, 600)

    def restart_as_admin(self):
        """Riavvia l'applicazione con privilegi di amministratore"""
        try:
            import ctypes, sys, os
            if sys.platform == 'win32':
                script = sys.argv[0]
                params = ' '.join(sys.argv[1:])
                self.log_debug("Tentativo di riavvio come amministratore")
                
                # Se il percorso di ricerca è impostato, passa quel percorso come parametro
                if self.search_path.get():
                    if not params:
                        params = f'"{self.search_path.get()}"'
                    
                # Avvia l'applicazione con privilegi elevati
                ctypes.windll.shell32.ShellExecuteW(None, "runas", sys.executable, f'"{script}" {params}', None, 1)
                
                # Chiude l'istanza corrente
                self.root.quit()
            else:
                messagebox.showinfo("Informazione", "Questa funzione è disponibile solo su Windows")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile riavviare come amministratore: {str(e)}")
            self.log_debug(f"Errore nel riavvio come admin: {str(e)}")

    def optimize_disk_search_order(self, path, directories):
        """Ottimizza l'ordine di esplorazione delle directory quando si cerca in un disco di sistema.
        Dà priorità alle cartelle più importanti come Users, Documents, ecc.
        """
        if path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"]:
            prioritized = []
            normal = []
            
            # Cartelle ad alta priorità (esplora per prime)
            high_priority_folders = ["users", "documents", "documents and settings", "desktop", "downloads"]
            
            # Cartelle a bassa priorità (esplora per ultime)
            low_priority_folders = ["windows", "program files", "program files (x86)", "programdata", "$recycle.bin", "system volume information"]
            
            for dir_path in directories:
                dir_name = os.path.basename(dir_path).lower()
                
                if any(priority in dir_name for priority in high_priority_folders):
                    prioritized.insert(0, dir_path)  # Aggiungi all'inizio per priorità più alta
                elif any(low in dir_name for low in low_priority_folders):
                    normal.append(dir_path)  # Aggiungi alla fine
                else:
                    prioritized.append(dir_path)  # Priorità media
                    
            return prioritized + normal
        
        # Se non è un percorso di sistema, restituisci l'elenco originale
        return directories

    def start_search(self):
        # Pulisci risultati precedenti
        for item in self.results_list.get_children():
            self.results_list.delete(item)
    
        # Ottieni i valori direttamente dai widget
        search_path = self.search_path.get().strip()
        keywords = self.keyword_entry.get().strip()
        
        # Stampa di debug (puoi rimuoverla in produzione)
        print(f"Debug - Percorso: '{search_path}', Parole chiave: '{keywords}'")
        
        # 1. Verifica il percorso
        if not search_path:
            messagebox.showerror("Errore", "Inserisci una directory di ricerca valida")
            return
        
        if not os.path.exists(search_path):
            messagebox.showerror("Errore", f"Il percorso specificato non esiste:\n{search_path}")
            return
        
        # 2. Verifica le parole chiave
        placeholder = "Scrivi la parola da ricercare..."
        keywords_raw = self.keyword_entry.get()  # Ottieni il valore grezzo senza strip()

        # Verifica se è vuoto, solo spazi, o è il placeholder
        if (not keywords_raw or 
            keywords_raw.strip() == "" or 
            keywords_raw == placeholder or 
            self.keyword_entry.cget("foreground") == "gray"):
            
            messagebox.showerror("Errore", "Inserisci le parole chiave da cercare")
            return

        # 3. Verifica colore (controllo extra per il placeholder)
        try:
            if self.keyword_entry.cget("foreground") == "gray":
                messagebox.showerror("Errore", "Inserisci le parole chiave da cercare")
                return
        except:
            pass
            
        # Mostra avviso se la ricerca nei contenuti è attivata
        if not self.show_content_search_warning():
            return  # Interrompi se l'utente annulla
        
        # Aggiorna le informazioni del disco qui, quando l'utente clicca su cerca
        self.update_disk_info()

        self.show_optimization_tips(self.search_path.get())
        
        # Reset per la nuova ricerca
        self.stop_search = False
        self.stop_button["state"] = "normal"
        original_params = self.optimize_system_search(self.search_path.get())
        
        # Aggiorna l'orario di avvio e resetta l'orario di fine
        current_time = datetime.now().strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")

        # Salva i parametri di ricerca
        self.last_search_params = {
            "path": self.search_path.get(),
            "keywords": self.keywords.get(),
            "search_files": self.search_files.get(),
            "search_folders": self.search_folders.get(),
            "search_content": self.search_content.get()
        }
        
        # Aggiungi alla cronologia di ricerca se non già presente
        if self.keywords.get() and self.keywords.get() not in [h["keywords"] for h in self.search_history]:
            self.search_history.append(self.last_search_params.copy())
            if len(self.search_history) > 10:  # Mantieni solo le ultime 10 ricerche
                self.search_history.pop(0)
        
        self.is_searching = True
        self.disable_all_controls()  # Disabilita tutti i controlli
        self.stop_button["state"] = "normal"  # Solo questo è abilitato durante la ricerca
        self.progress_bar["value"] = 0
        self.status_label["text"] = "Ricerca in corso..."
        
        # Aggiorna il valore della profondità massima
        try:
            if hasattr(self, 'depth_spinbox'):
                self.max_depth = int(self.depth_spinbox.get())
            else:
                # Se depth_spinbox non esiste, usa il valore predefinito 
                # memorizzato nella variabile self.max_depth
                pass  # Mantiene il valore self.max_depth esistente
        except ValueError:
            self.max_depth = 0  # Valore predefinito se non valido
        
        # Ottieni le parole chiave di ricerca
        search_terms = [term.strip() for term in self.keywords.get().split(',') if term.strip()]
            
        # Avvia la ricerca in un thread separato
        self.search_results = []  # Resetta i risultati
        search_thread = threading.Thread(target=self._search_thread, 
                                        args=(self.search_path.get(), search_terms, self.search_content.get()))
        search_thread.daemon = True
        search_thread.start()
        
        # Memorizza l'orario di inizio come oggetto datetime
        self.search_start_time = datetime.now()
        current_time = self.search_start_time.strftime('%H:%M')
        self.start_time_label.config(text=current_time)
        self.end_time_label.config(text="--:--")
        self.total_time_label.config(text="--:--")  # Reset del tempo totale

        # Avvia l'aggiornamento della progress bar
        self.update_progress()
        # Salva original_params come attributo per ripristinarli dopo la ricerca
        if original_params:
            self.original_system_search_params = original_params

        current_theme = "dark"  # Sostituisci con la logica per determinare il tema corrente
        self.update_theme_colors(current_theme)

        # Aggiungi questo nuovo metodo per calcolare il tempo totale
    def update_total_time(self):
        """Calcola e visualizza il tempo totale trascorso tra inizio e fine ricerca"""
        if self.search_start_time:
            end_time = datetime.now()
            time_diff = end_time - self.search_start_time
            
            # Converti in minuti e secondi
            total_seconds = int(time_diff.total_seconds())
            minutes = total_seconds // 60
            seconds = total_seconds % 60
            
            # Formatta la stringa del tempo totale
            if minutes > 0:
                total_time_str = f"{minutes}min {seconds}sec"
            else:
                total_time_str = f"{seconds}sec"
            
            self.total_time_label.config(text=total_time_str) 
    # Esempio di utilizzo per cambiare il tema
    def change_theme(self, theme):
        self.update_theme_colors(theme)

    def calculate_block_priority(self, directory):
        """Calcola la priorità del blocco (numerica, più bassa = più alta priorità)"""
        # Se l'opzione è disabilitata, usa priorità standard per tutti
        if not self.prioritize_user_folders.get():
            return 1
            
        # Altrimenti, applica la prioritizzazione configurata
        # Dai priorità alle cartelle degli utenti
        if "users" in directory.lower() or "documenti" in directory.lower() or "desktop" in directory.lower():
            return 0  # Alta priorità
        # Dai priorità medio-alta alle cartelle di dati
        elif "data" in directory.lower() or "database" in directory.lower() or "downloads" in directory.lower():
            return 1  # Priorità media-alta
        # Dai priorità bassa a cartelle di sistema
        elif "windows" in directory.lower() or "program files" in directory.lower():
            return 3  # Priorità bassa
        # Priorità standard per altre cartelle
        return 2
    
    def initialize_block_queue(self, root_path, block_queue, visited_dirs, files_checked, keywords, search_content, futures):
        """Inizializza la coda di blocchi con il percorso principale"""
        try:
            items = os.listdir(root_path)
            
            # Prima aggiungi le directory come blocchi separati
            for item in items:
                item_path = os.path.join(root_path, item)
                try:
                    if os.path.isdir(item_path):
                        # Salta directory nascoste se richiesto
                        if self.ignore_hidden.get() and (item.startswith('.') or 
                            (os.name == 'nt' and os.path.exists(item_path) and 
                            os.stat(item_path).st_file_attributes & 2)):
                            continue
                        
                        # Salta directory escluse
                        if hasattr(self, 'excluded_paths') and any(
                            item_path.lower().startswith(excluded.lower()) 
                            for excluded in self.excluded_paths):
                            continue
                            
                        # Calcola priorità e aggiungi alla coda
                        priority = self.calculate_block_priority(item_path)
                        block_queue.put((priority, item_path))
                        visited_dirs.add(os.path.realpath(item_path))
                except Exception as e:
                    self.log_debug(f"Errore nell'aggiunta del blocco {item_path}: {str(e)}")
            
            # Poi processa i file nella directory principale
            for item in items:
                if self.stop_search:
                    return
                    
                item_path = os.path.join(root_path, item)
                try:
                    if not os.path.isdir(item_path):
                        if self.search_files.get():
                            # Verifica file nascosti
                            if self.ignore_hidden.get() and (item.startswith('.') or 
                                (os.name == 'nt' and os.path.exists(item_path) and 
                                os.stat(item_path).st_file_attributes & 2)):
                                continue
                            
                            # Processa direttamente i file nella directory principale
                            nonlocal_files_checked = files_checked[0]
                            nonlocal_files_checked += 1
                            files_checked[0] = nonlocal_files_checked
                            if nonlocal_files_checked > self.max_files_to_check.get():
                                self.stop_search = True
                                return
                            
                            future = self.search_executor.submit(self.process_file, item_path, keywords, search_content)
                            futures.append(future)
                except Exception as e:
                    self.log_debug(f"Errore nell'elaborazione del file {item_path}: {str(e)}")
                    
        except PermissionError:
            self.log_debug(f"Permesso negato per la directory {root_path}")
        except Exception as e:
            self.log_debug(f"Errore nell'inizializzazione dei blocchi da {root_path}: {str(e)}")

    def process_blocks(self, block_queue, visited_dirs, start_time, timeout, is_system_search, 
                    files_checked, dirs_checked, last_update_time, path, keywords, search_content, futures):
        """Elabora i blocchi dalla coda in base alla priorità"""
        # Determina il numero massimo di file per blocco in base alle impostazioni
        max_files_in_block = self.max_files_per_block.get()
        
        # Adatta automaticamente la dimensione del blocco se richiesto
        if self.block_size_auto_adjust.get():
            # Ottimizzazione per ricerche di sistema o cartelle molto grandi
            if is_system_search:
                max_files_in_block = max(2000, max_files_in_block)
            elif files_checked[0] > 100000:
                max_files_in_block = max(5000, max_files_in_block)
        
        # Ottimizza il numero di blocchi paralleli in base al carico di sistema
        active_blocks = 0
        max_parallel = self.max_parallel_blocks.get()
        
        # Utilizziamo un set per tracciare i blocchi già processati più efficiente
        processed_blocks = set()
        
        while not block_queue.empty() and not self.stop_search:
            # Verifica timeout
            if timeout and time.time() - start_time > timeout:
                self.progress_queue.put(("timeout", "Timeout raggiunto"))
                return
            
            try:
                # Prendi il blocco con priorità più alta
                _, current_block = block_queue.get(block=False)
                
                # Salta blocchi già processati
                if current_block in processed_blocks:
                    continue
                    
                processed_blocks.add(current_block)
                
                # Aggiorna lo stato
                current_time = time.time()
                if current_time - last_update_time[0] >= 0.5:  # Aggiorna ogni mezzo secondo
                    elapsed_time = current_time - start_time
                    self.progress_queue.put(("status", 
                        f"Analisi blocco: {current_block} (Cartelle: {dirs_checked[0]}, File: {files_checked[0]}, Tempo: {int(elapsed_time)}s)"))
                    self.progress_queue.put(("progress", 
                        min(90, int((files_checked[0] / max(1, self.max_files_to_check.get())) * 100))))
                    last_update_time[0] = current_time
                
                # Aggiungi blocchi di livello inferiore dalla directory corrente
                try:
                    # Blocco attualmente in elaborazione
                    dirs_checked[0] += 1
                    
                    # Lista elementi nella directory
                    items = os.listdir(current_block)
                    
                    # Prima processa le sottocartelle (aggiungi nuovi blocchi)
                    subfolders = []
                    for item in items:
                        if self.stop_search:
                            return
                            
                        item_path = os.path.join(current_block, item)
                        
                        # Salta file/cartelle nascoste
                        try:
                            if self.ignore_hidden.get() and (item.startswith('.') or 
                                (os.name == 'nt' and os.path.exists(item_path) and 
                                os.stat(item_path).st_file_attributes & 2)):
                                continue
                        except Exception as e:
                            self.log_debug(f"Errore nel controllo hidden per {item_path}: {str(e)}")
                            continue
                            
                        # Gestione sottodirectory
                        try:
                            if os.path.isdir(item_path):
                                # Verifica se la cartella è già stata visitata
                                try:
                                    real_path = os.path.realpath(item_path)
                                    if real_path in visited_dirs:
                                        continue
                                    visited_dirs.add(real_path)
                                except:
                                    if item_path in visited_dirs:
                                        continue
                                    visited_dirs.add(item_path)
                                
                                # Verifica se il percorso deve essere escluso
                                if hasattr(self, 'excluded_paths') and any(
                                    item_path.lower().startswith(excluded.lower()) 
                                    for excluded in self.excluded_paths):
                                    continue
                                
                                # Aggiungi alla lista delle sottocartelle
                                subfolders.append(item_path)
                                
                                # Verifica corrispondenza nome cartella
                                if self.search_folders.get():
                                    if any(keyword.lower() in item.lower() for keyword in keywords):
                                        folder_info = self.create_folder_info(item_path)
                                        self.search_results.append(folder_info)
                                        
                        except Exception as e:
                            self.log_debug(f"Errore nell'analisi della directory {item_path}: {str(e)}")
                    
                    # Aggiunta sottocartelle alla coda con priorità calcolata
                    for subfolder in subfolders:
                        # Verifica limite di profondità
                        folder_depth = subfolder.count(os.path.sep) - path.count(os.path.sep)
                        if self.max_depth > 0 and folder_depth > self.max_depth:
                            continue
                            
                        priority = self.calculate_block_priority(subfolder)
                        block_queue.put((priority, subfolder))
                    
                    # Processa i file nella directory corrente
                    for item in items:
                        if self.stop_search:
                            return
                            
                        item_path = os.path.join(current_block, item)
                        
                        # Salta i file nascosti e le directory (già processate)
                        try:
                            # Controllo file nascosto
                            if self.ignore_hidden.get() and (item.startswith('.') or 
                                        (os.name == 'nt' and os.path.exists(item_path) and 
                                            os.stat(item_path).st_file_attributes & 2)):
                                continue
                                
                            # Salta le directory (già processate)
                            if os.path.isdir(item_path):
                                continue
                        except Exception as e:
                            self.log_debug(f"Errore nel controllo del tipo per {item_path}: {str(e)}")
                            continue
                            
                        # Processa i file
                        if self.search_files.get():
                            try:
                                # Verifica limite file
                                files_checked[0] += 1
                                if files_checked[0] % 1000 == 0:
                                    self.manage_memory()
                                if files_checked[0] > self.max_files_to_check.get():
                                    self.stop_search = True
                                    self.progress_queue.put(("status", 
                                        f"Limite di {self.max_files_to_check.get()} file controllati raggiunto. "
                                        f"Aumenta il limite nelle opzioni per cercare più file."))
                                    return
                                try:
                                    item_path = os.path.join(current_block, item)
                                    if os.path.isfile(item_path) and not os.path.islink(item_path):
                                        file_size = os.path.getsize(item_path)
                                        self.current_search_size += file_size
                                        
                                        # Aggiorna la dimensione mostrata ogni 1000 file o ogni 5 secondi
                                        if files_checked % 1000 == 0 or (time.time() - last_size_update_time) > 5:
                                            self.progress_queue.put(("update_dir_size", self.current_search_size))
                                            last_size_update_time = time.time()
                                except:
                                    pass  # Ignora errori durante il calcolo della dimensione
                                # Gestione sicura dell'executor
                                try:
                                    if self.search_executor and not self.search_executor._shutdown:
                                        future = self.search_executor.submit(self.process_file, item_path, keywords, search_content)
                                        futures.append(future)
                                    else:
                                        result = self.process_file(item_path, keywords, search_content)
                                        if result:
                                            self.search_results.append(result)
                                except Exception as e:
                                    self.log_debug(f"Errore nell'elaborazione parallela del file {item_path}: {str(e)}")
                                    try:
                                        result = self.process_file(item_path, keywords, search_content)
                                        if result:
                                            self.search_results.append(result)
                                    except:
                                        pass
                                    
                            except Exception as e:
                                self.log_debug(f"Errore nell'aggiunta del file {item_path} alla coda: {str(e)}")
                                continue
                    
                except PermissionError:
                    if self.skip_permission_errors.get():
                        self.log_debug(f"Saltata directory con permesso negato: {current_block}")
                    else:
                        dir_name = os.path.basename(current_block)
                        parent_dir = os.path.dirname(current_block)
                        
                        is_user_folder = (parent_dir.lower() in ["c:/users", "c:\\users"] and 
                                        dir_name.lower() != getpass.getuser().lower())
                        
                        if is_user_folder:
                            self.log_debug(f"Cartella di un altro utente inaccessibile: {current_block}")
                            self.progress_queue.put(("status", f"Saltata cartella utente protetta: {current_block}"))
                        else:
                            self.log_debug(f"Permesso negato per la directory {current_block}")
                            self.progress_queue.put(("status", f"Permesso negato: {current_block}"))
                except Exception as e:
                    self.log_debug(f"Errore durante l'analisi della directory {current_block}: {str(e)}")
                    
            except queue.Empty:
                break

    def _search_thread(self, path, keywords, search_content):
        try:
            # Inizializza i contatori di file e directory esaminati e il tempo di inizio
            files_checked = [0]  # Uso una lista per poter modificare il valore nelle funzioni chiamate
            dirs_checked = [0]
            start_time = time.time()
            timeout = self.timeout_seconds.get() if self.timeout_enabled.get() else None
            
            self.current_search_size = 0  # Dimensione totale dei file trovati
            last_size_update_time = time.time()  # Per aggiornamenti periodici

            # Determina se si tratta di una ricerca completa del sistema (C:/ o simile)
            is_system_search = path.lower() in ["c:/", "c:\\", "d:/", "d:\\", "e:/", "e:\\"] or path in [os.path.abspath("/")]

            # Per ricerche di sistema, adatta automaticamente i parametri
            if is_system_search:
                # Informa l'utente che la ricerca potrebbe richiedere molto tempo
                self.progress_queue.put(("status", "Ricerca completa del sistema in corso - potrebbe richiedere molto tempo"))
                
                # Temporaneamente aumenta i limiti per la ricerca di sistema
                original_max_files = self.max_files_to_check.get()
                original_timeout = timeout
                
                # Imposta limiti più elevati per la ricerca completa
                self.max_files_to_check.set(5000000)  # 5 milioni di file
                
                # Disabilita temporaneamente il timeout o aumentalo significativamente
                if timeout and timeout < 3600:
                    timeout = 3600 * 8  # 8 ore
                
                # Avviso all'utente
                self.progress_queue.put(("status", "Ricerca completa avviata - parametri adattati per ricerca approfondita"))
            
            # Variabili per gestire l'aggiornamento ogni secondo
            last_update_time = [time.time()]
            
            # Crea un executor per processare i file in parallelo in modo sicuro
            try:
                max_workers = max(1, min(32, self.worker_threads.get()))  # Garantisce un valore valido
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=max_workers)
            except Exception as e:
                self.log_debug(f"Errore nella creazione del ThreadPoolExecutor: {str(e)}")
                # Fallback a un valore sicuro
                self.search_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)
            
            futures = []
            
            # Insieme per tenere traccia delle cartelle visitate (evita loop infiniti con symlink)
            visited_dirs = set()
            
            # Coda di priorità per i blocchi di ricerca
            block_queue = queue.PriorityQueue()

            # Aggiorna lo stato iniziale
            self.progress_queue.put(("status", f"Inizio ricerca a blocchi in: {path} (Profondità: {'illimitata' if self.max_depth == 0 else self.max_depth})"))
            
            # Assicurati che la lista di esclusioni sia inizializzata
            if not hasattr(self, 'excluded_paths'):
                self.excluded_paths = []
                
            # Avvia la ricerca a blocchi
            self.initialize_block_queue(path, block_queue, visited_dirs, files_checked, keywords, search_content, futures)
            self.process_blocks(block_queue, visited_dirs, start_time, timeout, is_system_search, 
                            files_checked, dirs_checked, last_update_time, path, keywords, search_content, futures)
            
            # Ripristina i parametri originali se erano stati modificati
            if is_system_search and self.max_depth == 0 and 'original_max_files' in locals():
                self.max_files_to_check.set(original_max_files)
            
            # Aggiorna lo stato finale di analisi
            self.progress_queue.put(("status", f"Elaborazione risultati... (analizzati {files_checked[0]} file in {dirs_checked[0]} cartelle)"))
            
            # Raccolta risultati dalle future con aggiornamento temporizzato
            completed = 0
            total_futures = len(futures)
            last_update_time[0] = time.time()
            
            # Gestione sicura delle future
            if self.search_executor and not self.search_executor._shutdown and futures:
                try:
                    for future in concurrent.futures.as_completed(futures):
                        if self.stop_search:
                            break
                            
                        try:
                            result = future.result()
                            if result:
                                self.search_results.append(result)
                            
                            completed += 1
                            
                            # Aggiorna il progresso e il tempo più frequentemente
                            current_time = time.time()
                            if completed % 20 == 0 or current_time - last_update_time[0] > 2:  # Aggiorna ogni 20 file o ogni 2 secondi
                                progress = 90 + min(10, int((completed / max(1, len(futures))) * 10))
                                self.progress_queue.put(("progress", progress))
                                
                                elapsed_time = current_time - start_time
                                self.progress_queue.put(("status", f"Elaborati {completed}/{len(futures)} file (tempo: {int(elapsed_time)}s)"))
                                
                                # Aggiorna anche il tempo totale durante la ricerca
                                if hasattr(self, 'search_start_time'):
                                    now = datetime.now()
                                    time_diff = now - self.search_start_time
                                    # Converti in minuti e secondi
                                    total_seconds = int(time_diff.total_seconds())
                                    minutes = total_seconds // 60
                                    seconds = total_seconds % 60
                                    
                                    if minutes > 0:
                                        total_time_str = f"{minutes}min {seconds}sec"
                                    else:
                                        total_time_str = f"{seconds}sec"
                                        
                                    self.progress_queue.put(("update_total_time", total_time_str))
                                
                                last_update_time[0] = current_time
                        except Exception as e:
                            self.log_debug(f"Errore nell'elaborazione di un risultato: {str(e)}")
                except Exception as e:
                    self.log_debug(f"Errore nella raccolta dei risultati: {str(e)}")
            
            # Completa la ricerca in modo sicuro
            try:
                if self.search_executor and not self.search_executor._shutdown:
                    self.search_executor.shutdown(wait=False)
                self.search_executor = None
            except Exception as e:
                self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                self.search_executor = None
            
            # Ripristina i parametri originali se erano stati modificati per la ricerca di sistema
            if hasattr(self, 'original_system_search_params'):
                try:
                    self.max_files_to_check.set(self.original_system_search_params["max_files"])
                    self.worker_threads.set(self.original_system_search_params["worker_threads"])
                    delattr(self, 'original_system_search_params')
                except:
                    pass

            # Riporta il risultato finale
            elapsed_time = time.time() - start_time
            self.progress_queue.put(("status", 
                f"Ricerca completata! Analizzati {files_checked[0]} file in {dirs_checked[0]} cartelle in {int(elapsed_time)} secondi."))
            
            # Ordina i risultati per tipo e nome
            self.search_results.sort(key=lambda x: (x[0], x[1]))
            
            self.log_debug(f"Ricerca completata. Trovati {len(self.search_results)} risultati")
            self.progress_queue.put(("complete", "Ricerca completata"))
            
        except Exception as e:
            error_msg = f"Si è verificato un errore durante la ricerca: {str(e)}\n{traceback.format_exc()}"
            self.log_debug(error_msg)
            self.progress_queue.put(("error", error_msg))

    def is_whole_word_match(self, keyword, text):
        """Verifica se la keyword è presente nel testo come parola intera."""
        try:
            pattern = r'\b' + re.escape(keyword.lower()) + r'\b'
            return re.search(pattern, text.lower()) is not None
        except re.error as e:
            self.log_debug(f"Errore regex con il termine '{keyword}': {str(e)}")
            
            # Fallback manuale se la regex fallisce
            keyword_lower = keyword.lower()
            text_lower = text.lower()
            
            if keyword_lower not in text_lower:
                return False
                
            # Verifica manuale dei confini parola
            index = text_lower.find(keyword_lower)
            while index != -1:
                # Controlla il carattere prima della keyword
                has_char_before = index > 0 and text_lower[index-1].isalnum()
                
                # Controlla il carattere dopo la keyword
                end_pos = index + len(keyword_lower)
                has_char_after = end_pos < len(text_lower) and text_lower[end_pos].isalnum()
                
                # Se non ci sono caratteri alfanumerici prima e dopo, è una parola intera
                if not has_char_before and not has_char_after:
                    return True
                    
                # Cerca la prossima occorrenza
                index = text_lower.find(keyword_lower, index + 1)
            
            return False
                
    def search_current_user_only(self):
        """Imposta la ricerca solo nella cartella dell'utente corrente"""
        user_folder = os.path.join("C:/Users", getpass.getuser())
        if os.path.exists(user_folder):
            self.search_path.set(user_folder)
            
            # Mostra la conferma solo se l'utente non annulla il warning sulla ricerca nei contenuti
            if self.show_content_search_warning():
                messagebox.showinfo(
                    "Ricerca configurata", 
                    f"La ricerca è stata configurata per cercare solo nella tua cartella utente:\n{user_folder}\n\n"
                    "Questa modalità evita problemi di permesso con altre cartelle protette."
                )
                # Avvia direttamente la ricerca
                self.start_search()
        else:
            messagebox.showerror("Errore", "Impossibile trovare la tua cartella utente")

    def create_file_info(self, file_path):
        """Crea le informazioni del file per la visualizzazione"""
        try:
            file_size = os.path.getsize(file_path)
            modified_time = datetime.fromtimestamp(os.path.getmtime(file_path))
            created_time = datetime.fromtimestamp(os.path.getctime(file_path))
            file_name = os.path.basename(file_path)
            file_extension = os.path.splitext(file_name)[1].lower()
            
            # Determina il tipo di file
            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type:
                file_type = mime_type.split('/')[0].capitalize()
                if file_type == "Application":
                    if "pdf" in mime_type:
                        file_type = "PDF"
                    elif "word" in mime_type or file_extension == ".docx" or file_extension == ".doc":
                        file_type = "Word"
                    elif "excel" in mime_type or file_extension in [".xlsx", ".xls"]:
                        file_type = "Excel"
                    elif "powerpoint" in mime_type or file_extension in [".pptx", ".ppt"]:
                        file_type = "PowerPoint"
                    else:
                        file_type = "Documento"
            else:
                # Fallback basato sull'estensione
                if file_extension in ['.txt', '.md', '.rtf']:
                    file_type = "Testo"
                elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                    file_type = "Immagine"
                elif file_extension in ['.mp3', '.wav', '.ogg', '.flac']:
                    file_type = "Audio"
                elif file_extension in ['.mp4', '.avi', '.mkv', '.mov']:
                    file_type = "Video"
                elif file_extension in ['.exe', '.dll', '.bat']:
                    file_type = "Eseguibile"
                else:
                    file_type = "File"
            
            # Formatta dimensione file
            if file_size < 1024:
                size_str = f"{file_size} B"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.1f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.1f} MB"
            
            return (
                file_type,
                file_name,
                size_str,
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                file_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni del file {file_path}: {str(e)}")
            return (
                "File",
                os.path.basename(file_path),
                "N/A",
                "N/A",
                "N/A",
                file_path
            )

    def create_folder_info(self, folder_path):
        """Crea le informazioni della cartella per la visualizzazione"""
        try:
            modified_time = datetime.fromtimestamp(os.path.getmtime(folder_path))
            created_time = datetime.fromtimestamp(os.path.getctime(folder_path))
            folder_name = os.path.basename(folder_path)
            
            return (
                "Directory",
                folder_name,
                "",  # Le cartelle non hanno dimensione
                modified_time.strftime('%d/%m/%Y %H:%M'),
                created_time.strftime('%d/%m/%Y %H:%M'),
                folder_path
            )
        except Exception as e:
            self.log_debug(f"Errore nel creare le informazioni della cartella {folder_path}: {str(e)}")
            return (
                "Directory",
                os.path.basename(folder_path),
                "",
                "N/A",
                "N/A",
                folder_path
            )
            
    def should_search_content(self, file_path):
        """Versione ottimizzata per determinare se analizzare il contenuto del file"""
        # Prima verifica le condizioni più veloci (principio fail-fast)
        if not self.search_content.get():
            return False
                    
        ext = os.path.splitext(file_path)[1].lower()
        
        # Blocca specificamente i file .doc per evitare blocchi
        if ext == '.doc':
            self.log_debug(f"File .doc temporaneamente escluso dall'analisi: {file_path}")
            return False
        
        # Seleziona il livello di ricerca attuale
        search_level = self.search_depth.get()
        
        # Ottieni le estensioni personalizzate dell'utente
        custom_extensions = self.get_extension_settings(search_level)
        
        # PRIORITÀ #1: Se l'estensione è stata aggiunta manualmente, cerca sempre il contenuto
        if ext in custom_extensions:
            # Aggiungi log per debug
            self.log_debug(f"Ricerca contenuto in file con estensione personalizzata: {ext} - {file_path}")
            return True
        
        # PRIORITÀ #2: In modalità profonda senza estensioni personalizzate, cerca tutto
        if search_level == "profonda" and not custom_extensions:
            return True
        
        # Liste predefinite nel codice per ciascun livello
        base_extensions = ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log', 
                        '.docx', '.pdf', '.pptx', '.xlsx', '.rtf', '.odt', '.xls', '.doc']
                        
        advanced_extensions = base_extensions + ['.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', 
                                            '.vbs', '.js', '.config', '.ini', '.reg']
                                            
        deep_extensions = advanced_extensions + ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.mp3', '.mp4', 
                                            '.avi', '.mov', '.mkv', '.wav', '.flac', '.zip', '.rar', 
                                            '.7z', '.tar', '.gz', '.iso', '.psd', '.ai', '.svg']
        
        # Verifica il livello predefinito se non è nelle estensioni personalizzate
        if search_level == "base" and ext in base_extensions:
            return True
        elif search_level == "avanzata" and ext in advanced_extensions:
            return True
        elif search_level == "profonda" and ext in deep_extensions:
            return True
        
        # Verifica supporto formati specifici
        if (ext == '.docx' and file_format_support["docx"]) or \
        (ext == '.pdf' and file_format_support["pdf"]) or \
        (ext in {'.pptx', '.ppt'} and file_format_support["pptx"]) or \
        (ext in {'.xls', '.xlsx'} and file_format_support["excel"]) or \
        (ext == '.rtf' and file_format_support["rtf"]) or \
        (ext == '.odt' and file_format_support["odt"]):
            return True
                
        return False

    def should_skip_file(self, file_path):
        """Verifica se un file deve essere saltato durante l'analisi del contenuto"""
        ext = os.path.splitext(file_path)[1].lower()
        skip_type = "File di sistema" if ext in self.system_file_extensions else "File"
        skip_filename = os.path.basename(file_path)
        
        # IMPORTANTE: Verifica se il file è nelle estensioni personalizzate
        # Se è stato aggiunto manualmente, NON deve essere mai saltato
        search_level = self.search_depth.get()
        custom_extensions = self.get_extension_settings(search_level) 
        
        # Non saltare mai i file con estensioni aggiunte manualmente
        if ext in custom_extensions:
            # Aggiungi log per debug
            self.log_debug(f"File NON saltato perché in estensioni personalizzate: {file_path}")
            return False
        
        # Salta i file di Rights Management Services
        if "Rights Management Services" in file_path or "IRMProtectors" in file_path:
            self.log_debug(f"Saltato file protetto: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Rights Management Services")
            return True
                    
        # Salta file con estensioni problematiche
        problematic_extensions = [".msoprotector.doc", ".msoprotector.ppt", ".msoprotector.xls"]
        if any(ext in file_path for ext in problematic_extensions):
            self.log_debug(f"Saltato file con formato problematico: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "Formato problematico")
            return True

        # Salta file di sistema solo se non sono nelle estensioni personalizzate
        if self.exclude_system_files.get() and ext in self.system_file_extensions:
            # Non saltare script files in ricerca avanzata/profonda
            script_files = ['.bat', '.cmd', '.ps1', '.vbs']
            if ext in script_files and search_level in ["avanzata", "profonda"]:
                return False
            
            self.log_debug(f"File di sistema escluso: {file_path}")
            self.log_skipped_file(file_path, skip_type, skip_filename, "File di sistema")
            return True
                    
        return False
    
    def log_skipped_file(self, filepath, skiptype, filename, skipreason):
        """Registra i file saltati in un file di log"""
        try:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_entry = f"{timestamp} - {skiptype} - {filename} - {filepath} - {skipreason}\n"
            
            with open(self.skipped_files_log_path, 'a', encoding='utf-8') as log_file:
                log_file.write(log_entry)
        except Exception as e:
            self.log_debug(f"Errore durante la scrittura del log dei file saltati: {str(e)}")

    def export_skipped_files_log(self):
        """Esporta il log dei file saltati in un formato CSV"""
        try:
               
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da esportare.")
                return
                
            # Verifica che il file non sia vuoto
            if os.path.getsize(self.skipped_files_log_path) == 0:
                messagebox.showinfo("Informazione", "Il file di log esiste ma è vuoto. Non ci sono file esclusi da esportare.")
                return

            # Chiedi all'utente dove salvare il file esportato
            export_path = filedialog.asksaveasfilename(
                defaultextension=".csv",
                initialfile="file_esclusi_export.csv",
                filetypes=[("CSV files", "*.csv"), ("Text files", "*.txt"), ("All files", "*.*")],
                title="Salva il log dei file esclusi"
            )
            
            if not export_path:  # L'utente ha annullato
                return
                
            # Leggi il file di log originale
            with open(self.skipped_files_log_path, 'r', encoding='utf-8') as log_file:
                log_content = log_file.readlines()
                
            # Crea il file CSV
            with open(export_path, 'w', newline='', encoding='utf-8') as csv_file:
                csv_writer = csv.writer(csv_file, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
                
                # Intestazione
                csv_writer.writerow(["Data e Ora", "Tipo", "Nome File", "Percorso Completo", "Motivo Esclusione"])
                
                # Analizza e scrive ogni riga del log
                for log_line in log_content:
                    try:
                        # Formato tipico: 2023-01-01 12:34:56 - File di sistema - nomefile.exe - C:/percorso/nomefile.exe - File di sistema
                        parts = log_line.strip().split(" - ", 4)
                        if len(parts) >= 5:
                            csv_writer.writerow(parts)
                    except Exception as e:
                        self.log_debug(f"Errore nell'elaborazione della riga di log: {str(e)}")
                        
                # Aggiunge statistiche alla fine
                csv_writer.writerow([])
                csv_writer.writerow([f"Totale file esclusi: {len(log_content)}"])
                csv_writer.writerow([f"Esportazione eseguita il: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"])
                csv_writer.writerow([f"Utente: {self.current_user}"])
            
            # Aggiungi un link per aprire il file esportato
            open_export = messagebox.askyesno(
                "Esportazione completata", 
                f"Esportazione completata con successo!\n\nFile salvato in:\n{export_path}\n\nVuoi aprire il file?"
            )
            
            if open_export:
                try:
                    if os.name == 'nt':  # Windows
                        os.startfile(export_path)
                    else:  # macOS o Linux
                        subprocess.call(['xdg-open', export_path])
                except Exception as e:
                    self.log_debug(f"Errore nell'apertura del file esportato: {str(e)}")
                    messagebox.showinfo("Informazione", f"Il file è stato salvato in:\n{export_path}")
                    
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'esportazione: {str(e)}")
            self.log_debug(f"Errore nell'esportazione del log: {str(e)}")

    def view_skipped_files_log(self):
        """Visualizza il log dei file saltati in una finestra separata"""
        try:
            if not hasattr(self, 'search_start_time') or self.search_start_time is None:
                messagebox.showinfo("Informazione", "Non è stata ancora effettuata una ricerca per visualizzare i file esclusi.")
                return
            # Verifica se il file di log esiste
            if not os.path.exists(self.skipped_files_log_path):
                messagebox.showinfo("Informazione", "Non ci sono file di log da visualizzare.")
                return
                
            # Verifica che il file non sia vuoto
            if os.path.getsize(self.skipped_files_log_path) == 0:
                messagebox.showinfo("Informazione", "Il file di log esiste ma è vuoto. Non ci sono file esclusi da visualizzare.")
                return
                
            # Crea una nuova finestra
            log_window = ttk.Toplevel(self.root)
            log_window.title("Log dei file esclusi")
            log_window.geometry("1000x800")
            log_window.transient(self.root)
            
            # Frame principale
            main_frame = ttk.Frame(log_window, padding=10)
            main_frame.pack(fill=BOTH, expand=YES)
            
            # Intestazione
            ttk.Label(main_frame, text="File esclusi durante la ricerca", 
                    font=("", 12, "bold")).pack(anchor=W, pady=(0, 10))
            
            # Area di testo per visualizzare i log
            text_frame = ttk.Frame(main_frame)
            text_frame.pack(fill=BOTH, expand=YES)
            
            scrollbar_y = ttk.Scrollbar(text_frame)
            scrollbar_y.pack(side=RIGHT, fill=Y)
            
            scrollbar_x = ttk.Scrollbar(text_frame, orient="horizontal")
            scrollbar_x.pack(side=BOTTOM, fill=X)
            
            log_text = tk.Text(text_frame, wrap="none", 
                            xscrollcommand=scrollbar_x.set,
                            yscrollcommand=scrollbar_y.set)
            log_text.pack(fill=BOTH, expand=YES)
            
            scrollbar_y.config(command=log_text.yview)
            scrollbar_x.config(command=log_text.xview)
            
            # Leggi e inserisci il contenuto del log
            with open(self.skipped_files_log_path, 'r', encoding='utf-8') as log_file:
                log_content = log_file.read()
                log_text.insert("1.0", log_content)
                
            # Rendi il testo di sola lettura
            log_text.config(state="disabled")
            
            # Pulsanti
            btn_frame = ttk.Frame(main_frame)
            btn_frame.pack(fill=X, pady=(10, 0))
            
            ttk.Button(btn_frame, text="Esporta CSV", 
                    command=self.export_skipped_files_log).pack(side=LEFT)
                    
            ttk.Button(btn_frame, text="Chiudi", 
                    command=log_window.destroy).pack(side=RIGHT)
                    
            # Centra la finestra
            log_window.update_idletasks()
            width = log_window.winfo_width()
            height = log_window.winfo_height()
            x = (log_window.winfo_screenwidth() // 2) - (width // 2)
            y = (log_window.winfo_screenheight() // 2) - (height // 2)
            log_window.geometry(f"{width}x{height}+{x}+{y}")
            
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante la visualizzazione del log: {str(e)}")
            self.log_debug(f"Errore nella visualizzazione del log: {str(e)}")

    def get_file_content(self, file_path):
        """Versione migliorata per caricare contenuti, inclusi file all'interno di archivi"""
        try:
            # Controlli preliminari (codice esistente)
            if self.should_skip_file(file_path):
                return ""
                    
            ext = os.path.splitext(file_path)[1].lower()
            search_level = self.search_depth.get()
            
            # Ottieni le estensioni personalizzate dell'utente per questo livello
            custom_extensions = self.get_extension_settings(search_level)
            
            # Gestione dimensione file (codice esistente)
            try:
                file_size = os.path.getsize(file_path)
                if file_size > self.max_file_size_mb.get() * 1024 * 1024:
                    self.log_debug(f"File {file_path} troppo grande per l'analisi del contenuto")
                    return ""
            except Exception as e:
                self.log_debug(f"Errore nel controllo dimensione del file {file_path}: {str(e)}")
                return ""
            
            # --- MIGLIORAMENTO: Gestione archivi per ricerca nei contenuti ---
            if (search_level == "profonda" or ext in custom_extensions) and ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                all_content = []
                found_data = False
                
                # Gestione archivi ZIP
                if ext == '.zip':
                    try:
                        import zipfile
                        import tempfile
                        
                        if not zipfile.is_zipfile(file_path):
                            return f"File non valido: {os.path.basename(file_path)}"
                        
                        # Prepara le informazioni di base
                        all_content.append(f"Archivio ZIP: {os.path.basename(file_path)}")
                        
                        # Crea una directory temporanea per l'estrazione
                        with tempfile.TemporaryDirectory() as temp_dir:
                            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                                # Includi solo i file che possiamo effettivamente analizzare
                                supported_files = []
                                for file_info in zip_ref.infolist():
                                    # Salta directory e file troppo grandi
                                    if file_info.filename.endswith('/') or file_info.file_size > 10 * 1024 * 1024:
                                        continue
                                    
                                    # Estrai estensione
                                    file_ext = os.path.splitext(file_info.filename)[1].lower()
                                    
                                    # Include il file se l'estensione è supportata o se è una ricerca profonda
                                    if search_level == "profonda" or file_ext in custom_extensions:
                                        supported_files.append(file_info)
                                    elif file_ext in self.get_default_extensions(search_level):
                                        supported_files.append(file_info)
                                
                                # Se ci sono troppi file, limita l'analisi
                                if len(supported_files) > 20:
                                    all_content.append(f"Archivio contiene {len(zip_ref.namelist())} file, analizzando i primi 20 per contenuto.")
                                    supported_files = supported_files[:20]
                                else:
                                    all_content.append(f"Analisi del contenuto di {len(supported_files)} file supportati all'interno dell'archivio.")
                                
                                # Estrai e analizza i file supportati
                                for file_info in supported_files:
                                    try:
                                        # Estrai solo questo file
                                        zip_ref.extract(file_info.filename, temp_dir)
                                        extracted_path = os.path.join(temp_dir, file_info.filename)
                                        
                                        # Verifica se il file esiste e non è una directory
                                        if os.path.exists(extracted_path) and os.path.isfile(extracted_path):
                                            # Leggi il contenuto del file estratto (usa funzioni esistenti)
                                            file_ext = os.path.splitext(file_info.filename)[1].lower()
                                            if file_ext in ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log']:
                                                # Leggi il file di testo
                                                try:
                                                    with open(extracted_path, 'r', encoding='utf-8', errors='ignore') as f:
                                                        content = f.read(5000)  # Limita a 5000 caratteri
                                                        if content.strip():
                                                            all_content.append(f"File: {file_info.filename}")
                                                            all_content.append(content[:1000] + "..." if len(content) > 1000 else content)
                                                            found_data = True
                                                except Exception as e:
                                                    self.log_debug(f"Impossibile leggere {file_info.filename}: {str(e)}")
                                            
                                            # Per file docx/pdf/office dentro l'archivio
                                            elif file_ext in ['.docx', '.pdf', '.xlsx', '.pptx'] and any(file_format_support[fmt] for fmt in ["docx", "pdf", "excel", "pptx"]):
                                                # Usa il metodo appropriato per il tipo di file
                                                doc_content = self._extract_specific_file_content(extracted_path)
                                                if doc_content:
                                                    all_content.append(f"File: {file_info.filename}")
                                                    all_content.append(doc_content[:1000] + "..." if len(doc_content) > 1000 else doc_content)
                                                    found_data = True
                                            
                                    except Exception as e:
                                        self.log_debug(f"Errore nell'estrazione o lettura di {file_info.filename}: {str(e)}")
                        
                        # Se non è stato trovato nessun contenuto significativo, includi almeno la lista dei file
                        if not found_data:
                            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                                all_content.append(f"Lista dei file nell'archivio:")
                                for name in zip_ref.namelist()[:30]:  # Limita a 30 nomi di file
                                    all_content.append(f"- {name}")
                        
                        return "\n".join(all_content)
                    except Exception as e:
                        self.log_debug(f"Errore nell'analisi dell'archivio ZIP {file_path}: {str(e)}")
                        return f"Errore nell'analisi dell'archivio: {os.path.basename(file_path)}"
                
                # Gestione archivi RAR con supporto contenuto
                elif ext == '.rar':
                    try:
                        import rarfile
                        import tempfile
                        
                        if not rarfile.is_rarfile(file_path):
                            return f"File RAR non valido: {os.path.basename(file_path)}"
                        
                        all_content.append(f"Archivio RAR: {os.path.basename(file_path)}")
                        
                        with tempfile.TemporaryDirectory() as temp_dir:
                            with rarfile.RarFile(file_path) as rf:
                                # Filtra i file supportati
                                supported_files = []
                                for info in rf.infolist():
                                    if info.isdir():
                                        continue
                                        
                                    # Controlla l'estensione
                                    file_ext = os.path.splitext(info.filename)[1].lower()
                                    if search_level == "profonda" or file_ext in custom_extensions:
                                        supported_files.append(info)
                                    elif file_ext in self.get_default_extensions(search_level):
                                        supported_files.append(info)
                                
                                # Limita il numero di file da estrarre
                                if len(supported_files) > 20:
                                    all_content.append(f"Archivio contiene {len(rf.infolist())} file, analizzando i primi 20 per contenuto.")
                                    supported_files = supported_files[:20]
                                
                                # Estrai e analizza i file supportati
                                for info in supported_files:
                                    try:
                                        rf.extract(info.filename, path=temp_dir)
                                        extracted_path = os.path.join(temp_dir, info.filename)
                                        
                                        if os.path.exists(extracted_path) and os.path.isfile(extracted_path):
                                            # Leggi il contenuto del file estratto
                                            file_ext = os.path.splitext(info.filename)[1].lower()
                                            if file_ext in ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log']:
                                                with open(extracted_path, 'r', encoding='utf-8', errors='ignore') as f:
                                                    content = f.read(5000)
                                                    if content.strip():
                                                        all_content.append(f"File: {info.filename}")
                                                        all_content.append(content[:1000] + "..." if len(content) > 1000 else content)
                                                        found_data = True
                                            
                                            # Per file docx/pdf/office
                                            elif file_ext in ['.docx', '.pdf', '.xlsx', '.pptx']:
                                                doc_content = self._extract_specific_file_content(extracted_path)
                                                if doc_content:
                                                    all_content.append(f"File: {info.filename}")
                                                    all_content.append(doc_content[:1000] + "..." if len(doc_content) > 1000 else doc_content)
                                                    found_data = True
                                    except Exception as e:
                                        self.log_debug(f"Errore nell'estrazione o lettura di {info.filename}: {str(e)}")
                                
                                # Se non è stato trovato nessun contenuto, includi la lista dei file
                                if not found_data:
                                    all_content.append(f"Lista dei file nell'archivio:")
                                    for name in [info.filename for info in rf.infolist()][:30]:
                                        all_content.append(f"- {name}")
                        
                        return "\n".join(all_content)
                    
                    except ImportError:
                        return f"Archivio RAR: {os.path.basename(file_path)}\nInstalla rarfile per vedere il contenuto"
                    except Exception as e:
                        self.log_debug(f"Errore nell'analisi dell'archivio RAR: {str(e)}")
                        return f"Errore nell'analisi dell'archivio RAR: {os.path.basename(file_path)}"
                
                # Altri tipi di archivi supportati (7z, tar, gz, ecc.)
                else:
                    # Per ora mostra solo le informazioni base
                    return f"Archivio {ext}: {os.path.basename(file_path)}\nL'analisi del contenuto non è ancora disponibile per questo formato"
            
            # --- Il resto del codice esistente per altri tipi di file --- 
            # Mantieni tutto il resto del metodo come è
            
            return ""  # Formato non supportato o livello di ricerca insufficiente
                
        except Exception as e:
            self.log_debug(f"Errore generale nella lettura del file {file_path}: {str(e)}")
            return ""

    def _extract_specific_file_content(self, file_path):
        """Metodo di supporto per estrarre il contenuto da file specifici all'interno di archivi"""
        ext = os.path.splitext(file_path)[1].lower()
        
        # Word (.docx)
        if ext == '.docx' and file_format_support["docx"]:
            try:
                import docx
                doc = docx.Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs])
            except Exception as e:
                self.log_debug(f"Errore Word: {str(e)}")
                return ""
        
        # PDF
        elif ext == '.pdf' and file_format_support["pdf"]:
            try:
                import PyPDF2
                content = []
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page_num in range(min(5, len(reader.pages))):  # Limita a 5 pagine per gli archivi
                        content.append(reader.pages[page_num].extract_text())
                return "\n".join(content)
            except Exception as e:
                self.log_debug(f"Errore PDF: {str(e)}")
                return ""
        
        # Excel
        elif ext in ['.xlsx', '.xls'] and file_format_support["excel"]:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                texts = []
                for sheet_name in wb.sheetnames[:2]:  # Limita a 2 fogli
                    sheet = wb[sheet_name]
                    for row in sheet.iter_rows(max_row=50):  # Limita a 50 righe
                        row_texts = [str(cell.value) for cell in row if cell.value is not None]
                        texts.append(" ".join(row_texts))
                return "\n".join(texts)
            except Exception as e:
                self.log_debug(f"Errore Excel: {str(e)}")
                return ""
        
        # PowerPoint
        elif ext in ['.pptx', '.ppt'] and file_format_support["pptx"]:
            try:
                import pptx
                prs = pptx.Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            except Exception as e:
                self.log_debug(f"Errore PowerPoint: {str(e)}")
                return ""
        
        # File di testo semplice
        elif ext in ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log']:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read(5000)  # Limita a 5000 caratteri
            except Exception as e:
                self.log_debug(f"Errore file testo: {str(e)}")
                return ""
        
        return ""  # Tipo di file non supportato
    
    def update_progress(self):
        if self.is_searching:
            try:
                # Processa tutti i messaggi nella coda
                messages_processed = 0
                max_messages_per_cycle = 50  # Limita il numero di messaggi processati per ciclo
                
                while messages_processed < max_messages_per_cycle:
                    try:
                        progress_type, value = self.progress_queue.get_nowait()
                        
                        # Processa il messaggio (codice esistente)
                        if progress_type == "update_total_time":
                            if hasattr(self, 'total_time_label') and self.total_time_label.winfo_exists():
                                self.total_time_label.config(text=value)
                        elif progress_type == "progress":
                            if hasattr(self, 'progress_bar') and self.progress_bar.winfo_exists():
                                self.progress_bar["value"] = value
                        elif progress_type == "status":
                            # Se il messaggio contiene informazioni sui file, separiamo le informazioni
                            if ("analizzati" in value.lower() or "cartelle:" in value.lower() or 
                                "file:" in value.lower()) and hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                # Estrai il percorso se presente
                                if "analisi:" in value.lower():
                                    parts = value.split("(", 1)
                                    if len(parts) > 1:
                                        path_part = parts[0].strip()
                                        counts_part = "(" + parts[1]
                                        if hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                            self.status_label["text"] = path_part
                                        if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                            self.analyzed_files_label["text"] = counts_part
                                    else:
                                        if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                            self.analyzed_files_label["text"] = value
                                else:
                                    # È solo un messaggio di stato semplice
                                    if hasattr(self, 'analyzed_files_label') and self.analyzed_files_label.winfo_exists():
                                        self.analyzed_files_label["text"] = value
                            elif hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                # È solo un messaggio di stato semplice
                                self.status_label["text"] = value
                        elif progress_type == "update_dir_size":
                            if hasattr(self, 'dir_size_var'):
                                self.dir_size_var.set(self._format_size(value))
                        elif progress_type == "complete":
                            self.is_searching = False
                            self.enable_all_controls()
                            if hasattr(self, 'stop_button') and self.stop_button.winfo_exists():
                                self.stop_button["state"] = "disabled"
                            
                            # Aggiorna la lista dei risultati
                            self.update_results_list()
                            
                            # Aggiorna l'orario di fine e il tempo totale
                            current_time = datetime.now().strftime('%H:%M')
                            if hasattr(self, 'end_time_label') and self.end_time_label.winfo_exists():
                                self.end_time_label.config(text=current_time)
                            self.update_total_time()  # Calcola e mostra il tempo totale
                            
                            # Calcola la dimensione del percorso alla fine della ricerca solo se non è disabilitato
                            calculation_mode = self.dir_size_calculation.get()
                            if calculation_mode == "disabilitato":
                                self.dir_size_var.set("Calcolo disattivato")
                            else:
                                self.dir_size_var.set("Calcolo in corso...")
                                path = self.search_path.get()
                                threading.Thread(target=self._calculate_dir_size_thread, args=(path,), daemon=True).start()

                            if len(self.search_results) == 0:
                                if hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                    self.status_label["text"] = "Nessun file trovato per la ricerca effettuata"
                                self.root.after(100, lambda: messagebox.showinfo("Ricerca completata", "Nessun file trovato per la ricerca effettuata"))
                            else:
                                if hasattr(self, 'status_label') and self.status_label.winfo_exists():
                                    self.status_label["text"] = f"Ricerca completata! Trovati {len(self.search_results)} risultati."
                            
                            if hasattr(self, 'progress_bar') and self.progress_bar.winfo_exists():
                                self.progress_bar["value"] = 100
                            return
                        elif progress_type == "error":
                            self.is_searching = False
                            self.enable_all_controls()
                            if hasattr(self, 'stop_button') and self.stop_button.winfo_exists():
                                self.stop_button["state"] = "disabled"
                            current_time = datetime.now().strftime('%H:%M')
                            if hasattr(self, 'end_time_label') and self.end_time_label.winfo_exists():
                                self.end_time_label.config(text=current_time)
                            self.update_total_time()  # Calcola e mostra il tempo totale
                            messagebox.showerror("Errore", value)
                            return
                        elif progress_type == "timeout":
                            self.is_searching = False
                            self.enable_all_controls()
                            if hasattr(self, 'stop_button') and self.stop_button.winfo_exists():
                                self.stop_button["state"] = "disabled"
                            self.update_results_list()
                            current_time = datetime.now().strftime('%H:%M')
                            if hasattr(self, 'end_time_label') and self.end_time_label.winfo_exists():
                                self.end_time_label.config(text=current_time)
                            self.update_total_time()  # Calcola e mostra il tempo totale
                            messagebox.showinfo("Timeout", "La ricerca è stata interrotta per timeout. Verranno mostrati i risultati parziali trovati.")
                            return
                        elif progress_type == "admin_prompt":
                            # Chiedi all'utente se desidera riavviare l'app come amministratore
                            response = messagebox.askyesno(
                                "Accesso limitato", 
                                "Alcune cartelle richiedono privilegi di amministratore per essere lette.\n\n" + 
                                "Vuoi riavviare l'applicazione come amministratore per ottenere accesso completo?",
                                icon="question"
                            )
                            if response:
                                self.stop_search_process()
                                self.root.after(1000, self.restart_as_admin)
                            return  # Non interrompere la ricerca se l'utente rifiuta
                        
                        messages_processed += 1
                        
                        # Forza l'aggiornamento dell'interfaccia ogni 10 messaggi
                        if messages_processed % 10 == 0:
                            self.root.update_idletasks()
                        
                    except queue.Empty:
                        break
                    except tk.TclError as e:
                        # Widget non più esistente o distrutto - non grave, ignora e continua
                        self.log_debug(f"Widget non più disponibile: {str(e)}")
                        continue
                    except Exception as e:
                        # Log altri errori ma non interrompere l'aggiornamento
                        self.log_debug(f"Errore durante l'elaborazione messaggio: {str(e)}")
                        continue
                        
                # Aggiorna l'UI forzatamente dopo aver processato i messaggi
                try:
                    self.root.update_idletasks()
                except:
                    pass
                
                # Richiama se stesso più frequentemente per essere più reattivo
                self.root.after(100, self.update_progress)  # Ridotto da 200ms a 100ms
            except tk.TclError as e:
                # Gestisci l'errore di widget non esistente o distrutto
                self.log_debug(f"TclError nell'aggiornamento del progresso: {str(e)}")
                # Riprova comunque ad aggiornare (potrebbe essere un errore temporaneo)
                self.root.after(500, self.update_progress)
            except Exception as e:
                self.log_debug(f"Errore nell'aggiornamento del progresso: {str(e)}")
                self.root.after(500, self.update_progress)
            
    def stop_search_process(self):
        """Ferma il processo di ricerca in corso"""
        self.stop_search = True
        self.status_label["text"] = "Interrompendo la ricerca..."
        self.analyzed_files_label["text"] = "Ricerca interrotta dall'utente"
        current_time = datetime.now().strftime('%H:%M')
        self.end_time_label.config(text=current_time)
        self.update_total_time()
        
        # Chiusura più decisa dell'executor
        if hasattr(self, 'search_executor') and self.search_executor:
            try:
                self.search_executor.shutdown(wait=False, cancel_futures=True)
                self.search_executor = None
            except Exception as e:
                self.log_debug(f"Errore nella chiusura dell'executor: {str(e)}")
                self.search_executor = None
        
        # Ritardo per evitare problemi con l'interfaccia
        self.root.after(100, self.enable_all_controls)
        self.stop_button["state"] = "disabled"
    
        # Aggiornamento forzato dell'interfaccia
        self.root.update_idletasks()
        
        # Aggiorna l'interfaccia utente
        self.root.update_idletasks()

    def update_results_list(self):
        """Aggiorna la lista dei risultati con i risultati trovati"""
        # Pulisci la lista attuale
        for item in self.results_list.get_children():
            self.results_list.delete(item)
            
        # Aggiungi i risultati alla lista
        for result in self.search_results:
            item_type, name, size, modified, created, path = result
            
            # Applica stile in base al tipo di elemento
            if item_type == "Directory":
                tags = ("directory",)
            else:
                tags = ("file",)
                
            self.results_list.insert("", "end", values=result, tags=tags)
            
        # Aggiorna lo stato
        self.status_label["text"] = f"Trovati {len(self.search_results)} risultati"
        
    def update_theme_colors(self, theme="light"):
        """Aggiorna i colori del tema per evidenziare cartelle e file"""
        style = ttk.Style()
        
        # Configura i colori per il tema chiaro
        if theme == "light":
            style.configure("Treeview", background="#ffffff", foreground="#000000", fieldbackground="#ffffff")
            self.results_list.tag_configure("directory", background="#e6f2ff", foreground="#000000")
            self.results_list.tag_configure("file", background="#ffffff", foreground="#000000")
        # Configura i colori per il tema scuro
        elif theme == "dark":
            style.configure("Treeview", background="#333333", foreground="#ffffff", fieldbackground="#333333")
            self.results_list.tag_configure("directory", background="#4d4d4d", foreground="#ffffff")
            self.results_list.tag_configure("file", background="#333333", foreground="#ffffff")

    def show_block_options(self):
        """Mostra la finestra di dialogo per le opzioni avanzate di ricerca a blocchi"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Opzioni Avanzate Ricerca a Blocchi")
        dialog.geometry("560x600")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Descrizione
        ttk.Label(main_frame, text="Configura le opzioni per ottimizzare la ricerca a blocchi", 
             font=("", 10, "bold")).pack(anchor=W, pady=(0, 10))
        
        description = ttk.Label(main_frame, text="La ricerca a blocchi divide la scansione in blocchi più piccoli "
                                        "per migliorare le prestazioni e la reattività dell'interfaccia.",
                         wraplength=520, justify=LEFT)
        description.pack(fill=X, pady=(0, 10))
        
        # Frame per opzioni dei blocchi
        block_frame = ttk.LabelFrame(main_frame, text="Dimensione e Parallelismo", padding=10)
        block_frame.pack(fill=X, pady=5)
        
        # Grid per organizzare le opzioni
        grid = ttk.Frame(block_frame)
        grid.pack(fill=X)
        
        # Opzione: Max file per blocco
        ttk.Label(grid, text="Max file per blocco:").grid(row=0, column=0, sticky=W, padx=5, pady=5)
        files_per_block = ttk.Spinbox(grid, from_=100, to=10000, increment=100, width=7,
                                 textvariable=self.max_files_per_block)
        files_per_block.grid(row=0, column=1, padx=5, pady=5)
        
        # Opzione: Max blocchi paralleli
        ttk.Label(grid, text="Blocchi paralleli:").grid(row=0, column=2, sticky=W, padx=(20, 5), pady=5)
        parallel_blocks = ttk.Spinbox(grid, from_=1, to=16, increment=1, width=5,
                                 textvariable=self.max_parallel_blocks)
        parallel_blocks.grid(row=0, column=3, padx=5, pady=5)
        
        # Checkbox per opzioni aggiuntive
        auto_adjust = ttk.Checkbutton(grid, text="Adatta automaticamente la dimensione dei blocchi",
                                variable=self.block_size_auto_adjust)
        auto_adjust.grid(row=1, column=0, columnspan=4, sticky=W, padx=5, pady=5)
        self.create_tooltip(auto_adjust, "Regola automaticamente la dimensione dei blocchi in base al tipo di ricerca e al sistema")
        
        # Frame per opzioni di prioritizzazione
        priority_frame = ttk.LabelFrame(main_frame, text="Priorità di Ricerca", padding=10)
        priority_frame.pack(fill=X, pady=10)
        
        # Checkbox per dare priorità alle cartelle utente
        user_priority = ttk.Checkbutton(priority_frame, text="Dare priorità alle cartelle utente (Users, Documents, Desktop, ecc)",
                                variable=self.prioritize_user_folders)
        user_priority.pack(anchor=W, padx=5, pady=5)
        self.create_tooltip(user_priority, "Elabora prima le cartelle più importanti per trovare risultati utili più velocemente")
        
        # Frame per spiegazione avanzata
        info_frame = ttk.LabelFrame(main_frame, text="Come funziona", padding=10)
        info_frame.pack(fill=X, pady=10)
        
        info_text = "La ricerca a blocchi divide le cartelle in unità di lavoro più piccole (blocchi) " \
                "che vengono elaborate in base alla priorità. Questo permette di:\n\n" \
                "• Ottenere risultati utili più rapidamente\n" \
                "• Migliorare la reattività dell'interfaccia durante la ricerca\n" \
                "• Distribuire meglio il carico di lavoro sui thread\n" \
                "• Gestire meglio la memoria per ricerche su grandi volumi di dati"
        
        info_label = ttk.Label(info_frame, text=info_text, wraplength=520, justify=LEFT)
        info_label.pack(fill=X, padx=5, pady=5)
        
        # Pulsanti per chiudere/applicare
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        # Pulsante per ripristinare valori predefiniti
        ttk.Button(btn_frame, text="Valori predefiniti", 
              command=lambda: [self.max_files_per_block.set(1000), 
                              self.max_parallel_blocks.set(4),
                              self.prioritize_user_folders.set(True),
                              self.block_size_auto_adjust.set(True)]).pack(side=LEFT, padx=5)
        
        # Pulsante di chiusura
        ttk.Button(btn_frame, text="Chiudi", command=dialog.destroy).pack(side=RIGHT, padx=5)
        
        # Centra la finestra
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")

    def show_performance_options(self):
        """Mostra la finestra di dialogo per le opzioni di performance"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Opzioni di Performance")
        dialog.geometry("560x700")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Descrizione
        ttk.Label(main_frame, text="Configura le opzioni di performance per ottimizzare la ricerca", 
            font=("", 10, "bold")).pack(anchor=W, pady=(0, 10))
        
        description = ttk.Label(main_frame, text="Queste opzioni consentono di controllare la velocità, l'uso di risorse e la precisione della ricerca.",
                        wraplength=520, justify=LEFT)
        description.pack(fill=X, pady=(0, 10))
        
        # Frame per opzioni di calcolo dimensione
        size_frame = ttk.LabelFrame(main_frame, text="Calcolo dimensioni", padding=10)
        size_frame.pack(fill=X, pady=5)
        
        size_grid = ttk.Frame(size_frame)
        size_grid.pack(fill=X)
        
        ttk.Label(size_grid, text="Modalità di calcolo:").grid(row=0, column=0, sticky=W, padx=5, pady=5)
        dir_size_combo = ttk.Combobox(size_grid, textvariable=self.dir_size_calculation, 
                                values=["incrementale", "preciso", "stimato", "sistema", "disabilitato"], 
                                width=12, state="readonly")
        dir_size_combo.grid(row=0, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(dir_size_combo, 
            "Scegli come calcolare la dimensione delle directory:\n" + 
            "- Incrementale: aggiorna durante la ricerca\n" +
            "- Preciso: calcolo completo ma più lento\n" +
            "- Stimato: più veloce ma approssimato\n" +
            "- Sistema: usa comandi di sistema esterni\n" +
            "- Disabilitato: non calcolare le dimensioni")
        
        # Frame per opzioni di timeout
        timeout_frame = ttk.LabelFrame(main_frame, text="Timeout e limiti", padding=10)
        timeout_frame.pack(fill=X, pady=5)
        
        timeout_grid = ttk.Frame(timeout_frame)
        timeout_grid.pack(fill=X)
        
        # Timeout
        timeout_check = ttk.Checkbutton(timeout_grid, text="Attiva timeout ricerca", variable=self.timeout_enabled)
        timeout_check.grid(row=0, column=0, sticky=W, padx=5, pady=5)
        
        ttk.Label(timeout_grid, text="Secondi:").grid(row=0, column=1, sticky=E, padx=5, pady=5)
        timeout_spinbox = ttk.Spinbox(timeout_grid, from_=10, to=3600, width=5, textvariable=self.timeout_seconds)
        timeout_spinbox.grid(row=0, column=2, padx=5, pady=5, sticky=W)
        self.create_tooltip(timeout_spinbox, 
            "Interrompe automaticamente la ricerca dopo il tempo specificato\n\n"
            "Normale funzionamento:\n"
            "• La ricerca si interrompe dopo il tempo indicato\n"
            "• Vengono mostrati i risultati parziali trovati\n\n"
            "Comportamento per ricerche di sistema (C:/, D:/):\n"
            "• Per ricerche complete su disco, il timeout viene\n"
            "  temporaneamente aumentato a 8 ore per permettere\n"
            "  scansioni approfondite")
        
        # Max file da controllare
        ttk.Label(timeout_grid, text="Max file da controllare:").grid(row=1, column=0, sticky=W, padx=5, pady=5)
        max_files_spinbox = ttk.Spinbox(timeout_grid, from_=1000, to=10000000, width=8, textvariable=self.max_files_to_check)
        max_files_spinbox.grid(row=1, column=1, columnspan=2, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_files_spinbox, "Numero massimo di file da analizzare prima di terminare la ricerca")
        
        # Max risultati
        ttk.Label(timeout_grid, text="Max risultati:").grid(row=2, column=0, sticky=W, padx=5, pady=5)
        max_results_spinbox = ttk.Spinbox(timeout_grid, from_=500, to=100000, width=8, textvariable=self.max_results)
        max_results_spinbox.grid(row=2, column=1, columnspan=2, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_results_spinbox, "Numero massimo di risultati da mostrare")
        
        # Frame per opzioni di processamento
        process_frame = ttk.LabelFrame(main_frame, text="Processamento", padding=10)
        process_frame.pack(fill=X, pady=5)
        
        process_grid = ttk.Frame(process_frame)
        process_grid.pack(fill=X)
        
        # Thread
        ttk.Label(process_grid, text="Thread paralleli:").grid(row=0, column=0, sticky=W, padx=5, pady=5)
        threads_spinbox = ttk.Spinbox(process_grid, from_=1, to=16, width=3, textvariable=self.worker_threads)
        threads_spinbox.grid(row=0, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(threads_spinbox, "Numero di thread paralleli per la ricerca (aumentando migliora la velocità ma usa più risorse)")
        
        # Dimensione massima file
        ttk.Label(process_grid, text="Dimensione max file (MB):").grid(row=1, column=0, sticky=W, padx=5, pady=5)
        max_size_spinbox = ttk.Spinbox(process_grid, from_=1, to=1000, width=5, textvariable=self.max_file_size_mb)
        max_size_spinbox.grid(row=1, column=1, padx=5, pady=5, sticky=W)
        self.create_tooltip(max_size_spinbox, "Dimensione massima in MB dei file di cui analizzare il contenuto")
        
        # Checkbox opzioni aggiuntive
        options_frame = ttk.LabelFrame(main_frame, text="Opzioni aggiuntive", padding=10)
        options_frame.pack(fill=X, pady=5)
        
        # Checkbox indicizzazione
        index_check = ttk.Checkbutton(options_frame, text="Usa indicizzazione (velocizza ricerche ripetute)", variable=self.use_indexing)
        index_check.pack(anchor=W, pady=2)
        self.create_tooltip(index_check, "Utilizza l'indicizzazione per velocizzare ricerche future")
        
        # Checkbox errori di permesso
        perm_check = ttk.Checkbutton(options_frame, text="Ignora errori di permesso", variable=self.skip_permission_errors)
        perm_check.pack(anchor=W, pady=2)
        self.create_tooltip(perm_check, "Continua la ricerca anche quando alcune cartelle non possono essere lette")
        
        # Pulsante per ripristinare valori predefiniti
        def restore_defaults():
            self.dir_size_calculation.set("disabilitato")
            self.timeout_enabled.set(False)
            self.timeout_seconds.set(3600)
            self.max_files_to_check.set(100000)
            self.max_results.set(50000)
            self.worker_threads.set(min(8, os.cpu_count() or 4))
            self.max_file_size_mb.set(100)
            self.use_indexing.set(True)
            self.skip_permission_errors.set(True)
        
        # Pulsanti per chiudere/applicare
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        # Pulsante per ripristinare valori predefiniti
        ttk.Button(btn_frame, text="Valori predefiniti", 
            command=restore_defaults).pack(side=LEFT, padx=5)
        
        # Pulsante di chiusura
        ttk.Button(btn_frame, text="Chiudi", command=dialog.destroy).pack(side=RIGHT, padx=5)
        
        # Centra la finestra
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")

    def copy_selected(self):
        """Copia i file selezionati in una directory di destinazione"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da copiare")
            return
            
        # Chiedi la directory di destinazione
        dest_dir = filedialog.askdirectory(title="Seleziona la cartella di destinazione")
        if not dest_dir:
            return
            
        # Prepara variabili per tracciare l'avanzamento
        total = len(selected_items)
        copied = 0
        failed = 0
        skipped = 0
        
        try:
            for item in selected_items:
                values = self.results_list.item(item)['values']
                item_type, _, _, _, _, source_path = values
                
                # Ottieni il nome dell'elemento senza il percorso completo
                basename = os.path.basename(source_path)
                dest_path = os.path.join(dest_dir, basename)
                
                try:
                    if item_type == "Directory":
                        # Se la cartella esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("Cartella esistente", 
                                                f"La cartella {basename} esiste già. Vuoi sovrascriverla?"):
                                # Elimina la cartella esistente
                                shutil.rmtree(dest_path)
                            else:
                                skipped += 1
                                continue
                                
                        # Copia ricorsiva della cartella
                        shutil.copytree(source_path, dest_path)
                        copied += 1
                    else:
                        # Se il file esiste già nella destinazione
                        if os.path.exists(dest_path):
                            if messagebox.askyesno("File esistente", 
                                                f"Il file {basename} esiste già. Vuoi sovrascriverlo?"):
                                # Continua con la sovrascrittura
                                pass
                            else:
                                skipped += 1
                                continue
                                
                        # Copia il file
                        shutil.copy2(source_path, dest_path)
                        copied += 1
                        
                except Exception as e:
                    failed += 1
                    self.log_debug(f"Errore durante la copia di {source_path}: {str(e)}")
                    
                # Aggiorna la barra di avanzamento
                progress = ((copied + failed + skipped) / total) * 100
                self.progress_bar["value"] = progress
                self.status_label["text"] = f"Copiati {copied}/{total} elementi..."
                self.root.update()
                
            # Messaggio di completamento
            if failed > 0:
                messagebox.showwarning("Copia completata con errori", 
                                    f"Copiati: {copied}\nSaltati: {skipped}\nFalliti: {failed}")
            else:
                messagebox.showinfo("Copia completata", 
                                 f"Copiati con successo: {copied}\nSaltati: {skipped}")
                
        except Exception as e:
            messagebox.showerror("Errore", f"Si è verificato un errore durante l'operazione: {str(e)}")
            
        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    def get_zip_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome del file ZIP"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Nome file ZIP")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome del file ZIP (senza estensione):", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = StringVar(value="archivio")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            zip_name = name_var.get().strip()
            if zip_name:
                result["name"] = zip_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per il file ZIP", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea ZIP", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]

    def compress_selected(self):
        """Versione avanzata della compressione che mantiene la struttura originale delle directory"""
        selected_items = self.results_list.selection()
        if not selected_items:
            messagebox.showwarning("Attenzione", "Seleziona almeno un elemento da comprimere")
            return

        # Ottieni il nome del file ZIP
        zip_name = self.get_zip_name()
        if not zip_name:
            return

        # Ottieni il nome della cartella principale
        main_folder_name = self.get_main_folder_name()
        if main_folder_name is None:
            return  # L'utente ha annullato
        if not main_folder_name:
            main_folder_name = "files"  # Default se l'utente inserisce valore vuoto

        # Opzioni di compressione più dettagliate
        compression_dialog = tk.Toplevel(self.root)
        compression_dialog.title("Opzioni Compressione")
        compression_dialog.transient(self.root)
        compression_dialog.grab_set()

        frame = ttk.Frame(compression_dialog, padding=20)
        frame.pack(fill=tk.BOTH, expand=True)

        ttk.Label(frame, text="Scegli il tipo di compressione:").pack(anchor=tk.W, pady=(0,10))

        comp_var = tk.StringVar(value="standard")
        preserve_var = tk.BooleanVar(value=True)  # Nuova opzione per preservare la struttura

        ttk.Radiobutton(frame, text="Nessuna (solo archiviazione)", 
            variable=comp_var, value="nessuna").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Minima (massima velocità) Liv.1", 
                    variable=comp_var, value="minima").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Standard (buon equilibrio) Liv.6", 
                    variable=comp_var, value="standard").pack(anchor=tk.W)
        ttk.Radiobutton(frame, text="Massima (compressione ottimale, più lenta) Liv.9", 
                    variable=comp_var, value="massima").pack(anchor=tk.W)

        ttk.Label(frame, text="\nOpzioni struttura:", font=("", 9, "bold")).pack(anchor=tk.W, pady=(10,5))
        ttk.Checkbutton(frame, text="Preserva la struttura delle directory", 
                    variable=preserve_var).pack(anchor=tk.W)
        ttk.Label(frame, text="Se attivato, mantiene l'alberatura originale delle cartelle", 
                    font=("", 8), foreground="gray").pack(anchor=tk.W, padx=(20, 5))

        ttk.Label(frame, text="\nOpzioni per file di grandi dimensioni:").pack(anchor=tk.W, pady=(10,5))

        use_chunks = tk.BooleanVar(value=False)
        ttk.Checkbutton(frame, text="Elabora in blocchi (per file molto grandi)", 
                    variable=use_chunks).pack(anchor=tk.W)

        btn_frame = ttk.Frame(frame)
        btn_frame.pack(fill=tk.X, pady=(15,0))

        result = {"option": None, "chunks": False, "preserve": True}

        def on_cancel():
            result["option"] = None
            compression_dialog.destroy()

        def on_ok():
            result["option"] = comp_var.get()
            result["chunks"] = use_chunks.get()
            result["preserve"] = preserve_var.get()
            compression_dialog.destroy()

        ttk.Button(btn_frame, text="Annulla", command=on_cancel).pack(side=tk.LEFT)
        ttk.Button(btn_frame, text="OK", command=on_ok).pack(side=tk.RIGHT)

        # Centra la finestra
        compression_dialog.update_idletasks()
        width = compression_dialog.winfo_reqwidth() + 50
        height = compression_dialog.winfo_reqheight() + 20
        x = (compression_dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (compression_dialog.winfo_screenheight() // 2) - (height // 2)
        compression_dialog.geometry(f"{width}x{height}+{x}+{y}")

        compression_dialog.wait_window()

        if result["option"] is None:
            return  # Utente ha annullato

        # Determina il metodo di compressione
        compression_level = 0  # Valore predefinito

        if result["option"] == "nessuna":
            compression_method = zipfile.ZIP_STORED
            compression_text = "nessuna"
        elif result["option"] == "minima":
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "minima"
            compression_level = 1  # Compressione veloce, riduzione minima
        elif result["option"] == "massima":
            # Prova LZMA se disponibile, altrimenti BZIP2, altrimenti fallback su DEFLATED con livello massimo
            if hasattr(zipfile, 'ZIP_LZMA'):
                compression_method = zipfile.ZIP_LZMA
                compression_text = "massima (LZMA)"
                compression_level = 9
            elif hasattr(zipfile, 'ZIP_BZIP2'):
                compression_method = zipfile.ZIP_BZIP2
                compression_text = "massima (BZIP2)"
                compression_level = 9
            else:
                compression_method = zipfile.ZIP_DEFLATED
                compression_text = "massima (DEFLATED)"
                compression_level = 9
        else:  # standard
            compression_method = zipfile.ZIP_DEFLATED
            compression_text = "standard"
            compression_level = 6  # Compressione bilanciata
            
        # Log della scelta di compressione
        self.log_debug(f"Utilizzo compressione {compression_text} (metodo: {compression_method}, livello: {compression_level})")
        self.log_debug(f"Preserva struttura directory: {result['preserve']}")

        zip_path = filedialog.asksaveasfilename(
            defaultextension=".zip",
            initialfile=f"{zip_name}.zip",
            filetypes=[("ZIP files", "*.zip")],
            title="Salva file ZIP"
        )

        if not zip_path:
            return

        # Determinare il percorso base per i calcoli relativi
        base_path = self._find_common_base_path(selected_items)
        self.log_debug(f"Percorso base per la struttura: {base_path}")

        # Raccogli tutti i file delle cartelle selezionate
        files_in_folders = set()
        folder_paths = []
        single_files = []

        # Prima fase: raccogli informazioni su cartelle e file
        for item in selected_items:
            values = self.results_list.item(item)['values']
            item_type, _, _, _, _, source_path = values

            if item_type == "Directory":
                folder_paths.append(source_path)
                # Raccogli tutti i file nelle cartelle selezionate
                for root, _, files in os.walk(source_path):
                    for file in files:
                        full_path = os.path.join(root, file)
                        files_in_folders.add(os.path.abspath(full_path))
            else:
                single_files.append(source_path)

        # Filtra i file singoli che sono già presenti nelle cartelle
        filtered_single_files = [f for f in single_files if os.path.abspath(f) not in files_in_folders]

        total_items = len(folder_paths) + len(filtered_single_files)
        processed = 0

        # Lista per tenere traccia dei file (tutti, non solo omonimi)
        all_files_log = []
        # Lista per tenere traccia dei file omonimi (come prima)
        omonimi_log = []

        # Ottieni informazioni sulla ricerca corrente per i log
        search_directory = self.search_path.get() if hasattr(self, 'search_path') else "N/A"
        search_keywords = self.keywords.get() if hasattr(self, 'keywords') else "N/A"

        try:
            # Primo passaggio: analizza tutti i file e identifica gli omonimi
            self.log_debug(f"Analisi dei file da comprimere nella cartella {main_folder_name}...")

            # Dizionario {nome_file: [lista_percorsi]}
            file_names_map = {}

            # Raccogli i nomi dei file dalle cartelle
            for folder_path in folder_paths:
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        file_name = os.path.basename(file_path)

                        # Aggiungi al log di tutti i file
                        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                        modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                        
                        all_files_log.append({
                            "nome_file": file_name,
                            "percorso_originale": file_path,
                            "dimensione": self._format_size(file_size),
                            "ultima_modifica": modified_time,
                            "tipo": self._get_file_type(file_path)
                        })

                        if file_name not in file_names_map:
                            file_names_map[file_name] = []
                        file_names_map[file_name].append(file_path)

            # Raccogli i nomi dei file singoli
            for file_path in filtered_single_files:
                if os.path.exists(file_path):
                    file_name = os.path.basename(file_path)
                    
                    # Aggiungi al log di tutti i file
                    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                    modified_time = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%d/%m/%Y %H:%M') if os.path.exists(file_path) else 'N/A'
                    
                    all_files_log.append({
                        "nome_file": file_name,
                        "percorso_originale": file_path,
                        "dimensione": self._format_size(file_size),
                        "ultima_modifica": modified_time,
                        "tipo": self._get_file_type(file_path)
                    })
                    
                    if file_name not in file_names_map:
                        file_names_map[file_name] = []
                    file_names_map[file_name].append(file_path)

            # Identifica i file omonimi (con lo stesso nome ma percorsi diversi)
            omonimi_files = {name: paths for name, paths in file_names_map.items() if len(paths) > 1}
            self.log_debug(f"Trovati {len(omonimi_files)} file con nomi duplicati")
            
            # Per l'elaborazione a blocchi
            if result["chunks"] and total_items > 50:
                chunk_size = 10  # Elabora 10 file alla volta
                self.log_debug(f"Elaborazione a blocchi attivata: {chunk_size} file per blocco")

            # Secondo passaggio: crea l'archivio ZIP
            with zipfile.ZipFile(zip_path, 'w', compression_method, compresslevel=compression_level) as zipf:
                # Tiene traccia dei percorsi già aggiunti al file ZIP
                added_zip_paths = set()

                # Comprimi le cartelle
                for folder_path in folder_paths:
                    # Utilizzo la struttura originale o piatta in base alla scelta dell'utente
                    for root, _, files in os.walk(folder_path):
                        for file in files:
                            file_path = os.path.join(root, file)
                            file_name = os.path.basename(file_path)

                            try:
                                # Salta file non leggibili
                                if not os.access(file_path, os.R_OK):
                                    self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                    continue

                                # Calcola percorso relativo in base all'opzione selezionata
                                if result["preserve"]:
                                    # Calcola il percorso relativo rispetto alla base
                                    rel_path = os.path.relpath(file_path, base_path)
                                    zip_path_to_use = os.path.join(main_folder_name, rel_path)
                                else:
                                    # Comportamento originale senza preservare la struttura
                                    rel_path = os.path.relpath(file_path, os.path.dirname(folder_path))
                                    
                                    # Verifica se è un file omonimo
                                    if file_name in omonimi_files:
                                        # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                        if omonimi_files[file_name][0] == file_path:
                                            # Usa il percorso all'interno della cartella principale
                                            zip_path_to_use = os.path.join(main_folder_name, rel_path)
                                        else:
                                            # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                            parent_folder = os.path.basename(os.path.dirname(file_path))
                                            unique_name = f"{parent_folder}_{file_name}"

                                            # Se anche questo nome è duplicato, aggiungi un contatore
                                            counter = 1
                                            while os.path.join("omonimi", unique_name) in added_zip_paths:
                                                unique_name = f"{parent_folder}_{counter}_{file_name}"
                                                counter += 1

                                            zip_path_to_use = os.path.join("omonimi", unique_name)

                                            # Registra questo file nel log degli omonimi
                                            omonimi_log.append({
                                                "nome_file": file_name,
                                                "percorso_originale": file_path,
                                                "primo_percorso": omonimi_files[file_name][0],
                                                "posizione_zip": zip_path_to_use
                                            })
                                    else:
                                        # Non è un omonimo, inseriscilo nella cartella principale
                                        zip_path_to_use = os.path.join(main_folder_name, rel_path)

                                # Verifica se questo percorso ZIP è già stato usato (potrebbe accadere anche con la struttura preservata)
                                if zip_path_to_use in added_zip_paths:
                                    self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                    # Crea un nome alternativo
                                    alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)

                                    # Se anche questo nome è già usato, aggiungi un contatore
                                    counter = 1
                                    while zip_path_to_use in added_zip_paths:
                                        alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                        zip_path_to_use = os.path.join("omonimi", alt_name)
                                        counter += 1

                                # Aggiungi il file al ZIP e registra il percorso
                                zipf.write(file_path, zip_path_to_use)
                                added_zip_paths.add(zip_path_to_use)

                            except Exception as e:
                                self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Comprimi i file singoli
                for file_path in filtered_single_files:
                    if os.path.exists(file_path):
                        file_name = os.path.basename(file_path)

                        try:
                            # Salta file non leggibili
                            if not os.access(file_path, os.R_OK):
                                self.log_debug(f"Saltato file senza permessi di lettura: {file_path}")
                                continue

                            # Calcola percorso relativo in base all'opzione selezionata
                            if result["preserve"]:
                                # Calcola il percorso relativo rispetto alla base
                                rel_path = os.path.relpath(file_path, base_path)
                                zip_path_to_use = os.path.join(main_folder_name, rel_path)
                            else:
                                # Comportamento originale senza preservare la struttura
                                # Verifica se è un file omonimo
                                if file_name in omonimi_files:
                                    # Se questo è il primo file con questo nome, inseriscilo nella cartella principale
                                    if omonimi_files[file_name][0] == file_path:
                                        # Usa il nome del file all'interno della cartella principale
                                        zip_path_to_use = os.path.join(main_folder_name, file_name)
                                    else:
                                        # Per i file omonimi successivi, crea un nome univoco nella cartella 'omonimi'
                                        parent_folder = os.path.basename(os.path.dirname(file_path))
                                        unique_name = f"{parent_folder}_{file_name}"

                                        # Se anche questo nome è duplicato, aggiungi un contatore
                                        counter = 1
                                        while os.path.join("omonimi", unique_name) in added_zip_paths:
                                            unique_name = f"{parent_folder}_{counter}_{file_name}"
                                            counter += 1

                                        zip_path_to_use = os.path.join("omonimi", unique_name)

                                        # Registra questo file nel log degli omonimi
                                        omonimi_log.append({
                                            "nome_file": file_name,
                                            "percorso_originale": file_path,
                                            "primo_percorso": omonimi_files[file_name][0],
                                            "posizione_zip": zip_path_to_use
                                        })
                                else:
                                    # Non è un omonimo, inseriscilo nella cartella principale
                                    zip_path_to_use = os.path.join(main_folder_name, file_name)

                            # Verifica se questo percorso ZIP è già stato usato
                            if zip_path_to_use in added_zip_paths:
                                self.log_debug(f"Conflitto di percorso ZIP: {zip_path_to_use}")
                                # Crea un nome alternativo
                                alt_name = f"conflitto_{os.path.basename(zip_path_to_use)}"
                                zip_path_to_use = os.path.join("omonimi", alt_name)

                                # Se anche questo nome è già usato, aggiungi un contatore
                                counter = 1
                                while zip_path_to_use in added_zip_paths:
                                    alt_name = f"conflitto_{counter}_{os.path.basename(zip_path_to_use)}"
                                    zip_path_to_use = os.path.join("omonimi", alt_name)
                                    counter += 1

                            # Aggiungi il file al ZIP e registra il percorso
                            zipf.write(file_path, zip_path_to_use)
                            added_zip_paths.add(zip_path_to_use)

                        except Exception as e:
                            self.log_debug(f"Errore durante la compressione di {file_path}: {str(e)}")

                    processed += 1
                    progress = (processed / total_items) * 100
                    self.progress_bar["value"] = progress
                    self.status_label["text"] = f"Compressi {processed} di {total_items} elementi"
                    self.root.update()

                # Crea log in formato CSV con collegamenti ipertestuali
                current_time_local = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                current_user = getpass.getuser()

                # Identifica quali file sono omonimi per poterli escludere dal log generale
                omonimo_paths = set()
                for entry in omonimi_log:
                    omonimo_paths.add(entry['percorso_originale'])

                # Lista di file normali (esclude gli omonimi)
                normal_files_log = [f for f in all_files_log if f['percorso_originale'] not in omonimo_paths]

                # ---------- CREAZIONE DEL CSV DEI FILE NORMALI ----------
                # Creazione intestazione per il CSV con informazioni sulla ricerca
                csv_header = f"LOG DEI FILE TROVATI - {current_time_local}\n"
                csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                csv_header += f"Directory di ricerca: {search_directory}\n"
                csv_header += f"Parole chiave ricercate: {search_keywords}\n"
                csv_header += f"Opzione struttura: {'Preservata' if result['preserve'] else 'Piatta'}\n"
                csv_header += f"Totale file: {len(all_files_log)}\n\n"

                # Prepara l'output CSV
                csv_content = io.StringIO()
                csv_writer = csv.writer(csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)

                # Intestazione colonne
                csv_writer.writerow(["Nome file", "Tipo", "Dimensione", "Ultima modifica", "Percorso completo", "Link"])

                # Aggiungi i file normali (non omonimi)
                for file_entry in sorted(all_files_log, key=lambda x: x["nome_file"]):
                    # Ottieni il percorso della directory che contiene il file
                    file_path = file_entry['percorso_originale']
                    directory_path = os.path.dirname(file_path).replace('\\', '/')
                    # Formatta il percorso della directory per il collegamento ipertestuale
                    folder_uri = f"file:///{directory_path}"
                    
                    # Aggiungi riga al CSV con la formula italiana per i collegamenti alla CARTELLA
                    csv_writer.writerow([
                        file_entry['nome_file'],
                        file_entry['tipo'],
                        file_entry['dimensione'],
                        file_entry['ultima_modifica'],
                        file_entry['percorso_originale'],
                        f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")'
                    ])

                # Ottieni il contenuto del CSV come stringa
                csv_data = csv_header + csv_content.getvalue()

                # Aggiungi il file CSV dei file normali all'archivio
                zipf.writestr("Log_file_Trovati.csv", csv_data)

                # ---------- CREAZIONE DEL CSV DEGLI OMONIMI (SOLO SE CE NE SONO) ----------
                if omonimi_log:
                    # Creazione intestazione per il CSV degli omonimi con info di ricerca
                    omonimi_csv_header = f"LOG DEI FILE OMONIMI - {current_time_local}\n"
                    omonimi_csv_header += f"Utente: {current_user}, Compressione: {compression_text}\n"
                    omonimi_csv_header += f"Directory di ricerca: {search_directory}\n"
                    omonimi_csv_header += f"Parole chiave ricercate: {search_keywords}\n"
                    omonimi_csv_header += f"Totale file omonimi trovati: {len(omonimi_log)}\n\n"
                    
                    # Prepara l'output CSV per gli omonimi
                    omonimi_csv_content = io.StringIO()
                    omonimi_csv_writer = csv.writer(omonimi_csv_content, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
                    
                    # Intestazione colonne specifiche per gli omonimi
                    omonimi_csv_writer.writerow([
                        "Nome file", "Tipo", "Dimensione", "Ultima modifica", 
                        "Percorso originale", "Link", "Posizione nello ZIP", "In conflitto con"
                    ])
                    
                    # Estrai dai log completi solo i file omonimi
                    omonimi_files_details = [f for f in all_files_log if f['percorso_originale'] in omonimo_paths]
                    
                    # Crea un dizionario per unire le informazioni di omonimi_log e omonimi_files_details
                    omonimi_details_dict = {}
                    for entry in omonimi_log:
                        omonimi_details_dict[entry['percorso_originale']] = entry
                    
                    # Aggiungi i file omonimi al CSV
                    for file_entry in sorted(omonimi_files_details, key=lambda x: x["nome_file"]):
                        # Ottieni il percorso della directory che contiene il file
                        file_path = file_entry['percorso_originale']
                        directory_path = os.path.dirname(file_path).replace('\\', '/')
                        # Formatta il percorso della directory per il collegamento ipertestuale
                        folder_uri = f"file:///{directory_path}"
                        
                        # Trova le informazioni aggiuntive sugli omonimi
                        omonimo_info = omonimi_details_dict.get(file_entry['percorso_originale'], {})
                        posizione_zip = omonimo_info.get('posizione_zip', 'N/A')
                        in_conflitto = omonimo_info.get('primo_percorso', 'N/A')
                        
                        # Aggiungi riga al CSV degli omonimi con la formula italiana per i collegamenti alla CARTELLA
                        omonimi_csv_writer.writerow([
                            file_entry['nome_file'],
                            file_entry['tipo'],
                            file_entry['dimensione'],
                            file_entry['ultima_modifica'],
                            file_entry['percorso_originale'],
                            f'=COLLEG.IPERTESTUALE("{folder_uri}";"Apri percorso")',
                            posizione_zip,
                            in_conflitto
                        ])
                    
                    # Ottieni il contenuto del CSV come stringa
                    omonimi_csv_data = omonimi_csv_header + omonimi_csv_content.getvalue()
                    
                    # Aggiungi il file CSV degli omonimi all'archivio
                    zipf.writestr("omonimi_log.csv", omonimi_csv_data)

            # Prepara il messaggio di completamento
            skipped_files = len(single_files) - len(filtered_single_files)
            message = f"Compressione completata!\nFile salvato in: {zip_path}\n"
            message += f"File organizzati nella cartella '{main_folder_name}'\n"
            message += f"Tipo di compressione utilizzata: {compression_text}\n"
            
            if result["preserve"]:
                message += f"Struttura delle directory originali: Preservata\n"
            else:
                message += f"Struttura delle directory originali: Non preservata\n"
                
            message += f"\nCreato file di log completo ('Log_file_Trovati.csv')"

            if result["chunks"]:
                message += "\nElaborazione a blocchi: Attiva"

            if omonimi_log:
                message += f"\nTrovati {len(omonimi_log)} file omonimi."
                message += f"\nCreato log dettagliato degli omonimi ('omonimi_log.csv')"

            if skipped_files > 0:
                message += f"\n{skipped_files} file saltati perché già presenti nelle cartelle"

            messagebox.showinfo("Completato", message)

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            self.log_debug(f"Errore dettagliato durante la compressione:\n{error_details}")

            messagebox.showerror("Errore", 
                f"Errore durante la compressione: {str(e)}\n\n"
                f"Tipo di errore: {type(e).__name__}")

        finally:
            self.progress_bar["value"] = 0
            self.status_label["text"] = "In attesa..."

    def _find_common_base_path(self, selected_items):
        """Trova il percorso base comune a tutti i file selezionati"""
        if not selected_items:
            return ""
        
        paths = []
        
        # Raccogli tutti i percorsi
        for item in selected_items:
            values = self.results_list.item(item)['values']
            paths.append(values[5])  # Il percorso è nella sesta colonna
        
        # Trova il percorso comune più lungo
        def common_path(paths):
            if not paths:
                return ""
            
            # Normalizza i percorsi per gestire i diversi separatori di directory
            paths = [os.path.normpath(p) for p in paths]
            
            # Dividi ogni percorso in componenti
            components = [p.split(os.path.sep) for p in paths]
            
            # Se ci sono percorsi di unità diverse (es. C:\ e D:\) su Windows
            if os.name == 'nt' and len(set(c[0] for c in components if c)) > 1:
                return ""  # Non c'è un percorso comune tra unità diverse
            
            common = []
            for i in range(min(len(c) for c in components)):
                if len(set(c[i] for c in components)) == 1:
                    common.append(components[0][i])
                else:
                    break
            
            # Se non c'è un elemento comune, restituisci la radice/unità
            if not common and os.name == 'nt' and components and components[0]:
                return components[0][0] + os.path.sep  # es. "C:\"
            elif not common:
                return os.path.sep  # radice Unix "/"
                
            return os.path.sep.join(common)
        
        base_path = common_path(paths)
        
        # Se il percorso base finisce con un separatore, va bene
        # altrimenti dobbiamo aggiungere il separatore se base_path non è vuoto
        if base_path and not base_path.endswith(os.path.sep):
            base_path = os.path.dirname(base_path)
        
        # Se non abbiamo trovato un percorso comune, usa il percorso di ricerca corrente
        if not base_path and hasattr(self, 'search_path'):
            base_path = self.search_path.get()
        
        return base_path
            
    def get_directory_size(self, path):
        """Calculate the total size of a directory"""
        if not os.path.exists(path):
            return 0
            
        try:
            total_size = 0
            if os.path.isfile(path):
                return os.path.getsize(path)
                
            # Per evitare il blocco dell'interfaccia, imposta un timeout massimo
            start_time = time.time()
            max_time = 30  # massimo 30 secondi
            files_count = 0
            error_count = 0
                
            # For directories, walk through all files and subdirectories
            for dirpath, dirnames, filenames in os.walk(path):
                for f in filenames:
                    # Verifica se il timeout è scaduto
                    if time.time() - start_time > max_time:
                        self.log_debug(f"Timeout nel calcolo della dimensione per {path}")
                        return total_size
                        
                    try:
                        fp = os.path.join(dirpath, f)
                        
                        # Verifica esplicita che il file esiste ancora prima di tentare di leggerne la dimensione
                        if os.path.exists(fp) and not os.path.islink(fp):
                            total_size += os.path.getsize(fp)
                            files_count += 1
                    except FileNotFoundError:
                        # Ignora silenziosamente i file che non esistono più
                        pass
                    except (OSError, PermissionError) as e:
                        error_count += 1
                        # Limita il numero di errori da registrare per evitare spam nel log
                        if error_count < 100:  
                            self.log_debug(f"Error getting size of {fp}: {str(e)}")
                
                # Update the status periodically to show progress
                if files_count % 1000 == 0:  # Ogni 1000 file
                    self.root.after(0, lambda size=total_size: 
                        self.status_label.config(text=f"Calcolando dimensione: {self._format_size(size)}..."))
                    
            return total_size
        except Exception as e:
            self.log_debug(f"Error calculating directory size for {path}: {str(e)}")
            return 0
        
    def get_directory_size_system(self, path):
        """Utilizza comandi di sistema per ottenere dimensioni di directory molto grandi"""
        try:
            if os.name == 'nt':  # Windows
                # Usa PowerShell per calcolare la dimensione (molto più veloce per directory grandi)
                cmd = f'powershell -command "Get-ChildItem -Path \'{path}\' -Recurse -Force -ErrorAction SilentlyContinue | Measure-Object -Property Length -Sum | Select-Object -ExpandProperty Sum"'
                result = subprocess.check_output(cmd, shell=True, stderr=subprocess.STDOUT)
                size = int(result.strip())
                return size
            else:  # Linux/Unix
                # Usa il comando du che è ottimizzato per il calcolo delle dimensioni
                result = subprocess.check_output(['du', '-sb', path])
                size = int(result.split()[0])
                return size
        except Exception as e:
            self.log_debug(f"Errore nel calcolo della dimensione tramite comando di sistema: {str(e)}")
            # Fallback al metodo standard
            return self.get_directory_size(path)
        
    def estimate_directory_size(self, path, sample_size=100):
        """Stima la dimensione di una directory campionando alcuni file"""
        import random
        
        if not os.path.exists(path) or os.path.isfile(path):
            return self.get_directory_size(path)  # Usa il metodo esatto per file o percorsi non validi
        
        try:
            # Ottieni un conteggio rapido dei file (questo è veloce)
            total_files = 0
            sampled_files = 0
            total_sampled_size = 0
            
            # Prima passata veloce per contare i file
            for root, _, files in os.walk(path, topdown=True):
                total_files += len(files)
                # Limita il tempo della prima passata
                if total_files > 10000:  # Se ci sono più di 10000 file, passiamo alla stima
                    break
            
            # Se pochi file, usare metodo preciso
            if total_files < 1000:
                return self.get_directory_size(path)
                
            # Seconda passata per il campionamento
            for root, _, files in os.walk(path, topdown=True):
                for file in files:
                    # Campiona casualmente 1 file ogni X
                    if random.randint(1, max(1, total_files // sample_size)) == 1:
                        try:
                            file_path = os.path.join(root, file)
                            if os.path.exists(file_path) and not os.path.islink(file_path):
                                total_sampled_size += os.path.getsize(file_path)
                                sampled_files += 1
                        except:
                            pass
                    
                    # Se abbiamo campionato abbastanza file, calcola la stima
                    if sampled_files >= sample_size:
                        break
                
                if sampled_files >= sample_size:
                    break
            
            # Calcola la stima finale
            if sampled_files > 0:
                avg_file_size = total_sampled_size / sampled_files
                estimated_size = avg_file_size * total_files
                self.log_debug(f"Dimensione stimata per {path}: {self._format_size(estimated_size)} (basata su {sampled_files} campioni)")
                return estimated_size
            else:
                return self.get_directory_size(path)  # Fallback al metodo standard
                
        except Exception as e:
            self.log_debug(f"Errore nella stima della dimensione: {str(e)}")
            return 0
     
    def get_disk_space(self, path):
        """Get disk space information for the partition containing the path"""
        if not os.path.exists(path):
            return (0, 0, 0)
        
        try:
            if os.name == 'nt':  # Windows
                total, used, free = shutil.disk_usage(os.path.splitdrive(path)[0] + '\\')
            else:  # Linux/Mac
                total, used, free = shutil.disk_usage(os.path.abspath(os.path.join(path, os.pardir)))
            return (total, used, free)
        except Exception as e:
            self.log_debug(f"Error getting disk space for {path}: {str(e)}")
            return (0, 0, 0)

    def update_disk_info(self, path=None, calculate_dir_size=True):
        """
        Aggiorna le informazioni del disco.
        Funzione wrapper che mantiene compatibilità con le chiamate esistenti.
        """
        # Se non viene specificato un percorso, usa quello corrente
        if path is None:
            path = self.search_path.get()
        
        # Se il percorso non esiste o è vuoto, non fare nulla
        if not path or not os.path.exists(path):
            self.log_debug(f"Percorso non valido per update_disk_info: {path}")
            return
        
        # Avvia un thread separato per calcolare le informazioni del disco
        threading.Thread(
            target=self._async_update_disk_info,
            args=(path, calculate_dir_size),
            daemon=True
        ).start()
        
    def _update_disk_info_thread(self, path, calculate_dir_size):
        """Thread per calcolare le informazioni del disco"""
        try:
            # Get disk space info
            total, used, free = self.get_disk_space(path)
            
            # Aggiorna l'UI dal thread principale
            self.root.after(0, lambda: self.total_disk_var.set(self._format_size(total)))
            self.root.after(0, lambda: self.used_disk_var.set(self._format_size(used)))
            self.root.after(0, lambda: self.free_disk_var.set(self._format_size(free)))
            self.root.after(0, lambda: self.status_label.config(text="In attesa..."))
        except Exception as e:
            self.log_debug(f"Errore nel calcolo dello spazio disco: {str(e)}")
            self.root.after(0, lambda: self.total_disk_var.set("Errore"))
            self.root.after(0, lambda: self.used_disk_var.set("Errore"))
            self.root.after(0, lambda: self.free_disk_var.set("Errore"))
        
        # Calcola dimensione directory solo se richiesto
        if calculate_dir_size:
            calculation_mode = self.dir_size_calculation.get()
            if calculation_mode != "disabilitato" and not self.is_searching:
                self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
                self._calculate_dir_size_thread(path)

    def _async_update_disk_info(self, path, calculate_dir_size):
        """Esegue il calcolo delle informazioni del disco in background"""
        try:
            # Get disk space info
            total, used, free = self.get_disk_space(path)
            
            # Aggiorna l'UI dal thread principale
            self.root.after(0, lambda: self.total_disk_var.set(self._format_size(total)))
            self.root.after(0, lambda: self.used_disk_var.set(self._format_size(used)))
            self.root.after(0, lambda: self.free_disk_var.set(self._format_size(free)))
        except Exception as e:
            self.log_debug(f"Errore nel calcolo dello spazio disco: {str(e)}")
            self.root.after(0, lambda: self.total_disk_var.set("Errore"))
            self.root.after(0, lambda: self.used_disk_var.set("Errore"))
            self.root.after(0, lambda: self.free_disk_var.set("Errore"))
        
        # MODIFICA QUI: Verifica modalità di calcolo dimensione
        calculation_mode = self.dir_size_calculation.get()
        
        if calculation_mode == "disabilitato":
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo disattivato"))
            return
        
        # Se è incrementale e siamo in fase di ricerca, non fare niente qui
        if calculation_mode == "incrementale" and self.is_searching:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            return
                
        # Per le altre modalità o se non siamo in ricerca
        if calculate_dir_size:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo in corso..."))
            # Il calcolo avviene già in un thread separato
            self._calculate_dir_size_thread(path)
        else:
            self.root.after(0, lambda: self.dir_size_var.set("Calcolo disattivato"))
        
    def _calculate_dir_size_thread(self, path):
        """Thread function to calculate directory size"""
        calculation_mode = self.dir_size_calculation.get()
        dir_size = 0
        
        try:
            if calculation_mode == "preciso":
                dir_size = self.get_directory_size(path)
            elif calculation_mode == "stimato":
                dir_size = self.estimate_directory_size(path)
            elif calculation_mode == "sistema":
                dir_size = self.get_directory_size_system(path)
            else:  # incrementale o fallback
                dir_size = self.get_directory_size(path)
                
            # Update the UI from the main thread
            self.root.after(0, lambda: self.dir_size_var.set(self._format_size(dir_size)))
            self.root.after(0, lambda: self.status_label.config(text="In attesa..."))
        except Exception as e:
            self.log_debug(f"Errore nel calcolo della dimensione: {str(e)}")
            self.root.after(0, lambda: self.dir_size_var.set("Errore"))
            self.root.after(0, lambda: self.status_label.config(text="In attesa..."))

    def refresh_directory_size(self):
        """Aggiorna manualmente il calcolo della dimensione della directory"""
        # Ottieni il percorso corrente
        path = self.search_path.get()
        
        # Verifica che il percorso esista
        if not path or not os.path.exists(path):
            messagebox.showinfo("Informazione", "Seleziona prima un percorso valido")
            return
            
        # Verifica la modalità di calcolo
        calculation_mode = self.dir_size_calculation.get()
        if calculation_mode == "disabilitato":
            response = messagebox.askyesno("Calcolo disabilitato", 
                                        "Il calcolo della dimensione è attualmente disabilitato.\n\n" +
                                        "Vuoi attivarlo e procedere con il calcolo?")
            if response:
                # Seleziona la modalità "preciso"
                self.dir_size_calculation.set("preciso")
            else:
                return
        
        # Aggiorna lo stato
        self.dir_size_var.set("Calcolo in corso...")
        self.status_label.config(text="Calcolo dimensione directory...")
        
        # Esegui il calcolo in un thread separato
        threading.Thread(target=self._calculate_dir_size_thread, args=(path,), daemon=True).start()

    # Funzione helper per formattare la dimensione del file
    def _format_size(self, size_bytes):
        """Formatta la dimensione del file in modo leggibile"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.2f} KB"
        elif size_bytes < 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024):.2f} MB"
        else:
            return f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"

    # Funzione helper per determinare il tipo di file
    def _get_file_type(self, file_path):
        """Determina il tipo di file in base all'estensione"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.txt', '.md', '.rtf']:
            return "Documento di testo"
        elif ext in ['.doc', '.docx', '.odt']:
            return "Documento Word"
        elif ext in ['.xls', '.xlsx', '.ods']:
            return "Foglio di calcolo"
        elif ext in ['.pdf']:
            return "PDF"
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            return "Immagine"
        elif ext in ['.mp3', '.wav', '.ogg', '.flac']:
            return "Audio"
        elif ext in ['.mp4', '.avi', '.mkv', '.mov']:
            return "Video"
        elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
            return "Archivio"
        elif ext in ['.exe', '.dll', '.bat', '.cmd']:
            return "Eseguibile"
        else:
            return "File"
    
    def get_main_folder_name(self):
        """Mostra una finestra di dialogo personalizzata per richiedere il nome della cartella principale"""
        dialog = tk.Toplevel(self.root)
        dialog.title("Cartella Principale")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Variabile per memorizzare il risultato
        result = {"name": None}
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=20)
        main_frame.pack(fill=tk.BOTH, expand=tk.YES)
        
        # Label di istruzione
        ttk.Label(main_frame, text="Inserisci il nome della cartella principale nell'archivio:", 
                font=("", 10)).pack(pady=(0, 10))
        
        # Campo di input
        name_var = tk.StringVar(value="files")
        name_entry = ttk.Entry(main_frame, textvariable=name_var, width=40)
        name_entry.pack(fill=tk.X, pady=10)
        name_entry.select_range(0, "end")  # Seleziona tutto il testo predefinito
        name_entry.focus_set()  # Imposta il focus
        
        # Frame per i pulsanti
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=tk.X, pady=(15, 0))
        
        # Funzioni di callback
        def on_cancel():
            result["name"] = None
            dialog.destroy()
        
        def on_create():
            folder_name = name_var.get().strip()
            if folder_name:
                result["name"] = folder_name
                dialog.destroy()
            else:
                messagebox.showwarning("Attenzione", "Inserisci un nome valido per la cartella", parent=dialog)
        
        # Pulsanti con stili
        cancel_btn = ttk.Button(btn_frame, text="Annulla", command=on_cancel, width=15)
        cancel_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        create_btn = ttk.Button(btn_frame, text="Crea Cartella", command=on_create, 
                            style="primary.TButton", width=15)
        create_btn.pack(side=tk.RIGHT)
        
        # Gestisci l'evento Invio e Escape
        dialog.bind("<Return>", lambda e: on_create())
        dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Centra la finestra sullo schermo
        dialog.update_idletasks()  # Aggiorna per ottenere dimensioni corrette
        width = dialog.winfo_reqwidth() + 50
        height = dialog.winfo_reqheight() + 20
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima
        dialog.minsize(350, 180)
        
        # Blocca la finestra fino alla chiusura
        dialog.wait_window()
        
        return result["name"]
    
    def open_file_location(self, event=None):
        """Apre il percorso del file selezionato nel file explorer"""
               
        selected_items = self.results_list.selection()
        if not selected_items:
            return
            
        selected_item = selected_items[0]  # Prendi il primo elemento selezionato
        file_path = self.results_list.item(selected_item, "values")[5]  # Ottieni il percorso del file
        
        try:
            if os.path.exists(file_path):
                # Ottieni la directory contenente il file
                directory = os.path.dirname(file_path)
                
                if os.name == 'nt':  # Windows
                    # Converti eventuali forward slash in backslash per Windows
                    file_path = os.path.normpath(file_path)
                    # Usa il metodo più sicuro con subprocess invece di os.system
                    subprocess.run(['explorer', '/select,', file_path], shell=True)
                else:
                    # Per sistemi Linux/Unix
                    if shutil.which('xdg-open'):  # Verifica che xdg-open sia disponibile
                        subprocess.run(['xdg-open', directory])
                    elif shutil.which('open'):  # Per macOS
                        subprocess.run(['open', directory])
                    else:
                        self.log_debug("Nessun comando disponibile per aprire directory")
                        messagebox.showinfo("Informazione", f"Percorso del file: {directory}")
                
                self.log_debug(f"Apertura percorso: {file_path}")
            else:
                messagebox.showinfo("Informazione", f"Il percorso non esiste: {file_path}")
                self.log_debug(f"Percorso non esistente: {file_path}")
        except Exception as e:
            messagebox.showerror("Errore", f"Impossibile aprire il percorso: {str(e)}")
            self.log_debug(f"Errore nell'apertura del percorso: {str(e)}")

    def show_advanced_filters_dialog(self):
        """Mostra la finestra di dialogo per i filtri di ricerca avanzati"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Filtri avanzati")
        dialog.geometry("530x250")  # Increased height to accommodate the new field
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Filtri dimensione
        size_frame = ttk.LabelFrame(dialog, text="Dimensione file")
        size_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(size_frame, text="Min (KB):").grid(row=0, column=0, padx=5, pady=5)
        min_size = ttk.Entry(size_frame, width=10)
        min_size.grid(row=0, column=1, padx=5, pady=5)
        min_size.insert(0, str(self.advanced_filters["size_min"] // 1024))
        self.create_tooltip(min_size, "Quando viene imposto un valore maggiore di zero, solo i file con dimensione superiore a questo valore verranno inclusi nei risultati")
        
        ttk.Label(size_frame, text="Max (KB):").grid(row=0, column=2, padx=5, pady=5)
        max_size = ttk.Entry(size_frame, width=10)
        max_size.grid(row=0, column=3, padx=5, pady=5)
        max_size.insert(0, str(self.advanced_filters["size_max"] // 1024 if self.advanced_filters["size_max"] else 0))
        self.create_tooltip(max_size, "Quando viene imposto un valore maggiore di zero, solo i file con dimensione inferiore a questo valore verranno inclusi nei risultati")

        # Filtri data - FORMAT DD-MM-YYYY
        date_frame = ttk.LabelFrame(dialog, text="Data modifica (DD-MM-YYYY)")
        date_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(date_frame, text="Da:").grid(row=0, column=0, padx=5, pady=5)
        min_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        min_date.grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(date_frame, text="A:").grid(row=0, column=2, padx=5, pady=5)
        max_date = ttk.DateEntry(date_frame, dateformat="%d-%m-%Y")
        max_date.grid(row=0, column=3, padx=5, pady=5)
        
        # Cancella sempre le date precedenti
        min_date.entry.delete(0, "end")
        max_date.entry.delete(0, "end")
        
        # NUOVO: Aggiungi frame per le estensioni
        ext_frame = ttk.LabelFrame(dialog, text="Estensioni file")
        ext_frame.pack(fill=X, padx=10, pady=5)
        
        ttk.Label(ext_frame, text="Estensioni (separate da virgola):").pack(side=LEFT, padx=5, pady=5)
        
        # Inizializza con estensioni esistenti
        current_exts = ", ".join([ext.lstrip('.') for ext in self.advanced_filters.get("extensions", [])])
        extensions = ttk.Entry(ext_frame, width=30)
        extensions.pack(side=LEFT, fill=X, expand=YES, padx=5, pady=5)
        extensions.insert(0, current_exts)
        self.create_tooltip(extensions, "Esempio: txt, pdf, docx (senza il punto iniziale)")
        
        # Aggiungiamo un frame di debug per vedere i filtri correnti
        debug_frame = ttk.LabelFrame(dialog, text="Debug - Stato filtri correnti")
        debug_frame.pack(fill=X, padx=10, pady=5)
        
        debug_text = ttk.Text(debug_frame, height=3, width=50)
        debug_text.pack(fill=X, padx=5, pady=5)
        debug_text.insert("1.0", f"Data min: {self.advanced_filters['date_min']}\n")
        debug_text.insert("2.0", f"Data max: {self.advanced_filters['date_max']}\n")
        debug_text.insert("3.0", f"Extensions: {self.advanced_filters['extensions']}")
        debug_text.config(state="disabled")
        
        # Pulsante Salva
        def save_filters():
            try:
                # Analizza i filtri di dimensione
                min_kb = int(min_size.get() or 0)
                max_kb = int(max_size.get() or 0)
                self.advanced_filters["size_min"] = min_kb * 1024
                self.advanced_filters["size_max"] = max_kb * 1024
                
                # Ottieni le date nel formato DD-MM-YYYY
                min_date_value = min_date.entry.get().strip()
                max_date_value = max_date.entry.get().strip()
                
                print(f"DEBUG - Date inserite: min={min_date_value}, max={max_date_value}")
                
                # Validazione date
                if min_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(min_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data minima non valido. Usa DD-MM-YYYY")
                        return
                
                if max_date_value:
                    try:
                        # Verifica formato corretto
                        datetime.strptime(max_date_value, "%d-%m-%Y")
                    except ValueError:
                        messagebox.showerror("Errore", "Formato data massima non valido. Usa DD-MM-YYYY")
                        return
                
                # Verifica che la data minima non sia successiva alla data massima
                if min_date_value and max_date_value:
                    min_date_obj = datetime.strptime(min_date_value, "%d-%m-%Y")
                    max_date_obj = datetime.strptime(max_date_value, "%d-%m-%Y")
                    
                    if min_date_obj > max_date_obj:
                        messagebox.showerror("Errore", 
                                        "La data di inizio non può essere successiva alla data di fine")
                        return
                
                # Salva le date validate
                self.advanced_filters["date_min"] = min_date_value
                self.advanced_filters["date_max"] = max_date_value
                
                # Analizza le estensioni
                exts = [e.strip() for e in extensions.get().split(",") if e.strip()]
                self.advanced_filters["extensions"] = [f".{e.lstrip('.')}" for e in exts]
                
                print(f"Filtri salvati: {self.advanced_filters}")  # Debug info
                dialog.destroy()
                
            except ValueError:
                messagebox.showerror("Errore", "Inserisci valori numerici validi per le dimensioni")
        
        # Creato un frame per i pulsanti nella stessa riga
        buttons_frame = ttk.Frame(dialog)
        buttons_frame.pack(fill=X, pady=10, padx=10)
        
        # Pulsante Annulla a sinistra
        ttk.Button(buttons_frame, text="Annulla", command=dialog.destroy).pack(side=RIGHT, padx=5)
        
        # Pulsante Salva a destra
        ttk.Button(buttons_frame, text="Salva", command=save_filters).pack(side=RIGHT, padx=5)

        dialog.update_idletasks()  # Aggiorna la finestra per ottenere le dimensioni corrette
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Imposta una dimensione minima per la finestra
        dialog.minsize(500, 400)

    def configure_extensions(self, mode="base"):
        """Dialog to configure file extensions for different search modes"""
        dialog = ttk.Toplevel(self.root)
        dialog.title(f"Configura estensioni - Modalità {mode.capitalize()}")
        dialog.geometry("1000x450")
        dialog.transient(self.root)
        dialog.grab_set()
        
        main_frame = ttk.Frame(dialog, padding=15)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Heading text
        ttk.Label(main_frame, text=f"Seleziona le estensioni di file da includere nella ricerca per la modalità {mode}", 
                font=("", 10)).pack(anchor=W, pady=(0, 15))
        
        # Category tabs
        notebook = ttk.Notebook(main_frame)
        notebook.pack(fill=BOTH, expand=YES, pady=(0, 15))
        
        # Initialize dictionary to store checkbutton variables
        ext_vars = {}
        
        # Define extensions by category - EXPANDED with more extensions
        extension_categories = {
            "Documenti": [
                (".txt", "File di testo"),
                (".doc", "Word vecchio"),
                (".docx", "Word"),
                (".pdf", "PDF"),
                (".rtf", "Rich Text"),
                (".odt", "OpenDoc Text"),
                (".md", "Markdown"),
                (".csv", "CSV"),
                (".xml", "XML"),
                (".html", "HTML"),
                (".htm", "HTM"),
                (".json", "JSON"),
                (".log", "Log file"),
                (".tex", "LaTeX"),
                (".rst", "reStructuredText"),
                (".epub", "E-book EPUB"),
                (".mobi", "E-book Mobi")
            ],
            "Fogli calcolo": [
                (".xls", "Excel vecchio"),
                (".xlsx", "Excel"),
                (".ods", "OpenCalc"),
                (".csv", "CSV"),
                (".tsv", "TSV"),
                (".dbf", "Database File"),
                (".dif", "Data Interchange Format")
            ],
            "Presentazioni": [
                (".ppt", "PowerPoint vecchio"),
                (".pptx", "PowerPoint"),
                (".odp", "OpenImpress"),
                (".key", "Keynote"),
                (".pps", "PowerPoint Show")
            ],
            "Database": [
                (".db", "Database generico"),
                (".sqlite", "SQLite"),
                (".sqlite3", "SQLite3"),
                (".mdb", "Access DB"),
                (".accdb", "Access DB nuovo"),
                (".odb", "OpenOffice DB"),
                (".sql", "SQL script")
            ],
            "Immagini": [
                (".jpg", "JPEG"),
                (".jpeg", "JPEG"),
                (".png", "PNG"),
                (".gif", "GIF"),
                (".bmp", "Bitmap"),
                (".tiff", "TIFF"),
                (".tif", "TIF"),
                (".svg", "SVG"),
                (".webp", "WebP"),
                (".ico", "Icon"),
                (".raw", "Raw"),
                (".psd", "Photoshop"),
                (".ai", "Illustrator"),
                (".odg", "OpenOffice Draw"),
                (".xcf", "GIMP"),
                (".heic", "HEIC")
            ],
            "Audio": [
                (".mp3", "MP3"),
                (".wav", "WAV"),
                (".ogg", "OGG"),
                (".flac", "FLAC"),
                (".aac", "AAC"),
                (".m4a", "M4A"),
                (".wma", "WMA"),
                (".mid", "MIDI"),
                (".midi", "MIDI"),
                (".aiff", "AIFF"),
                (".opus", "Opus")
            ],
            "Video": [
                (".mp4", "MP4"),
                (".avi", "AVI"),
                (".mkv", "MKV"),
                (".mov", "MOV"),
                (".wmv", "WMV"),
                (".flv", "FLV"),
                (".webm", "WebM"),
                (".m4v", "M4V"),
                (".mpg", "MPEG"),
                (".mpeg", "MPEG"),
                (".3gp", "3GP"),
                (".ogv", "OGV"),
                (".ts", "TS")
            ],
            "Archivi": [
                (".zip", "ZIP"),
                (".rar", "RAR"),
                (".7z", "7-Zip"),
                (".tar", "TAR"),
                (".gz", "GZip"),
                (".bz2", "BZip2"),
                (".iso", "ISO"),
                (".tgz", "Tar GZipped"),
                (".xz", "XZ"),
                (".cab", "Cabinet"),
                (".jar", "Java Archive")
            ],
            "Eseguibili": [
                (".exe", "Eseguibile"),
                (".dll", "Libreria"),
                (".bat", "Batch"),
                (".cmd", "Command"),
                (".ps1", "PowerShell"),
                (".vbs", "VBScript"),
                (".sh", "Shell script"),
                (".msi", "Installer Windows"),
                (".app", "Applicazione macOS"),
                (".deb", "Pacchetto Debian"),
                (".rpm", "Red Hat Package"),
                (".apk", "Android Package")
            ],
            "Configurazione": [
                (".ini", "INI"),
                (".config", "Config"),
                (".conf", "Config"),
                (".reg", "Registry"),
                (".cfg", "Config"),
                (".properties", "Properties"),
                (".yml", "YAML"),
                (".yaml", "YAML"),
                (".json", "JSON Config"),
                (".toml", "TOML"),
                (".env", "Environment"),
                (".htaccess", "Apache Config"),
                (".plist", "macOS Property List")
            ],
            "Programmazione": [
                (".c", "C"),
                (".cpp", "C++"),
                (".cs", "C#"),
                (".java", "Java"),
                (".py", "Python"),
                (".js", "JavaScript"),
                (".php", "PHP"),
                (".rb", "Ruby"),
                (".go", "Golang"),
                (".rs", "Rust"),
                (".swift", "Swift"),
                (".pl", "Perl"),
                (".lua", "Lua"),
                (".h", "Header C"),
                (".hpp", "Header C++"),
                (".vb", "Visual Basic"),
                (".ts", "TypeScript"),
                (".scala", "Scala"),
                (".groovy", "Groovy"),
                (".kt", "Kotlin")
            ]
        }
        
        # Create a list of all extensions for "profonda" mode
        all_extensions = []
        for category_extensions in extension_categories.values():
            for ext, _ in category_extensions:
                all_extensions.append(ext.lower())
        
        # Load current settings (this would load from your saved settings)
        current_settings = self.get_extension_settings(mode)
        
        # Extensions that should be included in each search level
        base_extensions = ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log', 
                        '.docx', '.pdf', '.pptx', '.xlsx', '.rtf', '.odt', '.xls', '.doc']
                        
        advanced_extensions = base_extensions + ['.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', 
                                            '.vbs', '.js', '.config', '.ini', '.reg']
        
        # Per modalità profonda, usa tutte le estensioni definite
        if mode == "profonda":
            # Sovrascrive le impostazioni correnti per far sì che tutte le estensioni siano selezionate
            current_settings = all_extensions
        
         # Dizionario per tenere traccia dei pulsanti per categoria
        category_buttons = {}
        
        # Create tabs for each category
        for category, extensions in extension_categories.items():
            # Create a frame for this category
            category_frame = ttk.Frame(notebook, padding=10)
            notebook.add(category_frame, text=category)
            
            # Aggiungi pulsante "Seleziona tutto" in alto
            button_frame = ttk.Frame(category_frame)
            button_frame.pack(fill=X, pady=(0, 10))
            
            # Funzione per selezionare tutte le estensioni in questa categoria
            def select_all_category(cat=category):
                for ext, desc in extension_categories[cat]:
                    if ext.lower() in ext_vars:
                        ext_vars[ext.lower()].set(True)
            
            # Funzione per deselezionare tutte le estensioni in questa categoria
            def deselect_all_category(cat=category):
                for ext, desc in extension_categories[cat]:
                    if ext.lower() in ext_vars:
                        ext_vars[ext.lower()].set(False)
            
            # Aggiungi i pulsanti
            select_all_btn = ttk.Button(button_frame, text=f"Seleziona tutto", 
                                    command=select_all_category)
            select_all_btn.pack(side=LEFT, padx=(0, 5))
            
            deselect_all_btn = ttk.Button(button_frame, text=f"Deseleziona tutti", 
                                        command=deselect_all_category)
            deselect_all_btn.pack(side=LEFT)
            
            # Salva i pulsanti nel dizionario per riferimento futuro
            category_buttons[category] = (select_all_btn, deselect_all_btn)
            
            # Create a grid layout
            content_frame = ttk.Frame(category_frame)
            content_frame.pack(fill=BOTH, expand=YES)
            row, col = 0, 0
            
            # Add checkboxes for each extension in this category
            for ext, desc in extensions:
                # Check if this extension should be selected based on mode
                is_selected = ext.lower() in current_settings
                
                # Create a variable for this checkbox
                var = BooleanVar(value=is_selected)
                ext_vars[ext.lower()] = var
                
                # Create the checkbox
                cb = ttk.Checkbutton(content_frame, text=f"{ext} ({desc})", variable=var)
                cb.grid(row=row, column=col, sticky=W, padx=5, pady=3)
                
                # Update grid position
                col += 1
                if col > 2:  # 3 columns per row
                    col = 0
                    row += 1
        
        # Button frame
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(10, 0))
        
        # Helper functions
        def select_all():
            for var in ext_vars.values():
                var.set(True)
                
        def deselect_all():
            for var in ext_vars.values():
                var.set(False)
                
        def restore_defaults():
            # Modified restore_defaults function
            if mode == "profonda":
                # For deep mode, select ALL extensions
                for var in ext_vars.values():
                    var.set(True)
            else:
                # Reset to default extensions for base/avanzata modes
                default_list = base_extensions if mode == "base" else advanced_extensions
                for ext, var in ext_vars.items():
                    var.set(ext in default_list)
        
        def save_settings():
            # Save the selected extensions
            selected_extensions = [ext for ext, var in ext_vars.items() if var.get()]
            self.save_extension_settings(mode, selected_extensions)
            dialog.destroy()
            
            # Update the search_depth combo if needed
            current_depth = self.search_depth.get()
            if current_depth == mode:
                # Refresh the UI to reflect the new settings
                self.log_debug(f"Aggiornate le estensioni per la modalità {mode} ({len(selected_extensions)} estensioni)")
        
        # Buttons
        ttk.Button(btn_frame, text="Seleziona tutti", command=select_all).pack(side=LEFT, padx=5)
        ttk.Button(btn_frame, text="Deseleziona tutti", command=deselect_all).pack(side=LEFT, padx=5)
        ttk.Button(btn_frame, text="Ripristina default", command=restore_defaults).pack(side=LEFT, padx=5)
        
        ttk.Button(btn_frame, text="Annulla", command=dialog.destroy).pack(side=RIGHT, padx=5)
        ttk.Button(btn_frame, text="Salva", command=save_settings).pack(side=RIGHT, padx=5)
        
        # Center the dialog on screen
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        
        # Initial focus
        notebook.select(0)  # Focus first tab
        
    def get_default_extensions(self, mode="base"):
        """Get default extensions for the specified search mode"""
        if mode == "base":
            return ['.txt', '.md', '.csv', '.html', '.htm', '.xml', '.json', '.log', 
                    '.docx', '.pdf', '.pptx', '.xlsx', '.rtf', '.odt', '.xls', '.doc']
        elif mode == "avanzata":
            # Prima ottieni le estensioni base
            base_exts = self.get_default_extensions("base")
            
            # Poi aggiungi estensioni avanzate (senza duplicati)
            advanced_only = [
                # File di sistema e script
                '.exe', '.dll', '.sys', '.bat', '.cmd', '.ps1', '.vbs', 
                '.config', '.ini', '.reg',
                
                # File di programmazione comuni
                '.py', '.java', '.php', '.cs', '.cpp', '.c', '.h', '.rb', '.js',
                
                # File di database semplici
                '.db', '.sqlite', '.sqlite3',
                
                # File di configurazione aggiuntivi
                '.env', '.yml', '.yaml', '.toml', '.json5',
                
                # File di backup e temporanei
                '.bak', '.old', '.tmp', '.temp',
                
                # Formati di documento meno comuni
                '.epub', '.tex', '.rst',
                
                # File web
                '.css', '.less', '.scss', '.jsp', '.asp', '.aspx'
            ]
            
            # Combina le liste evitando duplicati
            return base_exts + [ext for ext in advanced_only if ext not in base_exts]
        else:  # profonda - usa un metodo diverso per ottenere tutte le estensioni
            return []  # Il valore effettivo viene determinato in configure_extensions
        
    def get_extension_settings(self, mode="base"):
        """Load saved extension settings for the specified search mode"""
        if not hasattr(self, 'extension_settings'):
            # Initialize with defaults
            self.extension_settings = {
                "base": self.get_default_extensions("base"),
                "avanzata": self.get_default_extensions("avanzata"),
                "profonda": self.get_default_extensions("profonda")
            }
        
        # Aggiungi log per verificare le estensioni quando si accede ad esse
        extensions = self.extension_settings.get(mode, [])
        self.log_debug(f"Estensioni caricate per modalità {mode}: {', '.join(extensions)}")
        return extensions
        

    def save_extension_settings(self, mode, extensions):
        """Save extension settings for the specified search mode"""
        if not hasattr(self, 'extension_settings'):
            self.extension_settings = {}
        
        # Ensure all extensions have a leading dot and are lowercase
        normalized_extensions = []
        for ext in extensions:
            if not ext.startswith('.'):
                ext = '.' + ext
            normalized_extensions.append(ext.lower())
        
        self.extension_settings[mode] = normalized_extensions
        
        # Log what was saved
        self.log_debug(f"Saved {len(normalized_extensions)} extensions for {mode} mode")
        self.log_debug(f"Extensions: {', '.join(normalized_extensions)}")
        
        # Forza un aggiornamento dell'interfaccia se necessario
        if hasattr(self, 'search_depth') and self.search_depth.get() == mode:
            self.log_debug(f"Aggiornata UI per modalità {mode}")
            
        # Qui potresti aggiungere codice per salvare le impostazioni su file
        self.save_settings_to_file()
        
    def save_settings_to_file(self):
        """Salva le impostazioni delle estensioni su un file"""
        try:
            import json
            import os
            
            # Cartella per le impostazioni
            settings_dir = os.path.join(os.path.expanduser("~"), ".file_search_tool")
            if not os.path.exists(settings_dir):
                os.makedirs(settings_dir)
                
            # File per le impostazioni delle estensioni
            settings_file = os.path.join(settings_dir, "extension_settings.json")
            
            # Salva le impostazioni
            with open(settings_file, 'w') as f:
                json.dump(self.extension_settings, f)
                
            self.log_debug(f"Impostazioni estensioni salvate in {settings_file}")
        except Exception as e:
            self.log_debug(f"Errore nel salvataggio delle impostazioni: {str(e)}")
            
    def load_settings_from_file(self):
        """Carica le impostazioni delle estensioni da un file"""
        try:
            import json
            import os
            
            # File per le impostazioni delle estensioni
            settings_file = os.path.join(os.path.expanduser("~"), ".file_search_tool", "extension_settings.json")
            
            # Controlla se il file esiste
            if os.path.exists(settings_file):
                with open(settings_file, 'r') as f:
                    self.extension_settings = json.load(f)
                    self.log_debug(f"Impostazioni estensioni caricate da {settings_file}")
            else:
                # Inizializza con i valori predefiniti
                self.extension_settings = {
                    "base": self.get_default_extensions("base"),
                    "avanzata": self.get_default_extensions("avanzata"),
                    "profonda": self.get_default_extensions("profonda")
                }
        except Exception as e:
            self.log_debug(f"Errore nel caricamento delle impostazioni: {str(e)}")
            # Inizializza con i valori predefiniti
            self.extension_settings = {
                "base": self.get_default_extensions("base"),
                "avanzata": self.get_default_extensions("avanzata"),
                "profonda": self.get_default_extensions("profonda")
            }

    def create_widgets(self):
        # Frame principale che conterrà tutto
        main_container = ttk.Frame(self.root)
        main_container.pack(fill=BOTH, expand=YES)
        
        # Intestazione (titolo e informazioni)
        header_frame = ttk.Frame(main_container, padding="10")
        header_frame.pack(fill=X)

        # Layout a tre colonne in una singola riga
        # 1. Tema a sinistra
        theme_frame = ttk.Frame(header_frame)
        theme_frame.pack(side=LEFT, fill=Y)

        ttk.Label(theme_frame, text="Tema:").pack(side=LEFT)
        themes = ttk.Style().theme_names()
        self.theme_combobox = ttk.Combobox(theme_frame, values=themes, width=15)
        self.theme_combobox.pack(side=LEFT, padx=5)
        self.theme_combobox.current(themes.index("darkly"))
        self.theme_combobox.bind("<<ComboboxSelected>>", lambda e: [ttk.Style().theme_use(self.theme_combobox.get()),self.update_theme_colors()])

        # 2. Titolo al centro
        title_frame = ttk.Frame(header_frame)
        title_frame.pack(side=LEFT, expand=True)

        title_label = ttk.Label(title_frame, text="File Search Tool.. Forensics G.di F.", 
                            font=("Helvetica", 14, "bold"))
        title_label.pack(anchor=CENTER)

        # 3. Data/ora e utente a destra
        datetime_frame = ttk.Frame(header_frame)
        datetime_frame.pack(side=RIGHT, fill=Y)

        datetime_label = ttk.Label(datetime_frame, textvariable=self.datetime_var, font=("Helvetica", 9))
        datetime_label.pack(side=RIGHT)

        # ==========================================================
        # SEZIONE DEI CONTROLLI DI RICERCA (organizzati per righe)
        # ==========================================================
        controls_frame = ttk.LabelFrame(main_container, text="Parametri di ricerca", padding=10)
        controls_frame.pack(fill=X, padx=10, pady=5)
        
        # ------------------------------------------------------
        # RIGA 1: Directory di ricerca
        # ------------------------------------------------------
        path_frame = ttk.Frame(controls_frame)
        path_frame.pack(fill=X, pady=5)
        
        path_label = ttk.Label(path_frame, text="Directory:", width=12, anchor=W)
        path_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(path_label, "Seleziona la directory per effettuare la ricerca dei file")
        
        self.path_entry = ttk.Entry(path_frame, textvariable=self.search_path)
        self.path_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)
        
        self.browse_btn = ttk.Button(path_frame, text="Sfoglia", command=self.browse_directory, width=10)
        self.browse_btn.pack(side=LEFT)
        
        # ------------------------------------------------------
        # RIGA 2: Parole chiave
        # ------------------------------------------------------
        keyword_frame = ttk.Frame(controls_frame)
        keyword_frame.pack(fill=X, pady=5)
        
        keyword_label = ttk.Label(keyword_frame, text="Parole chiave:", width=12, anchor=W)
        keyword_label.pack(side=LEFT, padx=(0, 5))
        self.create_tooltip(keyword_label, "Per la ricerca di più parole usa la virgola. Esempio: documento, fattura, contratto\n"
                                        "Attiva 'Parola intera' per cercare 'log' senza trovare 'login' o 'logo'")
        
        self.keyword_entry = ttk.Entry(keyword_frame, textvariable=self.keywords)
        self.keyword_entry.pack(side=LEFT, fill=X, expand=YES, padx=5)

        # ------------------------------------------------------
        # RIGA 3: Opzioni base di ricerca (checkbox)
        # ------------------------------------------------------
        options_frame = ttk.Frame(controls_frame)
        options_frame.pack(fill=X, pady=5)
                  
        ttk.Label(options_frame, text="Livelli di ricerca:").pack(side=LEFT, padx=(0, 5))
        search_depth_combo = ttk.Combobox(options_frame, textvariable=self.search_depth, 
                                        values=["base", "avanzata", "profonda"], 
                                        width=10, state="readonly")
        search_depth_combo.pack(side=LEFT, padx=5)
        search_depth_combo.current(0)

        extensions_btn = ttk.Button(options_frame, text="Configura estensioni", 
                    command=lambda: self.configure_extensions(self.search_depth.get()))
        extensions_btn.pack(side=LEFT, padx=5)
        self.create_tooltip(extensions_btn, "Configura quali estensioni di file includere nella ricerca")
        
        # NUOVO: pulsante impostazioni avanzate che apre una finestra con tutte le opzioni
        advanced_options_btn = ttk.Button(options_frame, text="Impostazioni avanzate", 
                                    command=self.show_advanced_options)
        advanced_options_btn.pack(side=LEFT, padx=10)
        self.create_tooltip(advanced_options_btn, "Configura tutte le impostazioni avanzate (profondità, filtri, esclusioni, performance)")
        
        # ------------------------------------------------------
        # RIGA 4: Pulsanti di azione
        # ------------------------------------------------------
        button_frame = ttk.Frame(controls_frame)
        button_frame.pack(fill=X, pady=(10, 5))
        
        # Frame per centrare i pulsanti
        center_frame = ttk.Frame(button_frame)
        center_frame.pack(side=TOP)
        
        # Pulsante di ricerca (principale)
        self.search_button = ttk.Button(center_frame, text="CERCA", 
                                    command=self.start_search, 
                                    style="primary.TButton", width=15)
        self.search_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.search_button, "Avvia la ricerca con i criteri specificati")
        
        # Pulsante per interrompere la ricerca
        self.stop_button = ttk.Button(center_frame, text="Interrompi ricerca",
                                    command=self.stop_search_process,
                                    style="danger.TButton", width=20,
                                    state="disabled")
        self.stop_button.pack(side=LEFT, padx=10)
        self.create_tooltip(self.stop_button, "Ferma immediatamente la ricerca in corso")
        
        # Pulsante per pulire i campi di ricerca
        self.clear_btn = ttk.Button(center_frame, text="Pulisci campi", 
                    command=lambda: [self.search_path.set(""), self.keywords.set("")],
                    style="secondary.Outline.TButton", width=15)
        self.clear_btn.pack(side=LEFT, padx=10)
        self.create_tooltip(self.clear_btn, "Cancella i campi di ricerca")
        
        # Pulsante admin solo su Windows
        if os.name == 'nt':
            self.admin_button = ttk.Button(center_frame, text="Avvia come Admin", 
                                    command=self.restart_as_admin,
                                    style="info.Outline.TButton", width=20)
            self.admin_button.pack(side=LEFT, padx=10)
            
            # Disabilita il pulsante se l'app è già avviata come amministratore
            if self.is_admin:
                self.admin_button.config(state="disabled")
                self.create_tooltip(self.admin_button, "L'applicazione è già in esecuzione come amministratore")
            else:
                self.create_tooltip(self.admin_button, "Riavvia l'applicazione con privilegi di amministratore per accedere a tutte le cartelle")
        
        # ==========================================================
        # SEZIONE INFORMAZIONI TEMPORALI E SPAZIO DISCO (SOPRA LA LISTA RISULTATI)
        # ==========================================================
        info_bar = ttk.Frame(main_container)
        info_bar.pack(fill=X, padx=10, pady=5)
        
        # Frame per le informazioni temporali a sinistra
        time_frame = ttk.LabelFrame(info_bar, text="Informazioni temporali", padding=5)
        time_frame.pack(side=LEFT, fill=X, padx=(0, 5), expand=YES)
        
        time_grid = ttk.Frame(time_frame)
        time_grid.pack(fill=X, pady=2)
        
        ttk.Label(time_grid, text="Avvio:").grid(row=0, column=0, sticky=W, padx=5)
        self.start_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.start_time_label.grid(row=0, column=1, sticky=W, padx=5)
        
        ttk.Label(time_grid, text="Fine:").grid(row=0, column=2, sticky=W, padx=5)
        self.end_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.end_time_label.grid(row=0, column=3, sticky=W, padx=5)
        
        ttk.Label(time_grid, text="Durata:").grid(row=0, column=4, sticky=W, padx=5)
        self.total_time_label = ttk.Label(time_grid, text="--:--", font=("", 9, "bold"))
        self.total_time_label.grid(row=0, column=5, sticky=W, padx=5)
        
        # Frame per info disco a destra
        disk_info_frame = ttk.LabelFrame(info_bar, text="Spazio disco", padding=5)
        disk_info_frame.pack(side=LEFT, fill=X, expand=YES)

        # Uso layout semplice con pack per ordinare gli elementi in linea
        disk_grid = ttk.Frame(disk_info_frame)
        disk_grid.pack(fill=X, pady=2)

        # 1. Usato
        ttk.Label(disk_grid, text="Usato:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.used_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 2. Libero
        ttk.Label(disk_grid, text="Libero:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.free_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 3. Totale
        ttk.Label(disk_grid, text="Totale:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.total_disk_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 4. Directory
        ttk.Label(disk_grid, text="Directory:").pack(side=LEFT, padx=5)
        ttk.Label(disk_grid, textvariable=self.dir_size_var, font=("", 9, "bold")).pack(side=LEFT, padx=(0, 15))

        # 5. Pulsante aggiorna
        refresh_size_btn = ttk.Button(disk_grid, text="Aggiorna", command=self.refresh_directory_size, 
                                width=8, style="info.TButton")
        refresh_size_btn.pack(side=LEFT, padx=5)
        
        # ==========================================================
        # SEZIONE STATO RICERCA
        # ==========================================================
        # Contenitore per le informazioni di stato sopra i risultati
        status_frame = ttk.LabelFrame(main_container, text="Stato ricerca", padding=5)
        status_frame.pack(fill=X, padx=10, pady=5)
        
        # Progress bar
        self.progress_bar = ttk.Progressbar(status_frame, mode='determinate')
        self.progress_bar.pack(fill=X, pady=5)
        
        # Status grid per organizzare le informazioni
        status_grid = ttk.Frame(status_frame)
        status_grid.pack(fill=X)
        
        # Riga 1: Status
        ttk.Label(status_grid, text="Analisi:", width=12, anchor=W, font=("", 9, "bold")).grid(row=0, column=0, sticky=W, padx=5, pady=2)
        self.status_label = ttk.Label(status_grid, text="In attesa...", wraplength=1000)
        self.status_label.grid(row=0, column=1, sticky=W+E, padx=5, pady=2)
        
        # Riga 2: File analizzati
        ttk.Label(status_grid, text="File analizzati:", width=12, anchor=W, font=("", 9, "bold")).grid(row=1, column=0, sticky=W, padx=5, pady=2)
        self.analyzed_files_label = ttk.Label(status_grid, text="Nessuna ricerca avviata", wraplength=1000)
        self.analyzed_files_label.grid(row=1, column=1, sticky=W+E, padx=5, pady=2)
        
        # ==========================================================
        # SEZIONE RISULTATI 
        # ==========================================================
        results_container = ttk.LabelFrame(main_container, text="Risultati di ricerca", padding=10)
        results_container.pack(fill=BOTH, expand=YES, padx=10, pady=(5, 10))
        
        # Frame per i pulsanti di azione sui risultati
        actions_frame = ttk.Frame(results_container)
        actions_frame.pack(fill=X, pady=(0, 5))
        
        # Pulsanti per la selezione
        selection_frame = ttk.Frame(actions_frame)
        selection_frame.pack(side=LEFT)
        
        select_all_btn = ttk.Button(selection_frame, text="Seleziona tutto", command=self.select_all)
        select_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(select_all_btn, "Seleziona tutti i risultati nella lista")
        
        deselect_all_btn = ttk.Button(selection_frame, text="Deseleziona tutto", command=self.deselect_all)
        deselect_all_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(deselect_all_btn, "Deseleziona tutti i risultati")
        
        invert_sel_btn = ttk.Button(selection_frame, text="Inverti selezione", command=self.invert_selection)
        invert_sel_btn.pack(side=LEFT, padx=2)
        self.create_tooltip(invert_sel_btn, "Inverte la selezione corrente")
        
        # Pulsanti per le azioni
        action_frame = ttk.Frame(actions_frame)
        action_frame.pack(side=RIGHT)

        self.copy_button = ttk.Button(action_frame, text="Copia selezionati",
                                    command=self.copy_selected,
                                    style="TButton")
        self.copy_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.copy_button, "Copia i file selezionati nella directory specificata")
        
        self.compress_button = ttk.Button(action_frame, text="Comprimi selezionati",
                                        command=self.compress_selected,
                                        style="TButton")
        self.compress_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.compress_button, "Comprimi i file selezionati in un archivio ZIP")

        self.view_log_button = ttk.Button(action_frame, text="Visualizza file esclusi",
                                        command=self.view_skipped_files_log,
                                        style="secondary.TButton")
        self.view_log_button.pack(side=LEFT, padx=5)
        self.create_tooltip(self.view_log_button, "Visualizza il log dei file esclusi dalla ricerca")

        # TreeView con scrollbar
        treeview_container = ttk.Frame(results_container)
        treeview_container.pack(fill=BOTH, expand=True)

        # Creazione della TreeView
        self.results_list = ttk.Treeview(treeview_container, selectmode="extended",
                                    columns=("type", "author", "size", "modified", "created", "path"),
                                    show="headings")

        # Creazione delle scrollbar
        vsb = ttk.Scrollbar(treeview_container, orient="vertical", command=self.results_list.yview)
        hsb = ttk.Scrollbar(treeview_container, orient="horizontal", command=self.results_list.xview)

        # Configurazione delle scrollbar nella TreeView
        self.results_list.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

        # Posizionamento con grid
        self.results_list.grid(column=0, row=0, sticky='nsew')
        vsb.grid(column=1, row=0, sticky='ns')
        hsb.grid(column=0, row=1, sticky='ew')

        # Configura grid layout per l'espansione
        treeview_container.grid_columnconfigure(0, weight=1)
        treeview_container.grid_rowconfigure(0, weight=1)

        # Impostazione delle intestazioni delle colonne
        self.results_list.heading("type", text="Tipo")
        self.results_list.heading("size", text="Dimensione")
        self.results_list.heading("modified", text="Modificato")
        self.results_list.heading("created", text="Creato")
        self.results_list.heading("author", text="Nome")  # Rinominato da "Autore" a "Nome"
        self.results_list.heading("path", text="Percorso")

        # Impostazione delle intestazioni delle colonne con funzione di ordinamento
        self.results_list.heading("type", text="Tipo", 
                            command=lambda: self.treeview_sort_column(self.results_list, "type", False))
        self.results_list.heading("size", text="Dimensione", 
                            command=lambda: self.treeview_sort_column(self.results_list, "size", False))
        self.results_list.heading("modified", text="Modificato", 
                            command=lambda: self.treeview_sort_column(self.results_list, "modified", False))
        self.results_list.heading("created", text="Creato", 
                            command=lambda: self.treeview_sort_column(self.results_list, "created", False))
        self.results_list.heading("author", text="Nome", 
                            command=lambda: self.treeview_sort_column(self.results_list, "author", False))
        self.results_list.heading("path", text="Percorso", 
                            command=lambda: self.treeview_sort_column(self.results_list, "path", False))

        # Impostazione delle larghezze fisse delle colonne
        self.results_list.column("type", width=120, minwidth=50, stretch=NO, anchor="center")
        self.results_list.column("size", width=100, minwidth=80, stretch=NO, anchor="center")
        self.results_list.column("modified", width=150, minwidth=120, stretch=NO, anchor="center")
        self.results_list.column("created", width=150, minwidth=120, stretch=NO, anchor="center")
        self.results_list.column("author", width=180, minwidth=80, stretch=NO, anchor="w")
        self.results_list.column("path", width=600, minwidth=200, stretch=YES, anchor="w")  # Solo questa colonna si espande

        # Aggiungi binding per l'evento di doppio clic
        self.results_list.bind("<Double-1>", self.open_file_location)

        # Applica stili alle righe
        self.update_theme_colors()

        # Auto-focus sull'entry del percorso all'avvio
        self.path_entry.focus_set()

    def show_advanced_options(self):
        """Mostra una finestra di dialogo unificata per tutte le opzioni avanzate"""
        dialog = ttk.Toplevel(self.root)
        dialog.title("Impostazioni avanzate")
        dialog.geometry("800x650")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Frame principale con padding
        main_frame = ttk.Frame(dialog, padding=10)
        main_frame.pack(fill=BOTH, expand=YES)
        
        # Crea un notebook con schede per le diverse categorie di opzioni
        notebook = ttk.Notebook(main_frame)
        notebook.pack(fill=BOTH, expand=YES, pady=10)
        
        # ================= Scheda 1: Opzioni di ricerca =================
        search_options_frame = ttk.Frame(notebook, padding=15)
        notebook.add(search_options_frame, text="Opzioni di ricerca")
        
        # Profondità di ricerca
        depth_frame = ttk.LabelFrame(search_options_frame, text="Profondità di ricerca", padding=10)
        depth_frame.pack(fill=X, pady=10)
        
        depth_desc = ttk.Label(depth_frame, 
                        text="Imposta il numero massimo di livelli di cartella da esplorare.\n"
                            "Il valore 0 indica profondità illimitata.",
                        wraplength=700)
        depth_desc.pack(anchor=W, pady=(0, 10))
        
        depth_control = ttk.Frame(depth_frame)
        depth_control.pack(fill=X)
        
        ttk.Label(depth_control, text="Profondità:").pack(side=LEFT)
        depth_spinbox = ttk.Spinbox(depth_control, from_=0, to=20, width=3, textvariable=self.max_depth)
        depth_spinbox.pack(side=LEFT, padx=5)
        self.depth_spinbox = depth_spinbox  # Salva riferimento
        depth_spinbox.set(self.max_depth)
        ttk.Label(depth_control, text="(0 = illimitata)", foreground="gray").pack(side=LEFT)
        
        # Contenuti da cercare
        content_frame = ttk.LabelFrame(search_options_frame, text="Contenuti da cercare", padding=10)
        content_frame.pack(fill=X, pady=10)
        
        ttk.Checkbutton(content_frame, text="Cerca nei file", variable=self.search_files).pack(anchor=W, pady=2)
        ttk.Checkbutton(content_frame, text="Cerca nelle cartelle", variable=self.search_folders).pack(anchor=W, pady=2)
        ttk.Checkbutton(content_frame, text="Cerca nei contenuti dei file", variable=self.search_content).pack(anchor=W, pady=2)
        ttk.Checkbutton(content_frame, text="Cerca parole intere", variable=self.whole_word_search).pack(anchor=W, pady=2)
        
        # ================= Scheda 2: Filtri avanzati =================
        filters_frame = ttk.Frame(notebook, padding=15)
        notebook.add(filters_frame, text="Filtri avanzati")
        
        # Filtri dimensione
        size_frame = ttk.LabelFrame(filters_frame, text="Dimensione file", padding=10)
        size_frame.pack(fill=X, pady=10)
        
        ttk.Label(size_frame, text="Filtra i file in base alla dimensione:").pack(anchor=W, pady=(0, 10))
        
        size_grid = ttk.Frame(size_frame)
        size_grid.pack(fill=X)
        
        ttk.Label(size_grid, text="Dimensione minima (KB):").grid(row=0, column=0, padx=5, pady=5, sticky=W)
        min_size = ttk.Entry(size_grid, width=10)
        min_size.grid(row=0, column=1, padx=5, pady=5)
        min_size.insert(0, str(self.advanced_filters["size_min"] // 1024))
        
        ttk.Label(size_grid, text="Dimensione massima (KB):").grid(row=1, column=0, padx=5, pady=5, sticky=W)
        max_size = ttk.Entry(size_grid, width=10)
        max_size.grid(row=1, column=1, padx=5, pady=5)
        max_size.insert(0, str(self.advanced_filters["size_max"] // 1024 if self.advanced_filters["size_max"] else 0))
        
        # Filtri data
        date_frame = ttk.LabelFrame(filters_frame, text="Data di modifica", padding=10)
        date_frame.pack(fill=X, pady=10)
        
        ttk.Label(date_frame, text="Filtra i file in base alla data di modifica:").pack(anchor=W, pady=(0, 10))
        
        date_grid = ttk.Frame(date_frame)
        date_grid.pack(fill=X)
        
        ttk.Label(date_grid, text="Data inizio (DD-MM-YYYY):").grid(row=0, column=0, padx=5, pady=5, sticky=W)
        min_date = ttk.DateEntry(date_grid, dateformat="%d-%m-%Y")
        min_date.grid(row=0, column=1, padx=5, pady=5)
        min_date.entry.delete(0, 'end')
        
        ttk.Label(date_grid, text="Data fine (DD-MM-YYYY):").grid(row=1, column=0, padx=5, pady=5, sticky=W)
        max_date = ttk.DateEntry(date_grid, dateformat="%d-%m-%Y")
        max_date.grid(row=1, column=1, padx=5, pady=5)
        max_date.entry.delete(0, 'end')
        
        # Aggiunto: Filtri per estensioni
        ext_frame = ttk.LabelFrame(filters_frame, text="Estensioni file", padding=10)
        ext_frame.pack(fill=X, pady=10)

        ttk.Label(ext_frame, text="Filtra i file per estensione (separate da virgola, es: .txt, .pdf):").pack(anchor=W, pady=(0, 10))

        # Inizializza con estensioni esistenti
        current_exts = ", ".join([ext.lstrip('.') for ext in self.advanced_filters.get("extensions", [])])
        extensions = ttk.Entry(ext_frame, width=50)
        extensions.pack(fill=X, padx=5)
        extensions.insert(0, current_exts)
        
        # ================= Scheda 3: Gestione esclusioni =================
        exclusions_frame = ttk.Frame(notebook, padding=15)
        notebook.add(exclusions_frame, text="Esclusioni")
        
        ttk.Label(exclusions_frame, text="Aggiungi cartelle da escludere dalla ricerca:", wraplength=700).pack(anchor=W, pady=(0, 10))
        
        # Lista dei percorsi esclusi
        list_frame = ttk.Frame(exclusions_frame)
        list_frame.pack(fill=BOTH, expand=NO, pady=5)
        
        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side=RIGHT, fill=Y)
        
        excluded_list = ttk.Treeview(list_frame, columns=("path",), show="headings", 
                                yscrollcommand=scrollbar.set, selectmode="extended")
        excluded_list.heading("path", text="Percorso")
        excluded_list.column("path", width=450)
        excluded_list.pack(fill=BOTH, expand=YES)
        
        scrollbar.config(command=excluded_list.yview)
        
        # Aggiungi i percorsi esclusi alla lista
        if hasattr(self, 'excluded_paths'):
            for path in self.excluded_paths:
                excluded_list.insert("", "end", values=(path,))
        
        # Frame per aggiungere nuovi percorsi
        add_frame = ttk.Frame(exclusions_frame)
        add_frame.pack(fill=X, pady=10)
        
        path_var = StringVar()
        path_entry = ttk.Entry(add_frame, textvariable=path_var, width=50)
        path_entry.pack(side=LEFT, padx=(0, 5), fill=X, expand=YES)
        
        def browse_exclude():
            directory = filedialog.askdirectory()
            if directory:
                path_var.set(directory)
        
        ttk.Button(add_frame, text="Sfoglia", command=browse_exclude).pack(side=LEFT, padx=5)
        
        def add_exclusion():
            path = path_var.get().strip()
            if path:
                excluded_list.insert("", "end", values=(path,))
                path_var.set("")
        
        ttk.Button(add_frame, text="Aggiungi", command=add_exclusion).pack(side=LEFT, padx=5)
        
        # Pulsante per rimuovere elementi selezionati
        def remove_selected():
            selected = excluded_list.selection()
            if selected:
                for item in selected:
                    excluded_list.delete(item)
        
        ttk.Button(exclusions_frame, text="Rimuovi selezionati", command=remove_selected).pack(anchor=W, pady=5)
        
        # ================= Scheda 4: Opzioni a blocchi =================
        blocks_frame = ttk.Frame(notebook, padding=15)
        notebook.add(blocks_frame, text="Opzioni a blocchi")
        
        desc_label = ttk.Label(blocks_frame, text="La ricerca a blocchi divide le cartelle in unità di lavoro più piccole per migliorare le prestazioni e la reattività.",
                        wraplength=700)
        desc_label.pack(fill=X, pady=(0, 15))
        
        # Dimensione blocchi
        size_frame = ttk.LabelFrame(blocks_frame, text="Dimensione e parallelismo", padding=10)
        size_frame.pack(fill=X, pady=10)
        
        size_grid = ttk.Frame(size_frame)
        size_grid.pack(fill=X)
        
        ttk.Label(size_grid, text="Max file per blocco:").grid(row=0, column=0, sticky=W, padx=5, pady=5)
        files_block = ttk.Spinbox(size_grid, from_=100, to=10000, increment=100, width=7, textvariable=self.max_files_per_block)
        files_block.grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(size_grid, text="Blocchi paralleli:").grid(row=0, column=2, sticky=W, padx=(20, 5), pady=5)
        parallel = ttk.Spinbox(size_grid, from_=1, to=16, increment=1, width=5, textvariable=self.max_parallel_blocks)
        parallel.grid(row=0, column=3, padx=5, pady=5)
        
        # Opzioni aggiuntive
        options_frame = ttk.LabelFrame(blocks_frame, text="Ottimizzazioni", padding=10)
        options_frame.pack(fill=X, pady=10)
        
        ttk.Checkbutton(options_frame, text="Adatta automaticamente la dimensione dei blocchi", 
                    variable=self.block_size_auto_adjust).pack(anchor=W, pady=2)
        
        ttk.Checkbutton(options_frame, text="Dare priorità alle cartelle utente", 
                    variable=self.prioritize_user_folders).pack(anchor=W, pady=2)
        
        # ================= Scheda 5: Performance =================
        performance_frame = ttk.Frame(notebook, padding=15)
        notebook.add(performance_frame, text="Performance")

        # Timeout e limiti
        timeout_frame = ttk.LabelFrame(performance_frame, text="Timeout e limiti", padding=10)
        timeout_frame.pack(fill=X, pady=10)

        # Rimuoviamo il posizionamento separato del checkbox e lo inseriamo direttamente nella griglia
        timeout_grid = ttk.Frame(timeout_frame)
        timeout_grid.pack(fill=X, pady=5)

        # Riga 0: Checkbox e secondi nella stessa riga
        timeout_check = ttk.Checkbutton(timeout_grid, text="Attiva timeout ricerca", variable=self.timeout_enabled)
        timeout_check.grid(row=0, column=0, sticky=W, padx=5, pady=2)

        ttk.Label(timeout_grid, text="Secondi:").grid(row=0, column=1, sticky=W, padx=(55, 5), pady=2)
        timeout_spin = ttk.Spinbox(timeout_grid, from_=10, to=3600, width=5, textvariable=self.timeout_seconds)
        timeout_spin.grid(row=0, column=2, padx=5, pady=2, sticky=W)
        
        ttk.Label(timeout_grid, text="Max file da controllare:").grid(row=1, column=0, sticky=W, padx=5, pady=5)
        max_files = ttk.Spinbox(timeout_grid, from_=1000, to=10000000, width=8, textvariable=self.max_files_to_check)
        max_files.grid(row=1, column=1, padx=5, pady=5, sticky=W)
        
        ttk.Label(timeout_grid, text="Max risultati:").grid(row=1, column=2, sticky=W, padx=5, pady=5)
        max_results = ttk.Spinbox(timeout_grid, from_=500, to=100000, width=8, textvariable=self.max_results)
        max_results.grid(row=1, column=3, padx=5, pady=5, sticky=W)
        
        # Processamento
        process_frame = ttk.LabelFrame(performance_frame, text="Processamento", padding=10)
        process_frame.pack(fill=X, pady=10)
        
        process_grid = ttk.Frame(process_frame)
        process_grid.pack(fill=X)
        
        ttk.Label(process_grid, text="Thread paralleli:").grid(row=0, column=0, sticky=W, padx=5, pady=5)
        threads = ttk.Spinbox(process_grid, from_=1, to=16, width=3, textvariable=self.worker_threads)
        threads.grid(row=0, column=1, padx=5, pady=5, sticky=W)
        
        ttk.Label(process_grid, text="Dimensione max file (MB):").grid(row=0, column=2, sticky=W, padx=5, pady=5)
        max_size = ttk.Spinbox(process_grid, from_=1, to=1000, width=5, textvariable=self.max_file_size_mb)
        max_size.grid(row=0, column=3, padx=5, pady=5, sticky=W)
        
        # Calcolo dimensioni
        calc_frame = ttk.LabelFrame(performance_frame, text="Calcolo dimensioni", padding=10)
        calc_frame.pack(fill=X, pady=10)
        
        ttk.Label(calc_frame, text="Modalità di calcolo:").pack(side=LEFT, padx=5)
        calc_combo = ttk.Combobox(calc_frame, textvariable=self.dir_size_calculation, 
                            values=["incrementale", "preciso", "stimato", "sistema", "disabilitato"], 
                            width=12, state="readonly")
        calc_combo.pack(side=LEFT, padx=5)
        
        # Pulsanti finali per la finestra
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=X, pady=(15, 0))
        
        def save_options():
            # Aggiorna max_depth dal valore dello spinbox
            try:
                self.max_depth = int(depth_spinbox.get())
            except ValueError:
                self.max_depth = 0
            
            # Aggiorna le impostazioni degli advanced filters
            try:
                min_kb = int(min_size.get() or 0)
                max_kb = int(max_size.get() or 0)
                self.advanced_filters["size_min"] = min_kb * 1024
                self.advanced_filters["size_max"] = max_kb * 1024
                
                min_date_val = min_date.entry.get().strip()
                max_date_val = max_date.entry.get().strip()
                
                self.advanced_filters["date_min"] = min_date_val
                self.advanced_filters["date_max"] = max_date_val
                
                exts = [e.strip() for e in extensions.get().split(",") if e.strip()]
                self.advanced_filters["extensions"] = [f".{e.lstrip('.')}" for e in exts]
            except ValueError as e:
                messagebox.showerror("Errore", f"Valore non valido: {str(e)}")
                return
                
            # Aggiorna i percorsi esclusi
            self.excluded_paths = []
            for item in excluded_list.get_children():
                values = excluded_list.item(item)["values"]
                if values:
                    self.excluded_paths.append(values[0])
            
            dialog.destroy()
        
        ttk.Button(btn_frame, text="Ripristina valori predefiniti", 
            command=lambda: messagebox.showinfo("Info", "Funzione non implementata")).pack(side=LEFT)
        ttk.Button(btn_frame, text="Annulla", command=dialog.destroy).pack(side=RIGHT, padx=5)
        ttk.Button(btn_frame, text="Salva", command=save_options).pack(side=RIGHT, padx=5)
        
        # Centra la finestra
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (dialog.winfo_screenheight() // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")

    def create_tooltip(self, widget, text, delay=500, fade=True):
        """Crea tooltip con ritardo, effetti di dissolvenza e larghezza automatica"""
        
        tooltip_timer = None
        fade_timer = None
        
        def show_tooltip():
            x = widget.winfo_rootx() + 25
            y = widget.winfo_rooty() + 25
            
            tooltip = ttk.Toplevel(widget)
            tooltip.wm_overrideredirect(True)
            tooltip.attributes("-alpha", 0.0)  # Inizia invisibile
            tooltip.attributes("-topmost", True)  # Mantiene sopra altre finestre
            
            # Crea un frame con bordo e padding
            frame = ttk.Frame(tooltip, borderwidth=1, relief="solid", padding=10)
            frame.pack(fill="both", expand=True)
            
            # Calcolo della larghezza del testo
            font = ("Segoe UI", 9)  # Puoi modificare questo per usare un altro font
            
            # Calcola la lunghezza approssimativa del testo
            temp_label = tk.Label(font=font)
            
            # Determina se il testo necessita di essere diviso in più righe
            lines = text.split("\n")
            max_line_width = 0
            
            for line in lines:
                temp_label.configure(text=line)
                line_width = temp_label.winfo_reqwidth()
                max_line_width = max(max_line_width, line_width)
            
            # Limita la larghezza massima a 500 pixel
            max_width = min(max_line_width, 500)
            
            wraplength = max_width
                
            # Crea l'etichetta con il testo
            label = ttk.Label(frame, text=text, justify="left", 
                            wraplength=wraplength, font=font)
            label.pack(fill="both", expand=True)
            
            # Distruggi il label temporaneo
            temp_label.destroy()
            
            # Aggiorna immediatamente per calcolare le dimensioni
            tooltip.update_idletasks()
            
            # Posiziona il tooltip in modo che non vada fuori dallo schermo
            tooltip_width = tooltip.winfo_reqwidth()
            tooltip_height = tooltip.winfo_reqheight()
            
            screen_width = tooltip.winfo_screenwidth()
            screen_height = tooltip.winfo_screenheight()
            
            # Aggiusta la posizione se il tooltip esce dallo schermo
            if x + tooltip_width > screen_width:
                x = screen_width - tooltip_width - 10
            
            if y + tooltip_height > screen_height:
                y = screen_height - tooltip_height - 10
                
            # Imposta la posizione finale
            tooltip.wm_geometry(f"+{x}+{y}")
            
            widget._tooltip = tooltip
            
            if fade:
                # Effetto dissolvenza in entrata
                def fade_in(alpha=0.0):
                    if not hasattr(widget, "_tooltip"):
                        return
                    
                    tooltip.attributes("-alpha", alpha)
                    if alpha < 1.0:
                        nonlocal fade_timer
                        fade_timer = widget.after(20, lambda: fade_in(alpha + 0.1))
                
                fade_in()
        
        def enter(event):
            nonlocal tooltip_timer
            # Avvia il timer per mostrare il tooltip dopo un certo ritardo
            tooltip_timer = widget.after(delay, show_tooltip)
        
        def leave(event):
            nonlocal tooltip_timer, fade_timer
            
            # Cancella il timer se esiste
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
                tooltip_timer = None
            
            # Cancella il timer di dissolvenza se esiste
            if fade_timer:
                widget.after_cancel(fade_timer)
                fade_timer = None
            
            # Rimuovi il tooltip se esiste
            if hasattr(widget, "_tooltip"):
                if fade:
                    # Effetto dissolvenza in uscita
                    def fade_out(alpha=1.0):
                        if alpha <= 0 or not hasattr(widget, "_tooltip"):
                            if hasattr(widget, "_tooltip"):
                                widget._tooltip.destroy()
                                del widget._tooltip
                        else:
                            widget._tooltip.attributes("-alpha", alpha)
                            nonlocal fade_timer
                            fade_timer = widget.after(20, lambda: fade_out(alpha - 0.1))
                    
                    fade_out()
                else:
                    widget._tooltip.destroy()
                    del widget._tooltip
        
        # Gestisce anche il caso in cui il widget venga distrutto
        def on_destroy(event):
            nonlocal tooltip_timer, fade_timer
            if tooltip_timer:
                widget.after_cancel(tooltip_timer)
            if fade_timer and hasattr(widget, "_tooltip"):
                widget.after_cancel(fade_timer)
                widget._tooltip.destroy()
        
        widget.bind("<Enter>", enter)
        widget.bind("<Leave>", leave)
        widget.bind("<Destroy>", on_destroy)

    def select_all(self):
        self.results_list.selection_set(self.results_list.get_children())
        
    def deselect_all(self):
        self.results_list.selection_remove(self.results_list.get_children())
        
    def invert_selection(self):
        all_items = self.results_list.get_children()
        selected_items = self.results_list.selection()
        self.results_list.selection_remove(selected_items)
        to_select = set(all_items) - set(selected_items)
        for item in to_select:
            self.results_list.selection_add(item)
            
    def treeview_sort_column(self, tv, col, reverse):
        """Ordina il TreeView in base alla colonna cliccata"""
        # Ottieni la lista di item con i loro valori
        l = [(tv.set(k, col), k) for k in tv.get_children('')]
        
        # Determina il tipo di ordinamento in base alla colonna
        try:
            if col == "size":
                # Gestione speciale per le dimensioni (KB, MB, GB)
                def extract_size(size_str):
                    if not size_str:
                        return 0
                    try:
                        # Estrai il valore numerico e converti tutto in byte per confrontare correttamente
                        value = float(size_str.split()[0])
                        if "KB" in size_str:
                            return value * 1024
                        elif "MB" in size_str:
                            return value * 1024 * 1024
                        elif "GB" in size_str:
                            return value * 1024 * 1024 * 1024
                        else:
                            return value  # Assumiamo che sia in byte
                    except:
                        return 0
                        
                l.sort(key=lambda x: extract_size(x[0]), reverse=reverse)
            elif col in ("modified", "created"):
                # Gestione per le date (assumendo formato DD/MM/YYYY HH:MM)
                from datetime import datetime
                def parse_date(date_str):
                    try:
                        if date_str and date_str != "N/A":
                            return datetime.strptime(date_str, "%d/%m/%Y %H:%M")
                        return datetime(1900, 1, 1)  # Data di default per valori vuoti
                    except:
                        return datetime(1900, 1, 1)
                        
                l.sort(key=lambda x: parse_date(x[0]), reverse=reverse)
            else:
                # Ordinamento standard alfanumerico
                l.sort(reverse=reverse)
        except Exception as e:
            self.log_debug(f"Errore durante l'ordinamento: {str(e)}")
            # Fallback all'ordinamento alfanumerico semplice
            l.sort(reverse=reverse)
            
        # Riorganizza gli elementi nel TreeView
        for index, (val, k) in enumerate(l):
            tv.move(k, '', index)

        # Memorizza la colonna ordinata e la direzione per il prossimo click
        tv.heading(col, command=lambda _col=col: self.treeview_sort_column(tv, _col, not reverse))
        
        # Aggiorna l'indicatore visivo di ordinamento (aggiunge freccia all'intestazione)
        for c in tv["columns"]:
            if c == col:
                tv.heading(c, text=f"{tv.heading(c, 'text').split(' ')[0]} {'▼' if reverse else '▲'}")
            else:
                # Rimuovi eventuali indicatori di ordinamento da altre colonne
                tv.heading(c, text=tv.heading(c, 'text').split(' ')[0])
                
# Funzione principale per eseguire l'applicazione
def main():
    import sys
    
    # Crea la finestra principale con il tema desiderato
    root = ttk.Window(themename="darkly")
    root.withdraw()  # Nascondi completamente la finestra durante l'inizializzazione
    
    # Crea la schermata di splash
    splash = create_splash_screen(root)
    
    # Inizializza l'applicazione (ma l'interfaccia rimane nascosta)
    app = FileSearchApp(root)
    
    # Controlla se ci sono argomenti da linea di comando
    if len(sys.argv) > 1:
        app.search_path.set(sys.argv[1])
    
    # Funzione per completare l'avvio e mostrare la finestra principale
    def finish_startup():
        splash.destroy()
        # Configura le dimensioni della finestra prima di mostrarla
        root.geometry("1300x850")
        # Mostra la finestra completamente costruita
        root.deiconify()
        
    # Completa l'avvio dopo un breve ritardo
    root.after(1500, finish_startup)
    
    root.mainloop()

def create_splash_screen(parent):
    splash_win = tk.Toplevel(parent)
    splash_win.title("")
    splash_win.overrideredirect(True)
    splash_win.attributes("-topmost", True)
    
    # Dimensioni dello splash
    width, height = 500, 250
    screen_width = splash_win.winfo_screenwidth()
    screen_height = splash_win.winfo_screenheight()
    x = (screen_width // 2) - (width // 2)
    y = (screen_height // 2) - (height // 2)
    splash_win.geometry(f"{width}x{height}+{x}+{y}")
    
    # Contenuto dello splash
    frame = ttk.Frame(splash_win, padding=20)
    frame.pack(fill=tk.BOTH, expand=tk.YES)
    
    ttk.Label(frame, text="File Search Tool V9.2 Beta", 
            font=("Helvetica", 18, "bold")).pack(pady=(10, 5))
    ttk.Label(frame, text="Forensics G.di F.", 
            font=("Helvetica", 14)).pack(pady=(0, 20))
    ttk.Label(frame, text="Caricamento applicazione in corso...").pack()
    
    progress = ttk.Progressbar(frame, mode="indeterminate")
    progress.pack(fill=tk.X, pady=10)
    progress.start(10)
    
    return splash_win
if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        import traceback
        with open("error_log.txt", "w", encoding="utf-8") as f:
            f.write(f"Si è verificato un errore: {str(e)}\n")
            f.write(traceback.format_exc())
        print(f"Si è verificato un errore: {str(e)}")
        
        # Se stai usando tkinter, puoi anche mostrare un messaggio all'utente
        try:
            import tkinter as tk
            from tkinter import messagebox
            root = tk.Tk()
            root.withdraw()
            messagebox.showerror("Errore durante l'avvio", 
                f"Si è verificato un errore durante l'avvio dell'applicazione.\n\n"
                f"Errore: {str(e)}\n\n"
                f"I dettagli sono stati salvati nel file error_log.txt")
        except:
            pass  # Se anche la visualizzazione del messaggio fallisce, continua
